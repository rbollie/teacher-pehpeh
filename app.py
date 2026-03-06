"""
TEACHER PEHPEH BY IBT - AI-Powered Support for Every Classroom
Institute of Basic Technology (IBT)
Built by Rodney L. Bollie, PhD
"""
import streamlit as st
import time, os, base64, json, random, zlib
from urllib.parse import quote as urlquote
try:
    import pandas as pd; PD=True
except: PD=False
try:
    import qrcode; QR=True
except: QR=False
try:
    from PIL import Image, ImageDraw, ImageFont; PIL_OK=True
except: PIL_OK=False

# === API KEYS ===
try:
    from dotenv import load_dotenv
    load_dotenv()  # loads .env file from project folder
except ImportError:
    pass  # Running on Streamlit Cloud — uses st.secrets instead

def _get_key(name):
    """Get API key from env vars OR Streamlit secrets, strip quotes/whitespace."""
    v = os.environ.get(name, "")
    if not v:
        try:
            v = st.secrets.get(name, "")
        except: pass
    if not v:
        # Try lowercase and variations in st.secrets
        try:
            for k in st.secrets:
                if k.upper() == name.upper():
                    v = st.secrets[k]
                    break
        except: pass
    return str(v).strip().strip("'\"")

OPENAI_API_KEY = _get_key("OPENAI_API_KEY")
ANTHROPIC_API_KEY = _get_key("ANTHROPIC_API_KEY")
GOOGLE_API_KEY = _get_key("GOOGLE_API_KEY")
ELEVENLABS_API_KEY = _get_key("ELEVENLABS_API_KEY")
LOGO_FILENAME = "logo.png"

# ═══════════════════════════════════════════════════════════
# SUBSCRIPTION & LICENSE SYSTEM
# ═══════════════════════════════════════════════════════════
import datetime

def _check_subscription():
    """
    Returns (is_active, days_left, tier)
    Reads LICENSE_KEY and LICENSE_EXPIRY from secrets/env.
    If not set, defaults to free tier.
    """
    expiry_str = _get_key("LICENSE_EXPIRY")   # format: YYYY-MM-DD  e.g. 2026-12-31
    license_key = _get_key("LICENSE_KEY")      # any non-empty string = paid school
    if not license_key or not expiry_str:
        return False, 0, "free"
    try:
        expiry = datetime.date.fromisoformat(expiry_str.strip())
        today  = datetime.date.today()
        days_left = (expiry - today).days
        if days_left >= 0:
            return True, days_left, "paid"
        else:
            return False, days_left, "expired"
    except Exception:
        return False, 0, "free"

SUBSCRIPTION_ACTIVE, _SUB_DAYS_LEFT, _SUB_TIER = _check_subscription()

# ═══════════════════════════════════════════════════════════
# LOGIN SYSTEM
# Schools get a SCHOOL_CODE and SCHOOL_PASS in their secrets.
# IBT admin can also set ADMIN_PASS to access all schools.
# ═══════════════════════════════════════════════════════════

def _get_credentials():
    """Return list of valid (username, password, label) tuples from secrets."""
    creds = []
    # Primary school account
    sc = _get_key("SCHOOL_CODE")
    sp = _get_key("SCHOOL_PASS")
    sl = _get_key("SCHOOL_NAME_LOGIN") or "School Account"
    if sc and sp:
        creds.append((sc.strip(), sp.strip(), sl.strip()))
    # Admin / IBT account
    ap = _get_key("ADMIN_PASS")
    if ap:
        creds.append(("ibt_admin", ap.strip(), "IBT Admin"))
    # If no credentials configured at all, use a simple open-access mode
    return creds

def _login_required():
    """True if the app has login credentials configured."""
    return len(_get_credentials()) > 0

def _is_logged_in():
    return st.session_state.get("_logged_in", False)

def _show_login_screen():
    """Render the login screen and return True once authenticated."""
    b = get_b64()
    if b:
        st.markdown(
            f'<div style="text-align:center;padding:2rem 0 1rem">'+
            f'<img src="data:image/png;base64,{b}" style="max-height:200px;filter:drop-shadow(0 4px 16px rgba(212,168,67,.4))"></div>',
            unsafe_allow_html=True)
    st.markdown(
        '<h2 style="text-align:center;color:#D4A843;margin-bottom:.2rem">Teacher Pehpeh by IBT</h2>'+
        '<p style="text-align:center;color:#8899BB;font-size:.9rem;margin-bottom:1.5rem">Enter your school login to continue</p>',
        unsafe_allow_html=True)

    with st.form("login_form", clear_on_submit=False):
        username = st.text_input("School Code / Username", placeholder="e.g. grandfield_hs")
        password = st.text_input("Password", type="password", placeholder="Your school password")
        submitted = st.form_submit_button("🔐 Sign In", use_container_width=True, type="primary")

    if submitted:
        creds = _get_credentials()
        match = next((c for c in creds if c[0].lower()==username.strip().lower() and c[1]==password.strip()), None)
        if match:
            st.session_state["_logged_in"] = True
            st.session_state["_login_label"] = match[2]
            st.rerun()
        else:
            st.error("Incorrect username or password. Please contact IBT for access.")

    st.markdown(
        '<p style="text-align:center;font-size:.8rem;color:#556;margin-top:2rem">'+
        'Don\'t have a login? '+
        '<a href="https://www.institutebasictechnology.org" target="_blank" style="color:#D4A843">Contact IBT</a> to get your school set up.</p>',
        unsafe_allow_html=True)

def _show_upgrade_prompt(context="generate"):
    """Show the free-tier upgrade wall in place of AI generation."""
    st.markdown(
        f'''<div style="border:1px solid #D4A84366;border-radius:12px;padding:1.2rem 1.4rem;background:rgba(212,168,67,.05);text-align:center;margin:1rem 0">'+
        f'<p style="font-size:1.1rem;color:#D4A843;font-weight:700;margin-bottom:.4rem">🔒 AI Generation — Paid Feature</p>'+
        f'<p style="color:#9AAABB;font-size:.9rem;margin-bottom:.8rem">Your school\'s subscription has expired or is not yet active.<br>'+
        f'Everything you\'ve built is still here — your student records, quizzes, and past content.</p>'+
        f'<p style="color:#D4A843;font-size:.85rem">📲 WhatsApp IBT to renew: '+
        f'<a href="https://wa.me/message/IBTWHALINK" style="color:#D4A843;font-weight:600">Get in touch</a> &nbsp;|&nbsp; '+
        f'<a href="https://www.institutebasictechnology.org" style="color:#D4A843">Visit our website</a></p>'+
        f'</div>''',
        unsafe_allow_html=True)


try:
    import openai; OAI = True
except ImportError: OAI = False
try:
    import anthropic; ANT = True
except ImportError: ANT = False
try:
    import google.generativeai as genai; GEM = True
except ImportError: GEM = False

# === CURRICULUM — Load Ministry of Education data ===
import sys
from pathlib import Path
APP_DIR = Path(__file__).parent
if str(APP_DIR) not in sys.path:
    sys.path.insert(0, str(APP_DIR))
try:
    from curriculum import (
        load_all_curricula, get_grade_topics, get_topic_details,
        build_curriculum_context, get_curriculum_summary, get_available_subjects,
    )
    CURRICULUM_AVAILABLE = True
except ImportError:
    CURRICULUM_AVAILABLE = False

CURRICULA = {}
if CURRICULUM_AVAILABLE:
    CURRICULA = load_all_curricula()

# === MANO LANGUAGE — Load Mano language library for bilingual lessons ===
try:
    from mano_context import (
        build_mano_prompt_context, get_mano_preview, get_mano_stats,
        match_vocabulary, MANO_AVAILABLE,
    )
except ImportError:
    MANO_AVAILABLE = False

# === IBT COLORS (from website) ===
C_NAVY = "#0F2247"
C_NAVY_L = "#1A2744"
C_BLUE = "#2B7DE9"
C_BLUE_D = "#1D5CBF"
C_RED = "#8B1A1A"
C_RED_L = "#B22234"
C_GOLD = "#D4A843"
C_GOLD_L = "#F5D98E"

# === KNOWLEDGE BASE (assembled at runtime only) ===
@st.cache_data
def _kb():
    return ("IBT RESEARCH (183 students, 6 Liberian schools, 4 STEM subjects):\n"
        "Overall avg 0.433(C-). Chem 0.494(C), Physics 0.399(C-), Math 0.391(C-), Bio 0.447(C). "
        "Scale: 1.0=A,0.75=B+,0.625=B,0.50=B-,0.44=C,0.375=C-,0.25=D.\n"
        "MOTHER'S EDUCATION: Mom HS Grad avg 0.449 vs No HS 0.418 (p=0.031). Physics gap: HS 0.438 vs NoHS 0.363 (p=0.0075). Significant across ALL subgroups.\n"
        "SINGLE MOTHERS (22%): SM HS 0.457 vs SM NoHS 0.399 (14.7% gap). Physics: 0.461 vs 0.310=48.8% gap (p=0.006). SM+NoHS+4kids: Physics 0.283(D). 29% work after school.\n"
        "DIGITAL: 58.5% never used computer. SM NoHS: 81% never. 100% of SM NoHS users study only. +0.041 boost.\n"
        "SCHOOL: #1 predictor (F=8.60 p<0.001). Best 0.512(B-), worst 0.354(D+). 16x parent edu effect.\n"
        "INTERVENTION: Gap widens +0.055/2yr without. Narrows to 0.024 with. Physics most sensitive +0.09/yr.")

# === CONNECTIVITY ===
def email_result(content, subject, key_suffix, container=None):
    """Render email controls with IBT branded template"""
    ctx = container or st
    with ctx.expander(T("email_result"), expanded=False):
        to_addr = st.text_input("To:", placeholder="teacher@school.edu", key=f"email_to_{key_suffix}")
        cc_addr = st.text_input("CC (optional):", placeholder="principal@school.edu", key=f"email_cc_{key_suffix}")
        clean = content.replace("\n\n", "\n").strip()
        # IBT branded plain text template
        branded_text = f"""═══════════════════════════════════════════
   TEACHER PEHPEH by IBT
   Institute of Basic Technology
═══════════════════════════════════════════

{subject}
───────────────────────────────────────────

{clean}

───────────────────────────────────────────
📋 Generated by Teacher Pehpeh — AI-Powered
   Teaching Assistant by IBT

🌐 App: https://teacher-pehpeh.streamlit.app
🏫 IBT: www.institutebasictechnology.org
📧 Contact: info@institutebasictechnology.org

   "Curating Personalized Content to Support
    Underresourced Teachers"

   Powered by Claude · ChatGPT · Gemini
═══════════════════════════════════════════"""
        # IBT branded HTML template for file attachment
        branded_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
body{{font-family:Arial,sans-serif;margin:0;padding:0;background:#f4f4f4}}
.container{{max-width:700px;margin:20px auto;background:white;border-radius:12px;overflow:hidden;box-shadow:0 2px 12px rgba(0,0,0,.1)}}
.header{{background:linear-gradient(135deg,#B22234,#8B1A1A);padding:24px 30px;text-align:center}}
.header h1{{color:#D4A843;margin:0;font-size:22px;letter-spacing:1px}}
.header p{{color:#F0D5D5;margin:6px 0 0;font-size:13px}}
.subj{{background:#0F2247;color:#D4A843;padding:14px 30px;font-size:15px;font-weight:bold}}
.body{{padding:24px 30px;color:#1a1a2e;line-height:1.8;font-size:14px;white-space:pre-wrap}}
.footer{{background:#0F2247;padding:20px 30px;text-align:center;color:#8899BB;font-size:12px}}
.footer a{{color:#D4A843;text-decoration:none}}
.footer .brand{{color:#D4A843;font-weight:bold;font-size:14px;margin-bottom:8px;display:block}}
</style></head><body>
<div class="container">
<div class="header">
<h1>TEACHER PEHPEH</h1>
<p>Institute of Basic Technology (IBT)</p>
</div>
<div class="subj">{subject}</div>
<div class="body">{clean}</div>
<div class="footer">
<span class="brand">Teacher Pehpeh by IBT</span>
<a href="https://teacher-pehpeh.streamlit.app">Open Teacher Pehpeh</a> &nbsp;·&nbsp;
<a href="https://www.institutebasictechnology.org">🌐 IBT Website</a><br><br>
<em>"Curating Personalized Content to Support Underresourced Teachers"</em><br>
Powered by Claude · ChatGPT · Gemini
</div>
</div>
</body></html>"""
        # Mailto uses truncated plain text
        max_body = 1800
        mail_body = branded_text if len(branded_text) <= max_body else branded_text[:max_body] + "\n\n[Content truncated — see attached file for full version]"
        encoded_subject = urlquote(f"{subject}")
        encoded_body = urlquote(mail_body)
        cc_param = f"&cc={urlquote(cc_addr)}" if cc_addr.strip() else ""
        mailto_url = f"mailto:{urlquote(to_addr)}?subject={encoded_subject}{cc_param}&body={encoded_body}"
        c1,c2,c3 = st.columns(3)
        with c1:
            st.markdown(f'<a href="{mailto_url}" target="_blank" style="display:inline-block;background:#2B7DE9;color:white;padding:8px 16px;border-radius:8px;text-decoration:none;font-weight:600;font-size:.85rem;text-align:center;width:100%">📧 Email App</a>', unsafe_allow_html=True)
        with c2:
            st.download_button("📎 Text file", data=branded_text, file_name=f"{subject[:35].replace(' ','_')}.txt", key=f"email_dl_{key_suffix}")
        with c3:
            st.download_button("📎 HTML file", data=branded_html, file_name=f"{subject[:35].replace(' ','_')}.html", mime="text/html", key=f"email_html_{key_suffix}")
        if len(branded_text) > max_body:
            st.info("💡 Email body was truncated. Attach the downloaded file for the full version.")

def check_conn():
    import urllib.request, datetime
    r = {"online":False,"quality":"none","latency_ms":None,"label":"No Internet","emoji":"🔴","checked_at":datetime.datetime.now().strftime("%H:%M:%S")}
    lats = []
    for u in ["https://api.anthropic.com","https://api.openai.com","https://www.google.com"]:
        try:
            t=time.time(); req=urllib.request.Request(u,method="HEAD"); req.add_header("User-Agent","TP/1.0")
            urllib.request.urlopen(req,timeout=4); lats.append((time.time()-t)*1000)
        except: pass
    if not lats: return r
    a=sum(lats)/len(lats); r.update(online=True,latency_ms=round(a))
    if a<300: r.update(quality="high",label="Strong (server→API)",emoji="🟢")
    elif a<800: r.update(quality="medium",label="Moderate (server→API)",emoji="🟡")
    elif a<2000: r.update(quality="low",label="Slow (server→API)",emoji="🟠")
    else: r.update(quality="very_low",label="Very Slow",emoji="🟠")
    return r

# === IMAGE GENERATION ===
def gen_image(prompt):
    """Try DALL-E first, then Google Imagen as fallback"""
    img_style = "Clean, professional educational infographic style. Bright colors on white background. NO chalkboard, NO hand-drawn sketches, NO chalk-style text. Use clear labels, simple diagrams, and modern flat design. Culturally relevant to West Africa/Liberia."
    # Try DALL-E
    if OAI and OPENAI_API_KEY:
        try:
            c=openai.OpenAI(api_key=OPENAI_API_KEY)
            r=c.images.generate(model="dall-e-3",prompt=f"Educational visual aid: {prompt}. {img_style}",size="1024x1024",quality="standard",n=1)
            return r.data[0].url,"DALL-E"
        except: pass
    # Try Google Imagen
    if GEM and GOOGLE_API_KEY:
        try:
            genai.configure(api_key=GOOGLE_API_KEY)
            from google.generativeai import types as gtypes
            img_model=genai.GenerativeModel("gemini-2.5-flash")
            r=img_model.generate_content(
                f"Generate an educational visual aid: {prompt}. {img_style}",
                generation_config=genai.types.GenerationConfig(response_mime_type="text/plain")
            )
            client=genai.ImageGenerationModel("imagen-3.0-generate-002")
            response=client.generate_images(prompt=f"Educational visual aid: {prompt}. {img_style}",number_of_images=1)
            if response.images:
                import base64
                img_bytes=response.images[0]._image_bytes
                b64=base64.b64encode(img_bytes).decode()
                return f"data:image/png;base64,{b64}","Imagen"
        except: pass
    return None,None

# === TEXT TO SPEECH ===

def speak_elevenlabs(text, voice_id="VU4qoZUtDRUWXB09nrAd", model_id="eleven_flash_v2_5"):
    """Generate speech using ElevenLabs streaming API. Returns base64 audio or None."""
    if not ELEVENLABS_API_KEY: return None, "No ElevenLabs API key"
    try:
        import re, urllib.request, urllib.error, json as jlib
        clean = str(text)
        clean = re.sub(r'<[^>]+>', '', clean)
        clean = re.sub(r'[#*_`~\[\](){}|>]', '', clean)
        clean = re.sub(r'-{3,}', ' ', clean)
        clean = re.sub(r'={3,}', ' ', clean)
        for old, new in [('\u2014', '-'), ('\u2013', '-'), ('\u2018', "'"), ('\u2019', "'"), ('\u201c', '"'), ('\u201d', '"'), ('\u2026', '...'), ('\u2022', ','), ('\u00a0', ' ')]:
            clean = clean.replace(old, new)
        clean = re.sub(r'\n{2,}', '. ', clean)
        clean = re.sub(r'\n', '. ', clean)
        clean = re.sub(r' {2,}', ' ', clean)
        clean = re.sub(r'\.{2,}', '.', clean)
        clean = re.sub(r'\.\s*\.', '.', clean)
        clean = clean.strip()
        if not clean or len(clean) < 10: return None, "Text too short after cleaning"
        if len(clean) > 1200: clean = clean[:1200] + ". See the full written version for complete details."
        # Use streaming endpoint for faster first-byte
        url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}/stream"
        payload = jlib.dumps({
            "text": clean,
            "model_id": model_id,
            "output_format": "mp3_22050_32",
        }).encode("utf-8")
        req = urllib.request.Request(url, data=payload, headers={
            "xi-api-key": ELEVENLABS_API_KEY,
            "Content-Type": "application/json",
            "Accept": "audio/mpeg",
        })
        # Stream response in chunks
        chunks = []
        try:
            with urllib.request.urlopen(req, timeout=120) as resp:
                while True:
                    chunk = resp.read(4096)
                    if not chunk: break
                    chunks.append(chunk)
        except urllib.error.HTTPError as he:
            _klen = len(ELEVENLABS_API_KEY)
            _kpre = ELEVENLABS_API_KEY[:4] + "..." if _klen > 4 else "(empty)"
            if he.code == 401:
                return None, f"Voice error: 401 Unauthorized. Key length={_klen}, starts with '{_kpre}'. Go to elevenlabs.io → Profile+API Key → regenerate a new key. Then update it in your Streamlit secrets.toml or environment."
            elif he.code == 429:
                return None, "Voice error: Rate limit or quota exceeded. Check your ElevenLabs plan usage."
            else:
                return None, f"Voice error: HTTP {he.code}"
        audio_bytes = b"".join(chunks)
        if len(audio_bytes) < 100: return None, "Empty audio response"
        b64 = base64.b64encode(audio_bytes).decode()
        return b64, "Teacher Pehpeh Voice"
    except Exception as e:
        return None, f"Voice error: {e}"

def highlight_result(text):
    import re
    t = text.strip()
    # Collapse excess blank lines
    t = re.sub(r'\n{3,}', '\n\n', t)
    # Bold
    t = re.sub(r'\*\*(.+?)\*\*', r'<strong style="color:#D4A843;background:rgba(212,168,67,.1);padding:1px 4px;border-radius:3px">\1</strong>', t)
    # Headings
    t = re.sub(r'^### (.+)$', r'<div style="color:#2B7DE9;font-weight:700;font-size:.95rem;margin:8px 0 2px;font-family:Playfair Display,serif">\1</div>', t, flags=re.MULTILINE)
    t = re.sub(r'^## (.+)$', r'<div style="color:#D4A843;font-weight:700;font-size:1rem;margin:10px 0 2px;font-family:Playfair Display,serif">\1</div>', t, flags=re.MULTILINE)
    t = re.sub(r'^# (.+)$', r'<div style="color:#D4A843;font-weight:700;font-size:1.05rem;margin:10px 0 2px;font-family:Playfair Display,serif">\1</div>', t, flags=re.MULTILINE)
    # Numbered list items
    t = re.sub(r'^(\d+)\.\s', r'<strong style="color:#2B7DE9">\1.</strong> ', t, flags=re.MULTILINE)
    # Keywords
    for kw in ["WASSCE","BECE","Key Point","Note:","Tip:","Important:","Answer Key","Teacher's Guide","Objective","Assessment"]:
        t = t.replace(kw, f'<span style="color:#D4A843;font-weight:600">{kw}</span>')
    # Split into paragraphs, wrap each in <p> with tight margin
    parts = re.split(r'\n\n+', t)
    out = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        # Already an HTML block element (div/h3/h4)? Don't wrap in <p>
        if part.startswith('<div') or part.startswith('<h'):
            lines = part.split('\n')
            out.append('\n'.join(lines))
        else:
            lines = part.split('\n')
            inner = '<br>'.join(line for line in lines)
            out.append(f'<p style="margin:2px 0 6px;line-height:1.55">{inner}</p>')
    return '\n'.join(out)

def clean_parent_output(text):
    import re
    raw_lines = text.split('\n')
    clean = []
    skip = False
    for line in raw_lines:
        lo = line.strip().lower()
        if any(lo.startswith(p) for p in [
            '# parent communication', '## key features', '### key features',
            '## why this works', '### why this works', '## analysis', '### analysis',
            '| element', '|---', '| **', '## message', '### message quality',
            '[spoken duration', 'key elements included:', '## format', '### format',
            '## tone', '# sms', '## sms', '# email format', '## delivery',
            '**key features', '**why this works', '**analysis']):
            skip = True
            continue
        if skip:
            if lo == '' or lo == '---':
                skip = False
                continue
            if any(lo.startswith(g) for g in ['dear ', 'good day', 'hello', 'hi ', 'greetings']):
                skip = False
            else:
                continue
        if any(x in lo for x in ['character count', 'characters /', 'this letter includes',
            'this message was', 'why this works', 'respects parent capacity',
            'achieves confirmation', 'no digital access', 'key features of',
            'this works for', 'warm, not accusatory', 'math focus',
            'practical home support', 'addresses low stem']):
            continue
        if lo.startswith('|') and '|' in lo[1:]:
            continue
        clean.append(line)
    result = '\n'.join(clean).strip()
    result = re.sub(r'\n{3,}', '\n\n', result)
    result = re.sub(r'^(\s*---\s*\n?)+', '', result).strip()
    return result

def tts_player(text, key_suffix):
    """Simple 'Hear Results' button — one click, plays audio."""
    audio_key = f"tts_audio_{key_suffix}"
    if not ELEVENLABS_API_KEY: return
    if audio_key in st.session_state and st.session_state[audio_key]:
        aud = st.session_state[audio_key]
        audio_data = base64.b64decode(aud["b64"])
        st.audio(audio_data, format="audio/mp3", autoplay=True)
        ac1, ac2 = st.columns([2,1])
        with ac1:
            st.download_button("📥 Download MP3", data=audio_data, file_name=f"teacher_pehpeh_{key_suffix}.mp3", mime="audio/mp3", key=f"tts_dl_{key_suffix}")
        with ac2:
            if st.button("🗑️ Clear", key=f"tts_clr_{key_suffix}"):
                del st.session_state[audio_key]; st.rerun()
    else:
        if st.button(T("hear"), key=f"tts_gen_{key_suffix}", type="primary", use_container_width=True):
            with st.spinner(T("generating")):
                b64, src = speak_elevenlabs(text)
            if b64:
                st.session_state[audio_key] = {"b64": b64, "src": src}
                st.rerun()
            else:
                st.error(f"{T('audio_failed')}: {src}. {T('try_again')}")

# === SPEECH TO TEXT ===
def transcribe_audio(audio_bytes):
    """Transcribe audio using OpenAI Whisper. Returns text or None."""
    if not OAI or not OPENAI_API_KEY: return None
    try:
        import io
        c = openai.OpenAI(api_key=OPENAI_API_KEY)
        audio_file = io.BytesIO(audio_bytes)
        audio_file.name = "recording.wav"
        r = c.audio.transcriptions.create(model="whisper-1", file=audio_file)
        return r.text.strip() if r.text else None
    except Exception as e:
        return None

# === STUDENT ID CARD ===
def generate_student_card(student, school_name="", grade="", subject="", country="Liberia"):
    """Generate a student ID card as PNG bytes. Returns (png_bytes, filename) or (None, error)."""
    import io
    if not PIL_OK:
        return None, "Pillow not installed"
    try:
        W, H = 600, 340
        card = Image.new("RGB", (W, H), "#0F2247")
        draw = ImageDraw.Draw(card)
        # Try loading fonts, fall back to default
        try:
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 20)
            font_name = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_body = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 15)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
        except:
            font_title = ImageFont.load_default()
            font_name = font_title
            font_body = font_title
            font_small = font_title
        # Gold header bar
        draw.rectangle([0, 0, W, 56], fill="#D4A843")
        draw.text((20, 14), f"TEACHER PEHPEH - STUDENT ID", fill="#0F2247", font=font_title)
        # School name
        _school = school_name or "School"
        draw.text((20, 68), _school, fill="#F5D998", font=font_body)
        # Student name
        draw.text((20, 100), student["name"], fill="#FFFFFF", font=font_name)
        # Details
        y = 140
        details = [
            f"Grade: {grade}",
            f"Subject Focus: {subject}",
            f"Siblings: {student.get('sib', '—')}",
            f"Mother Education: {student.get('mom', '—')}",
            f"Computer Access: {student.get('cp', '—')}",
        ]
        for d in details:
            draw.text((20, y), d, fill="#C0C8D8", font=font_body)
            y += 22
        # Risk level
        rsk = []
        if student.get("mom") == "No HS": rsk.append("No HS Mom")
        if student.get("sm") == "Yes": rsk.append("Single Mom")
        if student.get("sib") == "8+": rsk.append("8+ siblings")
        if student.get("wk") == "Yes": rsk.append("Works")
        risk_color = "#EF5350" if len(rsk) >= 2 else "#FFA726" if rsk else "#66BB6A"
        risk_label = "Higher Risk" if len(rsk) >= 2 else "Some Risk" if rsk else "Lower Risk"
        draw.rectangle([20, y + 5, 170, y + 28], fill=risk_color, outline=risk_color)
        draw.text((28, y + 7), risk_label, fill="#FFFFFF", font=font_body)
        # QR Code
        if QR:
            qr_data = json.dumps({"name": student["name"], "school": _school, "grade": grade, "risk": risk_label}, ensure_ascii=False)
            qr = qrcode.QRCode(version=1, box_size=4, border=2)
            qr.add_data(qr_data)
            qr.make(fit=True)
            qr_img = qr.make_image(fill_color="#0F2247", back_color="#FFFFFF").convert("RGB")
            qr_img = qr_img.resize((120, 120))
            card.paste(qr_img, (W - 145, 80))
        else:
            draw.rectangle([W - 145, 80, W - 25, 200], fill="#1A2D50", outline="#D4A843")
            draw.text((W - 130, 130), "QR N/A", fill="#D4A843", font=font_body)
        # Footer
        draw.rectangle([0, H - 36, W, H], fill="#8B1A1A")
        draw.text((20, H - 30), f"Institute of Basic Technology (IBT) - {country}", fill="#F5D998", font=font_small)
        draw.text((W - 130, H - 30), "teacherpehpeh.org", fill="#F5D998", font=font_small)
        # Export as PNG bytes
        buf = io.BytesIO()
        card.save(buf, format="PNG")
        buf.seek(0)
        fname = f"student_card_{student['name'].replace(' ','_').lower()}.png"
        return buf.getvalue(), fname
    except Exception as e:
        return None, f"Card error: {e}"

# === QUIZ BANK ===
QUIZ = {
 "Mathematics":{"easy":[
  {"q":"Simplify: 3(2x - 4) + 5x","o":["11x - 12","11x - 4","6x - 12","x - 12"],"a":0,"e":"Expand: 6x - 12 + 5x = 11x - 12.","t":"WASSCE Paper 1 style. Teach: distribute first, then collect like terms. Write each step clearly."},
  {"q":"If log₁₀ 2 = 0.301, find log₁₀ 8.","o":["0.903","2.408","0.602","0.301"],"a":0,"e":"8 = 2³, so log 8 = 3 × log 2 = 3 × 0.301 = 0.903.","t":"WASSCE loves log questions. Teach index form first: 8=2³. Then the log law: log(aⁿ) = n log a."},
  {"q":"Solve: 2x² - 5x - 3 = 0","o":["x = 3 or -½","x = 3 or ½","x = -3 or ½","x = -3 or -½"],"a":0,"e":"Factor: (2x + 1)(x - 3) = 0. So x = -½ or x = 3.","t":"WASSCE always has quadratics. Teach factoring AND formula method. Students should check by substituting back."},
  {"q":"Find the gradient of the line 3y = 6x - 9.","o":["2","-3","6","3"],"a":0,"e":"Rearrange to y = mx + c: y = 2x - 3. Gradient m = 2.","t":"Standard WASSCE. Always rearrange to y = mx + c form first. Gradient = coefficient of x."},
  {"q":"The 5th term of an AP is 17, common difference is 3. Find the first term.","o":["5","8","2","14"],"a":0,"e":"aₙ = a + (n-1)d → 17 = a + 4(3) → a = 17 - 12 = 5.","t":"AP formula appears every year. Drill: a, a+d, a+2d... Students write out terms to verify."},
 ],"medium":[
  {"q":"If P = {1,2,3,4,5} and Q = {3,4,5,6,7}, find n(P ∪ Q).","o":["7","10","5","3"],"a":0,"e":"P ∪ Q = {1,2,3,4,5,6,7}. Count = 7.","t":"Venn diagram on board. Students place each element. WASSCE tests union, intersection, complement."},
  {"q":"In triangle ABC, a=7, b=8, C=60°. Find c using cosine rule.","o":["√57","√113","7.55","8.06"],"a":0,"e":"c² = 7² + 8² - 2(7)(8)cos60° = 49 + 64 - 56 = 57. c = √57.","t":"Cosine rule is a WASSCE favourite. Students must memorize: c² = a² + b² - 2ab·cosC."},
  {"q":"Differentiate y = 3x⁴ - 2x² + 5x - 1.","o":["12x³ - 4x + 5","12x³ - 4x + 5x","3x³ - 2x + 5","12x⁴ - 4x²"],"a":0,"e":"dy/dx = 12x³ - 4x + 5. Rule: bring power down, reduce by 1. Constants vanish.","t":"Differentiation appears in WASSCE Paper 2. Drill the power rule until it's automatic. Use mnemonic: 'Multiply by power, subtract one.'"},
  {"q":"A fair die is thrown twice. P(sum = 7)?","o":["1/6","1/12","5/36","7/36"],"a":0,"e":"Outcomes summing to 7: (1,6)(2,5)(3,4)(4,3)(5,2)(6,1) = 6 out of 36. P = 6/36 = 1/6.","t":"Draw the 6×6 grid. Students count all sums. Probability questions need systematic listing."},
  {"q":"Convert 101101₂ to base 10.","o":["45","37","53","29"],"a":0,"e":"1(32) + 0(16) + 1(8) + 1(4) + 0(2) + 1(1) = 32+8+4+1 = 45.","t":"Number bases: WASSCE staple. Write place values (32,16,8,4,2,1) above the digits. Multiply and add."},
 ],"hard":[
  {"q":"Evaluate ∫(3x² + 2x - 1)dx from 0 to 2.","o":["10","8","12","14"],"a":0,"e":"[x³ + x² - x] from 0 to 2 = (8+4-2) - 0 = 10.","t":"Integration: reverse of differentiation. Add 1 to power, divide by new power. Then substitute limits."},
  {"q":"The 3rd and 6th terms of a GP are 18 and 486. Find the common ratio.","o":["3","9","6","2"],"a":0,"e":"ar² = 18, ar⁵ = 486. Divide: r³ = 27, r = 3.","t":"GP: divide consecutive term equations to eliminate 'a'. This technique appears in WASSCE Theory."},
  {"q":"Two fair coins and a die are tossed. P(2 heads and even number)?","o":["1/4","1/8","1/12","1/6"],"a":0,"e":"P(2 heads) = ¼. P(even) = ½. Independent: ¼ × ½ = ⅛.","t":"Compound probability: independent events multiply. Tree diagrams help visual learners."},
 ]},
 "Physics":{"easy":[
  {"q":"A car accelerates from 10 m/s to 30 m/s in 5s. Acceleration?","o":["4 m/s²","6 m/s²","8 m/s²","2 m/s²"],"a":0,"e":"a = (v-u)/t = (30-10)/5 = 20/5 = 4 m/s².","t":"WASSCE kinematics: always identify u, v, a, s, t first. Write them down before calculating."},
  {"q":"A 2kg mass falls from 10m height. What is its PE at the top? (g=10)","o":["200 J","20 J","100 J","2000 J"],"a":0,"e":"PE = mgh = 2 × 10 × 10 = 200 J.","t":"Energy conservation is key WASSCE topic. PE at top converts to KE at bottom. Ask: where does energy go?"},
  {"q":"Resistance of two 6Ω resistors in parallel?","o":["3 Ω","12 Ω","6 Ω","9 Ω"],"a":0,"e":"1/R = 1/6 + 1/6 = 2/6. R = 3 Ω.","t":"Parallel: 'product over sum' shortcut for two resistors: (6×6)/(6+6) = 36/12 = 3Ω."},
  {"q":"A wave has frequency 50 Hz and wavelength 4m. Speed?","o":["200 m/s","12.5 m/s","54 m/s","46 m/s"],"a":0,"e":"v = fλ = 50 × 4 = 200 m/s.","t":"v = fλ is the wave equation. WASSCE applies it to sound, light, and water waves. Students must know it cold."},
  {"q":"An object on a spring has period 0.5s. Frequency?","o":["2 Hz","0.5 Hz","5 Hz","0.2 Hz"],"a":0,"e":"f = 1/T = 1/0.5 = 2 Hz.","t":"Period and frequency are inverses. This is basic but students often confuse them on WASSCE."},
 ],"medium":[
  {"q":"A 5kg block is pushed with 30N on a frictionless surface. Acceleration?","o":["6 m/s²","150 m/s²","25 m/s²","35 m/s²"],"a":0,"e":"F = ma → a = F/m = 30/5 = 6 m/s².","t":"Newton's 2nd Law: WASSCE asks both calculation and conceptual understanding. What if friction = 10N?"},
  {"q":"A ray enters glass (n=1.5) at 30° to the normal. Angle of refraction?","o":["19.5°","45°","20°","30°"],"a":0,"e":"Snell's law: sin30°/sinr = 1.5 → sinr = 0.5/1.5 = 0.333 → r ≈ 19.5°.","t":"Snell's law is examined every WASSCE. Students must use sine tables or know key sin values."},
  {"q":"EMF of cell = 12V, internal resistance = 2Ω, external R = 4Ω. Current?","o":["2 A","3 A","6 A","4 A"],"a":0,"e":"I = EMF/(R+r) = 12/(4+2) = 2 A.","t":"Internal resistance: WASSCE Theory favourite. Draw the circuit. EMF = V + Ir. Students confuse EMF with terminal p.d."},
  {"q":"Half-life of a substance is 4 days. After 12 days, what fraction remains?","o":["1/8","1/4","1/16","1/6"],"a":0,"e":"12 days = 3 half-lives. Fraction = (½)³ = 1/8.","t":"Half-life: count how many half-lives fit. Each one halves the amount. WASSCE asks both fraction and mass."},
 ],"hard":[
  {"q":"A projectile is launched at 30 m/s at 60° to horizontal. Max height? (g=10)","o":["33.75 m","45 m","22.5 m","15 m"],"a":0,"e":"Vy = 30sin60° = 25.98. H = Vy²/2g = 675/20 = 33.75 m.","t":"Projectile: resolve into components. Vertical determines height, horizontal determines range. WASSCE Paper 2 staple."},
  {"q":"A transformer has 200 primary turns, 50 secondary turns, input 240V. Output voltage?","o":["60 V","960 V","48 V","12 V"],"a":0,"e":"Vs/Vp = Ns/Np → Vs = 240 × 50/200 = 60 V. Step-down transformer.","t":"Transformer equation: Vs/Vp = Ns/Np = Ip/Is. WASSCE tests both step-up and step-down."},
  {"q":"A 0.5kg ball moving at 4 m/s hits a wall and bounces back at 3 m/s. Impulse?","o":["3.5 Ns","0.5 Ns","7 Ns","1.5 Ns"],"a":0,"e":"Impulse = m(v-u) = 0.5(3-(-4)) = 0.5 × 7 = 3.5 Ns. (Direction change: -4 becomes +3.)","t":"Impulse = change in momentum. KEY: when direction reverses, ADD the speeds. Common WASSCE trap."},
 ]},
 "Biology":{"easy":[
  {"q":"Which organelle is the site of aerobic respiration?","o":["Mitochondria","Ribosome","Nucleus","Golgi body"],"a":0,"e":"Mitochondria: the 'powerhouse' where glucose + O₂ → CO₂ + H₂O + ATP.","t":"WASSCE cell biology: students must name organelle AND function. Use the 'factory analogy' — each organelle has a job."},
  {"q":"In humans, the diploid number is 46. How many chromosomes in a gamete?","o":["23","46","92","12"],"a":0,"e":"Gametes are haploid (n). 46/2 = 23 chromosomes.","t":"Meiosis halves the chromosome number. WASSCE tests: diploid vs haploid, mitosis vs meiosis differences."},
  {"q":"Which blood vessel carries oxygenated blood FROM the heart to the body?","o":["Aorta","Pulmonary artery","Vena cava","Pulmonary vein"],"a":0,"e":"The aorta carries oxygenated blood from the left ventricle to the body.","t":"Blood vessel quiz is WASSCE standard. Trick: pulmonary artery carries DE-oxygenated blood (only exception)."},
  {"q":"Sickle cell trait (HbAS) provides resistance to which disease?","o":["Malaria","Typhoid","Cholera","HIV"],"a":0,"e":"HbAS heterozygotes have partial protection against Plasmodium falciparum malaria.","t":"Sickle cell and malaria: a key WASSCE genetics topic. Discuss why the gene persists in malaria-endemic West Africa."},
  {"q":"Which hormone controls blood sugar level?","o":["Insulin","Adrenaline","Thyroxine","Oestrogen"],"a":0,"e":"Insulin (from pancreas) lowers blood glucose by promoting uptake into cells.","t":"Endocrine system: WASSCE asks hormone name, gland, and function. Make a table: gland → hormone → effect."},
 ],"medium":[
  {"q":"In a cross between Tt × Tt, what fraction of offspring are tall (T dominant)?","o":["3/4","1/4","1/2","1"],"a":0,"e":"TT:Tt:tt = 1:2:1. Tall (TT+Tt) = 3/4.","t":"Punnett square every time. WASSCE asks ratios AND probabilities. Practice with different crosses."},
  {"q":"Which nitrogenous base is found in RNA but NOT DNA?","o":["Uracil","Thymine","Adenine","Cytosine"],"a":0,"e":"RNA has Uracil instead of Thymine. Both have A, G, C.","t":"DNA vs RNA: sugar, bases, structure. Mnemonic for RNA bases: 'GACU' sounds like 'gecko.'"},
  {"q":"Oxygen debt occurs after vigorous exercise because:","o":["Lactic acid must be broken down","Glucose runs out","Lungs stop working","Blood pressure drops"],"a":0,"e":"During intense exercise, anaerobic respiration produces lactic acid. Extra O₂ is needed to oxidize it.","t":"Link to students' experience: why do you breathe hard AFTER stopping running? Oxygen debt!"},
  {"q":"Which structure prevents food from entering the windpipe?","o":["Epiglottis","Larynx","Pharynx","Uvula"],"a":0,"e":"The epiglottis closes over the trachea during swallowing.","t":"Digestion system: WASSCE tests the pathway. Students act out swallowing — what moves, what closes?"},
 ],"hard":[
  {"q":"In ecological succession, the first organisms to colonize bare rock are:","o":["Lichens","Grasses","Trees","Ferns"],"a":0,"e":"Lichens are pioneer species — they break down rock into soil, enabling other plants to grow.","t":"Succession: pioneer → grass → shrub → tree. WASSCE Theory asks for full description of stages."},
  {"q":"A man with blood group AB marries a woman with group O. Possible children?","o":["A and B only","AB only","O only","A, B, AB, and O"],"a":0,"e":"Father: IᴬIᴮ × Mother: ii → children are Iᴬi (A) or Iᴮi (B). No AB, no O.","t":"Blood group genetics: multiple alleles + codominance. WASSCE loves this. Always write the genotypes first."},
  {"q":"Which process releases CO₂ back into the atmosphere in the carbon cycle?","o":["Respiration and combustion","Photosynthesis","Nitrogen fixation","Transpiration"],"a":0,"e":"Respiration (by all living things) and combustion (burning fuels) release CO₂.","t":"Carbon cycle: WASSCE asks arrows and processes. Draw the cycle on board with students labeling each arrow."},
 ]},
 "Chemistry":{"easy":[
  {"q":"What is the oxidation state of Mn in KMnO₄?","o":["+7","+4","+2","+6"],"a":0,"e":"K(+1) + Mn(x) + 4O(-2) = 0 → 1 + x - 8 = 0 → x = +7.","t":"Oxidation states: WASSCE standard. Rules: O is -2, alkali metals +1. Solve for the unknown."},
  {"q":"Which gas is collected over water in the lab?","o":["Oxygen","HCl","NH₃","SO₂"],"a":0,"e":"O₂ is insoluble in water, so it's collected by downward displacement of water.","t":"Gas collection methods depend on solubility and density. WASSCE tests all three methods."},
  {"q":"The IUPAC name of CH₃CH₂OH is:","o":["Ethanol","Methanol","Propanol","Ethanal"],"a":0,"e":"2 carbons = eth-. -OH group = -anol. Ethanol.","t":"IUPAC naming: count carbons, identify functional group. WASSCE organic chemistry requires this."},
  {"q":"How many moles are in 44g of CO₂? (C=12, O=16)","o":["1","2","0.5","44"],"a":0,"e":"Molar mass of CO₂ = 12 + 32 = 44 g/mol. Moles = 44/44 = 1.","t":"Mole calculations appear every WASSCE. Formula: n = mass/molar mass. Practice with different substances."},
  {"q":"Which of these is an alkali?","o":["NaOH","HCl","NaCl","H₂SO₄"],"a":0,"e":"NaOH (sodium hydroxide) is a strong alkali/base. It produces OH⁻ in solution.","t":"Acids produce H⁺, alkalis produce OH⁻. WASSCE tests indicators, neutralization, and salt formation."},
 ],"medium":[
  {"q":"What volume of H₂ at STP is produced when 2 moles of Zn react with excess HCl?","o":["44.8 L","22.4 L","11.2 L","67.2 L"],"a":0,"e":"Zn + 2HCl → ZnCl₂ + H₂. 1 mol Zn gives 1 mol H₂. 2 mol Zn → 2 mol H₂ = 2 × 22.4 = 44.8 L.","t":"Stoichiometry + molar volume (22.4L at STP). Write balanced equation first. WASSCE Paper 2 calculation."},
  {"q":"Which type of reaction is: CuO + H₂SO₄ → CuSO₄ + H₂O?","o":["Neutralization","Decomposition","Displacement","Combustion"],"a":0,"e":"Base (CuO) + Acid (H₂SO₄) → Salt + Water = Neutralization.","t":"Reaction types: WASSCE tests naming AND identifying from equations. Make a chart of types with examples."},
  {"q":"In electrolysis of brine, what is produced at the cathode?","o":["Hydrogen","Chlorine","Sodium","Oxygen"],"a":0,"e":"At cathode (negative), H⁺ ions are reduced: 2H⁺ + 2e⁻ → H₂.","t":"Electrolysis: cathode = reduction (CATions go to CAThode). WASSCE tests products at each electrode."},
  {"q":"An element has electronic configuration 2,8,7. Its likely ion is:","o":["X⁻","X⁷⁺","X⁺","X²⁻"],"a":0,"e":"7 outer electrons — needs 1 more for stable octet. Gains 1 electron → X⁻ (like chlorine).","t":"Electronic configuration → group → valency → ion charge. This chain of reasoning is WASSCE standard."},
 ],"hard":[
  {"q":"Calculate the enthalpy change: C + O₂ → CO₂, given C + ½O₂ → CO (ΔH=-110kJ) and CO + ½O₂ → CO₂ (ΔH=-283kJ).","o":["-393 kJ","-173 kJ","+393 kJ","-110 kJ"],"a":0,"e":"Hess's Law: add the two equations. ΔH = -110 + (-283) = -393 kJ.","t":"Hess's Law: WASSCE Theory. Energy is a state function — path doesn't matter. Add equations like algebra."},
  {"q":"0.1M NaOH is titrated against 0.05M H₂SO₄. What volume of acid neutralizes 25ml of NaOH?","o":["25 ml","50 ml","12.5 ml","100 ml"],"a":0,"e":"2NaOH + H₂SO₄. Moles NaOH = 0.1×25 = 2.5mmol. Moles acid = 2.5/2 = 1.25mmol. Vol = 1.25/0.05 = 25ml.","t":"Titration: write balanced equation, find mole ratio, use C₁V₁/n₁ = C₂V₂/n₂. WASSCE practical and theory."},
  {"q":"Which compound shows geometric (cis-trans) isomerism?","o":["But-2-ene","Ethene","Propane","Ethanol"],"a":0,"e":"But-2-ene: C=C with different groups on each carbon allows cis/trans forms.","t":"Isomerism: structural vs geometric vs optical. WASSCE asks students to draw both forms. C=C restricts rotation."},
 ]},
 "Reading Comprehension":{"easy":[
  {"q":"'The government's education budget was slashed by 30%, leading to school closures across rural communities.' — What caused the school closures?","o":["Budget cuts","Natural disaster","Teacher strike","Student protests"],"a":0,"e":"The passage states the budget was 'slashed by 30%' which led to closures.","t":"WASSCE comprehension: identify cause and effect. Teach students to look for linking words: 'leading to', 'because', 'therefore'."},
  {"q":"'Despite the drought, the farmers of Bong County managed to harvest enough rice to sustain their families through the dry season.' — What is the tone?","o":["Resilient / hopeful","Angry","Sad","Humorous"],"a":0,"e":"'Despite' shows challenge overcome. 'Managed to' shows resilience. Tone is hopeful.","t":"Tone questions: look at word choice. 'Despite' + 'managed' = overcoming. WASSCE tests tone, mood, attitude."},
  {"q":"'Deforestation in the tropics has accelerated at an alarming rate.' — What does 'accelerated' mean here?","o":["Increased in speed","Slowed down","Stopped completely","Remained constant"],"a":0,"e":"'Accelerated' means sped up / increased. 'Alarming rate' confirms it's getting worse.","t":"Vocabulary in context: WASSCE gives unfamiliar words. Teach: look at surrounding words for clues. 'Alarming' = bad = getting worse."},
 ],"medium":[
  {"q":"'The author argues that technology alone cannot solve Africa's education crisis; rather, it must be coupled with trained teachers and culturally relevant curricula.' — What is the author's main argument?","o":["Technology needs teachers and relevant curricula to work","Technology is useless","Africa doesn't need technology","Curricula are already good"],"a":0,"e":"Key phrase: 'cannot solve alone... must be coupled with.' Author wants technology PLUS human/cultural elements.","t":"Argumentative comprehension: find the claim AND the qualifier. WASSCE asks 'What is the writer's view?' Look for 'rather', 'however', 'must'."},
  {"q":"'She was the proverbial candle burning at both ends — teaching by day, farming by evening, and studying by lamplight.' — What figure of speech?","o":["Metaphor","Simile","Personification","Hyperbole"],"a":0,"e":"'Was the candle' (not 'like a candle') = metaphor. Describes exhausting double life.","t":"Simile uses 'like/as'. Metaphor says IS. WASSCE English Paper 1 tests these. Students collect examples from their own speech."},
 ],"hard":[
  {"q":"Read: 'The policy, while well-intentioned, failed to account for the socioeconomic realities of the communities it aimed to serve.' — The author's attitude is:","o":["Critically sympathetic","Fully supportive","Hostile","Indifferent"],"a":0,"e":"'Well-intentioned' = sympathetic. 'Failed to account' = critical. Both together = critically sympathetic.","t":"Advanced comprehension: authors can hold MIXED views. WASSCE rewards nuanced answers over simple ones."},
  {"q":"'The minister's assertion that unemployment had decreased was contradicted by data showing a 15% rise in joblessness among youth.' — This is an example of:","o":["Irony","Metaphor","Alliteration","Flashback"],"a":0,"e":"The minister's claim is the OPPOSITE of reality — this is verbal/situational irony.","t":"WASSCE literary devices: irony = opposite of what's expected/stated. Teach with local examples: 'The fire station burned down.'"},
 ]},
 "Literature":{"easy":[
  {"q":"In Chinua Achebe's 'Things Fall Apart', what is Okonkwo's greatest fear?","o":["Being seen as weak like his father","Losing his farm","Being exiled","Dying in battle"],"a":0,"e":"Okonkwo fears weakness because his father Unoka was a debtor and failure. This fear drives all his overcompensation with violence and rigidity.","t":"Character motivation is central to WASSCE Literature. Ask: what does this character FEAR? What do they WANT? Everything else follows from there."},
  {"q":"'Weep Not, Child' by Ngũgĩ wa Thiong'o is set against which historical backdrop?","o":["The Mau Mau uprising in Kenya","The Nigerian Civil War","Apartheid in South Africa","The Liberian Civil War"],"a":0,"e":"The novel is set during the Mau Mau struggle for Kenyan independence against British colonialism in the 1950s.","t":"African literature is rooted in history. Pair each text with its historical context — students who know the history understand the characters' choices far better."},
  {"q":"What literary device is used in: 'The sun bled across the horizon'?","o":["Personification","Simile","Alliteration","Oxymoron"],"a":0,"e":"'Bled' is a human/bodily verb applied to the sun — this is personification, giving the sun human qualities.","t":"WASSCE asks students to name AND explain devices. Not just 'personification' but 'the sun is given human characteristics through the verb bled, creating a vivid sense of pain or sacrifice.'"},
  {"q":"In 'The African Child' by Camara Laye, what does the golden snake symbolise?","o":["The spirit and power of his father's lineage","Danger and evil","Colonial oppression","Wealth and greed"],"a":0,"e":"The black snake that visits Laye's father represents the protective ancestral spirit of the family — a totem linked to his father's supernatural blacksmithing powers.","t":"Symbolism in African literature often draws on spiritual tradition. Teach students to ask: what cultural meaning does this object carry? Not just Western literary symbolism."},
  {"q":"What is the meaning of 'dramatic irony'?","o":["When the audience knows something a character does not","When a character says the opposite of what they mean","When opposite ideas are placed side by side","A sudden twist at the end of a story"],"a":0,"e":"Dramatic irony: the audience/reader has information the character lacks, creating tension or tragedy. Classic example: we know who is plotting against a character while they remain unsuspecting.","t":"WASSCE tests irony in three forms: verbal (sarcasm), situational (opposite outcome), dramatic (audience knows more). Practise identifying all three with local examples."},
 ],"medium":[
  {"q":"In 'Things Fall Apart', the title is taken from which poem?","o":["'The Second Coming' by W.B. Yeats","'The Waste Land' by T.S. Eliot","'Ode to a Nightingale' by Keats","'Dulce et Decorum Est' by Owen"],"a":0,"e":"Achebe took the title from Yeats: 'Things fall apart; the centre cannot hold.' This captures the collapse of Igbo society under colonialism.","t":"Intertextuality: WASSCE may ask about allusions. Teach students to ask — why did the author choose THIS title? What does the original source add to the meaning?"},
  {"q":"Which narrative technique does Achebe primarily use in 'Things Fall Apart'?","o":["Third-person omniscient","First-person narrator","Stream of consciousness","Epistolary"],"a":0,"e":"Achebe uses a third-person omniscient narrator who can enter multiple characters' minds and describe the community from the outside — deliberately mimicking the style of an oral storyteller.","t":"Narrative perspective shapes everything. Third-person omniscient lets Achebe present Igbo culture as legitimate and complete — a deliberate counter to colonial narratives."},
  {"q":"The theme of 'chi' (personal god/destiny) in 'Things Fall Apart' suggests:","o":["A person's fate is partly self-determined and partly pre-ordained","Fate is entirely controlled by the gods","Humans have complete free will","Chi only affects warriors"],"a":0,"e":"Achebe shows chi as both given and shaped — Okonkwo says 'when a man says yes, his chi says yes also.' His own stubborn choices contribute to his downfall.","t":"Thematic depth: chi is not simple predestination. It's an interplay between destiny and choice. WASSCE rewards answers that show this complexity rather than reducing it."},
 ],"hard":[
  {"q":"Achebe's authorial intention in 'Things Fall Apart' was primarily to:","o":["Restore dignity to African culture by presenting it from the inside","Condemn the Igbo for resisting colonialism","Prove that Africa needed European help","Celebrate Okonkwo as a hero"],"a":0,"e":"Achebe explicitly stated he wrote the novel to counter Joyce Cary's 'Mister Johnson' and Conrad's 'Heart of Darkness' — to present Africa as a place of complex, functioning civilisation, not a void.","t":"Author intent is fair game in WASSCE essays. Achebe is the most quoted African author on intent. 'To teach' answers are weak — drill the specific counter-narrative purpose."},
  {"q":"In 'Weep Not, Child', Njoroge's education represents:","o":["A false promise that colonialism cannot deliver on","Pure opportunity with no contradiction","A betrayal of African tradition","An irrelevant colonial imposition"],"a":0,"e":"Njoroge's faith in education as salvation is tragically undercut — colonial education promises uplift while the colonial system simultaneously dispossesses his family. It is hope within a system of oppression.","t":"WASSCE essay topics often ask about 'the role of education' in these texts. Train students to see it as BOTH genuine hope AND colonial contradiction — not one or the other."},
 ]},
 "Economics":{"easy":[
  {"q":"Which of these is an example of a free good?","o":["Air","Water from a tap","Textbooks","Electricity"],"a":0,"e":"Free goods have no price because they are unlimited in supply relative to demand — air is the classic example. Tap water requires treatment and infrastructure, so it has a cost.","t":"Free goods vs economic goods: a foundational distinction. WASSCE Paper 1 tests this every year. Teach with local contrast: river air (free) vs paying for electricity."},
  {"q":"The law of demand states that, all else equal, as price rises:","o":["Quantity demanded falls","Quantity demanded rises","Supply increases","Supply decreases"],"a":0,"e":"Inverse relationship: price ↑ → quantity demanded ↓. This is the cornerstone of demand theory. 'All else equal' (ceteris paribus) is essential — other factors are held constant.","t":"WASSCE loves the law of demand. Drill the ceteris paribus condition — students who say 'unless it's a Giffen good' in essays score extra marks."},
  {"q":"What is the difference between microeconomics and macroeconomics?","o":["Micro studies individual units; macro studies the whole economy","Micro is about imports; macro is about exports","Micro is government policy; macro is firm behaviour","They are the same subject"],"a":0,"e":"Micro: firm/consumer decisions. Macro: economy-wide variables — GDP, inflation, unemployment. Both are needed for policy analysis.","t":"WASSCE Paper 1 always has a definition question. Train students to give a crisp one-sentence definition, then an example. Micro: 'Why does a Liberian market woman price her tomatoes?' Macro: 'Why is inflation rising nationally?'"},
  {"q":"If supply increases and demand stays the same, the equilibrium price will:","o":["Fall","Rise","Stay the same","Double"],"a":0,"e":"More supply chasing the same demand: sellers compete → price falls to a new lower equilibrium.","t":"Supply-demand diagrams: WASSCE Paper 2 asks students to draw and explain shifts. Teach the 4 cases as a matrix: supply up/down × demand up/down."},
  {"q":"Opportunity cost is best defined as:","o":["The value of the next best alternative forgone","The total cost of a decision","The monetary price of a good","The cost of production"],"a":0,"e":"Opportunity cost = what you give up by choosing one option over the next best. It doesn't have to be money — time and resources count too.","t":"Opportunity cost appears in almost every WASSCE Economics paper. Teach students to identify it in everyday decisions: going to school vs working on the farm."},
 ],"medium":[
  {"q":"Which type of unemployment is caused by a mismatch between workers' skills and available jobs?","o":["Structural","Cyclical","Frictional","Seasonal"],"a":0,"e":"Structural unemployment: the economy's structure changes (e.g. technology replaces workers) and their skills no longer match available jobs. Long-term and requires retraining.","t":"Types of unemployment: WASSCE Paper 2 essay favourite. Teach each type with a West African example: seasonal (cocoa harvest off-season), structural (artisan printers replaced by digital), frictional (graduate searching for first job)."},
  {"q":"Gross Domestic Product (GDP) measures:","o":["Total value of goods and services produced in a country in a year","Total government spending","The value of exports minus imports","Total wages paid to workers"],"a":0,"e":"GDP = total monetary value of all final goods and services produced within a country's borders in a given period, regardless of who produces them.","t":"GDP vs GNP: WASSCE often distinguishes these. GDP = within borders. GNP = by nationals (anywhere). For small open economies like Ghana or Liberia, remittances make this difference significant."},
  {"q":"A progressive tax system means:","o":["Higher earners pay a higher percentage of income as tax","Everyone pays the same rate","Lower earners pay more","Tax rate falls as income rises"],"a":0,"e":"Progressive = rate rises with income. Designed to reduce inequality. Income tax in most West African countries is progressive — more brackets at the top.","t":"WASSCE asks students to evaluate tax systems. Progressive (fair by ability), Regressive (VAT hits poor harder proportionally), Proportional (flat rate). Always link to equity vs efficiency trade-off."},
 ],"hard":[
  {"q":"If a country's current account deficit is increasing, this most likely indicates:","o":["The country is importing more than it exports","The country is saving more than it invests","Government spending is falling","Inflation is decreasing"],"a":0,"e":"Current account deficit = exports < imports (plus income and transfers). The country is spending more on foreign goods than it earns from selling abroad. Requires financing via capital account inflows.","t":"Balance of payments: WASSCE Theory Paper. Current account + capital account + financial account = 0 (in theory). Teach: Liberia has persistent current account deficits — why? Iron ore, rubber exports vs manufactured imports."},
  {"q":"The multiplier effect in Keynesian economics means:","o":["An initial injection of spending leads to a larger final increase in GDP","Taxes multiply government revenue","Inflation multiplies the cost of goods","Investment always doubles output"],"a":0,"e":"Multiplier = 1/(1-MPC). If MPC = 0.8, multiplier = 5. $100M government spending → $500M increase in GDP, because each round of spending becomes income that is partly re-spent.","t":"The multiplier is WASSCE Paper 2 gold. Teach the formula, draw the circular flow, and use a local example: government builds a school → workers earn income → spend in market → market woman earns → she buys goods… chain continues."},
 ]},
 "Geography":{"easy":[
  {"q":"The Tropic of Cancer is located at:","o":["23.5°N latitude","23.5°S latitude","0° latitude","66.5°N latitude"],"a":0,"e":"The Tropic of Cancer is the northernmost latitude where the sun can be directly overhead — at 23.5° North, on the June solstice.","t":"The five key latitudes appear in WASSCE every year: Equator (0°), Tropics of Cancer/Capricorn (23.5°N/S), Arctic/Antarctic Circles (66.5°N/S). Draw them on a blank globe and label from memory."},
  {"q":"Which West African country is landlocked?","o":["Mali","Ghana","Liberia","Sierra Leone"],"a":0,"e":"Mali is landlocked — surrounded entirely by land with no coastline. Ghana, Liberia, and Sierra Leone all have Atlantic Ocean coastlines.","t":"WASSCE Geography: know which West African countries are landlocked (Mali, Burkina Faso, Niger, Guinea-Bissau is coastal, Chad). Being landlocked affects trade, development, and climate."},
  {"q":"What type of rainfall is common in coastal West Africa during the rainy season?","o":["Convectional rainfall","Relief rainfall","Frontal rainfall","Cyclonic rainfall"],"a":0,"e":"Convectional rainfall: the sun heats the ground, moist air rises, cools, condenses into clouds and heavy afternoon rain. Classic in tropical coastal West Africa.","t":"Three types of rainfall for WASSCE: convectional (tropics, afternoon), relief/orographic (mountains), frontal (temperate zones). West Africa = mostly convectional. Students should draw each diagram."},
  {"q":"The Sahel zone is best described as:","o":["A semi-arid transition zone between the Sahara and savanna","A rainforest belt","A coastal zone","A high-altitude plateau"],"a":0,"e":"The Sahel (Arabic: 'coast/shore') is a narrow belt of semi-arid land stretching across Africa just south of the Sahara. It receives 200-600mm rain per year and is highly vulnerable to desertification.","t":"The Sahel is frequently examined in WASSCE. Key issues: desertification, overgrazing, population pressure, Lake Chad shrinkage. Connect physical geography to human impact."},
  {"q":"Which river forms the border between Liberia and Côte d'Ivoire in parts?","o":["The Cavalla River","The Mano River","The Niger River","The Volta River"],"a":0,"e":"The Cavalla (Cavally) River forms much of the border between Liberia and Côte d'Ivoire as it flows to the Atlantic.","t":"West African river geography: WASSCE WAEC students in Liberia should know the Mano River (Liberia-Sierra Leone border), Cavalla River (Liberia-Côte d'Ivoire border), and the major regional rivers: Niger, Volta, Senegal."},
 ],"medium":[
  {"q":"Urbanisation in West Africa is mainly driven by:","o":["Rural-urban migration seeking employment and services","Decrease in agricultural land","International immigration","Coastal flooding"],"a":0,"e":"The primary driver is rural-urban migration: people move to cities seeking jobs, education, healthcare, and infrastructure. This is the dominant pattern across Accra, Lagos, Monrovia, and Freetown.","t":"WASSCE Population Geography: push and pull factors. Push (drought, poverty, lack of schools, no jobs in villages). Pull (jobs, hospitals, universities, markets). Teach students to write two columns and fill them for Liberia specifically."},
  {"q":"What is the main cause of soil erosion in West African farmlands?","o":["Heavy rainfall on unprotected soils and deforestation","Light winds","Over-irrigation","Volcanic activity"],"a":0,"e":"Intense tropical rainfall on exposed topsoil (due to deforestation and slash-and-burn farming) removes the fertile top layer. Without tree roots to hold soil, erosion is rapid.","t":"Soil conservation: WASSCE Paper 2 topic. Solutions: contour farming, mulching, windbreaks, reforestation, terracing. For every problem, students should know at least 3 solutions."},
  {"q":"Which of the following is a consequence of the greenhouse effect?","o":["Rising global temperatures and sea level rise","More rainfall in deserts","Decrease in ultraviolet radiation","Cooling of the ocean floor"],"a":0,"e":"Greenhouse gases trap outgoing radiation → atmosphere warms → polar ice melts → sea levels rise. For low-lying West African coasts, this is an existential threat.","t":"Climate change is now examinable at WASSCE level. Teach the mechanism (GHGs trap heat) and the West African impacts: desertification, coastal erosion, disrupted rainfall, crop failure."},
 ],"hard":[
  {"q":"The concept of 'demographic transition' predicts that as countries develop:","o":["Birth and death rates both eventually fall, stabilising population","Population always grows indefinitely","Death rates rise with development","Birth rates rise before death rates"],"a":0,"e":"Demographic Transition Model (DTM): Stage 1 (high birth + high death), Stage 2 (death rate falls, population booms), Stage 3 (birth rate falls), Stage 4 (both low, stable population). Most West African countries are in Stage 2-3.","t":"DTM is a WASSCE Paper 2 essay question. Draw all 4 stages, label the axes, and place West African countries on it. Liberia is approximately Stage 2 with declining mortality but still high fertility."},
  {"q":"Describe the main economic activity associated with the savanna vegetation zone of West Africa:","o":["Mixed farming, cattle herding, and groundnut cultivation","Deep-sea fishing","Cocoa and rubber plantation farming","Diamond and iron ore mining"],"a":0,"e":"The Guinea and Sudan savanna zones support mixed farming systems: crop cultivation (groundnuts, millet, sorghum, yams) alongside cattle herding. The Fulani pastoralists exemplify this zone's dominant livelihood.","t":"Vegetation zones and economic activities: WASSCE classic. Match each zone (mangrove, rainforest, savanna, Sahel) to its agriculture, challenges, and countries. Make a table and memorise it."},
 ]},
 "French":{"easy":[
  {"q":"How do you say 'I am going to school' in French?","o":["Je vais à l'école","Je suis à l'école","J'ai allé à l'école","Je vais l'école"],"a":0,"e":"'Je vais à l'école' — present tense of 'aller' (to go) + 'à' + definite article. Note: 'à le' contracts to 'au', but 'à la' stays as 'à la', and 'à l'' before vowel sounds.","t":"WASSCE French Paper 1: article contraction is tested frequently. Au (à+le), aux (à+les), but à la and à l' stay unchanged. Drill these with classroom objects: au tableau, à la fenêtre."},
  {"q":"Which is the correct past tense (passé composé) of 'parler' (to speak) for 'je'?","o":["j'ai parlé","j'ai parler","je suis parlé","j'avais parlé"],"a":0,"e":"Passé composé = avoir/être (conjugated) + past participle. 'Parler' uses 'avoir': j'ai parlé. Regular -er verbs form past participle by replacing -er with -é.","t":"WASSCE French: the most tested grammar point is passé composé vs imparfait. Passé composé = completed actions ('I spoke'). Imparfait = ongoing/habitual past ('I was speaking / I used to speak')."},
  {"q":"What does 'Qu'est-ce que tu fais?' mean?","o":["What are you doing?","What did you do?","What will you do?","Where are you going?"],"a":0,"e":"'Qu'est-ce que' is a question structure meaning 'what'. 'Tu fais' is 'you do/are doing' (faire, present tense). So: 'What are you doing?'","t":"Question structures: WASSCE tests comprehension of spoken/written questions. Teach: Qu'est-ce que (what), Où (where), Quand (when), Pourquoi (why), Comment (how), Qui (who). Practice with rapid-fire drills."},
  {"q":"How do you say 'She has two brothers' in French?","o":["Elle a deux frères","Elle est deux frères","Elle a deux frère","Elle deux frères"],"a":0,"e":"'Elle a' = she has (avoir). 'Deux' = two. 'Frères' = brothers (plural of frère, add -s). Note: 'avoir' is used for family possession, not 'être'.","t":"'Avoir' vs 'être': one of the most common WASSCE errors. Avoir = to have (possession, age: j'ai 16 ans). Être = to be (nationality, profession, state). Drill: J'ai faim (hungry) — NOT je suis faim."},
  {"q":"What is the French word for 'school'?","o":["l'école","la salle","le bureau","le marché"],"a":0,"e":"L'école = school (feminine noun, elided article before vowel). La salle = room/hall. Le bureau = office/desk. Le marché = market.","t":"WASSCE French vocabulary: school environment words appear frequently. École, salle de classe, professeur, élève, tableau, cahier, livre. Make vocabulary sets by theme and drill weekly."},
 ],"medium":[
  {"q":"Which sentence correctly uses the subjunctive mood?","o":["Il faut que tu fasses tes devoirs","Il faut que tu fais tes devoirs","Il faut que tu a fait tes devoirs","Il faut que tu fairas tes devoirs"],"a":0,"e":"'Il faut que' triggers the subjunctive. 'Faire' → subjunctive: je fasse, tu fasses. NOT the indicative 'fais'. Subjunctive follows: il faut que, vouloir que, douter que, bien que.","t":"Subjunctive: the hardest WASSCE French grammar point. Triggers include: il faut que, vouloir que, bien que, pour que, avant que. Teach the 20 most common irregular subjunctives by rote: être→sois, avoir→aies, aller→ailles, faire→fasses."},
  {"q":"What is the correct negative form of 'Je parle français'?","o":["Je ne parle pas français","Je ne pas parle français","Je parle ne pas français","Je ne parle français pas"],"a":0,"e":"French negation: ne…pas wraps the verb. Ne comes before the verb, pas after. In spoken French, 'ne' is often dropped, but in writing, both are required.","t":"Negation patterns for WASSCE: ne…pas (not), ne…jamais (never), ne…plus (no longer), ne…rien (nothing), ne…personne (nobody). Students often place 'pas' in the wrong position. Drill with physical flashcards."},
  {"q":"'Je me lève à six heures' — what type of verb is 'se lever'?","o":["Reflexive verb","Irregular verb","Defective verb","Impersonal verb"],"a":0,"e":"Reflexive verbs include 'se' (or me/te/nous/vous/se) — the action reflects back on the subject. 'Je me lève' = I get (myself) up. Common reflexives: se laver, s'habiller, se coucher.","t":"Reflexive verbs are examined in WASSCE composition and grammar. Present tense: Je me lève, tu te lèves, il se lève. Passé composé with être: Je me suis levé(e). Agreement rule: past participle agrees with the subject."},
 ],"hard":[
  {"q":"In the sentence 'Je lui ai donné le livre', what does 'lui' refer to?","o":["An indirect object pronoun (to him/her)","A direct object pronoun (him/her)","A reflexive pronoun","A possessive adjective"],"a":0,"e":"'Lui' is an indirect object pronoun meaning 'to him' or 'to her'. 'Je lui ai donné le livre' = I gave the book to him/her. Direct object would be 'le/la/les'. Indirect = lui/leur.","t":"Object pronoun placement: the most tested WASSCE advanced grammar point. Order: ne + me/te/se/nous/vous + le/la/les + lui/leur + y + en + verb + pas. Write this order on the board and memorise it like the periodic table."},
  {"q":"'Bien qu'il soit riche, il n'est pas heureux' — the subjunctive is used because of:","o":["'Bien que' (although) which always triggers the subjunctive","The word 'riche' requiring a special mood","The negative 'n'est pas' requiring subjunctive","The adverb 'bien' triggering a special form"],"a":0,"e":"'Bien que' (although/even though) is a concessive conjunction that always triggers the subjunctive. 'Soit' is the subjunctive of 'être'. The indicative would be 'est', but after 'bien que' it must be 'soit'.","t":"Concessive conjunctions + subjunctive: bien que, quoique, encore que. Contrast with 'parce que' and 'puisque' which use indicative. WASSCE Paper 2 essays are marked for correct use of these structures."},
 ]},

}

PRAISE = ["✨ Excellent! Teacher Pehpeh is proud!","🔥 You're on FIRE!","⭐ That's the pepper spirit!","💪 Outstanding! Getting stronger!","🎯 Sharp like pepper!","✨ Brilliant! Hard work pays off!","🏆 You nailed it!","🌍 The village would be proud!"]
ENCOURAGE = ["📖 Not quite — every mistake teaches!","💡 Close! Read explanation, try next one!","🌱 Even the tallest palm started as seed. Keep growing!","🤝 No worry! Let's learn together."]

WASSCE_TIPS = """📝 WASSCE EXAM STRATEGY:\n\n1. ANSWER SHEET: HB pencil only. Shade completely. Erase cleanly. Check numbers match.\n2. ELIMINATION: Read ALL options. Cross out wrong ones. 'Always'/'never' usually wrong.\n3. TIME: Paper 1: ~1 min/question. Paper 2: start easiest. Leave 10 min to check.\n4. NIGHT BEFORE: Review only. Eat well, sleep early. Rested brain > tired cramming."""

# === DROPDOWNS ===
REGIONS={"Urban":"urban","Suburban":"suburban","Rural":"rural"}
# Sub-Saharan African countries only
COUNTRIES=["Liberia","Sierra Leone","Ghana","Nigeria","Kenya","Uganda","Tanzania","Ethiopia","Senegal","Cameroon","Gambia","Guinea","Côte d'Ivoire","Mali","Burkina Faso","Rwanda","Malawi","Zambia","Zimbabwe","Mozambique","South Africa","Botswana","Namibia","DRC","Angola","Togo","Benin","Niger","Chad","Somalia","Eritrea","Djibouti","South Sudan","Sudan","Central African Republic","Republic of Congo","Gabon","Equatorial Guinea","São Tomé and Príncipe","Cape Verde","Comoros","Madagascar","Mauritius","Seychelles","Eswatini","Lesotho","Burundi","Guinea-Bissau"]
FLAGS={"Liberia":"🇱🇷","Sierra Leone":"🇸🇱","Ghana":"🇬🇭","Nigeria":"🇳🇬","Kenya":"🇰🇪","Uganda":"🇺🇬","Tanzania":"🇹🇿","Ethiopia":"🇪🇹","Senegal":"🇸🇳","Cameroon":"🇨🇲","Gambia":"🇬🇲","Guinea":"🇬🇳","Côte d'Ivoire":"🇨🇮","Mali":"🇲🇱","Burkina Faso":"🇧🇫","Rwanda":"🇷🇼","Malawi":"🇲🇼","Zambia":"🇿🇲","Zimbabwe":"🇿🇼","Mozambique":"🇲🇿","South Africa":"🇿🇦","Botswana":"🇧🇼","Namibia":"🇳🇦","DRC":"🇨🇩","Angola":"🇦🇴","Togo":"🇹🇬","Benin":"🇧🇯","Niger":"🇳🇪","Chad":"🇹🇩","Somalia":"🇸🇴","Eritrea":"🇪🇷","Djibouti":"🇩🇯","South Sudan":"🇸🇸","Sudan":"🇸🇩","Central African Republic":"🇨🇫","Republic of Congo":"🇨🇬","Gabon":"🇬🇦","Equatorial Guinea":"🇬🇶","São Tomé and Príncipe":"🇸🇹","Cape Verde":"🇨🇻","Comoros":"🇰🇲","Madagascar":"🇲🇬","Mauritius":"🇲🇺","Seychelles":"🇸🇨","Eswatini":"🇸🇿","Lesotho":"🇱🇸","Burundi":"🇧🇮","Guinea-Bissau":"🇬🇼"}
FLAG_CODES={"Liberia":"lr","Sierra Leone":"sl","Ghana":"gh","Nigeria":"ng","Kenya":"ke","Uganda":"ug","Tanzania":"tz","Ethiopia":"et","Senegal":"sn","Cameroon":"cm","Gambia":"gm","Guinea":"gn","Côte d'Ivoire":"ci","Mali":"ml","Burkina Faso":"bf","Rwanda":"rw","Malawi":"mw","Zambia":"zm","Zimbabwe":"zw","Mozambique":"mz","South Africa":"za","Botswana":"bw","Namibia":"na","DRC":"cd","Angola":"ao","Togo":"tg","Benin":"bj","Niger":"ne","Chad":"td","Somalia":"so","Eritrea":"er","Djibouti":"dj","South Sudan":"ss","Sudan":"sd","Central African Republic":"cf","Republic of Congo":"cg","Gabon":"ga","Equatorial Guinea":"gq","São Tomé and Príncipe":"st","Cape Verde":"cv","Comoros":"km","Madagascar":"mg","Mauritius":"mu","Seychelles":"sc","Eswatini":"sz","Lesotho":"ls","Burundi":"bi","Guinea-Bissau":"gw"}
GRADES=["9th Grade","10th Grade","11th Grade","12th Grade","12th Grade (WASSCE Prep)"]
SUBJECTS=["Mathematics","English Language","Social Studies","Physics","Chemistry","Biology","Economics","Government / Civics","Literature in English","History","Geography","Agriculture","French","Religious Studies","Business Management","Accounting","Computer Studies / ICT","Physical Education"]
TOPICS={"Mathematics":["Number and Numeration","Fractions and Decimals","Percentages","Ratio and Proportion","Algebraic Expressions","Linear Equations","Quadratic Equations","Simultaneous Equations","Sets and Venn Diagrams","Trigonometry","Mensuration","Geometry","Statistics","Probability","Vectors","Logarithms","Indices and Surds"],
"English Language":["Comprehension","Summary Writing","Essay (Narrative)","Essay (Argumentative)","Letter Writing (Formal)","Parts of Speech","Tenses","Active/Passive Voice","Punctuation","Vocabulary","Idioms"],
"Physics":["Measurement","Motion","Newton's Laws","Work Energy Power","Simple Machines","Pressure","Heat Transfer","Gas Laws","Waves","Sound","Light","Electricity","Ohm's Law"],
"Chemistry":["States of Matter","Atomic Structure","Periodic Table","Chemical Bonding","Reactions","Acids Bases Salts","Electrolysis","Organic Chemistry","Mole Concept"],
"Biology":["Cell Structure","Cell Division","Photosynthesis","Human Body Systems","Reproduction","Genetics","Evolution","Ecology","Diseases and Immunity"],
}
DEF_TOPICS=["Core Concepts","Key Terms","Applications","Review","Exam Practice"]
FR_DEF_TOPICS=["Concepts fondamentaux","Termes clés","Applications","Révision","Préparation examen"]
SW_DEF_TOPICS=["Dhana za Msingi","Maneno Muhimu","Matumizi","Mapitio","Mazoezi ya Mtihani"]
# Topic translations: {English subject: {English topic: {fr:..., sw:...}}}
FR_TOPICS={"Mathematics":["Nombres et numération","Fractions et décimales","Pourcentages","Rapport et proportion","Expressions algébriques","Équations linéaires","Équations quadratiques","Équations simultanées","Ensembles et diagrammes de Venn","Trigonométrie","Mensuration","Géométrie","Statistiques","Probabilité","Vecteurs","Logarithmes","Indices et radicaux"],
"English Language":["Compréhension","Résumé","Rédaction (narratif)","Rédaction (argumentatif)","Lettre formelle","Parties du discours","Temps verbaux","Voix active/passive","Ponctuation","Vocabulaire","Expressions idiomatiques"],
"Physics":["Mesure","Mouvement","Lois de Newton","Travail Énergie Puissance","Machines simples","Pression","Transfert thermique","Lois des gaz","Ondes","Son","Lumière","Électricité","Loi d'Ohm"],
"Chemistry":["États de la matière","Structure atomique","Tableau périodique","Liaisons chimiques","Réactions","Acides Bases Sels","Électrolyse","Chimie organique","Concept de mole"],
"Biology":["Structure cellulaire","Division cellulaire","Photosynthèse","Systèmes du corps humain","Reproduction","Génétique","Évolution","Écologie","Maladies et immunité"],
}
SW_TOPICS={"Mathematics":["Nambari na Uhesabuji","Sehemu na Desimali","Asilimia","Uwiano","Misemo ya Aljebra","Mlingano wa Mstari","Mlingano wa Pili","Mlingano Sawia","Seti na Michoro ya Venn","Trigonometria","Vipimo","Jiometri","Takwimu","Uwezekano","Vekta","Logarithimu","Indeksi na Mizizi"],
"English Language":["Ufahamu","Muhtasari","Insha (Masimulizi)","Insha (Hoja)","Uandishi Barua Rasmi","Sehemu za Hotuba","Nyakati","Sauti Tendaji/Tendewa","Alama za Uandishi","Msamiati","Nahau"],
"Physics":["Vipimo","Mwendo","Sheria za Newton","Kazi Nishati Nguvu","Mashine Rahisi","Shinikizo","Uhamishaji Joto","Sheria za Gesi","Mawimbi","Sauti","Mwanga","Umeme","Sheria ya Ohm"],
"Chemistry":["Hali za Maada","Muundo wa Atomu","Jedwali la Vipindi","Vifungo vya Kemikali","Athari","Asidi Besi Chumvi","Elektrolisisi","Kemia Hai","Dhana ya Moli"],
"Biology":["Muundo wa Seli","Mgawanyiko wa Seli","Usanisinuru","Mifumo ya Mwili","Uzazi","Jenetiki","Mageuzi","Ikolojia","Magonjwa na Kinga"],
}
# Build topic display→English maps
_TOPIC_TO_EN={}
for subj_en,en_topics in TOPICS.items():
    for fr_t,en_t in zip(FR_TOPICS.get(subj_en,[]),en_topics): _TOPIC_TO_EN[fr_t]=en_t
    for sw_t,en_t in zip(SW_TOPICS.get(subj_en,[]),en_topics): _TOPIC_TO_EN[sw_t]=en_t
for fr_t,en_t in zip(FR_DEF_TOPICS,DEF_TOPICS): _TOPIC_TO_EN[fr_t]=en_t
for sw_t,en_t in zip(SW_DEF_TOPICS,DEF_TOPICS): _TOPIC_TO_EN[sw_t]=en_t
def _get_topics(subj_en):
    lk=_lang_key()
    if lk=="fr": return FR_TOPICS.get(subj_en,FR_DEF_TOPICS)
    if lk=="sw": return SW_TOPICS.get(subj_en,SW_DEF_TOPICS)
    return TOPICS.get(subj_en,DEF_TOPICS)
def _to_en_topic(t): return _TOPIC_TO_EN.get(t,t)
# Quiz subject translations
_QUIZ_SUBJ_FR={"Mathematics":"Mathématiques","Physics":"Physique","Biology":"Biologie","Chemistry":"Chimie","Reading Comprehension":"Compréhension écrite","Literature":"Littérature","Economics":"Économie","Geography":"Géographie","French":"Français"}
_QUIZ_SUBJ_SW={"Mathematics":"Hisabati","Physics":"Fizikia","Biology":"Biolojia","Chemistry":"Kemia","Reading Comprehension":"Ufahamu wa Kusoma","Literature":"Fasihi","Economics":"Uchumi","Geography":"Jiografia","French":"Kifaransa"}
_QUIZ_SUBJ_TO_EN={v:k for k,v in _QUIZ_SUBJ_FR.items()}
_QUIZ_SUBJ_TO_EN.update({v:k for k,v in _QUIZ_SUBJ_SW.items()})
def _quiz_subjects():
    lk=_lang_key()
    if lk=="fr": return [_QUIZ_SUBJ_FR.get(k,k) for k in QUIZ.keys()]
    if lk=="sw": return [_QUIZ_SUBJ_SW.get(k,k) for k in QUIZ.keys()]
    return list(QUIZ.keys())
def _quiz_subj_en(s): return _QUIZ_SUBJ_TO_EN.get(s,s)
# Francophone countries for auto-switch
FRANCOPHONE={"Côte d'Ivoire","Senegal","Cameroon","Mali","Burkina Faso","Niger","Chad","Guinea","Togo","Benin","Central African Republic","Republic of Congo","Gabon","Equatorial Guinea","DRC","Djibouti","Comoros","Madagascar","Burundi","Guinea-Bissau","São Tomé and Príncipe"}
# Swahili-speaking countries
SWAHILI_COUNTRIES={"Kenya","Tanzania","Uganda","DRC","Rwanda","Burundi"}
TASKS={"Lesson Plan":"detailed lesson plan","Quiz (10 Q)":"10-question quiz with answer key","Quiz (20 Q)":"20-question quiz","WASSCE MCQ (50)":"50 WASSCE-style MCQs","WASSCE Theory":"WASSCE theory questions","BECE Exam":"BECE-style exam","Homework":"homework with minimal resources","Group Activity":"group activity","Reading Comprehension":"reading passage with questions","No-Lab Practical":"hands-on zero-cost activity","Rubric":"grading rubric","Strategy Guide":"teaching strategies","Weekly Scheme":"5-day scheme of work","Term Scheme":"term plan","Remedial Material":"catch-up material","Study Notes":"revision guide","Educational Game":"zero-cost teaching game","Illustrated Lesson (AI image)":"lesson with AI-generated visual"}
SIZES={"Small (<25)":"<25 students","Medium (25-40)":"25-40","Large (40-60)":"40-60","Very Large (60+)":"60+"}
RESOURCES={"Chalkboard only":"chalkboard/chalk only","+ shared textbooks":"chalkboard + shared textbooks","+ handouts":"+ printable handouts","Computer/projector":"occasional tech","Phones/tablets":"student devices","Well-equipped":"regular tech"}
LANGS={"English":"English","Français":"French","Kiswahili":"Swahili"}
# UI Translations
UI_TEXT={
 "en":{
  "generate":"📋 Generate","chat":"💬 Chat","quiz":"Quiz","students":"🧑‍🎓 Students",
  "task":"Task","time":"Time","topic":"Topic","options":"Options","subject":"Subject","grade":"Grade",
  "country":"Country","setting":"Setting","class_size":"Class Size",
  "language":"Language","student_level":"Student Level","school_name":"🏫 School Name",
  "school_placeholder":"e.g., Bahn, St. Martin's","my_classroom":"My Classroom","my_students":"My Students",
  "gen_btn":"Generate","clear":"🗑️ Clear","hear":"🔊 Hear Results","grade_work":"Grade Work",
  "grade_btn":"Grade","students_work":"Student's work:","offline_title":"📴 Offline — Practice Quiz",
  "offline_msg":"No internet? These quizzes work offline!","practice_quiz":"Practice Quiz",
  "adaptive":"Adaptive. Works offline too!","score":"Score","level":"Level","next":"➡️ Next",
  "reset":"🔄 Reset","wassce_tips":"📝 WASSCE Tips","add_student":"Add Student","name":"Name",
  "upload_excel":"📤 Upload Excel","lit_library":"📚 Literature Library",
  "lit_desc":"Select a novel for passage-based comprehension exercises","select_book":"📖 Select Book",
  "comp_type":"Comprehension Type","include_img":"🎨 Include AI illustration",
  "img_help":"Generates a visual aid using DALL-E or Google Imagen",
  "ask_about":"Ask about","draw_hint":"(start with 'draw' for images)",
  "mic_hint":"🎤 Tap mic to speak instead of typing","heard":"🎤 Heard",
  "email_result":"📧 Email / Download this result","recheck":"🔄 Re-check",
  "voice_ready":"🔊 Voice Ready","no_models":"🤖 No models","offline":"OFFLINE",
  "generating":"🔊 Generating audio...","transcribing":"🎤 Transcribing your voice...",
  "audio_failed":"Audio failed","try_again":"Try again",
  "assignments":"Personalized Assignments","risk_flags":"Risk & Support Flags",
  "no_students":"No students added yet","single_period":"Single period (30-40 min)",
  "double":"Double (60-80 min)","half_day":"Half day","full_day":"Full day","weekly":"Weekly","na":"N/A",
  "differentiation":"Differentiation","formative":"Formative assessment","takehome":"Take-home activity",
  "wassce_align":"WASSCE alignment","local_ex":"Local examples","literacy":"Literacy integration",
  "large_class":"Large-class strategies","cross_curr":"Cross-curricular","ai_visual":"AI visual aid",
  "pass_short":"Passage + Short Answer Questions","pass_fill":"Passage + Fill in the Blanks",
  "pass_essay":"Passage + Essay Prompt","pass_mcq":"Passage + MCQ",
  "pass_vocab":"Passage + Vocabulary Exercise","full_comp":"Full Comprehension (All Types)",
  "assignment":"📝 Assignment","risk":"⚠️ Risk Analysis","creating":"Creating...",
  "see_all":"📋 See all","model_responses":"model responses","streak":"streak",
  "generating_content":"⏳ Teacher Pehpeh is cooking...","done":"✅ Done! Content is ready!",
  "ask_tp":"Ask Teacher Pehpeh","thinking":"⏳ Teacher Pehpeh is thinking...","response_ready":"✅ Response ready!",
  "asking_claude":"🟣 Asking Claude...","asking_chatgpt":"🟢 Asking ChatGPT...","asking_gemini":"🔵 Asking Gemini...",
  "chat_ex1":"How to teach fractions with no textbooks?","chat_ex2":"My students keep failing WASSCE.","chat_ex3":"Managing 60+ students?",
  "combining":"🔀 Combining the best...","creating_img":"🎨 Creating illustration...",
  "photo_hint":"📸 Upload or snap a photo to ask Teacher Pehpeh about it",
  "photo_upload":"Upload image","photo_camera":"📷 Take photo",
  "photo_ask":"What would you like to know about this photo?",
  "photo_default":"Please analyze this image and explain what you see. Provide educational context relevant to the subject.",
  "analyzing_photo":"📸 Analyzing your photo...","start_recording":"Start Recording","stop_recording":"Stop Recording",
  "transcribing_handwriting":"📝 Transcribing handwriting to Word document...",
  "transcribe_done":"📝 Transcription complete! Download your Word document below.",
  "download_docx":"📥 Download Word Document (.docx)",
 },
 "fr":{
  "generate":"📋 Générer","chat":"💬 Discussion","quiz":"Quiz","students":"🧑‍🎓 Élèves",
  "task":"Tâche","time":"Durée","topic":"Sujet","options":"Options","subject":"Matière","grade":"Classe",
  "country":"Pays","setting":"Contexte","class_size":"Taille de classe",
  "language":"Langue","student_level":"Niveau des élèves","school_name":"🏫 Nom de l'école",
  "school_placeholder":"ex: Bahn, St. Martin's","my_classroom":"Ma Classe","my_students":"Mes Élèves",
  "gen_btn":"Générer","clear":"🗑️ Effacer","hear":"🔊 Écouter","grade_work":"Noter le travail",
  "grade_btn":"Noter","students_work":"Travail de l'élève :","offline_title":"📴 Hors ligne — Quiz pratique",
  "offline_msg":"Pas d'internet ? Ces quiz fonctionnent hors ligne !","practice_quiz":"Quiz pratique",
  "adaptive":"Adaptatif. Fonctionne hors ligne aussi !","score":"Score","level":"Niveau","next":"➡️ Suivant",
  "reset":"🔄 Réinitialiser","wassce_tips":"📝 Conseils WASSCE","add_student":"Ajouter un élève","name":"Nom",
  "upload_excel":"📤 Télécharger Excel","lit_library":"📚 Bibliothèque littéraire",
  "lit_desc":"Sélectionnez un roman pour des exercices de compréhension","select_book":"📖 Choisir un livre",
  "comp_type":"Type de compréhension","include_img":"🎨 Inclure une illustration IA",
  "img_help":"Génère une aide visuelle avec DALL-E ou Google Imagen",
  "ask_about":"Posez une question sur","draw_hint":"(commencez par 'dessiner' pour images)",
  "mic_hint":"🎤 Appuyez sur le micro pour parler","heard":"🎤 Entendu",
  "email_result":"📧 Envoyer / Télécharger ce résultat","recheck":"🔄 Vérifier",
  "voice_ready":"🔊 Voix prête","no_models":"🤖 Aucun modèle","offline":"HORS LIGNE",
  "generating":"🔊 Génération audio...","transcribing":"🎤 Transcription en cours...",
  "audio_failed":"Échec audio","try_again":"Réessayer",
  "assignments":"Devoirs personnalisés","risk_flags":"Signaux de risque et soutien",
  "no_students":"Aucun élève ajouté","single_period":"Période simple (30-40 min)",
  "double":"Double (60-80 min)","half_day":"Demi-journée","full_day":"Journée complète","weekly":"Hebdomadaire","na":"N/A",
  "differentiation":"Différenciation","formative":"Évaluation formative","takehome":"Activité à emporter",
  "wassce_align":"Alignement WASSCE","local_ex":"Exemples locaux","literacy":"Intégration de la lecture",
  "large_class":"Stratégies grande classe","cross_curr":"Interdisciplinaire","ai_visual":"Aide visuelle IA",
  "pass_short":"Passage + Questions courtes","pass_fill":"Passage + Texte à trous",
  "pass_essay":"Passage + Sujet de rédaction","pass_mcq":"Passage + QCM",
  "pass_vocab":"Passage + Exercice de vocabulaire","full_comp":"Compréhension complète (Tous types)",
  "assignment":"📝 Devoir","risk":"⚠️ Analyse de risque","creating":"Création en cours...",
  "see_all":"📋 Voir tous les","model_responses":"réponses des modèles","streak":"série",
  "generating_content":"⏳ Teacher Pehpeh prépare...","done":"✅ Terminé ! Le contenu est prêt !",
  "ask_tp":"Demandez à Teacher Pehpeh","thinking":"⏳ Teacher Pehpeh réfléchit...","response_ready":"✅ Réponse prête !",
  "asking_claude":"🟣 Consultation de Claude...","asking_chatgpt":"🟢 Consultation de ChatGPT...","asking_gemini":"🔵 Consultation de Gemini...",
  "chat_ex1":"Comment enseigner les fractions sans manuels ?","chat_ex2":"Mes élèves échouent au WASSCE.","chat_ex3":"Gérer plus de 60 élèves ?",
  "combining":"🔀 Combinaison des meilleurs...","creating_img":"🎨 Création d'illustration...",
  "photo_hint":"📸 Téléchargez ou prenez une photo pour demander à Teacher Pehpeh",
  "photo_upload":"Télécharger image","photo_camera":"📷 Prendre une photo",
  "photo_ask":"Que voulez-vous savoir sur cette photo ?",
  "photo_default":"Veuillez analyser cette image et expliquer ce que vous voyez. Fournissez un contexte éducatif pertinent.",
  "analyzing_photo":"📸 Analyse de votre photo...","start_recording":"Commencer","stop_recording":"Arrêter",
  "transcribing_handwriting":"📝 Transcription en document Word...",
  "transcribe_done":"📝 Transcription terminée ! Téléchargez votre document Word ci-dessous.",
  "download_docx":"📥 Télécharger le document Word (.docx)",
 },
 "sw":{
  "generate":"📋 Tengeneza","chat":"💬 Mazungumzo","quiz":"Maswali","students":"🧑‍🎓 Wanafunzi",
  "task":"Kazi","time":"Muda","topic":"Mada","options":"Chaguzi","subject":"Somo","grade":"Darasa",
  "country":"Nchi","setting":"Mazingira","class_size":"Ukubwa wa darasa",
  "language":"Lugha","student_level":"Kiwango cha wanafunzi","school_name":"🏫 Jina la shule",
  "school_placeholder":"k.m., Bahn, St. Martin's","my_classroom":"Darasa Langu","my_students":"Wanafunzi Wangu",
  "gen_btn":"Tengeneza","clear":"🗑️ Futa","hear":"🔊 Sikiliza","grade_work":"Sahihisha Kazi",
  "grade_btn":"Sahihisha","students_work":"Kazi ya mwanafunzi:","offline_title":"📴 Nje ya mtandao — Maswali ya mazoezi",
  "offline_msg":"Hakuna mtandao? Maswali haya yanafanya kazi nje ya mtandao!","practice_quiz":"Maswali ya Mazoezi",
  "adaptive":"Yanabadilika. Yanafanya kazi nje ya mtandao pia!","score":"Alama","level":"Kiwango","next":"➡️ Ifuatayo",
  "reset":"🔄 Weka upya","wassce_tips":"📝 Vidokezo vya WASSCE","add_student":"Ongeza mwanafunzi","name":"Jina",
  "upload_excel":"📤 Pakia Excel","lit_library":"📚 Maktaba ya Vitabu",
  "lit_desc":"Chagua riwaya kwa mazoezi ya ufahamu wa kusoma","select_book":"📖 Chagua Kitabu",
  "comp_type":"Aina ya ufahamu","include_img":"🎨 Jumuisha mchoro wa AI",
  "img_help":"Inatengeneza msaada wa kuona kwa DALL-E au Google Imagen",
  "ask_about":"Uliza kuhusu","draw_hint":"(anza na 'chora' kwa picha)",
  "mic_hint":"🎤 Bonyeza maikrofoni kusema badala ya kuandika","heard":"🎤 Imesikika",
  "email_result":"📧 Tuma / Pakua matokeo haya","recheck":"🔄 Angalia tena",
  "voice_ready":"🔊 Sauti tayari","no_models":"🤖 Hakuna modeli","offline":"NJE YA MTANDAO",
  "generating":"🔊 Inatengeneza sauti...","transcribing":"🎤 Inabadilisha sauti kuwa maandishi...",
  "audio_failed":"Sauti imeshindwa","try_again":"Jaribu tena",
  "assignments":"Kazi za kibinafsi","risk_flags":"Ishara za hatari na msaada",
  "no_students":"Hakuna wanafunzi walioongezwa","single_period":"Kipindi kimoja (dak 30-40)",
  "double":"Mara mbili (dak 60-80)","half_day":"Nusu siku","full_day":"Siku nzima","weekly":"Kila wiki","na":"H/T",
  "differentiation":"Utofautishaji","formative":"Tathmini ya mchakato","takehome":"Kazi ya nyumbani",
  "wassce_align":"Ulinganifu wa WASSCE","local_ex":"Mifano ya mahali","literacy":"Ujumuishaji wa kusoma",
  "large_class":"Mikakati ya darasa kubwa","cross_curr":"Mtambuka","ai_visual":"Msaada wa kuona wa AI",
  "pass_short":"Kifungu + Maswali mafupi","pass_fill":"Kifungu + Jaza nafasi",
  "pass_essay":"Kifungu + Swali la insha","pass_mcq":"Kifungu + Maswali ya kuchagua",
  "pass_vocab":"Kifungu + Zoezi la msamiati","full_comp":"Ufahamu kamili (Aina zote)",
  "assignment":"📝 Kazi","risk":"⚠️ Uchambuzi wa hatari","creating":"Inatengeneza...",
  "see_all":"📋 Tazama zote","model_responses":"majibu ya modeli","streak":"mfululizo",
  "generating_content":"⏳ Teacher Pehpeh anapika...","done":"✅ Imekamilika! Maudhui yako tayari!",
  "ask_tp":"Muulize Teacher Pehpeh","thinking":"⏳ Teacher Pehpeh anafikiria...","response_ready":"✅ Jibu liko tayari!",
  "asking_claude":"🟣 Kuuliza Claude...","asking_chatgpt":"🟢 Kuuliza ChatGPT...","asking_gemini":"🔵 Kuuliza Gemini...",
  "chat_ex1":"Jinsi ya kufundisha sehemu bila vitabu?","chat_ex2":"Wanafunzi wangu wanashindwa WASSCE.","chat_ex3":"Kusimamia wanafunzi 60+?",
  "combining":"🔀 Kuchanganya bora zaidi...","creating_img":"🎨 Kuunda mchoro...",
  "photo_hint":"📸 Pakia au piga picha kumuuliza Teacher Pehpeh",
  "photo_upload":"Pakia picha","photo_camera":"📷 Piga picha",
  "photo_ask":"Unataka kujua nini kuhusu picha hii?",
  "photo_default":"Tafadhali changanua picha hii na ueleze unachokiona. Toa muktadha wa kielimu unaofaa.",
  "analyzing_photo":"📸 Kuchambua picha yako...","start_recording":"Anza Kurekodi","stop_recording":"Simamisha",
  "transcribing_handwriting":"📝 Kunakili maandishi ya mkono kuwa hati ya Word...",
  "transcribe_done":"📝 Unakili umekamilika! Pakua hati yako ya Word hapa chini.",
  "download_docx":"📥 Pakua Hati ya Word (.docx)",
 }
}
def _lang_key():
    """Get current UI language key from session state."""
    lk=st.session_state.get("lang_sel","English")
    if "Français" in lk or "French" in lk: return "fr"
    if "Kiswahili" in lk or "Swahili" in lk: return "sw"
    return "en"
def T(key):
    """Get translated UI string."""
    lk=_lang_key()
    return UI_TEXT.get(lk,UI_TEXT["en"]).get(key,UI_TEXT["en"].get(key,key))
ABILITY={"Mixed":"mixed-ability","Struggling":"below grade level","On level":"at expected level","Advanced":"needs challenge","Inclusive":"includes learning differences"}
TIMES=["Single period (30-40 min)","Double (60-80 min)","Half day","Full day","Weekly","N/A"]
EXTRAS=["Differentiation","Formative assessment","Take-home activity","WASSCE alignment","Local examples","Literacy integration","Large-class strategies","Cross-curricular","AI visual aid"]

# French dropdown translations (display→English value for AI)
FR_TASKS={"Plan de cours":"detailed lesson plan","Quiz (10 Q)":"10-question quiz with answer key","Quiz (20 Q)":"20-question quiz","QCM WASSCE (50)":"50 WASSCE-style MCQs","Théorie WASSCE":"WASSCE theory questions","Examen BECE":"BECE-style exam","Devoirs":"homework with minimal resources","Activité de groupe":"group activity","Compréhension écrite":"reading passage with questions","Pratique sans labo":"hands-on zero-cost activity","Grille d'évaluation":"grading rubric","Guide stratégique":"teaching strategies","Plan hebdomadaire":"5-day scheme of work","Plan trimestriel":"term plan","Rattrapage":"catch-up material","Notes de révision":"revision guide","Jeu éducatif":"zero-cost teaching game","Leçon illustrée (image IA)":"lesson with AI-generated visual"}
FR_GRADES=["9e année","10e année","11e année","12e année","12e année (WASSCE Prép)"]
FR_SUBJECTS=["Mathématiques","Langue anglaise","Sciences intégrées","Études sociales","Physique","Chimie","Biologie","Économie","Gouvernement / Éducation civique","Littérature anglaise","Histoire","Géographie","Agriculture","Français","Études religieuses","Gestion des affaires","Comptabilité","Informatique / TIC","Éducation physique","Art / Arts créatifs","Musique"]
FR_SIZES={"Petit (<25)":"<25 students","Moyen (25-40)":"25-40","Grand (40-60)":"40-60","Très grand (60+)":"60+"}
FR_RESOURCES={"Tableau noir seul":"chalkboard/chalk only","+ manuels partagés":"chalkboard + shared textbooks","+ polycopiés":"+ printable handouts","Ordinateur/projecteur":"occasional tech","Téléphones/tablettes":"student devices","Bien équipé":"regular tech"}
FR_ABILITY={"Mixte":"mixed-ability","En difficulté":"below grade level","Au niveau":"at expected level","Avancé":"needs challenge","Inclusif":"includes learning differences"}
FR_TIMES=["Période simple (30-40 min)","Double (60-80 min)","Demi-journée","Journée complète","Hebdomadaire","N/A"]
FR_EXTRAS=["Différenciation","Évaluation formative","Activité à emporter","Alignement WASSCE","Exemples locaux","Intégration lecture","Stratégies grande classe","Interdisciplinaire","Aide visuelle IA"]
FR_REGIONS={"Urbain":"urban","Banlieue":"suburban","Rural":"rural"}
# Swahili dropdown translations
SW_TASKS={"Mpango wa Somo":"detailed lesson plan","Maswali (10)":"10-question quiz with answer key","Maswali (20)":"20-question quiz","QCM WASSCE (50)":"50 WASSCE-style MCQs","Nadharia WASSCE":"WASSCE theory questions","Mtihani BECE":"BECE-style exam","Kazi ya Nyumbani":"homework with minimal resources","Shughuli ya Kikundi":"group activity","Ufahamu wa Kusoma":"reading passage with questions","Mazoezi bila Maabara":"hands-on zero-cost activity","Rubriiki":"grading rubric","Mwongozo wa Mkakati":"teaching strategies","Mpango wa Wiki":"5-day scheme of work","Mpango wa Muhula":"term plan","Nyenzo za Kufidia":"catch-up material","Muhtasari wa Masomo":"revision guide","Mchezo wa Kielimu":"zero-cost teaching game","Somo Lenye Mchoro (picha AI)":"lesson with AI-generated visual"}
SW_GRADES=["Darasa la 9","Darasa la 10","Darasa la 11","Darasa la 12","Darasa la 12 (WASSCE)"]
SW_SUBJECTS=["Hisabati","Lugha ya Kiingereza","Sayansi Jumuishi","Maarifa ya Jamii","Fizikia","Kemia","Biolojia","Uchumi","Serikali / Uraia","Fasihi ya Kiingereza","Historia","Jiografia","Kilimo","Kifaransa","Masomo ya Dini","Usimamizi wa Biashara","Uhasibu","Kompyuta / TEHAMA","Uchora Ufundi","Elimu ya Mwili","Sanaa / Sanaa Bunifu","Muziki"]
SW_SIZES={"Ndogo (<25)":"<25 students","Wastani (25-40)":"25-40","Kubwa (40-60)":"40-60","Kubwa sana (60+)":"60+"}
SW_RESOURCES={"Ubao tu":"chalkboard/chalk only","+ vitabu vya kushiriki":"chalkboard + shared textbooks","+ nakala":"+ printable handouts","Kompyuta/projekta":"occasional tech","Simu/tableti":"student devices","Vifaa kamili":"regular tech"}
SW_ABILITY={"Mchanganyiko":"mixed-ability","Wanaoshindwa":"below grade level","Kiwango sahihi":"at expected level","Wenye uwezo":"needs challenge","Jumuishi":"includes learning differences"}
SW_TIMES=["Kipindi kimoja (dak 30-40)","Mara mbili (dak 60-80)","Nusu siku","Siku nzima","Kila wiki","H/T"]
SW_EXTRAS=["Utofautishaji","Tathmini ya mchakato","Kazi ya nyumbani","Ulinganifu WASSCE","Mifano ya mahali","Ujumuishaji kusoma","Mikakati darasa kubwa","Mtambuka","Msaada wa kuona AI"]
SW_REGIONS={"Mjini":"urban","Pembezoni":"suburban","Vijijini":"rural"}
# Grade/Subject mappings (French display→English value)
_GRADE_MAP={**dict(zip(FR_GRADES,GRADES)),**dict(zip(SW_GRADES,GRADES))}
_SUBJ_MAP={**dict(zip(FR_SUBJECTS,SUBJECTS)),**dict(zip(SW_SUBJECTS,SUBJECTS))}
_TIME_MAP={**dict(zip(FR_TIMES,TIMES)),**dict(zip(SW_TIMES,TIMES))}
def _tasks(): lk=_lang_key(); return FR_TASKS if lk=="fr" else SW_TASKS if lk=="sw" else TASKS
def _grades(): lk=_lang_key(); return FR_GRADES if lk=="fr" else SW_GRADES if lk=="sw" else GRADES
def _subjects(): lk=_lang_key(); return FR_SUBJECTS if lk=="fr" else SW_SUBJECTS if lk=="sw" else SUBJECTS
def _sizes(): lk=_lang_key(); return FR_SIZES if lk=="fr" else SW_SIZES if lk=="sw" else SIZES
def _resources(): lk=_lang_key(); return FR_RESOURCES if lk=="fr" else SW_RESOURCES if lk=="sw" else RESOURCES
def _ability(): lk=_lang_key(); return FR_ABILITY if lk=="fr" else SW_ABILITY if lk=="sw" else ABILITY
def _times(): lk=_lang_key(); return FR_TIMES if lk=="fr" else SW_TIMES if lk=="sw" else TIMES
def _extras(): lk=_lang_key(); return FR_EXTRAS if lk=="fr" else SW_EXTRAS if lk=="sw" else EXTRAS
def _regions(): lk=_lang_key(); return FR_REGIONS if lk=="fr" else SW_REGIONS if lk=="sw" else REGIONS
def _to_en_grade(g): return _GRADE_MAP.get(g,g)
def _to_en_subj(s): return _SUBJ_MAP.get(s,s)

# Literature library for Reading Comprehension
LITERATURE={
    "Julius Caesar — William Shakespeare":{"author":"William Shakespeare","origin":"England","genre":"Tragedy/Drama","themes":"Power, betrayal, ambition, loyalty, rhetoric","wassce":"WASSCE Literature set text"},
    "Things Fall Apart — Chinua Achebe":{"author":"Chinua Achebe","origin":"Nigeria","genre":"Novel","themes":"Colonialism, tradition vs change, masculinity, cultural identity","wassce":"WASSCE Literature set text"},
    "The African Child — Camara Laye":{"author":"Camara Laye","origin":"Guinea","genre":"Autobiography","themes":"Childhood, tradition, education, family, cultural heritage","wassce":"WASSCE Literature set text"},
    "Weep Not, Child — Ngũgĩ wa Thiong'o":{"author":"Ngũgĩ wa Thiong'o","origin":"Kenya","genre":"Novel","themes":"Education, colonialism, Mau Mau uprising, family, hope","wassce":"WASSCE Literature set text"},
    "Murder in the Cassava Patch — Bediako Asare":{"author":"Bediako Asare","origin":"Ghana","genre":"Novel","themes":"Justice, rural life, mystery, social commentary","wassce":"Common African set text"},
    "Hamlet — William Shakespeare":{"author":"William Shakespeare","origin":"England","genre":"Tragedy","themes":"Revenge, mortality, madness, duty, corruption","wassce":"WASSCE Literature option"},
    "The River Between — Ngũgĩ wa Thiong'o":{"author":"Ngũgĩ wa Thiong'o","origin":"Kenya","genre":"Novel","themes":"Cultural conflict, education, tradition vs modernity, love","wassce":"African literature"},
    "Anthills of the Savannah — Chinua Achebe":{"author":"Chinua Achebe","origin":"Nigeria","genre":"Novel","themes":"Power, corruption, friendship, social justice","wassce":"African literature"},
    "The Beautyful Ones Are Not Yet Born — Ayi Kwei Armah":{"author":"Ayi Kwei Armah","origin":"Ghana","genre":"Novel","themes":"Corruption, morality, post-independence disillusionment","wassce":"African literature"},
    "So Long a Letter — Mariama Bâ":{"author":"Mariama Bâ","origin":"Senegal","genre":"Epistolary novel","themes":"Women's rights, polygamy, friendship, tradition, grief","wassce":"African literature"},
    "Nervous Conditions — Tsitsi Dangarembga":{"author":"Tsitsi Dangarembga","origin":"Zimbabwe","genre":"Novel","themes":"Gender, colonialism, education, identity, poverty","wassce":"African literature"},
    "Half of a Yellow Sun — Chimamanda Ngozi Adichie":{"author":"Chimamanda Ngozi Adichie","origin":"Nigeria","genre":"Novel","themes":"War, love, identity, class, Biafran conflict","wassce":"Modern African literature"},
    "Teacher's Own Selection":{"author":"—","origin":"—","genre":"—","themes":"Teacher provides custom text or topic","wassce":"Custom"},
}

# === LOGO ===
_LOGO_B64="iVBORw0KGgoAAAANSUhEUgAAAZAAAAGQCAYAAACAvzbMAAEAAElEQVR42mT9adCuW3oWhl3rXms97/ft4UzdfU6fnrt1WhJqkISEBGaGYMAydiwwoBQkTpw4VU5wUk5M4iG2N+VUyFTl/IhjV+yESlxgAjY4xMRAjE1BzGwzKJqlllqtnk+fcZ+zv/d51pAf93Xd6/nkrlLpnH323t/7Ps8a7vu6ryF9+NUPzzkn5pwoueA4DuScgZkw0gTGgKUEWMKxN2yXDaMPzDlgOcMsYb/usJyRc8Z+vcIsIxdDbwOlGMaYGGMimQFzICWgj4FaCuYEshmONtDHQMkGzIk5BpCAXCoAoB8dlhOQJtrR/eeVjNE7AMDMMEYHYGjHgVorUs7orSHnhNYGzICUEnofKKWgzw5LGa01lJwxAWTL2O+uyFuFARhzoveGUor/jN4xJmCWMAHMOVBywb4fKDVj8rvmUgAMWDL03jHnhOWM49hRa4VZRmsHMCZyqWi9wVJCzgWtN+RsOI6GUjN669jqBW10zDGQLQEAxpj8zIbWGnLOGGMgWUKxAiTD6A0DE2kCvXfkmlFyRe8HjsP/TM4ZYw7MCQAJYzSUUv09zImUEhISjuNAvVSMPn1NJH8+aQIwIMGwHztKLrBsOPYDwMS2XfzvAXDddwBAqQWjd2Qz/29m/qxbw7ZtmHMCcwJmsJQwx0AfE8kmUvLvW2pFbx3AxBwTpRS01lFr9eeQJlofyMlgOaP3BkuG1jvMMubs6L1jqxvMEo52wKwgmeHYr6h1w5jDP2/O2PcrxgRyzrCUkAC03lFKQU6GPjom18ycA1upmJhozdcyJpAswV9fwr7vqFtFThl9+DrOpeBoOwyGdjTUS0VCQusdic+hlIyJ6eujVJRcgDnQxsDRGkouyDkhYWLfG2asR3+Xc07MeGYbJrTGM9L0PbbvDcZ3k7Px+TTf2ynhaL6Pfd8N5GzIltHHQB/Nnx3Pid47zBIsF8w+MDG4DzvGTLCc+OcLMCda72hHw+XmgnYcSMnAR+Zrrk9/970B/HzH4T+j1IrjumOmiWy+tvvoSCkDGLFPjqPDkp8HOWeklHC9u6Ju1ddK6+jd12Lv/v7SHEDyZwLuiTkHEhJSNvTW/Zlh+hnXOzABwN97zjn2U29+rrbRcblsaK3HGtF+rnVDOxomBsxMC8jP0gRg+llTqp+DEwkpJT+zLaPPASR/z7VW7PsOS+b7Ofl5mfn7tO/248Bl2zDG4Oc2pDlxtMZz1tdeTglICfnRo4dPRvdDYA6/FOac/gW4UMGDFRjc0zrk/UtMTF94GL65LMHMkOCLBJh+KWFiTqDW6hdCMsw5MOZELRkpJ37I5B88JR6I/hIGL61sPCgTMIZ/ppT8ZyIl1JL9cB0jNvnlUgGsl9j7QMkZgH/vWiom4BdOLShmQOILKgXtaPEZsiX00VFqxWgdnb9njuH/PWdg+iVpPBgTd0BC8kOJl6Sl7JdxShh9YPDCrtUvsHZ0PofOzZIxuj+HlPzz1uobL+WEOfygnnyfY3QYN/fE9OcGoDf//CkZck7+WZPxwJywZNysA3NMgIdF7xOlFiTzw2yOCVjCaC3eWcl+kfp3hi/W2f3C4IWly8kPKb/QMRPXyOSh2WCW0FrnIZAwJ1BKweD788tlIhcvBHxzJK4N/3vqVrkO/UekhHgmOpj8vvK1ngDkbPH3WeZmywWl+JoxLkqu/rg86lZ5Ceu5+ffz9T8weOHADImHuPHAmHNwvRswwUvaCxXMiVIqRm9I2cB7jf+bvj6TNrxfbnNMJEvYLjfAHCw4JswStssFrXnR0nvztcmLdvQBmL9HXTgAkLjPhp+jvAwTkvk6QUoYvfHdJ0x4Uaq9YbGf/bKom1/cmZ9ZaxtzwHLy84HvA5ZwuVyA4UXb7P4A5lRxCr/EOy+z4kVRb41nGgvU4p/nsm3x3cwSzxB44cSdCiA+j4qUWirGZAFr5vsDE2O0dZbUCgwvdLJlWPZiINfivw5fn7DEg9ove19J3FeYGJjIZnw3vo9SAi/1Ev+ckHhRTWCm2Hf6/AkJaSYWUX7u5uyFtGWLc7C1jsI9od+nws9KQc6GfT+QuQdSUt1oCSln5Jr54f0AAw8wM76MPmD80BNgR+APKReDWcGcvhCv16s/gGTYNq/kjtb8xQ6/bQcGRvKKLuWM2fyF9+EXjW6OOScGK3gzw8BALr64rGRYTmjdN8jozatAVltWDKUW9DF4cfhBVqtfJKP1qMg6D8HMA8MP9OYvrOY4RMf06maOgdY76la885gTtW4wMxzHDlOn0HpU8WOsymiwyE5mGPy6c3akBOz7FQc7hpR8gcw5vUrmBTX65CbM6H0gTf/s/uz84AN8cyD7gRIViGVWNEBrA6P5hjT4AT3niIMgZUMyX3zJgHYcePbsmR9qPMRzqXGptOZVE+AVIWb3w5Gns5lXazlbVKKtdy8QLKGUDCuGZAbLBjNDKRVm/s+jNe/UWo8CI+eMWqu/u97jWSfzzdCOxueha9w/i6oxy+YHPiva3vySLrlg9onWelzAqsb8MM4Y3av1Wgr24+rrxfyS1UEOXlrqgMGj0YueVaHmkpFTRs4JM/l6zWbAmF5lJtyrIrMZem+YyfeMpQTM4c+3lLV3p18mYMV97P489B1K8Uveuyw/hP0dDsD8u3qBYCiV3Rb3zRgTfQwvYDBhxdenmcUhmcwcddg27vcUnYnexd52zDQw/SvAzLxYmx2lZBxHYxdl8Vx8Tfh+UgHYW8dxHOwu2NW05sXk0TCnd3SOcGQkFjFWih/o3JPbxddcrt6VWi3o3BdtP/zvLObvzoq/j+KHrq/n6e9l+EWqg3nwLJosCBO8u/azYa1RLwAnCv9Os4x2NIzR0Y6DF7z5wb9VlJxRWAyOOZBLWQVVNszu79CLreEICnxt51K4PllkmSMMffB7dD8Pc8nrOQLIDx8+fAJCMX6YZb4koPfmB+VcXciYgPFwnxhxMKu0693bXF+oxhfF9hS+iDshBCA5/JGA3tTaevU1hi+uyhbSF1qKVj+p35urikjmB20plZWJV+1mFofh4HcZY7A6SPH7vbVs0aX4LZuiBU3gIQD/e8x4e8/pv27qLlgNwA8z/Xwk8HKaSPCKEpYIITiUkpKxy/MLfXT/XLlkHsYzqiW/6r06VSveDq9QjRWgZX+eiRt6jskuD5jTIcXzJp5RZfpCV2cwOuGuZNFFePeT0brDZtG5+Mr0Z8EOodQcdVbOxSG14dVRqYWX5CQsiICNWvPDYw7fiF6XeaXfWmP1OKPDAOHF/dhx2TZWb/4zWuuExfz7ekdUWKl1r6zYuRjh0QSHWPy5Ce7QejBubn++6j4n/64ZlaM/3MKflYtFNW6W4L2In9VzLMgoWebPndw7pno8qkzLPCRMh6ofmmPO+Bl6Ln4ACOplR8Hux3SRJkPZMnpXJzuREyESwnbe/Vt0zSn5Hkq8TMz8gkl8j3123ZercxurOzVz2Fmfq2wVljI6D6mEhDS8Pof5uaILSBdpSgj4suYc63xobyLBWF3PMRZiwf0zfXs65KlLLwF9NP+zyffO6MPPyKQ97t/JYUU/D8H1Y5YxRyfUlwKtwUQc1mBHnYsXpCrC/GypUbToNycWD2eoG8nhY8sZncWaEUFJ8LXkjQEIGY/oPmb3y8ZUjPA8L7UQjja0RoRmAnXLmPyzx9GQHz588KRujms7lIJo/3LO2GrhDGPEA2/D8eckeCtb/FDNBWopvqBL8Wq7dX84c3qFmW0txtOG0D9bzl51jY45gFoKcQbHISYmtpsLoSVWV9OhjOM4ULiJwEO2FMdMc+ZBTGhLB44Wr38o4LJdotJ0fNZ4+3vlcLl4V6V/760BKWGrlW298TDIXtGDkBUP9N67V9U8YUpm63h0PsPJFtIhj1qqbzrCdoB3MqVkr9hm4vNO/HkD7Wj8Z3hFxv3SeanlYnzvghhTLGqzHJVeKSUgsJQSWju8Si7FcfFsa8bCQ9ihJq/sj95i3iE8eya/nBKx3cEOrghuOHzt9NEX7Je8DS/F1+Q6SP0QQVrFSa3F/31MFPPLXQf6Vismq+ZSimPGxaJgULWqi83xaV1CR3QwbQjj5nrcHLN2KKSgC5I07+h66wHL6Jn3PnAcBw/XjDm9aDNzHL6WzGIlcdZQoIWgOYD/Poc0jDBOMj9Q/UDy26JuFZjAcTS/AFMKWPiybbhydgj4vCNb5vvxz1u3LSBhMy/+cikx1zmOxp/hz7BYwiAkKUhzTqC1hsF1DcKpgs69K3UYvHF2mLOxcEwBGxV2CzMuCN+jmnUl/r5c/OxhC4+cy6mLn9j3nQe3IEnCtmDX1DpRF/+vlv2ySMNr5mLF3y/3GwjRq/PxvWpIxudTNxaR7BItLdhwIgrOUje0fmBOfz5azwkDKQmGb4SyOSPi5ZQsI2fD9e7KcULCEKROnCGb+SiBhW0j1AVwn8ERDENGrr6mBPOqyE0pIT9+/NwT/w8jsEodMOCiiBKGL8tSfFwO333m4POAsdpqYrqzd6TsG9G7iMGZy/SXUCorF29D23CMM6WEWgpKzbi7uy78OfsAVPihhqca+qUYcM2AQXxg7I/vuO5+0M3JQy3dG4zBUgwcB1+Kqq7JFq8fzSs4Pg/d8oOtaT96XIRjsnJJEyll9OMIrNz/3HTogGSAxNYxp4x937nREiZ6VNxDFUUuqyNhldz78A2VfRMB8Ium8dc5M5nCxJMv6DG9ekr8zMa/X4e4X6isxEsGLHFw6ISAXDKMP18LrbWDm0mdzUStJYgFmtNo0bajxTsBgNm9KIjDjgc6m02SFga2rSzMnM/0aD4YHlxrlYeUX6ZsDdi9TBYnXk37+hKpQPtqcjgZkIclpFTudUEO8TipQoN2M0PbD1wuN77hiRXmtIgS2Syqc82yLPv7z4Txjt5Pw1MnAKiq7IRecs4cwJbT8/NOs3eHIrbL5oiDZjmEiIUY5Fw4n/Hf4wf1iEr8slUAM4qH/XqgT3824FmS87q0i1kUdKYDMyBl7QsLZEAkAxVinUOXxj2ccsI4WJDkjGPf2el5p2w5+wOeEzknfrfCGYP/8xidM3lfy9pTiRWsGoxSSnTRmTCnOkd1LYPzlczuYJCcoQ6rj8bLNvszAzAGCFmy0E0pRgG9daIn6tgN+3F4gc6fOaf//Z0dphUnpejidPTI+CwQ8Oix7+yOwXVSgTFi3WP62AIGpJnQ+hFkiTmBUnlucs06hMVPNYnPW/Zh6uyDra2/SD8oMyuffm94uW0+kBWbJueMZBk32xaLsOSCqi6Fra7/rPVSvNXPi6lDHNCxc2/zBgZKdrwb03FPXWrt8DmLf+H1kC0btrIFtNb5sKGXWUoMzoSFOwjlB01rPR6ewZBLQb1cAt5qbcQ/GzsyQR7tOGDFMIjrlpIB/ky/VBEXmgbRk5fjzYNbtHbEUNtyRtMga6aYI0xMbgSvumbvXNy+aVIyr7LVnrPLiDY/+UVsJS/2iIaMZBjp/U8AtQiPHQEd6nmKLGG6CJLjrn0MlJrR9hYXhC7qNLW2uMn4ObeLM6rUCfGxEgbkc4Mv/sHiZTbHbfV5knkh0tvBr7WGpk4+sPi7fK6iA81BNx/EF/Qx78FZmBNWjMxFIzQCHrQpYGFnynDzW8JAh8ECltCcYvQe8F1j1ezPKQXhwOGPvKArFjhROOkgJbnEuyigsxuOQ5LzF99bg6SZQkhJDKlJOBpBJvDDmtWpJWfkkbxiZg7fGedX5hd/PxFh9LNVVIjN2NtBksaCqI3zNcyEks071uRsOYd0S8xaVJGr43NIq3EO4x2EM7IG14lfDAMjLuub7ULY1teB5kCz359hCY6OCzGYhCngoZwz+s5OTfNIvv/edZFarAljBy0yRTKdpZ0zU4fhihWiONVZelvluplBwNE6QFqDcn/vmbOXwXXsEJ3PgkqQjiynE1wr+C8FErXmuQP58ePHT/QFNXgcsxPS8HbIqbIpGFmjd0I3g1W8Hyqj+eEVtL1kjimycxh9oJESm7PTxyYPOFUcIDXOor324U9UbvBZx+gdffrCTGJXjYF6uQSl04vKGS/dqcYWcxGDvzANm515cMILOWTUQaYKSzRSr5qNl+/gvKaRMZPjYE4cgpZSCOWZd1kc7vGsJvsnx2KfguAIH/XWnZHEas8Hn7wMUor5ks9EcjCOVKmN3rlIfCgesB//e2vNK7jkFeOcjuXqYvVBpV9stZbAWo049hmbx/RKOiWfE3TOE6Zwf85QEmEpXYDBqtMGGxODEOZ5s2k+IybfsTf+bK8zSy3snibfFWKWps/pz3nNFfogfBDDRydFCHZRFSakU+woXydknxFGchYih7wp7uMTaxABC6v6TslwsHsoW/XPSyihsGDovTmrje95pgWHlVKwX/dgAdat4tgP72zPsxnuBeNzmpPvc+DEXPLP5cSLNVfEiRUUODpmsOEsZ9TqXd9+HJxxbQ65lozR2HHw/uzDq22fV/mF5zCrX2ogeUbV9ZptDO9EOVPUc8yWnbrKbtLXerqHSoi0kDjDUfGk2YPmXj7j8T2gGYUousO5BThIv82mYss7CQzEXDTBopF1iYMXeiA7aox2IvHw4DaLIjLzszoZ5OD5N2Mu63OYFBBfsrQu4WKnYT2hfhgJUYmUdAuiTmJRkThL09mpbl4wvFCb/OjRoyc6aP0QLOhHg6UMpEU1VBvtQySLG9QXUSHm2wOSAVuk89Bmu1z8duUtCVZ9c4iSh8A2Bx/wJHQ0OKh1CIJt6kiBcWoR6JBE3MQWVLbMCkeHLEdafnH2BpiqcccH1c2o4iy5sttoccn5QlyDKr9NUsAaRspvOmkyVpXgi9SyKmILeuvonb8+nc/Ni3MiRbWreVNrPYa1tW5kT/jlIwaY/+vA6H7RjdHjYJ1jBl2vbiUYIqL7avHqQPeqKwczCScIzTcH4uItlqkVmVGNjTGjgp28mMbU1YRT9X7SLehAYGGh6ug4rvH81Vk4Pur/P/Mw16BwDaG5HjQnAFBriaFw7+pEOPPJi/KrtV5y9sqeA3dBWLqMdJkOPjNVtb21ex2JtCzk7gaVMlnCDLzf4RgVFYO0ZsBQS+bBHrIDhxd7j4v5TJvWZdiiGxHbMkU1q/1pLNBmgs9/WN1OOCwrFCCfOg7XkAiSF810xlrSTOlMPklkJAbrTLCwwNsxyQTtuGw1iCwahhfBUITVt0t1qIuzhil6OJEFUekTJQGuA/OiQRIAEAbWz0iwWL8gFLZtlcVTwmiDFwMLHc38SIePtU8dCEmQ7H4L5x/Z/xtRA5/XDbTRUcuGmbwgKqU6kWl0XzOcEYv5hSHy0Ah9iGZDfu5YoA+ZqEbsO0wOzf1SUwMgspUulZQM+YXnn3tyNMfo6nbx2QUrlpzIxS+JwhpiocnuwUCixsLUxXSn8OayKvtscVD4yxwcci8mhwZVcyxcVDe6WFmtDdRaHFoyXibJ29RDkEou/jJ9OaE3p44luLBMVeuqvLovDkJttVZ/qTVjv9tZYeSoYlsf2CgsvN5d12CbOLqZY8leKZy6mQmy2lIM2CXQytkokLSoOPX8xmiopQbNGML0OfBsraNulTON0N/F0N144Q9WiJof5ZKdtcNDXlWwdAuN8xMJSZ0KmzjkzIuhxu7Ssq+TwXlawgyGjy/4FALEOXrAXPt19w0XrQthC85KkgRSZMlJZzFGd3okxWwzOdPpaI1agnxvLiX4VG1Eby1YflE9Fh9eT3Z5Xs2PoEKOPnngI+CdxsPc38GIA1oEC128wcDh+tDcRfOis/4g5UXdbs2hCh2EYGej6W7Xd+zdPyNhmGyFsIZF5dg54xADp9Z6KqT8xijlzJgkxJdmFA+9D2poUhSVnQVNj+++5kXx/AgD8QYKWFYcdhFc2rGen19EiTAa4jATLJ64fjtnQr111EvFbOOEmYn5lAgXLjaTukWHIVMMi2vdgjyAlELjJZKFzgj9r7WGXI203+oMMq6p1ttisXbCsVu5xzZN0y/LQjKLpBIO33If5nVuiO2nrrgRsq4bi2jM0HHkbDj6YBfp0gSz0zsgumGE6AdF3LrYILmBt27BZs1myLcPHjxxCMNCqOaUwkldAoK94OyAipwSdkI4YmK5clk/OGFjhSURobHiWDz4Ea2n4C0dgsahvF6wszAaMA1V+o9QkRpfjhHnz2toPBJbZbao029WYeilGvrRg+o54SyXKoYHKZdJpIAx4oGLHml5YfA6wHt3EWYIwlhlVkIjrvVIwdAotXBAmHFze+uHwZjYj91/L/JJRQyU7LCR04YnFyaZVYLZgJOgbtGjjS1sFkvOSvDFQ2hoKaiO/rAt1Lqj+0FslnAce1w+iSo9y4lq5oQ2BjaK+CY7nXGCUcBiw2zBpzipm89iRLAad9VyWTMUltzi2+fsegzylClGFF5PokgfsLIOaHBQKQaX6NoQdEWYy06q5cFKTZDXHBOzT6e5E0Iq7MbOlXgiDRXDqZd1yxyEzhPF2g/PzgvBQrOTeMHPEEkmzbSAEGVKBe1X8lysu1xWMXOClhO7YEGuCHqpCgVQK5ApjvVnWS/bScntF4nZXO9HOhFBh8l1HWfI1UkefjC20dghSKTslFmhDJl7OASAcTHrcinUzADIKkgoDhb0RkaXsxqdNDKTdzdzkPTCrlQbyeDnz1QXa2R8iWSDiT4d6nM9pSr5FHsoq/slfJuw5lUqsoUGFNK8RXQouUSHkqjEtxMzDbh/jlpeZ48KIV8j/lyloZqJsProyLUC3Z/ltlWXIPCCC0TCzOfAp5lefvDg4ZNEDD/XjP16xKF/HAdhJR9kugCvhh2ASZhUnbII4nKVLxqYtGDwNnRo4bAFci6/xGO+CXs7orrLpOQ1HjqhT+jEKs1iwDwJhzkrCDHw8Rfgthp1u/gQbt8JMflmLqU6s6r1YCfo0B0Uq9WtEJ4CylZcQFgKtpsLrnd3Ti8MVlaOC7f3htEncvV/3/crSQleLXv10wIWUsstUsG2bX7wc7Dn85UCGHD3/rNgjPR+hNDrOA7HljvbdOKcWnBamKDwcd+vQSEWRCNVrbGw0Hwg8zJpZGRpSOzFRYniQBeB5h5SZIuH7/CNL2bfVy6cE1wgKrkTHHKoqHOpFLF5RaVnKa2JhtuDVii1VuLnKWDSETqAHN2NqJRBRiA0JdglWEOkdmvwrq4oc6akQz3U/nzHc8w4NDVPFNynag8n7DpNp6mnKLwsOueSawzsLcR0FhU+6zjOSPxS2i5OMfbugcK0XIA0iYcvjD90XTx8fU2sZyW3COn2RbAIlX1KOK6Hr12+bwzum6PFTFU2MSlxHWi9DZCwYQtCLCoqJy7bRrgQQcDZto2iYEdJvNIfMRMohIJA8TEAQuEWcLdZInV8aaMEXfY2kClmba3TkoVQD/UhGmIPUfajwPN3Lm1RzhQ/c+Yxuacl4B19kCBAjYshXB9KdpGt08W5V3JeHW5KyOw4QYGsP1OKSSVA5Ty4txYaHVHSR18Ueenncs6AAfvdNeYoSEB+/rnnniRaILSjx7Q/08bED0GqfU3K17GEZlRBq3IVhtrojzP7QN0u9yCPLjwUwJg9rC6kJeG4yy8mbkJ5dOklibuc+YBHDLhoCcKXYGxDJ/UjXgiorE3u1zRduy5V68SMwyAOQtIIkSZGG7i5uXFV8nFgY5uXWTmKkmkpOeeeliOhgeGB0npzrDIlYJJifOxRtWfjgCuT/95GsNbcE6lGdVsEmZyqdr90jWppwkr0wClkL+nAGFxwW938oMtGPJdeXsMHhbqknMI7qID26oZyA+LajhmDCnQNWVOWRYrU+QiYRcPfSaaMr5cRinfpXxLgzgXCt2nnoEqJt1RcEFAFzU5Ccwd5iak4chbRjEtZUJAU3p3w5uRANnHYKzFXJ8QnlmEwl45Gl4cZ85g5RnQzot9KBW2cWbXjWEIyQSnTBaSXy4beGqFi4xxswcCgRYgG0hJYyWFCGia/n4bj7jwHRGBJybs+dQrtJLKdrhTl9+Y7yFwTrFABh3tVbLajhUZhsMsTqcCH/gX71RltIyi+ZB6yYHGGlXRby5poTu8AwU5PjD9B8vN8PpDBpzmFRM2Ds8eERBRiktZuwc5URZ5LwdGXfcmCgb0wwfBOLJNS3NqxLgZ5+6mw4KyizwkpOkJfZYuoYZy/inYcHQeFm0FkCVYYBdpT1keJwkbqrjg68OKBZ17CkiXwfBZ7zJ/5jIE9EpBvbm6eSPMhaCbnHKio/l1D7JIz5fqIKrtxxoGTUUQpmR+eVZI6j7loi8mSC2vMbze18E7HE0PJf6YbIOZQDgclcixbCHAYvt1sDlHl7GIvjarSuljO6vbCC4BKmRjciY474UNp4wzkcrk4vEHfH5EE6lbjItOw3wVAXgXd3NxijBHWAfI30lAv8/CR2aMqJdEEjfYCg3qaJIk5rR/qpXLoB3YzM5hSPociRdeWT9IcfbEw5BPF35+mM6jEEVdnIm8fJKCc/qwWaaL+wLIfPJl0RNiC+kLgxMNNw8I+RrDHwnSEkJosaTCdmTLmxGW7BIVNw289M6nDS3TEWBdVuu8s4K9fkF8PaMQsh4rfNRK2ZhaWSINMYT0yJ5CTlMopBt8OUcx7GggrTliRbqDUujqSuHy6vwPzzvlyuQTbTOJJMwsIVJBEModzexv0PJPiemgkEB3UHKBC3YfJKhJkaCr/OrFwlg0XmXIQspCXSaa0YSwEJtFQFQK5yJKGl6j2P/UVpRaf51BOcLlceFn7B4qhP5lT6zJe8w0dpEY/sYnl4yfRcYKht8PPtKDnjpi5CRUYhAZV9Y8xUKl7AuFv2S1lE7TldC1pYkCPstbcQNEvXXBmASea5Bw2Kxaq88VOHBiLTQbg5uZCvQjCBy86LhZzsktS9+ydu4VsICwDEnVDmxfumMsTLxefOzs7zDvQnAvyCy+++GSeKLkgtU/V0OAB5KpOx866qKitYY4TP5/0t8RDTth50D0JARm9i1T9T+L4qz32B7epWiK2HZL+NKPtFCVXbK5t22ixUUPZWeih5fQ4VkLJD5hCaKyNjlovAO0I5LoaWgV5FiWEoldVvpspZrcG0e7kAKoP4Qn+Z4/mjsZJVThdPOvGF65FnhzGGKNzKIx7egjRA+Wkq3fU2oH4oFPeOj4E7cPV75ld41YLwocNi421U8HuEOTBdTDRjt0vUjHBOODVIgVmQAryP0r0ZppIbvshFt9pc4h9A/69Z42HXHmTZg9NqvF0Gioj4NStbvdYW5o9WLFg/QQs1cc9fr4G05l6plwz2rWtDZ29CxQ+nuAmoJ0XerYcF2zAXmMsmCn7vnL8mLYf+QzfdvcAC41GDzqt0apisgCZWlS8gAUtOwW3hIJZa9jZNFLWj6jyc8peqZJmnlJy+m2aUVwJurOU451lWzRfQYiCNpPmLXP5SmmWqYt49OEeU2LQEVIqVOH75bmH9kiebpn7MNeCtjdam/hZImGze+BxbqJimHYf4FzQmUcNrTe6N7OoSfkEu43TAJwDZBbV4Fy1HUes/wvhZhEJxHQS6aIUC7FwP9aoILPbkddfCBm5hkRCShP39HLzJPYUXy0RYpMdy+hdyCgsOSQeZCMzduCucUszIRdD2/2ZCOmRJENreYwZ7Mr88NGjJ51dRnC/g020KLbRbnIwmk7sIImnIKPFaGERYja/8Svm9KpFTBVnXRzLr4cMIImRJHwScyvBbeUrWVBh2iY67his2jIaaWu9zziYXCEvef9YfGyZ750uJXHELedQoBb6U61bvSPljI1eQJPakrDvnn7U6LByh9cckIh+XXQ+tzkZcWFGhRWXMZaamur+Ks0G51GYCb0dYdkyOSBci83xaDE+GuEbFzrmOAiEBefgyOdY8Pu+03cpB0nAUl7Mj2Lh1ixBn8mQUZ5JHMo7ZLgYeC5mJYRAC5PB9ZJIeVYXoW7Psnk1T3ZXCutv1/D0vXFYSaNHHng4CbAsL/8ldW8pu+VLiS7cnG3GSAHRNJ33n06FFw8yQiMSMIpVYyktZZVEl9zQnd2+6K/lZNMfXmtzROEgmCuzQ0zJwjZ9DPofYYZI140q88nNdsE76eQYkE/7a07NGB02a4e7S2gInk4mqzjpxTIPRbkVx7lycpF10WQKcZs6XDeorCEsFZGkHT18owwpKuLFZuO7pY8ckg+iHcrjcUuHg4kezD7LS0hcSBhIWA4ZFt+MHnwn3chgMZXDjJV7baY4/86MMMxJxGKdX9LVHLszwnQWqVgeWJf1ViuWCiRFcSgihEOSc3VkZ4o2Cxqx7y71wmcwNfo6FXWLsDEI4eWceREmH6Lnk6157y08iSRGUovlXQQHh8Koxf/nMPiybV5N0cak66YjXq02+NgPYtGDrXe+9wFn0GJTqDGFudZalxU4XX3nkGsooovR//xAPLDR8yux6iy18CBY6u95OvA0OJTK/HJx763WGwqtjQEOukU/tiXO0UxHsJhszqfyQVihZ2KhslsRJTKgFwqfDgqzzoe2cP1GAkC616aPZXZX86lyosV0nzFvGH0pb0NgJOt4XmCaidnZrlyeVFyUnYdRoeo3JYS/1Widfmkr78Wy+XvRs4cfpEUaI0IaS03rG9wFdWnBMZMCQmkmSiYt+qQFWujFstQgtCSxpxxVBwe/yoBxSqz/vuO6B6nC/Y9mQEzOBErscqRJuMRlJqsON8ZMZCxuhIpTmCe6XX2JzlN7QTBpskxIr8ch1I4jhutIM9aQH+aTGh+E07IG3iDTR1RvHVwqcJbJnx8wRop7wN2kQcVw3hxKS3nBQjIsFH002F9iNbJwVMyAZgoyDNRcBrCThQ8ZRpIQMGdocpaQ8pqL+dpnx0W5qXceTiHeNqf+e4FnQYXdtkvQbdvJeHDyTCqlhI6n1oqjHy785UUs+DUYfEAowg2u45C4uh/LfiSF+HQ5LODkIq1ixgv3xXxzD7Ic+iqjRX8oYE8uEb0JDuvCrtmZe0eS6RkXGiso6+gkRHz08NGTZWMC+uCQ3mg+1I1FFXYPjZvehTnBSrEldJJWRFS7ZBmmACLS6XJgqBNH605PjYAa7yAuN5tz0SfcPLH1OMzCPqDPJS0+HaKZFWUpFj7/FjYhrkUwM1Q5tfLw16UzGGJkthb8cuddGhXh/st2PYqUwCb90qP//i8O9xHWSMuMTO8gwQNaODKHWyympYK/3Fzi2RUO71yL4g66Myw7LA7fwrAd4cdqu4XlJgo7e2s05nOVcRNzg3qFdOouMEFihV8Wym+Q9iCd1P+6REZ33dAgpKjQIKPZowv8RLXlrEnQlrn1eckOF3aKRiUuU9E3KL7Lgjs579BFrOdsdB8t2b2+LIIPlmZAkMXq1FMES615h1t8+xByBjvIzS1dVHqmEYMdth8EtLCvOdiF2UrQ2zU7ingFQWVJsKkPeQsPlcF90gKaTEssSHhK/ZCozB62xsNMeqY5wkcrBJfs1JS5MrQwzGLN8o7wC486ocT96Z1jDyaRKKl+2S+Fu84U47wUp3NGM7TWx6LPp7m0T6q4+RyMLs/6s1LVg1oh6ZB0zgxmH4lCXpUZoj9nOCnCbQmJ4Qwsf1Z27/M4mWisPzNmsDuN70cZSo6cjCCkyG3Z2YZe6ACJpodrBujCz8XQU7FgadF7NddZM5TJYs3Zq24g6YU25vSAMzo47PuB/PDBgyeW3J9f1DANn+R8m0Oh679WaSUciWXFD4rBNjWR7SHh4YhBkp1shWmdoTacVdSqNGXZPJCDZYXwKTINCQ/SdtvA5fYSrKTMdLXRx8lry8KAznh4F84egurJlnGQyms1Y9KxdgIeZEO6qYwahQWrxY8KUIaEzBNQKIzR5FDQUQxoi1hbyx04fMMyMXC5ZPLyKJZD0+F/rp2ghBTkht6FxQNdbBUKkIQRC4sV3VeW04KZJE7SQDaG1zBqbLLTmLOYW2RN5RWEA1bFqk4lGi+5RrfhmouTkSHcs0c1lNhRVkoQG6asNPR8qXaWQ63JDvsETXY6oBrdljtnRnLqlYeYZioKmUKa/kzC8iGf7E9Wh+i6lrOQLYdBqOjJoKuvw5WOkzeK0MILK/4ehKJ/zIFaLMwsSyn8beym+wz2nqyDzkFXZjiZj7qVTSI0mNOpo1DFyuAhrTMxh2oEMw0YLP4d1FgAEuV5MmMuoqd6DsWIEDDD5VKDMdU4N1jGlX55HG3H5Eyqks483Zd+iTB12WAZjwpeEqS+XWjzwjVaSw13ZA3dkZzkYgGvjTDhVCfijrY835LFwSzLGbEjXevVVygU57oJCOGumFTGTBtBjAnmIj/R3fsMfZuTEQY7nRkGr7L8SfLtP0UWLAv+EnNv72wZmkbHYolwVSDnYtGd6bLOj59/7okGP5U8bQ21Z8j9LWwVhjQB2U7ilRmHo1SZ4lULk9Wvb1QXDzG3ePhetho+VbIy8IqbbRRZEloE+jsEPWWyT0YfVPHOexVENosQHcc0B7JVH2apKtnPA2FZWNeVN8IXduzuC9P7dB2KYJ1Q2SNU2cl0mPnlcwhm4GKV1iAYECWHtqCQCCCNSC45HEHnHNjqhn3fA1MVxBTWJIaI2USCG+HZulQmW/NgtvXlxDyJX7v1NwfQKa0B8CkBzd+bX4jz5AqrCm3jJb0SCEG3UeLdtJYBZlCv1Y3ITXgwhEmzAUUEnIdHOS/3UeN6nWLTpIRDBnmkZlrg3oSs2AGsmYy08eOkWDwxeGQ8SNgkGEeTUbo5rURKKrIP0o4ttAcW1v5ac7La9xyGQg+2GfYeiQN29r1RdMlSw1mBjbYiB/3ZXP8wKHQLqCcYmCk6cLEqNYuRCHhQubxd6rIYErxFGGu0Hhop/Rycu+d7hJkZxBzjZd918Vta2plkYcUjaxGdS03WJxC0uYTBYTsymcFiFs4ZpjRFthIqoC1zFqbihxYgYrYl2oeI7hshdSygWl/REso8CQdzdrztaJFJo8REMRZ9dqyiyzuezkG+lOFhj28+D1LWjZ0IEyqmJmti7fm2tzjLEou56/XK8cISaloSZdiJH32u9WV5Pbf84PbBE31hxX8Ouk96i6dqftIvqAauKthoBr62DOIG7SGiR0yLJSAOvbxrnMabAtfTEHeSwZNzCX+YSc2Gi/suXvUmC/rn4C08eHvq8DhT2qSH0BuWdECH01Y318QURqbOs17C+eoyabPoRDLunt2tsJjFdCSLLBMe8gMpJ+pskpu/1bpFZKl/R58NSOUrBpC8aEAoSMpit7Q/TtbT5QR7TtR6cUvtPgMjFk1Y9OlSKr20UmSkT+KvmRWnZhQ4EQOOdjBQKJ1yGka8e7kPSK9TaglsXMNbHS7SViR42+6QVQl2iltJF16MnP2EIG2xwVJK6JhLl8QLMVhdrJpV2c9wO5yclY3Ayo3BTqWu+Upi/K8IKErbDHNFWvirKxWUu9Vla645G5gK55RJCy8xKYlTPCcOQ+XVpLndWAP1FJkQFvMxowXLxLLPyaVEoNTyyOoBg3baGynuWK2gFcNoMzoDQSaKR9YBqD3XeiMEZ5H7Leds6WFEd/amey7zVc3JypoJKAIg3GpHj+RSMerC9JVdlM9QS9j6rxjb5YW3nIwtihxBNSLvaDY7lSPjCzpgbyODb5w6W1fZ9yBXZOaXl1wiEE5znHB0kN2PsoqCIWtBVhoKY5Pod0ogiiAXad5RtnovxCsXwoKk9gb9njOpcZIzhDOBHKwZsyt7lfzcc54HYnJsJaPlHi+YXlSygD6bJoY7qfziiwW2LPaJvljizEBiuFo3lOoiQePQVoyORVX0APdJAZhs3cOjnmKoJFYHD1QxZDKVphNzDf7nsjjOYTniG3RnVKQW5BpSH9guW7SCuiTUegpfzWZhBKnPU0pBH40kgxqRk4IEaq3IyXDH7AcfdGsh+uU7+4zqcwKk0Y2gP076QgXOKQtoW9nOkWJXOOsYi5FlOriierVgvmTL0Zm4e8ARVaaqa/kRhWaFQlTMFLMIHcKRmU1CkckCHenegSaDvtnHPYuTCFYi+6bQSXbB7ymy3LdaV6jZcHxZJINI62PxoL3bWiOdllGo5kPm0DWcmC9iDGbay0SlLsooFBjm1OvBGd6YYDc6AwZ2N959GfdRzKWCzGnfmamfpgj51XUJbkvrO3YpmHkpSpArdhF4uapQmFjJgQ73+VzgrM7HSpt2f7G+rOTHdLFl0+wjORMpqn7Q2HO4GNKtNWzZsrCas5l4cSRSnJeoUMaHJS7sHGwoxeoiLcqvapnMLo8NM4524LJthLjXdzv2Iy74OBtPfnGZiavAoM6GZxfXdGsN9VL5rhHEHMGjheFfmlVIZJ1pcQJDzPxwoiYLRcic6UoNr26q80zxTmGZz653Sxt3zmgsKf6Zjs563pqtzFMmOxQEN7Hvh89C0kS+vb19IhjAp/JYU/kBetQfZO74wiuRHubye6nVs6izATnQXXfOOKx797ziUpcVd+eX6KHEXTdsOAAralS+QMCJkbTUvStxhYLE3j0prvkAMdhI3LRjDOS0Iiy9pW4LL59LrTvHSgyTq+iEsynkZdSZX+JpjtvyQTr5/cuzKmC36WlmG5W4QX2O4fuMmUrYPBPOWIPvQb+tGvoBMYhUQCohco65nldCpDq6LmDlv5zdU+2E87fjoDV0PvlRzfBPGidzPy8WOlob9y0oBiJh0GcUK0wonQKuZHAH4CR0zDFk733SMymF/YU8jJRTou6jEyJQzCfITkswwqBSS6cQsKaUY/Yi/yKchLGNMKrTxvuCrMiUseID0jkR7LuzVU8uxnfgXZuU4gxwDTpxKZXdu5TzIypMD/JyqqwsOwQf6XLrtAwKejihliXwXQw9JTuqyAkjvYhP9sq1M/RIf/4cSzxoSiibc3nsKRpZp7ou8Rl0UYR3VqOGK+fsLCn+uuyOFACmPdYpACUxiRniR2SigJ2x1rPSMOfop3yVhAtJGir01uxukU2c/eYzEtGAlRsva/hEEo+qpDkllJ3hdeYIwwz2mmW3gZFmR8JjvftCkoSet9wnHAbNYcPkZ96C9tURpxM7tR3KFQLRHoTjcUQS8Kar1Ukcx7U5ccZcm5Sfe+75J8LUFladwlsls+q7bJflgQIsIzPRzSJMxTF1/bnMmzxS0+RiOt2iOUVmdlmW2REoY6T1MeGNfkvSo+h2dUw54+gNuVRUMnLGnEG3nGNSxenccvnvSzeQc4pNi7TgONk5ZA5s7dReLmVqwgnFc6YGiQdSf49xdsOk1sFWWyzvsZXVTGU9BimJyyxNlgNh6MhDPlM13kYP7vuy305rAJ7W8K7mysprRmtsJDjYyvL0nAc6wiZ2L3GpBecQTBzc4kByUekILyO1w/045Zm3I7oWpxfOcJVNYTSnjHdeWmmFLYFV2NAlM5fmglZ93FQV+3Vf7rPKAQ+L8xwKcA8CI8U3kXhg8vdaHlSOW+eY8ZVSlwsAZgjZBq0/HOJhxdjXANlOEIqyrHPJMcQctOWWDY2MLRHU5OF2QMT8FVW7bPCpsbJ8YmHpAE73lOOjkywhyDkBtVR3XgBOeTkl3hHorbUuthmUfQ3etc9ysqD5n6nnJa954CqUfL02HpBB8MGKUA7KfgLGMWKWI7qyXCtkClu2JQMIx138IvcGDdRjj3Km2NcBm6lSn7S4kbZCWSDQGE1DaNMMyGKWoQgAzXbkx3YcO+q2oe07ae15RQOIEMKzOBeZq2ZKB1xLkhKQsTonN6sFzU4Rs2OJei2vnJ1wD+alq+ymWvKKRUjJEwkHTeeO4whny2LeTRTLbnk9R6SPuf3CiCzeBFWxKVLCzslVKQH1skXFIJZRKZWKdX95ctgdc6LRg8UEKeQVb2lK05vwyEXLPJDUwqs78YCdMcZp0OiV2XH0CALygeEMaEAU08T+d5CXL+ilk+qmYeBBW498GmKKsZEIXzgOq3wAQ9nukwmmKvo5T0N72prUGlCGZxKc1PE6uLQ4EgKOkiArcbiuxZCL3SMJiKLsGOlcorCBoPgWJja6Q/DqAAshyR5D3BmdQD+OE26/ILREVZYx+7qWEimTva/qXc4Bofg/vLKaoy+DQVINNZgvPLADJtNVnCworGHVorkd0wDDNkX55XmFOIVTNAuecBEmXdLne6AR3SSDpyz7nulzh0Yiii4dCd8qafGKME1YGLgO+0jknMrZTjHzKsUCNnOLDgVdIRxndbDlQk0V8xZDlzMZqpWMa26uQ6swcvhERIicdn5PuQFPrGRToQ3RLWDRRvX/5YwwwQMN7qisWZtRkJmY9w7aKMlMMTzGqKVJ4XfWAp6UPUkP8sHKn5G5oWJuNbeToWGo3y0vq3pSgMeYtJdJSzMhdtvK/XZmFtGJbdtiFgiJm8diOur3SG6Qaz4lnCroig7MKcf8cF2sy2fNI5Pl6NGD5DPHEvtaWL+AvoZKMU3B8LS0HMBlOWRmyI+fe/wkkYWyXWqYAu7HjpvbG2+LqeL0DoHRqLTuFpfdMpg2ZuG4Ok6BTde7K0VBy7u/1ortpvKfSe+TapWzjcyYU+V9i47YwmzMorJ2wZYz2ge9n9SGqWL0w54zhlBYkl54UtJrYfQx6K67LOYTlo/+wTYwcagsbrsfxHyJ+zVavgHXzqiSicREzShmQrm479Bo04eTmtEEHLYghdG7MzGKGBje9ZVacOwNM417memCLdRiW/JLQZ1QqTWS3gRHlepMoEHB1Fa9jabLIP2AFvNqp827FMGtd2ykMst5tNMHScJVkAsvb6VtuzD0itoNWARvieWnzeYMphkaE3lRSXcimmsjQcHY/QUtG4p+XVbm2Ypb6VsivLGEovp5O1ltZ/2PF0LLoba1vuw++N3FOgudg4G2PqsK7rSXaMdO94YcBqaCuXI2smNWit2g/xyoQ3DqZg0nATOPpQ3CinlOtixtsg6NdriOwNzHatBpQA7bRlgsl+zsIzpIJ9m3y9aE3QQG0YDh3VgtJYSVNze3TmWWuR+h1zODUAfWHJ0CwAVPd6ZoOpxHKjpZpZk/x+1/ToXJZISzSBSR/T5P378v4g4ZZjip/lP0st4x5GKhp5lB4HHi0dGWJYrmwm6pvnzPFHdQt82RHcKMGgcsy36iH7wYj+seNGJLhq1691y3Gp57+jMqCgv9ruRg3caKg3DI0QO5juu+nKoJT8q3DhPIt7cPn8jYrh0dffYVbDIVZ3hwYIp70vrAauUhY+slhBPqiZl1udTlPqpMiDnDSRQA7u6uMelPyXD0I3QQujiUWUy6VUAskcXOLuB6veJyc3FG0vAMiGAadMW1LrX6uataMAG7BFZLTplrweQxsjSy2b0QGlm065AXrj8J8SVSGPMpbVGhTK37vEXphylmSVLKOxbZh9PsSim43l0Dv57EsDNZQ5EBzotFDCZFaDqO20MToGF442XrvmIzMlx6uNTOEzwxIjkxKJuyhTAvOEYfkXQVuRasrBIhLGmRlN2co3VnlC+piT5PW6mFWhdiWnVSLcViWQl+KSzK50k30Xt3n62AyjgPoeW1KtV2tOWIionRvNJXdQpFGE+xtwBvLDmo5QWuKFnXlIyAUyqH3RHLylQ8+XXViyume6c1OzPnW+vAGChbObn6LqPLLCiVjCLRmHMpcQhqOOzdfA7iAQDufx1kFsQa7fXIZM8GTfeHaNsdzI85AkrqYcBK/7TW3emht8gbV5SB3p+SN0GX2dgf2bC3I8xPU2QFOYQ3Zo+Zg3Rbrss4fNivSGayxqSbMD634/Bna0WU5/W5ZU8k5t44aaXOcLyYjb2vRNZ+HGGjn2j+OSnivd7d8QyyU+qpMXzO6djBKCuZbE56YiWhQj32mM4zFa1hEMm/O52ynSSG7O2gOSTNIy/byruXlfyjxw+f6Ob3wWy5hx1G5d6PEGl1iu7mKcimaUEXC8sKGTBGDoREYsGOmEGftJMFMU7tkvIvkrQRssgITI+MHIZLTUxmbhhx5r5yt3nAyStrEtIpubqwTbnLp0jPnD1h7LgeoUifGNi2C5IuvyknWGKt4xfZUAAesSk9TO8oW46sYQ17M2EimeZNORaXwsH6CGGZlRz+TzFcTo6tHo0VGfFYcHYgTLf37tkoisnt7jSsDkzad6n0tbCP/Ygsc1MqmWxosFr2FPGtOeYPKTLTV/ZySuYeYmnljNdaPa+dOS5RKZO9pc7BCJ/JHh687I02EbXWuIhcpU6VPGExW2lpsb62S40sBF0uJed4lyknp0SKlX7Kq1ZnurQDtKivS8ticdGVsE/XAT3GRL1s2K97GNZpsKRCRPCxjPzkcivoqNTKznBlbwcswxwbBP697MndGr4HE0cwVVjwK3dDl71os9zHow13JRasWStdHnLMP9to94Ld1J353GlBZbLgKbVGFk2WCBfU7PQZSITyYGrJJzv3kxMEzwrNGUaQfTr3ly1dBbulNGXXL6JFQZ/LvSEyc2aimSKLmdABdc6HSDLg5Wnhk7VIHvLWEqkDdLL22GRZmyS3PcHS2rhexffn4Cwa9G4zRjZE6FY/xRwQzq45R7exfMZyEDdgDCXjc1TxaSIcaY8//8LzT5wgtAwJMymiRtxYjCvR7eQRb2VdKJGdMJynf1x3b01bC0w0V79Zt8u2LDzMvAo/sbuEvV2vd4tzL+MyVrslu7JTAx155ffWHP44XQYTfvgJ39fBgrS49cd+RHVR6FKbc3ZqM11zhasjpZP1AsKNOCV4trBeTsTBum2BLLInfZsU9VlsJcUto0nvDIwq4fDOQooOSYyhHDxyVzOLrje6C9Fi3lELUnII0a0eRD1dhnhjLjuJTBrzHIMXndH7q8bFX/R8wtI9RYymEiPziXZphKL2q3INesBhOjASW39RoyNZ7SS8m6dcjzkWBbEPDkBpTSNVu+CbxFAqrXNFNCvFT9XZWmtkGBkiGleWIEomdNV2iiJKoj1dWg6XWzg3z0SvuVPss+y1dSj7IeBrVBnygzEHi9CBGHLHrOBE6hC7TsXUQdFpYVEQuphTNLHe3XJ0XlRc+Z2pIlU3khIcsyejzwkUFh5hKhoWqyeFPiRFzko5XUrr3YpKvyr3xfRTKJgL7dw6xsxizUTC6jzHwVq45dZtY1e9DAJFN6/bxo5rLrsfZYsP71bVCYsFZkXw5QwoTxk3Ky9JoXs1WJxbrZEg6fTxFakgkkDvnLFOdwhOhDWPfY+5sdOkF0NUVihOfmlkTZJ9Sup1H0vcq1TEQFbIdFSBdL3ukWcUkdiPHj58Moj5i+anNrDHLe3WHjEQJwYLCm3GiZ54Twk64DiumFS0Pk4gH74QB2R1s3HY1kf3fINSvCsIV4sUYUpjnjQAJ6hD4rezfYkS1XL2IabAwEw8WhtC+LtmG4qFzOxOAOBmuyHtcl8mf2ROgC9xI2FAl4HU+pbcskNwgmsqWlxAE64NuBDXnXNBd6BTKk6CKrGANMMAL15je6wwGym4B11+xf+WuMhVtif6GWGJjawrCcyEB89In8scMuawTRe2L/PMXCoppCdxpa0WP1g8Ka3EPjL0JGYUduy+Q9If9YAZdNAJWpJ1jjrrgNqQwlJdwTlaL1M0cqwDypSHziwH2fvo/YlFpU7XzFXaIp9oE08a+fXenSYtJ4Wx/JoqjetkfTJPMJ+CvaTUX4UMIjkwDAPlBs05Q0Sh2opADpo03VYz41ETjJ3GuG9hw0vLnYEtoCylQMpDqR+DueOGy82Fz6ivdxEx1inQDdeg9KXgZ6Eyx8Dl4omlk6FiK1/EqI1akDkYIpVIkpE5qzvcYkHfLJ7EzJyjh3Wzoil6b7BSwrx00vZe56ExGliU40mmHngIaxao7ujoB9196fIbs49JGvRyPfD3mTjDqSFoxGmdyaRRmUuJBb+ICkc7Ip1w8LxzaQaiU+5iKJrME2tIAZBOHnCw8N6rGg9gOY3nhzRTzExBE0aqwd8Ig7cV3qLDLipAbUacFr8UrXPFMw5imGcoZ06GogC/KEqyhNNqOal5s5UVgjJ92LqU2UcM7y6XCxeDhFszDjG3CCHTRkwyZRiMpfSevQPMEJ5wZ1BnMqToekbnUJOtPPysCSffMKjjMFcmZDgvGqwKW1CMHDil8FYlpwW0wrlKkAXczDCFv812ubjfz+y8fJdwTpdaEcMj7Ju9QMinjI9+MtNMylWYI2ZGwETbD6rfl15FuK38wYRRgzknK+uFOHCa4QAgqino5islu567G1TKIRfuTiucH2tDF1rXlJN9eQprDtpxwJM1ddA6EYEHRqKnE7+HAqq0jrPlsPXQfMw3qVEQNk4kjhy2GeD8gWiBX0ZYOeK9uXpaa0hZGjL5bJzLSUPkFjhSortw0GHfHjYeHlSmLta1MN6FpViOWqORz66ApKEU0tU19YirTkG2iNlG0L+557YaMLVOX4WUyTlW3W6ik/BUcunUs5oRH+FdoF/WNzfbPdr9mYlYqnQTq4OfMSuwWH9a3/txBONKs8fEPSCEJliSokPb6tSn1uf0OdrR2olYYMEaW87LWAasSGEiu1CStNJDZZeD+7NVMR+Xvf+ygs+5uLsFf0156hBLDMrIWRkhou86+SQHs84NQdMSZ4+J/PDRgydZ7os1M0h9hHAlJ3Hi3Q5dO1MfXPS1GbnNIMxx8aFq2JZbaDKkNC+1BrvJssMh1+uOMRMKuwHNPizROnksDxpRIWWWprlKrYXsHkQs6mQuQiK8Mfgwtq1SI0GvLy7CWioZVj4c68PnL6XWCFeR6lzBMaJN6saeJwzbsmPS+74vuIBmbcKpI4Mk+2XjrLh18MtVc4R9hrQf3rnIFVnCSnlnIcRWaXnt0JW29xGuoHUry700oMkeWgXQyE1/59BFQFqunJw1LNKQMRTSHBCO3uOZufA2hWGhmXdaoqLKyVkhYq6+RVhiyNpD3PWAGYh1yzDOTh1qwG1piQ9Dg8NOzkw7e1lcn32YBImJkTTHZBa13AlGwKM4VZqTRIjR10GgA1oeTGF7MwZyTfcS6ITtC683y8GSTLzcwvU1CzbzPVOyX3D63v1E2V46nRHxDG55nqMqLaWSHTRYZNCgsXfvEnjoKaAp0zFi34+w/NGA3AsrRsiWBQ/LTmieKx1DxO+Cl6ecCTT3E9yTzgM1Foahcu+ap6RYn5M6k5wzLpfqNi0kAszRQ2yszJg+hot1u0Okxgpec0unAyOs5H2O5uaZkx0ckrI/cnSZ4SBM+xUREXwPINipmJPsVr8Ai9ZeYSfLmGsvdHLAkuFjxV/LtYhEGcmLRaxFFS30osslszhOqBdHTSYvxPzw4cMnMuwTNuhxjYUYqu+6bdt8wE4OdApL5RT+SmGFnBJGP8jN91Z6niytJ00JJSDcia15wiEvsqANIwReydx+QHQ+TzedkcOsIXo/Oi43F1dMk5etFvfMDtk2t1xw1lEK5XkYRLIbmvzvxo4gqjG21YLRwm+JO11QgiAPDWY1Qwk2BjeEhuTCmdNcgzZLSwUsltlx7A7JEbcWDupVoIJsEs4zrvC/WNC3L9rizK7L5UKNDWM62VWJAWP0NTvbgHhCoMUzSdRE5GzAAPrkRpDVuXjxpPfOCZyVmNtWqTWacbh5Jb3szGeklYH6oxwHWYLDiP5cT6QIusCqu1ICoVh+rrJvDOoaUSUL3pSNjB9aOWYmYRURliRYQTzJufqh4zHEoeFrbrkEB41ZkFE89/vOz+3oIXYcfJ6lFlz3/SSsdBGuOsPeVIVaUHF1eVhWLnZfeefJh+HHvsf+62Ms9iSJGXNMFNqch5M1oaC7687YYObBzOUkIY8zh6dGGHTKwTghoRIe0yxM8JD/X4m8HXew9uGzoMR+9LjsBbPlwqx1wmSuOF977u7Z1QtZzYBtuRdv3BdirEKhTWSoJSI4kQd0HgXwucVMCCCF3qKzGnNEqqPPdAg1djlieP5I5WVZSKwRMiFHAjH4FF0ApniWKA7TOotCgb7gT83rgkwTzCGiCnPBq2NM5BdefOGJkt+MQ7W6FYbm1CVWGz1mDXOlyi4nC8MyUFS7F0OuHrMOZU/7wl5dRmKy2LaVwKIlTlye/QuGGlTOC3aR1cEYA54s66rzELidYkuP44gWF2HEl1ZeCN07NYSWfkKYojQMOCWE6fAII3VW7ZnumwqKUSV7HI3WLTOGpr1JO+Dq5pRW9TQJGSiwKc0FKRUeNqHeJ1NH1aZnIDsTBEE97CGsnGSTbTcboQ2/DOQIPPqAVT+4NCPrrMCWOaZ3U16lOvtLsGc5DfK3SEnsAWWsDQcfavYRzDzZXfhFO0PvobwSWc8PiunO9NLoROxE77SExrx2nPzTBvMWRGbAyYVBUaiCZtXuT6yKNptbYOdcsB+NtusWG22O7v5g/Jk6cI/jCCjlbAYodo9ytDUfUVIf4AeD49n0IONnH8Pt5mXep1nXtFUouivCUuT3U2iR9oocmh1iPA3Tc4rslzADSLxcLcWsQh2nW3jMUGsfRyOhgzbwqQRDsDJp7wzTJlKKOzVBpZChRfZRPmHzExOVkLsrzt2vqnJ21EhvT/QAq7XiuO5ht5LAjonvIkL1yIrU4QlRyVno6CQsZF8d+xEO4pHQCFFqC+boOFoPjzCcSA2rixuLNqsC+JQh1PoR84pO+3UkCwJNsNiINEhwi6A523J5YGepWaDmmKW4952o11ovupzyw4cPn8gfx1JCG2tTLCPEEbqOZR+BGKrqwyVCJNEl8LAKC/HpJmm5lgj2Gdzgoi620WOQKdt1Gb85HAGKqYrj1MyRdhzYlskiHXQ1P5kD0cVIhBbiHEWGhkZ1hRVpGD7pB+OmeJ5FIsOzsyJ2KZ8penQPbL4ch0YEQVVuhEH4pnI24/5eOToXXW5KPpMfF5L5wjnBEMfemPWQ3VIE58Aqf67Ss7jlWQpV+FYrD9/TQiW2PcZga5uWy+vJ8h+GyHDAWN1GCLXmShNEzJAQojMZIzpGnSJPu1TNPyyqeAlSSylipUINvw72MYGtbpFvnnCq/GOel2Nw23v3fBFb4V3rAkfka0yaC0YYFCtfZafknGkKmYK9OEe/Bw+OJiU9wkTUKNZL7CwCY2Z1e8bpx1DHOcL2G5M51TSeTJZO9OVTZcsuYEw/aH1QjtDQdFr7eELesiC5ufXY05hhWArmpdZv53MsOUfqp3vFLZZQZ4c36TWVS2Vanx92NTsNXSrv8MVqq+oWBDTHmi2oQw2CAenpfdBIJtyU6TogXQHNHCNnfCwFv7ocn+XOoKTLimn07mjKSZaQZPIIL3jl5qwLTEXfCGJKCjbfxrmVvO10Lwmeu+fwwaIlyQWYTEcFhwWUx2iEULVn6U2wXKkjdCvHnFtDf+UL5ZL5GVIUbCkB+eHDB08UShNKUeLe45RuFXh15FvT6C0C21Lk6Sp8Rr9XnHI9iHD8ncB22XB3ty8LFQrdDho4ioZXJBTky3KBVQpoxf2tEIEn6yFnZnOMYIb4Acgqk7qBTEdXBc933saaBWiRN7KbJGKT2EdVmmUPkYlhYlpZx2YLZspZeg3/bNtlY7wqIp8jOqPkSmVdetvlEmwVwQtg9RH+S8wlH13tcdz5tDjpy6r/5Egro0Yt6sYNjLlYUvIoSryIxnBGjKMjfjGEPXTJJE8kpqllwoop4LXQHWFFno4+QkMhuqQuB9l773RN9cvXGXszOrEUpAzN9MIUMSy+5rrEwrl53kt7FGU4GwVWkRZMaISfEzzMBuY97YTIJeswGHSULSf90whiiqAPdS+TVu9JcBthG3VFir5dLgMWkFdkqGieki2GrFLEp4BGFyMnisdMm33O8c4VqOgfZyaW5hzJEva9IdlcoXPM+lYX5op9Pxeu12v8ni6tCdlUoulPQkpNdku0nZG4VbMlzZM0a5U7cyAIJ7aliBYp51PeiwgPiQWu0+JrKTEHOM8UZI8DFtkzTWbopAjTcqeJHj5o0kGNtqBzZwlOaqDc3sTf80JGztY30V2SaCBGFVR0qMwnxVy6qEgK5ZGRUwp37jlcKyQbp5RtnXGcw9WTIDFZQn70+PETP/SdsTGiTQKjXltU6XKJVaWYTjTIiDQFUMwrH4et4Fm/NXMjTf9n4qQ7D+R8Ym25dUWOAW3dLmQS3MflluAvhaPljDnAWJbMHNIPYn+CjSbN5WSjuFg4PaoaUNzoegm+eOkS7JQjQPgoSyh3ysA+NWwuxCYVdp5S7VRJyaX3PB85M6+6DO1OOR6RVU7hWLYSMJkPHk/0TwqSxjiiE5A3VjplJg3GlqoqUWRsZGdo0qONW2Tnf842sJUjcMrdkEIcOMX/zoFcCxqpoOD8JKxUIF8kC2GeOtf1vmX3fta4rOS9Hp0brbjH8ETLU0qbNnhmtK1nvyQcvcc6mBh04h3UwXDPiFdnDrcVEieUG7HGPDOCzkQUmDjnyMzoyCMGN5hraTHqFObUT+Z9Mim1lTneglpqzDmnRoq5Ok6g2CImV9RkKaclwC2Zl87JpyvxUBeKIZhkuxQAFvMGCYhDq2DqIhxByPSKcg1GPZkGrsPX00cdklUMQ+/j5CaxaMiZ9kduxcLcFM76lAHvkdU9DuZwK1CHQ6sPOxnIigavsLY5pkPm7IHloRahdiQaCDZSfo3OV5mqwpxRKZ1JIcLhENPknK+SgNEX416dNRXwYig6+2wsNwD4ORwwKdeTmLQ6E1NaMggN9ztz4l0suooLT0pUONTw2/K+Jz3FgWd7dW7gTuaTt7Vw6OQkJpPgKeeElDOu1yPgpXY01wak5IKY5EK+XJjs1ztSISjmrU4wBcSmEWNGEnwtdDdEzCtpbS7TwfPgl2xSZ6iYGBQW3jQRhW1+kGUzeh8Z42WN9N1VXQpz7/Qs0sCvqtqk7iWd3GhDXEVXFlGJw0snG8CBvozQyukCnVGNVBrqWtAfCxdvrVt4d6Uszyq/qFx5PONA1MJRO52zMSITQUjIcpQVtZRMFcEz7h5aPbBqruAoUabdoh6rGpzu+eW5EsqwcNpuKGnTyQdI4kIgKI9+gXGjbT6YPPYWFeNygQaO6zXwYqU/5sJN1jvaaPSAs7C1rtTqyD47gpHO5IvpdOJEA76YI9CmvdYSUIHeR6NnXGSWTL8AxxzOdOFaDg3KmcptBouZxAhzRw2UxXgK+rTYgkUQDNlrOqg5iNWFJTdlnAqGiHWI+ZdimnMEermVT4/3oTXeW3ObHsE9/LnSWGjmcRwH9TIWhaocAurmIuSjtQV/n8wnvZNWxRzoaqw3zdTasRNuqgG1Ky5bn3kmp/XqUFZMs4pLU/aIaOOyMKGGxoWKp3NEFlH895wzpgrKrkTG49532GqFmV8G+36Nw7zUzQug3oE0TumNvmf3u6t3CCp8SYrIVJvHXqbLtkz1R3fLkpJLkIM0Y9WcC2liDko8HtzePhHe6vx0rMzwk2pR7Zfji6uyU5a4VwTrxSlq1qlrmWIpGckRDiJmLsqZMh7mTJFN4YfNiHD6Qb2EK4L3GNppoRmptq0dbuB2c+Hgu4fIzqt+Cz6/fEnNDGmsCEpjelsbTgnOZVlEr+yBtNhGor1mOym7Bxex2DE5IjjlUuv51puri0/zqEEdivDVobQ70o6HIB+yN0Zvka8xZ6cmYeWTK8mvmCGVHNRhdWlrQN5i2Ba+SnQ5Hawal8nbSqLG9AOxloo5u8992A30rrAbh3tmUCyx2EQnKM1gkdmi+F9EPOo54XEdPJOMmM4qb+PcYVGtFwNHlhW62CSMlAeZDr3jaOGSeoYAk2i5hAwHCRDRPZpXfMkWDRhJlhYpMsddLLvEcFMHu9h8cwWhqftWtLOxMJEY1oKp5O4O4aPEuWM6MWMVIaAqXrDnIHQ7p2tjlliYoV6wYCGqE4pAJF7IguY01zBzI8NEckJri26MOWnTMcKDz/VbI2Jv1e0Y50vyhCslO7TLf3Z1tQVMKnq8TBYVfCXvOfmFWTb0vQXtOPRnchGW5q13utWO2CsROMXiWrCpD9h7dGXKMjJehpIepJN1vVzKZ58nu6cFYzP6PZh6oBpdTMRMc1QpbNM80b/HDP9AuU0kS97VCTZNroXb9wMpqQADcq5cI7QAIqOtHQ35wYNHTzCn22WMHsOoUnKYn/kLnnETq/oUlqzkspIzjnYEC2OjeKi1jm3zxLVcSuCR8loCPDRlu1SHLU6mXrKQkACoH32FpIwZDpQKsnKMvkdMqtTbg9Xb6IvbLSV0VPalRKaEZiut9fCFOsNio3eUbUPbmzuDyuKcC01vP4X1wwyTPa1QBeiISaOrLDKJjyNgBh22knAfR4t2OlICjx4zGxc0pjicLTxyRihOpSvxGc8pFyFZuNJq84LDZmMegHI6ZLiXYrEjZhian0mPMkgEmIIbMVes6sltQHYU51zrCICyvIJ5MJa9i1kQIGqp3gGokiSrzlv4jE5fpsG/I6Xk1WLkoQjay0uoxYO7Vk9CDMdVXSiifIqqzoKs9R7xxeOUiyGvqjDbo1klxjq4ZuhhlpnePTt0pQLGRSIx4op17UcHMh2T5coqNwCA5p7LrqJumUwrY3RpC8cFVaVSaK+wtaVJmaPHwYgTOcThHi84bm4uJ+bRibBDOFCU1m2rsc7lgnF0Rz0KhX66FUU80Pm5Ck3RW93Zoh3HyWPMYm+oO5C/m86HRlJBzHPLmi/GJEiQn8LFNofbRHSwZGGuupL+6Ag9RxRUYjaFjUnyqG8iriTcjChGGi/IlBwyU9EuskiTPoVi8LIVhp5lxlxv/xXkwM+usXzLKEe43NwsEo9sYkr2QKmUl8RdOKWsPNreqGTNMUBTTsfRGvrRIi5TVseFB0A7DsIrhdDKoFgGVJfaUmFywV2vd2FLoq5H4TBnNoKGifp1VSjyydfQSNj0xFiYqpxeSyGv2Rfgfuz0q5LF+onXbYaUln7BnW9xjzEm/rZUzCWvTq1u1THDfKJJxoKXvUMKFex+vfpg/6TIdrLAFoeysFXNLCw5+8UPqzX8UliP0S11MGPEFb0a7i1FsQblOS/7EvlF4RR5KmaaaMqp2KqAksMNGmI3aXXCGHKGniVgJg1yif07XdrCFludRuXCdpJEDWvxQhHVJIU74jyzV6rq8FpryLkG1z2srtVVk3FwpjlLDCkbEl0y7ura4/IRW0p8PHe4XSmV0vjIel8d62gO2ZXNo0rdf+l0UQf8OcMQ8WwPM0eLv39SRBqaEOZcl1IJkY3IHrGcgxWpQkr0zz6Ge5cVWxVzKPTSPet17f2jLZ+wRFcI0WBNMDSLAhVcRzvu5aycYS8ZHYpEofgG3R5O5z2bIo44WCXuC3ueU0Z8zAhLit8j+ncUQEEKciG1C5THyQ2cELg6dbElew8/QQQT1e5ZzuS6XDIUJTCBmEsoJ933TY8wsxEszYQzBbHzop0p3TNXHeO+lqrkQncJf1dZGfCE5iwtll9KK453tB7VoZyoe5/Ijx8/ejIDvxUFzHHvRr421SWBsUmAJZ8ndRLK6E7JeIBbDKzUjpec3UtpjIXbJnNdSClLFc4Xh6EIzR6sE0EXg8EmifGjGjhXBjAt1ohXvo5LlggXkoFisvX9FBI1+qD9wjIZ0+A6qhyK+BLzsX1Y7yaHddtiaCamxsR0tTvpf8tsbUbL3FsPmxj9UA3Cel8q3z56JPkJSw7vreHuv3ayfFCFrAswkd1jaZnWOXzWo6xTteLYfAcmD7+5BvchkBb9cAxWTWc9gGjD4+TnAz/w2d1oFqfPFIZ+gut4SHQOgDGp/GZFpOHfWb38i/UM0nKI2njZCnoHvZL6SZeQQhA65ylrnIfuOFG1ldWNgG3PVhHJha6aJWp4Oc8RvedOzvdhrfmkCVh0IdG7k4gEp7ztRjO+cKHV78+LvTZPPmdSMJ+dFPyxWWgBvCpFzAj1Hj3iYASrUBV7thRwiGZzoYeZIB07RagaIOIDAoITE3Ji2SqJzSbCh1AEMDY5WY53F6FPfUTSYtggUXsjXYgP6ikqpetudg78svLh3sy0xOnzlFVESr+LFbfQcqzLJ8XsZ5FIZggPxRBTdK4cE8bJey7XHExPo++g0IuQSthyLtflGjG3jA7XhezkiwXTRrIsPNU10ek6WGbBOJW5bo4I35wz8sMHD56szsIXQ+YikU+NKrEh1bbsxscI9sZZONSIGyvKFGkZnB37Dox0D2cOkz7ZlacV3egDsuniLOBeJTIIFYWyOCWHlZjwtRYalsU3QHaSsbtwPHsES2bEIeDBND08YdrR3CpAIVTZ6Pw64tb34elZQcugLSXlsRowOc4mWZinsHCfUBqghWeUhqO9deKseWVNTxmc4+RHlCOgSxt9MtktqNWEn3Jeg3ApxrIw6piTrIADiTHzKY9dwskiHyQNmkVg4Kxs0H5fGiBlHBjnUzlgtaV6tpRgsgGX8DF8mXJsrJxPynL+7ODNnx1aOfuSh5TWXOHhLJM6WcF7mNSMaANh9Fp7SpqrW6HyfcTsw7KhMi5gKAuGNOrCZ75xTUkRfk7u89yOFQgmZtGCjlJYzWsorws3s6DwYLMZMQgpW2DfkGh2AsVORnkSSLLA0YDfIbuVU6FAsCCCdMRllzUwTutSEeNs368R9eAHrlFUmRjp63R6pITEAy5yMWQ1P2fMKvf96ntTSviTw0FlJw8yRjVHUsqjHuaIbJB0zxdO4Uk4hcnNk5dUJ6ko/AJLJgkgLft2uGRBwlcxnVakbjqFoC3HYxEjRut06iALNucoZdIpH16fLeeCY1+w+LG3OGMVBZGY4pnLoheD2qSDF1XnfCn2eTYmidIS5fHjx09OVAsefCM0FqJsKk9ZA8VaJMLyNlzZ3xs1DBEJKbsEcswHDVjk0SKoIeiuqsJKpSBnqYjD1puiQvm0uJgwR5QsgmLouo3E0J1JJbcW6NktV51Vb+OUy2FBydWLw7xfOZ4V7sJuA1tl15DP3zcRQjF1Yf0UsWqn4dwMLnbhhZNiBJLYyYEZ5acgLw4Zs2YVp4xtJR8OxhOfEjcD+tOMqh9t4dHU38jRVorkieQ29ViusEc7ovOUQZxTcc+xmzMsZdwSo0T8aCdLp9hyP1X29mjN52MynmOmi+ZIYyrMZ8b7SLQOn6TjTKbVSQ8BOgv55V9OdF8OW9P6TAF1QTqAwcKIynDtFV6WocmALWdW0uKd7p4DZh2jk53UwznBTiwnQa9h5IgUPksWiZZ++EpHpAmbDBOdAlvQ9oZU1oVmLHgSDztBGg77WlhsqKuwbLhed0YC+GBbzswp4qBxiklYF54HlBXCtv7/w7jRLGamddvQqPIW61DvJSHRpmhlz5+r8XmaXcrUcz/2sM8fJLvoZ507+N57rM/EsCcDosCS2aG6sYNFiJ29Afv6PoMO0hJiDUbGOpSdTjHX/j6nRgFY6/Vc4KcTTClYMKI4BCkOh38L46pFM4/gM5GQ+jK9jWK3+zro/UA+wfeTnnzBQuNMJz98+PBJRLSmlV3trKktWAQuaW9+ix6ycqBSHUvZaMlQLxvunt2F8lNDYUw/8FxC30I41o7GyhExqMqnnImcFr9b+RsevJNCg1LpKYM4POZJgcvgKEb2RkZDBP+kFSdL1pc6iD7aGrD3VTUUHrTzZEGRKLhSHvSUBw6XpHIk/DO6mFBWKaKBjtNmlS38Su7jfy9uOmnmVZozYyZ/1qQL7RbaGFl/OEvLorVXBLGcjcOPy/6rQ1tLCtiZy76dNioaxsvORZ2fLN7TCbaKjdpkub5s97e6Ub+RwicggoC44bVBVFF0QqFQ1kKS4jydgoKWA0Gp2/qs2fUaOdNWHEt30pnJIsaghtBS5PfWYx24XX0KLZrEaDhBAaJIhkJ4agbB3HBmfMiHSiFspD4hC6Kh/9RMiLmUiqAmL7s+lm4iG8ZoQS1NtvQOKbvITAaQcRBznuEzGlXt6iboVEtYquR1mIUWpljE/woCHF2D55UQGP53tAfRbMi4dvSMceKPqbuS/5n0VkFfHyNwfaEj8kcDHYjdE83C2QI4nQeNgV25nkLumFh56opXZzUCbk+nIKd1NuRgfYoYkqMAXpEKTuawOFeyKOti3s1THDWUblkWxVtnaSAUtij+CtIKcfeak4D5MUEeSctfSwhCZCkpiI3F30RCfv655550Bte48GuGZcSyL8W9kB1BRqqewtwrBkEe/6ickFoyDi6UELsE2wi0HKlcjD4zWYp3X2SdgfBGuwaF4kx2Ote7K7OOZcWV7vlbGRPLwnF1gEmE7p5pkUvsVeBk/oZneJSlIA+6ZGNy4qQ53Qg7bxmwCV6QH5RmBPOU0KehZYihxNEP3x3cG97WrZKRseYjspRQO6s8gvNC17xEDBX5J4XTLNtas4R2PVAvWzDGjEpdWUm08O/hfKq3cCLVcHuMQXPADEszslUKnQXiSHANFdvos+uBLtQepAqfMSxB12gj8GTBaZEXzbnMJKTQWZEunjrCYE9VundgtvJPpqFsOew0fJif4+IIDy7qO4TfR9gh4Q13TnYLj8xDPwvmHCt0KRg0HbSGaYEcqnvrx7hnYy/KsnQIqiSzaQi02kwjs2qeurhBeqbcaRXtG9RZHlgzne3Dlw7LSCMW+UZMNGMapawzfJ+xM0a6F2KVc3FvNnpmxRownPLtc1i8C2Zah6rdo4DrMh1D9vWGshWnqF7vgjWnmVlJKRTg7hq+4Xq9htXHWcshr8BBnU0JiJwkEQXJya+NlGkwtlf2KDJJTUwcTCdr+Rmhagk70YrYt00UYifCuIs0L6dgy1qEqkV3RT+tmdaaU3E5xow89x7EEYRvYWEolgxJxdQccyDf3D54MmcPPDWNGbS0ySFoNqfvia+sPbim+pLI19AfaPgkSEd0UagqyXZPGDRGD3sQOXaa5VCs6kaNQSndHJXPO+dyvdWLS6eqUMwoFxpaZIVEK0uFbQ4qn91bzOfktnvDUHYxpUpAluOgXhbc8s0hVZB2z3aaT8jLywO1BtJ0loNoymOO0Ig0JrBpZiGzPfkihZVDuHUu0Z46R0EHZobC4KjIR2GHoE3m3al5iFZYJvjlPjCpEU/xvqV6rtvGDIkRg1NXyuew13Y4gwPOU4rbjEt0LiiIhYB0LfcS6noYldDva0aFdR40TgYAGbsR5adHprfEiAo3Y+SwCwtLJFmOcM7tIW5VoFWwlIghp6GM87lysc1cZ8F1vHzFEPi9JT1ParAEXSi87fS5Na+KbAosjyPBrBEnfJqRzNB48VkRgfCKGJFHkubKe5dfmGfvONNSgrylu3Dn1vSL0APh85o5GlX7K5EQ9yj2eheT+zrnjJzcO+ygTTzicO/sgMo9DYory9spx8MNEve7/RRUB4x2BMvezMLiZZDx2Ul1PWd16CxYe95JEEgyl12U78gxYpRxTkwg5J4/O1XEGjrlnPgvObS2bRV3d8+CiTpIDRaEKgKG+xAmYPgaOTjQP9hpzTGCmDNO6ZzLOSOFFmUyRWrStw4A8gsvPP8EYKYwF6u3j8KtiSMzezvS98hz10GsquvYHaaRh1ViJaSKqZiHL/XmbKnBtDoN0lXJSOwUuCNdSHMcgnOlzR0taHM6dHBWOktoGHnj81TBjXsbSAPAdvTTQH+ExbQUnQqvUs7IxETbm78srNY2nSrjs5YkczCmrIbwYFL64QTzUgzX6x0PGHYhZQ17nSVVYujnUNGievZ+LBM/eXKdxJQeL7sHE2pibeBS+RzIG+/dq3ktRoSfTw7m0qBVSTiQ9hFEht7X33XQ/kLD6NaOEEJJTCdWS+HQfsblPZaZ4slKJ7ofKZWH2tEZdhzKjOfAAfveUJibMZiimLh2G6ElnC7GySpVeQqTpAB1KYOOr8qa2LaC1hb1sjXSu7HoyoI6OoWgjSrrlLwzO65uMz7ayZyUa0qCykbaveWEtvdgR2lYq3Urz6kxEFDjvcIwNA/cB5znpKAl+1zDOzqEHblgLM+EmXHAembNOnCz+fvUzNSFm3O9Q8E2tIGZ1Cwp3O647iRgrP3pzwNozCnpMf/wgiGycoBg18k4VbC95pTJUnTtMevk+xBL68r3AdqD1MsW8wt1QMcxwssvoK95smjnu3VnZ8ScyBl7a448T1D82TW8Hw1WS8yMB2c1chxQtHGCn8n1wkt7uM3MCrnzTv3YW+i1BKnmnN1RAKJJ8/KtmWvafIjuXPMWDrSTjJK76zUsmJVgpi9VxE4g9jb7xHapK4PYFguhtxEcbNHscsox6DR6YbWjk/6K8IWSAOt8Ycw5sd1cwro7F6fXhYFiMkY0zoC1JDQSC8JO9txmOWigThluYfq2bRtT45htzljYScqjc9wZ2LOVuK0z239hnPKIEqmgFF+Mx67oypXR4d9xcb41MG6to5ISrY3mXR0iSWxhvCngi8JNZeyWzlG/ypIQ7TaJYGAW5nFemfucwZJ5hnMbJ+rscleVGCudXHjVjSZhsjLJlFg0La69FPVj9HtRsZp9bAr0OllxyJROKnJVYK7WnycTRQ7GsWAjuclKgCehZLYcENvoMyzua6lupZ8AjBQzQjPvWLdawt02kj3z5GWDEONKQyPyQsxqYIyNneEKHGSIswePoKVTlyzcHQknKjyjkPn51X1udVtdKt/9ZJ67ukMBjdJGhTBU/k7RyTc2UsYZUnWX39GRq7MclTHjB5vj/+6qXGJGhOTVuKz2l1nr6gZyKQFhArygmCYqxCBTFC1tUy4K3fJnfOzNYWmMcJtOgGfAxMglhTI7Cr6xCBdmmdRw/+gbiy9pXXIp4RjgkeBLhxaQJ19nKYbjeoR/1brkaaueFnswLF/SDLcBkVyqCoSzJ1xOMdP0iyHf6/Jd3F19DetilX3QBMrmUgpnii5xqBMQGvLt7e0TBLVTGgEsbrHiViWwoj9+GPkxoMVKOoUosZ0Xtx3LQntGmcMLxrKmSIAtLNqYwKXBDiKnIYUxYLAowBZ6DOKthr73sN8WW0nwgpvsrXZt9LH882lahsgU7mGk1psGaoMGa+Mei0vwg9uZrxlRDNKyRQ5ClyEgo3uPoxELXQmGERnM2YHUvSH0CSFWw+VyE+ySTMxSjpmTh4GxEqzb5pUpoTdRVrebDeg0+WMcQsklbEwwl5/YoFjPzpRSnFxlBXXwPVvY3Vj4k9mJsiultyrTqPYj5CutudEYgYdLERvsMEvBLsonIaVYUII6Yxg/V06FoKQEXYYrJCqEe8Ew4xB9KwGLGqneYvNEOmCfS43fR1i2D7oPJHawJRc6PjfScj1PRsaAK/woMYGQtHZ+HhkwagbmnT/iUBEgla1EGFYutNI/ORF05odLM5NzZuIdzVH7YIgYbeRpASPRZ3SqLNSaihaSTSpFu5OQjQ7WPnr4y0Vyoq0Ygn4K11pkAOmnBp2V+4kNmQOLiUjrICogvNwSySulOsVa3TEI5bWDuTpbjYLJksV+rbUwRiFFUmDCgn1k4lq2CgxPN50DSMhOoebZFImjU0F71wgvUwET7xendW6ndFBF6p4MMkUC0rm6U3MXVkp9rLmhbGHmosEbNX9NbK0IaTPkBw8fPilV1s4ne+nh1beSzIIua+uyOcj6cDtneG7w6SA1Dlh1AQxmAW91Qz/lMqgqFi6qi8shjMKYW38Jmd2QLrOSvTIX5bQQitOiEn1NQ+XOIR9SRq15CeyodtfB3+i0mYvDOK4wd2+us4grm3mVywevACANstzqhBGcya3l29HDVNIxdLJpJM7qbgq5bTXcU4NiLU0Aq2x1CI20ujnUnsjlFzQKHFRBl0gOLDKQ44BTz1UDSowlHp3TfaakvlV7HzipqhfSXzOH7BrYdaqoVSEHxbefaKsM6Ml15b5kChmVlnfWlUjUtmg0vDSVLthaMF76GMx+yBFfW2vl4ea27k322uwg1V0gPIhy6BmMuohwpD65Iy9GXfcQrJOuwwIuUhAZTgIxzvhoVqc5leKP05gsJFZqYW8TGvWN7gN2McjmnCh83+oeXE1u95g20eGTRp4IcalDnycDSQ1s3RNtkLFUIp1vDA9sGq2FdkOxwvp7JDZUVHYmw3Gje8TksHaSJp1zCbt5SxbiRLlEn0XALbzAJCBMp+iFBf+5ieYRzLRJmKlF2l/ywCbS4j0ygXM5LPHuOfZZnnwiFxzHER1DKXkxEe9Zt3TMk8RBl6wbfJaIQc6J8DcJBJJCaA6mWZjITA5X9UAjWl/+esbCapzcyycLeN8rfN9kZJYQhDv8Nk/xFPnx40dP+uHe8wqkkdjEK9oRnvJ6AJNDaZxyrovJHG0sR9+TA+6cCTe3N/Rdoisr0wflYGn8gApCSimFziExgayPjkxBYGIGSUoeYerl04zAmknmEE7JaPJISvTZkgW5AmqM1bsuD5zM1+b0zqoUtahqExnEwget/x/2AAr9CUMzDvVFTc2GdjARzlYGdiZbLIlBw+Fj6Gw4uAwND80ERVBQVa88EFEl23FAKkbLK1FRsEiKDAt6HSGFNUbkd3MBgmFLsmeQdYhRYCdopm75ZLI5w1yzRHjSqhRjUQdVMYWNtZXMLnXF2kZUZ3O9RhrLm75oUFxyFA/OYtG8K1HMdor7PEV/CoZNZKM0RYUCYaqpCADZ64R5XqaWSDqPMRdcNFdipy5KOVglyzDeO1J5YwCTTJtlqb4q6VLqulzZHWbRZLGSF89zxJQSxrHgv8ywsWANWsKgpXnOKeA2Qb4RdQr/3i18pc6DfU8fTRyWy2Z/cj41IGV2SPPJ3nMjy8t2YXAadUIsSI59XzTWseKQo5NM65KUSaPxoEkTLEwII42xQq46O5+cFoHhVKgtLU4KCCslKdxpZmorcqDQ2kaiap9VWmieMBYVXuxN75w4F/Xc7hX5YCns+UETSi9M88lmyk7RFvke2WKe2I/zdGGJ5KF0WJARpzlaZ7SGMRZYSEt+9ODBE8WSdi0WLryUEqvOHJnncuadvUdFIiAwfKxOyVlePSym1GIyILjmmkuk9Is2IFk7Httop1x1DXQTjusRHZPsz92wjxTcUsPX57JdfPEdR1itBJzd3fxRkFQpha3tIIRjixwg6IEHrjGOt15qcKa1sD1Ip8VbiwQ/LE2EUsxEqXW6Z45fn8KDUzh00MF3qbl9jkT/HBIOvKpaOd6iG8JA1kti608/nC6aIaISbSehYug5iEunU9KcBrHzlDGt9ns5mi7cONNJVeZwRkPOTNhr5SgQTiINOsuqAaArAKvaZMxgUNWcVhZHQghgo4JTkpy8lk525MpeP9t9B9+fVM6jHRR0llN3ke7luKvVWuwnWwd5Xrkx0qmMuaif56IApyTQpLiEvt7bcbTlBttHaAdUbIRpI324zt5qmocZh8PSnwgizEyEdDIF7TJOfmER+MZD0M0PyQaKXPkSKXfulIug9Iq1mE4GnpB3HQuFlYZG+wx2lSrKQrOABZnqchZlWqzKzrAnh4AQufdy/41wuLk6BJEjSs4cnqcIFwOLFYeZ0knELHfevt4H7edzzaTDu65OlGPZzDgK42t12y6Yo0fOvGBGJxq4427igC2YhkOsN4s5kJ6D7/MZKMk42ez4bPoU3RAzPDslprJLIZyYHz54+CRTCWzZv1RAQ8zWLjXzYETc7L31BUWk5GpcWmcczQ/oyaquCcYai3mQS8ZxPVyYSEsSVYFJGRO9RWavP5A1wIlozZSWNQHdLfsY2C6XqCaP3mOhHcce9iliLk1ooC6q5NrY4b5pxuHuYCU0YMSse+/8Ho32FDOYQtJI5KzQF1fFO6y2KkHvxJh73ZYFgVOF14Egw8jbB7fMOD95MClmlAaHorkGgaFYiER9M57sE1iduXcX7RfUBYmSTV7/6Mu8Tu9ALrilEss1UhNlXxG2G6xOYdHO+0VE3x5eHL4n+hKp8vP3kwmjckOMXUzY5CdD3Vbncs6pyMkCek2R9oeIts30JlIUrOZVurSVs2H04QphbCQ5ArlUQiOI1D11IaFoZlLjmKvIUNWrZ+Ziu7LmgaesGqnNdeZ6fTBjbSNilUtoFc66AM9f6XQtsJMR5JrdhV6Hl64lN/5rhGwCCiIzs3Dg3LuyTOxE+JiRuS1atJyh5R2lZ6mfHTAtcf1Gk8LeVwCScX0u/VBa9jVR6DijzudYKeZlDosvsWutdc11SZVHpGguVugYPsc69mNFG2PNcqLb5JvQObBSAUdAXdq74dSAZWbpXUm7x6IstGtJCTiuR6w9sQNFpRb9N9eCcnLozieDSRWnumB0EU/lCqUFr0pYLCahCDn58XPPPUlB/0xxmOgARODrZMFcKjG1fmJHeKu7JPcpctKdWot7h0pYPUN5zMs3KKxAaMo3+8Sg2jZhtVzyERJ3X0mGc/YIoS8nNXQpFuyweqnRYqor0KErmGvMGZbnB1twWSvLegAMwqqloPURuKVyAUILQhZE1vc0n4O4QaG3ycqO0EtFIkuHWOY4bxKqnefJyC5w0HDS9YjizkMqlcR5zAhfr0IrCL9oe5ip3TP7i+jZJTw651qEziCdrDDSCjZqbcSg3as8/zsaB4ejkbnFAyadoAdV3LXUlR1hJy8icy0HTuaTM0RdIyAVdUbeXU6nXfLyKFuh/xRCeOfKc7cJmcPnDxpihicZLwGHNCkiLCsPvdFwzmcxhz93svQ0F9Qhp8yBs57JNTYj6LIKvlppkAtyUSyx1qLmOIKTVNCE8jiM/NzIxYe5mQe6Bv7UFyjHwk4Z20M06uyOEHmlbwZrk5fH6A31UtGOya65wEpx89S8YnslUJYfVg7n3Lx0J0wtVKa5umWtf7dqd2WqhHmtHRRCphPVfEaX244Wa70x813zCkd2KEw8drqKp9B0yK5prKjJe5ECYXmjAmxO7PtOV+gcSEquOeAvmSVOrnUvIoweY2BH2IIarDC5sHk5PU/B740kjzGXF5nIEmFamziLnmAE9Zn45PO2Yz/CXUMzsXxze/tEQSGWEy0gVvUADU+VW9GoMrVljZAAVBqFSdhTL8VT72zdiqqSlWEhNpEP4GXGmEI7EiyKlEKlLo66tBujLxNHOYaGGIgceWcurerkXhqaJYYS5XtVTNnqyTxQlhZlUXOpmBeDTMlucyiz25ZdRVqhPV0h9WaoPPwcRlgi6azEMy4G/5wWw3fMEVWBLr/IETllgAt+SkxVjFRD5WdzEXss6wzygiWRJkasCwVLyS3ZaBGjTiRbjgREGNgdLYEbTvb9ueRwI15hNh4EJgW6DhRlEGTOCjQUXPMVV8bLkuYksibEgeiUtpstPqvTTAeZgBYJmQ65LmKFRTRpWVTZPpf5KD21pJXIufh8COPe91BedVjxE8ncqGVZuD2oibGAbPxdlbicOlPiRLkctFkvYutQ9Wxh97Jyzs8QiHIdYqgfs4vCA7nR4cA31aCQbV1hidTO1an7PMef58Cg1iBzbTngq8/FUFSHqvmOhcGH2Wow4QgrkuYrKGn0DkS0sVT2M+a0wSiF/OAsmHH4Rd5mVlZXLUgVyb3aXF+EMJXUDFZ5JbONSPhT5HU4IWSj/oluxgHPD3C0G9Tb8OljUaJiL7zEwpo+n9CH5QouOrnifZWXE9IKnnsy/tT5Jr2WnDBkbDqCJj7XvJKEqK1WzwNR96HBrvA+2UKr/ZcSVI62aq/FaFB2Acx8nnB0+vOkcKkdY+Jy2bxFbz26kNBwlEKrg4kTQoJ2dOTKYRESspV1+LIKX3oGivBoADiUMTyWLbGzdRBDq8LMCAXpQDDDcCxRqnYxkMboMdQSvnk0z34fNJhE6E+w8PaJqJR1OctgUrqX1pe2QNXOHK6zCZ67Bu3EbyWIHOEculhRIPNlzaQWz32eDB0Fr8lqxjuMxMSzGe+mEUeutYSleoLYTbhncic1cy5LDOoHGqNxI2kxn7If1vwpOsOTirn3I1hHgjhrreiEPEvNJ3x+RIseYrmxuuJ27EDKpM8esJRJIUW4qPbevVupG/brHQsaZs+wSgTdfW1jmt506NcobExIYYQoB9i47FMOA8ZkOewuZIuxXAd66EcmUhhXFiYLCg7RhbXvexRh+hx+eFSM2YO5pz0sq+52HK6QllIkiSY6Y/CqGenqVHrA21LVL1+1RLICAsbV728qreeMzuCc9TFp3zPGcOV5oRCP30PQatvbqXLnrIeHnyMIFvvQNS60sjeL/TCDWk/GU85OJuJZIKNEy57kB3hR0o4e/733FVKmea+69DkHc1+WW4SIBp0wtwLwMtmlipQdkUK5LEfmmMvPi3ZA7ti7oN/emzt+E24/mmtB9rs9rIuO/fAChc8qCnLCinbq7nUOVuWZPHr46ElQEcPZdB2Kfpk4wyIk+2FPsm42nC+BbGjXhsvNhSrQSTvlyYWZIlrUD2HfVOsCoEAxe/pZP45FC6Z6W4fYZOssN9XGrGJZMChkRgZjl5sLTRdTvJBjX5RDQRXnbmOMpWat2+bVncLpEXHBcSBOdmKeH7Ba7kRoRLMkzxYYoeYPxlfQqec9zyHZfcjepG6kMIeFN0VVZR0aTkfk0J90VAmmxMjpbYSgMy4PWosf+06h5ggmVL2syNE0Z7DXNF+RBbs6viJFc1J+uA9J9fxqzStSWEQCORgw2z0GhKxydSAj4gBGYL6iYEMW9tJcJGVbuJp/jIFKEescDZda0KfPswoPlRjql4o+DqTsWRqtD1iRdT6H47ayu+M5xozOa/yVkTIVQBgOt5rJ+fdbJpNBy12O9IuswDmuyCUD7l6r7iElC4GselyFsIWPU/LZp+Uc9t0xIyQUFlHME79Ih7RiFURRB41ABdHELKl1F4LmlXYac44JXC6KTqXdvysvOU+QNRLuzV5bbzzU17pXQp0gQs0hQAhcjKtEOxS3bF928VK+I00ajhIOJ917kO2nhMMzaUExwZF0KBouZ2tBf2VBp/naxFjZ7uHWLBsYWvzUHHOPyExpnjhZyArVhW9BRAGLpbRYWa3ds7OxVMJOxU7RyEgrw71xvttbw0z8PMWQHz/3+Imq1u2yRdWHE9YqpbCoqi5eGlGBHYdrFib9m1SlzDGwHwe27RImc8pP1q1ZthKdAa1WYlgOPiDd7rJwxurkQniU8waMgTZ7CNf2Y/dK5XIJ+5VBuER5yYOVSNlKdAi+wVpYh+hlZ+aZiKHgz2my28nk4deFFZNJArkWK3PgxC13iwZE9a12MiqLtFTexgrKB7YjdDLip7sAMtF0bSXDzbRU1sd+hEgsYSXn6fuMPpgr3mI43akYFkR1ZriY3FhZHad89g5baYtjTqJTI3JZdLmEQaYiAthJZEKHY3Y0zrt8dmIwrjuLoiXfi4A1M78cUkIZHcUyLqXgxoAyJ7aUcFsMZU5csqGmiZIMN6XgkiZsTtxWpt0lQ+H8LptbeNetutndfni3LlYRNT0p+UEjiC7ou9SQJM4kVBlr9qEUSJCGOwdZc6VizhRzMdnBh7vCjFM1EgzP7K01MJ1RbcrFODqXaveiliOiWcXalJnhsrCVrchaF0YbfkTksjzs5qkjEW3c4dASxddZgCd6bmvd9RlbjeGyz7/yovCmGRKBORdx2Svygcvlcq+oUZQsCI16VzbC0XvQYkdwqfE5d7oRCNZxVmqOWUwSrBlssunsLwbi1VLDvkgeYYk6FhnTSk+jLlpMP802RR1PTIssuQRylHIKAkBY54T1i50cPQYZthYaqJxSaKlqrailYpcTL2nH21bDDaP3gXz74PZJ4a3oPkCs6k++SUtcOE7eOTMWhGOXOQZzYiKEoE47Cis+8Ry0IiqwLqp0MkmszEWeZD6kuWItCymGR2unAfXKI6i07lZCV9kq6Y4pAmBE2zNyzC3cUi0+r/tmZRoXOswmjYbUoWHwR8xR+oRsdo/KaTmHwt5yCl3KpN4Ac9khOA7aY86gA6gde4ij5NVVGSs8IdFUZVqfX5zLDI7xsGTiRGzrqWOReFQCs0y7l9ZaXJDCUzUUFF1TZIBztKrU8Sktum86OedKsDijvHT6+BySSCMqeVEqI+wnrCUWTptTwqNiuM2GbQ48yMBtvfjFAeBSDDclo4yJrRTc1owCw+NacEFCxsSNZaSRcFsy8py49gaAnYwhIpq9yhdjMtFV4DjFsiIU3JoTFMsoKaEWgyFjSxOX6hffcVzJtHHMPnQ/pM721lG3EhBeMmecKVlQzDtdqoNZJoIpV7TmDPaVDnvZhUR89TwNs5VPD0UAEBbsbvWuBid87MRMoiuvJXZovCQWpKmYX4Znne3JM+HXsronQZ29EXoNV9/Tn01LSwXzQ/vueheFZae+7WidrhopyAZg3pA6f3X9k8WDrE+cKanngyAIKEPJdUc9qNNJqv9JUgbh1WI5qNtNRd+JNCAq8Yj5lMXlH8P+OU+GkljOClh9QAo3iRMNN6+ZsNiPyl7XGtF5ZrRK8myVI8gk+dEj98ICue6tdVoX5DD2UsZ0Jg1NQ1ajQGYaggo35eVimYElDDiiEVmSjxZD3xPt0Bf1tdAM2aspWYdLFYsTNt8bBUASRTWvNNfAz1vE/biimC9YzSRWlzADFrHAzJcrb+9tfTcYdQGDWoAWrJ2UgE6mz3bZ4keMMfjyaQ1CB1unFnob2MfwaE5W8NvlsiJTOUs5O6omMyblpXsWLUsc5PYaJRd+l7RmJK3R8nM9HwkDZegniIzIsBMQaO+RconhtaxF0tk6OiGG5xOr+MhUDE/6bTn0IzO5EcPF0TuFT169Hf1AzV61uZeYIRvumfx1zjSyJRQkPLKES5q4LRk3JG2YVfgZOVBJPmhe3GOXEZ1gpgnc5oyHxVDgSX2ag6l9lzGd0Qp8zBECuDAnlAX7mB7fGgLTjscl46EZHljCZXbcFL9IbrO58nd05JyQaW0h5qJIA6XUBRlSxCbXaRmKeu2wUiZlca91Kb+woHfOySiDBVNNDq37GKjMBMrSI0yl8PU49N1Q1CKq1ZEFo+WLIxaxDvj+fM421hA+Lzq2hKpr7mILl59LxOvIAk52RDkSPAPKYt5N6x0bNWEy8RTbSDo3J0GkMG30NSlTxEWPFmVXzNHK7ymdj+ZZHh09qW/xfRu/DyeTWriF/Th1Io3Z6Qo721j4rU5unTXnCN0xJs08+8ow8KgqRysmZ1MkQ/XesdGLUDOWBXtO5FyDmScILL3yyiszFyN2ibDabjpITOZ5bWUdp1N28ZzYjz2Gu6J3llyxH/vKDsHCUpHgUBCZJ8exuzgRWEZ9hYwgWi0f+87hlh/chcwUQTG6GNwSwKnBd3d3uL19cEoZ883S9gPbzQVt35E0lKdVS1gncxOllNwNltnmY3Y+6ItvJnUwBuyH8+pLdr8ZucWuDCF1bTl41HbCM5UY1uiL1ceM6jWxGjiOjkoqnyFFFOx+HGSS8dKanf5g/Jk0leyc7whGyxTVKefcXTot8gp6a5xhjIi+VYay+PhWchQgaS5BndMGcc/PbI6TPkBaEzne8gAZhAGy7M7ngGE6HGUFbRziGPtFgOSxsKPjuZpxmxLebR0PtoqH2ec0Wy245IzNgBtLeJQrbqs/+5tsqGbICbiphmMAr79/hy+8/RRfff8OWzbsI+Gb1yuOmXAdA+1kHhiOz2FIKev4HhRuYa43ueBxSci9o42Ji008rBnPjo6tuOj2vf1AKhNtANc+kErF3b4j54TrTOjTFn6/OfvRWWoNKRccd9ewhE/pfn52qRXXu2ceFztTBLv5bKu6o/bRnCBhCddr4wB8pfmFkDWTkhsIwco0ESVVnU9i8Sim0n53dUQBS0iZGMQmh27NOmXXP9MMI0vFKiu19Cy0kx1Jp2VNsox9v8PlchNdbmtu9CjvsqUvoaiTOpZIwyRq0MlW3Ph8x3T6eQTkUTwoPcuyOEFcGIK7QVjUo4TJgu2+9psyaIDw5pocMxzX/d4sbNs27Lt/VxlxCpkoch/mZ+lHi05PWT5upGj0uvKmX/NDsQ1DIGsncaIB6eUPvTKlbI0wIzlqamjENDEw+jST++zB9SNUj7IMGMOr310TfcJRA8tRkmxU5EL1asq+QE7piGc17vV6pc+PO/YCzoRQ2pgGrMINr/vuzqmYSOQtu+ldjiAh5ZAkGK03Oo0Q2dKNteAF8Sx7bOYR5AXRKChP4fU5LWsI+Ujtx45SalQclXBFfPbkVF8dThpWK/N7nL1rCP/UjQp4vqvt4nnIFoO6Fmr7QgX+cRwRIjNOmGw+46UniwicIl2lHu6jrSjRZKH9KaTpuvZjkmLonYqlhOu+e5JiAo7e6ZZc/HuTEj3b7syRlGDTse2OhBsDqddeXlzTwDPapT+sGdtx4NVHN/iNH/0gLgY8qGTiwfDes4Zne8c+B9591nGHhrs2cO0Tz3pDGw3VgEd1w8ce3+C1Fx+jA/h/ff4reP3uwD4m7sbEswk8Iy26tyVai6mhMmKYNd57w812gz4aHhnwMCV84Kbgt37qFRwwvPH0PXztnffxTh94unfsY+KYCcUSOiY6L/42J7oBxwDanGgz4VCQ2c6LgOh/LQU76cBprvzynEuwr2RimdgtpDGRWbGmpZIKzQnIOExjdSVprpwSzfhUpevQHexYlWxY6HmWIoN9xdIKc1ZnKTZVU+wBHRhlJHm9XnG5uQmhIGbyPA52vyK/eAxCDQFj6GHInvMZLrv86Z1mZ8hVKSWIQ2JMhqcAIWrNTM8W8dqfEgODDMPM739OgFVWTrKMXDMmobCIAx/gIL6vJNWc3ZmBcyJdTjuJS5lOFSFZKE7nt2LhEpCSFwyYE7l69xS2OnOZk6oAyjn7RcOCMr366kemRIOHBE85x4BG2RDum7S0DfvRHHYhQ0GpVk3sqmNiu1nWILWWsBBBcmxfKYVhfxLZvY4f3tzcoLVj5SrPFUAlhWxrRxySoremML8ruN7dhX24blLXR7ht9nEcyPVkt3EyD1R1adSI+GIcoVy3lOMBX24uOPaDoh3xz+GdS63OvkoubFJAk4ZncyZW2vPElnHVfGaFV0nR672jbJm02R2X2wtmG0GBdANHpzsGn/2UliYrlzCQMwmO5hIgERPvYfu+/KK2rdyrXIOHTux8nKwwZGxXzC1HjuMg5KMhKbBVMr1ywUblbR/ApRY8rG7VnjBxqRssJ9y1HciGmg3v7w2tTXR2t7c54aUt43d/8oP4K199C198+w6lZDw9OuaWkWpF2SpuH2x48OgBHt5WPPfcQ7zwwiM8enTrIUNHx5tvPsVP/sQX8dM/8bP43IOK3/Xax/AXv/w6fuqtp4BlvHHtuJrhbm+kpV9wtAOtdVLUHV+X8aE2423NeJQGajb83m/9GP7ul7+G2oFPPP8IH37xIR5sN3j/euDN4w5fePN9fHM/8I27K965Nlw0E8yG94+GnVk7d61hnwl3bUSinlQdtVxw3e/CzkRCu4202HDOTRbQivJcaqmxr1RJyyiyUAhoHJIPxdfSrbo1Jxkc7Yi9IHi81hqHq2tbCGWHrQ/iUhFRoxMuzMVwXPdlakkn5cb5I8hKK6WEX9QYfWmGGLEwYp+d8q0pHJ4Abm5u2D0iGH0y75SrQr2wwHbQE/txLPugWtgRgM96BrVKGTueGNjCZgcjhY37RMJoB7bLzfKhiqF4i32siIT9aEHkWOenZ95stFcyar00azoHi0nlj5FCCyN6/3a5YNA5ZL8eSwArV+SXX355CjKQdkALSQPp5Z2/ch6OYw9x09GO+waBySKsRq1UJsd4C+42s833PfDaXAwZDpOI/gukqPh762zX9khEE4Msk3PvrJiCulU8u7sLJogGUXpBiQNmQWa6gJZgzv/O425HvVxC45KSrMVXYmPdarR6B22owxo/XqguPfjQX6mJWMLJqPnGwLG7eMgi6IttbYRznZhcrEDEQJpwy4ODTqdyzwyKMGc9mgkpKlYGhPNkHS6GjqwfCnHQ1hq2bYsq5TgO5ruPlZ6HxXuXsDLnhAEgY+JihjKH95CWYaMBlpFhqGTi5OLV9sOt4uUHF3zbiw/xgc3wOGW8ezT8iS98A88awqbic8/d4EPW8IWPfhL/7H/3B4E8cbsV3F4yHj684LJlbNlQL75eUnbBaBoDid8tJUPbG37m576Cf/3/8h/ir/6Fv4F/8Xu/A3/xK9/Az7/9DNc58VYD7uYKLRpYBpBzJrcFojhR++bRVlFnwy/5wPP4zI3hb97c4Dd//+fwn//Vv4dvfOWbwPXAizXj1Qc3ePl2w6sPblBLxZGBb+wHvvLeji8/fYb3egc6cJMSrrPj6UzYx8R7dzvytlFHtLQVmm8oClWK46B2Rrb3JIQ0l2hSFh2kGmsOYXQ8EKFBmqliXqG619MIFbnmYHYyJZXiXBBc6x2XbSOslO9Zelyvezgju0aC85R76w0o1XUeJdN1uns8tWje22WLImm/XtHnRDVHD6yQhXSalbbmUN0ZwknJ33E4NpBdqJx37TvNRY6j4WbbAoE5jgMmtlx2p1zlx4j5JWudRPJBmitgSp50GsqLKTqxyBO99wjtAy8Z13EtcsNiszpxoB+Ey7PPpMbJ5Xf93rkYqCm5nbviJDOn+zPU19Nt121loMuFV61gUQYw29t2NFwul5gfzCGxXeaBP8P4sLEjkeGhlNya+ktApkP9XL0UsoEUkAIxsLBU6WKwUHUQFEN1PPLZab151nVaBmKgt49/hx6UvLAKJ058jjrVAE+8cAhTLN6GF3ZWGIrrtXAmngF/0SaFvzD7jIvUjAN7JCS6y3rWcWZ+xCQZyILWOfoMDnxnYJB0C5kzKc1JZPYmL7IUQ2+EQjjszqlEP6vnxfqxbJybOJRZsutItlpRMVHNsCXDbS54aIZLqbjJ2S8UDmBzTfjI8w/wnR98Ab/6Qy/ge156Di9N4Btvv4cff/1t/PU33wZ6wj/wocf44feewQbQJ/DaC7c4nj7Ft/+678d/5/f/VtxcgJdefIjLTfFgJwwcfcf71yue3e14dm14+t4zXPeGd997D8+uHc/efx/XdocXXrzF7/0dvx7vz4T/+5//6/hvvPZx/OS7T2FIeNonjpOjQaKhZiYZIAR/nC1dasEFA5sZvvuV5/Cl19/C9/7Ar8cf+tf+AH7rb/5e/MAP/Bp81/d9Di9+8hW8vm340ffu8Ddffxd/72vfxNfeeheXPvBtzz3A933geXzm8S2yGV4/dswB3AygE5q4ksoaGSiMKJb2O1P67MZ4LmIU1dhniD6HVFxzzu7ZJTacSDKie8s+Y0zSs89+ZxyAi9Zq4W7ganpBw4LClSZp1FGUk92HTDyDpkwiSeSbc1BtNBTtw+OrFQzmLs50ZIjsb1OCG9Kphusk/1gubpPURqAUEu9OXhpSjTuo5fOmqsIKM+Y6oYOKbKEe1GoVuMZLSyzOUkoIp4UCaD2t3JscokvND1nVrawjMljzCV0RW8wYYobI8fE5iTy5InyOcy9Z7wM+ly7bKQsEMFiePCA7tq14+AxtEYTFHzQVzKXg6D3ovOCNp+HWwdYquOJJIjU/UCv1Jvt+Dfpaaw22JVzvPDZyEs+VL4yyPjBmKMoF10RCIofSl8sFKQH7foQlRSnVb/+50uyyBGyUh1t2OwfpQkQm0GKaJ0peOGTSAt4A3O0HkjVunIH96pXPoBBPL7cfDdMoEJRHFAd/k3huqYlKcx+st+ZzmQI4bxvATFTXz1Nnd90jUEadVS4pXJfF8Ni2C1tk77iiqzKDmxzN6LbyOfXvZEXvrThnDSQ1nNvqPg0bGVE3MGzcuEMxHtkPjVcePMQnHt7iw7cP8HgC795d8TNvvI3/+OkzvFcTbj74Il77ZZ/Gb/qub8Fv/P7P4X/xv/wjePWtd/ChmwveuTvQ+8B7dwdeulzw5a+/gbffegfvP7tDmglj7DC6G7iZZYXlgn1/hlJuKJRrGMgY7cA4Ju76FV+9+xr+Z//Dfxw/+fkv4W//8OfxfS9/AP/JF76Oh9uGY2+4Ho0JfN7hmkSSAxgNnCsV2Ey4rX7QPy4Zf29v+JUvPsbrX/o5PH3vfTx4tOH7v+9b8et/7ecwe8d+PfDmW+/hC1/6On7kJ38Bf/fv/zT+7Od/AW/93BfwqgG/4SMv49d+6qP4K6+/iR956ymeLwVvjw4rBe/eXdGnOwWn0VFLQoEXOz1NdHbzMvXcr65tumwbxnCosTI86ZAgUVHCVqJTPWiVExbtfWKIYcdfs5oolrUI1pKGJRzASw6DyGM/MDLFc8MFbMr9mGNgcM6aSXM/5EXWfJgds5Cx6Ma9NbTr1Z0cEmeujGoA9RnpRD+2nJ1Iwitru3ENjplidyfKVnDdB9XrM1T46H7hjd6RCgEuElMclu6E1+nGTGZpyXLGdisgY56Noqv3/Ur90GkyrHPJJvbrHkWcjwqkHAfMakB+iqkIUa7Oyy7hygjhcmdRIGisWgkkxZiNVJS30XsH+EzlDukCQQqjDKe8AkNO95WPvTU8ePDAbzameCW5lZbsFcyYyCWhHS6ESzOfZhdAP9z4bN8P2rHLp6jjaAM3D25xfXbnAiBS63QxBVbfBqZ529raEapmS4bJ2UtjxzIGwjMqJ3I658SkN4+gnKyXGWwtZ0CFm2gS5c1ZEHUrvN0NQhOOo4d1Rc6JCWZAP2ifoDwVtOjQnAJZTpAZ6bSEyTAV1Sg/H3Zu7eBl46lv42gwskkKI1uFZ/fG7u5Sac+wnIoVr4vu0ECPXBTvTHNaz2i2gbptjvHuXvHIkbYPp85u6CgpoVvClhM+/eJDfOrBDT766CG2feLtu4aff+ct/KUvfR1f7x3lA4/x2c99Ar/zl38W3/tdn8WnPvEKHj26gRXDo4eP8T3f/e345p/7z/HidsE33rt6kZESygBef/rMp1UTSBlIVpDhBA+QGdP7Dpv0UEs++O3tDpYLYANmFUDDO++9iz/w3/tB/PP//T+M3/lKwVYM77c1nJQthuCArbjIsM2J1BtsDGzZgOEHyo0Bb/SEz37qVbR+oLcdwIb3nt3hvfefurTcgNvHBd/9nZ/G9/+Kb0f9b/52vPXWU/z4T/8C/txf+rv4o//JXwN+8gv4A597DZ96dIs/96Vv4DkzvNkaqiUMV4ShpoTbCdxm4Og+L9q3jCO5EWJrA9tNRT8arWBGMH8KBYftcOW+2Fr5VByqdHXBXIusCLEhc7D5Og+XNVcJTyfBsPR8GpwN9HGwo01hp19KCX2SZi+ju8uy3SzNmA5ZaZ/qVjD7REOnVxsrfsg3q7CiLkjzHITsNkruN2cxu3An3jVYHwyGs5Oe5th9OO+xGIzIoKtvG8C2XTAOZ4J6Pg3hUxYlclbW2eLhW5lh2oaGRoZWDvSld9+Hx75jq5uzabGyPzp941bW0oDVzeUEvftmad3nV6f8oUFXCImfZS1VRKcsTN5L5gZnUjd7TjUpgdlhmdY6aikoFya6kZfex6SPkAUtTzTfWiusWCySUPgl2ok0p70qPnLC4rYfc6LebLEYZGAGwi0lFxz9gCE7hMAsAq+QS1x+a6iMWIBGVoir3lswhVSlS3Q4xsDlchND+0SLidF9BtJnJ0f+whhdZ47kLGXswJgZpW6LUcX5yXE0TGteRWHGz7dCquLmSnf3Favk168Wc999yJ4oUupNcxmElmd0WdJ3ZPL561Y96KY69NRbR94K5j7QSRKQkrz3QUuUHj5EvTf02ZBrRZ4TyXwYWB9sOBi5erSO25sNlR3k7aXiBz/5Km6y4endjp95/R38nS9+A98cE7cffIRPfetH8du++zV893d8Gt/+mY/g4eNb7O2KMQ3XveObbz9Fbx0ffjnjE5/4KH7q7sCHHt3gxwCUMXGY8+2fvvsM7sjeMUdCyTdox9WNAkdDyhXFfMCfq+HoHfudOx2EMtgK9rsD788rvvWzH8Ern/04vv7Gm3juUvFW2zEmSKlmrHLdmPrG52eGlIvPVzBwKQUlUwh5e4NPvvoB7HduCU5nfTSZY86Eu7srrndX1HqL/XAW4i/9JR/D9/3y1/A/+id/AP+3P/EX8S/+W/8+/uBrn8Zv+cjL+LNf+Tou2XBFhmHi4Q3dZrPDwJeccZP9An/n6LizhKZRMkkmuhxqqSEyKyU7w+9Y7MXC8LDWmtuTmLRGNayDUkpoV6n2mfZnCbOnYPTo71DxsdP5OixYjGcR45LRjhhKS/MlEbKiWcdYThgppzB5FdTaOA/y88Ch1uN6jQgIffZj34OlluDMrEmBtLQfyRImGWVS4JuZ022Pw+etKUWEgCVDvhj2u91jkYVGtB7RuJDbOVb2iqUEowJ+cMAdOShEJoxeekNeWCQ/tdbi7Kt07tXnl3XTHAgTVjeQ9O+gGbWxKBmE9mdynVORN4rJQnn6BqiXiixHXOJyk1BKIgZ33Xe2oJV4e6dNMimd04e5XgmniI2Ua66wzjFIee0z3FDT7DE8T0iEIXoM7MSwYOLEvWpINGT9c6cat4b4bsSgCZF8OCKRbU5PSZvM/cBMuL25UVKsUxhJW6ybX5yTbq+GhENpiNl8vsIhn7jgUhBj/mJbcKfRdQnSpgu/RDWeyhJgm+t2JXV5SMEvfFUgpZQIYkoGwCZ6ZzU4QQERxVYQO2PBDDm8b1JEHIsy6OFjDpFsdCCtmLAC7H24AK847jqHi/FSmfiHPvwB/N2vfxN/95338Oilx/jM5z6Bf/R7vh3f9R2fxLd84lW8+MID7MOtS959723cvXW4pcw4cHN5GGmCd3dXfObTH8Z/dNnw8a04d70D7++Oe7/15tvY9zvknLC3CcyrV8fJYKVgpIKnz3b00XD3bEe2hFdefoS799+HleqHR3FGTed6+JZv/SS+9v/5Bbz44AF+6p07YIID0MJnvztsZI555+nuBjVNDOpHXnxwg/f2hpvnH+OF5y+YQ/CmwxzZLujNTR23bXMYZ8jGI+G99+/wztP3YAb8T/7pH8Qv+7bP4L/9z/7v8Ic/9+349uce44ffeuq6l+wHx4ceXPB7Pvkh1Az8/NM7/PCbT/Hm0x0fqoanKeFuAnvyA3NvTuXdNlauVFRn6poSob/ejshLv729jaFwPnUHilQQZVTwiswKFTsAJBZmLSKpR59LAU4DQj8F/GdWZnaLPt/COy3huHpFLwJBsYQ0GVZlDNIacHiJCIQCmXprPHPuJw+G9UkQSlIQEFqYI3bqlnzIvtMyCBIQMhDOz8IS4tsVy+x72y+eFkNx0O5kUsvUGg1hs0Pfx9Fw7I2uAmRy5oxCPRjbD/cf2xvKRhftk85jv+60i5JzuT+nfb+S0bqskIwECOl8ilqvsMaAW5m70ydbH96+ZobBdrNTROfWBvLM8Txq91JxX63GAPeSuUDhFbKUj2L9bNtGPvJAp0DnQr60Bw1ZDMVcMSo9hnsnXUJM4wM7Ha5hiwC71x11Dq2CpXGi8ynvQsPnS620CHf7Dxc+bhQVjoDo2uEwIBJjXqnz6PT4b+0Iw8ZQ/oeTpy6+HMNsudem5Nbag/nQwXRJ/E6cQ7iNwRaYsejYrXfUzKxj0hVzLbAMr4AjB9xospfDB6yPjmIbkk1f2OTJ15KQWscNbUBq8aryGdvcx5eK9/eGPjoe3Lh32rc8eoAvvvMOnn7mY/jj/+Pfi8cvXfD848dAynh23bG3hq+++bYzzSaRyumKdUsV7diRckItG/pMePXll9BvLtif3WErGSVPvHttsIcX9NbR2oF2HGSOATX7cPjx40f4f//5v4V/9X/9R/Hg4S3sOPDNN97FH/yf/m78/h/6jfjmN9/ywfAxkTcDRkLfGz760ZfxU9eO2wswiFsPxYTSQHLCU+/kK9e7X05188yYVx5c8MZ7O17++EdwuxneevcZkoSUCWj7FeVSMUYCGvVBsyPnG8zh7yonH4B/+Wtfw2/5jd+Jf+5/8LvxR/7Pfwq/73Pfih9++j4eMN1wWsLv/uQr+LO/8GV85Y138Ds+9WH83k9/CD/2xjP8F288hY0DL5ihTcN7aeJuOhGh8RDSMxwoPKjp1BuhYljeTSRolHqCpTuC2dV4scyhy0VJkR37vtNqp2Gasjn2laAJC8EgMN2ChLKAZImasz0o9ia9iPzl4HCSlRLMuL6v4tGzxYFcDYlCSXX3rnUBpnXU7UJtCijqPcKGfStboCXSb7SDoXujw3IljYHBezQ57F2uvF6dHpyBSgOjNNZhfilt24br3Qr/s2zox4G8XbyjoUhyK9mZV6IFA0CeS8UvW/vTZTQx/AwXxM2Avyox5RyYpGvDOD9ZgUQMbGkHmiCnlKOyjfQxDlnacURLKE/9cqkB+RzHEcNrM6/UciVdruY1SyHTQoFRCmcpnH1g+q3fjkZ//rRsjmt1FknyCsmoXJZX0L7vIaarNNuTzcJgmp0U8MrxGq3Tw79AVpayXpdISPRDWbL03l28Q4fOWjLzUBJDqugVVUtAV7KSvtxc3CaGAT3taIFTWq0ABtru6tFt25AMuLte/RNbCnt66Qrb4VqUUnNwzOXnpQ7tctl80H09UDe3ZFds6CQWLdzVshFCcnZVx8B2KXhQM5673XBTsh8CaLjdJn7dx17CP/DB5zFHx5YTHtxU2tsnPMwb0CY+9PghXnvtY2jHwOvffAvffPsdvPvOO+g7faCsUoleYLVizEbDLF/U6AMYCZ/8yAfwwoeex21OeP5CunIxbDnjenfF9W7H5bKxKvTNcjQvSN6+O/Dq7PjXftkr+EPf+Qr+sU++gH/r3/0LePObb3uHWzaUy00Y800Az7/wEO/P6TEAKWHvM7Q5UI4G2Sr7fgSDZcuGS3IPrpcfbHj9+gzf+tqrKMXQjivmOEhQ6MgVsDTx0kuP8dKLj/Do0Q3K5YJ9v4bvUbYE0HX5m2+/hR/6r/8G9A88j/eevYeHxYe7d5j49uce4hfeeBPvf/gV/HP/0j+N/7QnPPmrP4b37p7h9332VfyOT7+K526cxvycJXyoGB6XhEcZeJCALQOXbXP2jtJICa9KHzZYyIRVDRmFWm9WLAggtRTfH8T388lGXmaSctBOw399q5tfxn0s5hDnMoLbjuPq7L20qL9Z4mXOZS4X/x4+K/Vi1JlWRg1bxr7vRDSWJk1Z6j4rbRQ1HyxGU5hhDgDdKV7hettGh1lGrS6aTfzsy5opBbuzHY2HNy+0o2HASQyj+6zXCG1ZsTCLnGPASsZ+vQb1WlY/hQmkKVukWWaiHIoTH7zQ5XXoxCX/34VGtyLH9Nb9XRaLE9PmKZ97jI7LzQ0tlGUYh8gOP45jJXKdYlDdS5/Ol8lLx8vlhtxm+XLRcqQWZMuuzKbb7LZtvimoIpcXl3E+YjREc9y9R3Rj2L5zgKTDVHi/Etk6uxUJcYaC7FlJZS7mth/uCso0wCw3W6zIykR6nXK0leFeaib814jfIj6baMmYE6NNPHhwQ6ql05obB45SgU625x6RW5GrJyt2cvjdzI0dibojwM0ikzZqjiFh5nMctEZp90KaJvMeGOtJwzQkzjSYPS4a8sPtggfZ8Ki4CraUhNdefozf9ZmP4ve/9nHctB3f9ugWP/TxVzAt4ZIz5wDA037gpe0Wf//Hfw7feOMtmhCSUlwvjDH1y3Ik0rNHws3lFje3N/jAi4/xweceYUsTX/7yV/Dv/am/hJ/92pt4bttwtx8uEu0TD24rRut4dj2Clj0GA8YsA8Pwyocc1vn4NvHplx7gu19+jKdvvY9vvP4Obi+bM8uGX1ZzNjepHAPDDM/2xg7QD8taN0YJnBPbejjBpmkwOio/Xwq++ezAd3z2E171lgu1Px1WK5JlPHva8G/8238Gf+rP/w38zM9+FVupePjoAbsRF5g689K1QR946TE+/drH8cW33sHDreByqdhbx6/44PP4L77xNn7nb/81+P3/xA/gj/87/xL+hX/1n8Lf3G7xh/7qj+HzX3sdP/jJl/GPfcur+NjzD5DMTShfqBnPGfDIDA9Twg0mbHZUFou55LDySEybRJqReCdr8xEHsRszJumxBoflJ5di5b84UuSpoYmFkQ+fl8Gg9p4He7kUpLMoA/O+Z/IuVsVZH57MxwBrv7TGCIQjAbR3OWXQs4sYzv8JrzA/5Ge4ULjVSYv4ag8NM9xsm/8dPH9mRGmTwixT1Znw4MHDSDPNmfTdmbzYSxbefI3izAiTYqcgqwhFBig1sLUuNyB2PzPgxMqERSMN29XwCGZd40xJ69jlEQitVzFDSSdnwZwrDopXRMV14ZEhSbDHsJjp1J+wCdBtCCso24bWDxeRFXN/JB7EjpeuKNmcfFNfr+1kk+HQjh/WlQ9uubz6qmF2B0AMkypPZl0sERzcUIyDLrm4Xu+uQJunSwn3QmP2sYfq0m3h/aCQ/1Ob2iiTCuCMUihkLLQv4cJshLBkJ+DK1ontAVWivMB887mg0kixVN6OM7wKOjrtnVO06AM9+O+ZGPLds2voPNzXTBksiY7LMrPzOY0PRwca3YVrrqeY0oxUM2pOuE0JeXS8N4HvevkFfP+HXsQG4C/81Bfx7+8HftWv+aX40z/+M/g9D17Aa48f4Kef3uF28wv1F56+j1/14Yd4+6uv4+uvv4VXP/wC3n+2Y7QDuVSk6n5WqQ88enCDy+Uxjmc73nv3ffzEF7+KH/6JL+Bv/9gv4Cd/+osYb7yNm2fv4/d8y6v4/F3DW88aHlXDPgb2u45x7OgD6I6joNSLUyKrV2svPHeLd2fCe3c7Lkh48dYzQCbtva1kEjwcLmqHD2/73oCHXMHTu12I2cdDcvaJWjfuk4E+O3oquFjCZgn7tuFbPvUqnl13xqJekEbD9dl7eOHF5/H3/t5P4n/7v/mjeO1Dz+GtmfHBD7+Af+UP/hP43l/+Wbzz9juoNxf6GB0Y/UDJwKsffRlv/ujPAskP0leff4gXq+HLteDX/spvx89+/mfQJ/CP/CO/Fv/Qb/1V+Ct/9e/jj/x7fwF/6G/8CH7F8w/x2177ONLHP4i/9uU38NNvvYfNJgomWkp4NhMuW8Gz7kN+TJ/LGF0ojn2HoeByuUFvbh1yudygHfs9PF+zhDkMRz+YGMlETNLg6+Yzgevd1RmPQ7EDoEU97fAprlMOSaKafkzGMLTptPQ+A0KOfPnuM0VgINOh+26/otYcMbu9++ytVD9UMfz39tZXHgxpu6ZzpmxhQeR7XDOI4UaE8rWylTQqw8u7650XCJkZONWjF5Q3I21NscwZc8NxuMuFx1X3cAZINvhsDKV4rEEnu3VQB4Kxs1gvjOaY1EkhSDwurZiEvt1Ms0dWOzANKFKJj5lOMwnPrpb3vf7SFpRfZ0X03dkT2fyWyhS9WKblADehGxg6Xpem8nVbGBYCnYtoRpyu2xG0GNadoS7LPqBrsgJQzrf0DZgYNJ2FGY6d7qXsQI79yllDD/EPbK4FWdyCeb86be5y2bBfD4zVrPnFNjp9mdjBkEQgB1rF7fbR0JrCiUzkM9y9/z7t8Z1y7L7bXqnWreJ6dd8sNyh0aMlzNyYyjCwgbzOl7G1MeazFi4FsXKpzdWr+ufwAdCEjXUeRsIkllpwefKQcm8V4CX3g8QP8wMdewdEP/Okf+Ty+bobf9Ft+Bf753/Ub8at/9Xfh3/g3/zT+9P/pT+JXfeLD+NF371B56R45YXt4wYPW8NM/92V87CMvuFIeEwkdozt+XS8VP/pjX8Df+Ds/hr/5Iz+HL37+S9jfehfPz4nXnn+IHygFL33gMT77wVfxM9cr/vKX3sTztxV3dwdSMWxbxjw67p5dsX3wgTsXBOVzYN4M1Dzx9Nrxn37pbTztEz/xjXdwc3uLl55/RNggYRqJHq2jZuDZ3fvYaNvTObAcyWcNkJ15SmgY4Sg8JnBbfPB/UwvevTvQb27w0Vdf8v1khr4fMOZ4byXjx37mS/glH3oO/8df9zl88e1n+H/82Bfwv/rX/wT+wz/2BHlzhXWaEykVJBuYs+PRw1t88dowywXvHwe+9/Hz+Pzrb+Mzn/sMPv3xl/HGm28i5Yo333wLKSX8hl//Xfiv/YZfjv/y7/wU/uh/8Jfwv////h18Kg389s98FL/54x/Aj7zxHv5/33wHb75/h+rkOuzUB/UEWK4ndk+mRTpZJsNV3inB19N0CGQrW3QPNszt3dG96KqFeoQR8cHOtLLYcy5OHhQ8+kXivl0DyUbkgaScYb07pGmZRbFnC41EjF9OxjQqrJfqULFNZ8SRBeaM04rWZjg3HO1w2u0xg+Go6ITRdk99rJUOxB2DFFyfsSLEgzl7/IERHqvbuuC6ztOEiHKo2+aX8r6vtFFMp9+nidGBvBE6Z8F9yO06KzRrYqsbeltq90ldz6C/1WRGiUsl6Ew+SOaxlU6ZAJQ2B4q8/ak6lpeM9AK9d6RSgtFUSsFWL9ivV1f2Hm5b4q3cAWspgpBKSS44yitLvJND7ApRV3Eu50w6ZeYZcNkcE92mV3yl+uXVusdN9sE4yBxMiZIzaX/+Z9yWwB+0ZjRAx4TnWeRaw1Rw9o7J1DGRA1yz4dbHR6c2JC+Lj0Iq57KyYBV/oUttY184PEfd3YtzsBsEyUnJbnnQl4fpcaOjUMAlSvWEV2t+cZcgBPh8rCyVOK31j+4DuIObtbXuOHKexJULjFXnTBObJUzLkaCYOAysNeMHX30Jf+bnvohfSAk/9EO/CT/0j/5afPyjH8azfcfPf/FL+G2/+Vfij/2x/xgPJ/DSpeC9Z7tnP/eJbzy74gPJ8Pd//Iv4gX/w+/D20zeZDuiHw7SMlA3/yh/+d/HmT/4c/uFv+zh+0wcfYvvQC3j93Tv87Dvv4C++/g6+YcArr7+J3/rSQ4x90r7FsDPboCTgzXfew2vbh2HXQa8192W72xte/sBjfOwzH8F/8PZ7ePTcA9x8/Dk8+W/9drz4wiO8++7bSFYw++EmodOHi1/9yjexAdi7H/a5+DhGMOfe3JqiFjfzNIYBZQNymnh+K3j37g6vfOxlvPD8Izx9+i4sL+uYUi9IPeFHfuJLeNkMab/is89V/IOf/jD+D29c8ezZ1aOXjw4k75IM7otWLxsOXvx77/ilLz7Cn/vxn8c//I//NiTahBRFMIyGN994A3Xb8J3f+Wl83/d8B37u57+KP/ln/jL+r//RX8btj34Bv/1bPobf922fxM8/fR9/+yvfwM++9RSPi+HdMXHXgQanot/tV6eQpwQrCf0YkfveO2m0V4c8rtcrWYvjlAuUMaxHYNo8RSY7q5MXFGdyijdY7sJeqW+5hJ6kM+XTjRZLBCqF8LX3ZRzaOiFtw5QdCKFtD4DayOZAEEvGcPr2pMdcp529YLfBNL+76zXEmZOxCoDrxBoJP5fLxUkuW41kUo0Nsjnzsl13WEm4Xp35J+PHQZ8uCX1778i9LaJCNthgbHXymGazHDPYzARSdxhWiJ8P2t3yZYYH3vJNS8HSAoCSQ6A2AlM7glLafXbJl22MPwUm9SGG/WgR0yjb5jm6e1nJ4uAUYBKqdDKC2mwRA5kksJTCmoO44GFTvTkO3vjUehDaw0ZjNQ/CyZEt0Jm/bsVDoM6+X9Oyh+IozCZN2PTqM6UZMZWTxADLBpvmoilCZa5No5+XaMhpuJ29ZZgVuos4huqtIVXe0qJhBje+j5WhfTBrXsM3GdK10XB3t98LkOpwkVYBwXllsOcZNvHh408c2IfL3YOTUsclJxzDlb4leS5GG65reGYZv+aVF/HXvvYNvPjLvxV/9F/+J3H7eMP7d1d87ZtvMGBpw6uvfhDf9au+Gz/6n/0tfNuLj/C3rjsqMsZs+Mq+41tfeh4/8TNfxnH1Q0P0b6NP0M1W8YlPfRiXn/ki3nj7Dn/xy29iv6l4/oMv4Fu+65fg933Pt+J7vu0j+Bf+5X8b77x3oGbg2iZSnxgZuLvuuEz/tT5XnoFCp4524PZBxb/zb/4zSN0JE4+fe4i9dzx99w7lcos5OkZSRKYfWn/tr/8YXntwi88D2ErC3b6iXGUYOtoAMmj3fzgubgUww4cf3+Krb72NT3zPd6BWoz5qQ0o5OlYg44tf/jp+3QsPMfrAe23ip956Fx/60Efx6MEt3nrnilovHFY35oAY7u52mAHX1vDSwxvk1vDVnPFbft134/27Z7jc3jhMkSYsV0zG4jx974qn8w4vv/oC/uf/zO/GP/X7fgB/+i/8TfzJ/+d/hj/+l/5LfPcLj/FbPvtxvP/qB/EnfuLnsc2ElslaMvNApGTuIwVgZndIMCS0OdxFVuFgLCrraZA7aWyYGbA2Fd6mgbSCz3ZGx8IAuoePOZCnkZHVFJIYXlHb5YJBNqRn3dDzrw9GNLgiPU3g2ft32C4V4EHsMJU7KeO0r52gwpmkzAsJ088Od0GmYavQASuGlCp1IMvMNdG7D2nFF3uCaWIOyIj5qiNFOZIhU0I4gGfzZwI7xUtDzsR17S9bAXKWPI53MlXy4MUl55DWxorO5YjAUyA9pTCx2y7uK3PyhprLhO3/z9V7xul11Hf718ycc+66va96tyTLlivuGDdc6L0TeqgPPaEECC0htBBC6MV0MGDA3cbggrtly7ZkSVbvbdu9dz11/i9m5qz4P3nxfAKxtdq995yZ3+/7vS73Ay6WigYYZ6GJrvqfhokF6KVWOKSJQ3cikTlQTVkBVRiFZjbs4IBxauN2cyIfl/iYg8CdAAzMrX9Aah+qnoTUfHjDyLgL0izJHQiuZGTw+4Ikja1kybON5ATlq3wE5+Y9TlWpLO7Y8xQ6Ntc9gz7xEJ7pRAiLXHEeS5dec21sKTGYAkP8yxWUjnJqmrnY21RoTWdmCWZw7w5nbz6EsXVkCwG+F5AkSS7eUli9qhSg7QIxTUhjkzHX9oGdZgYSJ9DGk6ETijbg0BV4tNMUvyAY9QNm4pgZrUk7xosw2+xw8dqldPcW2b//GEEhMN0Be3Jsd9q88PJz+PTt9/Fs3+PBLKNYCgijlCPtkNXdPTy06yDT0w2UJ4xrHGOGMyDDmCsvPJ3vH5ig5+zVfGj1ElYtG2XeaC/FUolUJ5SCIktWLaXx+NP0+4rdnYSegqKTJJQrRTwN9UYbXxrnvdGfpjaFEpFa0ZJQkiTJmKmZyrqQJm0UxR2kUnQ6EQO9VR58eAuHth/g4lWLeHDvYZQ21IHQft+zWNuHiU01ptY+maSUrZ2urxiwuxNxzrJ5pEmUxz6l1EjPBEmOTUyz85n9rOspsH02ZH61zP2HZjj/0osJfFv60lnumnEa0KmpBl1ewKFMc+pgL08fnmHl2qUsHOni2NS0GU1kkGAe4NIGDqQVg7VaIfVag2LR502vejavfcmFPPjoFn71x7/yqb8+zIuXzePKBSP8YvcRAt8japvlcqHom3es1qbMJixdNxP52Nv3fDMmsd2PTqdjUDJCmpenZ7whSWypu/b33rGaojCc09Q6wZQ9sCauoCcVgozYlpqlEpZCK+0u1vw7DaRU5OY+J5wKLPbHNdNN4su4c5Bml2BSoTrfRzixU5rqvBvnmt8OlY+d2kgh80mGAxQaBLuJNBfLRaIwND0Nm4JU1lIppDnQpUmak6+NjE7lC2+dJnkfTQuNbw+HoR0lkoq8MC2EMjsee/iJo4SgGOQrgLx5bz8bc1548zJCk8fSvcA3RTgHwdOYuV6a6XzOn6WOxKtMAUVrwlZoUkhSYPiA1iPilIzasqUykx1GuIa3uZkYtLfKX0ZJbJlWDsLmEjnW44sQFosgLSvKPkitTEZYCZKQkiSKTJpMz6Wn/MCQVn3fz0tLeXlOGNd0Yjk1MOepdi8Px+xJ05REZwbtkGoC25kx4yBleV+WvJuaF6XBOGAFUUn+EtY6I0pSyDK8gkngaDkHQMOe5DL7wjAteGlmoUmcQy6VMpbD2L6QDWDR/L8wDE1TuKDs39sa4CxYzxNQ9BRVX5Km0NIZBR8uGurm5L4S9dkmvX6Z3x+vszvN2NuJWNFV5cZbH+L1L72IYjFA2nCBuznONlqcfvJiygtHCRsNxislGpG5eU2EMUJ6RDMt9h+dZPWKURrNDJ0laJvIajY7PPvi9Vx52RmUigHtsEOSadphh+nGcZT0GO6XLFw4wqN3baC7v2qR9ZLAV4RJSiAl9UY7T8L5QUAStQ1axqHMNcgkQ3oFc2ARZpQUxqGdw2d0V4p0mhGf+tIvuHjeINvqTcLUHJaEBE9bcKXIyEgRmfklTezYsRB4pNosNQNgItKsXDaPKEnw/YIlGkR292bYchecuYrrH3uG3x3fjxeGdI0M8LoXn09tdhZ0lveWhBT40hAaGrMNBgs+rXbE6mqRmw4d5Y2XvIhMG4J0ps3DSmVz9kKEiaJHcRuBR6FovtajExP4fsC5Z63kkgvX8dTjz/CuD32Nl48XKSpFM84oFgJ8ZURmPpJAaKLACM9SCdLzaEexGUvZMaxE5VhypRzrSeccOykVmTTjSIde95wZNVdgC5Am3uoenmmSmASilngWKpokWa7FFp7K6wG+b7DkmQCZWeCoULlmAGHKr6bzZOV3cUSmZa7NdcqLJE0s7imwUfUUbU//ZsRlacg2GOTc6cJOczI3ts7mPCqpVR6oghkzZ1obtEqq836Kkh5aJaTW1Ki1JaZniQ0S2eebQx3ZQI7bzaba9GV8W2FQgTAvdZsuNTIyOyGKzYhceTKnoBvWnXn+eIao6xFhZUvaKljzB9Ycvj2P7Nmrjuf5+VzSsydrbfcTnmcGxIVCwS6uI2xz38y8peHiaCMJNjP8JEV62Jmcyk/o7qBlxEjSVv2TnLGTpokdGaTmoWBLeJ7yc/ZMFpu9hBd4uWtYW3Ofse555rpsT40uCeZkWK5YNYdgAS3nwHEufuwW54avcwJGQMic6WNmjAU6nQ7YNrtAkunU5gXJKaNuDKUsgDJJE1RqzYVpQmY95Z5QpArIhG2Xm7Zo4PnmOuopfE/RiSKk9CnYwmKl4FFUkhjBWG+RC4Z6GPMV2yanufap4xQWDRPMNjmvWGRvs8PuWsi5i4ao7TjAY5v2cu7pK6m3Ouala1vDaZIQ9HTz7AvO4r6f/4lT5g/ztyMzVIQCJWmJjO4k5ent+zn5pLH80JJZBlEmBDrp0Ewls81G3jMwi1hDek6SjDWrFnOX5zPm1KtWhBPHCQVlBFJCeLa8Gtk9VJrPc/0ggMw6ipS9hQrPao5NbLU+0+KdH/0OI+2Q4eFeHtx5hILv0SYhDmPTenZR8sx8dpWS1pFjS3C+pOiZ8YXfXWHRvCE6nRitU9OGxvyCx1FItRrwP195BxPHa+w7MMmWnYc4bd1SRoa6mK03zQPEKU8tgbY222Ty6DTLyx4LSwFZO6RWqXDRs9ZQbzTy7pETNUkV0NVVAa2ZbbZQmJdlEiZIz3wWyaBebxJFGX65Qq+2e6rAx5MZyqoLejxFr+/RTjL6JNRSk9oSmP1HJzWFt4J1W7gyqjNtur2USzyaKKpJEyWWzeXi+oar5+UJS7BpqVyW5pnxrR2lpdbR4iYLvvKIw/gfdMsi3/GaiCo6I7M3SZ2mxmsUBFYNPGcq1PYk7/uufOuhtTvhm7i+wy25z1ia89JsIvMEynlGdkKs3gBgs9QQCAQC6c91yjK7bpDSeVbMCznwfeZ6E+TQSm0pIUYUp9AizSc6xklvO2o2DKHyQJOtKngGLFsIAqI4RmPWBGaqguUt2VlZnqawMVXP8+2i2csd5Y6ZYsT0sZ3ZmYdbkqQ5WkP5lt4o7RvQMy+DyJYQozA04yQbz1O2OZ3X750H2PmlM+MpcZKqwDcN2UKhmKPGlTTIeZ2ZWXdqi14ZGi8wRaUkja2G0rTRlZQngM70XBHJueEtG6dQKNgHsZcnmFy3IPADk6hQc8pHJ67xfWNGTOIYzw8QStJpt833RJ3AtpFm76HRBtZmP9xmFGdSK3NOdzMHdqgSk2mf6yZIq4TNTqhJOjWmiwUWfPOLECjBG5aO8YZFY2w/OsNXtx7g8d4+3vbxN3DDz7/AqnNO5anjk6zqqdJKUqaylFO6SvzxtkcplkpGLiQknheglCQoFpltNHneJWdwRPmMS0mpYGLOaZIxlaWMF3w2P7MPT3lmdIG05FHrg89S4jg0LW4vIIo7doZvugitTsyC8X5anocvFOWCb27DSUaxUqAkoN7sGJIqIm/qG6yDwLNR6UJQsMjyzAYgnMIVyqUCf73/CR6+fzNvOn0ZDx2eokf5pDY4Ygpi0lICUjtW1HMsLaCgPHxrIGy2QvpH+xka7CK2CUJjcPQtDdmULmdmapRKkvWnLOLNr7mYk1aM0mzVEWQkUWiLveZWERR8jk/V6RyvITzB6sEutk1Pc9bZJzM81EMnCu0C2YzcpFdAeR4/uPZmbr79ISpekb7eLhMySTVpFNmHtCTJoLta5ua/PEBXmtKWHlGSEghJpFPOGOzmjQuHeNGCPq5aPEgUJ3QFHoPlAiWrEPY8ZcMbkRnR2MBHEps9g/s8OiOftIGNKI7wfIW0Lw5tS7FxFNmfobR0YGkldeYBm9hDnmE9eQR+kMf3syyz4ygfdYJYzQFV7bjC7AlscsQveLlbw+wjPKuMFXmpUwqr3HVgJS1yXW0URuZA7lweFshqCB2+Ge9bz3sURTYkYBSzXhAYqbewvh/L7zLPiyTHEZmfXUIURzlLLEtT4jhBespEfZkzRzqyum8RVJ12lBeltdaEYQedJXZhbjouSWa+PmUP/p7vm5+d0y0KO9bwvcCyqjKLW47yk7CyjmAjibEObOac4OaX3S6bMruvsD+INM1sPNguwhODgzfXJGfpw1BelcptiO6l4FJK0r6MXN46tlep3K5lS0N+wbeJDFN6dKwYd1NI0tg0k63yslQqmHGbU1KmzpQo8pdRHMe55jLNTNtYZ1kOcXPX0dR+/9LExHwNcl5ZU1mEEBZJjzXz4dzyodWQqhMEUHMlQCHsySxJDU5FeTnoEftQdHNc9/Jx12QTwZP21GZ2JDKDqhS8Yck8Nhyd5N8370KdsZpvfP0D/OrbH+WFV51NO25zzcVnsDvVLJACT8Lm6SZnzOvnkQc3cejIJMVCwdwwHS4dw0lbsXSU1Wes5mBtluXVMnV7+tw9W2ewUGD79gMkic3+JzbynEb2lx3b7NbEYdvSSK3vXvn4nmDlgiG6+6uknQ5Cm2atJyX1VoiHoFFvkFoxmZSeTeeZUZ0SHs1WxMEj08ShpRRLZcuyBmRXb7V5wZXnsvZZq7l75yHOHe2lgybAlLCSLCXN7AvHvtBd4mbu5inoRAmD1SLTScSSpfPwLTlR+QXSTJNqY7okTU2/QnlkmaTVjjk+3aATJubnbIVobuSRJQmlUoGNm/dQajVJSwWGPZ/HJhu8+MpzrNiogFAeWkuiTkzB8zh8ZJJffO96/vPf/o8XvP6T/ODaG4lbEcODfVS6u3Klqe/7tFtt7rzrUdYM9vN0o2n3dAmjpYAr+/v48f4jfH/bPmbrDd5z2lJW9FUhy6jqjIpntANSawqBbxNQpttgiN0OvWG4W6434crCiW1KFwoF043wFJ5v94vWYZHmOHXzXEmTxAAktfk6XRBI2AU79kAqhMjHwghtk6U+nvTsSJj8dw2rz3X7B8+GAhzGRFiQqNDm+ZVp15rXFMvFHGqqpPncSCAoFXLPjvKUWV7bPZWSnunPCdsgt/+T6Tk5nzFGOs0EObLJtwVCYV0pURRZuKLO6wNZalzzYWTgoW68H9sbiTFpmoOy/P/VK1zi1mwnNFJYPWyapAYIlyYWVuiDfSCjzWIF13HA7DQS64TwA9+aA7GcK3MawF1Vs8ye/HVuyDKGMjnH4HfEXPvwE/ZD4ex74IpwGJ2ubeV6yqPdbpNpQ8vNbKErjmPTY7HEWVPyYg7LrgxyxIWgElsucq5ll/RwMi5XWsyyzBb/VI7Bz2z57oSaiMn0FwrmpGJPEML5n60n2bnIpRL5/sjE8bIcPCeVWb5pzKgttTws0x8xoQNtcSeZ1V+6KKO25Ewljf9EKM8+gFKkNsWhFy0Y4cajR4jWLuHX3/sY3/jsm1h/yiJmatNMHp+h0WyxfvUCepeN40chi/qqHGhHyFKBwbDDbXdtpKu7apwFcdsgHVKD8QiTiJc9/0KeqLU4patEkhqs+dEwpFL0mT1a4/jkLEJ6eBY5ktqHsFLGoSGQFIoefZUKC8dH6atWicKEjZv28Z2f3cFTEzWCgs9JvSViz8zMYy0oKcXUzGzOUzL5+sRGRSVBqcR///d1vPjFH+PTX7jWLEQTg/DRSGNOTAWV7iqf+/DruOHgDEsqAcPVAr4CTxrGUmIj2o7tFgQqj8RLi09JBfQGBQ41Yk5ascjcclOIO22ENPN0rDvdvMgEyg+QSuEHdj/jlczo1ffNC84RrRO4+Y6HWFz1qXZXaczWKY+NsH7NEuqzDZO4MUsGhFJUygG33fkoIg75wsUnc3nR47ff/QMvfN2n+dR/XMueHQfoKhWpdpWpVos8sWkntf3HGKgW2FlvUQ08Qp3xgvlj/Hz7Lp7/yiv5yfc/xV1hyrcf2soVCwZ4zYoRip4iQDBQ8Ch7c4QGz/dsk1qY8U8SUwgCs/dQc3FU8+DUFlGe2MNsZk2AGs+TeT8sjkwLHGkPwpYrl0Zpjtl3I2BnUkxtBNbphqM4sgqJuRSXqTRYvEqWIaVHHJkJTJIaJbcwJ2crRnToE5mrFQRzXg5tD6xxHNvnlIvHpvkYTpyoobDPlzCKclmcA7+6iQlWGuZwS27XI+zYyvNMhPtEWZ/yPLzAww880sw8S5VQ+MqzzyLzdYVxZL/3c/Ffz3qQEhtdVl2Vrs+4N522kTPsgtCNYIR06ZIkT2rlOkzLrHJvNqW8/IopbYZeKWXTRPZdKuZmzr5V3OYlPQtc05ZE64BnZjFm/7KW1Jm5IpyaM2o5IOFcOsEgqI161/zznjpBveqQ8nZhZpbhIs9EG02ll+OPle0smJmhKYgJy8bKUzUWPZBYjo1zPZs4ncpDAnNxO52nHNz4zkAuVY6BkLY05CKRIi9TzYl69Akx4MBTFHwfYfEugT/3z3meh6fAl/Cs/l7uPz7Jf37qHaxcNY/9B4/mJF/fD0iSiO7uCtPthI0PPsmK/h621loUPI/VpQK37DjEK59/vm3rexZ2aW5DSZayfMF8/njP48wLQyY11NsxCZpV3V3sm26w/vyTWThvkFargxLC7LLSCM9+7jQZBw7N8NBjz/Dbmx/gu7+8nR///FYeuP1hatv2cPFID+eM9bHMC1hdLTFpT4yFTkpzqJfnXXo6rU4HqYUF2kmEVHRXKtx856MsbDbYtv8YC5fPZ8WSYdqd0N5AzUOn3UlYvWox+47U2PD4Vp67YIAHjzcQGYR5H8FEJpWvSBPTrHajzJInKXiCtb1VttZavPxVlzM02E0Yx7kYKQnNi83zPevkAJ2ZNKEZh7liamyW6JlJBg4ND/D3+7fx/e/9npcsH8evBDy04xBrL3sW11xyBi17+nTji8xGZGdrbR7avJfrntxNRUletHohZw32seGRp/nRH+/i3se20RV4rF66iK995zq6j82gykWemqrjeR7L+iqsCALuzRK+/ul/ZmB8kJddcz47js3wP7c8xEjB57XrFnOk1uB4OzLPAWVw4CbumlhSrTxB0zpHmchd65ayLYQgs41umRcLhf0d1hZT7llpmhVBaTNCN0h38vJsak2KwvmDbDTYkRpc7weLa9IazCrDm1PB2t8zgx6xnvU0zV+Avuflv9cnqmfNQUOZ9Fc850HS9vDpXnCZ9THp1Oxs3V5DCPNnu4SUI0u4kkpQKJjJg7bfF2c+VHMhAmepjKMop3K472eaJf9gOTRpU3IuHhICzzcrCVtKVqVK6TNBEOQtdOMntmU1u+twitY4MugFZX+5pSC32bmRku8H9u9rZ+9WdpKmiRG6ZIktKJ6Y+z7hg8ScpcuMyJiL63mexZeL/JudJEZs4yQpQeDlO4BcxmLNX0rKXKjkBda0qGR+SklsLtwtkLTDydsCjul8pPl/l9p2s2DuA4oQdidhI3/uFJLqPGniAgPmuSPzeap7OasT+Fuu3In9vriXThKn1kVhkClh1Mn3Hr4UeEITCImHoOBJCkKghBEuSWFSWlrAimqFer1NreLxrFOXWdmWk7olFuugGRkY4De3Pciz+yo83eiwbbLBRUsGeXzXQZatXsqyJSO0Wi200KY9nGmSDPp6Kkw2Iu7/++OsHe7j6ZkGBV8yXinRbrbpXTaPs9YupRl2cqqtWZoLZmc7vOO93+Dan93MxvueRO7cz5lZyhX9XZw1XKXoKXYebXDLwUn+MlVnc73Jc7qqtDwop5oDhQIvvuIswo4pvErfQ8mALEno6+ni749tpXdyknmlAjdt2cOrXnAOUZzNLUrtzDpJM85bv5r/+9PdrCtIEgHHWhFaCMLE3mrsgUNax3WamJtFxZeUPcXK3hJ7M8HrX305KcYOqW38nSy1P3OoVop0dVXwCwWz77MGyzjsmPKo5WX1dHVz7ECN137067x+fJiWpykKwZ37p3n3e17FQF/Zpv60/Tl6+EFAFCWsWLGAV7zwQlYun8+GQ5P89vEd7JuY5tLl83nuvBGaB47z21vu5xc3/p3t2/bwqmWL+PPhYxS8gOmww2uXzeePz+zmJW98EeeesZLJ6WmEJ7ny8rM5+aSFXHvPU9zy5C6ev3iYbt9nV9MUIDtxCras58ZUmU4N+l6afWPY6eCfUApESHtKtqY8yDUMzseuTyB8S2VisM6eKqQki1N8C2Q0IyPfjnfTOU2t7Ys5La62hzhHKHdpUecdMXuHzHAvMjPK92xEPLHQS/cySLPEvqakDdeoOadQqvNAhEuAer6flxejKLSLbhND1/aBnlmKr7KmR1c6duyrXIurPLtn9u0hFIuG982USJgUm+vG6dS8tLU1SxrNhXmGOUWztNilJE5Q3V1dn3FeCk9JG7/Nckd6kqbmlyJK8INCPpMPw3jOuTyH08pvI9grZ5bNyegd8kBYJ7e21wZnGXPN75zPZflUuQZXG6Cg+YubE4jv3ODW+uVSDWiRi5zylIw9VeTFSCunMlx8fUI70yQP3MI6S7VNLFgGjn2huCiuixp71p/h5t+eb3oFsXUDmKIilq1v56f2YZ2mJxR33CLPnZLyjo7O9bsGJ53Y/o2ZXZY8j7KS+EJQUJKybwICqf1+xcK8oAt2j1QtBqQ6ZW2lzH1HJ3nRlc+yKQz7z0hBFhso3sL5I9yzcSfB8QlG+ns4pavEOcNd3LVtP9trbV56zbm0Oqbklv/shCTNUhaNDfGjG+7m1EqRPdaCV1QePUnGbE+Vyy88lVbLmgCzOVe2lJIf/fQ23riwnyvmDdOKUjbWmtw+McN9jZiZnm4WnbOWq593IR9+78s53GjzwIanuWhsgOP1NodLRV5y9bkmRitV3qbXQLlY4NHHt7Phwc2cu2CYGzZsZ9HSeaxZtYAwSi0xVoE2nZmh4V4q5RLfvfF+XrxyHo8enUELQUe7SK9FcNuAgrl1ajwtGOku060zWv29vO4Vz6EdhwZJk9nRrYCenh6u+929fOKLv2B6tgVIeiplenvKVIpFqqUK1UqVSqkCmeTGWx/hLf/2Ta6olFnRU+be6RqVEI4P9vGhd72ceqNuyalebrtMMyA1C/9Ma1avms9LrjmPyy88g7gY8KfNO/nrjr2cNNTHm9auoL8Tcd5gH5taLbbXO6Bg1XAPQ2nGQxr+6+NvoxW17ZjalIeXLB3jlS++iEjDDx98isuGejkeJkxlUAo8pCutKpPaM25zO4NPYgI/sLjDbA7nbg9G7nNv0kKZgR8K8zB19sTULo5dz0IpifKVOXEL8wKOoxihRH7zIB9Tm5eEKzqafasxl7pDpduNCqtAiC2KJLMCNm0LwErN3Yjc767GlJ0dtWKu/6ZN89sGgaRlEUopDIHCTlbcv9MRxqVyoNa5m0Wh4FsXk87L4M6jhB3he9anom3Pw7dTDTdZiaKIIAjmpFN2J3yCk9ZsboXEQxvgnGPQu1NsmqSkOrXuYoVfNNehLNFoNfew1plBNruRjftKTV7ZXgel+SEYmx3WWKZtXCy1p3kbH7PubZcmyk6IoGWZsem55qmL8hUK5n/XGA6P8uy+wuFYpHlTp7H57w1Wxb6hA7uTSW3vRWf5bNssdTPLoyJ/CWRpZmez6ZyEJc3Q0nK8/Lk4r0lWmb9XUDB8K3GCmF5aPo5UFhUgleX1/+OLw4DMnCFSWbClAcgppShIiQ8UhYUfKg90Qm8xYGyoysKuMl3K5/HJaR6fnKXLLyIzzZ52xAVj3Uzv2c+Tm/dy6rqlNJrN/OUuPc9ACdOUqy85kx999WneefIgv3liD+/dcYgFq5fzqmvOpdmJrPTI/MKjzS2n1W6xeMEg51x4Btvv38CyaoUttRYtCcu7i2zae4ROxzT2BZJUpPmttr+3yjlnr+Lrf/o7Y0vHGF8wwrpVC3n92iWsW7WQ4ZE+CgWfRGcUywWef8lZ3PWnu1EafClJ2h2SMDJ4GgRKW4ObkEgyyqWArBOxq9Hi4tE+vvOLv3D5xacbt0uS4nn2pqtgcnqKV77ofK674R4ePniMi+b1cfuBKQq+TxSnOTXakA5U3rdJhaYqMo7XE+atG8EPAuuf1jnfzRgkPTY8uZPKniM8+pu/8sef3ErQXWF04RDDo0PMGx1AeYqDR2d48sltTO0/xFsWLmC0u8Q3t+3ixasW8cC2gzzndVcjRWqTQMqcWO1YQgptmFZJTMEPqDeaCOmxZOU4/3rKa3n3m1/MAw9s5F++8VPKByW9XVV+s+8IkYDuoMjRTpMLu3v52RPP8KYPvo6uasDEVNM04zMzqmm1YoTUfObf3sKWfYfYt/Mw3eUC3mwTqaGYaRQJwwWP2QQacUYiUrTQFKzVNNOmIxNHoe2kZXmMX1uoqQO95uQKx8qyzxDI7FhL272ml6uapTIHQOlZbIodn2d2+hJb/LwZLUmkNGVA7OlfWThoqlPzgrC7GIMoEiRJZpNhbtep80a4Y/tldu+rM/0PaUsTdlEWcZLgewHtTtsiR5K8fe/wOc6+6ND6cZydMDLP7FjVvogzp88wLwIlZV7gTuz6we11kii2/pTYJkDtzU86vJRVeGj3TrGxriQ2cS13YkdaDWQaE0cxhaCQA/nMrsG4OJLYYQGs5ardoRAEZlYXmz5IHEf5HsTl5LW9/UhlNJ7mKmny34FNKpkXTWZmpvkvt0kKmCVzbHDqSKQPnXaYR5C1XXB3mh1838vHdAZuZmOz9mbiFmegrf3PnESNetazEeU5P0B+TbXMnsTKWhwo0lMeUnhkFlVgRinCfpDmXixoYy/U9urp4sDKxguFMid5IaTpuVhaqNMPl5XC0xkl36fkKTokjFR8zhgcZrwQcHCmwa6jU9SThCvHh+kOfB44Pktv0We2E7MnjTm5VODGux7jnLNXM1OrEZTKNkqtkKlmpt7i4nPW8BUC3nXPFi57zml886UXcPYZa0EJGvWW4RLFKWnUoViskqUddAbNToeXXX0+H73zQS7rC9haazPdjuga6Gby0CQTtTbVkhXzBEWSuIPOMtqdNu9+xwu4+qpzWLtsHqODvUjfo9luk6YJjTBmth2DTvFbHZYtHqfU081kq0NvpUjabBkrobWw6czI0jqtDkkSERQ9Sp7kQL3NqoFutuw5xA23P8JLn38BR49NkabC3kQkWRQRJxFf+Mhreee7v8SVxTKeNI19z5eW9WTEZmnmuj6GJD3eXWbvkSlOXTiSByGCwNn7zOek0wyZPF7j1WsX0SoVabZCGu2Y43uPM/H0PvZps5vrLgRc1t/F2LqVPDXT5Prt+wmURzkTHBCaK59zOq1Gw+wLlMJDkHkeSRja3z1NMSiSJIYBJ6VH2EmpTU9Srha45srz+f6Pb6Bfw01HJpkNDaFgRkecNtrPZL1JMtrLq553HjP1uoE5ZmnOOpPKo+grbr/zIXZt2MpZ48PcM9FAZJpzhvo4v6/C7labuw9MUNEa3/eIFNTDhMyhPZSBD2Z2p4g1fJqHvsxRI2EY4Xt+7gQ3SJjUyNxElrtGTPQ1meuPCNBpRic0pT1zlpR5nNqoeRUxMb5QxLYVrq3hMOxEhuVnaZkuHer71l2UJMRE+Q3HBZHygIzdV2ot8pt2kqWI1HxmTIjAYJfixNCM3ffB983tJomi/FYihSQTGpEZLbf07N8hMLtb3/eIOhHCcwRga49VMmeEGdSJIagKJfIYcz5NsDdALKvM84yL3hNCkaXGBBbHSb6INctZW4SzD3zt25SPjcClqUbr2KapbDnLNxjwQqFg/7Mkx2k4EVOcpHg2PeR8yaYhqdGpxXRYuq506QmlwF5wPE+RRAleQblSN1qTG+tMGsJZ+nTecNdkRJE5KUm7CzEKTGlTZA6+6CHt/kbbZJPSnt15SPuhE/b2oPOmZpZmuVHMYA1iu/zEnJTt15NpA1YzL2ljRDRoApn/YJWn7FLP3DXNct4WPO2foYAu34Aji/YaXfQ9Xjh/Ab1C87e9h/ntdJ2e+UOcc8F6zlsywrd/9GfeM2+cp4MWjU6CJ+Dx6RrPHRzg9/dtZHLieQZbHxsUP3bWmkYdKpUy3/jyOykGklNWLyLSmpnZBlmamvl6bFwFHh5J0jFsHS+g0Wpz5inLGFq+BFWbZElvkR0zIQSKUhSxa99Rzl63iHqriY6N6TIFwrDDUH+V+eP9RFHGdKtNHIZ4QWDBmRnKC9DaICnGxwdYuGIBx3bvZbRSJepE1JsderzinBLVqkEVGX3VgCTTlKXi8UaHVy8Y4qvf+xOXnr2OxQvGabdaBpCcpbQbHWoTNXwhoKdKK80oBpJ6ZCReTv3q8D+Fgk+gFDpNKGlNPU1Ys2ohcRojSY2v3aZdgkLA0ak67alpot4S1+86QDtKKRV8Kp6ip7+HYd+0rFtpxmOtiEPTh2gnGb5ULB3sZv90nUWrl7Nq6RgztRkyKZFZQiZlHo/t6ioRtWOSKCLWKYHyrOWuSKEUUC0Wuf72+5jZe4j2SYs4PjVLwfeIdEakY87p6eUnT27jbR9+PdVqwORUyx6+NJmW6CwBFOVSkW/9/FbOCHyeaHVoRSFr+3s4vbebT294kpcuncebTlnKrbsPs3emQbcwGJeG1lDyCRPLoPOMhsHzHag1JYsza840/70p5wq7J7XtaN+UQXPOlFXCOj4emSvW+lbzbK2J9vPr0k0SSWzBpzo17W3lS0RmIuUOF4LlBxplt/ndNT0TAyXEduewcEJtZke2yJzlIRqE00UnSGGqAp5dkhtVg2WtZTGZNlgko2qYc6YYyV3qer2W0GD1BDb1JSQoW8hGevhK5rFjoxDWpHav7VmFMZ4kDA1i3xy2LbQUG/VMnenKNXXt/FBrZdWVmRXfJGRICx4zuWcs4t11PNwXK2z8DUvNdSkTl4v2fN/sSERioYMGkKiFeUM7VLuyvgDlKWLn77UJg9Rm9rFIj5w6Ks2LSknzZxuLlqBQ9PKrn7YBcekAZSbtayKpYYjn+yZ6GgTmZWwRxlqaRIibs3q+W2SB5we5rQyl7ChB25OEst0Aw8/R7oXtechA2tODSVFFHUNz9X3DYzJguMRGco3kqSTBE5quwKeZxpw1MsClw0PctmMff52c4sKLT+PrL34Op5+ynFK1SE93H7VWzC2/vJlLF47zqz1HGfR8jjQ6ZKM+lQNt7nn4aa65/GxmG508FWeWgJpWu8XJ6xaRpTGHJyYJiiWUdTGkicXYKA+d2hKg8gD3d9S88sUX87sv/4TzF4/ydC2knqX0ZpotOw9y4ZnLabTM9zjWkSUFSMIoJjKNV5TnQxDYRBGmL5GG+EGJJDa/RGtWL2Lf7n3M7+mic3gKnQhGh/poNZogMqI4pt5uM5PG3L1hOz2+B0JyuBXR6Spziejwyo/8D2etXkKj0SaMEmZmGtTrDcIw4sjMLM8e7MP3PGaSFM+gFeyI1YAr86i7jYgWhKQBHD5yjIJczlBflxlHJhBZTM3BY7PIdodOT4k0lUgNs7HmUKtDrFsUAo/YusF9z0Mh8aXHbByytFLkwZ0HeMkrn4uwPhBfzaE0hBIUVMCmTXuYPz7AyMgAs/UWmXaiswSdQblU5La7N7CmWuXJepskSvEDQStLObW3i0PTdYLFY7ziqnPtwSEjwyDQpSygpUdvTxf3PriJg489wyUrxvnR7imKSnJBfw//t+kZXv1PL+ORhzey4ZFneOfZa9g+W+f2nQfp9RVFAZOxJhOaVpLiBdKCCDnh98mytoTxF7kdbGob58I22H3fz4MJZoxtdhi+r4hj8zkoFk1EOghMVF5YI6JLKbkehHuo5qPpzBxCHFAx8HybFBM5fDWJEwpFo4GwFC+beBKWFqyIrXfelbId5cMt7MnmSM05YNbeCAoFC47N3SQyF8El9gblwkRJHJsDs02wJbY4rpSHSBOiJDaUX7+ASNIcf4925WeVp90EhkKSCRNl9nJdkdUsJkmG8l3pxvwfzyUKzC8twlzjSqWiRYk42JYi6kR4BR8llFki+n5eJozjmCDwbSxRW9yG5cs4V7enbE9D5iUhd+XSmJeTIWM6Z6/5AJnSlIeyy/PYxtek43KlmS26aSvL8qxWV5gl6AmAsDBM8hGVgT4m1j5oP0DC7B8Me2tuAZfZH36mMzzpQzLnB8myxMAkPdO8NhY0nV9h4yTO8S1SYNEKmjQxCzZtZ5RpEuMLSY9v4oWeljTiiJctW8yAlHz8vkdZuv4kfvYf7+K0U5fQiULaYURjssVsvcmrrrmIP1x3O+enCV2+R6OTIJXH1pk667q6uOH2h3nJ1ReSpS1zm9SxSVXZr6fVaiN0RrFYsj+z2Lw0yMgSiU6NiEtaUKNLjdRqNa656BR++rMhip0O/YHPwU7MaMlj285DZEl6wkNA2ZCFcbMILcl0gmIuqk2m0aQgFEpCsVygu7vKgtFengxDBjzJaFeRD3/1l6xZNMrRY9M0mm2ajQ6dRpso7jAaxVwz3s/vGzGjBZ9bp+u8bLifsXaH/Y88RVV6lJKM8aJPoARBV5Fqd5npMOPOqRpKQGjJxkIppAs5KEkaZ2SeR+ArDjQjXrpghK997df84Fd3sG7tUk5dNZ/T16xgeKyXJfNG+e3WR+jWGTUUtVSTAHGaYWomkjjRoM3pOk41MSlFJemtlBiQimnP45Jz1jI7Wz8hmKHRWULgBTQaEW997zdQZZ+3vuoKXnr1OQyPDdGJEzN+9CT79h9ly2PbeOFILz8/NImnwVdGEXzxyAhfeWgj//b59+AXFY3puklEWgxKahM9Xpry7Z/dyOXDPWxqRiQi5dnjwzxzeIL+U5by7x95I4ePXcPXvvlLPnHdnbzl1OW8/9y1/H7zLvbWW5QyQWJ/npEN6qCtSttNPqwFMoojhLAjRgwvSthDYZLGdq94girXJvzMqCqxfh1QnivhefmJ3nlq8sW521tK7BLfkAB8GwByZWxp3T3aTjXMjSCzyTEzcRB212EXcjn9Ajvy8qRHqlN7G5IIm8QUlvic2iKxtK4eITDJLY2J4AdevmdzgaDMEhhSWyUwh37jNXHUhLgTonzPCKjs9zuTc0BF5xNy8WQpBapULn8GnA89wy8YqGChULAcJp0nnrQ2XJjUnprS1MwbpVA5NNB88Wlu2nL+ijSPthk8+ImSqNyva9+QrmOSZfbaaLP22pbt5rSLIt+L5P6MTOf5cNcfiW2JL9OZaXkKkQt8HEo6KAREYWTw4jZMkGmNZxMsQirjd7cxP+nkRPbNnMSxnX2bEYzBiGT2ZGrLYjYOaFIQmR3NmReUiz4LyxYzJNI5QrJpqkdI4VEOFIEUBEAiM9540jLa7ZDPP76ZD7zn5Xz502+mf6ibmUaLsN2xpygT0xsfG+SxLfs4vnUPw31d7O9ElKVgqhNx8bwh7t15gAsuOI3u7hKdTsuaEo1gTEhzS8syYfD/0hAHlOdTLAR0d5cRShGFUb7gl5YqnOmM4aFejsx0ePKRTawb6uaZdsgpxYCNjZAXXnEWaRpZf7PKOT+mkS5sSzwlS2PK5SI9XRWDL0Fw+Og0G57cyS+uu5Of3Xo/L+vvoZ6mXNhXoX18gomteykcn2Gg2WZJEnNqUXF2OWBtfy+3tWJqsRkrkgkea7TQUlEoFukoyWwK00nG8TDhQCfk8VqLR+pN4tQcqCw205Bg44TA9+wuwNzGi0JwJAwpSslVo710Nzrsf+YAd9+3id/f8iA33vEof7r5fu5+8EleMtTP47VZjoUxnQw6aYqWkmaYkACpMIm4ODHR3kRnnDzUA60mpdWLec1Ln81svYnwbAnWJikHerv5812bmXhgA69dPsK1tz7Mz2++n8njM4z0VxkcqDIy0sdPf3cXh+97iqHBbjbV2wQI6nHIOcODzE7PMjtvkM++/1XUmw2DKUpNyossI9PQ3d3FbXc/xu2/up1LFo1y477jVAOPy8eG+fn2vXzxs2+ntycgyTKuvOJc1p6yjB/85WGe3nGIV6xbRtHz2V6rm4W0trcz4YjSKdImNlPbhZHCvFiEjbsrO9IxB0xnFjT7Em3jq2makCSxnYY4IVRyQrlZ5RMY4Qq+lgclhO2x4HYPjhpubjCGAK7yCoJLXiUWmWK6L+RlQWmJ5aYJLsw42yKkDMVYzUFlLW/P/ZkaJ71KrQ/ElI7TLJkDqtqDaZJlJgwlBEr5tpdiIsvKM7cl7eoTmUPfW75XRv6ccxRjV/kQQqJ6ero+o7XZf/gFz9jXfJ8wCnPhujghP6zdg84Jo3zPxMGEyHPGDgdsltKOemsY9I5l74iwnu/ZFJIVLIk5VavylelRuH+3hjCMcjSHzjAZcXe3s6kG5SkT21RzS6zMCmuEbZsmqeHtuwixaUCnOaVWWt2rgHzRro0/N79OCsvdShJHD51LZJiYpvnzXJHHsJKkJYtmORrZdU+ksKwrzEPbAB6ZcxwjKHqKkqcoKUUrTfjn1SvYeXSCaw8c5Of/8yFe8eKLmJqdpdXuoCTm+yA9pDBI8cD3KAclrr/lHk4fHWDjZINKIaCdxSzvqdJotKkXfS48ezXtjrFKakwZMY06duRnXfdC0t3TS09XhaNHZ/jBz29nYqLG+pOX0GiYJJe03ychzThq6bwRfnDTg1xcDngyyTi9VODhYzNc/Oz19PeV6YRR7mHJMrPs9HyPUiGgt6cLX0oOH5ninvs38/M/3MP//PBmfvmrO3jkb49Q23mAF/b0UpGS3xyZYH8zYrBcpFAs0BbQDDwOZ4Jt7ZgtMfx9tkMzNrPpVmx2Ub4nOZbG7JhtcaAZcjyOONqJmE4yJsKEWEo8YZe5gPCExWgb2JzOXELRIruBQAp2N9s83YzJpGJFbxfPnj/K6V0lVuiUarPB1SMDHItTHp6qkSmfSEqzE7Qhk2LRcOViG/4o+GYnd15/F9smp7nyFZeybvUSOklskobStA6yNKNYDLjjbxu56c6H6asUeOWaJazvKnPXw5v5zu/v4qmte1g93MN3f3IrZxV8tiYJtU5MpRSQoblicIifbNvBZz/5ZhYvGqLT7hiIpPlwmsOJlJSU5L2f/A5XV4o8E6bsbjS5etlCNu89TN+Zq3j3665kYqpmiMdRzPIVC3jRCy5i90yLb9/6AOcvGKKn4LF1qo4W0poebddL2pKlY9BJV0oW+e+8K8olcZKDC93oxbDizIRDeZI00RgYuCMbS7DcOhPztXLFzHrYnYPDM2EalxR04iwXrVVSWXq2yrtyLqmZJObrz2x606WuoijO0ezO1WTQTGmOnT+RjutUF8ZBn+Udsbnumk1KCZ2vALRNlrpyYGY7aba8Zv7+9vN24sHfJFlFHiEWTgegzJjNc6clab/f5iSWnnAbkHk3APsGTDIjI1KucW2tfqn1qbuXgouMuXp8ptMcm+4icc7C5/j/xlhoZuBRZL6ooBBYvg25Z9yeO0iSlHIQEKdmjusc5wCZxa+Y5IJ76ZhkmXQ03iSZu95KZXwiSUamDVM/TVMzQkHnTgC3tDJLd9dATnEf2UynuU8kiWKEp5AauydJSLSBr3menyfahC32xTZhBnM0YmetK3segYCKFIRZzJtOWsruiSl+Nz3NzT/+FMtPWsCRiWk8pcyM1C4RU3v9FBomJ6e54OyT6Fs4RqvRorfoE6UZJPDUTJ0zBnq47c6HedurLkFKDbbElWnju9fafNC8oEClUmDz0/u49vp7efi+x5nafZjlp6zgykvWGay2dCBKo0BttzosntfP+medzFP3b2RVf5WjYUyx0WbnnsOMD68wOOnAp1AwqBmlJdMzTR55egf3P7aVjU88w9F9RymFEUvKRS6qluifP0CGJlGSZ2ZDHqs1KArF7jDhmXaTxOkC7I3U8yQqFSit0RJDxhWCTmwQLFlsFMdJkuX9HPezFbaPpBHESWZvGk4GBH7BeG48T+HbOGWkMcv3LGN3M2RLvU1xska3p+jzFeOlIncen2FXo00mAuppRmobxb6UZm5vcfiZhesJDaVA0iMFU8rn7FNW0my1zXjFJZe02RvVmi3e+rpLWLFkhOtv/juffXwbKwo+1ywa5nXlgFu37uPlb/8qfaWAeUvmcfPhCUpSUo86nDc8zIaDh1l3zjouPf9Upmo18zumjZdHaBPq6O0p85s/3kX58CT9Jy3kN5v3srS/zJJCiV/W6vz4zS+iFSU2mQZSaGqzhrT8uU+/ldPXLuSLX/wh71i3ir8eqRlEjH2amm5URqFUyEtzwn5/JAJlx0SO2eRoEdhJhKP5mgOldWHoDDJBqg0NI0nMsjkKw/wgmKWux5HZfVua97viODvB0idzMOscIcM8IwrFgjGUCm0L2GY34IyHxmwpbAnVQ5PYjpoky5+Lc8VrZZEi0kZwhfAtFsUkzZz/yL0UMm1G+VGcmBdnluIpP9dOpEmWlyOFjQq7G8qJnRO3XnBMryRLzf7XNDa1Sd6kCZ5TzOo5e51zPSilrLs7y529YRjatJIVOFl+S5IkFIvFfPnuSUWUxnkiJEtMOzQKoxxt4pANZkZnsvoObJgkqWlLaid+x6ZdCjm7xXPyJCulNwkM8xbNdJozi7DZbHfiz6weUli3gJs9OsSAeTlosjSxD36ZN2PTJEHHOscrJ4lBdHsWEaA8ZU4YylAxTZbbkIHR5hcQJXLXuXs5YRd00p5UNClFIej2zPfx4kVjJFGHaw8e5pZrP82SZSNMTEwbwmfgW0WwXdDYbku1u4RKNFu27uZgM+TMapGSFzHdTikXCjxTb/LswV5aT+/igQ1bOO/sVczONgnKFUSmTV8IC2RMUj75uZ9ww00Pc1a1wEtH+1m6fIzPPfQMjz+xi/WnLqMTmlGnVF5+MuokIa+88ll89I4HeelQF5MCRnzYuHknr3jhBQSTNTqNNrt2HOLRzbt58NFt7Nx+ADHTZFHZ59TuLkYWjqOl4FC7w9OzTY41a0y3IhJMCq0oPTqZSQb5PhAZc18jTM34x44SSwXzoqqUS5SrPqWCT7GgqFQKFHyF5wsCz4zTCsXAtKARdNoRYappNmNmZlq0w4h2O6YdJjQabaLQnCqFsKM/IPAV1YJPteDhaTMuDpXkQCdh62zbEJ0RZApie/uSSBIbdTc7wcS1UImSlMVdZaYbTQYWjjJ/bJDZ1qxZjDuSbFAkjkN0nKKk5prLz+KFV57Ltu37+e0N9/KtvzyIPz3D1fMGOHNsmFP6evj94eNoCUXlkaWwslzi68/s4YefexdRHJrfAelu45CmkTFuAr++/u/4rZCj9TY6S7lofJSbt+3ivOecyfqTl3J0YgKhM/xCwUwPYnPAqddnue/JHZze08XRVoeO01AnqXnoegZrJOwtPUwSdJJSKpWI4ogsjgy51oJNg4JnKOpZlpeS9Ql7XiEFSijb+lbWxGcoDcZt7vxC1vWCJmy3TRDHxmYty8HUDQpFuycwUWTP1gUKhQI6NepoqYT9V+kTxkNm/JTGmd0XZvmLwil3U7u0d67yNHU6gsQmQiXKl/m4P7C6alMCTigUCkRxnE89wIi2TI9NEwRFWwkwtwelrGVSStIozQWDribhe5IwNgEfIcDLLD3SGP6UNXxl+UnQ0R8dTsR6Js2oRpgXjcRcrZyjWwhz7YrCyOofyUc7bvfh/n3mdKfzEZlxK6v86hXFRunqrl7uupYlKSi7K7G3l7yx6ckcj+BbOF0SpXmj0rC93A7DvIwyQb6jwWKjlVXgpli/Cdog6e1sSglJJkzixviV4/zvnqYaFdi/h2dnqLatLqzlUNkbjlvCuZSDti9hz3ZTkkwQ+EbehdaM91Q5Y6iPD9/zED/8xkdYsXKMyelZvKIHqZm5SntL66oWKRV89h+Y4Ppb7uf2e55k9+Y9XD1/hFLRY8qisptpSiNMeLzWYn1vD3+4+SEue/Z6ZmbbJjpovzDp+cRJQrUYcOeGbbxyrIfxsV7+e+tBnssQ67sK/P6WR7jowlOJJmbJspQ0i3Jnymy9zeknL6Rn0RgTU7P4/SVOHejiR39+gK5CmQNHpnjq6R3Uj0wygmZlV5Fn9feihnuZyjQ7Gk0eODbBbJTQilIkgoKEUsEnjGKajZCabpPFhnkUBD79/RUW9lYYHq4yb14Pgz1F5o330j/YxUBfD9VqkVK5SMGDIDBuBSVASBd79OYYZEqRhQbpk2qIopQkywjDiEYzYna2w8TxBkcn6hw4UufQoRqHj9U4drzO7FSTqZmmsRgqgezE5qalFEHRp9lJLAPLjFux7Ch0DNiSq2Ux+QXFeLnEdKtOuRLQ21UmJqbdCo06IIlPUESbz+70bAO0ZuGSUT71kdfzz6+7khvveIBf3ngvh7fu4pzeHl62bD7HWh2eanV4/uJRHt57gLOffSYXnHcKx44fMV0qIdEixhM+SHPCboURX/zEG/nSt67n1w9v4qWLhugRknunZvnjG59nbkeeUbBqYTSzwhcUfZ8tT+/ipj/dzcdOO4mfHzhqHnL2Ael5HmkcoxCGeCsFnvCJpbLJRo84MtFk9zsUWfOfwdSa54byPUOMtnF9Nz6KI/PCSW0xV/rKeETi1O5RZT4FcaidKArNUj0z4+2w07G0bZO6VMpEXh2aXQWKJIzxA3eQsgtuV5xWJnij7Og6syN7iSDO7LMYL5fguWiz6YvY24eFZTpNb+L4hNbt5Eahnu8hkdYZlBEnoXnJ2j2lJz33NMaT5tmcJpmtOmQ5MURb6oIYHBrSyv6gtEVGGxyAzJe35nnpSn7mLVgsGkGMEfN4uUM3CkObyPJszIycreV+KXKFq22jR1FMUPCJo8TwVuxfRtmbgLvy5aIoX+V7D239w47sa0idKu+iZFlGqlPLvbJLKLAdC2GTZhadbhfnfmDEM0qZ/Dd2IYq2YJS8QGh++KmDTtripc7sbNx1TpyjxKa0dGp86WZ5pvNThbIRvTSzQD6k+UUSgrKnGC4owizl7aeu4lePb+bMF13EZz/6Wg4cPkqxWMqXXFlq0k993V1s3bqb7//2Lzz4wNP01NusH+xicV83e6KMh6brhLEmsaRkhabowWvmj/L7Q8f58fc/Qt9AN3HmPOpZXoAa6u/m2z+7lTt+eANrFwxyy4Ea88pFrugrcd1Und/+5N8olgN79RUksTmpttttFo4P84Nf3MVN3/ol6xaM0VWQiCjmtr3HWVQOWNPfxehQL5NJxt5WyI7ZFhNRTCfJ8IWgaMVdnTChHUWEkblx9veWGR3pZvHiQVYuH2XxgkEWLRhkoL9Mb3eZUqBQKiOf1wpAeGZ5lqWgLb/M8yCK7X+skUFx7v8+iiFQ+XgErSEIIE7Jl15Zah5eQkCYkgiPZjumNttm/8EZNm89zM4902zZcYx9+yeZmKrT6sSUyoFhudk9nImnZkbNan8f/cCHNCOQsLha4qqhHr66fReXnHsKb3v981i+fB5aZsw2miSJJvADA7ZMUpAKkZkbvgZ8D6rVMp1mzP0PP81nv/c7Fsy2ePbwED89dpx3jo/zpae3c+0PPsNJy0bptDvESYRQWFRQhBKGXpwmIV3dXQTS44Y7H+OnP7qJhx56ije89QV84ZP/xNRM3Sy9s8QGJIyCdaR/gPd84ruIjc+waGyAX++dwM+gk8RWFWHCOmGS0aUkvYFHrROSCUlHayKrW1XKswiPzFKzTULKJTKVfYi6XpnbS2apgS0GfmB/P7VJf9mFsrJQw9RyqJQwLntPeXbPJPIUnlOBG7RRZpwv+Ut/jqOlsXthpJVDkXuQMvu1ubIhVtwl3N1fmRCSQFs/DviFgrlpZuRTCzd8cKEkI/TWlgws8wSs2ZeJPKhwYiveCfJceTt/6VkWGQjE2NioTqxyVdmWdT5msiU8kQkrbIkpFAM6YWhO33bO6mJk2kLGNGax5T64BiMwd30Tjjhpy3JKefaLTAyyw7KtktQkG3zPN+DG1DzcO+020lrfcsmVk9hb/EixFBiQmUtBeZaVky/TTMRP2lOLY165eam5fdgEmtt12Ea4Y+44vH1qT5SZTXW4tIdj9yRpim8Njpm9hoocJO8Y/XambPHPuN6JXRb2KOgr+Czrr7C+q4sfHZ/g5l/8O1rqnPcvhcqJu5Viha9/90/85Je3cmalwhmjg3SkYONUjf1RRCvWVAOfONPEWUacaCqBIswSXrpolE17DnPpm67irW+8imPHZ83YQSdktljl+5L6bMhL3/QF/qm/i5uaEUebMa9f2Mete47wuve/mle+8DwmJmbMhzNJKZYDuisB9eNNPvDlXyI3bmNpXx+3zdS5uK9MX0FRV4qdsx0OtENqYYyvJGUbS+xEKY1WaDL6vs/IcDdLF/dz6snzOGXtfFYsHWJoqJtqtQyejzGoaogMxI8kJekYzH+90WF6psXsbEit1qbe7NBodGi0OjTbCZ0oJYpjkjRDWbGQEoZyXC4V8H2Pajmgu7tEtVKgWgno7S5Q8KFaDigGikJRWY63Rx6YV77zBxBGCYcOT7Nrf42ntx7lya1HeWb7UfYenGG2adr45WpA4PsEBWUKXvYzWlKCIFBc3N/DukqRa5/ZzZZWyHlnrub1L7qQc89ZR7FSoNFs0u4khq5c8HPrnnl4JESxSf4sX7qAX/z6Ln7xXz9i5YIRfK/Ajj1H6LroVL73X+/nwMEDNonnoaRpYEdphG8DGrlkTmt6uypMzbT49fV3cfnF61kwf5hmq5N3N8zvWEq5HLDtmUO85R1f4hNnrOTHe44zMdvBWIIVUmREScKicokr5/ezfabF5okZPKlopCntNCPMNM0oztOVzovBCdMEqQRRPPcSMS8Ce6hzGqhMo1OzR3FxfVfLMEnJOUqFEObFkQlI45RCwcJnLSpeaLNgdrcBtzsxO4/80Zvvfc1uNZurA1gPkMM7Cau18C2nyiGSzDNSW6iq23/Yl4H993meIYV4fmDGjWLONWREgl7+tWC/7ixNkTa2rHIvkfUdCUEcRgbvJAViZGREu1SRFEYE5HuGISOtLyPvbEizV0jSzCKrrebVpnSyzFz73UslTTOKxYJxdQszv9Rmez7HcLGLpDQ1P2DHgJJS5M1mMvPicLFhR540vhCMTMmCDj0H4hNz8nohRb7HSW3E2IyLMoMGsWyYJMso+L755tsEhbYGRBPbE7aJT36acXiE2DKq3ILLeaf9IMhHfk7U5cpNRtNqxjBRHFuMtchTHS4+V/YDhnyFzhJeuXohD+w6wPrXXsX/e9vzmZyYsqeCxOAagJ5qlU996dfc9ru/8aEz1rI7S7nr+BTH2xEi03iWB5XozDRa7W1OKkWAZkm1wOnlAg/6Hr/8wb/S6sQkSWxLgxilJZKBnm7e/rHvU358M2lvDw9OtTmjv0RPo83+RaP84n8/wEyjTalUoBQU2L/vCH+45QGu/cPfWJNqrhkf4SeHj9MKU9pCYrdAVAsegZQIUhqtiHY7Ikmhv7/KmpUjnHn6Qs5Yv4BlSwYZH+02UVJ7YiQoQJyShim12SaHj9U4eKjG3gNT7DlwnEOHaxw+2uTQsQaT0w3aUWR/2Qx4UgoDuvR9hW//vnO8INsjSjITDbceBk1G4CsC3yPwJf19RcaHqiwY62HhWC9LFvSzcFEv88f66CkHVHuKVoWbQdIxN9wkI1MBtdmIvQdneXLzYR594gCPP3WAvQenabc6eL5PuWzIDcXAo5hlRDLlpEqZSwd6qHdC7jo4wWOTNUYXjfCqF1zEC646h6HRARrtkDhJkdoegNLYOr4N/2i4v493/9v3KG/cTHGoj6GO4Lv7D/KHaz/N2GgvYZIRJ5GxGiJIU2PU0yd4LYSwfgkrdOuulmk0WjalGBAlLXMrwByyRgd6eOuH/5fiph0MjfXxpz2TFJSgHUaUCkUanTZLeyu8adl8fvLUdp6/eJhmJvjTziOUPY+pKKadpIRCEKfaPlDnXDpOAOXYT45QqxxY1R7+nKfdYYhy3p3A6nfdDoJ/UMYq20R3S2ZH+DbhB1MUdLtWbXsYSZzk9HFHBxZWEuZSrHFiXvg6zcxLKf3Hl5pbgHv2z8hFGPb2E/iGMGCGJlmOZ8eCX81YLrOOENP98+zawLYG7DNOz2GbMo3ny/xFbF7Y4Hm+TxKHSGny1W6pLOwDLE3NIgbPxXKVXbRDasdJcZQhfEng2weZO1tL8/ab49CLHNDnS8Owcg9h7T6HVhcqlaNEkmNW7AQsRywbX4bOH/TKsl38QmD95eZ4n9mHvrTJsSgxH3JtmThYRPpcV4McpZ7pjCDw597OOBCeuby6xrq0WkrnOTC3jfSEa62w/64gPy24ZXkGBHaclWVzwYHcq21/8KVSQADsjzTvOmM1cRjZkZIt7ilJMfB5cssebrnhPj51zsn87vgUW+ptysqjICRamVNJGKXEqbYEZsMdi5PUGOvaIb2jg7R37OPRJ7Zz1umrqM8aXzZZhNDme9jqdHj51efwxXs28PzBXh5Xmi0zLV4y0M2D2/fz1BN7OPPsk3n4kc386s/3cue9GxnXmvctnY+vPK7dvZd6mBFnEPgu7gxRJ2KmbfSZY6P9XPbshZx/7jLOPG0R8+dVCIoepNrFoiDVzEzWOXS4xrY9U2zZcZSNT+5jx65Jjh5v0AnNnLe/v8xQX5GR/jIXrR+gr2c+/T0+lUBSrRYpFgS+MPY85fmUih6+MJZDqTRZJujEmsmpBqnWtNoxjVZMI9TUmhHTMxEzrYjJWovjx2bZvvMYE9NtUrvYHugpMW+sh0UL+lixdJhVKwZYOl5l0bxuerorKFJ6KoL1awZZv3aYN7ziVGbqEXsOzPLoxkP8/b6dbNx8gMnpBmknIqgWqAqPLVMNnpquc1JXhbPnj3DF/HG2Tk7zk+/+gZ/8+nZecM35vOqFl7B02SjtTot6o0WaJgSFAlJJSkGFmdk2u57exdmFgGE/4Iatu3j1qy5h5Yp5HDpyDOkpAr9gxnSYeLhnb9fYEQ3CHOS0ND/DyTDG98xBRemUQlAyUwigUizw8GPbefKBp3j32iV8Z+dBPLuQDAIPITN6ywFvWbWEL973OKdfejbXPraFtw4N8OIlo/x5z2GGCgE1FdNOYSYxSTq3B4ltu9qczk1ISKcGS2SItL4tuWZI2zjHxoGzJLOHZ50frF1Z2o2JDdpc57ZWJCRJbCCEvk8n7NjRkJd34MiEIUpk5qbj/lmhzY4m0xqdZPkBXhgpk9l3JOZZ4hV9ksQgXWJpDpjGIqtJshjlmylL5kjHQUCamqKn84XYq3BugPV8i8CXMi9EC/Qcwt52VUyA1aSvMqt7FyOjI9o5vIWVoxjxknF7xHFMsVgijDp5W9LzPZLIXH+iODZqR5vXz1KTkXau8Dkfd2Lk9DYy61IF5sFMHmtzaa45dLK0D9TUok+ynGwbh5FNYM1Z/PIlX967EHmb3c2WozBGOECb1cU6U2BQCHLSr5lZprmNyyy9DejRfDgkZDbZRZbD0qRnZo7yBL6NsidOgczjhp7NWrsypUuY+Z5HGFtCrZCUlabf9+jyFVcuGOC6Q1P88Acfo7evTJKa+J4UkMQZg0M9fO9nt3P3d29i0ZJhbj54nB7l07Hx406UIk64DfrKxHsTnSE8RVkoWmnMc8f7yGoh0fplfO3zb+XIsSnzspGW0qt8kiSiGhR51Tu+zIXNOnelkql2xLMGu/AbHZ7oqjCvv5unNu5gdXeJy+ePUC4XuP94jcemavgaMi2JRUqjHTJb6+AryaIFA1xw7kouuWgVp506n+HhXsAkdohj0lQzM91m555JNjx1iI1PH2Tj5sMcOlIj7CQUyz7zB0ssm9fF4vEqK5YMMthXYai/TH9PiUrZR+mYajkgiTq02yGtVkRof+6eNuOh2UaHiVqHIxNNOlGKFJqFo90sXziIlIp2GNnboqKrWkQnVgYlBHGiCZOM2QgOHW2wc+9xNu+YZuvBJpNTLabrbaIkoRj49PdVWL9mnPUnj3DGKfNZtWyIkaEulAMuammSdQkcOd7i4cd2c8fftvDgI3s4dqyGLzTdlQJSSNpZwkDR48z+XtZ1V5lpR9y57zB7tWb9OSfz2hedz7POWImQknaUEnZCBgcHueeBTfzHv3yT152yjI2HJvh7s8Wfrv0MfkGQWNuelMIIszyz+G80mpRLRcrlErV608IDs5x7507FnhcQpyGelY0hJANdvbzxg19j4Z7DtKsF7j1eR8QJ0jPU5FYa8amzT+XXGzZRWb+S333v3/jzLffzlvf8J189+2SOJDE37tyPJwIm45iO1tSixMbzyfUPZs6g8/2oknPLAW1xRllqd0s2CuyICFmS4hV80jjNW9jKk3acrfJCXxJH+U1GWSZaFMXWSMmcyQ93SGbuUG0PxVnm9hNmX1gISmTaPEuN5tZFcFUeaMlsvL7gmwV/bEM8Zj9lyBtSeDbBZyYd+oQ/80QFh7Rq2067PbfDjuMcGYWlIEt7mPeUifJ6c+5vc4hITmyMJzGeTVY54JcfmHSVcLV4DD49/0KUQGpl3Rhpzrc3qlBt4rLuw5XZE0KS2bSXQihlkOVSERQK5mVk53pxFOcfZAdBlGLuJ+IKN07yEsUZvsWA5Bltm/dWliQchkZ2gx3PxXGc72fcfiNOzIchTcwCHksudrbBXEgVJ/jFwCSPEgOyMwXIDJ2Zpajh9HgoVSCMwjymmySZ7Sh4BnmizL4kyVIyWyoz2Xe3sM/sL4DJj8dRB5AofLJEMhlFLE8SSAWRTohskirJUjwpDWEUc/pKyBgo+rx8qIckSZnNNLdNzfLSwR6uf3Azu3cfYmCgShKbh5lLgCA8Kl0lzj/vFDb+8hYWDPYxmWkemWnxssFeglaTcF+L961dRAvBvdOz7DjYRiWC/mLAbBgz3WhSb3ToH6jw/KvWcfWla7nwnGUMDVfdnABkRlhP2L3/GA8+vJv7HzvIA4/u5sDhOkpqxkaqrF4yyNXnz2f5/G6WLhiktxpQrfgEntkRRYmB7gVSc/zIBFt3TLDv6AxJoumqeIz2KBbO70UJzZM7p3lk6zTb9tRod1KUSEnShE6YIggYG+7isjMGOffkUUqVMjO1BoeO7+fodIcDh+r5z2rBeA/LF/ZxxqphLlo/SL1pZv//++tNXH/fQdIsZaA7YNXiHjY8uY+b/7KJKMsYGerijHXzuOCsJZx75iJWLB+jR3koCfOGA158zSpe8NzVHD5a4577dnPjHZt4+OHd1Gcb9Hf5hFJw98Qkdx+dZs1gD89dsYiqzrhr43Y+fN/jjJ+0iNe84CKufs7p9PV2Uwg8/vKXB1lQCJBC8Jedh3jb+17OwHA3MzOzaJ0R+EWbsDQn3N7uCp/9wo+5f+NOvvSpt3DBeadQm20QW+S3K+yauGhE4HlEoSkn9/V2c/eDT3Fgw1aes3oxP9h/mCAVpBJkppmIO/zz2lU8uecwh7sq3P6Fd7L74AGefcE6vvy5d/H+j/8vXz//VK5ZOp8bdx9iqORzLErp0ppOCkmurpbmhSSEIfcWfCtjMs8DZf09ni/yXak4Qewm3cnfxm+dWVDYBGmapeZIaMt+iX3QJk65bcVoJ9YPsixFYnw8kZXhzaVBbajHD6wm29oMJfkzJrOGUWnscGhtgjY6dQiXjMD30dokajOVQWqLg7b7ltjnlrlkyFymJ60LJdPO/gqeVLh0jwO8Ot1E4HuI8fFxbZbNidXXWuKtkISJ2UEoOwMTnsk6OxImDreR61VFHjNDnhC9FYI4jM0PTWFx7H5+IhDCkmszg/lwXm83u5RS5bRN5ZmOg/vhB4F1bEiL3tZz0hNHDjaeDSttEiZud6LJMCdk2tuSsEuzIPByOFuWGPijcCiBfMmmc4VrvkhMjVwmTuN81JXb9qQgSxLLyDGpCd9T1tKW2FOLzpvpCEWXL+j1JAVP8vyFQ/z+wCTf/f6HGezvIordz8Lc1AqBx4H9E7zszZ/jn+YPcd3xGhnmgwWCdpLiC4FCYpE4lEo+7xzu5/cHj3FMpJznF1k7OMg+GfHQjsNc8YYred/br+HYxGxuC8y0+bAXPMmO3RP824f+m3N6ivypFlKWkEnJBT1V2mnKtnaH2Sil7HmUCz7NMKVWb6KUx7q1Y1x5xTouPm8h80e7jVNEgE5S6u2MLTsmueue7dx533ae2nqQVjtmbKDEKSv7OevkEVYt6mbRWB9jA1XKBY8kg2arQ5yYdFynE9Jux/T3FZiZ7fDHv27n/o0H6bRCesqCckHR7sRMtTR9vRWUhGa9ybJhnwW9gqFu34xqMgiKiiiGzQdC7ny6SVe1xEnzCsw2EwKZMtDjUyp4KCVIEs1ELebARIRUPlc/a5SzT53Hz/+6j607Jlg2ouh0NI/tCfn428/i+c85iQ2bDvHU7mnue3gPm7cf4dCxOkpJ1q4c49yzlnDpRWs4efUIo8PdSJGaJJhvTsjbdhznhlu3cvPtm9i1+xhlX9DfVUL4ikSnDBYKnDvQy/JKiccOTHLbkaOosX7e8ornctqKBbzto9/g/SsW8+jEDI+2O/z+l59F68hERB2wRUi0lnR3Vdm2fS/vfMt/8LzRbn67d4JzLz2bf/3Q6xgfGWBmetYELnRiug9k+BaQmiQJAz3dvPKdX+a0qRa7Pdg4UafomefMZKvFsxePcnZXF//xxFZu+cXnGJs/QKtpJFvDA338/s/38aFP/B//edYaJuOQG/YcRgmfWhYTZoJ6nIAwYNW5bpnLMZlpS5QkKE+gUNZxgTmUMff77WbjhcC3UwirudamQJpYDqBrxSepW3rbnWqU5Ippk9MxL45Op4OnPMvVmkuPuuW5w5VISxQXzuBqhVf5rsKCJYMgAJ3lJWTzkLc6XuVZPbadtNgQgEDk64Zcy2tvEfne2Hb9HDZeCfNsjOMkL9iKkZFhrZRHFEcEQcGcyPPl+Rxjxli9rINCqFwEL6zfIzO9SUvZ9SzJMcELfPt2TuwyPs1n+058opTJUjuSq+f7eEoSxTFpmhIE9u0uTIckDmP7FnTLX4MhUNb4l8ZxfuXMX0JWPpWmST5TzE1mSs1Zu7TOv/mJczfbCF5iXzYOX58kJo7s+75lcWX590Mpn4zUvGjsAj5LE9LMgRfdSUQZgJtTRTqTnw0JpGlGV8GnpDUFD16ybIw/bN3L5//7Q5yybjGdMCXN4tydnWUJo4MDvOPj36P598cZHu3nxkPTVJVHGNnRgjAltZLv0RExH1yxiJ88vpMznn8OH3jV5bzwg//NinbEe9cv48aDx3g8geuudit76gABAABJREFU/QStVjsvnDohTrlUpEyRK1//Cc7pRNzeiQmjlILyaNliVwGoFAPaUUy9EbFwwSBXPnc111y5ntVLByhUfIgjyDTTsy2eePoQt97xDPc+tJtndk3gB5K1y/s4a80IZ60dYtWiLoqezbQnmnYn5JEnD7Flb43D0zGdKKG/K2DN4l7WrZrHyECFm+7dwU1/287CvozzTu5lpKro7ysbd4vv0UkDHt06zY/uOMCCfsm7Lus3L5aGpt5JCWPwvBRfeYz3BXRXfK57pMFNGxv861XdXHLqALNtE3P27c6sVBQ0O4I7n6xx+6Y6vVWPsW7NuavLgCEdPPpMm3ueSfj+Jy5g/nAXvb1VlBcw3Yh4avckG7cc5cENe3lq21HqjZCRoW4uOGc5V1y8ivPOXsL8kS7IIhNN9grUagn3PrSf3/5xAw89vB0dRgxUS/hFnzaaSuBzbleVdV1lnj42xV+OHGe/1lwzPsgLF47zkfs38sa3v5R3vOn5TE1Pm3BLIJFYa6hU9Hd3846PfpPeLdtZMW+UUtHj9i372B3C+z78Wq658hxmZ2tmUpBElIpFYyyMIgb6e7n+5of41n/8hJevXc6P9x6hqKX5XZGCvsDnbcuW8OG7HuB/v/IBrrz8dKZm6jkOBJEx1NvPn2+6nw984lt8ev0qmlpz856DKN/naBjRSlLaiTChmNxfZB6uYScmKBbI7C7Et/tLKaV5Bp0w4kqTFD/wzYjeBQUsJl55HkqaCYY7zadJAraI7Vn6hpt8pPlYXOZ7mjTJ7JJa2BuAzPXfbqxy4jgfy4czYzBLuxWcwOUzB1lTd7DK4Cy1sem5MJGb/khbesz9T9kcbilJk1wLjpA5A9HJ/FI72hdjY2PagcbcA8sPzBbf8aeUNG+vjIzAD/IbhrmOiXw34b446ZZrGPiiESyZBbGvvDwW59kZr8DM9wrFAmEnzKNrrk4vhFXBnnCKJ/cqk+9O3O1BCudlyHIcSM6JsUTfLE6tltJg34XO8AIj1HI6SIdYznLAmLUBOtSynkOxZJlBqXi2x5JaAKFh+ZsleZrEdjYskGiUNB0ZQyjUxAYNarwmSMOEQuMJyUDRQ+iMK+YNseXIBM950zW89Z+u4eixKfzA/EKQma/PD3z275ng5W/5Au9ZMsZPj9UIowytNJ51mehI05IJ/2/VMu7fsZ/pRf384f8+iiqVOLD7KC/45y9wlid49fKFfOihLXz9P9/FueetZXKyRrmrQqno46N55NFn+Pz3/4jce5Sr+vr5ztHJ/M/w7O/89EwTjeDk1fN59cvP5qrLTmJ0uGoW4coswbftOMqfbnmCP966iS3bj1ApB5y1ZpgrL1zCmiXdjPYZPHccdkhSSZikiDjkjocP8ucHDjJd6zDerZjfr+ipeMw2I47XNZVqF8qTFHWHNz67D6ljDk+E1OoRfgHiMEb4PtVKN6vXLGCqFvKpX+5g18E677iowtr5ZbxSgMpS/GKBRivi2ExIpxWxdLTEzZta3PREyP+7rMqaMY+p2YhWlNHsZERaoVPBaWsHeGxXm789XeeKU8r4IkMGgrCjSJH87J4ZSp5ksGrwJyevGebiM+dx6uqFaFGkGaccOjbN408f5c77d/PQkweYqUfMG+nm0vOX84IrV3HhWYvo6a+QJqCKZdIk48mtR/j1bx7lzr9tZbZWp7tUoFgu0EzNDvCMaoWzh3tN1DUTNNKU/35iOz/+7ic465y1HD581LpzJEkSkYQxfb3dPLFlPx9513/xjjVL+MaOg1QKim9cejZv/cXfeNOHXs1bX38l0zM1SyYWJFE751KVvAIvevMXuUx5PBJH7K+1KQZmgd1KIj556jq+eO/DvOgNV/Fv738Vh45NEAS+SXii7f+fMTTYz003P8C7P/o/fPLUk0gl/HbXQbOb0dBKMzraHvJsqsfoBchJslJ5tstmpHOBb24LWMW0KSHbZ5HEghjtdMRaVzOb0HRRfm0fVFkSI63syoyjMlx5zVC1M+sdwYaVzE0iSeb6RDl+PU0oBmb5HqeJzQbJfGftDIxSOpx9ahUKkjQ29QZnVnVeFWEnHZ6c6384aq+bADltrrlEGESOkmb3ou3OWlUq5c9omzDI7BsoiQwJMrYtcK3n2FGpFZcIRP4Gc1dT98DVtsyHFlbRaBJVaZLkSxkHBJuzBlrxiZrLRDt0o+f7VnmK7ZA4Wq1FrssTZCn2ze32NtqShj3X4PV9O2s8YYuUFxxTfIs2yLQVPOUzUW0jiNr0xKwHUwiT/HKnGZcSS2ITe3Y8ncye2l2zXwElKSgKQcXuU5Rnltqer3ISaZzOcXQKWlMueIz5AU/N1HnhVc+i07FomSzBV4osjYnjhMULRjk01eSRDZu5dP4gjzVadElFJgW+ljRlwpsXLWRyssZdSZvfffP/kSnBbK3O6Pwhzl23ki9c91eWdZfp92Dj0Rle85LLKAQFolbIHX99jE9++Zf84NqbOA/Jc+eP8ftajVaU5F2B2Zk2UZxw5vql/Mv7L+djH7iCc06fR7Vkynv1esid92zhi/99O1/42u3c98guVszv4W0vX8e7X34Kr3neGtYt66a7IEg7IWmSEMYChabdbvOV327hnieOctFSn9dd2MeFKwLWzvNZNRpw0cndnH9SF/N7NU/sneV4XbN81KevLMHzUJ7ByXhSUq+1mGk0OXJ0hkKxwPPO6OHoTMSP7q1TIqLHjzg22WZqtkOaChaMVBgdKtPOiqwZNRHs791bY+24zymLu+npK1MpmYfR8XrCUzsb9JZh73RCmsB4v0eaKcq+pBml3Lcj5tyVZc5aGlDwUnbumeTXf9nDI5sOMb8H5o91U/Hg1JUDXH7eMq68aAWrFg8yNdPktnuf4ee/38Add2/n+GSHgcEuBnpLKB9G+itccdkaLr9sLZXuMs/sPMrBIzUTnMgEe6KYx2stds52uGNimrO7y9TSlD8+8ATPOmU5CxfNIwxj87i03KveaomPfP5a1kcJBzyP/c0OF40N8MTu48yM9vKfn3gjjZZRsIZ2Xi7ISOKUof4+rv3D3Wz/2wbWzBvinolZfIsrmU5C3nXSCm58ehd9p67gm//+Diamp03AK0tNRF16ZAiUFNQbdU49dTknLV/IR39xCxcO9rJyoJunanWK9nSe6QwD8RHWqWFO5KlVI2RW7SBPgKNiBUxGm2vb33ZioNO5WKsrE+fWTmcPtQlKqRRKWuuqI4Nn5EU8bRldxmkucyyJs6bavb9tvM/FiA3o1RIupMGLOD97mmkrhpojFHuWtpvreD2VT4NkXpmQuXTLLfeNSM8ewjGdPONzz/K/U5ppVE9Pz2dcvFl6hvvuRPDKfmGuw+AplatiXcNSeYo0Tm16yHfxgtwxrDPjI0/i2Cy0bTnPXNdc8U3kCYQsy8jIcrG9sKmDOIrzK6DWcx4Q1+2QlhtjWuXmz8ixxVIRRxEgieIknyMataWDzql/gIa5ZVmeBDPJ6TziZm4k9kMgT7iF2R+Of4IHxYm3pDAJq4KA7oKHD0RCEyUpPb5H2c5EkzQltf+sS4MIJSgKRStJOW/eMPft3M85555Cf1/VfG+9gDSJ7EvXNNbPP20l37rlIU5JMjKl2Ntq0+UpJnTKW5csQNRafP/gQa773/czumCQVrOFXwio1+qsWrOM4WoXP77tAV63ZIxfbdrNxGSdO+9/ii/+z2+54+b7OEVK3rZmCWkh4KeHJ+gkVn862yDupDzngtV8+l9ewPvfdTHrlvVSDAQ6k+zZP8NPfvUQH/73G/jBTx8kTRJeec1KPvqWZ/HS565m3fJBBquwbdsB/v7IPu74+07ufewghyfaDFaLHD42wxd/uZVFfZJ3XzbAwm7BgaNtdh3qMDEbs/tgnW0HQrqqPksX9vDc0wcIo5Qf3DXL7qmEk8Z9dJoyU0uJ04ye3jLj4934QlCfbSGDCs9eP4ok48Ynm3T7mr4CeIFkZrbF1p11ogxGx7oIKl2cv7IbnWb8+MEWvSKioBM6kaboKxaMlQmUZrKRkUQxGw+kLB4KGKx4KJnyx4c7jPYIXn1RP0Ul6Aoy1i8scdqiAnsOzPD9W/cg4hYrRwsmEVmA0f4iZ588xouvXMtlF6xibGSAnQem+eOtT/G7P21k89MH8X2PkeEqxWJAT9nj/HOW8ryrT6Ont4ttO49z7GjNJLc01KIED9gdxrxh6QgHj07z1d/+lbLvc8b6ZaCg3Qzp6+vmvoe3cv2Pb+aCBePccvA4xUDyvPnDfPfJ3fz7x/+J5cvHqTcaRnNrBU8aQaFYotOM+fhnf8iLRof420zTKoA9jnVaPH/RfGrTszyiM677v4/Rjjv2GWpeQI4AndlxrfJ8mvUm605ewvrVS/jIz2/m3L5eTurrYePkFF2+R5hlJFrkMV6D65ibLpjfbc88w7I0P1QGxUK+XJae2QW45FGaaRsWknN48yy1IjyT+pKeysf1ZvJmcCDK/jOuTWiCPK6v5ojmKn9RuB6bsiQMRxZPbfcmTZMc5YS9jZzYrVPKs2oIc8hOjBwd5ZndrrS0A+lJEot/cWMuaekkLm7sdr9JkuQ7cd/zUOVy5TM58yXNeSP5XzS1bzGX+BFC5h4PfYIPI0uy/FqV2iW5kCahlSUGkeJuDJ5FirgXV5qmlnJrFt9KGPBgkpgRGbaEkwMQrcheSvtNcIhiWxa0NLn83+nQycLykn2l8kW5UuYUmnO08mWTfftbjIAdJOaFG+y4yrlUMpu9dmYww+KSeUoq1RkFz6foSzwERfvCvnrRCGeO9tJOItpJSlkAUhFheyeZ+cFh8dAizVjYVSJrhuyOOlz5nNOpN1omJWHbuMKaEKtdFRaPDfEfv/0LH1i5kCOkNIG3zh9ndqLOj44c5udfex8nn7qUeqNjb3mGKttotHjO2eu44+HNjDaavHB8kOsf3ERz934u6Stz2cJhUt/jj0cneLzRYbBcpD7bYbbW4dyzVvDFT7+Yd73lApYt6kdmMWkMG546xJf+56/86xdv4i/3PMOaJb28/42n8c5XrONZJ4/QXa3gi4QnN+3m+9c9zu/v2MGTzxyjXmswPdVk0+4af3v8MNf9/SjnLCrwoRctZLad0Yw0w2N9jI5Wmb+gj8WLelEIJqc7lAKJlpJTlnRxwcoqWw52uPGxJr1BSkFGtMKY2mzHiHiEpqunhEDTDjUrx4rcs3mWjYc0RR+qfsJof0BvX4FOo43UCV4a4ZFxzrIyU7MRv3o0ZElfTBLHTNUjJmc6FAJJX5fH0vESioy/70gY7VNsPpCifI9/vmqIKMlozNboRHBgSnNkJmb9whJnLFL84LaDHDjeRCYxT2ydYuOWw2zcvJ+ZyWkW9vm86PKTuPqS9Zx1+lIyrbnj7zv41e83cO/9u5BSsnD+AKWCohRozjlrPs+/5jTK5TJPbTnI4aMzFIsBlWKBVhizpRXx4iUjnDbQzfeu/zu33PcUKxaPsWLlAgqqyHs+9h0u9n22xwmbJ+tcvWiYzfuO03vGKv7l3S9jarpGUPDtjD5DWEndUH8f3/vVbTQeeZrxwX7umZihqxBQCzus6+/jtGKFb27bwR++83FGx/rodELzILVcOZ2m6CxFSt+y4szvWrvd5qSTlnDWKSfx4Z/dwAWDvZw62MsTE9MUhCK1AYjMiqaMkEnbQzL55MLzDIpIeX6+eDaLcONod//jRhYu+i+t+M6zL8vUHojd+Mkor/28juBMryJXDYtcFuWixXnB0b4Y8umM85TIuRKz5wd5BSAvGVoToaPyuumLdoJAO21xZJA5yvQJO2U7eXJqBcdgcwlaJQVRmCCGR0a08ZinOY7dgA+lXbyYWJrQTliS5YVA5Sk7I4M0ifMCoSvYeHbmGEcmDox9CekMK6U3yYIkifMavtPzoLUBHFq0iOdbS51DFdi9h+P9Z6mJvAoXN7OIY/fyS3WKsh2W2JJ93ehNSGVGUllGZpfRrtUqpZmdupibS3aIPLtt2EK+p+xiPbU8fbcbMfPVovLwBQSepCgkmcp476lL2XVkij1TdV6wZgGHOwnX7zhEJ5PgCZpZRitKaXYMLbPse4wUAgaKHleMDPD9Pfv41Y8/Rt9An7W3WX6NAF8FRFHEYH8Pn/3yb7juJzfyifPXkoQZf955kAP9FX74xXeyfNU8JqYmKfgFm8QzJ7HB/m4eeGg7H/rUd/nAyACPRCmrqkWE77O52WTDZJ1GJ6Gn4BGHKa044ZS1i3jzmy7isvOWmXh2HKO15r4Hd/KDXz/Mn27bhO/Diy5ZzosvWcaa5QOQJiRZRqsVUg7gtzdv4S8P7ueUeT6nLyxQLfu0I027HdNVNCXUx/dF3LMjoassefl6n9EuTa2REkWJvWEGDA4aQCKeR1eQIISiv7fEroOzfO6GaTw0lyzPGO+OabfNgrUTw7GaJo5hqM/nd09p5vX5PGdNkQMTMWXaEIV0dRfwpWa2rdER6ECx6ZjZX0iRMd3OOGskolTyKZchDDUpPpWqT1+3z2QDth3WPLQz5oJVBV5+QR9eENBsNM0Mf7pJM4zYM5GiU+gpZXz51jpL+zxOGlJEtgvViCW1SDE62sVLLlnO1c89k6m2z+6Dk1x/65Ncf/vTHDwyzWlrF/C6V53Ni567muGBqhmXlkscPNjihz97kN/+7kFazRb9PWWElHTSjOeM93Nubw+3b9/PrUcmeds/XUOhEHDzL27jjasW8s09xyjolFctHuGrm3bzmx//GyuWj9Nqh5biYIpzYdikUCiRhJKX/tOneMfwAH+o1ZhsxnTShEog+cCK5bzvzgf54mffxmtefjHHjk0bcjXWPJmZ1JORxWkDe9UglSHwJknE8PAgDz38DK9773/woZVLiDX88cBhhAqYijXtMJqT0vlmIezsgjmJ1R4a0ySlEPgkWWYnH1a2ZIvHeXhIC/DM80tYE6qh4yqiMDZncfscFP+/XYUpTttDZ54YlbmfxGkjpLKpVItdc7UK5WoJGbl5UNnxmFRWIyHmxvS5w8k+g92eR7h+h7IpMmmoGq6jZl6kEATmeeJuLcbuqFDlcvkzJzJk3BLdjIFc5NbLRfOmyCLzmVqapvnLAIcIsd8s1xhXNqLn+QaYKKRphIscb2y87Mbgl9kiolE7StvgNgsiU2JxlX/3xnYjsbwAYxk0qU1JueKjy0kn1gFiIm/mz3Ivpcw60OMk+UdMiZz7O6l8ZilyI5nO5tgyrh1u/nxBUSqKvkQJKEsjgnrzSQt5+sBR/opi3rqV/PCuxxiQglevWYpGs63WQGaSkqfIlPm+xzYyOBXFLCgXaE3UOSQyrr7kHGr1pinRWweLzjKj0Q0jLr3oNCbDiK/f8RiPN9uc9tyz+fZn38nQeD+NZstQB7K5lnxPV4Xt2w7zog98jVdXuzkmBHfP1HkmjHlkusGBeoeSH+AhOHJslkWLRvj4h17Av3zgSlYu7s9nxw8/updPfOEGPvuNOzk6UePVV63mE289kxdfvICR3oBmvU1tqg5ZSha1+eIPNjBxZJJ3XzrIytEiByYTdh2LmQ1Tpusxe6Y0rQ6ct6LEGfNiDk1G/GpDxO6jLYa9EKkz29KNadTbHD3Wxk8TJieabNzb5o4n6vz0/gYXr63yrst6melItAwolkrsPhriVfsYGuzn0gtXcDQs8PTeWb74ljUsmdfD4qEiSxf2cqwW0Y4FPdUSS+Z18fj+iMcOCi46qcylpw/wnJP7aIYxB2YSFvVCo5ORpoa55Bc0O/eH+FJw2RldnLa0yN+fiXhwW5O+kqS7UsTDaIKLhYDRboHwFLWmZsmQZMMhzTWrA1aOlFk5XOCStT2cvcCn1Qj5+d/2s+mZg5y7to+x4R6etXaYqy5ZQ293Fw8+vpff/PEx7rx7O1EqWLZ0kErgUfLgkktP5vLL1tKOYdOWI7RbIUM9FbbNtHhgcoozx4d40fwh7n1oM7c9voWPrF7C7RM1DkYx548PsPXAcRaffzLvft1zmak3rIpVWrueiccPDwzytR//GbF5JwO9Xdw306TieTRJ+NCalXz53g087zVX8KF3vIBjE1M2JOP2EibRmSUppWKBSjUwwFYpjbtIevhBQLPVYsXKRVx45il8+Gd/5pLhfroDjz21OkoqQqsG9jyfOIoICgZZZPS0pvPgB77ds5oRdmK5fcqT9sZjyoBSCdtxMwbAHAUlpS1+2peMVHMvAvscc/UHF/7J7IHcwQqFTWK5076pNZgXqFR2XaBt0tUyrBzoNkc3aeNK0a4OcIJEylHO0RqhRE7YSJLYJsDMjkM6QZecezbmHRDrGpFCIkZHR3WaJiipcllJmqbGBOgshA44aOeHYRjlo6Q0dQ1qs+zOLEQwTbJ8YeP7Fryl+YdYrbbfgNQCFzWpwSELiOM0b2+6Zjt6TsSSC6HSzLKLzKkin/PZ249BwkdGah8YmJwR2RvUSpqkFrPuE8eJRUBn+e1La6O/jeMUPwjyqLLGdkEsG8iFA9xIz7PR6ILnE1jJS7fnUYvbvPKkRZSjiP/dc5i7fvdfLF66kLvvfogv/u91HH56F29evYSeapk/7D3C3kaHUsGjFRlwnEJTkoJqoHjZ8CA/OHCQ3177bwwM9pkyojDz1CyOc+2lH/h0Vcts3bqfIPBYvnScmWaHKAxNp8Ni9Y0YqMKhA9M8923/wXMLRZ41PMC3jh+hK/OJMo22h43ZWpPungpvfeP5vO4VZ9PdVUDHGULDhqcP863v38VNf9nKYG+Rl12+lGsuXMiKxf2kccRjj+/k0ESLalfAgpFu0ijhS7/Yysq+hNde2MfeQ20aYUK1GlAOBPVakyQTlKsVarNt0hT6uhT9JcHTB1r88OGQeiy4ZmlMj0qJM7e/0uxtSCaybkb7fVaPF/n53Yf4/OuXsXzZAmIk3V1lHnhsH4eOzXDVuYsodfcQhglv+68HeNul/axfPUoYm1Xshm3HefTp47z5qiUEBY+/P36MGx8+xvufP8b4/CGm24rusuTHN21FZCnPXSU4MhUy2UiotxL8gsDzfCZmM8KgTOYVmajFPLa9TYpgXq9gpFsw3iNZM+4zWEzo76swNdNGpxn/99dpjtYFH76ggKZgIDgZrF5Qoqsa8IUbjzHjFfjWhy6m2tNNqbuXwPfYvWeCX9+6letu2cT2vcc5bd0C3vWmi3jZ80+huxKghYcoBjz+2F6+8j93cve9W6kWFMVygUaSsqKnwuWjfawu+dw90eCWYzW0Tnn5WC+/2HWU73zr/aw6aRGdKDUPXbJ8rF0oBNRrES97y7/zvuERfjU7TRjBkaTNB9eu5C9P7qC9ch5/+Pa/cGxq0jpQ7INbBYavlcR0dZV5evNubv7LQ7z/n59HJ0ms3EwhPBP/FxL6+/t54IFNvPM9X+Rjp63mhn1HORynTIaGaou9OcRJTCEwwNUkM1UBoUVeL8iyjMTuGBxSBOEK1qntvDktg5eHZdQJ4MU0NYt2nWmCYiHncxmulDGZnnibcOMtbKFbWb2GsrZWz1OkcWb7GpYHaBW6yjfJUteQTxJDCJFK2sU6VuYX2wuXK/G5CZh5lrlnNELiByrvgsSxiQObF6TIn82qWq18xvi+XdbZ6mWFTSQIsxD2lG/JuvIfSLG+jYYJT+ZFPAcOM/N0FzdzADFbmPHmBO1OzIQ2y2gXIXON+MxGdLXA+sR13vNwsVmnEdWQj84838Op5Mx8X+QvwTTNcmSJ+2FgIY5JnBIUAoslUfmS3dwyzE7EUwrpGQmWiRunVgM552X3pKLiKXwFPZ5HI464cN4QJxULfPqxp/n1dz7BqpMWcOjwIVasWMQrn/9sugf7+M5dj7L70HFes2oBi7vKbK01aSQpZalMg1wJ6lHEeFeF7nbCplqNF115AY1Gy5xmotA6XDyLcYHZ+iyDQ11UK0VmZmfBnmKyxMQD0jSlXCrTrHV40Xu+yvpU87x5Q3zz2BGCCGI0CZJWJ6Tdjrj68lP4+hdezhWXryWNQnwJ2585xOe+fgcf/8LNHJuY5a0vW83H33Qqlz9rAb1lxQMb9vKjP2/mp7fu5LYHD3PP48d4ePNxrr/vEEXd4R2X9DE7G9Lf7dFTDchSzY69M+zc12E29Tg80SZtdfCUYPJ4h2O1lH3HIkaqgt01xRnzMroKEAuPSsXnwUOKYlc3r75onJdfOMZNj00xUNG86tKFpF5Ad6WLTTsm2X9wgpedN5+ZZkIShnz/T9uoiBYXndzL/v0TzM42mZjp8LO/HuLC1T0oVeCxbVP8+u7DvPfqIYb6q7SihKJK2HGgxl8eq/HKCwbpGxlmYLCP8fkDLFk2Su9AH7tn4a/bYu7Z3OaJ3U027WoQxSnS9wiqXUy0JHc9NcW9W9sIYLRg9i86TBmtaO7cpVnYp7hybYGFo2UKRDy5r8lMO+Utzxlg454G371pL2cvKVGbmubA/sOEjRpXX7SYFzxnFcOj/dy/YT/X/fFRHnxkF719XaxcMkTWajM+2sULn7eG+WO9PPLkQfYfmmakp0IjSXms3uaBqRZb6gY3M1QK8FohcuU83vuWFzA100BrR56QORV7sKeHL//oT/ib9tDV18WGeoOJZotXL19E42iNe3XCn777SVphhzSKTbE2KBhaQmQQIZ5UZAm844Pf5IY/3seByVmuueIs0MpIky0IUSlFbWaGdWuX04pS/nLPBk4bG+KRwzPm1KwUpCax6KL+jl7r+ltmUY9dPtvnhhRGZRvHpiBtD7NJmlpDquXy6TkklKkMiLndiC0J5xiRzN467O5C2AN4FJs0q7LaXrN4N0vrOE4sPdx5lWxfLU7zAqBzu2t720gzk3h1zz5lR/rCVjPcy0lrTPze6qeFpa8j5trqZrRvnnuuCyJGhoe1c26kNpXgaLcuu+y8HGkSz11tnPbAxXZ9nySK8HyfsNPBCzyU8PLCWZZkKLt0di+IOWx7ap3lwrbh7Rfv+P3aXfcsCRdB4PnEaWTw77YBar6ZSX6dc38XgyjARnB1Po6TNsFlkMvaLOnsaSLwzeLbbfpzxbLQNtVh/SmZuxHJnK7ppCvlwIx5qp6HzmKW9HXxmqVjvOP2B/ncZ97OG159KUcnp63cRuJ70N3TTX2qwVe/ez2//t2dXDM2wIWjg/xtaoa/Hp1BpeBLKAjwPME7li7kf7bs5H+/8SHWr19FrTYLIkOnMUFQJNGp8V2gDSLGs04ULez323zfioEi7mRc87Yvsmi2zj+tWMjndu6zLxgTpqjNtli1fJwPvvdSrrriZNJOhOdJpibqfP/nD/DNHz9IGCe84XkreO1Vyxgf7sYTgm27jvOzm59h6/YjLBvQLB8OUFlGqeTTTGDf8ZB905rjTZhXzXjuKSWW9CvCVkwiPOqROSzEzRZCKmrNjI37M5oxXHZGhfv3Z4z2+Lz9ymEOT8f4ZNy7pU69pXntRaM0U8ETexp85Q/7ed9z+ygGkplWQiMrcNuGGVaPFwlTzaHjLaJM046gt+oRRqm1YipCm7xpdwzsM0FQ9qBaNp2pwFcIramFcPriMped3odfqtDX20Nvd4kjHfjxLdt5eEsNnSXUJutUu0qcum6cyy5Zw/p1C1m2ZDE//eV9nHPGUn7xi7v54y0bee8VXZw+bvx3nVDyl+0NbtsW8fWrK4z3KMpVn66BPu7acJhOK+U5Zw7w4T9MMduC81eUmWxGHJ6JKJaLnLxkkJdevgq/0MX3/7SFP976FDPTTa6+/GT+37sv46zT5tFpNClWihyf7vC1b9/Dr377KJ6SDPRViTomat5KE9b3VijGEbu6y1z37Y+iA8VsvYUXGFyI1hmVcpXDh6Z5+Zs+w3vGRvhtq810q8PJQ91cUirxyU07uOVnX2TxkmFqMzU70UjNbgLzfEjSjJG+ft736R+w+fb7+dIFq/jo/dsZOXU13//qe9HKdrBIEZg+RqkkOXRwlve9/fNcMX+I3+6aINLQimNEZlhl2qYnhVUyoM3vsLJjnCSN7YzcMac0vu17uAeusO4fbTW3qVV157cYN+ZKzcFbZzYNZkMA+Y0EbZEqnk1hZbk9MLWmVD+w/TLLGZP5lMezB3dhD9vWyFrwSSMrl8q1GF6OQskye+i1/5kpWvr25pT+A9FcSrPCyGySLLW75kyDGB+fp9M0JYkjgiCwvBZboLPjJGHfhq4cKC2XSgi7KM8JueakHnY6JslgFzpBoUAUhidgk2WeNc7JmTbNFHYiyuVSDgYzRj/PYJbtvstF3bDN7yQ2Jx9l43lKGoe1soGAJI4RQuUN8NSSbgXmKusQ8L7v2/6Gxb5bUYy2C3npm2KOQ7sndt7pmqhZllmfiUk2lDxFT+ATkNJX8HnbqsV89K6HeM3bXsLHPvgqDh85niMGpG/wzVmWUSj69HX3sHHjLj799V+yfeNW3rh8Ht2lEr87cIy99TbzK0XiFE4f62MsFfxVwY0//iStKMw/HJ6VXUmhcpaOlCCll/89oyiit7eXmYkmL3vPl1jSaPKWkxbzme17aNZjY8prh4SdhDe84jw++J4r6O7yyGLDIrvx5if5ynfv4YmtR7ni3AW842UnccryPmZrEb5IuOWBQ/zhzh2cOQYXLisw2Ug41ILabMxgX4G+iqSvKPBUxuHpmHt3Jjx+VLOwR3D+eMriPh88yfFjLQJf0I4lM6nPsiW9nHzSCCjFJ368hf989RipzmjUQ47OpvzqgQYnjfnsmYjYdzzieDOlJ4BqpYDvaQYrHj1Vj4Gq+TkO9/gM9ZaoFAMUGhVIRJZR8pXRGWfmYdMxE1LiJKMTaxphTDvSNFqaWjOm2Uk5OBFRDxOGun26ij6Hmik7JhOmZhooNPNGhzn/gmU8/4VrGBvo59ChNrv3H+HQ0RpHjh/iA+8/lbiu+X/vu43WVIMrlyl6y6Y3VCpIvvNgi1NHBS9eFhPqgL6hbsZHqjy0aZr+Hp9Cb5lP/+4IH7i4i2rJQ3ow00x4+kDIE0fgonOW8s+vOYsnnpnm539+mutv30Icp7zx1c/iI+96NiNDXbQ7MaWBHh5+aDdf+MrtPPb4PvoGKviej0BTloK3Lxnk8xueYXTBCF/+9D+xbNUS6s02UdQhy2D+6Ajv/+xPaP/tUcZHh/nFoWMs7y3zlsWjfORvj/G1r36Q519xNkePTiDEnB9IKImwJ+6hgQo//vU9/Od//Yyfnr+ckXKBpODxhlu3cNnLL+cLn3gNx44dN4dCpZDCo1hQHDlW511v/iKXjfVz3a4joHyaUWx1stbpLsyfk2kLatVzZtQkNqRsp90+cT/hRE5OgSCFQChh/u+0A3qLnKHl9rRmoa6IomgOD4/7793kRtnDrrl9JElix/3attfTvC3vaCHu2RwEQV4U1JmeqyLYvYfypPka7c44sxOWzNI+wOx7HT3cs5Mh6ZkqQXYCr0sqm4AdGx3TwmLQsSJ5KV0UVOdLoSAoWLe3eaP6hcDAUZM4L7oIq+NzEiVhybLYF5B76+b9PQFKKMKwY65OlgkTeAFIU/IxVjAvp1g6T7sWNoFgvSVml2LidJ43N6d0b3IXYXMLdrQpKOm8W2KvrPZ66j7IprAo5vLQCLROQcs5mZQ0pyWlTGlHeoqCpygrQbcUBErwnjVL+cqDT3L6Nefx9c//M0enp8zuRCjLGxL/MIfM0oRytUJJFfjDDffy+W/9hupUjdcsX8BxDXdPzdKJUjxP8vZVS7l+03ZWXnY6X/n0WzhwdBIv8O1V3aXVyP+enu8RhWam2d/Ty5ZtB3n9B/+bZ/maV560kC9t3U/byouOT9SYP9bPv//rNVx52Ro6nQRfCTY+eYCvfOcubvrbdlYv6+Ntz1/OpWcOEnYiwkRQ8lK+9+fdbHr6MG89v0oYJmw8kFLwYWzAw7PEWuVJ4k5ET9WjpwgLBgvcuWmWX22MEV7AW9drZBwSZjAxmzGRlXj95fNYtHSckbFhPvGdh5g+NsVZS3we2N5m32TEVFIkTg36ZelwkQUDAQsHA3orku6SYqjHp+LZXwK/iNAhGghjM/uNrX0RIUnDBDyF8otkYQfheURurKI1ntWs6sSMFj1P0olSSkWfJ/Y2+Pptx9l9qEZRZZx+yhpe8rxzWbGmyoZNz7B/T5tmPEtHNzj3wiGeerhBh0Oc/CxJY6qLaDbjh18+wBsvKpN1wBcR1UDw+OGUAzMx/3pxkUDC7n1TDA73cfJJA2zf32Tl0h6+dccEIot5+foqh2ci4ihjyXCBSqXA//1thom0yOfedRbd3QM8tm2C7/3+KW65ZzOrl43z2Y9dw8uuWUccpXiBJMoyfvCTh/jWD+4j0xn9XSUSkXFSd5Gr+rr42TMHeXimyT+/+gre+aarCColNBlbnz7AW9/zFT68aJyvHZ5AeJIPrl7Ef/xtA69+5yv42Htewb6DhyiVCmaHYB+U2uLOq2WfjY8/w+vf900+v3aUc0a7aKQpRenzhnt38sLXP4+PvPMapqamESpAazNeLpWL7Ng9yYff/SWuXjDCb/YcoR0lCOXTjtp4ykNZKnCczqWf4iQx2PIoRtjFOTlbby6ur7UmCDyU9OamJ9Lx5kBkmQGbWnq3K/m5/phDhmSZCX0o3zON8cDLsVA6gyDw7AOfPKVqHpLkQFMnosPeQLDJLmNLNdwq36rJ89uHNiMybXttgWcOmKnVW7vumW+TV644qZR5rvpW8CakQFWq1c84vaF7a5r5neXHFIK5xbW1VXmeb2Z7Vn7iWFAusYB1mLu3eaazf/CBOA6+tLlp7IIntYsjE08z3J25PLLFl9ixkc6sWz3naZkFlpNcKaWs2tGi5pV5qUmbntJWQm8Q9JY2LO2LAgy3337TXCsdLefGXsrLJUtak3dNzN9Z4AlBWRgez4uXzePe/YfpOeMk/u8/38NEvW5hbQYR4WaWAkO/jGPj3G63WnSiFqeftozXPO9ijoQJ37t/E5Uw5DVLRqkUA/a0Ix49PsE7VyziB3duoKe/i/OftZparY608WMhJZ5fAJ2ghEen1SYo+Qz19PKnmx7kTR/5b14x0sPFS8b4xvb9xLEg7MRM11q87Hln8+2vvZZ1a8aIOwntTsh/f/tvfPDfb2DPwWne9Yo1fPLNJ3PSom5ajTbNVkzJ13z+51uZPjLBh68e4vBUzKHpiFOXllkyaHZpcWzx+iKlHEBP1aPqaW7Z1CZVPrVQ8Kpz+7j0/KX4PT0sW9bP4/s7XLiuj+Een7ufmuB39xzgzo2THJ5O2LgvQivF6sVVLlld4aVn9/Diswe4+JQelg/5jA1U6K0qSiq1XKYEpE8ShWSYDkAgde4+4P+j6rzD7Krqtn3vtdvp02eSmUmZ9N4TCCV0EGmCBREUBQH187VgLyi+KmJFFBFRVLCCDaVI76EklCSk9zq9n7r798dae09er8t/IGTmnLPPKr/nee4njGShYVxApCpvdUJStikbKaMIv+piGhr4UpT00amvb2RzT8C3/nGYA0eGOGn5TG78xBV84MplZNqH2TO+gz3dhxgY62P2SkFjU4SwPdY/PUzLjJADu6FU1Fi0NsfWjRU6GzJ87QNT2Nnn8eSuKs1pjSmFiOYU+I7DpBab7pGAgSGXrlZBTbMxvQoPb3U4b1GKpgwI02L/YIDre7xrVYG39o7yu0ePctL8LEvnT+b05R1MmzKJ9W8e4g8PbGD/oUGWLO6kqTlLUHVZe+pcTj5pLpu3HOXgoQEKKYujFYddxRrvnzmZVfV5/vjcm/z1iQ20N9SxePpUPvrNu1nn+Ixk0mwoV/jirE5+u3E7C85cww9uupbegQFMU0xonPKuJ3lLboCdstiy/SCPPLaRtozN4voMOUPn25uPcixTx0+/9WEJYVW0CNO0IQpIWRbb9/Ty2tMb6Cpk2TJaxtQEXhBiaDq2IZK+b98PFZZd6apBqEjj0iqcSqeOG1tJB6ZpmMo55SckcqFJjHsYKUdpFCYwV6E05rhWW+onE82nIiLRZeNRfRDGrEEd07LUWDBKsmie56nec6krk2DhJ1LssZMrLs8yNF0W8/mhimooOm8YKPuxIQORSmMW6rAumHDCxjeX2DGrTZrUFsWhGN0wpGU1TmKrWlmUxhGFsjTFTqdwKlW1gKv5mPpBnu/LAKFqrJIfhGRiSSODugoq5HIUOwBivEi8SZlytlet1shkMyrNHRCGKJRxKB1gTPSHyPmcrNu0U/ZEOZMSqyIl/kdJh7E86euqhjauvpVRfVT2Rc4hQyK1UfrJZhEnUGPNBU3y9gWQ1gVZXYfQ56r5U/jXzoPcdMunOPW0ZQwOD2KasqAlCFxFx1SVtFGAH4Jh2ISRhxaFhIGHpuu0NDSze3cv3/7lP3jzpTd496RGVnS08O/eQfq9gAvr6vjh5t385DvXcdl7zmSkVMWrubiugwz0a6SzNjnbZrB/nJtvv5+nH3uNLyydRc0U/PFgHwXLYni8ip2y+epnLuA971qB47ikDY2XXt3HV3/wCK9vPsZFp8/k01ctYGazoFqLqFYdvJpP2obb/r6XoDjCDafVc3CgijAMcpbG4Z4aI9UQz3OolH38SKBbGiLSKNciDlcMMoU0y7ss/vnaGD+/YS61akhTg81rO4b53RM9TG4w2NXjohmCzkaLE2fl6ayLmNSkY6r0bGQI3JqLpulSx1L9EJmUDKWVqy4hUMFgtBLKUivACaHkRhTdCDeQ3Sl+FOJ7Uh+LrYtCRBjCRIQBKUsnkxLkLY1MSjCjNcOhgYhb/7aPSmmMmz7/XubPnsmwux8n18u+Y0fwI0HoRhBEOFUZos3k4fl/OMw8yeLgaz6nXZilY1qKlx4fZ2y7z6Wn1rNpR5FrT8uDmWdgeIzq0AC9x0ZozWdIpWDbEZ/OBp3OKXmslMn/PjTKvDadD67KoNc1kNJcthyoMjJe5aSFeb7y1z7MTI47vrAWMvXUZXPs7/f5zT82ce/fX2VSS46v3vhOrr5iNV61im6n8TSdn/zsae68+2maChlSaZOxmsO5rQ2c1d7MK0f7+eexQVKTGmgt1biuo43PHDrKjTOm8crBo/RNauSh39zEaK1K4LvJd09oSkPUI7RIon/QNOrrMzz+1CY++sW7uWxSjqxt8IdDIzz426+xYlkXo6OjCvAq1xTPc5jUOomf3vMYb/zpISY31fNU3zgiCnH9AFP1wqd0DSeM8NBkriaKJMcqVFMAZH97GEQqpiDXA8MykgNtFEU4jiOrJAxZhhYbduSYSh6GZXQhRpagJPMo6RmJXV+BH2CnUkRqnYu1DT+QUxMhBJ7jYdlyVCWrLnSl/0pzAJGGbaeUG1Yox6jEsBBX7qrmU9eVG6BpmBIGqTRmP5TCjq7paLp0z4aEaqMXiUZjWiZaa1trpKlkpW1b+L4nT/Bm3Ochg3Se62PZVoIViZ1Jim4mcSGeCs8oPstEpiRIhCqZTJdCtqeoufLNlVc7z/eUrqAn7ZGyTcxPAoASgOihGyKZ5+vSnoA47lYU34CiMIYequRlpKmyFS3hugiFm45vOpo2Mc6TPcqaYuYc15uuG0ki1TSNhM4bBJJ6miZERCEXzmzn2NAI9rK5/OS7n2BorEiIl2BZDE0niAJZGekE6IYFkSdPvcIkDB2V7NUp5Atk03mefnET3/3F/YzvOsQHpzYTpmyGDANRq3HH1v18/IrzuPbq8+lsb0ZDznM9J6S7f4T/PPEav3/gGeYHcNXC6Tw+NMLz/eM02hb9g0VWLOvi1m+9lzldzThVlyAM+fndz/CDu56jqc7ia9ev4oITWxkaLPHnR3axu7uCF4AIQ2pBwNsHKpwzz6Q1K5jeDM0ZDd2wqFRr2JaBRsDQmIPQTXIZg2e2Vdg3pjN7aoGPnNHCt/52jFNnW7xjZRPrt43z0o4iW464dLaYzJycYUVXHbNaI5rzBl5kyMZBDZyaj2YY2JZkkNV8GPc1essBfeWQ7pLPcCWiWIsYd3xcNGpeICFxho4fyNNvGFu00dTCIpE3Qt02I2VJR40gNCETu1bBwChFHN41Rs2p8Kd7Povv6Dz96pPMXWMwXC5T9Xwcz5HjmiBCN+VmYtoaL/6jTKoVKt0677g6TS6V4sDOCo//eQwvCFjUJrhweQPLl07DSJuUxosM7NtHuWccXdcoOT4HRkw6W03qzIht/T5/3xbwk4tS6IGHadlM7mhiV39AueozvRmu/XUvH7ygi49cvIhSyWHK1Enk6+r5z3MH+c5dz7N5RzeXvXMR3/36JcyY2Uqt4pEqpHn40bf5yjf/TmWsSiGXxgEaUgYXTGpmuqWza3CEGfV13Hmsh7VNTXilCo96NZ7//bfRCzbVSlkJu5IfF4SanEIo/l4YKTq2gLp8ls2b9vGlW/9MseTyva9+gNNOWczoWFHaaj1Pjn8iH7dWo62xlSs+dRvzevrZG2nsHC5jmfKk70cBdqghIqkRjnsRjhZRdX103VKmG1VAF0rirK6MQyI+7WsTwbw4XyYzdJFiVUkMuy5EkpmLJzQJ3kkXCScwdk8F/6eON0pAi5o+UfOtALzSWhsogKuhK0qHIoirEKBuyANPXBAoO3xEMh2SyXk/WVND5QiQwWuScHRcyCfNAIYKEcqQodbS2hrFzX1ahHJBSHtuXLkas1NkgEfqAIahq7paDdtWuHWlIYRJ7aUqUVJ/Lk6ox9Y3z/fleCvSkvBZpFxDckyksCnJi0AVvpBUNvpBiKbqpYWQf95xXLnZHbchxHmPMIotbPG4DOXIkH5pobhWvupx1wRJj4DneZiWmSDpJcnTx7ZT8rWEIZYS6KIwJGcK6nWDzpzF6Z2N/P5AHw/c903yjQXK1Yp0fPhecvurVl0aGxuoOQ6uW8MwJkQzw7QJI43A8dFNk/r6AiLwuf+hF/nRbx6k0DPM2dPaWDdlMgfHK3znrZ2kC3kWzumkta0eP4o42j3IgX3dNDsOl82Yip1N85ej3YwHESKAYtnluqtP4caPn42GIJ3R2b3rGJ//zn955sWdXPHOeXzuink05wS/fXAXj64/Rkr36ShoNNcZFNIKwhZpdBdDjg74HBxwyWUM1k4XTKkX5KlhGRp6xkYQ8d+tHgsXdHDu8mYaCxZv7B7l+w8cYO5ki30DLgEmC6fYnLkwy7zOHNlClmLJxak6CLdGOqtjpCyqJZ+qMOitwMERnwOjPt1Fj6GqT9mV47IgAqFm7LZtJH0o8jmJ1KxYUhU0IRCRhjAETs1VYVR5YgvCcKId01cNj3pE6ELv9nFGKyXu/e3nscjw0JOPcsmV03lp09tUXBfNBF1Xdcp+lLDbcg02Gx4u0dPt0digc9qlNnqkMzxo8Pq/a1yxNsc/X+hjd0+FRsPjslMnc+KiZqLiKK9vG+LIKLieT1uDgRsISqWIlnTIY3trnDvb4NzpAb3jAZphMXfeZJ7aNMbUZpMRJ+KWR4dZNj2PqYXUZwUzp9bxvnMWkCrkue3+Pfzm729QyFp87+sXcsXla6iMVsk01bF7Tz+f/eLfeHvLQVqaCuhaRMn3mZROMyOVYsNokbQpeG99HT84epQH7/was+d1MjI2piYaahSsTsuGau0M1TgwXjg9p0Z9XZZiqUoE5DIpKtWaEnstlW+oEfgeuUyKA4eGueFjP+AD0ybz555hUmp8VA0Czp/SysrmPH4YcP+eo/QWPQJDUPRDSjUHXROYloXj1CTcUB2qRdwlriYgruclt5AYhyKR6CFq8KH6fNQNQpfUW0M30ISqwjYM1RWiJX/W87yEgGvoMiwZr5uKbKJS+ALfdRMclGlaqi9JnzAJiQn9Mwjk7USovEiMUIki1UMkNAzLVEgqCANN6cKRquc4jryB/Pme56PncrmbddVPrsULarzQKkuaTIdP7I4aEwwZef1SsfwwUOHAieChbAaTCekwUlgCle0gKaohyU6gxSEVJW6HiqcVyl5fX423AtWbHpfNc1x2RG6GUVK1G6fKg0D2A/jKkhl3FcvFQFrwZFdylIAdk+KnZPcNE493pNKe6lXIeaIS46WfICJn6FQclyWtDezvHcZsqmPZkplUKmWiUJoT/DAkn8/wwgtv8bcHX+TMdctJ2YaEGxoGUSBPaJpuydOHCKlUKtR8hzVLZ3PFBacwpAnue3Ubr/X0srQhz/WzOsjoUOsZZvxQP9UjvUyt+pzdUMf81kZeLJf5V3c/hmFQHq+RStv85Jb38pErVsuMi+/xt4fe5EOfeoDRsRJfu3Y5X/rQUo4dHuBLd7zO7r0DXL4yzaUr0iyfYjK33aarLcWMyRlmtVssm2JwwiyLFdOkhfmhbT6T8zoLOnTsQpaGOpuDgz41LK45u43NB4rc9cgh/vBcL7m0SX3e4uI1Dbx/bYEzF9dRyJlUKwF+zcc0NAyhUwkN9oxErD/q89ghj//sKvHU3gpv9tY4OFKj5IaEQuJjdC0kZesYhoZhSNRGpAwfofq8g3iEEjd0qnyPpmkyYa3Q5ohoAmchpBPHqwUM7CvTPzTKD79/DQvndfGb+/7NNR9fxr7uo/QPjxFEPoYZqcVSfnuDMMTzJVY7DCM2PVlhxkqNXEbWAIwMhYwfjrjh/CmcMa+OVdMzlN2Ih94c56FXetiwt8ywY7B2bpZ1C+qZPinF4oVtTG+3aWsp0NWs01uMmNcmD239Qw7j4y7TOgvsOOpxztIc247VmN6gc8WaBnI27DhY5DcP76dSHOF/Lp/LmacvZePmHn752xcYHCxx6slzSQlBfc7msktXMzxW47XXD2DZBlk7xZjjsq9cwwPW1BfYNDDEyrVL+PiHL6K7t0+OcSLQdRUMVpmkwAmouR52ysLzAjUOkhbumiM3cV1DldMJwkiac6IglGMvTaOprp7v3vkPmg73U7RNuqseJlDyXD42fyr1QvC/G3bgeT6Xzejg5d5hUkIkjZ2amCi0M5TWETeQJty/UCLPYyagRL+bisKrficiNcFQ1RJBXEuhEagNRVN6hxyHhwlKxTBMNfk4bo1UbLwomlgvhZDMQCGvwIo1qCkDkxw3aSojE/P9SGjpIgl7xxEE1AYlHa8KTa/F1eK6quGIoYsCy7bR84X8zRqaahlUob04Qq8EFcMyEYjkSqWrhz06DvalG7ra1USSbIxZVZqQVthUypIbkZCz6Rh+GIVhEtCRtlw5/4zizSxGpuiGQurHziy5ucXY4vjD1HUt+Xt0XVceaeXwUJhnXXW2R6FMdsZCUxCESUthTJ2UOhAJrsAw5UMm0KSvOwrU72AoQUqTvScCTKEhQtn2taS5jqf2HuPSC06Upy6lo0gzgcbUKe18+5Z7efixV5k7dwrzuiZT8yK8IMIwU0SBjxYFKikv2V0jY2MII+T8s1Zw4VkncHCwzJ2vbeFYzeOSqc2MRiEN+Qxpy2JA03itVuapgREGqx71tsWx3lFWr5jOfb/4MGuWdVIplhkfq/Hl7z7KzT95hrNP7OB7nz6BU5a08sKr+7np11tY2upyw2l5PNfljb1lesZdBsc9hoZr5BpShJqgVvOxTY3h8YBNhzzamyw+/o4m6qdOJ1NfIKqUeHqrQ13O4o8vDfOH5/rRNXjXmnquOr2Fc1c20tFs4XkhlZKHJeQNobem8Wafy5OHHB7cXeKZgw5bex2GHR8/0tB1yKZ1UilDzm8VfjuKIkItShA2vhpPyv+jMjHKLRPTmnVJY40Zb7qtJ+KnlTLUgqJhZQz69xTpPjrER685h4999Fy+/u0/8MEPL+PI4GG6B/vRTB/Plf3bvhcy0O1TLgY4Tsj4mEdxxKWuRTDaH1Ko12jtNEmnDd58pUxzKDhnZTOVWsSk5gynLmrmrMUNWIZgZ4/P4eEaERH1GYMpbWmMdIqG1mZydXWsXjKZY33jZLQAz4WMrVFzIprrLarVgMmtWVKazzN7XFZ1WBQyBqfOTrF6qskjG0d5+OUjXLjK5oOXn0ItsLj7Dy+y/pU9rFzRyaTJDYRewDsuWEJbc54nn9tN2XWIhJBdMEBGF5xSl+HezXuY3VHP0uWzKZU8zJShWkx1XM+juaGBX93zED+4/a9cftmZWCkbz5EaqVBTkCAMVRZNntxlNQNJN3k+m2XPvj5+cftfubCrg0e7BxBBSGAKvrpsJjt7h/j5gX5+/9Mv8OKBY3i9w5w6tYW3B0bJWQYe4AURQRSiG/JQ6bpuUtkQw1L9wEcLtcQRiiDBB2mC5KAtdE31Kgm1Pirgq1TdVYeImAgXqlpx2T9kyIyJOljHCJFAacxC7cITyPX4dhImofCJPhFZNxFHBuJ1LR7fR0GUSAkT+g4TvUyQBK01oanOFJng19Op9M1CCDSlEegKnx7G1EhNojpky5Uni+WFXHyJr2nqShNvPjH6Pe7ONZTaj4ThSquZEqxjfSKO5AcKCRAnH8NY0FEOAKGRjMpk+5stF3nPwziuj0O+Di1h5fuBj0DlWtSpIN7lhdAUoj5K3GZxUUu8+Qg124zUcmKbtkQuKLxBfB+0TFMy/6MANLl55SzBuOtxQnsTr+07wpz505g+rZ2a66hbnww4NjbW4UcRj/zpCdY//xa7hsqcumoh2UIBx3HRDeUHj98vdULQhUGpWKG+Ic/F71zLumXzeHLrfg70jrIwZfPnoSK7yzWOlKrUHA9bSM/74HCR6z54Ond8/30Ucia6CNm+o4/3f+LP/Pf5PXzxQ4v43IeWUpezeO3Nw3z7ni1cuzbFKfNyPLNljKNDPlNbU3S02NSnBZNbMnhewPCIg1MNeODFYQ6MalQ8jZMW1tHV1UplvAJulSfeGuPxXQE7+gLa8gb/7x1tXHVaK/OmFPB9KXgaylLYW4FX+uGfu2o8tHucDd01Do95eH6IbYJpaNim7G4WppzrynGVmhtrUSJW6kqwNExdmR9CIk1gKCOE0JUga2gTLhkhR1moZzO2QkZo6KagOuxzZMcgc+ZO5k+/vpGvfOP3LF6TJd8SMDg2jGFHbNlQYaQ3ws5GjBcdhfQJ8WqgBRH5go5tCmYtsdnyDDRMFiB0nnlghNVrWpnXZpONHELfp1TzMIhYM7+OsxbnaczbPLujwr/eGmfX0TIttk9nvbwV3ftCP2MDo8zMeNSqHq7n47lgpzTSNrhVl5zu8/AOh8lpl6wBG/ZW8DyN963OUq0G3PbvoyyfYXL5RSvomtbKI8/u4vd/Wk/X9FYWL51GdbjIipXtrFzWyVPP7WJ8vIqdsjCiiMHQZ14mzapcii/941mWdrWzeNlcSsVy0hve0FDPprf28INb7yVT9Xjs1bd513nrSGVT1Go1hKnJjUNIHTAMJ5hSvueoBU82Jf7PN+/itCDgSAS7xio0pC2+tnwWf3/7IM/5Ef/6zZdZs2Yui2ZN4ebfP8RFU9oY912GynI87yPzb37gK53TVFkN5IIp5FRE6BPNgjJAHLcHhkkdtaagrCSLcJToDVqiRWjqX6s1T+U/5EFdPrMcF0fQDSMZ7yeOVhVEDNXaoKmbknZc31GctEcBEpMbC0qmkbERDDWZMQxd3kiE0pA1OU4TuvyuxMl7vVAo3Ox7vvoXsvQpVIynMJQ7v24Y+L6r+PjSqWWaMrEtOVeSJy9brfQk6OK57kTJk6bEbYU+0Y+rtQXJvI9HWPFoKFTYDlk8pRxiMdpEQ/mnfTQRJW4r4zhRSlezxEDdCuL5oGmYKuGuXFy6UJ3kisWlKzwJ2v/tDlC+bl2TIlnsBrNMEw1pu0PTCbUwEfx9IgL1WhptiwIRW4bGufAdJ1KtufJ2pkRcQp+mlkY2v7GD/+lsYuObe/jFv15mRlcnXV2TcWvqQdUieetRaXtNyDGb7weMjY4zf+EMTl0xj3v+8SRrmxrYVvExhXw4I6DquriOz63fejdf+PTZOMUK6bTFPx/dyhWf/CvVSo2ff+lE3n16B1U3or93iK/fvYWPrDWYMynNK7tKzOq0WDItTc0J6B0OONTjcqC7xuHeKj29ZZ58u0ShpY73ntbCC9vLfOC0JvxA46mNPfzyySF29GqcsyjPR89o5F0nNNFcn6bqyvcpm7YYqoa81u3zr10VHtxVYUu/w2DFRQjIpkw1htLwoxDD0vE1eSjRTf2405eWVB9L2oAUJOO+aV3ZrE1VHBaqZwZNw1f1BLrSy2LMT6iySlIgldUA/TuLVN0ad/30Bvq7x/jzfx9j5boURbeC60b8+/f95BoE05cIylWP4hAc2xNxYEvAobc9ju0J6TkIR/Z6VEsaVi5ixxs+uzYUwdIp1aVYv7/CUVcaSepNyGuBdMvoOvPa01yyupHZk1Ls6A3484sjbDlQ5JXtY8wveFy8KMX4uE9Tk8CydUbHQ0YqkNJ8dFPQ0pjhrQNlBspw6bIMHa0WxbLLq3sc1s1P05qDH/7zGIu7LE5fPYOTVkxj9+ESt//qaXzP4/TT5uNXq0zvzHH2utm8sbmbfQf6qW/KY3gRb1arrG6qY019HV/6y+Msmd7GosWzZUWyoWOEJtd99sdc0pjjm6cs4r9v7+OBpzdywdknUGjIUVMtpXKMbErnkUKLSC3Wp2NSC7/842NsfeR5zp05jd/uP8Lc+hyfWT6Dn6zfTn9TPQ//9qu0TGqgv3+QadMn4fgBv314PR9dNos3+och0vA1qHl+UuqkG0IBXHXZ3qf0T6GIGRI7oid49ZiJrutKbBbxuiLrGwLPk4tvklg3FIBVuZwUhSPGrMS5kUitYfEzHSh9Wig78vG3hpgzGJf1RSrXJtTvrxsKz6LaFvXYYnxcNW6krMWJu5QoCRTGgromdPR0On2zYRqyd0ObaMnSdeP/1LmGgTy1RaEcJXiusvuqlq+4R93zvASuiHb8LonyEcvQXRTfVpDU2jAMsW1bxu7jopeY9ZJc22RUP4pHEmGUjB4mZpVGkheR4pChesz1JBAoBfT4n8ugo6biQoayMMsQjie9T7FNTvUARwpYmDdNMoaOSYitqw9XBxNVG6vsb7ZpkBKC8ZrDyR0tPLZtP2efuYpcISd3e9MkCjzcWo2u9mbW7+thfO9RGgpZNtcCrnrP6TQ15fFdV2JKQk+RfuUITApm8pptGWALg+tv/g1Lag7pXIpXx0tEngzElcoOKdPkvl9+hHe9YwHlkRKptM2tdzzLp25+hGVz6rn908tYOqueodEqqajGl365jdO6fJZPtdl6oMiahVka8yl2HCgREFGXjmhtELTWCxZ06vSXAjpmTuXTVy7m+W3jjAxXSNs63/t7N28ecDh7WRMff0crJ8/LY1sa5VqIRQC6xp6xkH/sLPP3bVXWHy7RX/ExTR0R+ZiWjpkyZbhKfXlNS0cz5DOmGQZ+IGfomgqmClOTqGo/xEoZyYFB8vHU5i+05PQWLwCmacoDQuyAMSSwLp4pC11gZ0yK/S5HDwxx8QWr+dJn3sNnv/ZbFp4CqaxgaLDMv37Tz9LTc7TPCSmPCXa8GrDrFY+hIz6eLxEVesrEcTRKYyG9BwJGBgKqFRcjm6ZzcQORGzBWctgz6rJ5NOSVIxV6a2DpOg2W0h8FTG1N8Y5lDayYkWNHb8Dbh8tYOpg6dE62KboafmiSzQr2DEQcHQMtDDBcHycIeXSXw4o2Hd8JmDktS11e5+ENo5w4L0dDNuLn/+5l3cI0rU0F3nn6QoqOz+33vMS27cdYd/IcbHwKeZv3vOcEjvZUeP31feimRUYYvDQ0zJK6HGe2NvLFPz/BtNYmVq5ZSC6T5is3/xpz/2EumTWFe7ft59pls9h2oJc7H3yO09csYHJnC+WqK00txGMdXwm5AZNaGnlh/Va+8a17+PqKedy+5xCz6/PcMH8qX3r6LdqWzuFvd32ewNAolcuYtkGt5nHyCYv457NvMto7xNrOVt7qGyYldELTIDguF6YrR6hk+2lJtipCNoXG7s4w8AEh+3/8QFZgC2kqiiUB+d+pQLQsfJhoMFTrDGrxlnin4w66cVW3ahaMoonuI+24A3I81teFwA9DoiBIMFKSHSjjGLEzNUpqMpQbNq7kUMHFSF0kPMdNqm+lZhOi53O5m0MlEMdJ8ST0EhciMQEHC5QeEI9+ksCeJpu44k4KXdl6JWBQBv+EKnGJHUwki7+8RwUqtq+rIvsgkMx+efsIsJTzRVOOLCO22Kp+hJiTlfSkq9BNzPQKwgDTNNWuGiW3kHi+GIaBCiFK11k8DovLa0zTTpg1GdMgpWkYaj7o+RFNKUOeahOcsjqt6hqW0Ki4HjObChRHK1RTNutOWkKxWFI+9gny8cyOKfzvfY+yG52//+4bzJ45mXKxJFEu4UTdry4ElmmrGk7pEKov5Pn8d37H9hfe5PqFs/ltby+BGyIsk9GREh2T6vnDr69hzdJJBF5I1Qn4xJf/ye33vMINl87m9s+sorM5h+NCQ1bn53/bS1Qd42On5RktByydV0/VCRkfq9GQFxQKJpmUvHIXsim2HPPwcy1cccECTNtg++4+nthU5M29JdYtrOMzF05izew8GDZBFGHrgiDSeKXf596tJR7aVeJYUSJz0imdTFqOJXVLKFe8HBtJ6KfcHAKlxUnXXlzBrGGYehJqjcVY3TQmrvgiNoyo51AXSpAEOetEgfKku0YooKYw5AFBAH17Swhd8LMffJT9+wZ58vWnWb2ukWLZ47E/jTBrucX0BWkGDoe8/nSFA1uqZHImzTNyNM8oUNeRJtNo0dCeoaE9S+PUDA2TM9S3Z2mcmsUth5hpk8iPGNo7TL7OxkFja3eVl3sdNvd5eGFEe8Eib5uMlwM62rKctyTD4q4CL+2p8M83x9hytMp4BY4MQ6aQYcWCBgr1abYe9djdH2ILwfbhiAbTp6sQMljWmNxo0tlssWlPhfNW17Gvp8Yjrw1y0YnNpNM27zhxOpZh8qu/b+LRx7dywurpTJvWQuD4XHLhctxA5+VXd2MYkDUMNoyNM7OQ4aLJLdx0/5NMbqjjhVe289+/P8Wnls7h73u6KXk+bw2Ocvnc6ZSLZW7965N0TW5iycIuQpWzIpygQ7S3TeKVV3fwoc/8lG8v6uJfvUNkLIPr5k7hk0+9xboL1nH3rTdQ9h05xhaGOukHmJbJysVz+O59j3LO5EY0DYYqnryBRhH+cWuV0HT8QI7RhOJLJeG+2Kk6MTQinlzFuJH4RhHzpmLIa3zC1xINRUuYevFCLuLD0HGaRLwGxFpIIpKr6U6cdP8/td2xuSl2fcVatRJf9AThJAkivnLJogC1mlDSvJCTF4EmXVixV9lQ83iJGzZkG54vcRlxS1aSzFUOCXnD0FXZiozCxwAyQ4G+IiV4hcpeJllREqEeKgyzhiLxKv5VEPiqHMpIrmexMKUJXc6sw7hRS54M4lrI+Cpo6IbsRVNF8ZZtJ/AyiT5XoDQVdgwimQMQQmoYkUqVakJ+UBpy5GFokNd10oYg0iLOnNrKSS05RryAsZpL1jbwNJlu1jSB6weEQkfXNLwwYlVzA0/sOsAF561R81BVSxlGhJpOZ8ck9h3u57or38HJaxcyODSkgkoRURSQTqVoKBQIohBX8jeoVau0NRf41o8f4E/3P8lPT1jK77qPcXTcwTYNhkdLzJ3dzp/u/ggzp9UT+QGHjw3z3uv+wLMvH+COL67iw+d08OJbAzy+oZvnN/Xz8MvdPLKhj3ktEQcHA8YrPpFu4I2XCCsViZj3fcbGahD4bD1UZtNQlo9d2MH2vQPc8+Aunto0xslzs3zygnZOmV+HkbaouZA1Ybzq8ewhh/u2lXn6kMNoLSCfNtEIiXSO6ycIiDTZp6LpMpOjCQ3LNogIkxFnoEwZuqknDZMSzKduwUJD06IkIKoLjcCLESQGoSfxNYZlSCspUlNBU/WjIXLx8HyslE5l2OfonmFOP30Rn/n4u/nBHfeTnzLC5Ck5Xn1qBDSNxScXGDjqsvHJCn37azRNSzNpcRNWRkeYUvsKVSV0FEQIpXNFQUjgBmrkJp01fXuGKA3WyNSnyOUsDE2jv+izbQS2DAb4YUR7GjIiohZoNOdMLlzVyKKpebb1+by8t8rKGSkuWNOMkcrSNrmRJVMsUlqV7b0uFQdCobOkw8SIPHbtGWfGtAJ1WY2qp3PS/Ax/fH6IfUeKHDw8wJHuIc5bPonz1s7g8Y1H+cV9G1gwq5nF89pwyz5nvGMpzfU5HntyMwiDgmWyaWSMNCbXzurk7ide5s3NO/nCqrk8cXCAw6UakWHiBCGvdQ9yzpTJdKVNvnv/U+w/OMC86VOY1NxINp8lbZoEXsSv//QEn/rWr/nB0i62lh2GXY8ruybx8afe5OprLuHWr32E/rGiGmXHIWd5YCyVK3R1TSKfzXLb357m2oXT2TI0iqlpRJp0dTmul0xlLMtSxp1AjsGFigfoEvceBJE6tKDMGQr9YUpGlyY0PMdVBU0KOqmpP6vriegtF+8gmdjEm4bv++owK/vpTcOQqXRXjcWOc63F3SSxdhsf8uNpTqRGsoZhJZpxXGseRVqi6WjHdcDHAr0Wj/nR0PO5/M1xLiK2q8aZiyiSJ3bDMJMxQHwtSuBUSduVmNjqlIMp1jMMNSbwfR/ixivl+IqUBUE6ulCNhyK5QShhQcEL5WwyCiZslIauqx4PbaKl8Di4Y6S62TW12cXdItKaGSZvkK7rGGpUF4v1gQKIhQoz4rgOadsiqwuyhk7N93jX7A7m5jO8cKyPi6e2kLNNdpdriAjJVFIQNS+KSFvyWrqyrYm3DhyjqbOZxQunUylX5Z/TTbRQ/ndnnLaM6dNaKY6PoesxnTgilTIZ7Bnmi9+9l+mdbcycMZVSpUpbcz2/vu8Jbr3jb/x01VIeGRvgreESzbksPQNjLFk4jfvuuZ7mxhx64LN52zEuuvo+xsar/PTTy+kdKPHVX73Ni5t6qI6NkAnKOOPjLJqkU/ME+/s8dvaFvLTL4dWDAZv7NEwtxKwUGSqGTJs7mTd6dRZOz/LytlG+98Ax8vkUN5zVzDtWNGDYJiUnJKVLCOFjeyvcs6XIS4crVH2fbErHsnT8UKZuddVAh1CgyQhMBTbUNA1DjUKJYveIPL1Fsf6lvghJn30ERjL6jHtpUJsKE9brID71xfWiEye5SB1OhCGdeEOHy1RLNT73PxdRsDP8+eH/sOzUAgN9PltfLrH8zCzVos76/4zSs79M+/x6WubVE4UBgesTqAyaMPWk2iAM4lY7ZRd1pW3DSBnUSh7lsQpOyaN+Ur3kLQlB2hIUayFv9Dhs7HYoOgHT6i3SGlQjnbY6g3euaKCjOcWfXx3jodcHaUwFNOV0tna7/Heni22ErOmA4ZJPOnLImAIMnSM9VWZ22pRKHnUpgRdGPPrmMFPtCq/tGOYf64+Rzxp87r3z6Rt1+fad67EMnXVrZ+BXSqxYNZ35i7p4/IlNVByfvGlzxHdwopD3TZnEipZGnj82xM6RMr5p0F/zqEUCXcDm4VEmpWwum9HJK5v2ce+Dz/H8q2/z3MvbefC/r3L3fY+z7YU3+P7yLvYGEd3VGusac9z44na+9Lkr+Nwn301P/wAICDWBIeT77AcBMWqiVqtxwuqFbNrbzca3dvHeuVN5c2CEjGFSCTwiTagIwvGzf0WhULpofHOIGwIty5SWcLUmxbwq2R2iJ8nzOLgXw2sl+Vwcp1mIxMKrq0ZEhIYWIll3vkSeaMpg5CuEU6yvaDEhHJJwZJwN0ZJO9yiZaoRqUqQpaGKgOF/Ebi0ixciSYzCEhl6oy98chhPRdJlzkBWG8RgnHu8kll1d9qebpplY3CKiRLiW/4lyqsRfwiBMAGCa2likZqD9n5J6LZ42xhwWMdEvHIUhlm3LQGAYp4flVVTiEFBjhihxi8WohMSNoPSYpNhRI2kli/vE4/S5UGEmIrn5mYaBrUFK7gl0NeRY11zPTZv2MG3pDP7z1h7O7WxjWtZmU6kivf66mGhR1ASm4shMT6d44XAvl71zLaVqBSuTTVDKQt1xA98jClwiZS20bRO34vHhz97G2Fu7uf+5N2muL3Da6vn869HX+OoP7uOnqxbzUmmcZ3oGqTdtugeLrF09h9/+6lpSpoFtmzy/fjfvu+F+2ltSfPdjS/jzU4f59/NHuHSxzWWLTVoy8nTTnDZYPMXk9KU5VkwxWDlF44KVWUZGPfYOhrRlA1YtbmPyotk0tdTx39eHeHpTkQP9Pjec18T150+msSGLZ5qYvk+oCZ457PCbt0Z5+XAVdI1cVo4LY8ttGHfOxDRTIfUoTeUmhNDQTS256RqpCTyEpskwFELdMlROSRPyGTNNNTcmkrcM1SOjaRAJWcFqmDG1OUK3dHWiCxUPTSfSQqy0he9H9O4Zo729jq/d+D6eevotjtZ2MH1uno3PjdHWadM5I8trj42xf8soLbNyNM3KEzkRfqREeVVlKoQmw13qUOarA0/gB2im/CJblo4WRlTHHGzbQrMEmfoUoRqPGIYgnzEpuyGbBn1ePVYj0iK66iXUc8yF+V31nLeoQE8x4FfPDPLitnEGRlyWdFicv6zAvJkNNOiyh6TiaOTTOscGfUSkU5+HSOg0Zg1ePehy8eIUp89P05Hz2bJvhEc29vGR86bS0pzju3e+woEjg1x49nzCWpW58yZz4glzeeL5nYwVa+TTNsdKVXYVa2wdq3CwVMUROuNhiBOBp4egchXHajUGHJezprSwtqlAtlTG7B2gbniU1SmTUya38FCxihUJOsKQ727ez0++cz0fvupsuvuHEtu+aapqCd2Uk5IoSMRuL/A4c+1y7njwJdqCgFmtdewaKZFNmbgROK6f9GBEqkMjzmHEho1k5B8FE9WvYsJ9FREqQO0EFzB2iMY9SUJBDIMkSmEma6XQdAW8Veue76n1S0umL3E/e6jKr6RRKUowTJrQFCJKadtKXwmTcipViKWEdhG3I2ok1AzdUIdv35dacjabvVk6p5DBkqQcPpBd6KrtSkOo6scIEWlJvF8WF2lqB5vo3U1m0orBfvwNgChCGEYcv0PXTSX66KohK1QfupnQajXiK5WQPm2hqfCjtKDpppkUpAihJxuQ1FXUzUaNvaRdTVf1jgamYUjBXddlAj2ZKWqJOwGhYRuGJOyaJl7oce2S2dz1+g5OPn8tv7/985iGxm2PvsSCfIYzJ7XyfP8gugIkhopeaQlByXVY2drI+j2HWbFyDlOmtFOt1KRd+rjTDspAIHR5Hc7bKa783E+pO9TPd9cuZ7Jh8OOHXuCRV7dw/7+f49szp7Mrcvlv3yCt6RRH+sZYd+Js7r7tCjIZk5QN/3l0Cx/6n3+yfF49X756Lrfeux29PMJXzquDIOS5PT5FT9CQNchnBUYUIAyLAIOmRotD3VWESLFydo72KfUsXDYTS4Pb/rKD9btrrJyR5ouXtrN8XiPD4x4GIXoYsrHH41ebRnlmfwU/glxGR1O3BnkjlF8027bkM2GKJI8klG9d6Dq6KZIvjqEYRREaIn5OtPhwqTYPXZb3GJYx0VUfKS1F1/D8AE3T1VhLIAkaEnSpxc6s4xwr0vUSEjgwdGic1Su6uPxdJ3HbnQ/SNLuGppls31hi2Sl5Du10ee2JIQotNs2z6wm9KEF0TIwLooSYILu3pTs7sXAKgSEMvKqHldKoDLq0TaqjXHLJNGYIA19moELJcoqCkFxKpxYI3uj12NTnkjIEXXkdt+rhBQHnrWllxfQCrx8o0zdW45SZKWZMsjCydTS0NVKtOpTHHfQwIog0ekcjCiZUai51Buwd8inWAibldYSuc+pMm7acxh2P9rB6ToGTl3fwsz9u4Y239vOOM+eRFj6d7Q2ce84S1m/Yx4Ejw2RzGapByJjnUw01PC3Ci6Q0UPMC3DAk0EBHUPJ83h4rcajq4qAT2iauZfF2zePfw6O8v6OV/v5R7ukZ5o93fJ7zz1lF78CwNMQIXTLsAl/x+pR+qGkIYeL5AYHvU1efZ+Xi2Xzzd49y2dQ2RhyXYVetBZpBkPSpm6ocSqUhiRLDhqYOz/E4ScJWRRKkRlUrCNW/JEf8ymkVxSlvVXOhpiFBHFwW2nHittJuj8OyT5RQ6QjVd6RFWnIYjjMrYRAma5pQeadkbK9iC6apIIu+lwyaDLVpBmGEQOpAYRihZ7O5m13XS/IcUSwAKYyDpsmYfNwxoQldOa60ZEdEUxbRmNcSh2KIFNZDBl50Q+EAVKYktvzGVN7Y/hsRKZ7LhA0utlKGqgckDv4JITeUcKLdSvGopI02vqXElbrS2hugIaSgrjzUcapYjhLkidfzvcTKZumCjIjIGyYBAad1tlErVnjeqXH39z/JuFtlxYr5rJw/g9sefhmt7PLeqe1sHBuTADI1f4zUXLHBtsm6Hvsdj3eetRrHl68RwsSLreu6AkSGtBQa+MTNv2bwte18YE4X3zx0kDYrzbsnNVMbGOKq6e3sDgP+3T9MkzA5cHSI886Yz69/9gEsw8cWPn/5+xtcc+ODnL9uKv973Ry+fPtm1rS5XHdqHU9srXFg2OeMBTbLOnVaczaFtCCVNgiETlPBYPfeIQbdFJddOI9ntxU5c0UjB4+M89W732agBJ88r5XL17WgmSaeLyhYgn1DHr96c5SH91cpB4JCSj5TmqGpg0qIUCgaTfGBdFOmZEM1S1YMTMWcmuhYiDcWogjDkjC7QOFINUNuEjGDKFLEV7lhyM8jdvTEY8wgiBKMhCbi8UOUuG0CdXvUDUG132XoWJFLL11L++Qm7n/0CRauybJ/l0MUhrR3ZXjhnyOUxlwmL2pIbsmaSk5HoSZN1crlImtWJclBKDOA5wUKuy1nz6m0RXHQYebUDkZKZXJtaSKVJo7iumdlezF1QUqH8VrI+iM19o56tOVNOjOCsVJAcyHFO5cUiAydO58Z4O1DZaZmfVqbMvSXffpK0DPsEUYaFccll5Wfl6mB64Zs6fM5eYbFaDVge1/IoqlZTp9nc/sj3cxuz/KBc2fxs79t5Ynnd3HeafNoLKSpq89x0YUreH3LEXbv70OzTAKlK1ZdD1P1f+tCxwsCfOWsc4EgEoz7EQO+z7Gyw5GaQ02Hby2ayzO7D/OMH/DvX3+dZcu6GBgak2MiNISmJh4BKlCqUchncWo1tYgaaJpOpVph7pyp2Jk0v/jHs3x8+Ry2D4+SEoJQgKOqX4PAl7fT5PQeJmjzmC8lT+iqlCmujVX/Dk3ROJiQAuLDgxx1qcN5KDvLTdNKKiWkDBAk+gRxeZ0QspJXE0rf9VXhlTR/mLqpNBaVcVLJ9JjfF2vGsfAeZ//iTvRIVWbETfXJNCqKZBI9LlTShcQcOzUX07YUqEu6C1zPTfzM4XEpRpm61CbGQ3EZiybbATUVBguCMCmO13Ud3ZRiZDyLDkJffaE0dfKONyctsY6JGGKmdAnp9hIJbiS28MpRmLKlaRPVubFVzdBNVcvoqT+D4mSh3FyKl6X6UMJII6UL0roEsmUtnfM62rj19a18+6ZrWLZkBuVimWqtxvQ50zlv3Urufmo9es1lWj7D9lIVQ0AQyRYvAhj3fE5paeSZ3Yc575zVpDNpgsDH0G1J1VSOiiAMaCjk+N6d/+S5fzzLZ5fM5Y6ePoQTsLNaZa/rMzOX57VSmWf7hmm0LI4OjPLui1byix9fTuT7ZNImt//qRf7f1x7jqovm8MNPL+UzP3idRfkxPnBSHX97uUh9JuDsRWlGRjx27q/SM1RleNyht69Gteyydd8YbqqOd5w5nT4nxd59QxztKfLDvx/hpIX1fOm90+maUmB4vEouZVAq+/xtW4nfbh6j34GUFqGLiMiQoyctkjWusfEhikIs21ROKVl+own5JUDTEIYS81BCuCKFImRxWKCCVtLfPuFMiQt/fAXY1IR00AhdwgyFNtEuF+dDNF0j8ORibNlGQmQQujR06JZGqdfFrwZcdeU6hoeK7B7cRPu0LH3HAqbOStN/xGfT86M0dmao78hJZpb6maGyU4aBTKabKZ0glO4+Q0W45aFuQpuUoV6d8Z4SJ52wiJ37jpBrS0vBNorwA3lTMgzZkOirbgdNkxbzvmrACwer9FcjZtWZ5FMaJQdWzyxwweoWtnZH3PfSMJsPFMkJk/lz25gyvYX65gLbjxYJEeRNQbkGkxptnttXY0aDwaRGk4asyVNvl+hoTXPJ8gw/fLCbGR15Pvm+pdz3yG7+8fBmTjl5Ju3NaUxDcMnFK3lr81He3n6UfC6NF8hDA0KGQGUuSxoJPEJ5+hcakS5ImRal0GdSxuZL82bx8w3bGGhv4F+/+iqT2hsYGBiSm3UcOla1FBK7ZJBK2by0fguT21vV4TNCiyJMU2e8VGLdyUvZuOsor7+1iysWTOPVo/2kDItqGOAEMQZdmjdiiGycUI8/ryCQh2FTNyb4VsmBJFC30AjLtJJRplAu1rhONwxCLMtSnejGhJlKjZRiKcBQJiPZrjhhuYrD0DE9N1IeRrSJjc7zJPNNU0FxiZMncZeF6rZjWKZaq6OJqnDkZqZns9mbNbWbEGkTtrK4vzfWIyLFjFcoYN3Q4yVaObbiXl7UbqXQ7Ap5IoFcExTKuOcjZg1ZahOIHQeGqSfzNhmcUWKQ2qUNQzpvYs9zUv6i5nRhXMuY+Pj15KQZ37bi8VoQSieFKawk0BgHDqMoJGUbZHWNOsOkFnq8u6uDh3cdYMbJS/jKp97H8Og4qZQFmk6pWKStvZUTFs7kj/95gdkNdewYK6Ef12ui6xqOFjK3sUBlpIjXkGPVktlUKi4I6XGXgqpLXV2OZ17czLe/8we+e+Ii7h8eZMQJcSJN0kTLNbZWKwx7PnWpFEd6h7nk4tXcfuvlErkh4Ac/f4EvfPspPvyuWfziy2v42Z+2cmzvYT5/cTN7u2tMaxLMmZrjcI9HqVSlkNOoK2i0NadICajUPIZFlovPm0Vnx2Te2NrNL/9zkAMDPjde3Ma7T5tMxY9wAp+CbbCl1+NnG8d4c8DFFhEpW6ApJx8gF2XFTzN0aSuUoVB5KAjV7SKMVJxcPWnxFySMRcCICXFbutknkPqRHDWJWMOLoYlxQEwFVOMTf6QpPpAXJOJnIs6rxL/cRORnWOz1iLyAKy8/haHRPob8fdi2zXCvR0uHxcYny5TGfSbPq0cYMgQXBSrBHE5YhJNDkjpVBkpoDX3QBYl/P9eQZuDIGG3ZBhYtms3GLdtp6qinVnWT5zimJOimloQdDUMQhGCKCFPX2N5fZeOQR8rUmVNnEAiLdCHLBcsa6GiyeHDjCFXX59zFOWZPa6KxMcvcaXW8ddhh95EyumVRsCJ6yjBUqtGS0bFtm5ktOut3u7S35jhxmsb3/3mU957RyeXnzeHB5w/ylwdfZ9WyTrqmN0HN5V0XL2Pbrn42bztGJmvJG6GiHgNYpiFb8iJJpIg0gYVg2K2ypLmOj02bytfXv0nHCYt44I4vIWyd8dJEpa7i46oqbYvQ98llbO769X/4ny/czoKZnaxaMZ9KuSbzE6Gc+zu+yznrVnPXgy+Q91xWTmlha98Imm7ghJHqSY+pBSJJaR9/Mpc6BwlCJ3ZZxb1HxLEHTfZ9xOuaXIP/b1AxHpvGBIpANQZKdBMJ/TeuBI8D0bFrVo/Hw2p0rx9XaBV3uBuqByQO2wolP8Ttr5om5Q3DjEV5uX6apoFITmpoiWXMVG4hP/Bl8joG26jWKk3dLkzTxA98lffQk95fkZQ/yci9YZgygSlikqSRwBQN05RvoC+7dwNFipRMfdl2GPiSjIumSQa9Jt+c+A1CTLi4NE3DcRw5O0SeBDRdp1aVlbrSIhwks+e44lKO0Tx1uwLLshUPSyOtaVgaVNwaq9qacD2PjZUq37rxAxTLZYSIVC1viKbpGIHPvf9+kcWFPEta6vCJyNo62bSaLUZS5H1tcJSljY08/fgGPNeV7gbfw9QtDCHQdZMgDEmnLOpyWfpcl+6KS6T6SYo1H2wDA0nqPNg7yEUXrOT2H34Qv+Zj6Rq33vEsX/n+k3zmQ/P5/qeWceDIKJu3HOMjJxgMDAdMa9Bpr9Opjjh01EV0TdKxNY3imMG2PRVqgcagI1i3pp3GVMRfHtnGDx84wLLpaW69pouVC5sZKnmkLQ2/4nHflhI/eGWUoUpITo/QDIm0DoNIPoRCx3dla6OVtlEBGOlfN8UEyFK5WHQhiagIebWWn42lEBFqWhqGpCxTcsp0nVCB8azjRq+6QktoQm7gcad0bADRIgmR09XvqymsRIw7jQJ1aJLfcdyiSzafJmMZVL1x8o0pqtWApnab0aGQY3urZOtN7LwkVWuGwEipZzoKCLxggpTqS0qr7wZydBFCOmehp2WlaYTOcHeVnl2DvO995/DWG/tJ19kEro+hqg1sy5I6jS6x4oYhT46eJ2sCgiDCCwOaCiZOGPHbLWP8eNM4YwHko4jBSsipi1v47cfn4WJyxU928s+ndlEbGyefz3DZifXUN1vsGfZ4fJdDuQqHxgWeH9DbN8a+o0VOm2fz95cGyGVtPrKuwE//tIklXVnu/dZ51GdtLrvm9/z3ybexUmD4ZX738/dy/tkLGBwsYukKUSJ0dE1uuBXHTcjfKSEYClxObW3mQ5Mn8+mX3+C0d53OH376WYpOlVKpjGXZ0lFp6Bi6hq7J2X6tUqGukOU/j77CT396P989dQHfu+0v7Nl9mEI+SxQGUiMLApyqh50S/OzW/8cf93RjagaLmwoYgU/eMtGFhqO6PsIgUpWzSmhWI1XDMmN+eILOCUNJGHYdLxlb+X5ArVpTXTNaYiKSafBI2XYDXMeRt6gwUrkyI3FvuZ4nrbmqPC/wJLw28IPkABzF4csgkKM7faIDPoh1YuXYkjdw2ask22fl328qa7uu60mgPAgCReNV45w4KCJP39K1YCqUb1zfmrRcKQ0kUeqZGBElVy31QmI+jC4kjM6yzOQm4bsSi+KrL7IeI8zVdU+o3vXA85PgjKZJKqT8vd2EaBlb0uINIb6ySViYpnZWQwmiQraZqZBhEmhEI0TD81wMU8cSgqxuYGoC2xK8e3YH33/5bb74hQ9yximLKZaLiSfUDwJamxp5/Jm3+NUv/86nF0/n7YFRzm5t4I2RcWquTJZquoaOxmjgs7Spnt2H+pg8u5M5Mzuo1RwCvybbwRTKek5XO4+8ug1rcBTLtjlacxQNWL7eSOgMDoxwyTtX8fMffxAcj1RK8N0fPcL3fv4S37/xJD77vjnsOjDOv547wo6dvQyVQl7eU2P9zjKbjnn0jUmBddgxGC2FpIVPa6tONmNxdCRiwew8P/nnMe55opurTm7gExd1QATlik/B1Nje73LHxnFeP1ajkDExTJKCHWlL1dGF1Dmi+CYdxVncSGWQNMWYkmMsISTyP84dhap+dMJKqdrfYuC/OmkZhqGMFqCr056u60RBhO/KEh7D1DFSurTuulIcRN10dRW28j1pv4zUl8n3VE1yAGM9FerzGS6+YDl7enfg6GWGel0ydSZHd3kc3laheXqGTGMK15Xo78BTow4VhgzVQQKh4VRcDF2QypkQCY7tGmVgXwln0KU2WCMTuqxesZhCrom///txupZPknb2KCIKYm1HWehNIyFmC/SJmXmEPKREETnboLsW8kpvFQONec0WfgiarnHR6kZSls5PHulhz5EitVqFf7w6RkEPWNYRkUsZmCLC8TVyZkhdRufoaEDV11jUlebprS4Xr0jz9NtjzJzazPz507nknHm8vOkYP/n1C8yZNYnFCyYRVku868JVHDhS4bWNu8llU3ier7q45QQgCCNMTWfYc3hnRxunNtXzuVc38f+uu4xvfP5KBseKhKGb5MVCOQMi8LzkQIkK9rY11/PUxp3MFCEzbIM7nnqDqy5aRyjkQqmpLJvreXTNmExbUz23/fkJ3r9sFnuHRhlzZQ4qVBBDTSXB4zFj3O0hbb0K1BqGMldx/PoYqR6iCCzbJAxQorkskjIsPVkD4wuqrpxekgGmCMCqs0n2OYVJqDEmB8d6RhSP6MOJWu4wCNQ6rycW3zDWiAOUduwnANmY/yWNSWHy9+j5fP7mMMGTR0xsAYozbyjUQzSRf4zn0roKv4RqN4zfRBQKXhz3ZU+uc7r2f8QoyzYTEnBiXdSOK05XNxb1PUvmiGgC27LkAqNGZhoyTRy/uBg3EIbyTY7hYfE4IlI/V8TNXWrup6v5qR5pFAxB3tTwtZDL5kxl/d4j2HO7uOUrH2J0vKg+APmGG0KjOl7h6i/cwedmTObZ3iHu7RsirDpcO3MaW0tVAi2SIai4uMY2mKbpvDk4wjvPXkbVcdBNW52spRMun80xNlbi5Q27WNJWz45KDUsTOJ5PNmMzNlrk/LOWcMePrkILPVJGxHd+8ji3/uIFbv30KppyJp/5yUbuf3wvfT3DTGkysQ3B7M4cnW0Z6jOC7rLGs/tkL3h9KqK9CbJ5m55Bl02DgvueH2a46PGTj87g1EWNlHwDS8j36oEdZX6/aYyaZpBL6xKfoGzWsZso5q2FqpM9NjSggpQSk07C5jEM+d/YlpVUHSekU2XQsA1JOTBTBr4vLd+mJcOwqZQ5wQpDo1JzEJEcDTg1n95DwwwcGsfAwMoZCtFjJFRSX1GddV2eTOUtV5MNkS4MHx5n8uR6TjttAbt7dqNnHIpDEdmCyfZXK5RGfJpnFOQYV7VWxqKkpubHWgLSk/+z0iaVIZfd64+xYv4sli7rwMxXed9107j4PV1MnpTltlsepX1hI6mcTRjIbIFlGkmtq9B1qrUajuPhuQ61apWqU5OBuEiSp8Mowo880pYUkDf2VdkxWKM9FdFeMCnWQpbPKHDakhYe3TTKvc8Pct58m/esrsO0bObMaqA1q1FvBRwb8cmmLNl5f6DM2jlZKr5gZnua3T1VUvksp504A8u2OH/tVF7fMcRtdz/HgpltLFk2nbDmcvE7l7Nn7xDbd3VTX5/DdV0s08DzfFKmSTkKOK21kZWpNN/Yuosf3nQt13/4nRzp6UvGlmEYyPBl6EqtS4/hrxpWysKreTQ2ZFm5cA6f/d0jfGhmOzsO9rB1cJRLzjuZ8VIVQ9eTgivHcThh9RKO9AyxZctuZrQ0sb9YkbRutS7GoFdDkXtJkuMTsMFk3KQ0CMMw5bOg+osmFnoFoTXl5iE0Dc2QC3d8s5H28lBpeLGt1kgyHCLuX/J9xQxUbEGlyUShRsq2k7UuDm3HUQftuPwcUYgQxoS5Sukuhq4nmrPQNfRMNnuzQEsW4bhcSdPERFpC2UljOq/nyQ8pbtqK7boxP0WgSQuYkCKOqexqCa1Xvcm6IU/ZgcpY+LGXOq6VVHO6JMKvKcuuoatKRicphRK6XKgMxWyJFCvAVMUqsbAllNjvq987UByaeNanqU3ENgyyIiInpNg0tyHLjEyGe/d38/uffZZc3sZ1a6rkRVooG3NZ/t9N99DU209XJsPfaj4b/vC/PLH3MM9s2sOnF0xl63iZshsgiEgbBr3VKisacmzYdZTla+bR2tqQdCEYuiUX1SikMZ/nTw+vZ3l9gc0jZUwhsFM2Q8PjnLhyOvfcdR0ickmlNP73R4/y/V++xHc+sYx9h0b4xf07OGOmxg2nF7h0TR2rZuVYOtVkwRSbzuYMJ85OMzDm0FKX4YbzJzNn6SwCBJpb5fEDEU/tcjl1Xo5vvb+L1oY0lSAibwr6qhG3vzzEKwM+WVuODSJDoBvSCqgJkmpiw9QUIsGUHnVdMsWEriEMkyAKEvKxLBWKQ4IicaGksrZKnUv9S9fl4cN1PQzdwDQt3IrHyMAYpeEqmi+ojXv0HxylKdNI/7FBho4VGT08xiVnncOJy5by8rObQINsYypJrvvKEhlrdEIzlaAvrZFEMHSwRGtrnpNOms1LL22ndZpOqeSiaRGbX6wQRVA/JS1Pf7qYCIjF7ZwxxM4LwQ+xsxYjR8oc2tjLt795BTd942J+fudD7Hh9mP3bSjRM9mmfmqKQS/HG+n7qp2Rxq27Ssa3pmkyuI5jZOZ2uzql0NrazavEK5nfNpSFfR2tjC60NzTTXNdLRPBnLthkaG8IMI46M1HjpYBlDgwWNJi4aaeFz4Yp6fDT+8sqwTLrndI52Fyl64PsRxbLLcCnAFhqlSNCQFczsyJBL65TdiMc3DrFrbx/DA8NMbhB89F1LeHlbPz/5zXqWzm9n4ZIOglKRCy9exabN3Wx++wCFuiyO6yV02Lyl866mBn506Ai3ffN63v+e0+gdHEmIAjGWQ+hK1PZDDMNC0ww185cmCdf36JrWxqSmBr76x8f53xPn8ounXqe1rYHVqxZIrJCmDsCRRl02x+PPvoHbPUTBttg5XJSHXDVemsgfSWFdN8REc2VCGVfuJnW4DkI/6VmK1zdD5TViV5SucCmJ/qGCiBI1JYVzoigRv+ObgyTmKs1ZBRZjaq9QcQvf8yZ+N6WZSEOHEsxVsd//aUckRsqLJIAoe0IC9Lq6upvlDE1uM6ZhJIHB+AWaChMR/yXSWiZPmZZpqRchRetk7KUChkLthL7ny/h+NBGSEUqol9qLfEGGuhHE8MUwkswVecWUaBM/8BXPSE8oqfFVy3Vd9ftNVFJqCs8sxESfSEwIjsuzhFAtY0JDILA1QUqDrGlgmoLL53Ty89e2ct0n38t5Zy1jfLyoWtXk5tHUUOBv/3yR++57hC/O7+JLb+3jpzdfzYJls7ho3Uq29g7wjxff5v/NmsquahVfASc9H+y0Rdapcbjmcf5Zq3C9AKEbygstqMvn+Mu/X2LrGztY2zmJ10aK6BFUqw7TpzVzz8+vJZMS2EbI9297jG/+5EVuumY+h7uLvLHxGD/5QDPnLMpyrN9hy74y/YMVfC8gMk3SRsSLO8d5ZKvHJ981nTBTh5XNUx4b54cP9bG5Bz73zjY+dNZk3CDEDwKyusbT+0r8fOMwA65GXVrgB/LuGl+J4xNQEEl9wldXcl3NusPY/q2otvL2IA8RklKgYdoGgR9RKdaolD3Kw1V6948y2l1mtLvI0OERaiMBOb2e0Z4iR3b2UqCe808+m5XzltF/cIDDu47w8Q99kB/ffBPFkRI9R3q57+7vc90N1zA20s9f/vEwRmiRbUmp0YFMg8tnOO6XkRqOxLsLAjdi+FCR5tY8p6ybRo0BNDtkvOji1WDnazVyTRb5thRu1Uc3pEssjskmyfY4rCo0nDGXva9286NbP8LHP3Eu1eEiUybVcfUHV7LtrT62vD3K9LmCfKPBoR1Vyn5IKmegheC7yj0j4PMf/h+uvPC9rF26mtNPOI1VS5aybN4CTlm1ljNOOoXTVq1l3ZqTOeek01k+dxnbdu9k0cwltNc10NN9lPUHi+wdKLGwSWCGAY4bsm5OlqmtGX7z4jBvHiiztDPFnI48bi2ksWBRcQIasjqVWkBD3gDHJVuXo+J4HO6psrrD55W3+/nDo/sp1Ry+cvUKth2q8IM7X2Dl4knMXziFoFLjwnMW8+obxzhwqB/btkibJjUCTmpqoG94hMXnn8yXbriMo7396LpkU9lWCj/0lf1dLoT5fIHieAnTsmQttLKBm4ZFsVTk5FVz2ds/xj9f2MQXV87hG/96kXeevoLGlgZZRWsZNBcKfPl/f83L/13P+5fO4bEjPZT9gGqgrLtRzFPTE0eq5/mJc4p45IqGUFwpQkU4CEjcolJfCxNnVaSBbiheXySzGJKVp6teponnJxmXKUNGnFU63knle16Sbgjjmg1D3aoVdFbWh09owpqqMdDV4dnQjWTa5Lm+quVVLth0On2zpjDBcVAl8P0k/2DGyUXFS/GUnUvTJtR/z5sInMQY7Xgh8H2Z94gFSCkoBclsXJaqqC6Q0MdzXRVokSe+WBBN5ocidiJwXKoykDN0dSOKbWaBenN1RcI0kqSmiu4H4XEk3kiNSTQsIV1XBcvA8V3Om9nOpsMDVKa08YNvXsPoWOn/nB7sVIq+IwNc/eU7+dr86fxxdw+rLzyR/7n6XHoHRggEXHzeWrpLVf705EY+OquDN8fGmFuX49T6AtNzNtObcvz9zX288+zV5AsZfM/HdT0mTWri2We38Olb7uFLi+bxxMAgveUagReQy9j86Z5PMWVKC2bkc9fvnufGbz7GN65fRhCEPL9+Hz++vAld03ly4xhVN6RrsklznYVl62RSsH5nmV0DAssImdsCaVNjcLDI53+/n0qkc8vlUzlxXh0jZQcTqbfctWGIf+wsS4NBRsdzAwxDw7It5dJQlaSalrznURQpxIMv9R1luzUtE19pBGEQYdryJlEZdTi0pYfBgyOIskHKz9CcbuGUlSfwkQ+8jzkzZtFYV8/113yI6656PyuXLqGQz/E/H7+Wd553BitXLuaSi87jkovO4fR1J3LwyCHy+RzTprQzVqyw+c03yWRS9PUNM+qP0zKtSd4GIul8cj1X1TXb8veLu0SI8KshQ4fHaG7OceaZc6lqg4xVK1SqLrWyxr7NDpkGk2yTLRETrp9UnIbJwUt+iUM/JJOx2PHyYS45fxXf/uZ76Tt0BE24TGrzaaqHljaN//z9MONFl85FKQa7A/q6XeomZ+V4TAiqtSonLl7Fe869iHKlKA0CWkSlUlbjq4AwCPB8SUXwQo9UKsWZJ6zjpGWrWTJvASsXLeOMVScQ2C1sOFyh0axRZ7hUfYO5k9OcuSDPkztKPLStQqWsTDEq7DhelYJv4PhMb0tLurYFG/dX+MDaOpZNsVkxzeKR17r5z2sD3HT1fPpGfW79xXpOXN7JnFkt6FHI+ect48lnt9E7MEYumybUdTosG9N1iaa0ce5JSxgvlYiiUI6wXScZBQZBRENdnn/85yWu+9ztnHrifKZNaaNWrU0QLoTAcR0uXLeC+1/eiTs4zOq6LL96dhMfvPQ0dMsmY1p86ut3sumJV/l/qxbzx10HGapJAq8XSb3MMEzCKMR13SSzJgOLgazoVmsQym01YQrR1VjLSIToBJsuuUsJm0pPwoPKpRUEspHV0BGGORGDiLFNmrwZyXVNS4R90zATa6+marQN00pMRXEPUtIRoinjUpJRmWhnlWupon8QYcTArEAp+aZlHre7yV3VUEntIECJ66H0rRsWUSRHU4bKdejqCigEeK5HqBhStqlmb+rvjn+eUMVPQtMQAQgFXYwZMpZl4npucrWM/Imie9/3Vf3thG9cM9Q1TY3OIiIpfsahM3Vzkbde+XMsy5LuCE3WOppmSBqd0I+YWp+lPZvinqEx/vqjT+L7rmqwM3B8H0RIXhd88tu/5TTboFjzOdiQ5w+fuZzRskcmnSIIQkbHS3z9xisoWAa33/cYn1w8Dc8JuOXtQ0zOWfzwhFm0hx73P7aBL3/sEirlGk3Njby9+QAf/todfGpGFz2hw5ZiiQw6RafK7+/6KLNnNELg8Jd/beSTX32YT314BSvmN/KVH67n9qta6R2J2HdkjLVL62jIC/oGHcZdnWrN55k3x0g353n/OZP41b8P0JLXOdg7zpcf6Kc5K/jOldNpKNiMlmoUshaHB11++foAB8dC6vOWClEFGLa8KXpx8lY9dBJPLQVMSxiKz6NNVGQm3R06pmVip0z6D43Qu2+UvJXh5GUn8I53nMaiRfOwTYtsLksmbREEIafoJ0qBMPApl8ssXbKI1atWUhwfR9d0XD/g5Zff4MmnXuDZ5zcwNDKCU3MQuqDqugggl81QX19HLfAYO1bErrew0qZibMmxgu970mIc+EiCviybEqaQz13oS1t2WQbKyqMBgRuRzhsKOx6oLIpE+0cRmLY10TKnw1hviagGV3/oDLxqCahwtHsHo2MjNORtXt/Yww3XzqRzSjPf/+kbdB8Yo7mrkDDaMCGjZbn0zItxPAdDWKpFMyKVzshO7tjZ6LmqllQQadIRVPKq6LrGnFmziUI4cdUKhseKDA0Move/hL//JYZGdFoa0tx1zUy+958j/HHjCCfPznL5CfXMmTeJbMbGcxyeX7+X0ZLL/CaToh8ShBq9fSVqms3giMfnz6nnpX01vnLH63zhynm4ns97rr2PR/98DSet6qJBuNxz2xW8+5rf4XoBpgYHPZerJrdyy+Mv8+aFJzNnwVRKpaq81Qqd0HcJI2hqbODfT7/BAz//C+enDT7x9d/w2O++ip0x8NyAIHDV2DMgTOvc+4NP856PfZcPN+TJ7jzCt275A1//4tVceeMPCLbt44trFvK7w8fod1wMoeMGsm8o0pHZFRWuk2FXQ90GBGiBGjdJk4unshxxEDoRv4WGwQRqSeJIDOUw9QlUuBZNORSVDuL7AZYa2QuVM9E0PdHYDMOQBwWl+YbKvh4EIUJxCwPfT3IeE4acCNu2cZXLTCOa0FTUd1kRmVTtr4Yh8xsohLlGrVZV1NpIlauHCEtPTmO1WlXaJQ0D1/NkNF8lL2MefRQzrAxByrCpOVUilR1xnRqpVBoiWTAld1uBr5wMxGEuJKoEEzWbM9XcUfrF41lyEAQEYSjT8oGHrt7IMAjwAtk5LqtKg6TKVlOgMF3XCZJTgBzD2IYga2ikdIl0vmhOJ3e99DY3XHcxixZMY2BgSDFn5GiiLp/lrt8+xt439vKZE+fx0dd28ovbPkU6l2G0WMQ0LTzPAdfFT1u0tjQxXqsxxTK5+pXd3PLND7F5xzE+8eCzfGfNDH76yHquuPAk2ifVc/hAH++/8cdc2d5CW0OaOw4epVEzOTIyys++dyXrTphGUBvjufV7ue6z93PFRQu56YblvOczT/KJ0zKUxiOO9lc5YWGWSjng4NEqbrmGpkfU3IA57SnOOGcqL+92qLM1Dg56fPIPg5w8O83nL51KEEWMFavUZSw2HKvx6zfGKDoheVsQqMS8bmj4yrqdzsiHz0pZcqbKRFeMEBqBI0VCEUn3k+dJAkC2kKY8WGPvq93M6JzOBz52OSeeuJzp06agAeVSBQjxXJehclnawoOSJDxHEelUirGxMoVCRDpl8fhTz/Gr3z7AqxveorUty4IFnVy2ciFNLRmyWRPbkk4834W9e0Z46cVdbH1tH+lMhmxrlra5dZgpk7DqKRREiIZOGPhoyGc7vtWWyhXQAtmx4ISURmWLm2HrchbuRklBWxSQOMtUrox0zqZn+zBTpzayaP50+vv2ceDwGzhuDUOkqDkh55/dhEDQ3DaJ7375DP771E4efHQfTdMLhCFUnAqnLF3LrGnTGCuOKW+/DlokSa2KDuG4jqyC9mPBOVT9N7KmtFKVUE/Hl0iM9vbJMO2DlFrmU9pyP+WxcexMA1+/oJO8pfGv1wY5b2mB0+tDKprGpNnTsKjxyDMHqbgV5nWY6IT0DtVYMidNEAme2F5j3WybKQ0+P7xvJ594z0xGxmtcdvXvefaf1zG3I8+82c3c/ZOr+MD1vyGfTtPnOgxEcFY2z5dv+wP/vuvrGIaJ49Qm0OlC4v7rUxbbnJCPzZ/Etk0H+fT//o4//fTT9JWHsFIpdXKPGBuv0jWzlVs+fyXX33gbd61bws3PvMYFG7ZwmmGwcHEXP997mLITYekGNc/H0nXcwMcSGlnbouYFVOOaingtEzqaup0QytmGaemqzMlPJiaS6qtAiTHxFiZYW4paoAuNAKk7C6EjIi1x1CW2YWV0knEHTZE89ORW47mewt0EBOjSxEIEfii1R81I+kdq1epx/SJCaSSyo911aqr8TzIGA9+XSXSO0wMSdK9iQUl4IUljVnxjiFTgROhaskDHY6fYHRCoNKxQBJqYOKmhoRlasmNrqlfRME3FwYqSOlHtOIDj8cRgYQg1U1T24ri7QY3WdF0kHeeGYchTl5KEfEWUjO1vE0EbjZxpYEYaNdfh9Bnt9PSP0F2f49abrmG8WAJDV/3wAtM0OHKoh6s/cyefXzCFvx/oZs4ZK/nSxy9jpFiUeksgXV+5rMW+3T28/7M/4zvLurhn62HWXHQS37zxStaduISdvSM8smEbU8sVyuk0i2ZP44Lrvss602RtexO37z9MUzrLsf5RvvK5S/jYNacRVIts3nqYd139O9Ys7eDvPzuX2+/bROnIYd6xpIGjPSVmd5h09zqMjNWw9IhsKiSbEui2wZzl08k31PHGtgF2Havx11dLvPfEej5xbhsOcpPN2TqP7K9y92sjYAgyKV01MR5Xz6muyRES+mZaOkP7S/QdGiKVsUmlbXmdDzWiUKM4XqZarpHNZfGdiKPb+jm2Y4D3XHIht3z7y6xcvphMOk1pvCwrTXWRBEmFkAjrOL0eU0+bmxrYvms3X/jyLfzmt3+ma0aOD3/kVD7xiTO5+JLFLFs2hZYOg1yTQ6FBo6OjQGdnHatXzuADl59GU0sDI8UBCnaarS8fkdyzlqxMsTt+cnONInlLHT1WJl+wOPHUDmpBiZ7ecfwoou9gQP+RGo3T82gGEpInNFw3wLJNPPW90dSNRPga+97s5tyzlnHhObPZuvVJ/MDB0AzZ/6LsyEYExbLH7Bn12Cmdhx8/SMOUHKDh1lymd0zhhKVrcD1pB45HJuiaOmFKAGWkvP5RXH2g5uMS3Cc/z1CZWYIwxHcq2E3TSE9ZQTDaTTB+FC2d57T5DVhmxA8f7qZWqXHynDxhqJNNmyyYmefFrQPs6/U4Mq7R0SjI41OXgs4Wg+e3OSycZDC9JeLuJ4f48gfn8cjLvfznsa1cetFiUpZJ19xOOtsb+dfDb9KQTrG1XOSjU6fy0OZd9JeKvOOs1TKDpYOmGWhCrjdLFsxEGCZf+9ezfG/5bH75/GaEZXLmKUsolSpqcdbIZCzGBsf4+Nfu4qz6Ah11eU5vbWClnWbcNPhTTz+6MCXRWYMgjCj5PvPqc+QNmecxVWjVVyDWSJEQPD9Q4UCVxVCZNCKUO2wCcxKf/qVGq9bJ42ptw7g2I6EsaOpwpqm1UlHJ4zCgoh3IBL5yOjKxKUlqryoDPK6/JCScCD2qDc1Q06CEyqtkAyHkM6IJDT2fr7s55ujEye4wkLTSOAMc23tdt4amXDC+56Mry1kYyC5wQ83lhNIS4vrDWHwP1JcnVD9P8XAnmFeqZc+yTKWjeMlVKvBVxaPCo4cqABgLkYYuJmaOsWda6KrrXfFdknSlSPrXbbVBhmFI2tAxNbAFdBSynNBc4Hc7DvDDW26gbXIDnuerUnu5UEaawDB0GnMZ7n1sA72ZNL/94ScBacOLraGB55OzbD74hTtZEwVkNJ1nDYO//vTTjJVreIHLxWetoduPeOCFt6gc7ua+Jzcyrehw6bypfH//YeoMm8M9I7z70rV8+2vvwauV6O8vctEHf0M2JfjTbedRrda4989buGieoFJ0aEmH9A85FGsuodAIAg3XDan4IdVMHR2TG/BqLnf/t5vNByt84V2TePcpLQyX5a3QNAx+/dYo/95ZIps1E8R5bB1UJwYMXXF6PJ90zma4p0LKS/PVGz/F1k072bPlIOX+GmN9ZcZ7iiydtYTmfCtvvLyNsZ4aM9qn8fnPfYyrP/RefNenNF7C82XPi6YJ5YgKExpvwk0LpY+9oT7PvX94gE99+mZydTBnwRQuumwB554/B2HA0Z5D7O/dTu/YYbbv6mWsMkzZHWW0MkJffx/lyhjnnLKaXK6OhctaWb1iFs/8dwvVYYdMg4GdteWtV9mNQz9krKdGLmtxyplTGC2N0tM/RhCFHNvjUh4JKHSkVWYkTJou/TBIAJGe52OmDIaPlCkPVvnyjReQSR1mYPAoupAhmpia4Hu+dNqIiCef7KZcq/Li6wM0duYI3ADHdelobufUVSfJW4au2hSTg2CkKgtEchBEWepjy6iuSwePaZgJBFB+NwWB52BaGVJda/HdKv7QbmphipPnNzN3ss0vnu7nQL+D7lbYeXicodEKlYrDoT6HrBFii4DJOZ+RUZe6jEVnu8V/N1Y5c3GGI31V3jpU44tXz+dX/9jNlq3dvO/dywirAUtWz8J3XZ59YRdm2uKoV+XjXV3c+tR6lsydxoL506hVPcxUKnFXFitVTj9hEZv3HONfG7Zy06KZfPXRVzl52SymT2+TpIW6PG7F4ZIbbmFm2efEjlbu3nWQAyWHtx2HDcNF0ppOzQswDY1R36cxZfCx+dOYX5dlZXOeEd9n1FH1FCobYSkNJDaH6LrKcwh5kIi141BllWRgeoK8ER/GwkiNwBRvS1Z5S35b4Ev9WDeNiSAtMhYhJzOhqiUIk+ZYVMlfoIRz6UINCQNf4eijRGCX1l8SxI9Q5Iz4cBGp/vcoUnzDXC53s4ZkwSRd44p/pXJJmJalCuYNlfQN5D9TuAch9GTe67qeHFPIDlpJd7TMJOSkCaF0C03VQaIotCIZLcWJzLjTXHao66rSNlT8F3VTUuGzOFkZs5VQ7YRxxa7coZV+ouzI8Qw/jvmLCGxNzvcunD6Jv+/Yz9p3nMBHPnAeI2NFdKGrTUviRjQ0bNvizNNXkG5t4B2nLWXV8lmUyxWZNAZCP6SlqYHbfvkQTz32Mp9fMJUvbzrEL79zHdO7JlEpV9CAqlfl7FOWM2fuNP7y1AbmE3F552RuOXSEFIJSxWHRwqn88qcfwa956IS8/7p72H94kMd+8y4m11s8u/4wr79xlD2D8OI+l9cOR7zVDa8eCzlassj4DtmMoGvpVNZvrbBuYZpv/PkI2w4U+cGHu1gxq8BQSSIfXC/i9teGeK3XJWso2JoqWxJCetRBIWgI1WZpoQUa3duH+fhHr+bK91+KHunkM1mu/sDlnHvmGZx3zul84L2XcsapJ9NYX8f733sx119/JXNmTmd8dFyZLswkLR5bHITK9YTKQhn6UkTNZFJ85es/4O5f/4kbPnsq139uGctWtLDhlW4eeWgXDsfws30MVcrsfLvEji0ezZOgWK0xVi5SDkqMVcboG+xh1aJ5PP3UXq76wKmcdMocnnvqbfoOjFNoSqOnNLyayrL4IaXBGrYlWLqyETeq0Ds4jqZHDBz2GRsJaJqel5cuhRKPsROEckFA00hZNse2DnPC8i4+fMU0DhzYgjBM9VqVgSQKQTfJ5NIc3FfiyGEHMwXPvtpHQ1sG35cj3A9e/H7aW1vxfA+hm3GCEyGkRqiJSDknpcU1jAL5PTZTychXFwY93UdpqG+Uvn/PU3wtBeoLPPLT1xAaBYpHtlCpBSyZ1cyJM/P85tkejvRXuPr0Zlqa00xvMZnTojE159Mz4NLRnCGX1ymXAwoZgW4ZvLW7yrkrctz7zBDvO3Mya5Z28rM/bmG4f4RLLliIM+6w7qzFbNl6mF1buxnUIzKG4MzGJm7659O8+6w1FBoLeJ5szgtV2ZMfeVxwykruemIDfrHMhS0N3PLMJq668CSa2hoYODbExR+7hRljLqdNaedX+w9imzZDnk/FD0kZOo4fMO54lHyXC6e28+72Sfx79yEe6B/mwNAo75nSxs5SDSOKiHRBEGqJLiCEnpzt42Sdhkact5NMxQlih+8HSUMrCtAoVBmfXHdFMupCaIn2EoMbZcgwlIcGNX4Smq66QkjK/CTMVk9svEbs2kL+LoZp4gfxWFNhzlW+RB7m4kphkh53PZeTLKwgZkqpMijpagiTetg4nxEE/kTaPNISh5Suwn5xSCxpWFMv3kgYLBLPELcNhmoBTzIoylomT04iCZbFb0yoECbxCCOKIhVsUYKnmKibFELakuM0cuz8kdZg+fEGqnReFwITqa9MKaRosDQ2ll1+9aNP4akaX03X8XwHXViJAB/4IeWKy5oVs5k1YxJjY+MYuiVJm5pGXSHP5td38cnv3sv3lkznjm1HOOWCk/n0NeczNFJUIZ4Q3wvQtQi/5nP/wy/wlZnTuWd0mDE3IHQDUvkUv779Gpoa89i24Avf+BuPPL2Dv/z4Anbt7eMzt77Mf144ihOEpG3B1CbBwuk2y6fbbO2JmNuicfpCm+lLOhko6RzqqfLU22O8vqfKnZ+Yw4yODKWq4mmNuPzolSF2DXkUMvLkE8/9I/U5EMSYZ0HgyyR5Km0xPlilIOr49KeuZXh4lFmzuzjj9JOZOqWTjo7JdLRPwnGqhISsXrmUyZNa8BxXzt9VFXKsy8XPgjo2oR5jwigkl8ng+x7Xf/wrvPLaq/zorvew+vQm9u8/CJrHyhNamTojwysv9rNvZ4Vci0UYwpJVKcyUclR5SNuu8BmvjuNrVXJGjhde2s9HP3Qmixa188qreziya5BccwY7Z+M58lkYP1ohbRusPrWVYnWc0bEKwtA4usulPB7SOC2XYL/jsGCQYCUglTYJnZDurcN87KMraG0cZny0mlQaSA02VLd3A8vUePaJIc44o51nXu7jWDGgvtnGdRyyqSyXnXchmZRFFKlnXRUJBYoojXLyaGq2TjjRShejLyzT5sCB/Tz37PPMmNFFJpPGcyeqrYk0fLdCtn0h6bZZuL2bGR0eZkprgbMXN/CfTcPsPlpk7cw8fmSSn1xPfd5kfHicvhGfQkrDMqD7WJmuaVl29YR01Au6R32Gih6ffO88jHSWH/32DRoyGiefMhOnFHD22Yt5/NltjPUW2eVVOXNyK0GlygMbt3LVpWfjBpKtZ+o6fuASBSGprM3ZKxfxlb8+wXnt9fjDZd4YKrN4UhvnXvttVrqCU6a186vDBzEx8dVURNPA8X0GnRrtGZvPLplDdbzK97bsYNYpy3jgjq/w0o4DjB06xrKWOnaNVhChRk19VjFWSQaV5ftr6kaSpdD1CVdTLAfEeQ7UYVkcZ0RJNhYmCvYMIz4Mq3GTOsijbgzEHEI/TJ69KM7TqYuB5/mq7ZIkrR4x0YMer+sx4SMG4moclwExTPRMOnezpq6rhnItRYF0y0TqSqQri2X8QmInQOyACkIJrdOiCF3TkwraKIpUz7S02P0fOFgUKvqsrnj4JH0NMc8GldGIT2wTL2wCO6+LiWat2A4XBiGpVCpJLAvV9WtaVsKYiW3FliWxF6CQ7YaGHsGJ7fW8fHSArqktrFyxgFK5pkZtemJBBnniMXSTiuPhOQ5EvqwBNi1M08KvOLzvc3dwQc4iCkJe0Ez+9MPrqIU+gecqKFpEFPhkLZMrv3Qnp3qCim3weN8QeWEwWqnyo1uuZNXKOehRyK/vfYYf3fk03//cSfzxv3v4zV82c8pMiw+eUsely3KcvyzDkikG7QWT9nqT1w4EnL/Ypqk5T6G5gd7+Inc+1s9YNeT2a2cyqdGk4kBGg51DAbe+PMBgJSBjyQdMaBqGpSdZnphYiyp2Q8EGbcti72uHed9lF7Fu3Voq5TKe71Gr1vB8D9d1FA5EZoSqNQfXjdlAIulw0TQwdOm+E0LH0E1Z/RpFuI5LXaFA/8Ag193wRYZGjvG9n19KwyQYGBxAM2sUnRpbNvVxYM84/X1V3nhpnO1vVuk+7LNtY42DO2tUKz7CgmydwErphL7GSKlIR2eBt18fpW1SltlzGlm8fBJPPbmLsX6HbEtKjvE0wdjRCvmcycpTWqh4ZUpVF93WOLrDo1IMybWmMG1dGVNc5eUXRFrcQaMzdGicjoYGzjq9DlNUUOVvUudRs2th6BCG+G5E56Q2mifnuO2urdRPy2Kn5bjD0k3OWH0KuWwucTgKQwm3ytASKSghQt68tWjiFEkkswdhGNA5ZQqb3trMfx/6LzNnzaSjswPHcRILvyw9qpFpbMeavASnexOV8RGamgqcPr+Oe57t5vENvazoNEmJgOGBGrapcWzEJW3I777nhWTzNpGmY5kGdTmT1w+WmVKnc+mZs3EinW/9fD0rZjWycPFk0uksq1fM5s//WI+Nxsujo3xy7iye3rqXnvEy55+9hlK5IkvGVJW053hM65rMjLYWPv+XJ/niouk8u3kft9z/JFfW1zO9uY47Dx1GuGCY0kRgCI1qFFALfK7smsZpTU3cvWkbmyyd73/9o3z6ukswczazp0/hZ39+nHd1NnO44hBG4BPhBKqzSI3HDbWWxTEHy7ak+ykI0XRDlUcdVzcg1OejjEgxWikG28rsS6CCkfIwPcHBkjcRx3HlGh7KGmYp2Slmlxrha0Jq2bFTLM7syRGmjucHyTRqoiJ8IrYRqImQhoaeK+Rv9lwPO2UReH6SaiXUME0zSR/qqg4yUgE4IaQYFIvR8XUrCBX3ytAnkt5xM5cWJXO0MJJcqDigFQs2MTMmTlXqup68OTEjJy5SEerLkMx2A8lQIppwOkRqBKKabhXnSCh4pLwNGcJAiyIMLZJwujAibeuc3V7PV+57jPqMzamnLKPshBhCEPiOei0SES2/hEGCkogzMZahc+O3fk9p9yE+NruNz20+zF3fuYH5C6ZTqpSxLEkJDYKAttZGfv+3F3j2P+u5Yf5M7urrocHOcKR3mOuvXsc1HzoDLQp4/sWtXH/jA3zyqoU8vaEbf2iAu2+YxVmLMoRBxIGeKv2jPoZlkhMez26v8vr+GuevaSBblyEIPH76cD9l1+dHV3fR2ZphrOTRUsiwud/he8/34ukaKVuf6FYOIzRdnWBUCZMuJugFlq1jCJOdrxymY1IHn/309apEDCWqmxOguFBSnWP0tWWlJorIojAJhIahL/HpcZ+5ygA11hXYe/AQH7r68xQaBL+478OIzDDjpRK7dg3z4tPDvPrcML09NZyqw8zFWaIQRgY8pi+ySeektfLYQZddb7rs3eLguNAw2SSXE+hEjA2GFMcqtE0OCDWHuuY06586hGkIMi0pQj+k2FulUGegGwZj5ZK8kRoax3Z5ONWIho6suhUra7tyRvlegGUbGKZB37YhLjx3JSljmOYGqDm+PFGq02IQhBgWhG5A1u5g5doV/PX+bbywtYcp8wu4lZBIRJi6zQWnn4Nl2hPuHIWjQBNKxJVSqZClN0RaXFutFqhIVr0GUcSKlUs5sP8YLzz/EgQBcxfOk5RigSx/Q8OtlrFzdaSnrMHp2U51tId0Nss5ixv5z6YxXtlfodEWFKtyE6w6EeWKT13KxPHl4UQ3dMrVkI4Wm1d3FBnqHePpTSNcd+E0iq7ge79+jfPXTWNSnU3rlBY62pr42783YqcsDvg1Pjd/Nrc+9iKzp01i3rxpjI2NY5oiSYc7NZfVK+dTLDvc+cQGblo8nVVt9dQMwb1H+0hrhrTHBvJEOh66LKjL8/Hps9gzMMTt+/dz6fvfyS++/XHmzO1geHycaqXGjK52tuztZs/Wvayb0camgTF0TcOLIFBH/rglMCZLJ7dCNQpWl2rlgFWn/BjiqcebkOrhiGsJVOhWkwAuGcZVPyc2QcWFgIYhJsgeqrxKCIGuxSYnGeQW2sRzaVtW4oSNb0kRkerEkfqJZdlJFQJE6LlsTjYSRnE3uIFAqDh7kDiqwjAgDHxp3zQMefJPSJOKA5NU4EZKZAkmRDxNSCiiHyQNXLEXPlCugtiGFveNCEPInIEukSeRmiHGfdcTADHpiNINhXKO+yLULh1FEbqyn/mqaCqUpb74itVvGUJ9cQ1yhuDIWJnWXJZLZ0/lW395kqHBUc4+czkBocKOG6qLWCcIXCmqx3pKFJJJW+zccYgvfe9+vja/lT/s72fxWWv53LUX0j80gm0b6uamI0TE6HCZj3/jV3y2axrP14ocLdUojZdZvqKLn3z3SkLf4+ixQS6+6h7OO7UDH0Glb5DbPzyFiIgX3hhl5xGHmdPSdDQJ9Cjgv1tD9g9pzOuw0KKARV15bnmgh80Hytx+bReTGtOUaz4tdVleOlLmxy/1EwmBqUegx6EnMGwjaToLibDTRmJKMCyT4cPj7H39KPPnzOUnP/46ba3NVKu1ZNOIb4ASZy3t2Gjy0OB5rjrdKMNDEvIMJvA4gON4NNbXsf/gEa7+8I3MW9TAz+/7CDv27uOBP27muceOcujgOHbGpWuBzpLVeZavrcet6Gx4aoSLr0szfV6KyVN0pi+0mbsyzaylNrl6g6N7PLa+7FAc9+mcYdDSkmPzq8PMX5FhYGicrrn17N8/zv6tg2SabExLZ/hgmebJOTQtRM9U8b0IIyU4ttunMh5SPyUHWiiLrpQ10/U89AhSWZPefSOsmDmD1atnUynuZ0pHAS+Q9ksUekfXdHw/IqjClJnLGB70uemWJ2ial8OwddyqjxcEtNa1cMHp58jTrEZCVtA0mVMRhiFpv76XCLhx051ppRJacTyvNw2bNSeupqNzEg8//DiHDx5h8aJFEvcdyjEYCiVk2jbpqatwR3ZT7O/GtjOcNi/Pk9tLPLOjTDqTQrcsLNvi8JBP4NRI6xq1qk99zqBY05k9Lc2GPQ4fOrmOaqnELx7p4WPv7OBAb417/raJy981j1Tks2jpDIYGS2x8bR/VlIFPxCVtk/j2I8/zrjNW0tRaLx1OmrTWAtQ8n/PPPIGn3tjOrn3dnNDezI/3ddNsmjiuh4gixnwXjYhrZsxknpXih29vpTy1hXt//AWuuOw0xkslKk4tqQKIhGDmtHZ++tcneEd7IwcqDlU3JIgivChKaL2RuuUFCpVuGJLZFnefx4gRmWeTHUWhcvlFkbxBAHi+l+gYSTu3okubponneZJZqE0Up/lBXOEcHVfrLac4umHguW7C1grV5ylNMmEiWcQokwRfH/+ziMRWrOfzuZvD2KalZm6havAzhJ7wgBJHltIrPNUSqGsKn6w6OAyVBZHUW8l0iTs3YriYrsuEujhuthcEgURJqFi/aVlJP7a04iqXVUyVjNHucVGUJvAcV33C2nGWNF3h42NKL/h+mLR1xXNEW4s7FSIyuqA5ZbO3WCKXsrly3lTufnwDb761k/NOX0m6Lo/nOGgqkIiyOYvYWaHrOI5Le0cri+dN4Yv3PUutUOB3P/p/aIbMvOhCx/dcojCkubGez3/n9+QO9LOws41/Dw6SRcMn5De3f4S25jy+U+PKG+7FtgUXnjWDxx/dys9vmI5IZdi8pY9I1zhxSZ5yDd7eV+PlPQ5bD7t88b0tvHXQYUazydNvj/PImyV+ccM0utqzVByfxlyaZw6Uue35XkxLx0rpyalIN4wErmnZCkdCBIHq0PA0Dr7Vi+WluP76K/nijTdQyGUpVytJkZPsTFBFPIaedCbEjKv4vZMkAJH0xwjdUOVCBkEY0dxUz569h7ju+i+wcGkTX//fy/nnX1/nkYc3kG1wmbdaZ/Zyi/qGuH5TMD4ScP/P+zj1Cig0WIwPOAQBeE6I7wQSQTNZZ+4am2nzLPoPhbz1Yo1C1qTnaIW2KYIQD8d3MFIBr78wiDBDso02PTvHaZ6Uoq4+hVnnoUWQygt6D3iMD3nk29IqcCu/0K7nY+hSC7RSJkde6+fWW64lZTpk7REMIUmw8Rxb13TQBEP7KvRuH2Lxmrn85OfreW3vETrm5vEcuZFXnRrvPvsSli1YTM2pqps9icsxpjrElc0xUTVeGILAS76DcWW07/l4nsuUqVOZNXsW/f2DtHe2yxFjIG/28cw9ikI0Q8duX4Y7uIPIHaK5vsAZC/M8u22UrAnXnNbIrCl1zJueZWevR2m8yuR6g8bGDN29VWa0W2w86NCYFpy+qI7JmRq/enqYG97ZziOvDLBh0xGuuGgeXsXhjLNX8OKrexg+Nsq+wGVNcz3BWI0HN23nigtOVgS8uMFSMZ50wTtOPYF7HnyW1iCiL9Ao1xyErjPuu5wyqZmrpkzlxUNH+NvYEJ/6xPv5wZevob4xw8DwcCJcR4F04bl+yIxpk9l5qI9tm/ewdmorb/aPA+BGcu1yXFc5T6Mk4hD6csKiGzJAGx3XBSOdUX6iccS7lRb3FgklYscH7OMw7bpQa25sA1aOV01oeK4/UUalbh8SNSRT7r6yGhODZHVdZtx0gR+ECZJeiiRSupD96WqTqasr3ByLM7oud8C4Gzcm7MZXqjjsIhEVunoAw6SIPfB9NR9TjVXKH+0rHcXzPNkRIre6xEUgT6Pqyq7GZnFHcBhpCB0FSBN4nv9/am/jxSj2KMvfecLpFaiofqhOZ/EITheGpMAKQaOpoxMxvZBmdn2a0ZpHLfBpsS0OjlcY8X2uXzKL5zfv5YGnNnDO2kU0tDbgVF0CrzYhz2hCpV0tIvX6VyyeSdeMDs4/ewVzZ3dQrjqYpq4+oIi6QppnX9zKHXf8nc+vXMyvjx5DD2FoeJzv3XQZp5w0BzyH7/z4Mf7yry3c/d1zuOev27jmFJNZs5rYs72fVCpiSnuGLbtqVD2NQjogo3ucMK/ArBl1rN9aYl+/wyObK3z/Q1OZNyXLWMmjtS7LY7uL/OK1YdK2DAXGi3qMlzYsfSI9bugIQ948gqrGwdd7OePUE/nRD29i3cmrpdbheYkgKB9RkdSUxnyeuHcmZqVFcc5HOVBkla8BagSZz6TZs/cAn/ifrzJrbj0XvOtEXnzpNbzMAZaemgHhMdwf0XvAp1IS1GohZgae/EuRwuQqrdMMxoZCfMD1JGbTShukU5IlFfmCTMZk/po07bMFW98aZetrQ7TPtWmbatPfPU5dq8HIcMjBHUUa29MMH6vQ0pamUGdh1bkErkauwWC4N2DomEd9RwbNiFQwTN7mhKZT15xh27NHeM87T+Ij155Pa2OayB+iWhtHi0KslEkUaISBRnXIY2D3MKFvsGfY50d3rGfmmiZJmLUMBodGWDF/Gde998O4vqN0DvneCV1L+rnj+ue4vjdO38Wn2HhBkAVfRoL29nyfQiHP0mVLVFbBnRBa1eerCY3I99F0k7quE3D7d+IUh2ltrufEWTl++XQ3u46UWD07R0Ndms6CxuExjcB3WbZmFiJfIGsL9vS6dDQIqrWAjqYMbfmIv706zg0XtfPj+w+StiJOO30OhAYrVs3h/gc3YvkaG0rjfGLOdB56cyfjYcgZpy6nUqlhWJbsFwpDUoZNEMDjDz/H1DBia82n6Hm0Z2yumzGNyHG57e3tTF27mHt/fCNnnLqU4VKRmuMna1RM6Q4D1YNh6kxpn8Tdf3uKE5vrOVarIdBwg0jaehUrKr4Jyl70Cf04tvr6viet6qpuIFTrolB5ujjyEG/ugTqkhwokGYerwzAiCiRiRB4QQtUmKMkhE+aUKJELJCVB1QpoqnI3Zm2pyZGl6L2SySVrGRIzUhBIFlZsa41FE03Xko3Csixcz8PQpXYRBwIjFQgMouNa5IzYlitBiLGSr6nZmyz8MZOQoOcpsU/Xlb1Qn+C6RGHSRhgmZEt5S9BNE031b4RSfU+sbUEUKSCk6kU3lJ04KRNSkDDfxdCgYFhYaEyvT3Hh9GZc12VteyNNWZt941UcP8DVYE+lyrtndVIbK3HrHx9n2eypLFg6m2KpLGeamtoU1UJhqM24VKqweNFUpnU2U67UZFVwEOK7nrTN1Xw+9PmfcnXLJLaFLnurNUb7x7no/CXc+Jl3IkKfx57Zwce//B/uufU8xscddm3axeeumsPB/UXGR0aYMrsdX6SozwbkrJDI9xkohUyZXs+k1hz/fWOMRzaNc9O7J7N2fj2lQNDSkOGJPUV+8coQWVsowVd9VoZQ77OebI5Cl5uH1DxSHHq9h7UnLOeHP/gGAhgfKya1ndIqbSQnL900lG1aznTjLmmhKzKvpvS1MMI0pMNNQ3YjpFM23d29fOz/fRnTDpk6u5nu8UPk2qqMDcCWl1zefLbK2xvKHNxVZddbo+x8q8j210qUxquURmHfGyHdBwPGx0N0W6da9XGdEDstCAhwPZ9iKWDvzjLdRwL8MGR4wGXX62W6lgtsyyRwQxomWbzx4jDpepvxwSoNTTbZvEG6SbLAzIxgtNun75BHY1dWbsgqpW9YOpm8zY4Xj7Cwo5Of3fEJnLLH4NBBRob3Sk0tkhZhzTSpDXuM7hmjv7uGmNbJ93+xkVynSaE5jV8LcQOHmV1dfO26L5DLpqg5LkKYKufhJ7N3ScFWp1h1Og0ViTpGIMWd6nHjaVy/ABqOU8N1HGnVjuf1IUSRHEMGCkAqVVuL7JQT8Yd3MT4yQEtDjlVdOX7xeB/bD5cpjtc4NBTSnDd4aUeVbFpw0tpZGHaOrXsG8MbHWDAlzfodVdbMK3BsqEolMDh9RRM33bmJ05dNoqs9R2tnIxkrxX+e2IxlWxyqVvnMnFnc8uQrnLR8Dl0z2qlWZU9Qc2sDQwNF3nP9N1ldcTimaeyuVblyRienNuT51aYdbLU1bv369Xz++neBETE8WpRYIyQBI27gE0KXtQFaSKVSZcHcLnbtPsrogaPUZVMcGSlLVp8mGzBlUJBEAI/CIMlY6IYhLer/n6r3Do/zKtP/P+etM6OZUZclWZZ773FiJ07vIZ0kEHpf6gLL0hZ+LIS+313KAgvL0kvopEA6Cal2Esdx4t6LLBfJ6hpNfcs5vz/OeUcm1xVyGcvW6C3nnOd57vtzJ3NTY8yl7rszZGolEbZdBzQmitnEaOi4jsEvJcRzrc5SCN1Sh2lasWVroYHZJFwvyXqaTnxNBErJXFoYE7jjOEaYIerXQ2ice/YuJVXd86FxwhrDnpyqk4cnkfLKWOIa6WwCKLMSqadl1WW+YRAYdY1uG+kWnIX+Ems6w8DsSGEYasS3pU+/rqt9F9RZ9HpjSLwCVlIFybg+uLfEtO1ftxAw5EkdLq9Mzrlt2eQ8hxYLQqG4fX4nf9rfz/GUz6vHBpnhe1wxux3bd+gvVJBBzM7JIpf1tjMn5fGFXz5E1nG4+LLzqAURKoxRKkZYbr28T+KpqtWQWhBiengI4RCpmBltLdz1nT9SevUI58+dxe8GB7DLIZ0zGvm/b7+dlGcxMDDObe/8DTdfOYdPvHMFn/z6c7z3mmYyls3A8WFmL5tFaGewghKlkQlOni4TozgwLLnknHZePlTi/x4d4mM3tnP9ee2MlRVtuQxPHyny3eeHyGU8MD6OBHmhJLi+W8+RV0bypmJJKpuiNFgjFXl8/WufxXM9alUNlNMzHcdkuhu3rNH7RlFk7quhIpvwGy28COt+jyR2Vpq2Vy6f5Zv//X888+wWotim7+gQ5YJix6YpDu6scPTgJMWpMo4tyOZ8WluztLTkyOUawIJKscpkcQohbSZGYHwwRoYWMYKhkwGnj8QMHo853RdTKShSlkVDymbVRQ0MHK2w89kSay5pYGokomt2lhNHSwyerlGrhTQ3+zQ3pZFOQCpjgy0Rkc+pwwFNPRlsx8JNO/gZm7AiOfT8IFect4yf/vwjEFmMjx5ieHA7tWoRGSbRpw7hlGJ49wgD/SXcuT389L5DVIiZvbaZqKbfORlGvOu2d7Bs0UKGR4dpyGZIp30asmmy2YyGBtZCTW0wQWqxDOuqHs4OQ4qkqUis6Wz5JEdH2PXOhEKZzUfPq2QUaLOkIXYrJXG8NOmutdTO7KJaLtDdlGbJrDQ/feIM1Vhw3ZoWLMumq8Vix6ECu/aP0OSHVKohpWKV2S0OmZzPtkMVLlma4t4t47zjqpkcGQr5xX27ecPNy3GjCuvXLeDA0TEO7htgwhU0eBavaZ/BVx/dxJ3Xn4+fyZLLZdm/u4873v1leianmN8zg2HL4r1ze9gzMMoPDvVxwx1X8KOv/TOLl8xifHKSWhBqI7PQ6jeEMhWyU3+GFZKOtib27+3n/377N85vzjBQqjBYC6lJqKEDvhKfRWywMQkGJAG5WmajTtRNSby3VY/9NkuGaZ0lNgiNa3cMWkeYdMOoTg+QUtU5gMoYWadbUXqjcF3XpL+q6YAoY5lIZsRSqrptQwhR5xxistjDINQzEGFypjVhEZRxSSSLc0KOVBKiODKYZGXgefrDJxI6jOs2wZFYInEOm4tm24RRWG+JyEiXR5ZRjuhSTpdWruPW5wWJ8iowTvdk00pgjyhNkIyjuM7ST4aKvudpVUukIWUS8B2bvG2jYsnijjxtQvFKLHnkj19lwbI5/GXbAZ7Y08/qlhznt+Y5E+j+5I6xKeY2Zrmmdyb/8afHGR0a4YrL1iIcrWTQrRhplBJxfUDlWC6unyKKakgpyWazbN12iK9/49f8f2tW8eOBU1TLIZVKle/8xx0sW96LjALe+/F7ODUwwp+/ex2bt55i69bDvGVjniNHJli7LItwHGRxkolTwwyPanSzl8vQN2GzpNfnw//Xz+s2NPL262YxMFymszHDCydLfPO5QVK2jeUYS7nQpa0lBLZnJKSuTWyUbd5Zw+Djr5zmQx96DxvOW0thclLfdyG0QscoRqS5bwjbYBiseuSvVpOK6euUOHeFRSj1ocNOQmuML6JULnHxheeycGEvg6fGOTM8SvesPBdfspjXvf4C3vjGjdzx+g3ceMu53HLrem68aT2vvfEibrrxIubNncXBQ8eZN7+N8hQc2jfC1LDDSD+Ux2zmLUmz8fo8y8/J0NHp0NHj4KVslmzMsP/FEsMDEasvziNrFrGl2LV1AklMe0cDc+Y2M1Uu0dCozVy5nM++rWWyLS4ykoyeKDGwdxJnHN72xku49LLl/PW+F2jNDzAxtouoVkFF6GoFizBQHHjyBOMjgmJrJ3965hgjkxVyHWlOHyowcGycwePjFIYqbN70Eg//7QnuvfcRHn7oCR544O88/fQWdu05gFQxMzraaGtrJo4k1WpF6/Yt2wTC6cwV3SkQ9TQ9Ycy6SUS0fl8DE5zk65xxFZuc7KBumKtUQ4qTJTzPwU1n8duXUO7bQmFiikU9eeZ0prn/5Qk2zPe4enUTs3pbWDkvh1IRL+8ZY/eJKu25GCcK6JyRY6gQk/EkE8UqgXJ509Vz+fGDxzjWP8brb11FHMMl5y/lwSd2ExYq7KqWubyznWwx5Odb9/Ke2y/n0cdf5k0f+U+WuyluWryAVINDcy3gO6/sRyzu5btf+QBvvuMqpqpliqWyUZklno0ISySgVt0Wj6KIfC5DOpXmh796hA9/8rtc3ZJlTkuWTacmkEJQiWKiWJlIbf6Blpuku0qDbUeYDdxcxCSe2zHxAUlVaNsWMtI+IsvIfEEZhatdb8tTz9MRWCJJjq1Pr41VQv1DfG2yliZflVB/ldKzxGSdjpUkMo74JBxQyhi7IZO9K2nvCNO1RkmUiusnj0RSKwRnKToMC8X2dPtIJZuTqP9dSbpgEtbu+h4CpfPUPa3td1zdY1embaWjKTEZwlrTrUOsIgOw08x9pXSkZJI3YhnVVRKXqqSqz0Qsy6o7LBNZaEoomh0HieSang4ePtjPDe94DevPWUxLW57bb9lIc3Mjd2/azuBYgZvnzmDneBERwpFSCd+xeffKJfzyiS0899Ierr38PHJNzdSqQT0JUclIt7RsI79LZMyWwsXhn/79h9yRytJvS3ZXqlQnyrzx9nN415s2IpD89O6X+NaPnuFHX7mac9d08Kvf7YJSgYMnq2w+XOOZfUXu2TTCtoMFThdsqqWYFcsbKdtpTo7G/Pq5Mea3unzi1m7GypK2xgz7xwP+46kBPM/FcbQowfUdLEsD1iwnaW2YMJwkQEYp/LTP5GCJFr+VD3/4Pcg4rPdm9c9nwrVQ2MIhjMP63Mzkkuk8l0hXHElPV8u9kxAtk+OshBnohZy/fg1Llszl5W27ePHFHczsaeRt776cd737Iq67bi0LF3fiNyjK1XFKtUlGxk5RKo8i7RJ+JubccxdxzeUbGR0f5z3vuoY1K5dzvK+f8akiSsWcORpzfHeVKJK4WY0NiUMIKpKlF/lsuqdM4wyHjm4dvbx3+xTFYpWOjiZWnp8jUEWctE2lFNPa7nPkVZ1dUh2qsbirk41relmxrAPH8Ti09zBrVjpk0uOEtRoyMH1o4aBkQKriMzrRxI6S4N6nDzA0Wka5EJRqzO5t5IIN87j++nXcfNMGXnPNejZuXMvGi9aybHUPXbNSVKMxHv/bc/zkJ3/iySc3MTZeoLd3Jt3dXSilW1J1tU4ickDWOxD1mALXQcm4Lg3V988k5wlFHFWxLIUQLul0llMnjrHlhedobe0infJxGppwm+dRPfEiQRhx3rJ20q7F//vrSVZ0urQ2ZVDZPH5YZM1MRf9gEdsRdGShVoOmZp8zoxVSvsWposudl3bgpjJ89zfbmT8zz9qV3WRyGRbMbeeeh7fTlcvy/OQkb5vbzeSpUf7rL09zzwPPckdnJx9av4zdo2Pcu+MwW6TgM599F1/41zfS1JJlZHxSq9EczwRY6bVMyojphG4Hz3WY0dHOjj3HeM+/fJu9T27hE+uXYdsOfzl4inIQU5aSlG0TxhIcq96Cx7SsEswJaIagTmulbiScJoQr7QcyMReJ0TmZDYsEyWRrYKZVF6VIgyAxreKEFSi170cY+4JS04P85CCu/SR6to2w8H2vjmEBgW05moAe6jgGx9I9Azufz98lZcIYMg7vWOKn/Xrspu3YRn4ZT8eTJpx501KyjAQzCiMTG0vdB6CSnBCpwYrCTty5qp4Z7LnaoBaHOolMWFY91znhvwgxnTmRqAwscwN0YlwCI7NMjKQJaYF60JGwLISSNDoOec+lI+exKJviuUrI5//t7cQyoFKrooC1axfyxpsv48V9/fQdGySfTnG0HOBbNv1TJYZrVf5p1RK27uvj5395kssvXENXT6emx6oQYbnTeHN0vziIJTPaWvjOTx/ixKYdXL9sLr85M4yaqDBrVp5vffV1+L7Nnv2nef0Hf8/77lzOx96+gnvv28MvHz3KrpNlJqqSfFpgo+hu8Thd8alJxZ2X5WhsyXH4VJlvPzzEjLzNV97QQzmyaWrM0V8I+fITp4mVwHVFkhGsKzbT09QDUi3dVWeHHgl9jwb2j/OOt97J+nNXUyqW6zkelq1fFN0ysc1hwDU50xJhY5RnoXlIxVkZCnI6oTJONvyIhkyabDbFf3/vp3zkY18CO+KDH76Gd73vYhYubENIRf+ZfnYceYWdh3dz8NQRDg8c58TwafqHBjg5eprTQyc4fPIg2aaQc1es5sknd/Om113LB99zB+1tOUZHJjhx6jRTE1VOH1IMHo/JNkFDExA7+A0Os5YKnvj1FPPXurR05Dm8p8TIUJXuGY0sXNVAVUyBZVEqRsxekOP47pCBkwXe8Y6L+N5338L5G2fR3ZNj/lyHCzbkSDlFyuNFwlqM5TgI16E0Mkne7mbbwQw/fvhVHtu8l4a8z4YLZnP7refy/vddwa23rea882czZ2Ej7V0OubaYxpkxM2bpe3ns0DijwzXOWbuEm2+6EMuJ+NOfHuIPf3iI0fEx5s3rpbuzg3K5ot8S0xaxLde0TSJDHk5mIdongkrc07rCtQAZacS5EhDGkqamZnp65xBUqzTk8sRRDSc3AzvfSXTmJYplwYaFjZQiwQ8fO8W1a/I4BEyMFqgFMfmUYqQQ09nsUi6F+CJkvCRpbnJ44NUSRwfKvO+meRwfjfjB71/l1isX09LgMnduB+VyxFOb9pNL+Tw3WeCWznYapkrcMq+X3nSGz2/axpODo9xw2+X84L8+zJpVCxibnKRaDTU5Q+gNw7Y8YhmiVIiFXYcGNuazhIHkP/7793z9qz/m5o4cty+fy6OHB9gyMEw1iJnZkmVVeyPDU2VsRyBsh5rJX3FsmzAM9CxQ6RmzMG1hSxh/WywJQ+2viuNp37RtDLbTYiDT4pdx/UCtjE0CUwjEptXmum49zlsBrm3X1/ik9ZXAZTGYeT3Y15wsy9Hzs0T5FRtEPCgNjLRMJrqSUqujHEs7yQ1yGrMwa2DidFa5U/dxxNO9bGNKsW3HnF5k3c8hDHTwbFmnrA/JE9OYbkdNh5gYl7gZoIuzkClBoJ29cRiZvrsw0bSW7uuZSX4cxvXhYbIwhXFM2rJodCxCJbmst53Nh09xzo2XctPV65ksFXU8pOVSrtRoyjfy579uZr4MGJSK06WAyTDCdhwmagH7J8d4y8olyGKNr/3yfpbM6WbNuuVUKjWNpIgiswprDEC6IcOxo0N8+ks/5lPLF/LLgdMUKzFhGPDVz93C8iU9VKtV3vqRP9Lgx3zt/av54Jc2c/cDezlnpuK2NQ1ctdDlwsVpVnR6nL+0kVNnQlqbLFbOzdKQSfPrZycYKoR8593zyDT4OLZDoab4yt9PM16RNOQ8wkDfb8cg0THyQiNo0ycVc8MTfHO5UKVBZvmXj7zHqPR0lWjXKbVafpvc2yRjwKoPyaWR6sb/gGdQSpmqNjC8JmhpbmR0fIIPfvjfeWbzZr70pXfwyX+7la5On9GxSQpTU+wf2M32vl2MFkapBBHVamhcvkbAIQVBFFGLqwxNDiG9cS48bw3/892/MndBO9dcvZZrrjyX3p42jvQNMDAwRFT2OLY7Jgqha56FIwXpvOZNvfr3gNWXuex7uUT/sSmWLuki12gTeVPYnsX4aI1Zs5sQQZZ9u87QO6uFKy9fzvDwII5XZOL0EY4f7KNWCgAttCgNF7GnFOnMCn58/0m+8dO/4qQs3veBa/mXD9/Aa286h/mLGzkzeZwDJ/dx4Pg+9vcf5MREH8dH+nj8kd38+n+38beHX8FKFzhvYzvrzp/JihULufWmq7nlpoupRSV+/asHuP/+R1EozjlnNQ2ZDKVKxaTY6XiCJK5UKm06VELPrPTpTSIs814rfU+FcFAkHgGbTCZPNpejVqtiOz5hpUDDjAXgNRIObqMa+Vy6ooX9pyvc9+IoKzscfFeSdm1krcpYIWZouEzKs2hssAkNYbcQSA6eKvLczgKfeetSHtkyxAsv9/P6W9aAstiwfiHPbTnC0FCJBsfm2eIUG9taUCrmX3fs5PwNK/jWV97Hm15/JaValYnJCWzHwXU8ojjUij8swrimDcZCe7vAoiHts2nTLt7/0f/ixNZdfP6yNQxXI+45eJLjo0Vm5NO8ZnEPMzMpXKW4qLeVPcMFDMIBJfRhyLan7QzKtPITNVzi73Ad+x9yi6YhtHEdQissYdhqemacVClny6qtfzB9M71ZmFlyogpLJJeWIe8m72IcRYajpuchURSZ2FtpDNx2/VBpN+Yb74pkXB/oRGGoOw6xLl1dz9W7jumHStOW0sHxVj3WVWDpOUcyBDUAQ+ssllX9YiFNwpXWaeuBkahfVIHA8/TCFhuulVZ0yel5iNm5k5Ns0ipJQoviMDJ9dl2Gx0bT7DsuDRZkbJuGBodVuQx/Gy/yuc+8E8/XO7xWmQham5t5ZvNO7v/DY1w8s43H+0cIsIiBEEkcC+JI8crIKNfM7mFhJsuXf/0AOc/mggtWUzU+D9sMGhGCfCbLez7zP1yiBEXfZdN4kWiqymtvWcs/vf1ShIr40a9e4Bd/3s6X3ruUz31/J255lM/c2MYMV+BIyOVd0vkMrg3PHywxqyvNibGYDQtSPPrKBD9+coxvvWMeszqyVALwUy5ff/oUxycl2bRdTwfUATYWdpLrbO6HbRtku9lUdDVhMTVWprd1Frfdch2VcrWOi5MqxhJJy2NacqhkXBdhJLJgO7lnUtWjOhPjqaWTumjvaOe5TS/zjnd+gnRW8oWv3EH3rEbSacmpkdMMF4fYfXIXJ8ZPIaSF3qMVvu/UGWWOY+sgHl+YEaLFyOQ4sTPK2hXz+fLn76exewrHr7B8xSyuu3YjtnDYtecQlq0YOBRx6nBE52xJQ9Zn9tIGdm+u4NoKhc3BPUU2nj+fkfFJ0u0VLBwmxkM816F3TjM7No1y4tQYl25cSEeHx9Cpw0ycHEbEiqgUExYjhnePYlfaGfGW8IX/fZKnX3iVt7z9Wr721XdywfmzOTNxhJ1Ht7H76G7GisN6/ufaZBp9Du8Iefy3BYZPRay4KM/FN7fTs8BnqjrOkcEj7Du1myOn99Le1sprr7+aCy9czu49R/nNr//Krt17WL16BbNmdlOt1ox6MdYLqVQoZN2prKXcrhaiqOkWl87oDrEszyBoIu3hsh0819U9eiVQURWvdTEKRfnkq4Qiw5Wr23hyzwT3vjDMVNXhxIhG0IdRyNBkRCwFvTOyDE3WEMJiJPT51Otm8tS2IXYcLvKRO+fy37/bTy7jcuHFi7AUrFg6kz/ev1X3+RE0OjbPHx/gzW+9ge989UP4aZvR8UntvLfQxkmhf279zBo0h9BDcsd2iMIajfk0P//9Exx/YQf/dcel/O/Lh9jSP4jnWly/ZDYXzezggR1HeWK8zJ5SiS5gfluefaMFPMemGmj/zdldFakkvucbM6esq7AERmhkZOxKSsI4qhtq6+EnhtQhzhI1JW0xaYjnSXS3OittUJ2Vh5Qc3jF6H2HimvXQX897U6lU3cirzFggoR5rOrDCbshk7krUTEpiXN4JydPEvp5FgEzMJYmiKo7CaQCY9vDVd86kfRMZZLScdrQYZsw0FVwJtMoqSc5KAliUBvYJYVpgRr2lYyFdVB2sOH1aTvLck4uUODZjGZO2BC2uRywU67ubOXx6lI4Nq3nL7ZcxMTFq0g51gEraS3HXN3/HqlrISBxzaLJCYB7QINIblee6eMJm2/Aoi5rzXDO7ly/94RHOnBriuqvWg61LxCiSdLQ187M/PsmmvzzNnUsW8P1jJ0jFglze5Rt33UZza4a+vjO85xP38bbr5/D41jN0uRU+eVMnm1+ZoCHnM3deGse2KIxVOTpYQaTTnLe8hR1HK3iO5PP3jvDBa2dy1Zo2yrGgJevx/RcHeb6vTGODizI3X8+vHHMy0SdISyiEY9XNS5YJ76pLvKXF6Ilx1p+3hubGRl06G8QIKC3RFtPRyLaRZwuDW0geeNtyUELoMtmeli7GUtHe1s6f73mYD3zo37ny2mV8+/tvZ8aMPI88vplNW3YiGgscHTlKsVLCVm598fM8hyiUxCbPXGIqUkMcEGhUxODQGMIt0dXWxk/+5yXmnlPhaP9RwrjA9a9Zz3nnrmHby4eZKIwRBQ57Xg5IpQQ981yaOm2eua9MqGIKYxFveMOlbN95iMauANtyiJSkMB6yZHkrfXurnDoxwemBAjdeM4eTLx/m9LZhRvsqjB+PqU151Bpn8cxpybd+/gDdPW1845sf5OrrF9A39Aov7X2egdGTTFWTECgbL2NTLkU88IsxDu8us+7qBjbe0Ehjq8XkaEVLVxVYlotjWVSjMocH93JqoJ9zlq3m9tsuxvZcHnpgE3954DF6erpZvWo51Wq1fi+FYxkjbjxNhVWGiGwWWS2QcAwmSPfXQeC6njaNGpObxt14xEGJdOdqZG2SyvAB3FQj585Nc/+WITaunsFrr1pIpqWdUi3mlaMFMrYg54RMVRRBEDOqfNav6GLDfIsXd48ShhYXrmrj23fv4ZqLZtPRlqNrXjvVYo2nNx3CSXvMSPk0BCFBWyMXrl/G8PCY8YZYBnApiFWEEI4+xAitkJRK6sokirQCC8H8eT08+uhLOFHIiwOTnNfdzhtXzefE0BjffGkPiy46h5985+NcfeEavvHLB7h1fjenKgGx0m3hUJrhtUHaWLZFGIQoY2JOvB1Kga1VLZqQYZtDMkk8tGNGKPodTgIxbFsnuyatcjNM0euvbdejEGxTXSoDckxgubriwWSoa3KBbSqfyPi6BLr7YzsOURjqFETLws40ZO5KwpQEEIVJ6p+uDFxfD5dsYZn+mYEXokw/26nrhmOT9+y67rQb0pxyLdtgh6UkCqJ6yIl2vWubf2yCWJLSShj/RpJcmCxEjpNkleheXWwidXX5Ker505hyVEqp1VmWRcaGnJEcXtLVzgP9A3ziE2+nrVVHfyrzsuSzOV7ecYjf/fRBbpzfxcMnhkHYRMIiiGNcyyKMpVZeSEVKWOwvTGFZivcsX8pvNm3nkc3beM1FK2lsbcR1XQb6J/nAZ7/HpxfM5YGJMaawGR+b4rP/cg0Xnz8fGYd85N8fwFMBcztd9u8b5t9vaeWpV4osXZRj7fIMxcmQwYEKU6WQEemxblUbnu/z7O5xvv/4CJcszfGRm2cxXoxpa0xx344h/rRvipacrxdTg0iIkkAoCzMbMvcn1qV3HMnpoZ0Q2L6DUC4jx8e5/jVX0NrSTLVWMy1C6q1MGSvNCzNKDi3GcIlkqEcuWEQyqFeSekal3bS5XAO/vPsevvD5b/LRj7+GT37mRk6fGuf48QHiMGbz5iNsevw4SgpiCY4PblrWxRJSGFaZlLiejW0LXAdc18L1BemUg+9aFApVuma7TA7C7lfGWXK+z6mBM5waPUpXdxM33XgVw4NTHOnrRyA4vh/2bSuxfEOOg7tDDu6ZZPmKLl5706U8sWkzc1emIBYE1YixkSozZ2bp6Gphz4sjDIxMMjVaYlaDQ7GkqHmNDFoNbB0P+dOTu3h110E+9P438qnP3khR7eeV/VsYHj9jIlhVvb2RzsLxfVWeeaDAvBUeV9yZp63VozgaY3yb9TA3hCKSBogX20xURjl4cg+Nfhu3vOZSVqxewMtbD/Gb395DLGMuvfQCYmmCg4zCJ5GW2ta0XyShYMcy1PEJllNPJFXC0uFHJg5BnJVHIiwbZEBm1rnURg5THTtJW3Oennaf/77/CNed18WKBR30duY5OVplx7EplrfbRBIqkaAgXBZ3OuCnuXRVI7999DSXnNvGVDniL387wltft5agVGPDhoU8/sx+Tp4YRaU9NrZm+fOrB3nNVefR0tyk4x2SOY85yQsEwnaxlEVDQ4qgFujfs3QXJgojenva2XVshJO7D/GvV6ylybX5weYd7PczfOvrH+FfP3AbSgTMntvJ3iOnObb7IBtmd7HzzASOENSkIpRKV1kGHmo509BZhCb3atRJNB1AZZzmUsXG3KvhmcpkdnBWFK6S0xu5NDpaHSAW1wUpUsb1YCmSJNhIdxdAz05s1zYttumDoLCECdFLoI6OyWSS2Nls/i7b1kOaxAiYJAdajkMQhBqclSirwlhXJEayC8Is6AokpE0G+HRwkySKIvNvbFLlpqmOOnxKmQG97vMl5MhkYJM4ZkVdaqaI4tD09pJM97h+wo3CuG7WSb6PEpK049Dq63J8XU8bY2OT1ObN5GPvu5WpwgSW0ItrHEfksym+9q0/saBaY1SGHB0r4jgOVQWxkLi2pdU6UhJJSVkqGvwUp8pFBqsV/nnFQg4dH+Rb9zzFhiXzWTJ7Nq/72De4IBak8hkeHhqFSsSG9bP5/Kdeg+NZPPDQdr78nWf58gdW8ujzg7zpHIewZpFOw6Jejz07CxQnA8qVkBjJcOjQ25nBcQQ//NsATTmH//fW+VQjQUYodg6U+e6WUdK2AFe3Eh1by6U9X1cLsdS588m9T6ItfUPdjY2ML+V5HN8xwKUbL+C222+gXC6ZnAj9gNuGxSONksRx9a8tR+P1HcvFsRwTceuag4OpUIOIjhlt/Pm+R/jcF77JO//pChav7OLuXzzP00/vYMvLe3nhxSOUCyHFQsTR3VP07y1zcn+VscGIhhbIt2mgo+MKMjkLN2VTKceUipJTh6sc31fm2P4yfQcrjI1Ixs5UWXKux4t/m6CxA7rmeoRKcurMCcJ4mDtuvQ6HPC9u2cNlV8zn3LVLePBPJ6hFEaVyhTtu3UhvVzdPvvQUc5dnqNV0m3RqMkBKxYYL59C3v8jIUIG9R0Z56WiRvWMhW/sLPPdqP7t2HaGttZlvfvMTnHNxipf2P87g6ACWKwmlolQOkUJSCxTCtTiwq8rOrVXOuy5Nz4I0U6MRcaSwPF01Oo6Nn3K1LNekicaBrv5s4VEJqpwuHKOrcSbnrF7G1decz4kTo/zy53/k1JkzXHHFRbiOo4OuPKNIMrgd3YJRemGywHF9lIq1lysJo5KKxmyOWhSafr5nWHh6jYiVVkGmO1cRDW9ncnyMpb15To5W+M2Tp7j+vHbSuRzr5qQ5Olyjb7jCklab4ULMUCC4dFUbrx6aYkaTx4Jul3ueHeMjd87lR/cdIw5DLrtkLrZtMW/ODO5/4FXGVMzy5ixj40VKvsMVF62iUKroU7RtERnlppSK5pY8lWKNA4dP0D2zjaBWMcF52jeBspg1awa/uecZKuMF/nfnMW57+6388D8/wry5nQyPjhhChmTh/F7+5+7HuHJmG8NhQDWWhChiAdVaWE/00zMOPU9OXOtRHJmgNFUP5RMJKTFxTakEVRL9QwZ6bPD9yjJKzyQS17ZNJoyqR3FjLA7KZCglaivrLLyQVoZRX5cjI+8V9fhdZVRYjfm7pDLBIeKsvF0hTD6DACSOZT6kpUwbSeJ42jxmC0FsPnCSHGgbgF5ShSRSQMedDn+36xJPUWchJcqQZDeeDlXRshBdCp49YzGVkWXhmhaYbR74MAo1U8vgU1IoGiwHJSRX9MzgweOned+H38j82R1UqtX6jp7yfQ4eOsEXvvxrbl7cSaPjMBRIRmoRGVtvllFsAJGWCZoHQqVzMSZlzM6JMe6Y3UtHCN97+Dn+9OzLiP5B3rhgDv934hSecKmGNf7ny7cxc0YTg2cmeNM//5kbLuxiTlcDe3b3c/OqNCNTEauXZNm9u0AtkkQ1g4NJe5TcHBuW5HnwxWH+unWMH7xnIbmsjo+tKsHXnx0kFBaOIwwaxAy9DcZby3TtOqpBD91kfWFPrkc6m2LiZAlVgC996RN49jScTUvVLYPXt+rgTDDRmZaF63rYlqijp6VBWAsjGW9uzvH8C6/yL5/4IuvWzeFzn78TIp/W2TXSPePMXOixaF2OecszzF+TZ8G6PIvXNTJ/WR5buPTvjZkYjmmZaRFWLA5sq7D92TK7XyxxdHeZM8cDipOKSklRLcNAf8ChXSX2v1oiJuLo7hoDfQGH9xYIQom0ipwZP8BNr7mYtcvW8vu7N7NoaTNeCvbuGyabdfl/X/8Ae/f0c6y4k5Z2nzjUp8pyMWaqUKGx0WX5Od1s3zyCsBXFYsjYaJk4UtQqVRYvnss3//tDeC0H2XFkM5Wq4sSxCnu2ltm9pczeV6oc2FXl8N4ah3cH7H5xisiKOH4gZPuTRQ7urLB3W4kjuwL69pYYOlmjNBXg+haptI3n6VOjkuDY2hhaDUPOFE+SIkVn+0yuuvJ8amHEr39+L7v37OXyyy8im22gFoR1YGLSysa0V7QiRy8fWtkoqNZ0Ds+mrS8wc8ZMsg1pKtXKWWRYPVtDRrjpPFa2m1LfJqTlsWFJK/e/NMzBUwXmNMe8sGuMrrzFM0cCojAkm4KJyOayVa1Ulcu+I2Ncfm4HB06UCEKHi8/r5is/3Mo1F81iRmsDsxf1cOp0gRdfPEzFs7i6vZE/7TrGa69bj5fxdZvetomlItWQpq25mU2bd/LRz3yPv9z3JK+95XJSKZfYtG6SIfjsnjZe3neCZ/vH+PX/fZY33nEpk6USlUoZ27G0ATSSzJrVwaH+Yfbt2M+Fc2eyZ2gcx7GpxNLAEqdbwpaw6r4pDNvPMp68JMOljkgyEEvLcerrdJIlVB85CLteqSSH7jiKjDvdMmNoYTpKtvHa6Vak53r1iOjpPcBQKcz8xve0Ete29SElljF2tqHhLhlJI+PDmGam1VKW0HLds+W7KmGnGM5RKuWbsBrOWni0XnyaFybq3HutiY7qdMhEVqaVW/r0qpDGe2LVNyB98636EFgl3BaTC5KoDyIDA0sc7gpI2TY5x0EIybzmLJ6MGczn+dSHX081rGkpX6wlrAiB46bwLYtfPLUDP4h47dxOmn2fU0EAQtBg28ZRL6hFsYY6opCWTvqqBpKtkwUWN+a5uqmZpkqVO+b28rPBM4yWFSPjU3zg7Rdwx01rECrmC998ks1b+/jNVy/lj48eYWXjFD3tOTqaLPoOTTBZjClWI0qhBbbDjFl5ToxBV17x6btP8U9Xd3Px8kZCLHxH8N0Xhzk0HpJxDGtM6J4pQk0bh5Q0IVGYbBS37ukAjWxuaE5THg44vn2AL3350yxbspByuVKvDhF60B7LWGcwW06doea4LrbtMDx0hqlCgXy+UR8aTI50GEakMykGhsb44Ec+Cyju+uLrWbJwLuVomB39WymXatg4KGnSL21BLueQ9l3yrT49C1OEKmbro+Mc2x1w4kANGQu6ZqdYdVGeleflWLG+kRXrm1m0ooFFqxtYsjrDsnNyrDw/z9SEJN/ikOtwOLa/zCvPjPPy0wUO7p5gtHqAi66YyVVXXsyPfvAMh/pGCIOAW68/jzffcSU///VDqMYBmhozBDX9zNWqiqAWUy7VmL0wSzbv8cpzw+SaU6QaXIqTNVaunM13fvRWTkxt44FHX+KVTTVeeGKS40eqSBHT0i3oXmAzc4HD7IUu81b4zF6Ypm97xIp1Ka54fZ5ZizLMmpMi3+ziOA5To4rjuwMOvFzh0KtTnDlVAyVo7vBoyHkIYZH2fWpxjVI0Qq08SSqV4oqrL8B1ff7w27+xY/cerrryUjIZn0hOKys1vSEyLWe9uEklyOeaePyZv/P//u+7HDhyhB///m5GJ0bZsPZcPNetRzXorB4Ly3GIwwp+40wUisn+V3H9DMtmpfnfh06z+9AEOREwXqjQlYMDBY/tpxUrO21mtQhamlIcOxMTBQHzOtP88ckB/unGHp7YNsTjzx3nba9fjQoCli6dzUOP7OZ0pcaGmS2cHhyjls1yyYblFIoVLNuhtTnP+GiBL3zjN/z8h/fxT70tnBib5Ewt5tqr1lOrRWYt08Y/KRXnnbuUd77xKnpntTI6PqlTScVZA25LIKOIhfNm8p1fPcols9o4Xa5RjiJqEoLY2BZMKymKQjB+NSVjPamLVZ2o7Bgzrc5Goo56165z6jNi3WLCKK/QaWRCaSyOLeoG4ORnSRIToyTiVkAYhHXck7DMzMvE4NqWY1SyRomrBJFRc9kN2Ya7HFdn5zoGBaKUrLeOprlVGg+SxM6qs+YUURjpE64x8YWh5uTIKIG4JT94XHeqR3FsMAEmz9dE0E7nBsv6CUYlclimc+DjWEdyOo5DZH4wvUHFunVlNNOJD8S3wLcsylHEFXO7eHhfH2/+p1s5d81CpqamsMAMzy3iICTT4HPlFeu49ILVPH30NL/cvJO5KY+betoRQnAy1MlntgTL5BMrI6ML4xjHdXCFxb5ykSKKVi/Fn0dGGY0UQS2kZ2aOb951B5kGmxe39vHPn3uQr374POZ2ePzpkQO8cX0OH8XRE2Ue31Nl67DiyLhFsVBhw5ocVjbL0FjEb58fo7PJ5zN39jI2FdHoO/x13wSPHS2S9x1wdAUiLH3Cqb/MiTAByyiuEuSLfoh9z8VPpRjvL3Fs+2k+8fH3c8N1VzIxOWlOLZiMGEzGh21yO1QdDhfUQvbt28Oxo4c5ebIfhaKtrYMoCkzOi958PvbxL7J9x34+8bGbuPnGDRw/eYIndz5KqVzWQzuTv4w2UOO4Fl7G4vDuIn+7e4TJ4YhLbuugVoq59s0dLF+fZUa3i+9bCJNXXZ4KiENJHMTISCcCTk1EDJ0IufCmFlq6fOatbmDh2gwt3TanjwU891CBJ5/exZylDk3NOV7afJyerka+8Z/vZnysyO8efJDe5TZKChzfwbUtpoohYSApFkOqpRor17cTRw59eycJaiEr1rTxybsu5sEHX+UPv97F0KmQhhbJ8gsyLD/PZ85Sj6Y2h2yTTTpjk0pZ2JaisdVh/gqfHU+VSKdh1pI0vmvT0eMya0GKecszLD4nx6wFDXi+z9DxkH0vl9j3SolaRdLQbNHakSLb4OLYMaVwnEL5FGMTA1x57UraOtq5/4/P8sr27bzmNdfofO/EY2C8YAnWJuE2pXyfvYcO8vdNzzJSHKOjvYOde/ewdecOLt5wIbl8A4FBm6O0ukcpQRRWyXYtJZg8QmVyiN72LEEU8sj2SW5Zn+PiZTnm9+RZ0irZNyyZ1yaZ25nGT6dwbXj14CRr5+XYvHeC3nabC1d38+3fHWBGs8f562bRNKMFpeDRx3ZScmxumNHM77cd5OZr19PclqfBcfj9vU/x0S/8mKb+Qd6+pIf9Vcl5Xa38/KltXHHpubQ253WgmePW0xszad2yrdZq09HZRoyCqajDWkBvbyfHTg6z/9UDzO9sZc/QhEbJCC1QELbQpk1j1o1jg/kxoXqadaU9U5ZRrmozspk7mAF40j1IZNRaPGQbPFBcl9wmHaVpztY0eVfKGKSZawhBbKCJKFG3a8SmfWVZQhPITXVq2zZ2Npe9C5NDoAdkeldLpvQJXTcJQ0k+bOKwNJkmdXSJEEn8qMa8x2GI42lJGpaWV+ogKo1D8X0PWScB68/hJoMcw8RK3Jv2WQofTG912iGvy8EojrFtDROzhL6AadfGA1xLML8tT5sDe5TF5z/1ZuIoMJWXMczEoXbLxxHlcpWOjmZuu3EjK5fP4/7dfTy25xjn5Bu4srOFKSkpWJASAkfYhErf6Mj4I0IUcag4Uw3ZVa5QCmOIYapU4eufvZHVS3sIaiHv+tifmdHq89UPruP5rScZ6B9GKMUPn57i3t0BkeMyXLVY2WHx2vNzzF/Vw9BwjZ88McxkTfHNd8whQuBZFscLId9/eQzPtYzsOZFdaqmga9u63aj0NXJcu461t10X17expGByuMLJPUNURwM+99mPcOst11CYKNTlhbYlkEyjSDAqHf3g2riex5kzg+zatZuGTAPFYomGhiyzemYRRRGxjGlpauRb3/kxf/zDA9z+2o189F9uZHKqyBM7/sbIxASe7RLGMaaFi+vY+Fmb0TNV/vaHEQaOVrngxlYuv7WTTCbF4R3jLFqVpjAWUJyq1QeKCfhT2IaKoMD1LJ65f4j11zZj2RZTEzVqRUk649I+I83ay5pZsDTPwVdL3P+n/Rw8cIY4DHnL2y7h5ptX85vfvcjJaBetHR5hTRmWFNRqevjt+hYjIxWmJkMWrM5y/ECJbDbFxmt7uO/3exkrDrLmUo9VF3rM6EmR8T2EBKks4lo8PbtTAkvp05+dEfQsdnnq9wUauyxSWZfiREAUxFQrWhmTziq65zisOL+FZee1kkl79O2pse2ZAmdOVLFdl5NHahx4tcaLz4yy6al+nnvyANgOzU15nnvmFc4MneE1111BGBmRjNQhczJS9UXL83xq1ZBv/vB7FCaL5BobiFVIJt3AyYFT1GpVNq49T8+/jKxXmuFsFMfYrovfMp/y8ReoViNWz8+y9fAUO06EXL48j5PyyWVsWkQJCTQ4UAtjmvIpTo/HNOUcJsqSwydLvPbCGRwaqHLf44d5w61ryPgOK1bO4sUtfew5Pc7GOe2MnxnlaLHK0lkz+NC//4i/3/sM758zg+62Jv58Yoyn+4dY3N5Mo5Q833eGW16zkUq1hoojozjTPLck0rkOlDVpfpbB+SgU6ZRPT08Xf/3LMyxuzXNoooxCUIslgWFWJfPbJLAtimIc19ahUknUhFHEattD/A/xskl1kgifLEfU1zBlOFYakqjqSYQaFW/XZdj1FrQ4u/2sZyqa2Kv5g45tlJnG3S5lZP6syURPHIK2pfub9ahLFBYWrtkAojjGT/lYQhAEgf5AnvaB1Bcj075yXW3d1w71s4oHBcLRVUnyNSJprZhsEaWoZ5XXM4Nt2wS6q3rCXcLaEpYeBLmeW5ftJoRXVEzWtml2PcI44qp5HTxzoJ9rX38NV1y0gqmpIo7rmbCryMiPHTPHsSmWCoRhzMJFs3jDzZfQ1tXO3Vv3sKtvgKs6m1nX3sTJSo1JKcm4Lr7tIFwtKw7DGGVpebKFTgEbKxR5zZVL+dgHr8J2HX57zyv89Ldb+NGXr6A9rfjZAwd5YOsom45EdDYJ7lzn8sZzfY6PQWeLw4bzZyEcj93Hi/xm8yj/+Y75zGpLERsJ9v9sHWW8qvAdSwMRhb52jqt12xhFhWVc5yqWJq1MYdkOo4enmOiv0Og0smrpYj796X/m0os2MDk5qflU0gza1bSiI1HiKWIc2zMgREm+sYmOtg5yuRy9s2ezYP5CY6aSNDVmeXHLq3z+rm+yZFEPX/jSHXiez/P7N3F8+AQpxycKIpQliGsSz7dQVszWJybY/OAEi9ZkuPpNXTQ2uriWw7anh5kYLzFvRZqpQoDrOzjCIjDxnG5iRpUW+WaPv/9piFnLPJo7bIrFKjgKxxa4no1jKapFKIwKiC3OOWceJ0+ME8mIN73tXMaHq/z8zw+ydKNNVNWblOM4uI5FtRpRrcXYjm6JVmsBp/uKnD4cECvBy8+NUJyqoKRD/37JYF9IuVakEpUolmKCWg1p68osiSvA1ifTqAyprEs6By8+XGDReSkc4SAtjYxPBCRSCoIgxHYkXfN8lp7XxIJl7QydCHnmgVPsfHmUTNbjqpvmcO76Hi5av5QL11/A7a+7muJkkd/8+n78dIrLL7+QifEJrcSUOo1QGwh1B6Ax38zewwd58YWtpFNpbFdXo7nGBk4PDXLp+gvJZbNUK5W6j8uyBI7roVSAnWrBcl1GjryI5+dZMMPhV08N0dXqsrjVoRKCIyTFYo1cWlEoBjRkUigEpUqE7zu8fHiK6y7oYvGsFv7wZB+TkyWuv245nufSks/wxN/2cCKIefuCTh55+QC/uv8ZFpeq3LlsNs+OlHhubIrJSpUm1+Pi2TPoymb41r2bWb92IYsW9VILA22oNYdXPUuw6rYDFUuTTgpxpGhta2J8tMRnv/YzemVMo5/i0MQUkYBQJFVEQuCI6+0l/VtWXS2K0JnmJMmn5nvqJE8tIkqMidpiYdWVrUnXKBkbTDOyFCqSOlrXdGbsxAuSFAYGRuvYTj0SXJjqxjY+E6mmlbB2Ppu/S6JbRCIhogqITZnqei4oocOahDbYJaYvy9a5FkmvNJYxYRxhkSRdSV1JOHZdDpiY+qalutQHsEnYiU4/VJiBRH13VFLVJWJJ6YZxzdaVB+a/URRhAa5lkbEtXKFoSrvMz6bZUizxxc++C6VizfWRappIa6J3VQIUc/TNrdQCYhRrV83njbdcRFlY/GTTLkZGxrmlt50FuTRjUjEW1MjZLlGoiAzTS5NMLYJagGML/vurr6VzRhMDp8Z424f/yD+9bjlXn9PM7Z94mlcOjHHt8gbevsHj9etzzG7zGZiUFKRDNu3T1ZUjn/f48u/62Lgwx50XdjA6FZL3bf5yuMyTh6do8BMjkQaKOp6WSLuuQxTE2K5lZlo6kljFEsd3GT44SbPdyBf+/V94+1vv4KYbr6a5uckg6/UMSs/XJIIEDa1NZAiBpSwkcf2Bc2ybfD5PU1MLjY05jaCRknTap1oN+eRnv8LE+CRf+eqbWL1yAdsOvcz2IztwcQml7kG7toObthgeqvHYr4YojARc/44OFi7PUS3FRLUY4cATfz5F75IMHV0+Qlp1NI/j2XiuBZZGRfhZwba/TxLIiNUXN1KcCLE88F2HdMYlrCl2bymy5fFxtj55mmNHxmmfkWfmzFb6jgxx9MgoL7yyj4UXlrGwiGoSz7fxfB3tPDlZI6rq1M1Mo83kkGT74yUKIwF+2mLOygz5Fo/ylGJ8MODU4ZDBw4JiIcbNBsQipDgZMDFZRkpJOuXieTorQkmQNUXn3BR9u2sUxiJ65vsEValloK6r5bOWwHEFkpigGhLLCK8xZNXGVmZ0ZjhxYIyx4RrDQ1VWbMjR1ePTkMqQyeS58OJVHNh/mt/85n5Wr17O4oXzKZVKuK5nciX0UFmLVyLWr13HCy9tY3RslGw+baa+FpValXKtwuUbLqIWVgmDQG9yJl9IKUEcVEi1zKc2epSpkdPM7m7FteCPL4xxzTktzGhpwBUhA6NlGvMeoxMRE5NV2toaGBwLSDmSLYeqXLaiia6OPJl8lh//cTdXXTiXzq5m5s3v4uWt/ew6dIaC63F7TxOrGhuoOB73nRjhdKlGFEVc1NnK6+Z188jBk/zx2Bk+9O6bufyyddiWbq87piLQB2Vlkhtkfc0DcD2H1tZmHn38Zd7+ga/TMl7gNQt7+dvR05SjmHIMlSgmDCIdA66EQZHovzeJ4U7ijGVCAlZoTp2FMUNrQZEQlplvxFpOjf4smo1pTNUm9VVKTcFIVFQywcFLqcOjjDoySbLUvm1Vt2ug1DTy3Qzr9fDesLASsJqw6vF6JoFLGgZ+iON5homDkYbpIbWS+mI6tlZD2UmsrYmwTf5uvWvqQZz+c3E9cztxs6OoG1cSPLRl2fUefdJflzIminXmgeO5uI6rW2RJheK6BrAo8G2bnCUIZMwV82fy4sHjrLrqAm697nymSkXjuBd17buSscGv2GAphOUiZWjwBqF27tpw6cWruPHKDRwZL/KLF3bhhCHXtOeZmU5xtFqjJGNSjjbMBbUIx7WYmCzzzjeez+tuWI4jFP/f1x/h5Klh/u2tS7nzM88xM1PjU1c1csEcl44ZGY72VxAqpm9cseH8+fSfrjK/J8WzuyZ4ZMsQ//G2RSghSDs2B8cCfvLqOL7BeFim7PRSGormp1LEcYTrO/Vcet/3iGOJm3KYPFmhOhTyre9+ntXLl1OrBVSCgDAIcQ2VU9iJas6adiTX6WPCLCqybkKM45g40hh7XWZbhFLSnM/zP//7c+7/y2N86APX8+Y3Xs6OAzvZtO95qkWFn7FB6e/lZgQvPznGc/eMMX91A5e/rg3Xsokjnc2eytj0Hy6yc/MEG65tJpXydDaGa2M7gIgIopjhoQKVapWDOwoM9gVccH0zlSntQrYF+A0OJw5UeebeEfZvL1CtRMxZkGPlOZ1kmyxmLfOJQsWhVydYd5si3+owNRYTi5hyUVILIgrFGuWCXiAk0L87ZM/TVbJ5h7VX51l3TRO9yzJ0zvNYujHHonNytHenmBiLGDgqKQy5qDiiqVtBpGWdU1NVZCRJp10jiQUZOTQ0WWx9dJx5q3w9sDY3w3MtYqUXikhG2I4+PYe1iPJkmZ55PsOnYnoWZhnon+Sh353Cb7No7CqRwsFPNXDhRWt5/NGtPPbYM9x807U05rJEoVkjFDiOSxyHBEFIS0sjixfO5577HyDXmsNv0HO2dMZjx969jBcmuHjdBmJlwuaMwkf7zfSszG+ZR/HESwSh4twFOV7YP8HWwyWyImCqFDBVCpmcrKEsi8mqorPFZ2Bc0tHscv/WCbK+zdplLcye2crft5zk4LFRbr/5XCxH0dGa4S8P7qQsY14Zr7J7osLhQonJcpklbU28ddEswkqZb2w/TMuyuXz/ax/gttsuwbIhCKr1GYVCd2aEea6FEkQyRiFobmpkarLCp770U374P7/n3ct6Wd/bxZ92HWUyjKgKi0ocE0vLnPBDXM+tI9jrBF1LbwrqLMQIxuibJFVivHaYw2niUlcGG6UH6ZpflkBuBabFr1R9c9ItMqeOmUo8PijdyYiiuN7OklIZfL9Jo0zUXFLpQCnHsepJY7Zl1ym6yvTILEuYwCdNUq0z7U2SHOaDGVWmTuMLQhBKl79Ky3q1/V3/na7n6HCb2GQR2NrL4TjuWXwkgxB39ODGSlzttoVjOfWNYjpkTdRVBLYlsIUiJfTwvDHlsSKf5unhUb7w+feQ8m2EkPXwFcu2cWxBKuUTBwpsU5UJC8t2kXFowmV0RVEs18g3Z7n5ug1sXLeC54+f4ffb9tHt2Nze1UrZFpwq1UiZmYOSipamBv7r89fT3urx0rY+Pv7lv/Hpd6zim7/Zw7r2Gp+8tpVa7HBmrKxT9YyC6nTV44rL53Dg2CRRFPLVP57gvVe2s2ZOjkoEKJsfbh9jtCxJ+Vo2KSxLAxGldpBKFZk0MYWQ0w+JbVkgBQN7xvjgB97OJReez+jYGI5t0O62bfq+mn8UK11ZaNm1qBM9ERDJEMd2tUcnAbYZNZ5UmqnW1tLEphe28sWvfJvz1i3kE/96E6NTZ3hq55MUKwEpz9X3Ie0SKcnDvxrk+L4KV7+5lVUbGwkrOs8limOIFakGh00PjmI5irUX54mrEj/tEIQR5VKN4eESI6MFYhkTxdC/W7HqsjQyVPWXJJPz2f5kkZefmGByrMbCZU28+d2rufT6mXTMEXTMcXEbIuasSuH7NnufDUjPkDg+TAxCuRpRmQqJYoWbspgcVex9MmD0pGTp+SnOuSZPe5dHWIrqaBsbQVOzTe+CBpatywMWx/dMURhwKE/atMxS+CmIq4JSOaBajkinXOJYEVZimma4DPVHFMZiepf4yMgkd5oqwRLa3xOGJkpaGoNhpGhq8znw0hTv/Pd5CAF/+Xk/x44WWHRuGsuKyGcznHvOOfzhd49yvL+fW256DYHpQEg1De9zXIdSqcSMjnYGh4bZs3c/zW15bfyVCle47Du6n1TKZ0HvPFKpFOVqpU4BRgqCWgU/246ybeTwDnK5PPPabe5+epALlrXQ3pJlPHTYd7JKWK0RK5fOJj0DiBUMjpbYdazExGSRa87twM/m+OV9e9hwTg+9PU30zmrmlR2n2X/wDFnfQUURgYp43ZK5LM1m+PmeI+ywbb782Xfzbx++k0zOZ3Rk1DjqE9DgdNyyZQ7ZUoLvujQ1ZrnvoRd4379+m8yZYT572WpOTpV56HA/lUgSCMFETRuOlZiGvp5NQ048HEmSqxalGLSTZp/UGWWu55jYZ93ytxyrjlw/m7IsLJPi6jom4m/a35HEbQgs0xKTRkJs1SuQxA+m11+DlzemigSHgrCwGxqydyUslKS15HiewY/omUQSU5WUbFJJEhBJEIRmjmEyfQWkfV/38KSqg/XCMDLqA71rhoE2AsbxNGk3ibuN47OGJqbKkSYr3Xb0rCJREGhKpVnoDMdHvzASD0GTq7Eka9vyHD01yKwLVvGG26+gVCnX838916VUKPKJT3yfSMLqNQtR9UhfBWeFvygVk8iea9UaU6UKvXM7ue2GjaxZuZCHD/bz4O4jvKm7jQkZM1iLsC2L8fEy73/7hdx4+RziKOSDn3mYjiabUgByfITP3tzOtgNlanFMR6tuS1QDyfHhmDlLupnVk2dgqMhXfn2E3hlpPnRdF1M1SWtDir8eK/H3QwVyvoMS2vzjep7pEiapZ9Y0wtnc/DiS5PINDB+eYFZbF5/85Acol6qclQmmH1ZTtmsnq0Mkz4K7mf6wnkfpCs62HOI4mL5mEjNcTFGrhnz8U1+hVizxuS/cTvfcVl49+hLDhSHSXhqBIpN3GRqs8Nf/G8BxFa/9UDftMzKUJ0KkiOuZCX7aZXws4IWHxzn/mmZm9mQYGysyOlFkvFBlZLSEUjGplEMqZXNgi2LuKp90gyCoKhpyHjIQPP+XCY7uLuOl4bJr53Dru+YTOQX6B09zemCMiakSNVkjjENmLXfJ5G1efTwgIiDXbuG4Aj9tEVQkx14N6d8OzV2w7jUZZi1qIC4rJsdDSoWYSllSmowoTtSoVmCqEBDHMQtW5fEzFhNnQopjMHbCIpVV+I0xMrSplmPGx6vYjkUq7ZB2U2TbLF5+fJxF63Jks049hTCJKdUdATMLRPfog2pE19w0h1+tUC1GXP/WGczo9Xn12XH++vt9zFqWYu7sRuZ2z6W7cz7f+/4vaWpp4eKLz6dSqWgVjmZqm1OwfndvvO4aTp04w5ZXttHe3YqMNNLE9Vxe2rmdJzc/Q0/nTNYsW2Xc1tPJe1FQoqV7FZROMzV6mvk9LRQqVV4+WOCjr1tMd0cTM5rTPLRtjKhSYX6HjyMDxsuKjO/yxvUZfvT3ETJ2zK2XL+DJV4Y4dHiI225Zi2Upersb+fODu0jZNuU44o0LZtE/NsUPDvbxlrfewPe/+n6WLZ3DyPikWcNMBHOsXf71YbeS2lcRCzraGzl1coRPfP5H3Purh3n/qgVsmDuD+/b2s+vMBDY2hVBS0jpcYkPJSFpKCRi2HjxnMspVLI17X4tgEvWp53lYliAIon/I5KhHeceyrrZyXFcrW60E0Z90+qfBt8kzkaDfE25WXR2rpjNKklzFBJWiwBDSFXa+sfEuZRQSSYpVbKB3li2IgqCOfq5nNiQ7WILgNsYw7Up3jYHQ/H/1RQZkFNelwDqPV/sukpwQSKoQuw4eC4NQK4RsLf2t77bmRUHqIb82OeobHUcxjmWRdz18IfGEZH17E48PjPCpT7+FjvZGglpFD7qDkObGPA8+uJm//vJhnnt2O9t2Huac1Uvomd1NtRJorIHSTlOMRFX3/TVsslouU6nWWLiwh3fecRXjpRqPvbSXy3s6eHmqhiuhtTnD1z/7Gloabf7+7BH+8wfPc82F3Wx++SRfv7ONvUcDXFuxfHGaqaLi4GGd21AILRavbqe5weO5V4Z4aOsYX3nLPJoyNr7n019UfP/FM2RTXt1M5Pm+vgYmwMYSAtc3Um1bl61xEOL4LsXhgJN7B/jXj7+XhfPnUjN4kkSckAzfNFtHp6A5lqt165bu+8cGGW3VVU+qLvFNUCgCQb6hga/95//w9yc28dnPv5Yrr1tF/9BRjpzZSybjE4eS0lTMq5snePbeIWYu8Lj+3TMQoU1QjlEOSCVwHYEnbIIoYutTo5TGFedcmeXEsXFGJooEtZgglCbvxCKVERzZLmnpStM916Y8IUmloTAILz5UYORMlXSDhaWgebbNKzsHeP65Uxw5XmJoKGJyKqZajokDSRAp2ua4dM916d+pOH24SlCD0ROK4QOCXLPN0otdZs5NUR6W9O+r0He4xqnjISeO1jh+uMaJozVOHAs4ur/CqeM1jh0ocfJImXyrRdssn8JoTHE8ZOiYhZKCxhlSV15KUZgMcF19Ddq7UoycDiiMKhasaaBSCJFCEYcawCcNMVkfwPQp2vYshBJkmh22PTnB7MVpGtIu517ZRK0c8ecfHyHbLJi/IMeGc69gaKjAL37+By6++Hw6O9rrc4Dp2DkNwIul5LKNFzA8PMYLW14hnfZNhK6OrC6WS2zZ8TLDo+N0ts2grbEVra+RtLd2cPhIH798+GnW9cT4DTlW9ab59RMnSHmChb0dtHe005yKeW7fJAubFbm0zVhFgu1yzqJWlrRH/OTvw2xcnKV3dhff/+2rnLeyk3m9zXTPbGX33mFe2n6clZ0ttKc87hke5b7vf5rX3X4ZE1NFiqWqybcRhkel5cuWcHSgFrol2NTYQMr3+dndf+Ozn/shCypV3rVhGTtHJrh/bx+lakSMTSmGMopKDLUgMG0mUW9hWXXcjzS2BJOU6thnxYEznckeS6MwtU3lYgbbJmspjqfNgsmBX88sk5kG02txYt5Wsr5uKymx3envq2fXTl11Fxt1qy4WplWrdoNhYdnC0rLMBMiVyB8NqiCOpwPZtcJDh4uQaJCN41wY455EGfWtpaMzYx0L6ThOXVZpCQ37so1dXxqpb4JCjs2JN9E+O0bqm0DbtGxUIk1FkkAalZI4wiItFBaStS1ZCoUCDSvm86633kC5Uq5ztmzbIg4ln/3sT/jAxUt4w/pFPL/1AD/+7d9RMuacdUtIpdI6ntfSZkuMUU4JgZ2wYVBMTU6RbUizfXcfY3uOks1m2DtVpVqs8pH3Xso1VywkCCM++G8PMaNJUAgcNvYKZuVszowFXLIuz+hIxOEjNYPUDykLwdIlzVQLZb74yyNcsqKJ285tpRSA77l8f9soZ6ZifCPbdT3b6MgdUwrreYWUsr65Cxsc3yWuQf8rA/zzh9/Bba+9nqnJqTpmIWlRUvfmCBO5adVL4YTRkwDcEoimVNIo8mQ9brNzRge/+/MDfOfbP+Hd77mcG29fy5atu9n0wi4O7ZniwMslTvTV2LWlwLG9BVIZwdVv6ULW0FwuXd2TzjhENcnpUwUGz4xz6CVJ12ILkQoJgljPeCx9MJEGTXPkVYmXspi73CEuQ2NTmtN7BS8+PK4VNLaiVIuJHYvhI2X8MyGdFZcZUw4NE4LqmYiBkyH9fVUKAxHBqMJN2cxZ7ZJvdpBVm5Y2l8UXZGie4TB4MGTfriqHD1UonQnJFxU9gWCecljheax0fZZ4Pr3CJluUqCnF+JmQA3tLTBZDWvM2lnCpBhFjJxTlCWjp0QmRSkpqtZBaLaK5NUP3vDQvPTJBe69NU4en42ktjQ4RQoDBCsWR1OIJS6EiQWuXy/6XyuRbHbJNDuVJyaJzssxZ6nP/3Ydw3BzrNyxgzYrl/OX+p9h/8CC33HwdtWp1WtGT3H0rydyGyy++kKZcKy+9vA0v4xDEAXEs8VMeUSzZcXAvm7e9wOFjfaxcvBIZSX577z38/Hd388y2g8ya0ciyGT6phiwtWZsfPHCcq9e1IZXDgnYP4ohtxwosarfZe7JKV3sK37VozfoIAY+9Osr7X7uQzbvH2L53kNffug5hWXTPaObeh3bS25pFVKr4M1r5yAdu5vTgsMkNt3WctqkUoigwcEnbhClZtDbl2b2vjw9/8gfseuxF3r9mPs35NH/Y08fJqQq25RAoxZSUTMWSSqj9UbbJNNdJhwrP1aFpwrZ0dpFt1avqZE5hYWkulVT1tRWzFms2oHHUG1mvMIgny4wCkoU/gdYm+KiksnBcLcyw3ek5tTDiJy0KcM3d1V0FjIxZmJFF4j62cw3Zu6jn4doIrRWrL7D17HPzaycptUyOujBJWhiDkZfyjZrKLGJKK6Q8w0YyVHMzRIrPIvLq7x+Goe4FJpuXUYLVJcEYKZk5IWjnvP58Wm2gyzBXCHxLX4CLZraw6eQQV910KRdfupbJiQJKKmrVGq3Nee7+w5McfX4H53W28cyOE9x+7kJWdOT56T3P8uBj25g3p51FS2YThToyV5nPmsx1pNQXPpVymZoo8skv/YxbOlp5vlplvBTR2Zbhix+/imzG44mnD/HfP36Wmy6fz0vbB3j7+Rn2nQjYsNxndCjkWH+A60gCKWlqTXGiIFm1MM9jW0d4cucEX33rArAUzfk0T/SVePRggQZfgG3juo4edsUxjqeJmpieayJEqM+HHJv+ncO89fWv5f3vfwtjo+NmE7YM68bQCFRCYVVGIi2MuMGqy7L1pqEJA0rGWLZDbGSGURTR2tLC08+8wCf/7atcdfk6brr9XJ78+w76+gaJo4jmGT4Lzm2gWAgIqorzrmzFTQsWrspBJLAcTWUtTYWMjJQ4c6pEsVJkdMCiMmqx4DwLS1p1nE5QleZ5EowPKFTksGCVj2c5OMJlx1M1DrxSJLZCRkfGcS2HC3ubeGtXljfNzHFNS4r1DQ4X5n0uacpwUUOK89wUHZZH30jAvmMlhg5VqIwosi0urT0uKd+mb2eVlzYXGThepTeyuTWb4dZcmptaM2xs8FiX91mc8liYdljj2yz3bFZ6Npe1prgw77Ey51GtKA6NBqSyDr7nUAoCpoYUUyOC9rnJQFSgLMHEeIVcq43jCXY8VWPR+hSuKwhrWgLvWIIo0rw2P+XqxcaIUDIZm6AK/YcqrDw/TxxGVIuSxhk+6y5q4oVNR3j5xaPccP1qent7+Pa372bmzHbOXbe6nnWjksQ9qdsjURRhCUF3dxd/vu+v1MKI3p5ZVGsVwlC3vrP5BsIgYP/hw7y4bSsvvrKN+x94mGw+x4Y1K3li+3EuWZwj5aVYMMPnpX3D3PvcIAMDwzy7fZi2rMVgIaRSDpgMoDVrM29WE0fPRGxcmuOxV8dpywjWrezlf/+4nXNXdDJ/XgfdrVl27DjBM7tO8pr5M/jb9oMsXNLL/LmdVCqRaVUJbOFoL5y+0tSqAY1NjVjYfPv//spXv/wzLsm4XL+sl2dODvNc3zCWsqhJKEUxJQVVBNUwwnWceia6ZdmGrGsWYktXgsK0nB3X0a3fswzTyRzG8VwS00ccRWZWYlSpZlasD41CKytdTQxIDv2WJerRuq7vaqm30p2eehKh65o5t6u/nzk8aruIwnG9ukipnu0uhDESquQDq3poiYyVkY5p048goVcaXHuCxjAfUBh8e4JzTiz6UZggu5U5vVh16qvt2IbfMo1M0N4CkxUspblQ031Dx0TeJuTeBA6oefrUfzCUwrMsbCHozKRY1tnMzx57gfamDMtXLtSYY2JKpRqf/dyPuH1xL3tPjXGoUGb/0AR51+H2VfOoThX55k8f5uTgKOesXUxHZxvVSk1Htxp+mEASBSFNjVl+fc+zHNm0h3PmdvLY4Ci1YoX3vuNCrrhkIVE15KNfeIAFs/LgpMhGE6zrzZFOCRpTFvsOFLBtGJ6KGJ6UrFqbZ6KqPTj/cc9J7jy/mQuWNILjMV6z+J+XzuCYTdr1NXZbJCl/BtmMVOZBiQ05Rt/8oBTilnw+828fNARWY24yvdWElGvZNrGMTFStIFba6xHLCMtyjOBBg+liU/Eope+NjCLa2prZ8tJ2PvSh/4+5c7r4wpffTEu3T2unxbJ1eeasssANeereYYaPVXnnxxZz4NUivUsdmpt9KtWAyUKNM4NlJsarVKuB1sZjc2SLYv4Gm1xW0wgS2nBCbUZBQ86mvcsll08xOQDbHqty/FiBUm2KjmyWN99xNR+7egVXVEZpk7EOKTOLR4TSGRw25IRise+yriFNRtocmwroGwwZOFplpD9m744qlYEa53seb2nMcGtLA/NTFo4QlCVUYkk5jImk5MVSwN8nq5TCiJwFngJbKLqJuKwtw5KMz5liREeXTxjDeKVGNCWZGLToXmDjp2ziSGH7UCqEtPU6TI1Ldm+q0DoLUjlXC3fiWCPFlQ6EcmyB5YDn2hBZpPOCAy+Xmb88rWegIiYKtM5/xXmNDA4NcffPXuK1t25kbKLIn+99gpuuv5JMJm3mpcaHZRvZtNJO86lSkQcffwzX8fjvz/8new8e4MTACVIpj7AWYGHjOi7FSolCsUAUSd77lnfy9jvv5Ld/fYhyCOfNsIiVxYyczd3PnOE1axtZ0+0wEboEtYg9ZySnC5I1M6G3u4kDR8eIw4juVp/7nj3J+25ZzEsHC7y0/QRvuGEFlmfTmG3g94/soi2foimWvHRqjNtu2EAQRvrknaDOhaVZX65Ne0sTL716iPf/y7c589IuPrJuMQVi/ryvn0JVz5hKSlKUiopU2iwsBK6luyVRHCHRMmtppMCO6yKNkEW7xzGmZ6suIEr6RNbZVHPTPq9Ps01nJzE2RlGEZ2wXlmVSWY2T3LItLaoIjPLLJLRapqCwzqaIJCxE4yeJY735h2FUp6EnAgC7Md90l1Qaw+15LlIYlY4hbEqDeRbo3OBaEBg1gJHK2sluF9Wd6kkeR9LfcyzbsFO0mzhx7SamGa3Zl3Vjmoo1Yjg2qJIk19cxILEoUURZlu4LWjroSgcg6ZaX49i4AtJCMDhVYVVHM6vbsnzn13/jpW37Wbq4l2XLF/Cf3/g9U/uPsbSjmWdOj2L7LsJx6Rud4sDoFOtmtXHdgi4ee2YXv7pvM61NWVavno8CwlqoCcJGxhzVIj7ztbt5TWsz24MqZ6ZqtDen+NzHr6W5JcfzW/r4j+89zSffsYZnXz7FNYscZAQLO2wGTpaYKksGqzaFis3SRT6LVndQrSi+95fTeJ7Lp147i3JN0eC5/G7PBLsGKqQdCzflGcm1xrFEkTS0XbN4GIaObcx0lmtTngzJkeGWm67FwtaVomWZr1dnBbVYxmRq1wO9YhWZdp7J9lBorXwdLa1rw9bWFv7y18f59Ke/Si7TwJe//jZaZgYMjhxkKhrgyImT/Pmnx3j8d2fome/z9o8uZPRMjcMHR1h1QRP9xwqcHiowMVHR2e2ufmi9tMWhFyVts6FznktYVeBqQFwslW7VCIGlwPUthGVzcEuNnU9XGB2dpK0lw7+84za+/IUPc3lnmvLfnmFkcJxqLUQ42oCoAqlP+rHSwgnLohjGCKFYnPK4pLmBXt8hkIpKWXJdc5p3tjRwSVOatgabyXLAVBxRk5JiYBGivVWhFHRYFn6k2BfB8+WY0WpEtyVwXIBwXgEAAQAASURBVIvRSki3b3Feg8va7mYuOGcmJ0crnCxWsSoWYycFM+Y55Foswy7SVNfOJS7jA5LtzwZgR0R2jalyzNhYQLEUUAliSuWISimgWI6plSI6Z2Y4cbSC7Vs0NOnhuLB1u6taili4vIlUPuB3P3+ZCzeuYNPmHQhlc8nFF1ALDFnWtnV4myE7h2FIZ0cHI6MjvLpzJ2+/804O9R3h4PEjZBuyRJH2eLmOwEv5BNWAlQsW868f/AjKsnhhy7NsPzbB8pk+Pa0+M1ozHBgosevYJNetaWb+zAyLZ+dp9yuoIKCj2cNxXbrbG9jdV+LcBQ1sOhQyv8Nm8bwOfvCHXVx03lzmLOhkZlcTL7xwjBcOn+GOpd08vO0Q56xbRnd3s/55HA8Va19YYz5HUA348rf/wPf+89e8tivHxQu7ePjYADvOTCCwKEVSbx5xRCR0eJTjJpkpRl2KjWMMehgLhOYAKp2dZOIMYqkH5jqjXmeT6Bb/NEZfoeO9bcepr7XJAm8JsGzHIE2m02KTFnZC6UhyQpK1WZiE2cgIqBI2Xr3tZXpetvF2YTyDeqYjsTOZ9F3amzHdKzvbAZ7k5jqOOXXamuaqh9tW3ROSSM0STbOMkzaCY3rjArBNAJCs9+SmJbhmtmMiM23LKMDqcEZlLn5scMaxpgTbevCfQBbrrBmlzGDJIeNYHB4v0Jlu4JZlczh86BS/+vNT7D3Qz2OPvcgblszmyb4zTEpBWSqKocJ3HcI4Zu/QJMPlgDtWzaHbtfnObx5n+94+zl2ziJmzOqmUq0S1gMbGDA8/8SrPPbSFi+f38PT4OGE14pbrVnHTtauQccynvvogOV9w5w3LeOyJPVy9PIfrwN6jBR7aXeWJ44pXhhRXrXBZt7YZx89zYqjKj54Y4t9u6WbuzAy2ZXN0KuLXOyfI+g6Wa8JeEu6NCX7BnGrMvKwOYpNS4qcynN47Sm93FzfddA2hgWQKk9JmYdWlmnrQZhPLsB5XWz8FCQeldMsqCSAQQtDUlCeKIn7wo1/xta/+kJamPF/79juYuaTGwMReTgyN8PC9fTzy+yHctOCWd3eyan0O3/fY/MQZct2KkyeKjI+X9M/m6k0qqMY4jsXgQbA9xdwVLkFJIRxV9xx5rq6kvJSN7Qomh2HX30OG+yWVYIo1S+bys+98jpvfcDNDjzzKgf/7BZVyjaYFs2hcMJvymRFkqHT/2YT1+A0ZrVoLI4SMCSOJ5wgWpm0uyKe4LJdiTcrGcSTFWBKGEkcIcp5L1nFoSln4liDt6b50hYisJVjiCjoswYEINgUSN4YOYTMZKaqBRA6XWTCrhZsuXsTEZI2dg+N4eJzps0jnBK1dFg0NHjLUaJ+ZCx3CimLfCyGDx3UMtd8gCGNBEMaEQUwtiKmWAipVSa0aIEXM0PGYjnkupakQqSCMFWEYMDleJdNoMXdZms1/76daUBw63MeNN1xDOuXrA1QSpSplvVUaRTFrV67k4ccfZ+O6C9h35ACHjh2uC3MUEtfzUVJSKdT4zIc+Sc/MGYyNjfHAIw8RBDGDUxFr2/U739Pq8Yu/nWT1vAbyaQfVkMcJqti1SSqRlqM25tNUA8h4ikIg6Ds5wfXru3ly1wQnBia49eqFOJ5ABoo/PbqbnvYcnSpmU/8ot19/AZVKCMLGT3u0NOd57PFtfODj38M+eJwPnbeQ49Uaf+s7Q6EYoJRFUUlqQlCKNEA1jJWRzJpUVPPsJNGytjutdKpHB8dadar/nKgrnvS7Z1peps2VnOv0ZoDp6ljTnKyEhi0l2Gijom39g+rSEpoWIuOE+KupG9p/Yk13HcQ0iFYIc/g3Bmsd8OdiCUWs0C0sabwVUimiIDRudKMzNi2lRHmlTxDTlYdWBRi3tVTGsaq/kZ/yp3eyODYXRGdRKCO/1QNerdqx7ITqKqd3XSnr+mYtG0yUYHqGkpBltd/AQBZN/8420L/Y/E//RIHTxSobe7tY3pbnpe2HuKC7jdEwZvdEkQiHEKgqRUUJpLCwLUGhGrFjpEBPPsWtS2exa89xfvyHp/Bdm3PPWYjtWhBKPvcfv+HChizH4ojjkxUcJfjyv91A+8xmtr3ax79/7VG+8s/rGRiNqQyeJqhF/OHFKbaejGhvFGRToGK4ck2O5vYsvqpw9zMTZBx439UzmCiGpDM+P351gjNTMbY4a+M0Hh6UPoFLQ+RMTH7C0XnxqYxHYaBC6XSJL3zhX2lraSGMQt1ySuZKRkaovS8eMg5xHU8nCJoHWCfHShzLNaWzT2Njnlot5MUt2/j8F77NIw8+w3U3ruJL37mRfGeRzS++yv2/PcZjfz6NJOKK29u58JoWsmkXx7HZv7PA4b4R3DRUyjGOA7YrEEKf/HNNPinXo3WGR9tMCxFZeL7A83S16qUs4kAhI5uRvohjr8Yc2x4jI0EQl7li41p++r9fpHPmTI7+8ucMP/QgyvMQjkPLuuXEgWLi8Eks19G+F0ugYkg3NZBtzpFtzeNnG3BdDxFHVGshEQpLKSoyRikLITUx2LMEVbRJNeNbeLZFyhI0pbQ6sKokI7WYVtdibdomrwRPFAIGwpiZOEzWtHemcOQMLS1prtswH0s6vHJmjJTrMHQMimMOvu/gpiW+rzN6WmbauFmFim0qBUFDM7geeK4+1Hk+ZLMOnq/d7w15h8KQRSZrEQSK8dGYWg0qNYtiIaZcifB9l+55Kfr2lzl9fJQ5c3tYs2o51VpVC0nMXBQhEFIRK0lzYyMowY9+8SuOHz+OkxI6MtU8o5ayGB2Z4L1vfDc3XHM1E4UJWppbKJUq7Nqznb4zAa1ZhxUz0zTn0hwdqvLKoXEuWZzHzTVDXGZwYIrmXAbLtpgs1GjKOQxPhjRmbLYdLnHJ0hyNbc38+J5dXHvpArq6W+jtbuG5Z46wZ2CSWxZ08+CWvSxdNodFi2eT8l1Ghyb4ty/+nN/85H7es6ib5V1N3Ns3yIGRApa0KUnJZBQzFUlCpdccqTBZKrX6JmFZlqlep0GGkUke1Gh7aQgeNsK2ITbkDuMZc119WJNGXSVIKk4LIRRhEGmvBnpckPCtpKy3DupVAoYXmNgtADzfMwfOaQWXTu8Udc5hQl1P2mDCfDYllOlWKOymfNNdVj1PXC9I08NSDILArrvJdYZvYjLRYeuO7RCYCMY4nC6TojjW6qUECS11vnIcxUTG82HkBeaDyvqp2bItoiCsM1c4K8c3acXEZqBjGdWJhiwmaV9mQ1S6D2y5DrYQlCsBO0fHKUvJmp5OhqshL54aQ1kOsSMohzESqEQRlTimbMySacvi6NgUA5WAGxZ3sTyf5uf3PcfjL+xh49rFbHp+Dw/ft4kr5vXw7MQEE2MVrrlqMW+641xEHPDFbzzGyMgU/+/jl/LY04d58qWTDIcO5y/wuPP8LGtm2oShxaw2l7a2NN0zUgyNh/z3/Wf452tm0NmSwnM9Xhqs8siRKmkb3LRvHkpX91VVhO/7GmBnWnxKSVzXq88zkIJjL/fznne/iRuuv4rJyQmEbUi6CXfM9EOFaVtpUmhcj+Q8Oy0yimPSKZ+x0Qm+97+/4Ec/vZv77n8ULxVxx1tXc/GVs3lh82H+/LsdbH3+DI0dDhdc28iaDY24AoqTAdghkyMxW54bomW2QMYaxaFCkKFFdcpC1jyiSorypI0MLUrjkrBmUZqA0pigNGpz5pji1CHJ/q2TnD4SENUE6QaHcqXCRecs48c//ApWLDn83f9mdPNmlOcSV6oEMQTjRW0Oa8oRThS1cinUQLvqZInKaIFqsYoFpPIZ0k1ZHYEaaj+Tbdna9a10rIFSgruOTtDmWvT6LhWjMAyNEqY57WErQTEIkRa0WDarUy6FSNJsCyyhOB0rAmVhnxyhoaWBK89biIpjXjoxRlPeozAZcXR3kcJQivEBQaUIkyMRqQZBboakfbajg8KUhYpsCGzCisPkgM1Iv2CkT3Bsd5XRoTL7X57g6J4S/fuK9O0vcWJ/hVMHq0xNWkyMaYrAknVt9O0uMjFZ4MbrryKMtJM/MrQByzwvALUgZMO6dZTKJR574mkszyaQIaVilWwqS6lc4fw15/Ev73kfk1NjqFgSRQErly1n0cIFnB4Y4Nmdp1g+w6KzI0dPa4rfPnOGRV0pFs7MkMtY7D88hFSSIIwYmazR05nl1JkK3W1ptvcHXLCkkbm9zTz+0iCnhyvcfO1y0vkGqsWQR5/cR0M+xfrmND9/bg9vuPEi7n9kCx/42PdoHh/lgxuWsGd4nMf6zlALFJGEopIUFUzWQs3xk4ogCEwKqmFl2ZphFoZBPQxK/1qvZUnnRjiO2WhMp8ey/wFOG8vYsLcM7dyx6+gRKbWhMAoj03kQdUWcY+Yq8ixlVYJpT/rS2qoR1anZqk6IUnWCb7IXCOP3sYwSV6GH9knuuujs7FRnG1uU0rCtWjWoww5jA9xzXZtqNdDKJzOzcBzdaw9qQR0XLmWMhV1P2Ur+qWel16mtllE16R3NNSlcGlOiXzR9avKMfV5/8JrZ6ZOLZAm9gVimDIylxHM9HQcqpGHFgC9jfCXJ2IKpWkRk6YthOy5TkebMBFFUP1VJBb5nYylFRkKja5OyBEJIlrflWdPewlNHTrJzqsykJbizu5O+ao3tEyVK5Qp//PG7OP+SBRzYdZwLbvkJn3zbCt501QJu/cSjrO2WvOeqdlKeZHA0YGQspG84orU9TS20ueXSJr7y+wEOnQz47rvmUAz1hvGVTcP0j0f4vlX3aug8FM3sSmSAqq6gNnMix8b1XYYOj5GpZvjZz75lqhfqLcJpNLQgrj+AyX2zNTsMhRAOKG0mjKSguTHPR//1czz2+DPM6Gxn3uI8c5fkKZdjRoZGybYIZszyaWx1cC1BbUqiLJ1wJzwgcjm6q0pDT0Bp3KYwKCiMxQyfKlMu1YgDgxE3kmLH1lhtIZTBpMTaoOo6NDZlOfecFSxeMo/f/f5+xsemWLGgl9/d/U3sapm+73yLyuBpqlWN3WlYvYzs8rl4zVk8S0sbj/z6YcrD4yY7mzroDinNQqerYsu16xp9JdBCA5MnEysoRIrWjIetdN/fMnnzScXoCZuJco2BqQDXsogRpFBMRQqJpjq/Gkbc0OwxtylF69q5tC+ex7/8/gX+vvcUK1csontWN5uf38LQUIFMyq9j8pPAN8uxkUJRrQZUiyEWgtCIWTLpFK0tzfTO7qGjrYVUOq0Pd1HExESBU6eHOD0wzGShQCpls2B1K6rm0H9onO9/98usWL6UUqWkr40RwMTmkJksXs0tLXzh6/+PcrnMVVdcxsFjh7jmkqsIg4CWxiaaGrMEtUAH0EmFVBHZbAOlqRKPPbsFp7CfdY3HyWXzfOeefby4b5zr1jTRnY8R1SmmiiG1WFBTHhevbWP/sTLzuhv4+XMTfO4Nc2nKpvjjc8P84MFjbP7L+5i/qIeBk0Xe8OYfM1YL+NCybl44foatpQimyrx+eTcd2Qb+eniQQiXAsWzKsWIyCKlp8I0RD1GHEAoljKjErqsTlVBg3NtSaaK5Hs4rvVa6jpkL6xLGMcF32nQdYls60E8aU7U0VYlSms+VhMMlSkjOIupaQg/eXdchigItkDECJ01a1xggSwjcJEmU6apoOsPFJjbz7eQgH0dhfd0FcMRZu6dOjZNUShUj4ZQGs66lXMkuqqf1idxWEgZRXe8fxzru1XaMXrnuB6FOzbU8vbMlsbeJoUUIjXrGEthodZay9deEYYjA0DwdRwevENc9Cp5r5GkmITEMQ6SS+J6v5aQyRgqLQFhUJMSWDbZFRUqsOEJatva3mGFzHEkTwSqphRGRYyGVwAkUTZ7FttMTHJwsc0l3O+fPsSnGMFoJ2DM0Rakccs7Kmaxd2oaoVfnp71/Bs2MuW9HEO7/4BDevsHnTZV2cGi5TDRVBoCiUQro7G8D1GJmscmoo4OGt43zxzl6U55IWkudOVzk8XCOfdutoY8exkGiIoW87xMnQKxmUGcm0kApLQfF0kfd+9K00NuaZGBvD9X1ig23WZiEtY3SMcbOeTohCWI7OAgCUSMpqXXqvX7+OZ557CYSi/0iJyfGQtlkWLTMzeK5gaixmZKCmOUApAbYiCgVKOgSTAiyL4a1phk9WGB6dIOW6zJs3m5XLl9LV1U4u20AURYyNTTA4OESpVCGKYtJpn9aWJjraW5k9u4e5c2ezbPF8PvqJzzM0NEZ3Ryvf+8an8IMyR779DcpDZ6iM10jN7qbj6vPxO9uhFkAUgieoFUsEU0UzMxIQm1aW1BselsBJabqBjGJUqLBT+vfjSJgoZm2a7c64hComiqRRu5iWomnN1oSkOeNSjhWjpQhLCQpCy4+jSNHoWiyMLV6txMzLxYwfOImXyfKRy5ay48QoJ/pP8/nPfJoPvuvt/OXhR3l602aO95/Etl2UsLE8jatJuynmzOnFT3l0tLXQ0d7KrJ4uemf30tvdRT6XxU/7Jk1UPzthGFGuVBkbHedo30l279rPiy++Sl//CYaGR/nbE89xzjmrmCiEeK6n6RJRYHLSzWlVSsbGRvnn97wLC5vmpjyXX7CRIKzh2A61WkC1UsFxXEIZavGMZVGYmEAqwe3XX0wUX8aZzd+nOHyA61bn+fsrQ7S1ZPFTHlv6JY2qQD5lcepMlfGhIkJJyuUqlRAqpRo5z+Gy1e189579/Oqe7Xzpsz10d2W5YOM87n9gB3/sG+WdC7tZXCyR83vYN1rkr4eP41kWYSwoiJhAKWoKYgNVlDJGGO+XPgRLHM8jDiMTIWDyyRO0ksRE1EIcxvieRxiF9efKdjVvMLEtJHzAxCiYqKukyTJP+FVJa8l2rDpdN6kzLDN8VwZPQmzUXibIyjWoktB0jhLuVt2SIfSz49gOwtIwVktodaLraT9SHEtEZ1ensk1ClWXptCuppgGIrmsTGlKvNEFJoHHuyRA90SYnAfFxpBf6OscKUVce6FmHU9eOyUReKs7i7EvMBTaRukIQBjqbWZh5ijAxi4nkVNYjdfXvua5bZx0JtBZeyhjX9c1Or2+qkpIgCnSGs1JYQtWDc2whUBZ1nLyKYzyT/5ECUgIs16I7nSbtOOwvFKkoxYnRIj/6z9t4/etWc+LwKOff+hNuv7SH/qEK2eoon75jJqcLEicuc2pUEVYjxquSRYvaKdTg4NECxyci9p0O+c7bZ1KVLo7v8eVNIwwUQtIpx/T8fYNzdupGriSzOAxDHNfE+QrNxRo9VsCv+Pz0Z9/CcYSRS5tDgpRYppRVZxk7tdtf51JIZF1soW+15vJaQpDN5Hjkb0/zp3v/yrG+k0yVy0S1EEvpFCgvZZPK6hOqZTlYjtA8qkC3QGIpsE329IUXnselF21k0cK55HJZk12gpmddsXl4LVurT0wfWsqIbCbD3b+/n89/8b9obGzge1/5CNdctZHd//EflI4eo1qJaTl3Je2Xn6v7xEGAEDYim6F0cojTf32WYHjUbJDCpMPJOionoSuYEBssSxJFJjnTgigx1AoLGWnlWFLkWcIikjGeLZBK/z1SgS3g+EREJYx1wJelsJU2heUsmxcqIetyDsvaMjgdeXqWz+cnr5zkW3/bzpWXbuC7//2f2BaMT07xua99lR1HdpNtyaAsXRF52HzirR/myosvpVAuGFKrJiNHUUQcxJqL5th1rX+CK/c8F9/3NfeqWGLT5m3cc98jvOF1N3L++asolis4yfphDh+20IquWEV1BQ9AEAX6OpkTruu4+j0XFqCl4EFQw/c8grCGjCUN2Qbi8T6Gnv0Wjc2t/MefD1OuSL7/z+s4Nan49h+3Ux4ZYVajz7wWyLU24KUc/vPBCbpb0nzhDb20tjTzX/ceZvPeCbY89BEaOxp5/rljvPNDv6Ux4+ILwZycz7HxEoVahGPpA2YtjgikQjoO1SCokzakjA1ZI6p7pjA+i+mwDWUMhGYobdSatsmRj6JQQzCVmmYPmiG5lFp9qgxgMQxq2p9lUCeYcDihMNWwfo4sk5Vum88kjZXCdKEMIkrbM4JqaAb7er6dTEvBJBk6ttkw9NqppNSzb6GrmTjW44s6jdeybcIorHswtBnNJTAZ3IJE0jVdvrmei2Xr7HJp8CW+55mQ99iUXJrJ47muUSWc9Y+ibudP4F3JAxfFiQ/EDHRMqI1lpG9JyZb4EJJFP+G3SCnr0bax0UprP4SsO+t1ApjCNcYcjDosMhulMBBJ26R6SSmJlCASgOtQjbSDdaxS43QtQHgeY+WQxYva+ew/X0Q6m+aHP9vCC9uOsXheK7v2D/Cxa1s5PRZgxTVODoSUazB/VoqD/RXmzMnR3p7ioS3DPLarxEdv6KKzNUXKcXiyr8wzfUWyvl6E/bRHnJCHlTTyaBBGJui4rnFkK3AEtnQ4vv0k//TeN7Hu3DVUy6WzWjNm6Ic0CTDTyZCYayvMwM/0aOo30DZAzDAKWLFyMVddcREXXrCOtauXMW9uD10zW8k3NdCQ9nGEh2/7pPwU+XSOpnwTs2Z1s2TxAq66bCPvfMfrePc73sSVV1zCjBmthLUalXKZarVGtaL/rVQqBIEOTwplRFALKJerBEFAFMVMTk7xla9/h7GJAq+/8RLe/4E3c+j3f2T0hRcolyNaN66h66oNqGqNqBbh+B54LsPPbmfkoU0EU0Wog+moe5wS9I46SxYtFWf1lf+xjWCb06dUZwXhKAy12mz2RuWoFOR8m8lahG9BpKZJxwJoFBa7AslCV2DbesY3K+/x7Ikpduw5wrJli+jumkEmk+HIkeNs27udbD5FHEpUIAmjiMefewpLWCycM5/C5CRRLAmrOmVUI3A0PVv/vIavZNrI1UqFYqkMQrFo0VxuuuEqurra9BDdLKrT2Bp9So9laBa9hK8kcCy3zutCTQcgIYz030AAYxkZ1p5GxqdaeimNHCOeOM6crjw/evQk7RnFgt5mFvU28tDLY8xtjJjRaJNLCU6ORaR9H88WbDlQ4KpVzXTOaORPTxxh8YIZLF81i862DH/fdITB4SJCKE5OVVEowlgxpRShEFSVJFASaX4223ZMeqDQuBjb0Sd5Q+sVxuOWKE91aJvSWTyxIWgoaa6VVptqx7gWsFim+oDEfGhYdNKw7BLvHQniRHuu4tgQsI3KUre2bdNSM+THepSGbke5vqfvhaB+n21jk5AmyiMhDMgo0nSDKDY5QvWZBI6USRhKVPdhRMb1HYWhXmSkPmXqqFg9WLRUwq3CnMoElvGTyLMMgImiQOOb9UPjJTuzDZ55UGMT1i6VBiK6xhAY1AJc29W7pYEzWgY1b0ySYFkafOYInbsc6O/lOM60rd9cdNtWhrBrm+opJkgIw6Ytk/J9s+lYeJ5DGEYmjEkLXGMJk5UAx7IIIvBtm0hYBMqmHETcdt1SmpsaGBss8Zu/7Oaitd28tGeIN13QyFRF4qUFA8MRpchi3QKH4fEqubxLOu3guzYjZehtdVg9O0OsbCoxPHGiStpzkIJpDLvB4NtOIqvV5Zs2EmolnePYWJ7DyZ2DXHThBm675QaKU1NYtlPPN57GkRgpX31DUQbdruXRqs4800PiZEifPFSjo6M4js38Bb0sWjwXx7qCWEZIBeVSCRlFhOZQ4LraTGZZNp7r4vsutTAgqEWMjo7Ww20cx9Y5tsqcUm2rHqGbpB86nq176i0t3H33nzh6pI9ZM9t577tvp3D6JGPPPU2pEtEwv5cZF6wknioiEXi5NLVildMPPkn1yAmE5+GkHJTSkl3X1c+pbWv1kDRCBN0+rYNSdYqcUHqhqS+KGi3huJZOhbMEliNMa9Su6+uTFkRKCNozLgOFGr5rm+AhiAS0ORYNtYiDlYjlo2WKqRE6etq5pLeRnw6M8pf7HuT889bh2DEyVFRHQ+IOkwXjW8RVRaahge/84v9ozjZx+cUXUSqVjMNZv9s404ohIbR81DHR0JbjYAuLOAqZmiqeJXix8QzJW6FwLMPSU1r2r5Qy1ZvJ+sHQYWOJFNJEM5iWq+PW4Y+RjHGcFJaMtHxaxuTnX8bIwA7mdrSwblELv3hqgOXzm+lqaeT69V28/MoxlndbhKUy1apPxnV42/oMn7l/lEdfHeG1l8xm+bxmfvHHl7jlhhV4mRRXX76IF7Ydo6M1h7ItJkMtgY0ti0otMJ/LMfMdx6g/nel3T2lWHyZO2xIWMo5MZaLnuQ56jqGvhTbkaQmtY5hxgkgmSlWlfSSJC12aNrWv44WTKiOKIlzLMt0bo74UulJIOFnTqixVT3PVi7ljqhapN0Hj+3BMp0ECzln5M5alFZxI7ROsV9M2ulpOcjQSrLAyN1S/vM605t+guzFS3zCIiGJJEIQ4lo3ruriuq1EahqmVkFwjs6DrGYlBWpu4Rizqedd6zpGEVoFS+uV1LBehzkKHy+mdG7OgOUlKoTR4FUv8g7szmb9EMqozohJDjGcWYmFpN2ccaflcGIWmVSLqPV2t9jK57EBoWZSVRRBLSpUqnR3N3HLtCkQ6zVMv9DM8Mk5HW4ZGJ2JFp8v4eIXJiZCRks05K7OEwqUa2ihhk3IFpWpE35kKt5zXgoXARfDiYI2ByRDHon5idcS0J0YoQSyo9zA1Pka3WaRUBMUAaor3ve/N+oENA/0zSIWMI7BNMI7CaL6FkWXHZhPR9bUyQ8Nk6K5fCmXIxaqOq6lWAgqFEhOFAsWpMpVSGdd18VIp0qk0+WwDrqlc4yimVCoxPDxKuVhBRvo0ZRtjYpJvkMQNSKV/rY1aVj0DwfV8popFHnzkCSKluO7Sc1iydDanHnyA8ug4Tj5Dz1UbUEENKXXLtTo0wfFfPUD5yAlIedrMFRv6gqmCLcciUrHeZ80zjS10bLEw7VPMwSmUhqhAPQYhDOO6ACCKdAsgjGKELYji6cjhIJbkfRvPs4kiVZdQxlIxFYcs8QS+VIzUYkbOTFGZnOSCFot8JsXmF7dx6MgRPM9l4fw5BCVJZSLE811kBLZn4fgW6XSGPz/2ALWgVk/+VEpX6lomr1uDcaSHtHEUG3gi9e6AZZSQCL3JxTKu5wApJeuG4oQGnGyUGIqBiqN62FxyirWwiKLQzJlibNvTOTKafY6sVvDbF2Hlu6nUAu64uIfTEzH/8buD/Ndv95B2LcYjl5NDRbIph1oMXTlJpRTw5vOz/OLxkwwPF3jbdfN5ecdptu/sR6U8brhqKc1tDUwFMRUpKStFMYwJa4ncVbcxo1CLeqIo0hRxS7eOzJdo/4vBKukKwqrPuerdEUtn4khD4HDOwpw4ppKxhW71x2fNHi2zUQisui8pidewHZ3rI0yFq4yxOzbfN7E8JO8SZh5tGROisDUl3RL6oJSot6Tx4iWtThMySxhFxjysh++O42JJqNNzHbMoi7NgemFY08lXxkkemx6c47kmF1s/KEEQmgsc1IGQiUTNNq0sgcBzXIIonKY6ymmdsW7BGGWBUNqQ5dqA3rmxtEtW7+K63SLr/bnY/ByRqagkQSIhlorQ+Ft0cpeOk9ThKbpf6zhuPffdsoXZhPRmmGBSLNvSQMfEJCksHGu6H1ipVLnu0vnM7e1AFQL+9NftrFncynhBceGSDCnforlBcXokYNHsFHHkUClFjExUmSpLMrZky+5xbEtw0fImamFENVI8fmiKlCMQjoPreUgVI2y7nhJnOVotZBuGWd2Baj5zdSqgMZejpaWJSrVSzw5A6Bx0jMNc1cGQ1FVvtmljJfDL+mNonipL2IbJI1DKqoMyXdc111WHaoWRNP1gQbVWQymBY+v2m8YseKa6kXWTlUiO9HWlWGJQteobIOiDR0NDhpdf2cmhw320tzZy601XUOg/xfjL2whiSevqpWSa8kTVANuxCSYKHLv7EYLxIpaJbNXzMt1+qsuXjbFKmVOiZWmpuuto937yMmKJunbftk3bVaDNWhZYSRiaET5EYYTrmE0y1l4S24LmtAPCzECEJt7GShAhcFAUwpih0Son+ydY1ZphVWeWkYkiDz/8BJVqhYsuOo+vf+X/Q5UFYycndSUvJUEQk/J8jp88Rf+p03iub/xczv/P1V/Ha5adZfr4tdbae79y/JRbV7Wk3S3pSJOOEEII7s6gGb7oMMBgCTaDzSAzDMwwBAvBQpCQEIGmk463V7tWdbmcOn5e2XvJ74/nWfstfh8+0ElTcs579l7rkfu+btkBGdsWUa4s5RnT7O4UQivTNZkaYV0bF+BM0RJ6BYORsTaOwrh2RFIYmfmLUbgro3EreTPk8CIdtbiylO41enzyFN0p3MK1rK1vcf3Baa7eP8XJlZqrdxsGG0NuvnSGB5enObIUWFn3zBSJtVHgwLTh2t0l7/7XE9xz4w72LFa8530PYoZbXH7pNl558yFWVjapG4/REXbVLfUcoHV7pxw9q2eLVPdRO3jZQcQQNUqh0DG7aVEjRENRFSQjl6yg8XNUbtQ9sryDXtMIm6ZRvLxSPZKkagYtzIN2IBjTXvbOGEITWxBihrI2vsH70DIFrVJ+L76QjNOOn4vRVKn1tggVxDEe14zHdYtyag98r1kfXJSFUbhKk+akIi87VYv/NQg+wlinYe2pbfOctS2gUUiwThcyIsNs3ZQZ990ur+UD9nVoOTBRscbiCJwccikq3thaVbcUGDsZO+TFZ1G4Vt+ck/rQpZ+gUbzM83SPIGz9LASQi01yzaW6kH2OLK7rumkv1rIs+Iq3XkOqHC+e3uBzjx3j9uv2sL6yxvV7CpyFYW245Zo5ts1Ynnp6mfsPb2CKgmGTWN8Y8f5PX+DNN8yzMF0y1at44OyY0+ueTiE1TopJXkQTVZMt46qMSpAlrOQnG+ewhcGPA+ORx9e1Ol21RTZy+UnaotdMlthemrJoFfJnbOfU+lChKq+UlJOV2sO/VZdYowFHDot0jT5FXFG1iWrZWyKfuZGiQAGcEkSl+QUqlTTW4mOj41CRHpokMZ+f/PQDbG4NufWGV3DTLddw5jOfZbC0guv1mLvsEvAB50p8nTj6d/fRDIbYSsd9uVIWU7cYxJBDNO/8gmYrFIWlUcVao6FcMSWaFFp8tszAE9LYJc2tUd6cmr0kXzqBBa+L6+nSaWywpn3KzcY4JsYhMWoSg2i4sF7jYuTu/bNU3Q6f+tRnubC8QvCJt33RG/i93/wVds3uYXNtKB1iSlTTBcPxgJdePkan272IoCyocpKMcGMIgCOkILN55aFJNSppIEnZUTkAKXei8g4H/TWpvWRklxbbizokLx2wHJ/tRWP1+Wp9S0BVlpw4+jL/958fwrgKR+DL79rJ+Y3AlZfM84Zbd/ANb76Mt15Z8ukTlmdWHf0isH37NKtDw9vvmOeRF9c5dX6dr3jdAf7xo89x8ug5Uhn4srdeT6dT6sh2knVhNccmj76LslDhjmnjubNvgiTjn6qq1GAd2s9WVSl0OpX8+yQHu3NOO4s89hRBj0lQFiXOFkrZyF+DmRz6undMJEIT2+RAicgNuEKmOiEG6SaNuShGQyCyGDm7IeF9I+9m41u4YlIWYSumKir1+KSWbdgI51A1wqjlHUNZFa1xLyukCo2NLZyj1CV5jlsM2Z1uhd6b83it+kCstQqlk8o1S2yjjliyazMzWzAGW04MMk1Tt8tMe9FoKmvwjaoSXKnzcUWSu8K17vTcXkbNHJblYWI0GolkWX9ddmG6sqAqRS5snW2z2KN2YDk8K0vuBsOaqy/bxs1X78EA//jRp0jAjsUpOq5h/7zhwoWa/Tt6bGw2/NG/LvPRZ0acWglsWyi56do5fu+fLrA8CLz11gXqCMFYPn58SKcSJk5OQ3RVgUnafRnZF+V9eSL7QCZa7RiE0Cm8/yQY58JA1FFN9NiiFLSC0XhafXkNikhJysNJSZfvotpIGbZlFMKZoqjZ7GTskVMnk/5sYvRk0n8yqW2nc+dpyAdn0pGQykttqfQDWQAa62RmWxYsLS3x4AOP0ut3uOe1t+CS58LhJ/Eh0ZmeoprqUW/VlDN9zn/mMINTSxTdQubVxrQKGvlHvMiEJV+TVbpv02gutaJ4ilKSEydmW4TBVTilNk9eepAqL6YkEFAjQxL018cY6ZSWsrA0IdJRbYez0LOGcZTLq7CReqthdXnAHdt77Jrtc+Tlkzx2+Gn6U32Wl1a44oqD/Nd3/gRmJH9It1dgncGWhqWVpVa5aBX7U7qKZGJLgXVapSaksi1dRT7urS0IKaj/yLadfQy+LSStGuN88prpbQmhwWpHGpPHFWV7QCczSdKTnZqWeRb6vT733nsvH/nc45walvih59ZDM+ze3uOjj17A1zWr48j+2ZLX7Wi4dNrjjJw5GE/XJG491OX9nzjBF966i43Nmg/967MY57n7tj3s3j7NWIOanJPuMAFOvRo+yKjVFq71yuVFtpwNqmoLgVLZfQl5plFPSMzkXB1dRVILW9QEcPnsdRxV15IAKZYFKagwGTpr9FJQNWMIavKWCQ5qHLRODYRtvnqaTB/I75zug4v8/EZVa5k2CM6YxLget/iponA4IwpZmxVKUcOaRNon8MQ82sr2+qrq6O3s2u29pPMJkVNgar6d3eVcX980oNGqVVW29vim9hp6FNvLoLlID13X49bN6etaDxWrH1honZExyBirGde68NQQFdXGOR335MSvGPXS012K1e9fjDJRsclRc090N6RsLWONiAtMom7GUrEHGRG8/U3XMjVXsbE04H3/9AhXXzLNjvkuCz1YmK6o65oPPLHFT71/iXtfGBITXLrd0pupuOJAnwePjbj+QJ9dsw4X4anlyEsXarqlJekTVlUVPjQ4J5JoWeSGdlEme4LYdkgxJOa2z3Lq9Dn+9M/+FleUWv3m1UYUqa1oCHUdqiFRGsMZkjDQZCxB28W04TMqiRYDqaa3xYm4ohVZmMxMs5O9ho7GcsRAFgOIAkxGiHrGtu13CI28yHk0Uha88MIRTp8+w8LcLLfefDWrp0/TnD1H2evJM0PEdQwbJ86y9Mhz2E5F9DJ6kuQ5lwnW2pGFtvrKrUkMnlIzGPJyUg6AvCOTyy+QaEJsIaExTuJAZZc0UQyiRVJQ9zox0iuVKZUgIj6Ch0aBJ0cNZ4YNHStjj421MXvxXLvYYXPc8MjDj2HVx7G+tsmhA/u5ZN9+ttZH+LF2atby/AtHRWVZFHgv+x0fG73kBendNLVKpOUSiSkqftzi/RhnynaG71zZZvfI8lYC12IKSovW/IqyIEavX0eBD7XKViWhMGo0Q4ypNQOfX7rA0vI6J86epb84xWeXGsqqYna6w5fdtY/7nlxnNPL0OxbfLdgcBfZVgboJnFkasnvnHMsrQ15/zTSPPH+Bnmm4+RWLvOfvDjM6u8KuXV3e+NorGIw8zoAfN23Cn9cReJGZVlHIHCmvXo0Ut0lNfU3tJTrZy6VgjahGSaJ+yhG5WQKf9yMhSExE3j/q4aUdW46u0AkNUUUtqBUhT4uiejOCBPcZWcpbjSNMUVYNeS/jdd+dWV1Z5CS593qZF7a9TJwi46Pmk2TklOD8vRw41k1iEp0erEZHPCF4kgb0RFW/GCsH2mRPJh3JaDSEZMRgmI0xlZA7c2ciqYPgg6dWBHKbMYEcSDmLJB9kWfIXtFJ0+g3nSlcuoaQMmkQzbtr5NJrhXVaiqrCFVeSEopKdtKyyXIzt3Da7bFOiRSBnV2ahfpYYAzNTJW943ZWkqRkee/YMTzxzmi97wxUUHcN0FfjYY+v84sc2eeBE4AfeMMVdl/ZY6Dl2biuZm+/ywukh57YCb75hHlOWGGf46Iurrbs8KawtBpUHpkkuBzh9ECbO2KgXPwk6MxVTC1O88PzLgnm3qO5eD0jVnmfvTjvbVF2RHKhNGxwW04RLZvR0D8FPLiD12eQFtzjYVb4YL9Jw5yV8TC0AOKP584LM5BEXsdXIO1tI261Eg6oqeerp51ldWePQ/h0cPLSXzZMnGS+viT+piRDBdTssH34RPxjLzzDD7YwcYhnPkGf8YqycyHXF2xBa2gEpXbQ8nXiRMuY8Q/Wysz8vRAsjf04umtrUR+T350pPKftUwLFx5FwD55HbZhwNG6NIkSLXzxXYouCZZ55nOB6r6N9QVhU3XXsdTe2pOhJ525vu8eLxl1ldXZOMFw2BMzmSOCWCryWKQS9RY4RCkGNecwUrP9+ohuNClvHKUov6GeSfreDS9V1q399SGXkKRDVOFT4O4xL9Xp/3//0/853/8Yc5/NzTbFuc4bHzW5wdiRT1jTcsUPR7fPrIkA6e3XvmsFXBKCRGmyNWLgywKTAcBuZLw8xUl5eXPV/yqt08+MQZHn1mmdQv+eI3XEW3K0y0oiq1OEitETCTB2Tca1tHuFGlXVJ/kux8J14s3wRSLiJSpCgqUZDm0a81rcinbuq2y8gsQauThRRpz78YJ5dHPa5FLVcWOkKUs6nQczNfLkHPuE6n057bYpb2revcWEP0cbJIT7Et/LPsWAgKQjzPEwqbtEvAaqtuDWh2ec7vLVwhC0JF/iYmrVATlB+vphdnXetWd6Vrb+Log6p/oiKIo/4AIpUe6saZtpINwbcKAmdFDZbQWFytUGMy2nZqxrq2va6QsUHVreTXqUszo+ljUsWBIgbQea2rJnJfWlVs0i5G9dlqtLS2ENGBtaxvDLnuyu1cdfk8ZlzzDx9+nE5lecsr97GxvMF9Tw35Px/f4FvetJNf/Pod3Hb1HB08e7cZ5mY7XDizxX2H11joF9x0+RzJFrywHnl2OdDrOEwhi70micEnVx6uKAgpTpbnmuUMSszVzzMrMF77ulfS7XRawm6MQfZRRDWCSntqMJhkMUnVJTFhjQaLadueYqb+ahiZLXKEvULaFK9ictaBa7PZs3T04lFk+89W5TMxehqTBBefUotUEcd8oCxKUjIce/kEycIVl+6l3y/ZPHmKqjRUlSNsjmiGQwbrAzaeOUZ3uie7E2PV2yLjsJAyn013eTldUUeCttCxjsY/Z9VbzrQpnFMyskItVYhhkqRu2ojkVUeR8k6IAalNnPMRes60c2bEDM+ihbGFQYKtIJLh4TiwvuW5rOeYm+5y4tQZlpaW6fe6JKSDftM9d9Mrp+Rn7Az9qQ5La+d49KnDLC7M430jO0AdMcpI22lH4VrvlzOOgJo3i4IQG1XLCSdNvu9S91axpSNL4VYo7SIKOBADRpfLVjthonS6+tnFKDyo0bjh/PJ5GttQdCyDOvKJo0PCKLDQs7zy8hn+/P7z/O6HzvLhJ2rKfiUHoSlA0xtnZipGoxF75h0PP3OeV13WobKJf/jos5ja88pbLuHQ/kUGg3GLYRFMU9GCRfNB64qCEFKbn1RWlVz8bTRtbJEmrrATmT1Gw96kYHBl2aZ9WldQFZJ5HvWCdvprg5q50V1ZVRZKylVRj6YWlmWJdWIgFhCO1M1yTtg2ZiNoqmhmclmlnufnPitW80VWFIX+WYaqrNTgqKNlEjZGVVbF0GK9QyMPVFGWbbWVGVd1I1V9o7PjbreSy0XTAEXzLJViPR7rXsFRdiqKohKXtLo6rXUtzCvp7BUmIwRxaUZikhaPnMPeJuzp7a+HWMgVr9a4TSMQuKBmHK+y3TyflXdUW8gouujs+cp6b6xowY3iKbIKBZMYjxv1gkTedPcV9KYsF06f4wP/8hRvf8OVbO8F/vqjL2IL+J3vP8Rb7tjN6lZgsw7snHZMzXQx1rC+FfnXh1d5220L4gLFcv/LWzJvz9G9IVBmI5Mu1zJ8sr3ErdNUOPfvLnvxtkUWt81rVkHSg9lKdGcyRO9VTqlZBcRW349JOiZLk4fcmcliOMkYLbXqPa+VfGh3HDHmJX3QP0+q1KQcqZAm3KCkZOfMSyNOxocTdYtvZb2195xbOoezlksPbIdYM15apvaSa2vHNRfue4TzH/0csfHgLFErRREPOHzT4Fwe3Si4To1xXnMacrCOMYYmeu3CpJvzF2dbx6gFi3yfAfkcvBYuGX9hWp6bjH+zeKVbyYUTEDqzN7DgLFtNZBQMK+oH2KojTYjs71oWux3OX1jj1OkzCuYD33h27dzBbG+m3VGGOtLrdfjdP/0j7v3kJ5mdmaXT7bQGNx9qLZ6cPgvaARv0M0fzPxw+NrozkkJS3n+VgSsbKhIJQTxTIeZiU6W7MX/WvnWrR71EBA9ScOmlB+jP95naOQXOUCTDp08NWA2J/nSP1163CBjecP0s+3cvcDwsYJxhbRiIFIwHNUTPxvqYTmlZGUT2zpe84dY9fPjjL7J1YcDCninuuesyhqOxpHoC9VjQKoUezJlw0TS+JeSmlFrxSWJivvRex/767jjrWvRNNpAmfd6zpyaP4jMSPosIjCqs8iqhaVRcoB1+uMjSkHQ6kk2AST0i1riWlm6MnJ+y8JcCKXk1Etps2xDmW05BREUYo/GIpN9f0JGZLUuF8EXTKg+qqpTlkfetvBcj3UWjm/qiKDSHV8fWDl1s6mhDHZlJcSBZzYReRhnS55tGDzs5uH3T6GVkWj5LjJGicm3a4Lhp8E3QG1xmcWVVtv4BXzc6DpAKoNCoXFuIYmyybwlt8FUbnqI3sjVW4nljltNpkIua7TKHv24ats31eeNdl0JR8MDh07x4co1vfssB/vt7n+DlU6v8/o9cQ7fb4eGHj3P5FfOcOrbJtZfPslVLRXF+KzL2cPfVs1SdDueGkQeOD+gUlrJTtS0xKdFVGa8rnXaNgpLObbFTR79xhjwQtKWlM9fl0UeekLYZ1yIl2mW508AoNVlm2WZLQdavIZKNmWYy5lJyL+rfyY51q5nPefmnfeoEXWNNq3JyGchnC8mNyQo+3XuZ3I3oiCjn0xdFyebGBqdOLTHV63HJnp2YlAjrG4TG09m1yL5v+kLqrRHDE+dwPZEtF4U4762TA0yEIRJUFr3HFFZfeiWhKqdJZsSoGVD+HVHGqVEZYhEJGGNiUxJhglUcShTMezu5y2oeUsvM8jqfjgY8sLtjKQvZtZwLgXETGIdIUwcWXMGOqYph3XD8+EnKqmgPqKlen9mpGcZjT2xELl9NVWzVW/z87/w67/y1X2VrayTm1NhgrU4DFMVhkZ9NiIGy7GoBI++HszI6KWwp3DgCzkrRmXIcRJIuPaRA6co2ryamKPJfLUImsdjgfU2/16UoK9ZWNxWbpHunGFkeRx45N2a4PuamS6fYvtjj6ZObfNHt2/i6L76GQdFnFAyn1mourI6oRw3jWuSvppC8oLfcusjzR1d4+MnzJGt5y92XSSGMdINFKVy5uq4JwQt1Q8fuMMGsF65QM7NRc7TVzI0gmRzWtWdaiB4fgiobk6rNYrt/sEXRYpMyMkcOaduOycqq1BGU04tNzmOvhXomKORzM7Tjfdd+xj4InaDJ7K6ULlJPRpqmVjWtU7SVqOdyTnqnFF+esQ7bqpZ0pCB45iAvgH6Q2VSSkrjKFfLbcqiMlA46P81W+pyXHHWOJ0+H14PXOfkwyqrCFGbCotfRS9R5/sT6b/RClPCgwon7V9AkrqX6xjQJip9gOBSMp5jyEOPE9KhVoLNOmf4eQezrLZzVTHlfQqaNyi0+HHmuv3o311y9E0Y17/vg07xi3xSDzU3++IMv8avfdSmjUeDwY2e49vYdLJ1dZ+AbXnHlDFPzUyRjePJEze2Xz7Bve4VJiUfODlmrE4UeGC3B1VoaH1t580R5OmmzjUktgiB3Z74JzG2f5vCTT3NheVmcpib+O94/WYapqi20i5zMsEWmm5k/eeaadBAqvgH9fHUk1P65xhC16sLEVgKZNGNAionQjhuzqzuH2+QKLT+DMnmPJBOwNjEcDWgaMcdt27FADA0MBhhg8Q2vYu6Oa5l/1XUkpR0Y02ojJ+OrkEPO5BCQCFKwhe4CnJmYJp2MXp0rqEMjKBBxeuhSOYMoJdHQqqKlvSDQ5bizKgHO0QRJn9FEpysQ05BgTGRWuyZrHRvGyh7QyP9O2cTuSmKiTx0/3VIavA90uj1279hF8pGyqz6LEOhPdVnYNsNnH/scP/mun+P80rJITcXwIpJpJ2uLvK/xvsbZUg2CIqNXkfJFtOQg/z5BE8ZKhHatKz3/MyvuklOVUluISPTqCy++xA/+yE/wLx+7l/kdszSNdDZlp6S0hgfOjggm0e8W3HP9PB99eInjZ1dZXOxz6y0HWRsFNscNy0PDsJY19IlzW8x0LSvrI248NM3B3T3++f6XMCHyylv2ceWlOxgOajU66o5CabmtUc/KTiinpCbdBdEuuO1FpKakTn5FtVtDp9MBDXuqqlJ2FmpxIAY1l8poULLSS1KQUb4YGgNBcUwpKg5e1aAuUwGMdjdOiuOgGB5BNElQVdWpVI1qZC+jUnNXCJ3cWNt2UkZ/DXq+N8oAc85ivY5v8sJFNMoFGXHimxxoogeMVoFe9wqNrxV9UlCWpczkVM2UNcUx5G/Utki8kFIrhSUh87Uy2/udwsKkMutUVStHzaOMvDQvSuH35IWfVWJmUZbt4Z+hitZOZpA+eFFCeGkbs3KiaJMS5cIJF7WrdV3L5afz8kSi8YHXv3I/3U5i+fwa9332CG9/9W7e+9GjfNldi+zf1uXTnzvDtTftZHWpZunEGhsbYzpTfS7fO8NDL27xzJkhb7pxnqqqGPnAAycGFCqxtM4QU2gT3FxhdCEp6AhRQmWdtyhZmCCa9GA39Oe7rCyvcOLkGUmXrGVDG+OkU8uJhHmXZMyExGvVBNdCL1uxg20lvllKZbWb4GKmk3YNpkW2qdxYD0GXESWFbS/q/LM0RsJrsKY1p7aKLmvZ3BwwHA7p93t0OxV+3LB1fpXU61L0O6T1LYp+FzTz3FqNYZ0EDZBUQJDHC9iMZ9cdSBsP2pIRCY1vkRQhqTrRGY0Xdu3yNavITEotLTkl0fCHGImOtruJGg9s9Ofh1BNiQmLBORpgrYn4KHHDzVjMtts7Ms++sLSknZ3uGUncdtONrJ7ZwliRa9rK4lMk1J4dO7bzwvGX+bXf+Z9iCiUpjbnU5DpRWWV7cSLq3kzkuYWVXRQkStcRYzE6DjEu69fEL6a7yzyGzV6Z0laiQMNQNzULC/P8632f5N5PfRq3WNBd6JCCXMIi3LE8e3bIsc0GHxJ3XNZndeR45uU13HCFO6+cwXYqOp2Kc2s165uBuoksDyKrm0KfXZjr8wU3bOdfPvUCo9UtZndN89o7DjIcefq9rhqCFB5o8uGty3/nNDJWqvWy6shEJuW419jCZfP5mRVuvqnVx2Zo6qaNks6jL8HnCLYkd745JC8v9H3wjEeN7mUnuR++8W03axJYxQ2J6dlRVR187XW9UF/kY4nqrJ/AaJNGeMhuMlOIVcgSYqt2tc6KPS6PFJqmUcmZKAxktibALmtoZWhgCI3850K/4NzeipxRxillqcwllTrmnUPMbHqTu5kJEdTXjSgAFPku0jSZkzpXiLpKkRHeN3jftBnqMUogVJYCZrRHSobQNFotK67FFhNOvtRRNCHgo8D68t+ZL5+qKgleZuOF7iVmp7u8/q5DpLLk4587ztLyFjvnHC+e2ODtt89y3+fPcddr9lNVBaunVvEhEYsuK6s1V+0vOLYB68PEtfs7+CZyYph44dyQnlaLVmOBm3o8UQYpJyomlY46qRKxtItvwUjLMjyFSOOjRlJOfja5c5GlunrML1JDtR0NtsUsZle4aVVvogrLI8qkskWjCqXW16HLkmzOkj1NaC+rvLzLl/dE9WHaNttmw6H+TPMlORwOGQxH9LoFvame7M2ahjiuRdNf9Vh//AhJ57tRTYDpIrSnSRMfVE7Wk79Xd3ORVrKtX8BkH4QidNIECpiCxAokNX8aNU46Y2m0Wy+cxZFYLKCXP/gsYS2kM9EVC42BWyrLQQeNBVNYTAyEZIlNZL4QL8W58xeoa/n+yrJkMBzy5je8nqsPXc3po+fp97vEWkYnyRm8C+w5sIsnX3iGP33vX7F/z146VU8nAvJz8jrnzzk+QmUu1WgqqrhIwodxq0bybQRyav0juVo01soS3smesdF3vVOVzM/O8vjTz/HcCy+y69LtlPMFyYuC0jlZ/pqUGIXIZ45v4gdD9m3rcN2BHv/3Y+f4jb94jo8/do4rDi2wOkysbTV0KkOyjlEd+McHLvDkiRE0iS+4cTenTi7z2NNnSQTuec1lQGI4Gml17/SAl3ew0NjXxjctUDRfJkK+lWC3LAgpyqIdJ9p2ikHbRZRlSdRCG/17xqOxiGByXABJxvPOqfpJZMVGdcRSpNgWPxPzM2jz+S0Xf95nRZ3q5FWCXDqGqizU0zNRjeVRcx6lFWXRYqd804gZPHs/dLqg7ZRAvZxWl1G13qUu1TNGwpWuZaXUTa3YdTnMC+Ui5YsoBN+aY/IN7PQLEpqvqla0cM1ZvYnYLh6jmWT5el30i7HHUBaFfGgahGKMyFnLqmzn/SA7nlL16jkQJueSCIbaYBRP4PXWDSGohDe2XQxA3USuvmIHN1yzDxMdH7jvKNce7HPszBaHdvUZjRIz22aY7ZW8+NBRilgTk2O1LmXgkRJLG54bD86wba4HxvDZE1tEJ4Y/kw8glfSJWkpMdJlLVbiiHVnEID/DaLMUNrVeCxClxvTUtKaaTaS2GaiYYmxlmPJ7ozYScVKZG5nLR9WPX5xok+LFoGWV3l50UeSLSkKHJPcl+5DE5JpUojxBS6PL+DwWkp1Fasd6DqjrscoOAwTwI+GGlVhWP3mYU3/zL4yPniAWWk2FqL6SHOGrKYJqONUzW6o47UwyJdUgl4vskgxFu7txrblRPhIpuFxmyDlHNIbae4pCMyKABzYDf3+u5ug4UlkoSolRtrr0JEFVSE5531r2OEcRIemSdFAHGp+YsbJA3Vjfkn2f5kA0TUO32+GXfu4nue7g9Zx84RzOGjplId6LOhBoWNg9x4c/fh+/+tu/y7MvPE+lz9mEvBz1YCsVqkmbJBqjdsNqKET9IsJXy8o5dZ/rRWS0Oy3LDrYoKMuKjY0h7/ql3+B7vvtHeP7YS0zPT9GMvb5rXqO1C5omUSXLE0uRdW/o9ypec90851aGvOrKPi+fafjos56j65HKGaY7jnGITPc6fPXN0/zmP51k4D1XXTLHfNfw0U+8gDFwx427ufyS7TRhMv+vykrkxYZ2Z5r9T0b9S2KuK1omVdCIBKOvR+Nr/b22LVBCEDWnjFRt65PKClUxLsqIs3BOo8U1PA30HY66UNdAKC3em6YWP54VgVTUgj00OYVQSBrOOg3rk/2dbxpKtVugmU+TQEABgfa6XdnzqTTb5oVNIpFUl1xWJeNxnSPtZHTgJAPEZbKtZmOnkFqeUWbphBgZjsbEDAZTDHZR2PaDzg5YuS1de1gKlkNnqVlLrq1/vsldIbM7owdOUTiNzrUtej3/vV65Ra5wRBNbxpdTPLLRH6IrSzGK6c4n460LBS2ii31UGleWBYPBiDtv3MfUXI+Vsxt8+qGT3HJ5j5fOjLhiV8HRpYbLL53n+EsnwY8YDSMbA8OpVUiDEc++sMIzx4fcc/0sdRNZqwNPLDVUTtLuikJkpslI5YrN6haFqlgJh3HGYZOm55USW2us7JJyYRsaT7/bY2amryMvWXA73MTwp1HBgvd27diidQqrfdykzLmTAxl1DueXK3O2xMwUW6KptUUrXc0HU5HNZ7qjyhjR/GegMtqgLuekPy+YkBJSiAq5s4LO9h4aT6/fxR59Gf/sS2ASHQV42kLDemzRSrmTQblQtpVJZ95RUVgdkeqXZA1RhRqpfRdCi+GwuuQMWhHOdh3TRBY6BVOajb4cIv98Yci5YSR6w+PrDVsJCqUbV4XVLBwYBhlF1EHGHTPOsNVEBk1kqIfNrIFOVTCqR4xHtagPvaDHx+MxO3cu8lu/8i6+6Uu/huUz6zLTtgbbsYL7Lwy2D3/4nvfwx3/2XoyRRX+KXoqUhP781K/h5TkUGoDPSyWiieotMS0NwgdRsckuJGgLKtkoGX20e+cOlldX+ddPfIq5AwtsP7RtomRKhsqVuMLpwWjpdB3Hl8c8e66mGdTcedUs3X6HlUHgO96yl2++axrbKXnwbKKL5+x6YLYKfMWNfeaqxF9+6ix75ipuu3oX9332JfzamJ07e7zylksYjJpWxts0jV6EYrDMZ4PVi7As9dd434qWBMTZqN2hwFlJc0WR78ZNTNZOvWVGrQckDYDLWHelo8sZlLFMrkU4OV0vWCVCW2cvigHPo2LbOuCryileXy6AZlyrWVBd9eoHAttOOqwxmjYrmBfpnoTQbKPa5KPSZr0PilafxCOiQLV0UYkZNaDJZiSzLppCI8qUspRlTDYnZUZQdl0qqUJu3ay2SdrV6KxRMBmG9v8G+U1RccoZ3VGPG4za8XPIVVVVGAyN97oMUsXMhF3eOtmD9xCk+m5qrwdU5ueHNhUsaCZxNhqWruDVdx6CXskTL5zj/PIGV186x4WNwK45R3IFw7PnWTmzzux0n8FILuQDuzqcPz/kqaMDZvsF1+zvYq3hpY3I2XVPVUpV6xWbQrL45PWllEPMx0TQh7oJXrXZmmxGHhNq0FIdiE1kPGpaKbNcyoLVTrll1cWo7LiadsGedOkZSa0vRjpTOQAE/XzRf1bzYEweY5NSQeXyy65t1KUb9PdFL99fyuC4LBxQUqvTVLRMTZh0j8LPKrSgiDFA8lQFrAwDp0cBW1qmK4uLFls4xmMZQTRNI9LnJqgqKuiSl9YRbdBs9LzPSRDHnsWrD9KZn2oho9nLYPXnNvKeCvGLPD5K/PkFz++eGvGRzcSLw8iuTsVbt09x83zFtumCiOHZTS+OdJ80pEkzRIDKZIilBH8N1bMxbgIRw0xhcAYGgzF1U0v/FsX5DInBcMTm1oDv+c5v4rpD17JyepV+v5L+1EoEb2e6ZO/BPVxYWWMwGEjxkkxrEvWaox1jAlO0VXBhS3ysdTeVMSexTSbMxkgs+mdK0dntdZlfmGM0HPLwY4d5399+iLntM0zv7ml8glFoo5HEwhgplaMnCjXDk8sN0dfsniu5dm+Hew8vs7w+xMbI979uhj3bKh46E3j+fMM9r+iystnwNTf3+eiDpzmztM7dt+7j2RfO8PyRNegV3H3ngfY9y2dcBsFK+qomCSo2PSg8sQlyyaQoUuZ63LRx4CFEfH3RuPaiuAzxFEmhKKMv2aPEKHLhHJ4Wg28VrknzmlDibuup8xOSiKwRrDreta+2WZhCazLOo+eWnO697Gku8qRkYCnE9mfuvRjLbdkp2irV6EzPqA0+z3xl9hmV0igveMyGmcK06YCFK3HKjk8xexhol51Jdw9lqcZERUE3dSNCRyOyyZhiGzAv0tLMp3LU46ZdE1/s2SiLciIHzAHzGmKFgbqWxZTRDgVaBpo6q2WElYiTA04D6ZtGmP+dqpJK2ogaau/uWW69cR8MBtz3qSPM9wsW+pYOCesbwmBAMxgxVVoeOek5fLZhesrwmpvmOLVmePzlITcc6DM/26NuEo+dGcuh5cSAllVxkvnjwMY2N6FQZL0PjVTxJqkxsJgYH9O/ZwuNx2PG41qlsyrN1e/RWjXGqYzTFWVb7SdMC4GTLkx+v9Ola37489y45QTZos1Yyd2hQBaLi5oNxVTnnUwOZTKTsKFCZZ/WFYIkyew0LW6ctTR1bE2p4401nnphmaEr2dXpUI89Z5rEIEZ6MleiadV6QUaeivlAnzeNW29BiSklYh3wtaeY6rL71derGVGc4c4aCqeHcTNieqbH+Kor+e0Lgd84NuBDyzUPbCbed97zW+cafufcgJMxcsN0yWumDW/e2aVblWwGVbdY7caBQllKpVUuVgKHxcaA9wlfB3pWgH21bxgOR2L4LDThkwhJ3jUfAt/89V/DcKlmsDYiBbmwUsrMsqAR1dk86SZgc50SyF5QM06sFBZiFpS9XHHR8lxk77Z1pRudsU/NTLO+scXP/+Kv8Z3v+GHe8UM/zkc+eS9z+2dplO2UDaXGJiqt/sUcC74OlIXhmQsNG1Gw5LddNs2xNUTGbxKNtXz11ZbTGxHjCq7f12erDhyc8eydSnzwwTPceqiPrWs+/unnwBhec/sBdm2bEbMgRrrSDIMmKStLLuXgfWvMK6xlrPGw2RFujQAKrTEiLsqKw3RR9oqVsX9oZO+bi6KyFGMhFt3/FBSubGXxRSW5HkkNhTnNsCylaLeFpa6bFkMSFdcelAgqO9So6qq8njKtcTLz67KZNDMOU0JZiTLytr72Om+3CjyTsVNZFRpOMnHm2pZMKt+Ac04eQN1TCG49YJxtD5KgHovkdQSks/IMUnRWqkJrZUE/4eCbSQRuZj4FUZzILelbHlZGsExGIwoq0xs+xUhVFBBo/SchiJmp0g88BA82aUBR1brMnbN0O1VbOZeKHR8Ma26+YQ97dnQJw8jHHzzFDZdNszBXMdeXw2fXjGR1/NkjDX/y4JB985Ze37J/d0XVNRw+PuauK3qMNxtqHM+cH1GVpsWjY0S9IotzMdtkDlRE/ATZ9EYSqWzQTHmvcli8KIqGGyNm+lNMz0xrJKfuG2J21QaywDTp4jh3pvYiyFtui7N50OQEy5jzOkJ7ObSmwJTaCFWZ10b9/brv0FGRac2k6aIcAquJmGr6wypChdb7Yq2l26uoG8HixJC45fo9XHvpLv5pM/GL5+GXTgV++UzN+85sYcsSEzwhNcr2StiYMDGq3LFp9zhNI54OkYdb/Lhh561X0owa6pUtqqqQUU+QeXS9ucGe66/Gv/kL+eWHj/PEuU1mioK50tErYKoAEyMvDhK/enyLj67W6u6Hm6YcC/0Cn8BhqJzsW3K36QyMYmJ/YViwVg9TIZN1SJQYGi8UB1cIncFl1YwtKG3B2uo6t99xC6++61WsXdig06lo6tCaxkxlmZ7tUZbKUvMRVGgimTxeycSaAZKJwXEyYvYqwhFD6KRjd2qU7FQVTz75LN/9jv/EP/7zx0gzhj3X7+TQTfsoOlpkdArhNekOxQf5+wipHXeWNnJu4Dm2DmGcuObANHU0/P2DqwwHDbOVRC7UyXLtNumkd+6cZmNkePONixw+ssJiP3H9Zdv4+GePwMhzySVzXHvlLtY3B61XSDI/9PAsS3WYS+Wew+ryHijoJEP8Fl6LAelyc+JYhpMKJipMMsmtbQnhufvOOUiy9zStx0TiwtXM2PjWGd926T61hGGBQ4rfxJVFy6PLohXpSjJ63raEgMneUqc1yrULXr6XEKQbJgY5ZAQP7nBWlE42L7y1pbeaRy47hnw4yMVQVSVBXZIhyAgFzUBA0Qx6fSjUK6kmepI/UrQGNbksspJKiJhBwWaFHPpWVA458zdXV7kzsYpOcIUk33kl7Rp1O1snjIgcgNOCzpI86EZld61fQP9gr279pvbcdt1OisJw5PgyTzx/nldeM01TR7plotstWVod89MfXOeJc4lD2ysWZgpGDSwvj9l3ySzTnYKr9nSoOpajW55zG56OyqBFQigvTlV1FLSW2oerjVBVDI2gQTKJF632xTTlSsf60hbXX3sV2xYWaMbjFrrY5kAYO+GYqzQwK6M090sWpzapTJTJzuOiDBWrUm2jcluTywYzCX81KR/PqYXHZeZYWwzYSTaNFCM5eVEXlM60S1zpDA2dbkWyJbP9Gaq5aX7qsTP83fkxS6PIuI6sjCLvOz/mDweOnXfcRA+DtQnjPQtX7Wf7zVcSRo0iLORyKgsn2nxn8SFhOyXTr7iE9ZdOiYmryao2S726zsHX3cn6a1/Pz7z7Q5w7tcx0p2RUNzQx4lWR5DWutMTx7vMjPr3aMGcsGzFSe6+XsuJN5PWhW1gtogTrAuJCNlH+nYuJ0smIYjQcSbaFE1px7v7FpyWH2o3XX0s98sQQ6fQLXGHxTaTbL9mqB0KptpoPERRwaS2WAh+9uPY1cKzAibs/iVrLKMkXY3BF5l1ZHaOIGfY3f/N/szpc4ZrXX8nMrmlCCnjfEOrQYslFyCMFY6GyX6xw4SSVFMYh8dSqp+gX7NzW5eCOkr/5/AX+/OGafzpcsxYci1OWa3YkVla3cIVldRA5NB9ZXxtxdmmLm6+Y5XOPnOD0yxcopkpeddt+pT5r+mSyLSXcawZHUVYq8/eqHrUTCq/+3KwEwYgsWmnYWZVXONvy9wrF8RiFxkq4nm/9SiIrTxKype+Bs7qbKAts4TSWG+00gkYrJE3RtBSFktKzQTVK8dcpO1Rl0cJDc5QAbSCab2XzZVnIheiEBGKtxeaXMappJMQgjlpVadR1La2KyfnORl9ypTm27mcZQ4gqaqK+MEp+vXgxGrXCFB20fKApiLLI+6Y9THLEbkpC5sySXp/dl8rFamrf3uht1lGC0pbtrL+sytaX4hSPEY1kYmRmT1EI38qHPO6yE+NcjC0wMnhPt1tyy3X7SA4ee+YMm1sjXn3jIlu1YbbneODFEe/64CqH5g0/9gXTEA1NsvT6JVPTBU+8tMaeOcdcryDUkcdOjQlmUi2il1pRljShFu19mkhc83Wc5dBRR0zRS6Uoc1HZSflRII3gjW++WxaaGlxjFX6YycnJZK2abTEpF5EyW/l2VpeAcqSMvaigSO3XlZlieeeUK6LsBxCJbzYnhjbxr52Z8/9nWLzIpyEqQ3mhi6Kg6nUIdcAnQ1zb5Dc/f4wXN8ZMRzFkpiRSyYVOyUNHzvIbz5xj8Q2vpl+UJBPo9DrsvO0qeovzjAfjluBbe+VueVGxTO/bQWemz8rTL2NKixVbOs3mJpfe82qWbn8VP/7Lf8jayhb9qS5N8JTOUqonYK0JDIOIMlzyVCHx22e2+OzamPnCEoIa7Cw0AV3Ea16IRo/GJGPEwjoCiWjlUjYRfBMvYr6lFuMjVYUAFn3dsH/PHjkUqhJXlNSNqIJSYzl95iz/+vH7mZudbV3SOZ/cYrGm0C5QDhOfQot5j+o3lZ2awfumjVudnZlibm6Gzzz0COfOXWD3FTtVjeRbgJ8rLbEJ6kORTtmVBT5KiFo25MUQMUVB5SzPX6jZ3PJ0beLaAzMs9C0/+aU7uf7qnTx4znLjDhgNPE0M0Klkxu9rOp2SpfUhN1zSY3W95vHnV8BZXveqS+l1JiPc/L60hbPuRqKSi51mkWMkffXiSj4F6VpCmKCUlDusplPbBlSlKCZOa63EAcRJoZ7PnqSj47quW6Ni0NGXLOhNSwgIQdh2hXUiwtfztN3vRi+gXN3z0KohY5vZnhVfQXH1uXitNcrX5sVga/IyUrWkKIeutaZV8jn1DKQ8XrC0LsZG5385tD0rDZImrYUYWod06QoMwtPKf242qGSlQtIuwGl2eYzC5XdFoeoomU2O61pSC7PbPebqXKNqdFR1cXUbvTyUGdzWaqqbMIm+bRrxC2TzYcZ5GKibwCV757nqsu2YkLjv08fYu1iyf9c0pYV/fmyTDx1e5z+8doave+U8U53E0sqAlVFiut+hdI5PPrrCTZf2qUrHCMtzF8bYGCmqsqV+og7twhX4FLCuxCspVoKvgqIv9NDRkWC6CEFeVAUXTq+zf+cubrzhGoaDUSuZTTm3PlckuiiPOrsOMaehTSqWTDqWFyq0eytni1YOjMpwkyrXWny1dnQhKp03Bw5pi6JT58llzUW5EDnONntGLtqXlFWJwzFuGqYryyPHz/Cp0+tUtuBC3bA5qqXbjIlgErP9Do899iK/eP8zdF93J3Mz05x78RTN1pD5q/ZhnKUJSQ7KENmsG5JzpCYxe/1lrD57jMH5FXAlCcNgdYt9r7yVk9dcz4//wv9jWEc6XaEKdAspfDaSQBLfsq3LD+2f4scPTPHjl0zzY/t6fPv2Ps+MYdiIsS74SAF0CtNGmZY5R13WABTWtEKEQjPpc8ZDWRjKwhGaRmTGiqrRCpCNjU1uvfUGdm/bzZHnT7CxuklhS3nWbWJ2+yx//v6/5dkXj9Dp9vR9kko8Jt9GCFg1CRsxImGN7pJUzRai2AHmF+bo9jp89N6P847v/wn+8396F8W8pexY/KjWsWXUjBva3CCbg5u0yMm8pxRyEmSkU1mOrTec2WxIPnL7FbNsNYZza4E33bqDL7p2ioUqsjU2FKYg1SMGw0bOMO8ZjQNXHZphtm958PAZoOCaQ9s4eGAbw8FYDk/lw5mW5i3VeVYQ5vC9SKLb77TdnozqpZPtVKXkaCguKisjQwgS3KV2isxJ80F9WzmsTfdhKDm8KEqICjotClEh6o4ixGwCNO1/D/4iQK1z/84oLCpYe1HqotHYB1rjcArx38Ub5C7dhjY3NztFdfyhGv3MlM+Gu6CJap2qkmyNEElBAtdDyBpl+YK9kmuDD1Ixq5co70KMYgAikx1L4QpNRorKYInZfkRROOqxogZ0iV9V1UWgsgk7K2E0NU4YT0ExFiGIKsfXAd9MAnBSihNUvTPtbd00dRu1W3svnLCYuO6qXezZOc3W6hYPHD7JLa+YY7Zn+ODnlvAJfv4rd3DJYklIkY0tz/55x8NHR9iUOHJiyMnlwPWHprBVyalh5PS6p98vCU1DVQpPyV0ESZN0RVme51z2PF5sdz9OujtRX8sP25UFq6fWeNWdt7MwPy9L0Byopa7a1L70qXXeyoNWqpRaq5fIRMl20eGenw1ZytsWGZP/aZjMsGwm4IpQsM0GyQvxHPwlzCnbpt7ljIK2I9JnNaXAVL/H1FRXzGyN58TqkJ6BAy5w/ZRj3ibWg6ew4s4NSbDmTz51gnf+6xNUd91Ovyq58NSLTO9epOx3CU2gDoFoHZ8eiVTYVQ47GnPmc08QjcM3kdFwxM4rD3Hummv58f/2p2wOGzqdQoGBhlGMFM7wVdsq3nVomm/Y3uWaacdiYekbw/YOvGlnh6+cL/FAVUmgFCmxVQclqkKjhs2qNEonlqulXzmakBgHg4+BorCECMPRiG3bFrT7kiyTzD1LMdLv9fjxH/1B3nD7F3DZtldw7qUL+JEo7HpTHTZGW/z9Bz9Ev99Xs6jSBowhEcR1riIZgxBlQ2gmYMwEs9OyDHzvX/093/sff4J3/tJv8PTJF9h/yy7m902Dlf2C0FPc5HkxExKxyTyuvIcqhE7dYnhiZFAnXliRAujy3R1megVPvLjM5mAEzrIxaJiZcpjkBVDqOth+j7E3FIWhVySu2T/Fpx48RloZsn3vLNdeuZNx7el2O+3X1JpCtcPOZIZWVquet0molsij63Gtkxx9TYtCili1NeTitiwdvqmxOv4yulIoi0p3gQrsjFGtFZP4i6QFRVmWLd2jzWPXkX7wsQWOOetUuCKdvXXy7tVNre+xTBYKQSGoYELzTfIITsQDFh8iRSkztTpFjbWMuDInslldhtFmdGRTz7huME5TyprYArvKotA8c4f3Y3zI2IeouA1RSWX9dObjhOzsVYd6CMKDid7jVFFgnKEZ15Bnu424NaMuhDOdt1NWQmvV0J98EGeJXkYDyO4lCpHWqkyR0GY6FFXJeDimKCzRe3wTuf2GPdiO5eVn13juyDLf+sbLuP+xC/zlfaf5v9+9T/Y0rqQfhhxb87z55jkeOBb5zDMbzExVECM7pw1Nkzi2ERhHg2siriezZWcKJeCK4oV0kW86Iklkug/IEb5tymjSJWNVsHp6jb7t86Vf/hY2N7cUeCjta+lKvB/Ly5+CXunSaaCpcVYz0UV95nQx6NoRU/uytw522d0E38iLkoOZ1FyXomnHY4nYjh6jLqFT7lRQ8UaafO9WXwpUFWiMyLSrqhLZdowsb455/XWXs3P/DEW/S+EsG2PPZ4eBf1gZUViDi/KQTk/3eOmls/zSJxy/+LprWH/8Wbb6PXqLszQbQyBxukmsNpEZW3CmCTzyr4+yt3K4bsVoOGTbnm2sXPMKfua3/oa1rTHT/S7jUUNlDFsOrpqyfPv2Dgcqw1oIbKjCyQQIyCGwNpJ9YZUS0UOlysGqVE9TEv2b+jvlcjfQmIQhsN1VrBvDVuPZ3Bjxk//llxjVA171yjv4sf/0A/S6HcbjWhl3YmgcDofcfP21vOq2W1hf3+Tbv+sHOX7kFFfctI/ReMzcrhkeefowyxdW6HQr6rE4pK1ixrFRpfAadqTmuZig1xNz4NPPPMuv/tr/5NkXn6e/OMWBW/ZQ9R14KUKjRlRnt7qN0gEln7CFGGOtcp2iV5xOyEoYoNB/HxPPXGh406EeRQhcuWeKJ86O+OZUs3vnDGePdCmMZTBo6ATD/IxjfVN8YlNFIkXDPTfv5A8+dppzF7bYNVdwx417ef+HDuu+Vg2CQQkM+nMqK9e6tnNyZs6FiT5CoSPpwrYj3MZHUhLlX9N4KPOGUUpqZwu1MaTWCGiMXEoxyi6iTUhslyQ63AkRr4FPOR7AODXOZvKBK9RK4QXIqlgeg2l9NjHpzyB3nko6t3qJC4tPLMqaByLf/DhnaiStZIJQXYP3Akm0Tjfwiv7IVNOEjKPsJEQ+w7iytj47G03LN5JRhlVmjNUwmUmgjNFFVka3QyqERRT8BHWcNFmrrhvVSEfNFImqABNJXEubzPZ91RyJ3E4ie+VhEMla0uW9gBobysqpKSrS6xTcfuMecIkHHj9NiIlrL5vlv//VS3zz6xbYt9jh7MqYfbORleWRIE9me3zd6xf5yKOb/OPnlrhid5eZmQ7ROZ4+PyLFhrJTkHygKEraIkwNblHzUkIUzHV2xaeg5Ff1LBiN0DTo/NU56nrEmbNn6XU7utwUybUYKqtJO6tC3Tw+y/9MBpJ1Sghwk6At49pdBBlEKNZkiiJXqEJyluqlEFNkXpC3uxCDLeSiMlnychHw8GImV86OLnQ/ZTD0ux2mp/psbY3YWFln4cAOprbNYoxhOKrpV5Yvnq34sd095ipDcBaHpUmRqekODx4+xm9+/kW2XXWQ4clzMBrhnaEAVpNhvnJMWcsDY3h5LHypYVMztzDLyhUH+dn3fJzzy1t0qlIuWGMYlvC1O7v82qE5Luk61n3Epcg0lqlCjKJlKbJc4xOVk8/QONtieqxWmyZB6QylGnKx4E2i7xz9ouJDZ7f4T0+eZXZxju/85jfyzh/6En7yu9/GJ+/7BD/8oz9LjALAKzSNUj58y3A45MLKCp1+yTt/9seY682yvrJFp1vRmypZ2VrmmZdeEBZYRKSoQQ11UYoWSeszisBI7N61g2effYHve8d/4kd/4p2cG1/g8jsOsu+anZgE9bBpY4EzNdaVqpAsVRJdyM+2yBUv+d9JpklKQpgOTaCoLN3KcWorshoSnX7FDYd6PHmm4X99bIlTy0NmZ0uOnxuxObCMNscszlcsrdUcOddw/NwW1hhuunyetbURjzx5Akziluv30utUk3GwMRRWpiByZk4KJgEWpnbsYxEZbWs0NU4+e2PkudVCtyjkWS4VptjoCB2NAcjPt2pqJPVV/XhVWYoMXWOfcwJrWWgek9Xdje7DhIicWt+VK8RnlIKchyFEmtrreiGPv0uqXlfUqGWhggyNCvBeJhYpS1o13CfzWoyxrcIqu3J9U2OcoagkIF4yOiYLUqP8lInD3OqiXRzHVuewwU9icsVprKH1F4H8gg8UVSEPpo+U+nfmHOJsQBSscTG5uFQ9lQ+d2CpPckypEIarqmxb0Zx7XFWdtpouCgHjhRBb/lfwkWQce3fPcmjfNMTI5w6f4dK9Uzx7bMTLp7f4yjtnOHxkwF03zbO61jAew/Jmotd37No2zde9dpEHjzZcs38KAqyNAkfXPN1O5swoW0fzjY0CHZ21gr9wBY33CrwM7WFjFMHuVE2VkIu1tziF6zt+/dd/l9NnzspyLu+xFE9BpqNm4q1KdXPATEaQJwUkTsi4/Ls5dfYBCOuoFgly9LqfiSrxVYaXyoCxk4VrUsmi6kBa7pZQXvOzpVGxKbRqvW6vYn5ujsYnNjfXiaWh6RaMRzXGythtbANXdQ0/vq3g8inDgEQRDePGMzfd458+/Rx//PwSB/fv4djqgI8ME9v7FetEtlkxBj48ihwPUJWJhbkZhjdcxX/70GHOX9ii1+vgfUPTBFxp+L6d09wz2+PvVsb80dkxf3au4Y/OJ37vbMPnNxN1E+gZQwfZEUad7TdBFuYxQeNTC5JMbTEB887St47HNkb89HNrvLezk+/43rfz+//prbz9phnWnn2YxeEx3v1L38BjjzzM3//DP7OwOE9Tj1UKjXKTRJm4sb7F9ddfzbv/929xycIBtjYH9DpdhoOazz3wICZJFZzUCR8zgFQZWb2ufP4pwZ//5d/zzl/+dV44/RLbrphl9+WL2MLiR17H164dM2c0fwpJAZdCPw4hTrAfGM1gj63C0FzkWROfA1wYeM5virjmsl09ihS4ck/F00sdPn9SKvONoTDu+l3HQ0cHXLKt4P0PbLCxPmDnTMHiTMWnHz4JIXHl5fPs3jVL08QJvNU6Go28tYXTOFtaOngMQUZIzraq1aoUMYAPvlU7djodOUuYCFVijHpBh5bSLL67SbgdmjUjqaARX/vWKpErjrw/TCTN9EALgHjR2a7Jrl7G9jmDxBWTlEI0r11SXRWzXzdtfEGnU8kZbYylLI1eIrKFd1UpjkuVbKHVkCmEllkWpYbF0LLkUxDOvylFGdM0Db2y1xoRm8ZLV6PKohACZZkfFlFrpCS3Xmhk9ljXTWvc0XlbG+rjrDiUfSMZDq4oJvCzi2F3ShouioLxSJZWsuswrQM7K1MCjaJSBF5YFIUqXAw2GVzl2NyqufLSbezbM8dgecBDT57h9itnePqlDV571QxNMOzZ0WNjeYvz58cyb5/pQDNmZaNhz64+C9MFl+7u4ArDydWGc+s1ncrlbrLNXDZJIIZZkYEBq/wrkR87NB65xaHLwar0zxChSOy5ci9PffwZPvDBj/GO7/kP1KsrsiexrpX55Uorz3kzaiaZjOOOrTrKqpFPOrpw0cXhcKZQdlVJ9I3ElkbfKuuM5oGL+qsgeW2NY8Di2hlN0lm9XFChJRoYRY1Lqx/16yo4eHAfdTNiY22DsuxQ9nrUrNCpelQpstTA0ZHnxpmSn93V4X8vjfjURkMXQx0i2+Zm+KP7nuLS11/L0BScq8d0MKz7yFVTBRc6JZvR88AgMr0UuebG7fzpx55kfWNMr9tp8Thb44b5GPm3M4H/sVXjfUaCyHMZfOT+c5vsLi3X9B03zFZcM9dlLhmGqkgKVYkpCjByQLnKQAh0jHTBDy5v8okBLFx1KT/0ZXdy5/4+Lz/zDPd98GGePbHKmfWG+W7g29465Dve/kr+9u8/wld/9ds1qyRztVU8k2Qatb62wYEDe3nXT/wX/vO7fo7Tx08zOzfFE88+xageCxU6RpkYmISvG2ZmZxiPxrzw4ss8dPgwH7vvE5xcPs383ln2zewlhcBoc4StHNGKeKWoHMlruoU0p+hIX4Q0TaCoXBtWZmxq435Rf4MrrIzLNSKYlBiFxIl1z6VThkt3T7FjoUe/X/HNr9zJe+4d85mHTnL74pD1pQ02NhxHVgw/8sU7+aP7L3Dv4xf4li/cxnWHZnn4yTPEumbH9ikuOzDH2aUNSk2FzD6wyTJaeFKNnyyw67FkppdlJUFkOuUYj2tsRy8EzfAJquSSXYZpd54ZNivolKiBagJYtNaImMbLKM/p+Cvl+cG/o3cYRaXIeFC6IrSAjxhntZOQ4jOGQKmJpUIYkQkT6gvLe9hIIkY1NxrV+lpjiEqQrOtassFjUhNKavcHzhakkGclKrGtOu3+ITShjUbNKixfNzrfVlQG6s3QP9vESefhx2NsUVKPm5a5j+LXxRUcWkijKywmFZM0sJzLrTNJn+XHPjFsBpIBYiZz29wuGr3U6romWUtiLB+66r9DCDh1odeN5+Zrd2M7lmMvrHHk+Brf+JpDPPD0Mrce6rI1gh3zBWdPbeBMZGqqotNE1pZq+sU6D76wRVVaDiyWRAzHt8Aj7nUJsdGY2SRcoRw5KyrMsiVpZphf0phRE3NmhgR4Wc0DMNawtbbFJQcP8Na3vImx/iyISRfu8hlbJiyqnIGQ3eVGtmiqVdeKKIc1WzkBCltKIJEpNP9CZbwpYo2mrJmgTDEdUeW/Ax11ReEsNdFTGKP03CzdpfWuFLZsdeukRK/qsLiwgDGW//X/PsD7/v4T+JU1quCxNEyVlrO1oUzw2FbNYqfihhA5UY95cpSY6/XwBPpFwf964CWunypxQKe0zFSOYAypSayGwIwp+NBG5MOfeA4XacO3vI8MRg2H9s6xb9cCU4tTfMd8n7l+l26vR6eQIKDRYMzG5jrnNzY4vTLmo6fW+IsXLvDKmZKv3D0ru7g6MmwC4wCFTcwmC87x4Momn9ryzF55OT/29V/Aq6+a5amHHuOv3vsMjx8bsDyG9cbQqfq8fGaNh587yyuvfgV/9tGneO6FIxy6ZL/M5vVIzvnYKBdteXmN7TsW+JWffic/9d9+kQuj87x05BiHn3iKO2+/lY2NLUwyzE5PkQJ8/pGH+Z+/94ecOnuWIWMKDIvb5yk7siskiCEw6Si6KIx4akrb+hFUvS3wzSix0j6I8CBpDo01BqfCdadQ0RxgV5SWFA14z5FVzxsvnaJAfFd/c/8ZRoPIa6+Ypog7ePjRE7x+PnDvczVXH9hOrzC86tKS+59a4+teO+Ca/VP81afOcubEMnsv3cH1V+/hox9/js62aXLfLZLXCe+trseaRCiXnfhXCppxTVE4XVzH1kOXadJVIdJe63LOhr/IiBmUc+fpdCpSId+rWBykiysKjaCwssQPBI27iPS6FcNRQ7KCS8krhXpcC4ZdR5hJydISESEFq7OG8SjgSouJRqXEOq4yhZqINVrBWIqcHJclW7V2Hk6r+YiMN2zp2lwP66QFlxfYtowrkMzn0Xgs+8+L5G25G5CLJhBqT9XtalpX0ICmgrpWfIlmsGevQvCimHKK0zDalnkvkklaHfOkk8qmuzbykYkioSicjq5o6ZntXPAiO38y0r3U9ZiyLKmKgquv2Aa9Hs8fW8UAl+ye4v5HzrB/toMJkVTXNOPIbL9keQhPn/EsdB1lCrx0vuGSxVL8H9Hw0tK49dL4WhbPKUhyGtltnRk6JrZk3Sx3tTm/2ujO2YhiSf2AuNIyGo65+tLLue66q1lZXsUZQ5SN9GQBqMo1kxEUTHAGSblKBv3vuXLVjtDZAh+aNsLUWoP3XrArwStXV3Kyg/ftficiYoGkA/FMcHVGAsZsMm2HEXUXIAs9GUmORiO2LS7we//nT/jd//1u7rrlaq6/ajt7t/eY6VzCaLjOxsqAUZ24uS9Ik5G1bPrI6c2aS7emqdaGHDk3ZGMkUbJNSHxuveGWKcvYJ75opsufr4/4pmnHLf2CB7ciCw7qVIBSoW1pmZ123H7dfm67Zg9TnZJOWeKc/Lz6UyXew+yMFAbnLtTsXGt43bXbwF3C8fM1v/OPjxLPbfBNu6cYEqmbRFehoZ9aHfDptRGzlx3gP37f3bzl9gM8+9hh/uQPPsYTR1c5MoA4jGyMaua2LfJz7/xxfvKnfplnjq5yzx2e3fMlTzz1HNdefQXLF1Z1HKwVfBCTnrGJwhhGwxG7di3yu7/63/gv//WX+OhjH+e+T36au1/9aopZx8bmmE9+9vN85JMf59HHH+Hcy0sMBjXbFmeZ7lSceuoMUy9Pse36OaqpsiUdGN21Flg9R1ybPOqsiCGKspJ3UQnISSX+Qc+doEmb+Im/N+Tq2BqOrwW2Rg3zMyU75irqUeCynSXv/sQKN+wx9PsVp1bHHD5f8mOv7NGYgqv29bn36SEvnB1w+9Xz/L+PHOP5YxvsvX4/N129nU6no5ONTBWOBB+pCteeeW3ueUy4qmgX/Vm04MNFJF4dW6UkyKKgMbfZW2dsIWmX+v4GP+G+WZt/fRYzBJKVizVo7LKxMBiNhJihY8+m8ZTOUGn3JJYJ+XzLSvZaWYgzHI1l99YordtkVZm8i4JmyUTySEHecRgBJlqNnQ3Rt/GOVuNE2zl3m11uoUhy27ZZvsoSalVPFu/FQxEya8tOwlWsK6TbSYmm8To+qkHn/G2X06nECdqo9V+9K2VVatcx8aoYLNFMoH+yQC/a0J8YfZv/XpXinG+8VOYxZPWB0ehOxQ+UFclAr+O4dN8cjEc8/PgZZruGqa6hWxl2LnRYHwTGW2O6BZzeTPzJ5zYZeMOVO+DAIHB6NXHVLlGoDY3l+Lqn4yxBA2SsswRZTlG6DiEFVZ8VUpkrGkHUHU5dvsqKwkjLrwiQZKW7Wtw5zzPPv8D7//GDvObOO1XhJg9EoS/Fxex/iaLLaHJVWxhan4ihJcaIRyVKXGlOrhNxRQefNPMh6ufuvYRZRa/LxqJN8Iv/Lr8dLWCCiiuEB5YU8RGSIDvmZmd44PMP8eu/8Xv80Hfcw5vuOsQzj7/IyVNnOLq2KTTVUaRTgQtjrIHtMxU+Jtx2h99WcpnvcsdVi5xeq3nipRUurEdMabikV1I4y2yEeyrHcVvwk/sqfur4Oke3pDstTaJTFYx9wzhEHn76FB+9/1mqqsAwZnG2x0yvpOyCUxhj7YX0OvQltio4OGu56ar9fPvdl/PuDx7mrbv67KwKmhD56No6n1jepL9/F9/59W/mC28/xOkXnucv3v1Jnjy+wYvLEdsk3jjlePX2Hu884qnrhkMH9jG/sMCFtRXqrYYdcxUvHjmh1G0tCFVRUxadlutUVR2Wl1eonGP//p381A/8CC889zL3fvJT3HnH7bz08lE++eCDXNg6h3GG8UpDs9XwXV/1ar7+S1/F3MwUn370JX7pdz7E8rOr7Lp1u8BVS9OaXItuKQqsOrSxwD4IkymZKPsWpfwGVQ85J9J15+TCQ0UjLuPKraXbgaVx4PxazULXcsWuKT7+2DL7dnX4mtd0+MiDF9jyJf90OnHZzooDCwWnNhuqomT7bMEjL63zxa+eo1MaPv/QSb7gnoNcefkcncq10tvRcERZFJSuaIGvVafTqjp94yVK2lgZdV0Uzy2fucF7uRS8CW1gUxtCZiZxtmHctJR0ky7aT6tsOvjQ8qlilF2GUWl91KmKV1Wp4J+icgF1X21kDJVNn9k3kmlVOcMm2whklSHfV7wIW1Nko0s+FUxGeacJhynoIeV1Zpd9Xjl05uLAIGnpbauNlnmhaRdf1hh87SX/16IVqZhymtC0s7bgVXqryXVB42XLTglJ2Czuor8/qO4+YVr5p4yrQtsyyi6FCX04f83WUjpRauFS27YRA7iC0XhIt6pomsil++c5uHcGRg2PPnueKw9OYYInjhuiDzgic1MlD7044A8/t8UtB/u85aY5PvD5C7xil+f82pivuH2GFOH8yLMeEkrokK5Auwq5x72Ys4ygr0tbqqFSHr4cbtTSaTVJzGbDJ0ZIsh2DLUre+fO/wXd/5zfybd/0daxvrMmlnmmjRtDdUu0HyfFQrpXJgU5WjUs69nJWJIXybDQYZEZqjCOkWuS8yrHKMcAhyK+LSVzGVkUDKTvgdSeWKQfRprbKs/rzFCFBQ9Xp8Nd/92FecWCaSxca/vTP7+NzzywziomZfoeZ6S51k8BZStMwrmvq0Vj2Zi4xN9VldrpktleyZ1uXN9+yjfPLDZ87tsm2EBiNxlwIln3WUBWCnPjy+ZI/sUMOzJQ8eSoQQ0On22F2ZpqdOxe5+3U7OPzEs9z1mtt45NHHefyJo3T6fRlNRc+XfeGrWV1b5cSjz3NJ1eHxkyMwx7ls3wJlp8ODm57reo7ffPIMS1PTfOe3v5W33XkJG+fP8E9/8w8cPrLBiQFsjT03OsNX7ao40C853Ri2Rg2djoyPq6rEdSqG40i322F1ebXtGjN51eafcVvFW86ePcMjjx3mK97+Nvbs2c6P/cB/5Kd++Zf41f/z2yQD09NT7Ny9wPKLm6yfXubnvustvOP7vpwN3+DrwLd+xT6cCfzYL76fuZWa7mJ30vF7z/kjF0hjmL9kDlsaTDBtfnrSuOsQJfrXYjQETZ5ro7vAHA+Rs1vQZ2OzjpwdG652jj0LFck4XjozYtvcFF9+904eeXCdOsGtOxMbY4+rumxujNi9WPLkyxt8zWsT+xYrnnx+CZqGQ3sX2bt7jlNn1+hUrj2ws8HWKCAypKiEXWHVGYUgCvrFtkbrTL7N4/18EDvrwCs5oA5tpEYmo5siYbxRg3VUUZHkIKUY8VEC9oTRlem7pk00lFUBKrBJimVXHLqqX0UwY1SurFaKJOrIpvFU+q47xQqRIBoocvtlNMzEe0/jPWWhiHV1hja+odOp2pzcpGA8o2iB4OUFKcvqonjGSlsrMfWhwVWuktQuP64nDvQCSh155HZNzGo5yU6coFGD66tK3JdlVbZSOB8mHgajzvGsIhoncZbLSES+31KNNMZIip8PjcaWWgpTElMQqaCT+fW4Dlx6YI6F6ZLVpREvHlvlDddP03hPaAJbW3IBvvszF/jMCcP33TPH4vwUr75xgefPNXzy6JB+17Fj2tCZ7nDi/IjNTc9U38nILkaSTUrhtAQilb3IiY38YK1WYsYVaM4naExo0vxpq3iFhGVqvsMZf5433fMa3vGd38rZc+clo0MRLvlzcq6j46GideFn1VR27bf4ExlCy9+pUEc0u96SSKbQ5DrJC/cptIl2iUBhC6KAOFqJrnwNgm6RLALB4hgn8mYpUCS3pdvpsLm1yZmzSxzY3uUzD73MIy9usdkEfvB7v5Y333MX4w//Lf/2z5+lefMb+Jqv+RI2toY89dSLfOzjD3PJ/gP88Xvez9yopCw95zYH7FpI3HbZDFskto5uMY6QgqdOCdM43HjAsWHDjhnLXM8yHg14x/e8nS947SuZ6neZne0z3e/zq7/1R/zkj34/P/MLv8YTTx6lYwwxwqgZ8y3f+FU8+uwRPvzwY/zcwR2874Th3o0Rl9Gwe6bk/cc3+efpLl/w9tfwda+/htHyaf75g//Kw0+f55x3DGvDVaXhjq5jf9/RKSybMbHReBpjqBK4suSOO27k93/3Aegc5cHnzvJdd32hoGfQKFJkgVpo9oN1lsFgwA3XXcv6+jrHT51mcXGR2akZ5hdm6M536fa7YBIrR7c4+sxZfuQb7uT/+/++nhefeJHxY0/DaMDWbdfy5ntu5uo/+wQnz2wxu3eO4fqY1TMXGK2MuOrQlbx89BijtRHzu2dpvIyrM4xSAuFy1LLM+xOxfT9cWRFjg3HgG3nOXQnJR5K1HN9K1A3sWOhCiBxfGrF9+wynziwTm8Tl/TEEx2AEIz9ksLlFr7A8tTSiqcdcvrfH0y+cY7QRWNw5y+WHtvPyyWWqSrtyJyk6Ob4gC1vQKUBRlBovKwF7SeMnrBPVqK8b2XOWus8kidIpC2e0Q6A16crvqapyElfrsmJRRnchavCZj9hCXeN6xhlSG1MdYxRRjLVQaBpsjuwNTSvxzXeCVcwK2RqgF0yhXjofA0XVUUluCjRD+ULLImeiW1wnp7AFxuPxJE7VGTXaiTTTFeLMDT5QdSpNstKZed6dlIVqjhMpZMeyxTkYD8dU3ZLCljRNQ6esRKWgJkP0G5dMDsmbsIVtl/fexxZRnPXTqjZvWzejGeCZDjoxsQUlYSomXB3yuXrOBikfAtdevR0z2+PU0bOcPr/BKw7sYqZfUrrEoIbf/sgF+vM9/se378I0Y46tR2zH8qWvWuAX/2KNua5j97Yp6pHnxFogpkl1VjiB0vkUKJMu2lKkMKq2yvniJJJ2iQmhuxrr8DFohUArow6KfJ7ZPs2jjz3Bhz76Me75grtZW1uVv0/FENZYdd1bgq+lw4kR29JJchiXl0M8ejUvSoSwVUm0swLbM0p4TingkyzgZHSVsLbEhxpnK6G7KlI+mUQTGgpbElNDwoiqy2dDozy0zjkaX9PrVPimoVc5Xr7Q8IrLdvLIk8c5uGc7u/fs5uXNmnnfsBlH7No+T2kTV1y+j098+iliSPSqkl7pGEc4vwGRMWuDLs4aVkPD8dXEpjMM68Boy+MxPD5u6PYd68PA1saAmX7F/t1dLiyvc/bMEvXCAtiCn/353+Dj9z9Ir9PBWJFxSjJnzZlzS1wYBc4Mau5aKHn47Ji1UaB0kSsu2cZPfs/djFbX+fTH7uPzz57nhdVIEwxXF4G7pioOdCBay5YPLCTHVKcgNBCto1tVDIYDvv1bvh5jS+7/1Of56q97Dd/wdV/G1tZme4kXtlDfQmz3bMbC1mDAHXfcDibx1DPP8b3f/yPM7Z5npiPvZdj0vPzMOd5w+z5+4Ie/ipc/9RgX3v9BxhubRA8bTz3PwttfJzk9w4bTz5+lm0puecWNfOk7vog7X3Un3/cDP8LJjdM6XZiMUmOOUNZ0Uhl9yikd21Qg2RN430glbxK+AazDEDi1KcmnRUgszHV44ugGtx3qM2UblldljJmcpdO1bKwMwFhKZwmmYGucuPrAFJ998QJnlrY4NNfn8oPzfPBfPfOzfcapEep11H1eTNS1qFLlYBWGXT0aKRNQLg8JeEqaCGja3YexTHh/RFKQmYtvmnYKISouMQt3qs5FAX1aWBuDLQpBwVg7ofCqDLdSRa2Yh2MLhC0KuaxNdvy3YXmWsqzEs4d0sx5DXXspuGPEFBNRSzEe1ZRVIfnNJurG3bRsKR+Tck+S5kvo8rZ0hNqr+im2wTH54InacjonSYBVVdLkVipLz1L2NFlcskrfFdhf1DlfPsHKspDFj49U3UolfqZFhTtdqhcXpZZF/WemS8aExOh6T9Wp5MJRymzObU+a95tn/GLeE7dzSokrDixiyoIXji5hY+CyHVP0uiWdwvEbH7nAVZcv8jNfs4eXXt6kMzvDnIlgCnYu9Lj2sllOnR8x1bUEA8c3a1FAxPwwiCKtcDk7JbW4dAnW0dl19KpsknGdM7YdMbX6ellSYYB6WLOwf54Ty6f42V/8Db7rxEm+81u/nuXlFX1JpRubme3LLiMmtgZbLarf5KWfKqNSVCOnyrsz+FGc53IBJJParARnLSE1rV8opqBdTiOubE1Hs1p95bxsQcB7nJFuzBorC/YQNK/C0OkUWCIhWubmZwkJTp46TzMa48cN/akuz51eZnMwZGVjg/m5aZ57/iU+9OGPs21xDpLQbedn++zdvcDKxhLWex4vHM2BeemYjCFahy0Da8+v0E+ekAylMTz1zDG+8I23EVLCuQpXlNx4/eX88H/+DXZu3wYx0jRiyt21cxuX7NvNseMfpltUGO8xzrGrKFgbeMrKQPJ8/r7P8eDTSzy/GhgFy+XWcFvPcmmvIJrEVjJY7yXkSiN2NxoBP1Zlqa5hw/d/z7fwPd/xtbiioB431F72mp2iw7geiUw0m8asYzweqwpzBNawa88O3vrWN/Ohf/lXegslvV6Pl588z+45y8//9NcyeukUJ//oL6FyxKIgYGm2xrzw3o9x6twq+y45wHd86zdy7VVXcvDgfgAGwyHDrVEb75Az36ONE4eHVb+PdtFGtrkq+IstOj3G2I5USYmqhHODyOrAs22hw5UHpvjk48usb17g1gOyy/PBYFyB8yO2NhumugVbWw2lk53KgcUOW1sjjh5d5dCNl3D5oXm63VLjBSbU8KSRxQEZA2MMMXhSI5k9VrHngmCRKYatHN7Xah8QE1+36GnKnzwLzpUUlaEe1bIXUq5WUUoOkG+alnYsEmJPp9ORjkSD2wpNScxx5FUlF4JVEogPQQQqirWIPrRkB9JE7ZVl+2VRYIOM3KR4861TvnCFUV23b+NmxbNRtIHySZce8v+3GCJ+3LRYgzyC8CFKpoKOkJrat5Tb7BivOh1Jn3MWoz7E0DT6YWkl62W04XQ8BbJgz3DAnACXD00bbMt+aWqZJRbOqnNTUAJRyhT5Yaj7PEvSsqJJWF9RcdK+veTAMByOmJ3ucGjfLPjAsy+vMNcvOHTJPH7c8PljDdcenOIXvu1SHnjwDAcvXaA2jri2TtnrkhhyennMVXu7OGO4MPCcXG5kH0iiqCpZmJuE0ZjZJnhKK+5hH4LKBYNA/ohtfG/SL9RaQ8wAXGuJjVI+A5gKDt64lxNPnOFv3v8hnnnmRZ5/+iXqeoxzlv7UFLt2becVlx/kppuv56Ybr2NudobNwYimrmUcGX074osxgC1IePWJeHXDFySjQWSuaEGNEoglD+EE++Ba6XF2mUtWiW3z0VE0RJ4lo8DIoO2/hG3JqLMqS6b7gr53JMJ4zHRVcPzIaVZXVpid6lKUJTddfzkvHjktIxNrGQ1HXHXb9cwtLHLm2ZOMxpGbb9rJl7zuSs69vMxoMMAUlo2NAUf8BWKSAK2qcJw8vYz3BpccpirY2Bpx681Xccut13H4sefYuW0WHyJnz6/yNV/2eqZ7XY4ePcVC6fAx0jSGGQPnfWB+ussLx4Z87OmapQ3DJcZxy5Rhl5OxxWaMxGgobaIEyo6lcuBSYmQMAcPs7Az93hRN03BhNKRTlYThSEcVjkLzTlxR4n2tWHJH7RvKslLJuKGwjqqwvPNnf4zXve7V/NO/fJgHPnuY1Hje+dNfxb7+DIf/+/9Rao6hHicwIr0/NoSloef1V1zGl73tC1lfX2dtdU2Cpwod/1itnJugIhoV5mCIJmGiUfZVcVEyXo5eTRn2NElITAmbYGmYWAuWndax2Jfn9J6bejx3quGpU47bdidGY8/5pQHr6549C47Ta2N2L3TBRrZNOwoSzx1Z5fWN56rLd9ApLI2XPJLxeCxjKpND00xrpg1edgTZBGt1b5ANykYDqnwjKkLb5m0g2HvtGmw0Leg0c+FimwDodHwsF6fTKNugZueycu2aAZJaMor2MvBqqp5E8wp1utB/Z5wljpv2bEwxQSn84HyZF1UJXgbQFl20iqlPGCsuu4udyRYeYko0vlYqpW3NZzmtTaSdUnU2Gv9aaAeSlLybEfDGGAmE0lERCY3epPV05BZNAq60+izKVmWQggIX9cAnSuA8rZM9KhjQqNlILAuFBlRFxWRn+TJmwnlJScZWtJkZBmtLdsz3ObB/EULi0afOsmdHj92LJX/ygWfpOM9PfO1BPvfgWeYWOhy4bDvnT6/S7xWMNgesr65zdtlzcHuHonRcqCNbwVC4pEFcghEQ1lNqAZdCxkVGQ16jYYmQXPvfU6botrCs1EZX2mSwlWH93DrPfvYYx46c48yJM9z/b59gsLUKeJp6wOmTx7jv3+7nf/zOu/nu7/lxvuM7fpD3/Pn7aOqGubl5rWhU2qwLcR8aGXmlRneZQfMPZKQmh0KWIafJTku/hyyrFnOiqrWMbWmvcolkLldoEwxDVHx5ku6raaJiLwzdbsGzLxzFdSrsVIfF0jJYWuHFI8dJXuCHX/LFr2Fmpstga0SMicFgxLd+3duoOhWDrZqIY/XEKvf91YN89mPP8einT/DJe4/wmU+cYWtTzKd1jPQLw7kzS4xq4QqlFMFHLA0/9H1fxe5dC5xfWmN1dZ03vO5Wvu+7vpqXT5zmuSMnONhzjJJ0wE6DibqlZX0cuDIZvm3W8dZpx76eY0yiDqLpAIOPMIqGsdcL1BrWa3nfpqf7VKVAMDudLjFCVXZajpoIRDT3xslcPYYJ9y1DMoOPND6yub7BV37JW9k1s4sTL5/l27/iNr74njt45n//NcGPSVjqOonSyiUKY3lq0DAc1mxfXGRrMGR9fVOqZpXs5rC2xqsJNURZCKv8PMtPCy1UhKij2TwZvRaYjD7bpysx9pGzmx4/HHNozxRNk9i3a5rvffMurj3Q5bGzhuAbTi83zMxU1D7y2LEBVx+YZlgbts1X7N85xXPH1qDosXfXNLMzXaUhQLfb0zwcXZTnwlqzx23eIzcNjUZhNE1NCIEmCA6/peHabI6VKUeKsU25tApLjLo37nQ67RSiJWinqPyqyV66qevJ9MRYTUKkvWTzdAkln3svXjynme31uMaq5DelhC1Eqp8Vl1G/xhyJW7TizARl1ZEUP80dAPnGglJ1i1IrSibSOqsodOfE7CT8f2GypBAvik4sFH8iPCxhJCmErXCyiMqXlfKe8szQFYqa1hvQqIu1qT1NaIilHCgCZ2zaqMi8D5FUMEGABJW/tohoaOFkRpdXbcKiHuJlWbC+OWLnjim2zZSMVjd56egy1x+a4sTZNT7wqTP83DdfxssnB7hOxSWv2M6RZ05w+uUlbrjjUvzmJqfPbeKxbJsTGOXSINAgWn/JUhHRQlV28bGRwCErVUnGqztTtDI7K+wBhSDKg8tFSzLTRIrKUY9qzjy5xtHnTnNo9yz/8Ru+gJtuvIp+p8APVpnqFvSnp6mpWN8a8eLR0zz21Mt89qGn+bl3/Xf+5m8+wA/8wHfxpjfezdZgICMAI2o358QxXChWxdmSmMQlnpDL3CRF9UePVWlwjIHCVKLcMhfRDtQEKvsq1x4mVnc+GQHvNCHNlZZup4NpEoUV2eJ0r2JrOGY4HDEaNEyXBZ3BmH/75KPcdsPlXFjf4rJLdvLzP/tdvOcvPsLayib/4Vu+jFff8yr++gP3sTjTZdk7OL7OKxb7DKyhU0DtCmY6hvus52wyWAc3LPSJ401eOnmOm688xNrGJtZGhoMBrzi4g//1Gz/EAw89x44d23jlbddQFYaHnn2Ren2TyxanWa8jM1XBoAmkWLDVRK6YKfixQ7N86uQGyz5SeaMmT335gxBsPQm8pW7A9R1j52h8pNvtSVTBOPt4LD56qrJL7UcSvNXpUfuxSq6DRmw4Mhc7z+Xrumb74gLvfvef8/u//yfcee0hfvC7vogX3/NBNo68TDEzLfkz1uGI0ATWreXwlqdXOm659eZWAuqbhpnpKR597Emef+kIB24+gDOIZN0K/yurmnJCaoxRx7Ni0jPIwjqYiCksJshphOadWGOJFpZGciAu9gymKDl2apPpUPG6K3s8c3LA2CfmZg3bZy2feWnE8WXP2sYYYw1V6TiwvcPxUxswHLJtrmJuuuTC6hYz0z3FpmtCH4miKOWiIIlE3XtxyquHqiwLxnVNVVXt/skWjqYe0/iGsiwpjCzey06patOJTN4qiFQkzTLKDwo8JMjlJXywSehVltm3OSU6BYpR1Ju+DYkyrdl7PKpbmq9YKpSLGGVV0MZXW/EKZnhpYfQgtS3kUKVpKpv1qsYprMUkubFK/aAkkzu2IT95vGEySTclnKIpospLmxAodD6YHZIxmzguijfNv14kbxlMNgkvihqbWZWyC/C+oeyUbdxqYUQRhIn4RvAlVVW1Uato59R4397MeXl+sZw5S5SDj7zish305nqcOb7EiVOrfN3rDvGB+8+wd3uXg9sdTx4f8YZ7ruDZR15i+ewqg62Gk8cvsHCgx9ETA6rSMDctmuwL4xxor5eCHpxNaERulww+BUwEp87rojCEoFW8dikmJYxiImwhksbUiMxv6/yQ5x88SZEC7/yRr+BbvvrNHH76BI8//AAHyiXieIvVuiHh2LFtlu2793LZzQf44tffysr6W/nwfQ/wnr+9j//vB/8L3/j1X8YP/9D30e/32NoY6Bw200GDcsN0SW6S/OwV0hk1Bz2EpkW+h+TbvMJcPWX1SYqpBV8mNVElRFYeYwRn8I2gc+YX5xiuR2b6HYZ1YH5uhtXlVSDR63VYHzRcMd3lw//2EN/85fcwvTDLcDDkzluv5LYbXkEyJdt27mVzY8j65hqzFjYHgYOV5ZJewdhZUhAxx5Rx7BsXnLeJ2Z7Bdh27x0MOP/oUd1x3pezaCodpara2tpibrvjyt72KoqhYXd9gutfhr//xfq4sEgemKtZ8ogyw1cj71kSYNrBRQyPVne5WRDxSGLiQIoM6sKssKSrbxpeu1hKrMDs7rWhugzBI5QIh1S0BoPF1y4HLAWVZSutjVrsVTE31OHnqLL/7f/+YnYuzvOvHv5T00gnOfuYRiuk+zVhm/KSAiYYiJj6z0fDUhU1uufE6brzxegZbWzKrr8X/8olPf5ZgPJ3pSqJYW6mu+fcsNq3EY1aOTCJu2swMc1HeUNLldvSR01uRzUFDQWJutsODLw3YN2NZ7Fl2zjmCdWyfK1jbqrn/xZpvu2uBZ49vUTcNZbJcurPLJ164wHBlk8VtfXbt6PPiiTXtlKRjqDodxuOxeNOCGAabplbjoxTWxgl8tiolzK4q5YyN6isxaouo67rt+lKKNCEX5kEymVKi9g0WaGr5zI3KcXPgXghRVKxRGFtFZorpZ5jxKCbJ2CtlJH2arCdSnlxYS1VW4r1RN3tyGnUdJqmnIQTdgeatvXYfRh2NUY2AMc+brdHoRLlYohJgQTN7EVRx/lCi/joZedBGk6aEILg7HalOvJ+orPTBEC204NibplYWvrohfa0ZEvohBHFhopYWpyqHwml6Xoa2pURIoeUwkfNI8vOry2JXOJ1b8u8u1EP756F0HD+9zrBOHNgzxWefWObuG6Y4eWHMgYPbOPvCSTbOrsufYT3PvrBGPWpY9RXzXcN0vyKVjtPrvoVUmsz/N1aR2ZO8DqPo9EJnpDJZM+1LVijN2DnTareLbsmFI+s885mjXLpnhn/+q3fxzV/1Nv7rb/0F//CeP2Bm9XE2Lpzh/PkVTp9d49SZFf7p/me57yP/xhP/8vc8/6kPMFds8gPf+VW8+3/8KG983c388R+/j+///p/g9OlzzM7NtsUGRg7y7D8xOlMVf4fOfW1BTJ6y6OjC8aLcB52hG2NaB7y5aEeS1YAm5awQ00rDDYa52QUiBXPTHYZjz+Jcn3E9EuNUafGF5apt07jlNX7rT/+R2V6XTq/L1saAxjfE4Cmt5eTxUzz3zPPMz3RJhcWVJQ0wGjdgYVwHhiGwWBq2NiNVYRjayIGpDg//6wOcX1miXxU0oyEJK+gbU7CyMuDC2iZ7ti3y4Xs/zyc//TBfuXeeBmFCBWBNn83BqGFKjbw+yEjSGtl5gMz4V3xiPRmhCUTpFoYhsRoSMXoWFuY0zji2scfOltrFmrZYNCoxtbZQ307T0paNpm72ul3+9M/+ihdfOsq3vP1W7rhiH8/82Yepk9MoX3nubIyM6hHLo8gnBgkfI9/6zV9Lt9PVvAnD1PQM55dX+fj9n2Hbnm2qKDSqkryIzgzYJB214EFsS+nWiknyKKxAJtsRrpIKLImzW4FBDVP9gsW+4aVzY/7+kSHvPzzGe8O5zYiNgc8fHXHFnh7fcPdOHJEjZ0aUVcElu3tsbI45uzTAVpa9u2e1qEySrGoto9FQRlEqMpLlte4WtDv3jW9ZVMYKRzAnsKJiAOm+47+LkO5WHUji4whRgtCcym1Ny/ez7TgqxtSuILIM1zcNVVWgOQrYQkZn7aVb6p73IjqlybkrOm7zjYziUrYR6I44RyMXZSlCjqhE3kJZUa3aQHcRzjp1O9O2PJKuNSE55uo9JaOdQymsFDXfpKjLVA2MClG+iNhmWaQ2QjJnTOSkL6v4DO9925lYM+HnC5ZV9yFO+DPGaX6xYhEy/l0U8LZN1pL2zl6UZ6F6ac0q8CEoItlycM8sGDhxbp26aWjqwOnlEQcXC86vJ3bM9zl3Yolt20q21sai9cbx8gXP8qZnvmcpYqQOsNxIXgc5A16Xy4UrpWJElCZJH74QfcuAsvo94cTk45wj6MtkS8uFoxsce+I0N1yzmw/91a+xtBp4xw/8NPvrp3n15VM0PrG22TBqItE6mgjrvsNnTpY8eHTA+OTzHPnEBzh8/4c4tHeR//lLP8APf+9X8ehjT/G93/ujPPvMs8xOT4uwQfcTeSE+iRUWFU07mDQWH2uN19QhXCsNDqr0UkOqsKtl2asXqIxYkjqS5cIqy5I9e/ewvNEw03UMBwPqFGh8pBltMd4Y0XWWR0YN37hnjk/e9zC/8Afvx5qCbTOzbJufZ9vcIkvL6/zuH/wZ/cEWnamS0aimD+LBsRaTLJ2qIBnYXjqsiXQLw9hBpzQcGm7w67/zV6TQMN0t8eMxsZF3YfviHHsWF3jvBz/Du37tj/nuvbMcnOlSKy4npMTSuJbLYbOh9J6xClM6Thak8hgbqtJSaIdeh9BmNwcDA8Sz0u939dk20rlqPnZGxvjg24W1U3SMCOkcIH9miIFuv+KJJ57hL//6A9x5w2V877e9maMf+CzjzS2Kqmyx8sF7XMcxquEzY8snTpzjbW94LW9+0+vZ2NzQIi/Q63V49NEnWVq+wMK+RZWNpjZcKRsIjZUlOjG2asScWNom5DkVVGSvhEYRJAyFs2wGiJ0OM9MVM92CmR583xun2TFjefBkw/NrBWfXPc8vGb7q9nkGdcP+xZITS2Ompkq2TReMm8DSyib0Ky7ZO09T1zIZMFJMVp1ue9EKPl0ktBh0DxLVV5dHXYWqDF0bIW2SQmgLKZREhSoFd6nucasgzsIV4iK/iAs3YUyadvRlrVwCTq0YAmOVsVXUnXGbcJgTX4tCkSlScJfdjo7vJQrDarhgURai2NXpgEGNhKZFsMsLLAyqgqaR1so6K/nlJEoF4yUNdMmjqMI5aQHLUsi6UUKeMvfflBOHeEwJvHo7NErRWPDRY9Hc5iZfLFESBHPLmn0f6o7P46Vup1KdQ8Jr/rQsVTuMiTp6K1tOk8wQk5KAPWUpkkaH1bazUj26XCjT/Q6X7J0FE3jx6Aqz/UqRAVGyDFzF8uklvJcUx7KAo2fh7tvnee5Mw6MvD3nNtTNAYhhhvZExW1Y0FVpB2GTlUjQoQ0z3Q06W1mVR0gSvOd2BsnRiKgSKjmNwbszJJ89w+aWLvP/dv8w/3/cY7/vj3+Nr79jBcBQ5c2ELaxQZrqDEqnTcfamj2yl47FTDB54a8cYrV9htH+b54QUuufEL+KHv+nK2LUzza7/7V/zIj/4Mv/1b/5VXXHkZq+ubbUaCdEX5MlFFm3YpOd9ATIiyPxNqsnQaPtSURQcfxhSuEimvLQRvrdJfayxNUBmjK6gbz+WXH+IDlMzNlCw9v8rhF05x7eX7MMEz3NykkwzbuhWfGA/48b3zfOifPs5PPH+M6265gZmZOcxoyIlnnubqMyfY2NYnVYYwlpyaz27VXF86RiHSLUWuvNM49nZKSmcw3cjK2HHzXI8PPvo4P/Cz63zvN72Fqy7fS69bMmwSDzz5En/3D/dy7uEneeeBaXZXJRveUxmRyG+NE0tNZI8zDHxithSYYBOETeURr5AUDeBToq9U1W7HUVhLnWBpq6awlsXFhZa2G0NoR0PGWUyMlEVF04xVgVMQaQQ9ZApibCA1oEmf7/nz97G1sc53f+PbKc8ucfqTD2F6PVyUWANnE8nXXFjzvDS2/MPyFgf27OKHf+gdDAfDdhIg9OvIR//l36j6BUXX4cceWyj4tFDya8agI515vlDy4WiQzjtHTGPBhEj0qY0YcAY2x5HNoWd7t2T3Qslnn9liZavhi+5YYNfUiD/99CbvfQTuvmaag/tmWVpaZ/diydn1hhACC7MVLiVOnFzh9hTZt3tWzkY9OG0e3+gkwDcia606ErxlnGmZes45NROq1Dd6jBEPjrFBF9ROzZNhEial3VmoZVUgMtwJDTuqJ6qyhqosGTdNu5gPal8w6iERQkXClg5MbCdOMQe7tSFUIkm2+u+dm7jOS13WW53mCOPLU2TfRlEYBeXFdo5XFG5CQC1LGl/LJeFcm1ucqZR53pZ0nCXGtIAhtS5SQl6oCfsm53LbwmHUbCZ+DrlaB1sDyrKgqZNmrVcqR1NndF4YpUjdNG18rTGi9kjGKF5ZOqYQAqWbUCwL/c+5Cg7K2XfW4etaLhTrGI0bprol87MVjBqeP7rC4lwFiIwSW1A5T70xwFloxpHPnEjs2d5jx2KX5c3AiZXAtpkKWxZsbAU2RxFr00V6d5WxqnFQ07/lMHDSztu8PFbPi7StQbPjDYPVMccPn2F62vKnv/lT/NNHH+A97/49vuPuvSytjdkYBnxIOCuQyMEwEYuAJXE+RuZIvO6yirObHT703IA7B1vcyMsc+ewH2HHd3Xz717yJsiz4pd/+C374x36O//v7v87OnTsZj2oZh+DaeGKTNHIWJgt2TaKMSSpin2pKVykOpdQuq4OPjcTpZlEAyteioCgqku7lah+49OB+TDlFHRKdsmBjc8yr77qNqdl5SpcYOsfrpwyzxTSfHdW8aucMB4crxMMPYXsd3HiLO23B08kSFiQOtjuOHG5qDkRLp+MYpyi7JyuH95UdxxNrNdvnKp457ZlvAnfOdDn88hHe9a7/w65DB5iZm6He3KBaWuLuKccbX7GdtRRZGnmMT3ib6JYlz49GbJhEtywYDDw7dvXQnFiaKBykbmGptfM01tIlSuxtEOn5IMLQCAdubn5uMpbQ9MX8fKUYCanRPG3woaZTdqhDIyFytiDEhl5V8MhDj/Ghj9zL6++6mje/8mrO/L8PYAonBzZJLnZg2BiODwJ/N0ic3xzxh7/yi+zavZvlC0uURYkrK7qdihdfepmHHnqExcsWReWoog+RacvzksdsSTT67SQiQwxbOq+V/38OpLJOxmRRCQzjOrE6CJhtJbOVYX3giWXFlimYLiJve0XD+5+peM3ljkEjQXpzMwWHjw8ZDBtmpjt0K8P5lRqw7Nk1S7cqqDV3HNBOWkbxeZdQj8faDaGcsUIT/uRMNbq7Awm+w+b4bXW1xyiSW+dkMZ8iFIUQlBV+naMxXFG0ZsBGA/KMgk+tlb0MyjdM6qBPMarPR85wp34SH4Th5axtF+vCKEst+dz70KYcZrS/5IHApDLUWzMqrltAd4r9NbKk8Rk6GMXc4kPTLpcwlrpWaq6SV3VDLJpoLx2GK217WzZNQ8dZOdhUNRBDkAO+cK2ssKwqOVAy3VPDhqKaDNHWLtNjyXz97CGIk7/Pap52TlULyvcqiqJNJUyKbXbqT5mf77G40IUmcWZpwP5tHUnvi57NzTGdXqQz6/CN4w/uH3Bk3fJjV5cMhw37d/XZta3DdLck+sgQS7CWwijEMDkiijeISaJkbSG7AStY+jJL6IySa5XCizXtwbL8whqDrXV+7Wf+I+urF/iD//v7vOONB1gfeTZGAe8TkcDJ5cCpLcPZNRg2IgKeW6iYqSI7+g1Xbi9545XT3Pf8kGEccsclkdOP/Aup3uIbv/wLGI4j//W3/oyf+dlf4X/9z1+VhzNEIkFnsqGN2pTxiNB6XSvtNfhQy7jO16rM0s4yirkwV3it+91amjDG2Y7MoQ2MhmN27dzGws49nDx3jH3bu7x8boov+5J7GA3HNGNhlPl+xRv7fe7yXT6z1fB4HbjMGPaVhkCfzY0RH6kb3nzNHj768HkWQ+I1Vyxy6/55No4tUW3VYB2FTQQL11aOz488V+403L9k2FcHdpaGm7ZNc2gw5sKZE0wtWa5b6HPFgVmSK/jTtRE3dwx7K8fAJmprmC1KPnV2jbmZksEoEsaRXZVjy8voQZJQ5aLMxV2IOQYs0S2g16s4tT5kbeTpdipm52ZkSW4zVzkq50rGqRIjbFTpp6ZO64h4JQjI+Pm9f/l3jDZHfPPX3I158RQXnj9BZ366DWeTBFDL6QH84xAeObvMr/7cf+aVr7qdpaUL8meTiE1DZ3aGhx59lHEYMr39QKvsjnrJCWXVZp1yC21NWv22NoGM09Clv5zUE9WksbJcbmJi0wfqumHPzh6DJrK6FZnpR2IMdIqSr7oWOh1LjeS1lCYy9omR9/R6Jb2O48SpdYiG7Yt9ZqZ7wmyzqRX9OGtbCm4Mmp2Td3TqHak6VbvkJiEL9caTBGOrSBNN2jTKuVN2mAJvBTuU882TkaA9DSBzZdGq85wu7fNi3FknKJTgZfiT+zQlPsjY0rSBVT5Eut2Kcd1owuEkAz77wGSFUMgqIoH1CuKzVjqG/Jdb6ySLPAj7vdHFJWq6kxYzTHwBTljy1hj1kegsrygpy6Jd3sUY9CWQ261TaZQsGvhu0uTvyRebU4WCzLyY/I/ORK3IDr33itOw+nsCIUqb6DU4SYCJk62RV9WBy8ZJI1p25wqNtg3ECDsWeszOlDQ+cOHCgF0LHbqdkl4lgLNuCWeXRvzKP69jOl0WZyoCjk6VqIMszOZnCpyznFodM66jpjAa5V/ZVgKYWg2fVjtOLmqs0DrTRZckKeFKy+D8iPMnL/CFb3glX3jbTp74+Hv5+lftYH3oWVodQ0osbdXc+1LgI89Fnjhj2Kx6VPu2UR1Y5Pxm5NiG5bELkfc+vM7fPTbg+l0Vj76cuPfZDdxolWMP3ctTn/prvuLuS/nWr/8i7rv/s/zhH71HXi7kGcDknAmrLvRCFFSuzJorba0rgl6U7eGAhHmBEnq1yMgoYOdKIiqY0IVib3qa66+/nudObPCKPX1i0/Dhex+g19R0xyN6HUPXNwyTp6oMX7K9zzft7nPblGUfkat7JR9aG7H/ilk2N2uee26dr3r15Xzll93BnpsuZ8+NBwVb7QydyhBt5JKqwyXWcXbkuWGf5SNjeGGcMMGzs1ty3WyXaxanicby0rDmT5cH9IC9KkipbGJ3YXlgdZMHtzyvv347J1ZqFjqObaVlKwR8TDQRfIokJ1V2QhAmyQi2f7pTgA9sRRg0gal+xfZt82IqC0EKpKIEHAHf7jMT8uy5osQrykR+TjA93eeJp5/n3+77FK+58wpec91+ztz3AN35Dn2bmC8MZQr0y4LTG56/3gx8/Og5fuR7vp2v/MovZfnCskpLhWRdOMtwNOLe+z7OzPYZelOd7FrCVprn7YyaQvOCP6cvqhqSiRDGqLLMJNPmw2dydgri/wohsDZoICXmZ6TKXl6vmeoaBgPP5kZiRxcijqLjMN4T6lr2iNGQ/Jj5qYpzKyMInsX5HqVOaISm4BiNR8LtsxrelIwm/+nUQKkRIYkvTsQm2dJgNFbakUVMKInaFa4t6A3yedhW1zsB2OYdtYSUBT37JvsWV8gYvG7q9nxOgFfFF8hnhu5ljO536nEjY2gmFF4xNTeTBEQ1JBoj3Dv5g70oBgoNXc+b9lwFWGtVimdbL4YsN+WgCE0z8VCEJBePkmKDIikkuraUm1ERHcFHVV4pythocJFNms6VWlmnzHU1HlVVDJkWDImqI615DLHl7KOekBzJm1IS1r1V2qTJRjZRNUWf9Hu0LUpkNG7YNtelV1jWVgecvTBk20xJvyNKi8pBYeG/f3SVXYsl3/aqPsPBgAdf2MLGgrVNT+kSU5VhNPIMo205VSKfznukQrs0ZasENQ1dbMzLu6o2/90SxonTz5xnx+4FvvWLb+XJT/0js1Wk4ywbg5rR0PP02YYPv2B4eb1g/tAiB1+1n/137GHnNQvsu3EnO6/dzuK2aX7ll3+Bt7zti3ns3AZ/fniTvXMVj5+y3P/imG6sOfv8c7z4+Y/wHV96C3fdfiN/9O738tDDh5menhYjJ+oS1/yAqKMHUg7j0RchNNo1Bt2JCDLEey+8oaD/Xkd0Qh6WhXbj1cAYI6PBiLvuuo1zGzAaNrz19l381v/+c77p+3+BfzyxxQtbkbVxYrw1ZmtrxLFBzXJynB9GHjuzwc8/fZr1fSU37+/z9586x+27ZnnlnZezPvbEuqZanGNuzwLOyrNhUsIXcE93iudODanrxCuv7PBwkfi7Idy3FXh2nDg3DgyS4VhjuKdX8OZ+wUo0nB7Cka3In5/c5C/Wa772NdsxPvDMqU1umCqYriw+ExD0RZdEOEMTwQTxDtUhYFKiKiwXmsDIe/rdHr1+H9/4NjwtBt+6ikWmq/wjfXdMygpMGQk5V/DeP3sfw60BX/22O6hOLzM4dZ5+VfLC0PPwZuDIEB5da/jLDXhoeYMf/t5v5gd/4Ds5e/aMFIc+iH8syF7gyLGXee65F5nfvUBTZzyR0TApnQ7o9xxyVkgKbSSroE1SO2WIUfNuMiqiNQTLzsQWjs1kSNbSLyzTXcfnXhyzubbBuUHFKFgKIqONmrL2hFqwH2WnwJSSJb44W7F0YQQNzM506PdLRrXsSJqmoVI2mEAV5fCXcbNc7tkekBTzJOKc0CrKcgJoLoUlVkG6lxC8fiZSIIoSKrWR3Dlp0PugLn2oOhXGTQCISUdMWWjkdcxlXTHhytW1UEHaKPOEKSQmsqiKNgpXzINJVGYxElJoL/yijS5NsnRLmkqYfRKFzvGqqmA0GmFt0O5EDC1FUQpMLMF4JDGwWRIWNWcjI1CiQQi8WUmhIUYig8t4Y0GRpEYVDXYibzPafmOsosgTKcqBYwo7kSFbizN62SgOPsWGaCw2TZj5Vi8ruRx9m3JYe01ojLRQwMWFPqZTsbq5wubAs3dHj14ZqbolU134o3svsDhj+dKbpllpDG+5rstDJyNHztSc3xhTWUO3NJRVycZo1OaXhxAU1xChtFicRE+qtNi6gpA8pZNWuPENZdWRxZwP2NKyemSDtZUNvvYtb2S+Oc7Tp85isWzWkc3NhifPBR48Df3pDpfduoeyW2DabGXDaDRidtcs5zbO8Zd/+zf8wk//JG9+49380H95F3/52BJfccM2HjoygDTmNZf3GGytEpae4Ae/8+1834//Nr//f/6U3/1fv9ruPsjYkYzjt4p5z9ZdZTgRZb/ThLF8f1r5SMSniCokhKrBlWWbrVAVpcQNKIZk//69uN4cj72wwmtu2sU73n45Tx3b4LEtx8PeUIwDsx1Hp3CCHCESHdjtXa64apbZFPjLj5/hzOkNvu+tB7GdknKrJjZyYM/smmdrc5NmHOlaxxMbQxbKPm/uzPD3Jze5anGKOy+ZggKGybI8hpMDSRbsxsTjjWfcBKrSYSsLhWH7ZYt827aK8xdGfPDhJWYLx5vnu2yMPIORVwkreM2nMVZMd1d1HUdrz3TlJB6WxJaR0KXFxTm6nQ5YQWI4NQeKoq+UzBtrsVao1zEKMTslUSzOL8zx2GPP8C/3fYqbrjnEG199Nef/4sO4whCaADHx5Ei69ONbnqe2ArNVxdxcn8efeIJLDx6i8YmmFs9JRDhM9/7b/Qz9mM50Rzr6IM+tBEsxYbepyTc0UQjMKvC52C4gO5EcSKYzepLEFmieBSmwPpYx11TPsX3a8czRDf6fj3ifuHERxkPP/A6RPg8HnqE3TFU9upWlcJa56ZLnzg+pNwb0+wVzcz1OnFmjKsuWx5YdfyGrQ50YIlvHtlb3TgGuyViZ1KlYyetl4azD63ti1KleFE6nPo6k0RbWiT8jT1mSejcgh/fRfgZRhQjRR0KUy6F0ZSvFdYUjetkxpiCMwqKQd80aFS2pajURWmVfoUDcQtmG/z+q/jvatuwsz8SfOecKO5x4z7k5Vs4qlUpSSSghhCRkwEQbDLYbN8Ntu7txwsPuxthtg9vudrex+2fjhLHBAhthRBQCSUgoZ1VS5aob6uZ7cthprTXD74/vW+uUGYOBgKp7z9l7rTm/8L7Pm0V1c7YjHNTU1fogUhJ1RF03B6EjxuntJtlDdS2BUjGJYayV4LpCfkgtp7ob1Do9xDEa5B70YMiYTitcLumAKPIkU524bw5kvSjfBSOyuKjVPG31riHzzirNVUOYTA4uSdUT8JrOFvTLzqR1U9VElufkmWM8mbC8VEJm2duvqWrPykLJ0mLJoMj43PNTbk4jf/ldS+zNwIwmnDo25PTpgo89ucs9xyAzlkG/pK4D22Ovv3/qqJ1Z7mh8oOwVhJjINJjGR6/KmaaTR7d7KGss9aRh49Udzpw6xrfcfYiXn/2yVOkhMtqreHkr8Y1bMFjscfbNpyRWswoEIjIJysBE6lnN6plVnr74HN/3o3+Wv/eTP8lvf/A/8Jf/xv/O7zz3Et//wApPXa2o633eef8Cl16+wJ2PnuJ7v+vb+PXf/hh//Okv8r5v/1a2t3fIipIYBRUuaqvULfeSfi8mSqXW5SGgXSCWuqm7UZ6whAo1rDrNukoHWvRBj7WLW0ymM14a19z69HlOHV7kjfcf5v0rcywfWWQ8q7E6VvEN7I9nTMZ7bKxPeeXCHs9e3eMuk5hzls21fUyKhCZhY6ScG1APBsRKkhVn3nA8Gv7bzR3+0ulD1CnyqVtjXt0tWOkXnF4tuGfVUTgnB29IFKUjGUdTC1E9hsj+uOELL2zz1Kv7LITEX152zFvJP0lOwtd8agO2bLu9oOccxgb6Tv6OEAIbjez0jh09wvzckN3dPX03ZP8UUxQTp+4EfWiUjWXwvlafjuTpfOhDv001mfD93/kB5rb22Lh4FWctTYLbho675zIKH/npcU1l4Kgz/P/+2b/n/+UXeP/73s5P/MRf4uTxY+ztiwBmd2+XP/qjT3PszBHyXo6fNRJDm2SH0EJb0dTTjgiR0EA0VOquCiFjDsB/OvEgoiqm2B3qkzrim8igJ0qi972u5OFzPf6P35xyLAusLkBvkHHj5ojCRa7cDMwN5slUsTnfB+8NIRr6vYy5Qd71C3LW6NhBg/Vsbrscjc7xrbvcvJC9h+0gqfrzW4Q23aFkgsZhBOmyoXvvrea/x3ZH60NXRMcksFqijLjEexNIJutSRjOFyMZ2v5RMl7uTZZko/oI/oA513LIkkR656WwVbc6StZYMI3rx1gwo3CVHNa2EDmtlVpfpyCfFQHIqwXUWX3vlsTSdASx2nUsjSYIxdeqcljHVek1S69HQD8s5h022C4KS29hSzeouIMW1iqpCgIHEpKouDvYsGkwUdTTXKiea0Mg4KIjL02CxDnztBQKoS3tM6tReIQQOLZSQPGvrYwyRYQGlc8zNlXzk8S1+8ntX2assj941x1efWcfFHm+5c4HaG/7rH9/g5CGHNQmKHG8qYqiJwQL6d/gg2fL6xcQETimpMXo6+YeOQ0PjKfsl1VbNztoef+JdDzPducrW5j5lv8R72Jwmnt409PoFpx89qVJtvVjbJyVFzck2BDzHbj/Jzq0dfuKnfpo/+f7v4B///f+Nv/eP/i9+54Xz/LmHD/HcjSnJjnj3Q8uM1y/xfe97lI9/5qv8zu/8Ae9+9zsUxyC0z9D+bqHB0godUrdIT1EPixjYHu0zGMyRZ46FhaG0z8YQmoYQZWfS7/cwwO7enqa7ybjx2vVrbG5s84Pf+Vbuv+MIv/+RL/Brf3Se+UHJUmlYPNTHuYw6RMbTmqqOjCeepgoctYbvX+zzjsU+//t4gycvbfD2L79AeduR7jIP6zvSqTrhx64UGUWCr+yOeP98yX2F5cm9mlf3Z7y4PeNJawjJi4dLFXTJJvCmS59LdWA5h+8dOt5/51Hmjy5z9ZmLuAhNjBS5CEVCTORWhS3Gsl3LIVI4kYHb3LHeVMQQOXb0sIxBUiJXZpgPlS5AE5nVfIoUMFoht1jyxfk5Hn/8m3zik5/l/tuP8z3vfz3r/+X35XErHNQJTyK3jo/vep6rIv/3W8/w6Pe/k9oaPvO5Z/hnH/wUX/7qk/zD/+Nv8c63vw1r4bc+8nGuXr/GfW+7j6Qxqim1MdJSyAljPWpHksnzjumMgqgkuY06bs1zmERSL9zBFEW8EzMPTZPoF5bDSwX7TeLdbz3Bq+OS//ax89z3iMN7w2hnxullxwvrife8dUhUQOHCIGd3f5/pNNCfK5grMw3SkwNbkW4qNJK9HWp2TSpzdc4oOFG6itdSy1tYqEANXRdu572XxEE1CTvF2YtaVP68vMxfE9kASUkGtjWO5tlBga+THfPaaGwV47SBVwB108g4P8RO6NTUTacmi20IlU4I2t1TZq2h8XVHym2aGt8EivKg3TH6B5su4PRgwZLpQ9BeCHkptv0uO1clua0rOelcr0Md64JZsPBBlV2hU2B5H5hpUD2tCUaXTL5udFeBcrYUj6LmmKaqu8ujqeuDoJaO72c6v0KWO5rad92RNeIGr6qaoshYWupD2WNUifx2cb6krhpeXqt410NDBjbRm3MQAlujQDYMbO1F3v76wzx5ecTO5gRLogqBSSOqKmVuykLRSXiSdbbLiNJCQHOJFZuv+RmZOly3r+2ysNjnjlPLfO2Jr3D/kYL9aWBaeV7cgv1R5PbHTlAMM6r9mYDTtPoE6SzzTNVfIeGrmoWVOW5/6A5+/4//iOdeeoX3vOddXLl6gz+6MOLHHlniEy9PMc/s8IHBkLN3zPPed38Lv/8Hn+L551/ivnvuZFxNSKrxT16c0D7KSKWd+UZlMIWUCEEMVLmzzA0GvPzKRb72jW9y4dKr7O3tUteSt37qxDG+9V1v4y2PvZnZbEJdSWf56qs3aLznfW97HT/yY3+S733wNh7/Nx9k99hJRodXWHv6RerJjAzL0DnmipyFfp9la1hwhmAgZZaH53M+tzPlAy/eYHVjh95inzBrmK3vYaLBh8RGEziVO4JG1O6nSBEj3zrIITNM68hegJ2Us1sHxjFRR/E2DHPLQmFYtDlLzrBkInPHFynuPMH8ow+y9tI1mmkte4oopkFjDI1PBCvDhMYa8IaetfRyS2Msu8niEhw/eVzC4TRBEgN53mJwcsXjq1w3iTkzqMLSWssv/IdfYbw74kd/4gPM3bzFzecvUg5LKh9pfKLMIfeeT40b/sRCyZubCv/4c6zcfzv/y1/4Dt7z7kf58b/zi/zE3/j7/Muf+1ne9a5v4UO//pssHV2iGBbyfqkIJms5btptHKiulJ+nUa4pJcX1qEQfEa0IR0/GXyYJVl24eQZTeWo12lpjGGaGWROp3Dzf/kbLJ762zqWdPQ7tzVgeJC5sBJ5ZM3xPz4l/LEaGhRhsd/dGHDo6x9LyQOGPBVkmMlmj4/dWeITVbsTJWZKI0n3UnmCE62VV/NGemSlJQFNrJwghkOcZhXYwUUdlre0ghNilGxprVMklF0bSji7q2L4bq0XxAaWUiOqvE1VsIJlIiip8auMX8qzjZ3WQW1W5uUz+3bqqIQSJtJUWTPK1c+cIJnU7DFkuJ8qiUCxJ6lzjoCFATaOb/ERTS6ucl052EkFc4Zry1GmmrVYNkgEsY65MH3a9jNWfIotq8Yy0o3NJHhSJHNqtGLJc9hZEsJnwaqJq5Z1Ki+WSAOvyDo1iFKvgnCy3IeH1w5RqwbMwV0KI7I8qcmM4tJDz9Rc3efnqiB95yxG2x5G33bfIxZe3yGPgM0+PONyzZDayuVezPFdqC+uoYuouiNxpxd62wIaudZRc9yjGOpULiuRXXcyzwO76Pg/dc5rcGPZ39uDIEonA2jhx/lZg6cQ8/ZU+9ajGOEsWRSVVNxW5zShtyayekmUFRItxQdD+hePuh+9i7fo6v/X7v8eJo4d49dXrPLXe8IMPD/nwU3s8cbXhh9/m+Y5vfTu/95FP8JnPfp6HHrxX0CwGQpKDwJKEnxa8ykmFnRRCoigyhsuL+Bj43d//OF/6yuM8/+J5fPIMlvr0F/tCug2Bi0+/ykc/+Sn+wo/8MD/2536YURpjkuHWrVvM9UvOnF5lb2MTM5txtm84+ua7WfyWx7i+u4lf36EYFMxqT5OgwTD1qSuQSJFvWejzB1sTfuXGPn/eJBZ2Z9S+Yd1bvjmteaZKbM0S9/QiWzHxJwqHD4lhmcvSVN3rp0o4EQ30nOCxPTgnhY/JRFzXhEAwht5bH4b1TfzOvtCTrXwuMXLQDSPKwAy4UTfMZUZjYi0bdWBzFhn2C46fOC4iDEW3txkWLss0atVpEUenhmu859jRw/zmb/0BX/ryN3jrI3fwve9+A5f+xa8QjXi2mNVYA3nZ44Vx4MIMvvtIwTgExk++zM6TL7P7xIuc/cH384v//H/mh/7K/8e//8Vf5cbGBpdevcr977hPFZ4a1tbF0OqI22o2kJPxkEnqXUn+YJ9m0Hya1BWAykrtDuJuEWygbhDKhc0ZzBXs780Ybe1SWMe7Hlrg+ae3eVPy7I09v/es5f5TfT7zzC4Pnc7o5ZblpYK6DvgEmMD8wHWU6xbr1Ap52jFWSKkb32d5RlPXolwyKD0hdGgko54ua+x/R0F2LnZy/e7cVcFMu5e2epm0/7md1kiSoNIfXoN790l2NJmzkDndScqFYI0UGyamDoHju+I+gJHRca48r/Z37fVK2YHId2M7LXBscQJB+CqCkgjU6kqXpVuSm0lTrjqPRnuYO0OsVZ3Quhl1K+ycxUb5n16X1dbpLZ1CF8faGhkJIlONChhzRqVyRtq9Fm7WNI38HCl0irGosEVJMEyKzNAQ+izr5nktfM4o/C+ESJ4XcnjnOXkemevnkDl2divmhiXDvuUTj29w91HJLj683KepGvb2xszlhkM9z8ef3uV/WOkxGnlOLIgJcjSuqL2MoPJ+QVSpc11JsJd4VTxFLgu1ZK0o2HJZgmYuw+ilWO/XjHZnvP7eM0z2x2Q0VLWMpM5vROpkOX5uCYJ6ZzSnOXoxVCZdVBd5D++FsFlXqopTH9Dx244z3huzf3mPQV7wxOU97lic44ceXeRDT2xy9/0v8853fSd33H47n/vsl/ixP/en1L/iKHNNt2w17BVahAR6vYL5Xo/dvT1+87c/yq/9+u9w5eoV5g/Pc+y+w/QX5OJo/zs0DWWvYLxX8e//8wc5ffIE73z7Y0ymI26trVMUjoVBKQFau3vsjRuG0wn51jobuyOaJmInM0xCctI1y6FwcmBXKXFmWPAXjy/ywfUR/+rGmNcNC2bGcH5SEULipIG7+pY1It95qMdCZojWYJMht5oamSXqGKjDAbcp+ISNSQF6UonWtae/OGCwMGQ0mrB1/irN/oRs0Ac8uRPDZQyyFA2Iq3w5t5QGyX0JkfUqslvVzPf7nDhxjLqRfOyYBIUT1CTYGvGsPuMSsVBxeHWFz3/+K/yzn/vX9Ij8nZ/8QSZfepqdy9cZLg7YnQR+ew/61vLmuuFxn9HLco73HTaCsQXJGfaeu8j5f/8h7vmff5T/7S9/Fz/1cx/ml3/hVzl8fIW8lErWaEfVshMlXdB0AWVRg7CImnSq77rNnIAXrcQTJCtdq80tsRacjs10R+CVl2YSTeXJjDj2r4wbknUE53jjHT1eeaFgtFfx+es5C4tz/Pi7hvyLP55xa7vi7jMLFFbUVrt7M8gzFuZ7OI3B9iEotThTnEvo7AAJubiMlYsxy3N517P2XQ60joWsjW/Q/aC4yGW03jrOve4t2vwiwe/LJZmpp6Q1FJLo3OIiJJLLpSVaeB907BYOUgg1IMuojcC0cjh1w8dkCERRrpp2/CiCCussGTZ1t1ArHWzjDWNsZaOy7I0x4HV5k+VOU50kwS8vcjl4Mplbt8vhNvGuO6yNJRrZUYQo4yCvxpu6biTTt5WnatUd9KEQeqg9MBKq4srrPydhNBxIR410LRIQpdIw2nZNmfoxkiJkRY5JkoaXUuwMU03jyRzSgcTIxuaIub5gl5+5uM877xkwmkROHXbcuLpPURhGdeTtd/b45PnENy5OZJ9hIr6qQReSVm4yknWKWVHoms6Gm3ZJZRzRiDPVZZnEXRqDMY7R9gzr4HX3nuP5555l2M+xzjCqA5e3I8OVIeVCIXsiwESrCBiZfda+7vZWEgpUiQ6/KDr8RV03DBeH9O/vsfH8Gns3PS9teE6ueP7E/SUf/tg3+JZ3vpv3fttj/Otf/K+8+OIFvuWtb2JtY4v90ZimaZhMp/T7JadOnOhYWJevXuezn/4CH/mDT3Jt/SarJ5e5/533kRUi4QxeQsu895gK8swyG03pL/Q5eudRfulXfo23v+0N9Hslo9GEXi+j3yvIigI/qbBk+Cs32SszwnTWzYXzPMdYeVGiD/gkxUNuZHn93sN9bhtY/mBrxjennjlruC83rPQz5nPLSj/jZG7p545JE+T5TabDZ5ssEpuMnlPyurX4XEYu1kIIhmAiZZFRDEqqzW02n3iF6fYeWBmh+CZo+JMjmkSI8sD7BAOdleROJOgb0TBpAkeOLLNyaFk6Kg4IxjFJ5LFk10ftrEWhs7ywyEc++nH+7t//v4lV4P/6qR/i/kMDnv6Xfwy9gkklvoL7eo4vVI5/tSPR10MMExPJbCRvY6B7fSaX17j2G3/AO77j3Zw7dYjNnZojRxeJenEEIrl1OvJpLxKhUsg7L670pKFhUfO8k47Dg/oc2kVg9BGcJXk9S3SE5auId2BLkaOWmUZTjPdo8h7zRcO9xzIev1rzzCb87fc5XPCcGAYu3qq585in0J1qo4XwcFDQNEFjtqXzKIpco7Slq3utKMBZS64sslZxlRQfZHX0JPtaEUA0dYPgrmy3K26aNsXTdAV0yx5smlrOYifBZL7xlEWu6BLtTPQzLZwoptDvIfpAXhQdvLKqah2/vSYdFSMFrfdSfBkZpTuFOIakXXErKbMKRms1cmWmnUYbtJ6SAMIs2CxTxIls6hsveRrSyzjFFbtubyJLYdO1onleMJ3OOhS4s9KmOufkB1QrfdKlT665vu0iyljtGHSH4jTx0Cv7RzoKFZIqZqWtvJxxXbaIeAvkn6+rqpMbC8JF9j5JWVHDQQEhMKsDc6Vha3uMj7CqeHbjIzt7FdNpJLeeWXJ84NEhH/ryHjc3przxzgVcnjPdrWh86hLJoiYmtoqMmCI2WoGnZZIeZ1F/jl5qSZPDmknD8tIcRw4t8pXtbRYyR4qwPkrsV4ajdw3JjKMKNU5JrFlrWtSLKMuks8GCRQ4weTFC9/vH4DHWcfjB49TbU1641XDf6oy7z61wbnOdD//+5/jud7+J//Qrv8Wv/Npv8nt/+HFevXCVummog6eqajJruP++e3jowft49psv8OQ3n2Vvts/K6cPc87a7yQpLbELn2EVTBo2+lD5EMpsx2684cvIwL3zuBT78W3/IY295jO2dLU1Zsxib4yczsl7O6Oo6O+dvkqxgM6xzjGaNCja8PF9NZK6XizWh8cwyy10LPe5bKNlpIqM6URuwROato58ZqhCZ1V6d8rJYbwuWqo66vBQj3WRaY4Aydx2aJsVEcpbp5j5Xf++LRM1UzjPHqJGhn2mzF6oonqhougVmCDDM5WC5Oqnx3nPs+FEWFxYZTcdd1kMIQQyIyRw4irOcuq5Ynl/iY5/4DD/5t3+GQdHjH/6N7+ZP/+l38ewvfJjZtKY/KDEpkfKMB23kDfOGq2PHz+8mtnS5nXOQWGpDIpZ9ZlduUtxaJ0QBN1599irHs1MMV4bQqPs8anSDs6TUynltF0pnLDoOUiG4CgJEdi8VtdFOxavfyyW5YAzSjQQf8LW4pvuFYzRL/PFXN3ndGQcpUhD54nXD9765z5HVHhu7gbNHCi6vB/JC8tVNSlQNkOfMDUvy3GnKoFTjdV2Lu7woxesWojjGlethMHoGyrhNjHtRVwAiM259ae2+2JqWfNsW7kZGynmmO4lGjbdWx1deIayufW0kiTTX0L8QSFFah6iTBexBFkjQTiKqGquOKmTSfXGMUBSu8xbZ9unWIjxzuryK7cwsEzlu6LAC6Avsu6QsmVWqCiJJetUBekL+HFl8mi67vF1S11WtTkk1qfioOeS+Q1S3WHeXyULZx9YFrwZyPQxn04qizDtfSUuaTK1ZKrXsKFkKNVWQOMzkOvQyMR0AytQ0k2lQklQUGXkmlxgYRhPP0kKPKMshejZR9AumtWdc1RRGTH/b+w0P3b/AO++Ff3Nxn56NUq3IZhTXtplaYbTYmNaJ2jmrELR+iqL2iEj7nkWY7U04eXyV5bkeuztbLPQFMbE+tXgb6S/18ZUnzzORLau0tqk9Nssxisa34n7CGCkIfNCcBw5Qz9F7in6Pw/cf5dJXL/H0dcPhxQn3Hp/jC1/9Co8+9hZuO32SZ156keWzi+TLJVlWUtgB8waqac3jLz/L57/2NQZLQxbPLnLqxFkJy5lVNFM5GJx1ioYIAgA0Ei/Q+CDOgiiV7PG7TvLzP/+f+Pf/4VclC8Y6prNASDnJlYSgO7w8p5pWxGjIctkpOZczTZ5Z0zDMc5U2tlnXiXHtsQZ6RvJbRlWkCRacYRISjjaNTg5642Qh2dSxq0K9j/hgmHnfZdaXpUh7U1QoYIRgnByA1soLPK3IrLyg3kcJZtNhv3GGqpY8HUsiRLg6FfTPudvOkhcFYW8P5zJ9R1SwkryOfiO+buj3+rz04iv8/X/w/3Dn6aP8y5/9szz04O3sX7/J9vo6//xWYH4+8CPzhlUb2LaJfg07yXGzDjzUh9fPWSYhEeqaaDRhL5ODb+faOjc393nvux4h1A0f/eST3PvOuxmsDDthS0rxYKLQvrNaGLYTiJSCHLzWaAKfxDmb5DXeQUc7sSOwdUvnYFrFkOvUTV96peKPn0v82FsynrhuuPd4wd2Lno29hiblLAwM59c81TSQGcHp102E2NAr6MCy1megAVdyZoTunTXGYpPACQPCrTPW0NSNHMJFriNEwY8YI2mDTdO8JtLCd6Rcg5gxJS3Wd6rVEIOO9Yz6qw4ylZyG+bUdT4t1at/nIpeC3HTdIBSZowlexUzxIPXTyMI/KzPp9IzFWpSHJesrSZxS3Hk7E2tHTC2nP88yQbQ71yEFjH55XpMDE4mmlnl9a9MPQcZe3h+knaGa5xQTWSHVdWYdeVHoBygMGWGumO4LarlUKUacsutFdnwA92rZXkZDqJyVSrCpZUdi1M3qm6gfcuwkceJx0YNHnc4xBIrMkpkAdcNoXDPXzxhNKnKENmuQL9ommBtYNqaWs8f62H6fx+4/xPHVkjyTWXblE7WShUWpZiCkA8xyy6xRo6B0ZKa71I0VpZvBUdUNd547zvaV56hnUyKG5Cxbk0Q5V5IPShKCJ2hba6/yzaR6caF9J8GGq4SwHT0mTcHztVTY9agiX+rTm+/zzVs1F9am9Ac9Vuw+l69e48F77wAMJ+87xdLxBQbLffJhTt7PGS4NOf3AKe5+293c9vqzzK/M0UwrYlWrA9bq4Sv5L1nuNLDM45uGzGQdYaCpG/qH+tz+tjs59vAJirJHVUc2difkeZ/86JEuBnhSB0Y+slN5fExkznYXZJ4pf0gPIx8F5496lKYBdqtIkwx1apNADV7Nz41XTI5GKSdjMSZr00KF4UXXUOEchCbqstd0VaRzFh8EphfVjFnVHoNcIkl/rjZiICTIc2FmbQYJEDt54oQiTqyi0UMXCZ0595qALul+/79/+QtM90f8o//9h3notqO8+tsfY/bKK5x416PcfnaVb9wc80/XG677yIKCTz+y0zBpIu9ZKuTqVMqDcBEjoW7IjeH8pTU2t/b5ljc9zL/8Fz/Ft73jfl784stMtyeKLlcpq14WRicQLXjzIF5XZflBUwrRiFuV+KbQBo9JZnebOyTvksXmmW7aRen0t//Ucd5wZ49f+KJnu8l5zz2O/SpS1579cU2Rw97Y0wSlHSWYTWYQaojxAImUaVQEigfxjf7cLVncdsqlSKSaVjgrGTPOZhgtklyWk7lccsudIzMO6xxZJhilmA7YYKlDtljd0ebCwdI00nZvLUBY+TS9EjfaMKmoy/i6rklGPHimA81GJUDErsNq6koMp0Cog8ZJe0171fiIFrcu8i8ZHVWzGS5XyaV+KE0dupzqdtE+m83Ut2C6PPW8V3S2eZRW2SWN6S+OXgI2s/pLewWitemIcmOLXC11mJPOmZpkBCFMfNMpGKw7SETEHrCwmrrR/cuBmaLdy7T/OVMpK8lKRRmCAsjk3yt7JUSYVh5n5WLq5fL3VNMZ0/2KhZ7jG9cC02g5uTIUr0IyNEEeNWMSTUjEJLrusuzho/zuRV5o7om6a8OBKi1GQdcbqwwgK5jnLHPksWL9+sty6adISIa17YbBUq/LQ84LmX9bhRC2rW0MgaJXEhHlRV6WGl1su/wCH72yzAwmk8tr7tgCJHjqesXeuOHonGXt0gXuuPMc9f4Mv18TKw8+YpI8sCF6fC2faT2rNUSn3Y1ZVesJK0j2IHJ5O2e0+4uKcEmdAao3LJVMbAkmsba5R54XlKtH9GCWw6UOhkBiVkdFw0jxkju5FPqFo5fnwi9S3lg0hmQNHqi11Kp9whr5uYOXLtJa9xr0taP2kZQEN585Q54ZDIFeoXkQun+IKojPnNPY4pZ9JeZdUcjIhRdTwhkRkBTWkRmDS4a1WWBt0jDXLzhz9iR1XQveX+fuxtju7/SNjCnn54c8+dQzfPozX+T973mEt9xzmm/+41/k8u98iYu/+kl45iJ//298D//kL307/SzjH1wY8eRu4KXG8nQVedOc45H5nF0vB2+eu25OL4mMGd+8uEZ/fo6lL3+B6de/wT/92Z/g7lOrXHnhOs6aLorAWuHSJX3vWxptF12tcdhGw8ParkouCtcFLhkjI6wWxpope8/Xnno6w2khtD0zfM+7jvCe+xMPHaqY71lGE6FfTMc19XjGZOaZ1rL3bAsJnKXXL1WRCb5uKPJMlVHS5WS5EDvaAjtGGcGnkOT904wcKWpbWHtUvI+qyIzaEUJQ8q/pyLzSORul4urkJAT5M1psSgzdiNsH4eo5zXFqOwkpSuQcb6qmK/BdlqnyNgoxwTrtZMEpQNIo8dmoPkRzmcwB5Ta1kZGuI8G2y6ksk0rYB0+MHq9L3VbNlFp4WOOF4kjsjDPtBySVltdqX4KfnB4G3e5BwA3YzGhHYHUhL19Ul2euFatTMjDdZeA6JUFeSHiV0Zf8IEVPPgGrFVSrPkjKw2nHcFkmEbkSomJVtmwY9HOWFsoOwTwcOobDgst7iU++6HnhVs2LF3eZbe2SQoNvEsOhOEVd64h3lsbXwqXpFBZ0l6KxIuuV2abpkAfS+YqPJkWDZcbe1j7zZU7tEybL8Dh6wxwTEya2YoiDJVmofRfc1fgG6zL9HIK2sEl5NxY0QjeFJF2bjywcGVB5y9bU88rNMScODUg71zl6dJVSk9gShtpHGbdlrvvdWsRDTHIoxyDhQTZr+VdO5YlaDCTpSGMQIYe1TkQBWsiMtypuO15w9nDJ+Yu3CNMpvSOrkPfwjTCZMsRYVwXx3/SKgkFhKaxjUOQ6W0dlmtD4qGPGSC9z9Jxh3lkWnWHBGhZdYj6HwkicajTy7xgSeabeJ3UHl3nO/KDE2tfgOLQbtKrKke9SpLtVLbuazAm3qg1USjGpyzgxLEUwcrNObE4rFuYGnDhxlNms6t4D16Z/It22sXLY5lnB5z/3JSajCe9596Nsf+JzbF1ZJ1+Yp8Zy/fNPcemXPsI7z67wy//kz/Hnv+cxPjw1/JvL+4RG1EdNEzCatx5i1EwXS1bm+GT4xrUtTg8LVne2ePnf/Ufmbq7zP/zZ72CyOaKeNsqmay8C28VQk2TBexAol/57D1pMWl2/Zmyle4MWMW5dRmjUIJtkwmEyR91EZrt77E0jx+ctxweeZuqJKWNxaXgQXGWECmF0kuyD/H65jhnlwpCwJpvZLumv5fUFHR213ipBj2i1bqSCb5qGoBOTdmEetPJvVwEt+FW4b3JJN03Q/YTrElptiwfScwQMmbW655Vc+1YJaaw0CKLuaroAvhRRj1ALeHUd0NJYWT0Ya2j05w3K2dFNj1REsqyWDb/Tlkb+AKMt5gHQzBiDK3Kcs9SzSqJuFakeO26VkiSto9/v6WFt9WZz6stwGp8r1UZVVRK4grgejXEawCK3ftM0okLSsY73XqmQmjFsBc3QarNdlouMWNHPwp6R9rlt7ySoPtdlui7l20pHZ7MheEyoIQaqWUNmYHWxYHWpZDILBA/jyvNbT074C29f4I7VnN97ZsJoawze45z4HWxmqatGHlDFxBiTJAsjBKJJ3c+upHch1iKXWVCHfZ7Jv29JVNOaGOHQfIFPasqyVtMkrcwuoyZG6uxRzJ5KJ9YX02Ri0AohCNtM89mLLFdpaKLU+W1v2MM7g837XNuNZEVOXm8wKByDXp+d9X1iErxEspLQJnWBwt80IyYET5a7TgRhDOL5wGrAjRhKYpBQohA8eZlhtfNzzrJ9bZ+3PzLH6+8c8MzL1zF+wmB1nvzYKqGlI0RRjJASsybinBxMwzKjXwq7bVJ5jLP0rWWoL/talXhm5PnqfuDT+4GPbFZ8ehT48sjzwiSyI5BVCs0l3k2wUQd6ZYbs7VM3KjNYgledvvpfmsbjveSPD8ocHwONfv9twJK1RsKrLFgjBdegdAx7so8Yz2qOHD3MkaOr3VjYmoPRijWZGsBycldQzSqeffZFjhxd5Z6jS4xevYIp+zSTmtKBHQyY7O7z8gc/wuhjX+Cvvv8+PviP/wx/9y9/B++88whf3hrzq2tTimRUlix7vNyKvPTrV3Z4Yn3Ce4/NsbI8JJAYf+OrLPdKjE/UE99dICCdsODKk54DpiMzJz0knc78jVbukvonl43LHTHQBZQZq4o1r4iUIMpGH+HwyWUmW5tsbjVgSrxPDGzAT6bUVcOskjOHOkDUKO8YwEmWRzvlaCMgQgjCCNPJDbpLk9hXKbyrulIYYdQz0MkBneiW2G0Hd5BY6DqMlLNC4rDugGjextp25kWVQrtMRpXim5NzXIyLMvHJM1ktiJEz78691lZfFFlHvGh3wSbJ2XqwQKczhMYQyVKSTsHlpsOStGYca0UiR3Zg1hHYWSAliV3FSVXrMpHlZVlGUzWv0TKbjg7pQ8RpII3T27KpPUWRy6gyJpqoSPOUlNja5nZo7gUyC7eZEcekSBK06lXVVyau9hY/3OZrO63moUXGGw11QtPF5DZ3xiilUgw7ZZ6RaTfjg4zeFgaGpaFlewKHlyz/7eu7vPlMwWP3zrPQgwvb8B8/uclf+o5cRgiTihgGRCNtqMsthZKPjW1wximSRUO5NB0xhYjJ6FpSou2Aa9YYmpiokmN13rK+NZJExkaqYWyinsniHMS7k2VOsfEZoVbNuJPDNRlZbjbRixnTWOpYS8BVSuImx5H3c8phSYgwi4a6NuT1Pn1bs3LoELfWtjl0YpnZpNIMiwPjl5gpITYylvPB6+8jn710g1IJhhAgyvcZGk9hc2nbYyLPYbK5z3hrTOGWeeTOef7Fb17g1o1bLJ1YIT9xgr3nL5AP+mQ2UgdPmQu2owkJrCUkRAGTWRZyUec9OW74wt6MF0cN+1H2C4WF0hgJC2sSFYlpEymc5WjheGSx5M2LQ4JN/Jsbu/zUqXmsExRNMm0xpOrApFRpHyRKGSsZ6s7RRN+NYGOS7izESOxGN/ICz+UZ+MTFScOsrjh37jT9XkE1m2GMo/a++zyb2GCTVZOZGGm3t/dYWppjqWfZXtvDYgTMGBJ5AlPmhGC59cx5dl6+zOoDt/O+N97Dd/6N9/OvfvWz/NJnX+Y9ywVnTRC2l0/c2K9oJpGP7jQslo73rpY00ZNbgUHeurJBMJD3cz1Q6TLGYxO6ZM72hmgztyU2+QD30ZJsWwRI9J6W4RGJJDXBOSN7otA0+CAZ6sVwQNM4vnlxxj1HM9a24exJy2Q8Y1DAC5szBmWf/iBjezfgTC7YkJA0Dljel+m00tGlXhZ6QTjNRkfp5uKMz7pltkGI0vrLd2OgqCF40Ucl3nrxKTmHVwahTB5kz4Lus1xmEUuGjsCiqGZtq9byIr8N2tW1PK7WnR4Uooh65Jz+LFZNY7UPyJJD4zxU+NEm1MrEMsvwaOuobWVdN6qsit0FJcmAEkJic9tFwVpNIfSNl3B3mwsjSx9iQYIIXiTX1qoJQRaCRjkrukQ2tsVXS0eCpvVZNdehkjiXyRxPfuEEaozJbN5lBRsn5iOR4eYY66iqCqdsq5SipKMr+TJoYpuzokKjnXcGj7M9VXgIJdRlhqwoOXtygauX9/nc+ZrhwPJ9b12hNzdgZ7zHD73nDL/sb/KfPrMppsq8oGln+LItlDGJIh2sVlptEqEKs+WLb6mcQdprn6J0UMkwntbsTw3nDmXkZck4OPJc3M8EWewLNkT14KHphAo2s6pG08u9yHDJUs/0EveKmSml8/BBNfApMVjsM13bx+PY2psyX4LZ2+T4iWNcfOJy5xFqfTtJHfBtRroY20RBYpzpjH2N91h7YCZ1CtI0ucyKGvUEWGPZurzHyVWLDYk7jhhu3rzFUy+d59tOLjO88xRrn7IMc8egjJimzWhAq3nZT4RgWLCOb44rfmdzwo6Hs/2c96/OcXKQcfKekywdXSLP5MWbzjwbGyOurO3x3LUdnriyxa/dnPCbt0b8xNkVUrB8bt/zvasle61HoPNeGBWZOB2NSPVYe09IiXEVuriCrN3ZWaMsI0Om/7nMEk1mubRf40Lk3vvuIjSiTMoz142bIxlOtfu+1lCupqaqJMSLKLL9phYqc6/U8LeUSMEQTE6wlutPvMSNJ15i9exh3pwSv1YWPL1Tc+pIyfa0YXvqWXCO56Phjyeev3JijmM9x860Jo+G4Ho8d/6S7NK0KvZ6WSaNlW53YagkOi8KvMZHiCM6/fej8CZABjTqD0Moy+2oJc8MJgjQ1SiYtLSJxaUeTbTsTgLDPGdh3nH+WmBl0fHyi57bT4razjlZwAcfIBMum9XdKCrgaSnhNsuE2+ZDR7XFSQ6OUfijU9hj0n1FXjoxO8fQ2SFSRxeWcVlsIllRdMa+pO8NXjoVknzLVv1DCe3mNAe9S/+0MgHwrQHc0OWsG80zCSF0MuiWfC3ufsHjx8aL7Dga7a7lUspEbmu1tZI5dVGWnY48pkjS0Hus6Sqptj1v6oMQpCxv//8iI8tc1p4jslhqqbfGCpBRWZHGHmi420WYAZwm9LW7AOOk4pCOx8k+wXul7Yqj2HYBKfpzkDSjoNbAI8khcM5pIaCKghQ1kTB06lUfA8Y6ZnUNviKGnNG4opkE9vdrvvUNq/yjp7YY1Yn/9V19yn7J3NDRc4lpzPmJ7zvNf/rYVb5xeUw9m+HsPHkpuHZhiHmwjtzmavxSg6Fqd1IIoJpsp8TU4D04hytk2d80nr0xlEcK7jjSxx25g9w9xXhnzMKJebko9Pf2sdGUMqk4TUzcenGd8fqIol8QCRy6a5Xh4pDURGzpsJTMqkruM1XjuMzQX+yzcXWbSeXYrx3LQ4ffu8mxo4fYXd+jHtfS5ZkW8Oa6bGtrLUF1j7LQbwgmYqwEboUQOxl2XdWURUEdAs4a8lzJzHVkvOf5vsdKsjjjrtML3Hm84Lc+/jW+9a0PsHTPWXorS1T7IzJrGeaZuvtFr+9yQfsvuJyPbU54blLx/qUhd/cy5geFkHAJHH3wFHNnj9PMhFFVDnvcFRJvN4Z6NGH9+cs8/vXz/NaFTf7tzX2WkuXxvYbvXCkUWtd2xKLIMrnTWGEpHAQUKSKBKtDFiobUxisfUGdDtAxKx6CwrM0ato1hYTjHHbffrtkUUeNOXQfltFZ8WllRyp+XORYW5ljfWMPXtZAaSNwcexZqy+rQ4KKY+DLnSD5BVpIVls3LW1iX088zpsayPa7Z8pEV5/jSJPHLe/DmhR7ffWLAfuMpMkOW99nL+zx96TpzS0OyvqRyWpd1u5OkE4mOzxUPIl9TfI201Vpim7DnBDFvnbDBpEsRCGlIAZcSvV6GrwJ15dmeBn7r02v8yUcdd54quXF5jzN3ZaxvNaQmsV8lLm1HfuAdJT7BtGokRthECPLMhBBwtsSoT0rIHLnsCHRf0Fb8YlNoR+RO3fEigIkKRbQq8nC64JbxrYzuohXZMoFuv2xVKNCq19p433YJHkPUxboU/LHNOrdWikljcaVMZ4KeoQcXnBTW0ccDA6NRYoOxlGVB3dSdGs5mTsbUbQKYLC4zqroiz7Ju52ERN6cE5FmNNdQlqEbNipIgSlfhZOHuw0GeR1XVB9ruEMlLYTs1TSNRlskf5H6YA5WU94pyR2aOmVMOlF5GYunXTHUNuJLwFtvtToyxEjpf1/Iiqrve2Vxkqm1oi28UOS4LpsxlZHnGtGrIXIlxOVnuKAorOvSmZnkAOzN4zwMly0sli4fn2Lq2zWR9l2o8purN8aPvO8PHHt+RlhXl+CvSvM1gaelETl+YkLS7MoCXAB4Tg14CkBnxnrgiFzxKE/HRsDqAleOHWFpe5srGOtEf7g7t3FmCEcqmnwVG2yMmmxN6teMNjz2Es4697TFPP/EK3H+E5AOb13fJ+zlHb1uRfAJ1yQef6C30ZL/kciZBZInV/gbHjjwoy7tZRTlXdO17tFEl3EnnxgnX07GFy+R3V0SDNYJ0yFyGU1261bGXcYbB3JC1yzeJdc3Dt/e5cqPm5tjTX1zg0198lqtXb7G6Os/wzAm2v/EMbjggc1HdzVL5xwg9a/jGfkXuDP/ziSUqExmHxGwywzpDbhy3Pv0UO8OXSI2n8ZEGQ144QgLXyzGzmgfKxH33HOKpfc8vv7rDpWnDpWng9p5lv/KYrA0QCzLEUCWLLC91wV/LZSJKR1USpkQjdnYKY6mj5DLkyXCzDmyMKu44c5J777+bsizIMkuW95hWotqZHw4ZT2eMxzUx1vgQGPQKDq0s8fzzL3Fjd8ypxTmMX+fpqeH0LJCMY3UIRSbPmqpgqWvpvqsgPqTUeCaNYWgsfzw1fHA78oah4c8est0YuhlXHLntHF+/uc2r1zY4ft9xnEMOZj38mnacGmQH1RY8wXspGFVqbzNd5mYOkySILMscvmo6dVICkgox8lwP7szQJMPth+HmVsWHvmg5PICF5QH9vuXqeuC+U44PPt5w+/E57jrZp24CUWMi8kza5/bC8M2BP85Zucysfe0ER5ftAo2iaRVbqU1qpYPCSiRtoSTiKAFbTrqvlo/V7ieCDxR5odYCkf8bLeSNLuJb82FZiDQYYzSCQ86cpm6wiBqTTCgbedFePLEDPbZFpnANRd0VW9Ns0oRY3VG28pMulTBzMiMXcFbq5K8JucHaRUpU1UqMoZuxSd5Cw2xWy6y/EWmvs69VTAWaulZDoWDkQ+M7/XTjfffPWm2vYozkGnaSZVmXcJfnGZnNOoVRy+C3zhxoyU1StLseQEmXcsF36WKND51Er82niK1eOsm/H1J7zIugv8wNv/HpW0ymDQ+fLCErsTFw88o2w9wwbgrKXsZ4Jlwrm+ciGNAxXFAsvekMQ6ZDvTjt8oyOcnLrNGdcFU1WVWRFxnTaMKmh8YlkLHm9xR23HWVvU7qAsl/ImMIHsryg2q9Ze/w6xVZk4B1nzh3l8Sde4ulvvsLm7j4P3HeO6tIew33H606d4WS+wMZz6/jKiwEhQTPzlIOcYETaujuVOXM9G7E0cBSuYLI1ks4nQVY4NYnKAepK1ymQjO5EUkoUw4J8LqfsleRFQdGTJW0KGgpmZG8RAzS7E8azhhPzkWHZ8Hd+4yQbvUe4ee0qv/uJrzK/OMehN9yDdUalow6fZIbdKxxlZpilyOm+5S1zBVt1TRXU2CrDTVIT2V8fc+P8LS5f3WZ3c0S5OyJu7rN3fZPti7fYuLrN/rRhv2p4fc/wf9y5TG7h+bGIMXCmW6wa7SYNEJD5dJE5YvAi342RXp512v0yk4rWqhooxMT8sCDLLZemkVkdOXZkiSuvvsrHP/Fp/vW//SV+6u/+DOfPX2B3Z48P/9bvcvXaDY4eOczCcMDi3BzzcwMefeRB6jrw5MtXGJ5eZiHBSg5fmDbsTAKXt2q2xp7ko2BAUqRqPNFYdmKiwrNIIjeGL4wj/2U38IEly48fNuQGbuxMRTacO8Lhw3zkG88TUmR4dF52eq3XRlVAURfkIuXWnWUr9TXthEILkBgwju5Cts51uRtJTcuRROkM+ETwkenME03Gj3/7EpfXan75K4k8S9xcazi6CE9eCzx+HY7N2zYo+wCSqKmmjUCwD7hT0ElwRf1lyTNV9CGeC9+IYCbp+x61ULXWUBaFouvVMKu4/hih8Q3T6VQhtOY1u5LXiGz0Qm39Yy0X0Frt4NVsbTqjZdK1QyC8BsrYghtbFiJteJkPXYywV6it1XF03fjuYhdxSpSDqZWX5Vmm7WDLgtfZnM7qqlndOZWds92BG8IByAuSfLHRKppdmFnGWopMNPdJtf7GSphTizWxCmDMTN49FK2jvPWYpI5rL7+k0pslx7ydTUYZg0hAFAc5421oin6IIj+OXYCNb0GOOlZrgqil2gTNVv31x0/t8vZ7+uQGFhZK9jb2Ge/PWOxbrl3d4dgwUNeeFL14H6LBmSThN0aWj8a2yJS2FQWTWW3DdXaprK+W1S9zV8f8kTnq7T28gZu7M25bKRmvXeHhU3fwUZszWt/j8KHjkGoaYc5x65mbPHzfHcwtzHHzxiYvvniFnb19iizj+s11HrjvNh588HZu3Nxme3+fO+88zez5ihsv3OK2R88wnTQqMzSYImMya5j6nGnlmYx2WDk7ZXl5gcn2hKN3Z4RGgn2sbee8UAwk9yRUnmoaWTw0YO/WHjdeuEWsPfMnlsmLjP3NEfMrQxZXhlQTT68sRB48bqhngXfd12d5rmBYbrO23+f4uTPsXDjBf/ntz/Aj3/8O5s8cw83Ni5Q6RPJMvBQ+CJrEWUMvwSTJQtx7MaVaDDfHNb+7PmZfK9uAHI6HcsM7Vua4b76kVqFC4yMZsG0Sx8qCNyyW3PARkiMa4U9ZHb8mI3j7FCT5rQmRwhrhTlkrcmAj1WqjYw6R31qSNxQORj7x4iywtDzP9Vtr/PRP/wz9ouTUqSN8/fEX+MpXH2d19RCZgQ996Ld55NHXM5tMWVtb56f+zl/l+77nffzSL/86n/riK/zAn30D8ytzfNu0YctHfne35m19y9hDmUWOzcEgM+TGUvrIpzamjKvAWgZ/FOGjY8MjA8t3zkUak4GN2CSy2RN33saXNid88ZmLHD59mHK+JNaRFLyMsKyTBDxN8GvNuzKJSNrtyq7OtlESRnh3SdVLEelKCQkn1mmaWaBwhsyCrwNVkGjfhoK/9N5D/L1fq3h2reIDdzp2p5H/+rTjz799yJXdjL1xzZFMECtGMUYkw3TmRRauZkWji+U8zzuKRdU0srj3SlBWiWwXRxtkpNs0DVmWdzioFNEuLHawxKKQ7j1zTkgdGrqXF5limERY1O5MjO5NUuogTq9RbAlYxarxEMXjy6hYgLGtIdy9JpjLpINU5WpWdYiTXK0bIYQDDG8MsvTu9XrKFVKNsM7lW3dyR+tNB4lVUUNeijzX5bdR+VkmC5yUDqp3RM0QlSzpXCa/AJEiy3T5Ix9YjAdfVouQzzL5wtolVqM3r7jgG+0iZGZpnHhN6qaWD1uML6KMoEvHlQyAzIhyAcgzp3Jh0eOnZAh1DaGiyCxFkbO9W7E39tx7PGd3FMgw7GyOxJg26OG3b/H1Z9YYOMkvzsscGxvBJmur7ZxUFz76ToLXxlCaSPfF6m1Llhei0rEyb51b7rMzFTLmxY0pvcKysTPmSD7mvrtPs/7yDbxmafSHBXvXdlnqzTGuaj75ya/y/AvnmUymDAcDXO4YDPq8/Mo1Pv+VZ3np4mWef+kin/38E/Tn+hS1Zf3VLYpeRvKBsixYWF1ge79ibxzwWEJo2N/c5Mypo4x2xsRGNfGaQBdCoj9Xsr8x4uYza+w8s8n+c1u8/NnzNFfH3Ld6jIdPn+NoU7Lz3Bpn7ALNxTE3X16XKryqyV1GPfYsO8+fe0vBZDdSzzxhdJWiKDh250M88+I1fuN3P8+R245THlogN4lBX17YWZv54izJWKr2BY+t8VDzWXLD0TzjTfND3jI35K39AW/olxTe8muXd/mNa7sEL62wc4bkDEXm2G0871/p8Z7FjIqES4B1KkXOSVi8crliQkUmloaDEW5KohqymRUfU0oSnmUTvcyyXUde2ZvRM4mH7j3LmdMnOHzsMDE5vvWdb+bRh+/jzMljPPy6+3no/jv58he+zMsvvEg12edv/e1/wCuvXOSxxx7m8Wdu8MTWLqv3reB95E8e6vG60vKpKvHlKvDSJPC1WzMu7Xlu7Hn+7Y0Rn9xvOJYZPlNFPjUOxKbhkRKCzboiyZhAagx7cwv8568+j3MZy7eviKAHMRPHEDvydlTCQlQCrJw3MhZ3anq1aqgkGZKPsqB2mhffqbjkMLZAaRMmBYxLjKvE3DCXOAfX50ceK1gfR6bB8otfbXjfG/q899EFmqZhv4JeX4QqmbXMzxWAZVb5zgTdxtn2yrKDzlqj3p3Gd5OC9sxrC4J2mtN6LFpnuFEEi3UCS8y0gG4X2VaJGyF6/edRT1SGdRm9skebI96ikNqk16ZpRMKfEnXlhaume6e2w2mqWvzY1oo3rFGptbOkNp78NSKnqJObBGSCVFBDl5IvjZGZbV03yLgs02xxuW/EIe31FnRK6Y1qotEvXT0HokSQ5aXNncpWxQAj+RcRE4SV1TL1U2gdlWhMrqUoi85i315qmcshBJqmxkQx52XaOSXFOLS3sLFGMihixMWkX6oTLbkGukg3EJS2LC9z03gSOXVImDInywzTSc3WfoOzidW5jIClqiLbmxX9fkZVN6yUlue2PK+uTRiWjv2J/N6Dfg5hqkl9sctwb7sM7xN5LpWLIA4yYuOJGcTQUCi7qa5qBksFr3pRR+1Hx7XNGYXLsPUm733rHTz+9CtsvHCLlQePkQNpFDAp8sILlyRfuRDFRou9MNZKAI73FFlGOZcxq2t8E3nHWx/mi48/iz8SGSz28U2gNyy5VUNjsw6zMBmPOHX8MF958gVCLUZJZy35UEaPo1tj9l7Y5J6zJ1k5ssRwOOCVl19lMOyxvrnLdFZz5vgRHrj7HPvjKffcc5qnn77A9fUbnDp9mOvPX2ZnVPEtpy3HFjwvXopcvZkxTBtsbO3TW76DQysr/Lv/8im+871v4tDZo2xurYlQz0hUbIixU/9IAZHIC7lc8kxUOMdswQ8cy5kRmfrAfh3xxnKqn/FgU3LDB8Y+sNpzQmzNRPqdnOW4E+XKrKm72N4kXJ+O36TYWVyWMdPM9FxfaK+hNkYNpYLESMwPSvp5xuVZYscnFkvDMy9eZnt3RN00ECPnTp1gMp2wO54QfWJ+rk+/P+C2204z7BXcWNvkn/zTf8XS4jzz831+6aMv8cifvp+559cJNvLuWHB8v+GbdeTFILkpX9yuWfeJvg388ELOkX5OzzmeGHn+YLfCKWW3tpY6BFZzRzPo8Y8+/QwXr2xy/L5jlAsFdSWqo0LHdG0iqbOyOLeZTAKsEsFTmzxoVYKq5lpjDXhI0auyzirwUy+b3NAvHEVPPCL7E0+eDM7PwDmOLWbcfdjyi19veOD2IT/4LYusbU6Z7wW2xlH9QoaQjCJNLFWdqJsaa0vyvJQRuFLIW+ae1QCr6IOm96l02jeavBnVHEu3t0g+iJxWMAsEHWN3GR9JkEdoYRFUimusBrYFVVeFoOzA2MFv23zaRs+4RISkgVwtfsQIHJeoBQzi6WsvjaRSJ5cJDt7oqMeYNqnWyWIoy7LuZu3YU+0Xlh2oISTlSqTMctQq/l3zvW2WSeZyi1vXPN26ki+2ZVnFEMmsw+bif1CmqZoqrd68hqgoa/FqWH245MFqaLoHsB1LBb0pnZrkxMofu6CrEDxJ8yBEESGXXp5nij92WAx1U7V+fZomUvsI3jM3KKhnjlkUyKJJER8ss2nDbFKTFRbjIvuh4M4VwxdembG1V1M1kWQtvSLSLx2bU0+vKIhNoCxKkQsbQz4oO3owCmUzKs1zWS5JTV6qE+Zz5pf6vHhjwrvvnuPrlyvef++Ane0J95yJfOu3PsJnP/k1Fk8vEvo5ca/i8LFDrO3sURQStSpLM9MF5aRkDgB1mn/wysWrVHXDQn/AlSeusnxmmeUTC8yvDmgijCqZNUcfqcZ7nDt1VkZFM09vecB4a8Roc8Jsv8LvNrz77Q+zP5vy6S88zsryMgsLQ1569jzTSQUkrt9Yl++sCiwuzPHYG++laRoGc33OnjjCU89dYb6/RZknNvYb8qJHv1ljPN2mWDnC4aNHuHnhZX7xv32av3Jmie0QiIVmO3d5FEZl6lK8SDKjI6SEi9rREhlkhtJYVkoHWAKCqHijMTQhqGudLrciRtmXmUjn1jWKFUFpszGKbNtolHPlhVHWjnCKzFHVnn6RUTWRIrM0IdG3kcIZntuvmDWRYebZ2NmlzAt6ZSH021trHTCyyC3jyZTJZMaXtneoqpo7bjvFG1//ALfWNlma7/HchV0+c2nMux88yfpzV5hfKbgtwumBpU6BcTBcqzyTJnBu2MMYw/zA0NSRNy5kfHbm+dheoCRxonQsFo5n9iO/vTHi6saIE3ccZfXeI9QTr0DJAkKbOghYMbgmi8p2c3xourF4Gx4iI0AZ47SHdVRGm2TcyJjZ6l5u6EReO5t6miYwmdb8wWeu8vAdBUXjubLpOTSf86ffkLMfc8qeZ2EYGNWWzBrqRg7kfj+HKHuUIpf8ntZC4GOkqWvyvOw6DmvEt5NlOSE0hNB0Mc5GobFFkctuVyXJnQJKvTHCPDOdRaKFvZJQn0gbb2E6f4jTfaPLc0wMMnlwMiItdSnfgpC7KOMQuy5P0CmpGyGiJHNi6kycKUJZ5EJESMLqy1ILcYzt7FWypq3KCY3JlDuFJqbZg9AnL7LSrMghHMzLrMu6KkKwh6lbcnXGQOvIivwgyF1bvaQvb/tStuDDLo9ENfUt+qCNv2zbf5c5mqruzDZodK1IkPVhzJ2aZdSroAouWTYZMQtqlnWWHbihGfZZmC/Z249MRoHoA3UT2dlvODywWBoKV3KzKhkMDQul4XVHEp990ZEM5LnBNB6XwLhcVDkxybgsc+As2zd2qfcqcJblU0u6EzJkeYEP4v+QLBYhvh46fYjzXx2xMfGsFI6vX2143TFH2r3Mj33gMS5cuMa1J65w/OGzjPemTJs1eVA0bVJiTX0XX2mMXl6IcS2lSK8sePXVmywvz3PmyAqjm1Ou3brO2cdOkPUyNnZmjGclw37G9a1t7rntAYrMMdmfUS702X5xh/vuOI2fC8zfO+ClS9e4fHUNawxbm7vcurWJy5yOBJI8b9HQmy+ZzGY8+cx5qlqQMKfPHGFne0SzFInGUTWJYS/j8IJj/ebL5MfP8Tdf32dycpX/83e/xNvefy8nlufZm0oWi3EWX0tWO1bwEFmhUZ1a9SYkQ9qShH2WicIlRI8UkonkXKeSMS0eWys4octGbJaT1IMT2oTL16Ba6jrgikTdPcv6noREr5DlcJFZddLDQikX3Iv7FU5pvNaJP2mmBFc5uJTeEGS043SUW/RKXr1yg82tPeq6YTqViOM//Po1vu0H3oC9vourpxxaDGxuB+ZKS9ZEFpxlFAz93OEsZCrxXMgc37tY8us7Ff9qK3Iug7me4fzMk6rE2XtPsXj3ErNRpdkoaCZJVHyMVLNRTXe2Ted0MkKS+AZZmosfQ0m3yahUVgxzwYfu904pkQIsloZmUhFwbI4D3/u6eQ71PL/zpV3yLGPYK/mRtxbs7NbMz40Z5GLM3d1vKHoF41lFnhnm50Ultb8/69zlVmdQTo3AXoU/SaO1jREKQGtTcFYLFg1liiQdw0kBGxr5Xa2TLKUszzEaJ25T0jhviDbQ+EYxQKa7VGJHG5TLJAYvu5qkZ6rhNZBHo/9/22HhRZxkX5NwKAbI1lqRY3G5PI91a9RUbItFH+ZWWVA1jbaX6oLV285Yg4+BSFTt8cGHZpMswdsY2nbBUlVVx8Rq/4xMPSNYuvhaUrtjOZhrWtTPoV1FS4HMslz2IJkSL1OQuWqWySWjD2CW5x2jxgc1JKnTs11eBe+7LIKoc2nvfWfYaRlbs1nNdNpAUbAwV1IFGA6sXIDRQpaxuz3GZjkzLFdnjqUSesOS08eGPHimx85Y8jV6A0fh2g2VFYxJIz/D5vlt7K3E7YdOcqo8zM5z2/iJJ8wi9bjGlWJ6TK2xpg6UKwOWV3p8/oUxw9Iwqxq+fr3hqy/vs3/1Aj/+o9/B1to+2+c3yOZL9vZGUtlEEQvEmDqlmelcqYkY5ftGnfNleXAYPvC6sxwbLrJ1bZf+Yo9RnRjXER9ga2/M6tBy9OgK460RGy+s89ZH7mN+aQDOcP7iDS5cuE6vzMmdE8d2Ji+vD6LuyF2bfSB7sb39kRgdneHKqzfYH09pksxqq6rmzCHDiZV5qu2r2OtP8eBixXc9dJQHFh3/+vFLuOGCHDpB2voQRMbrQ8RaQTNgZFltEzitcg2pE2WIClE+txYwmpJ+PupUbiXTGu5CiP4AV2Fc5ywHgRtmmezfRpU8cz6J0C2ofNdZS61/Zpk75no52z5xpYr0nOzihM8kiXMxRZxK16WjNF2hIdr3SL9XUDcNW7tjHnngHH/3ux7l6pUNrl3Zpk9iNg74mNEbWAKG0lkmxrDYyzjcM5xc6rM8HHB0ccjCfMF3Hhvyf96+xJ85OsdWiHxzryL3hlMPH2XxgZUOqZ5p0RZDpCgzsp7FlYasVO+DEYm1UXKFRUQHFkgKXvTt4alyVR8k3a/DDQHRy0h+kIvHf2dryqxODAY5b314hePH+2zsN3zv66SwnTXShVdTwcxPq0CWW7zJKAtHmTkIgc3tkdAYGhkt+cYfhOV1+4r2v2Uv3O5ZO1VUZjWrXM5RoZyXFHnRyX9bWrFR9V6eCUsrEmXNoJkuKcj+1zoJ0aItgFUSnOVibsxbPLtGArR4GIMWNUCWFd1+uq59h0xKej6YljOmuSNGx91ZkWPbX7z1XljrtPpuuS6ha7ktEh0btLJJOsOt6poQfHf4tBLaVj1l1NkpAfKZOCBbAFeHsU8HwDQ1RAWFixV51gVAtX+30VCTiNHDQOTBYroxBN90AD9rRNNsnMH7SFXVGpjCAS6+7X6Sws28HCDyfzPUUQx+hxZ7bOzMyE2iCbA/qQnVtIvnfGo748bMkOeWWd2QlxmFg1klEbuEQM95mqohBUF3ZGXGZH3CYprj0Tc/ANYyf2iesydPsP7cBtWlEaPndlh74ia2pYAaC85iSRy+9ygJ+P0XKo7MWx44lnFyPtCLe7zpnsO8/71v4dWXr7C7M5JFoBFlWotPNyZTFpt8P43OUzMnJrS24pjNGp5//lU+94VvsnhoyGRzRjlwTJrELMJ02lDViV7uue22E9x6dZMjgwWurq3zR5/+Bk8/8wrrmzuSp1w1KoIICvOMnVrv4L/kZ+zlpe4MJFEws7A1CbjMEPKcdz5c8de/B37uz+fceukbPHVxwmQj8qeOlzy+OeEP1/ZYKnsEE6lrZRsZWVSHCEUhM1+rWe4pRrJcwJniX1KJpWa2WNP+bJlKzJ06pBXsaYQkYBVo6JywyFzm9JkKncQ5GkMdJOYWA5kzRCPsLgGZQhMilkAJXJ5GNmtP4Q6yr+VQjRIeVtedhDPrlD1RzajyS65v7vHYXYf5pz/9p/m+c/N87zw8/sIV/PoeNzZnbOzNmAYpDL+27/nwZsPhPGehLCgNDHLL3CBjuW8o8sQbhjkP9jNM5hjGxG3vOsvSPYeY7k6x1lIMCvrzBVmekTvHdHvGjedvsv7SJmuvbJCVFoxU4inJOYMWKz4GsO1n77oJRZ5LMJQIYWy3dLaFvO+LPQuZpcJSN4n5MjJmwIPn+rzlyIyytExmDXnPMp16MiKTmSeoJHZ9c0qWW8p+RmoC01rzetq47bwNhBNqrlN8Usv5a4vmloocvbK6lLNmrRUUU4e8TJ1XzVg5pFsxU1D6gtM883Y0munINL127G8OMu/bPKbW5mC0w82c07NDluSzutLxve3kutYcsNxaT4tYHKTQkh24xxonF0bLVzIYcUBrbkb7MPrgxbBn7EEilaJSOiIlB7njwjOSH8g3XowtKlXt7o4kYUGtvLZVoLQqgcY30r4GcYkbq4YZnYW2F0SuXYXA2NS5qV2VUX9IXuRqrNGo0CjVZefBUE10lstCryxLjHUiywswqyPkjqU5x+7UM7/YY26uz/ao0VbSsu5zvnINrmwEHr8BBE+1O+JQ37I7bT/wyDDPiMaSF4ZYRdZe3GLzwi4nTh/mK19+hieefJ4vf/lJbq5v8tibHuT06eO86a0PcKScY/2lNXE2e2FeeZ/Il+Y4+8hJdvZmfPSFwGgWGfYL9rf3uPHi43zfex/mxIkjhMbjVBZq3UEn1kbahijdXpnlyqGCQpMpW1FE0cuZTmt296cMez2KfsZ4GhjXAo2cVJ4w2uOuk4c4vLrCqTMnuHztFnPDvsZxyoHaYq/brArMQd6FqHFaZ7G8IEXutIsI+MazM4mQGnYrMTgOhob3PFzw+lOJ//rilI2NCa8va/7kIcu/fHmDV6YNQ/UMxRRpIgRUUq2dpm0lj9Z1L548UGqAzf/7DGvvA5lSpa0xBywnXZqnltcUUoejiIrskO7HUTURr2YwqaQTzlgarQ4t4j1YHpaUznK5TkSsxI2iCX96MFXeyx7S2I7u2qVdKt345q1t3vuGc/z8z/4w6fFnOP/Fp3jPqSXKjW3+9fUJrzSRaA2zCKPac6gwvFxHfubahD+6OeHJWxOubI3Y2Z2wsTNjfVTzSzf3+dkre+zvTDn15lP0TsxhvWFheY5Yeda+eZOXP3WBvWt7zHYrtl/c4r5Dx7l97hDDzciVr1yTZL9SAK2CEnGYzIpRuZZgqbbKNki+D87iQ6OMMM2y8JHCGAYy7Wd95Jk1EKNhOCxJZOAKMLC/W7NyuKCqEnlp2d4PLAxysI79mWdhvmA47FHViY3tifLFRMJOarONrPLnasWCyKK7KIXW7YyOR1U1GpGiO3OOFATNHjuX+0FGjNXvkSjg0ZZ4bjTrMGrGh2+Eyt2SqtvddYwSnyHvU+hQMQJdjwqxdKo0deq3O1g1GCP2jBYNn9l2bO5xOgnCWDKDhrRbo8Y5gd9lRSEHi46lspbTj+l4/S6z0hrphZG0YvRqLJSQIs3i0M27aQGCxuLUSdnK3apZ1RkFXS4XRUunTIo5aLHLIE5Mq+qBzBnIHHVdd2ZIZ4UGG0I8CKppR3bR0zSKkNcPyuWa/pa1hi+5ZMaTMTvb+xACq0fm2Nyt6RWWk4cyrm5Puf1YQUbkD14KvOfujDPLhl/40oTUJN73gGOplzHxgcnUi6Q2Fynu2nNrDOjzwLlzuNxx8dI1RqMxc/NDSIm93X2++ex59vdGLMwPefRN97Pz+PNM9ioyjTTNehlNVWEHJcPFgrX9ht99Dt50ynGk75ldusSJcsBdZw/zlacu0Bv0pYoLHp+SjilkXJArK8sol8tHjzO5hnBp3KazNCS2tvdwA0t/qcDHyPbEc9fhHsklJr7ijtNHiAa+9I2nJFLToEiaTJVtHJB+YxTdu5WCJEbpTHpliVegZqUEUCkIDKOZoxzkTJvAtc3IPUMPyfE9D1l++oJlzc+4M4v85A8vcfUjM/7p+R1+9sx85yFIAgkjzySIx1lL3WjGtxHBlAqFCDGQKcNNUDu8pnuVsWton6s2NlRR5T56nBXacGhZa616MSSmPiq+35A7caRLOqLsZIwRyfdCLvkk52cNmYIS20PBGjGf5XrxBx+VtBxJQd7TaeWZ7I/4H9/3EH/rJ/4kW199lmu/9wVBuvQM7zw+z7Pn9/n1Pc/3Zhn3ZRGXwd2546dO9vjIbuA/73qO9DL6+54UKjyJGzPPdpXoBcO5t51l6Z5VXv3qVabrI3plifNweHGON7/5Dm6sbXPxpZs88sDtrG3usDua8MbH7uXmjW0e//J5sp6jyCQ6YTyr6C2XzK/OUQwcdePlPNI01CwrqOtKRjyh7V5kAtJ3lsWhyLbXdxsGecPjr0y459Q6+9tT5udziA15v8dwLuOFV0bceeowt3YmPHCnpakDO+PEsG/JnGF7d8bO3qxzvFuXyd5SL7P2fIrdIR2Ewp3J3MaHphP5CNtOkfZJx/XO6A7Z6vnZ6PlpdC8SuiRWmdDomN66bv/bNPJcBM0hMfq+pWC68VkKidxlRGT0ZtIBJDd4j1UAbqZekNZ4iE1UtYRixaAXjRooM2m1rLYrWRdmExWv3UoQjXUq/YrCQNHqsE14srm8RFhDqD1lL+scoiSjEjwPiNzOKQ4YY7rWzmlVatWHkamRKGmOdxdc04LF/IG6RRafqYOvRR09JB9VHSNdCdYSak+InrJXHmQLaFqbtPpO6ZPaQmOYVjJDXhoWBBxbWxPecf8CH/z4Hu+3gU+/3HB8KfGBh0r2ZpEfeKTP1T3DR56acGzJsbnnmVaBXhZZMLDx/BqP3n07J04eZTyrGe1P2d0bY5VqbKzD5kL+LIqCvb0RX//GszQTT9wcs3B6EQusvbjG+No+84MeR1dW2N7e4/JWzbUdz1IvsdhvaL7xdW7u1py881Gyco7gZzTVjDDbhuk2dV3hshKb5HcNKWFU6BBjg7NZR0mNIZCXmWZRWw4dWyTr32JzBrNJDdFy+domdz54Jw7B2JRFjq+aLjzJqvm0lSpmWS656Zko4YxGCtS+1oWk7/IREolekbEzTmzuzCid48Z2ZGfXc/x4wTsfKlj5ZMPnZ5HHfmARe98cf7Fy/JV/XfFvb+3xvxxbYpwiWZSI2KaJav6K9IpcRnaaqaCZjWRqdHtt7rRrs3I02jhpTotzVv0dWqXalm7tcEZ2iliR6hYuY1qLfLcOARPFmR5SlMs9yTM813PM9SybTeDCuGFQOqIRySoJ6VZU7tuaa1MLfrWGzZ0xq0XiZ//a+/mOb3sTV37vc+w//hym53Ax4iyc3254X9+y6BK/sVNxuzM8OpexEhIVMMgNuTHcmtVMZwr3NFAMS87duczSXYeYOzrHK3/wIkeyPmceuoednSnOJoYLfV66cJ2FhTne/Ib7ePaFS1y9uSZS9SZw/MghHr77HL2y0FGLdJrrm3tcevomS3ctMTw6pJ7W6rMwRF8LRFBHoJkTpVoA5nuyL2pM5MZuw+nljLfemfP7X93hZBnJa0/PJY6eXKAJBatLfao6sDOOnD7ah7xge1Rz4twK9HPG4zH740Y8yMbSVLV003phdVRbdZwXRS6R3skoSVxGszaXQ7vxB/SOZCWlMuuovQGL7XKLmqbRrkT2J6FpOr+M16wlkHwg4woxhFdNJ7WVI0/PM40Jjj5gcgPK7opRCqjgA3mZ01S1gk+TBrOJknU2q2XE1Rb01pDBgXkqoQRa36gBLO8Iqq3aIc8ll0IAeVnHrplVM4qikLGVZoOkdJBPXmuOhbVSeSVLJxkTIJmoZEIM5EHkYvrRCGIlyyXYSPOkjSoe9F6RzOc2r8QaMQDqItI3ssCOxkiyW2ZJQRymoommC7MSs45cpkYXpz7Axm4F3rM4yAghcWO74d2PHuW/fmaTb16veWkb/se39THDeQ4vgr+2y1//M/fz4c/e4NNfOU/dOPankYXVHq4ecXT5MIePrvDZLzxFVdfkzglzrDU4trA0Y0kmUJQ5o/0pzcxjtizlucPceP4ag6njdY89hM2gmspDUtUbWOfYqwKbk4hvKuYWFgj1LrOda7g8IyvmsXNHcQsnqXau4Ue3SHlP1HHaoWGFfyYIEg39UjTDbFazemwZksxmt8e1yFeRS29lPmdxcZ6trS2R+Kb2swzk1ukoQgKvjKPLbRcdusSlei/mPqeLajQnxjnL3qjmxm7Bvacs66OcxX7DrcuJP/OuOb77ocSHvjDix/o5x4D50PC2U5ZPXJpwLBvxQ0cG7IdEpkgI3wRhNDWNZj2IOk5eaqNhPFYjnpNGf2r0MNKxO/MaKXobzuWE8WSdIyWvpjIF1EWoQmDgr5wvAAEAAElEQVSmF6ZkUsuIq2XfO2vwwdArHcZaXp1EtqpEmQm4tGkkJll09abL8rGqvNnbmxKmFe9781n+5l98H0cHA175z79Lc/EaKcvIEQ/F+c2K6cTTz+Hd/R4n+oHf3fX82k4gx7A7ayBzzA8L5pdLDi2W5MOc3kKfrMzoL/fYfXWPFz78PPedOcapc6u8+PINssyxubHP9SdekrgHLFluCSEx6ElG0Nb2HjdvbnL4yBJzgz6j0YSq8RxeXeLksVWGwx7PXLjEYGUokQxWjMJJ/V1WIj+Y1Q2ZBaJlsXRkQaKKN0eB3BjecNccFzb2ePbyPm8/Z/A4jp1d4Gtf3+XUiuOFy1N6A8fRQyWjCnZHgdct9aAs2d6vmM28YkxixylLEbJclGB5lsuyO2mImkZHuMyqv0yIykFHoMY4SRb1UpSniPrCNMzOqswZkdKLYfpAbWU1Vx6d4JDLBSDnm/rtFOyYXKYL9KQXk5LFo8ZGZA5jdM/sYwc59bXHZrZD0xd5obqM2DUWWaaxsq380GQG40UbH4LMVEMbAVpIWxjqg9lvWyFmTg7kLttbb6+oBEkL3YtllEIpeODYIVKwhlRLNRdbxYARPIn8PUnHGArjMra7ABIH7ZywroK66SWW0aakbP6DL6ljlbmWyOu7UKx2GZW5HOcaNrcmABw5tkBeFmzuNRxdHfD+Ny7xwT+8xrc+vMiZY3NkC8v4nVvsrW9yaRv+4g8+xPIg8HO/d5XdaWS1Diz0RE313PMXqJuGfr/Uz1/0+yF4jHV4X0vsapIdlLC/HKPNMS99/iWOzi/y2Fsf4vkXL/Dq1VuqAEEr/SgvrTEM+gO8n1LvTSVuFUM13RUXT7HAYOkU1dxhqpsvYbJEQugB7Wy/I8RqpTsZzegNeywcmaeaBhYPDcQoWTgWejCpGo4fX+Do0RWuXb/BoFeS5Tl146UbaTzRS6VmTJunIFicJkZGM080kOuSO0R5rkqN6I1R5kuvrDvecofjq5cK7jrh+MOnZ+SLK3zgDRN+5TOJf/5re/y1dyROLvb4S99f8up/Nnzo1oRjzvK2pZIpCRcRp3wE55R/ZCwuaz1FdGNZ5zp0nCzW1WGedRk6B7BOqz4Geflbw21U02igyHPG04YQIXdOgqQ0wC11cmBLQN+1mHhmv2YaIz3TOpdRv4DRi9UQEuzsTfDjCa+7/TA//qfexbe+/i52n3uVl778JHEkCXX93FHVDc+tT2gw9B3M6kBV5jy6NOTrzYwXxhOyfsnR25dYODVPsdjH9SHV4GfiuLYRrn35Gma95o0P3MZ4VvHRj38dnKGeitGuNyg6D5eo2Q7SGDNn6C0MGY2m7OyO1IUN126sc+HidR575D6Oryxz64kbkFkWb1sg68vzWVeNcqCCPEsxUk0aVvslvb5jVMHWKHG8jOxO4NvfvMLv724xmcw494bbmTUl+zs1px9Z5d/+4jXe8sgJhv2Mjd0Ze+Oa0yfmIURurI3ZH81YmOsJg09pvFnuqJumw45432gon2vTSWRXphHD7biq8QI8bBTt5JxkLOV5LsqyJOmUSZVTYqBu8f6iFPV1fRDfbVrDo+kmPil4jVJwajw8mNwYZ1+zmnAa5mY6QGSLdjLOaIKhLPbRczeFpOiVgG0aL1pkpdv6OuiHEToKpFFe7Gs5LcJwUZmXFQd1iL6DdoUQ5aHO7X8Xa5vEjyW+iohKZ4WJ39RNZ8WPSuAMjXDq67rplARdvreC161WtFFZ9yFE6aJUbmehy1qO8cCY45VN1CojOFDxa/cUdUmX2NiaQu3p9yyDQc61W2Man3jXG47iyoKHzxaMpwlL4NrVXc4sWZ587hpr1zd5/V1LlFlirzGUpWV56Dh5dIWlpaVOiik7IttRP0Pw5C6jUZdqEhES02rK/MI8Jw8dYunQAp/67Nd4+eJVjC6mizKXHHcUb+8kjyJzOcZlIjjAYVyOMQ5bjZjefJY43qJ/7F4imSwEjag0YkLlvHJo1lVD2e9RYLnx0iblXMnK6pBxlZg2GQt5YDybMd3d4vZTx/B1g8EynVaSjKZhXjJPleVoVQV2ZxWjuqFwhtOHBtx3ZJ75fkGZOW5f7XFiuUfjA3WQ1L7MwsUNx6H5nJ2dGYcW4epm4vr1Pc4dgx9+V49f/+KYJy/OmFtOnD3d4ye+v8e8zfnF9SnPjTwLmQgmvI9EBf9Efcm8qn2MMVSNVJ9exSVRBSBN546O6iNKGg2tOd/OKUZH/Qw64rBOgrFqr0TmNtJVexqcdDIhRfqFYzjIqG3G83sVub7oSdDRkv+tUdFbG3uM1rd4+NQ8//QnP8B/+Nkf5LGFPi/90ke5+odfhDqQ9wpyl9gcTXl1eyoqN2eoicwVOfvJ8vdvTPjozRFHji9x+h3nOPKGI/QO94mNp9pvqEcVxsvvO1rfZziz3PfQWdZ3Rzz/8lWMdZR5znA4kDRFLR6jEgAwKrbhNTsBA0WWicgls+TOUZYZe6MJd507zVsffZCH7zjHjaduMt6cKREjE4Ktyq4TIso5NnQ4JGVyYxQ5eaTPCxf2yQrDiSVL7PVYve0M3/jKLV73wBxffXHMxtTyntcvMQuGioxZnThzYgiZ5dbavkih9ftrJdJN3XQ54q19QDploWO0GRsS/ibx0fWsxlqRzTs9t5ICEtsxU9BnLwXdQ4SgNILYYUZaKkhoHexqrqzrRic2SrF2sp4IPhws4lV85JwYaLHtTi3ig9cdDWQ2OzCPm9d0uM7S1LXQn4VgK2ZAWapo/nYmv5y037onCamLbkTbX6PqmaKQ6jDL825TH6PsQ4zect4ncAmSE+NPE5TzI2wlk2WSA6w8mbaLSbqn8cFLkl1eQAqKgNa5YEz0emUHEstNITVAUJOhhvlYRdB3fo+YaIJEtkp+sSdz4uqVHHVR3Gxsz4hVxcJcn0PzJbf2avLcMAuWhZ4gomVP4lnfqjmxmPE7z1/joTN9zhwpKC2sb9WkMwVHD5Xw0oRbazsdRiUvCghGs7OtyGyNIVrZC2RZTl1XLMzN8br776T2nqefeRGfknx32hLL6DCQtJWN2kXWTSOiAieLXYvFGUswAWsL0mwd4gS3eIpm+4K01SJGBZeom1Z9Ial+d9x3jIuXN0j3H6O/0KNqEuujyNGhZTt4NtducdfZExjrcLklCxKIlXReW3u0aoKTqyWHBgv0nOPoUkmOpSicuP+lQcTHROUD41pAeZ9/ZYsr24mVlQWG5QaZ8UxHkRdenfDQPSv8yHfP85tf2ud3n418+zssVWj4tnfN8dd34Z/9lxE/vz7hH/ZzjuWWWdLupwnkVsY/JMVzhEiZaaaGc7oI17x6Y1Xunrq8d6u+DKcGOZcVXedhUuoWnRhD5cUo2BpMMsUC9TDEoqDuBCiwG8T/MSgyHZEY6rphMp2SvOfoUs7733EH7//W+3njXSeoL9zk2q98hsmtdVKW43ql0GmdZ2PUsDdtyKQ1JzOGhTzn8WD41Stj9kcVZ86tcvrt55hNZjSjGpxSKXTfmDnDoF+w/vw6D506yaXLt7hw8Qbz80Pxeqk3xmjcb9PUZGrkbHl4Qq/VzHZjtIO2XYFojeHCpWtcuX4Lawzf9o5HuPPsKa5dXGNWOuyiZXB4gG86oCzWwNFhBj6wuTVja9+TZSWTacPO+oilYc7w0AI7OxUPPLCE9VP+wx/u8uPffZoyM/i6YWsv4rGsLvcgJK7fHEnxmYKeg0FwTlpwOuvwip5Bw5lSa8oOAaej19zleuj7LkkVXUZj5MLQARehEZl69NKFCA0kKZVbmFlNHfQMO7AreCUDt2iibodobBuE2K0d0msvdqVwGCsiBTGi+gOBlAIujUItOzlzUqlZq4ay1uBwOGepvFeHuFUHZVTeijLobUbtazLrOhNh8FIxh3aO+5o9SVI0iEmtdK1FSeQqnzRdOFSblWGtEE1jJztTE5cuCI1m/7apWq13wGk1EG1nKZCUMZK2jqJ8aBeN7YjG2UwNTXKIpChfxtrmlMnUM7e0yOnj81xbv0UKkWs39hj2DaFJNHg2N3YELLk45PaFXf7bZ67zv/3ACY4slLy6PqMolxhg2dtaZ3/kmRuWRKvhL0bS4EKUCqL2YkTLnKVpBPPRK3t887mX2BvPKHOpJrs4yyCz/JaoqQwAkUiqFDDGhDMHDCzjxN1L0cNXU0JzFdtfoZltKr5Z0iOdiguwBl/NmExqelnGaGdKUwXZC00MZxcdRQWj0Q7nTr2Ohbk5qlnTdXlV5clN5M3nllkoc0iJQ3MluRNXeKMjqsZry28SPsqib1BYBrljcZBzYmnE1Y2KWV1z7FBOk+DYUs1Tlw2PvT5wajjgb/zQCn/zX1zn+VfH/GH9RrJ1+OkfvUITSv75f1nn/70x5h+eHLLUg/1KdjOeQK4sOB8DuSK3W9Jpy2PK9ftpd2/WyLgp6vjKB7lwBHtxEAUgVOmcOnimjceqtDgplG8aDL+2M5bDyFhGCXpbNY7Irf2Kphal2KF+xt3HBjxw12le/+A5XnfXUZYxjF+8zqX//Amq7X056PulqBdDYn3asFtFUhMYlhafIkNrSWXOr255PrVRsegS9955nM3RjPHuVDhsJmHRcCwrHX9mDZe/foW8EsrD5uYeC3MDuRBalVr7/GnyaIiSMe6ME6KrKFtUVCMZQqHxGlYUdXQqRs8QPU9+8wIP3nuO44eX6RU5X/7ac0ztlN5SSfSRaKDMHCtDhykCU5szKBKX1yLN0JNd2qHwFb2mZi7uU/mKn/mtXc4dH/LGOwfsVQ1zg5JbG2OKPGd1ZQCThis39iQELpOcjbaqR/0XdVO9Jsc8qTmvHbV43Z2pJ84YrHGq5ksd2QNjyQrBOjn9s2IKXbfe1EFG2ErCJUmGedDAqNB4kk3kZUFqlVLq/G8z59sOuaUXZAr3tMaA/jOt8TPElu1lO/Nsy3FrYebOWrIWp+6MGANDrQEvTeg2NkZHQU4X07G9PBWBkaJeBhwEzwuxV+Wzzh44IPUmE4Ng3kkoG99Iax/BaZKd7QilbYbHwUy+0UxpUUHELksjy4uufYwYikLGBV2gb0eXDIJgUdFA7vLOum8V5NgutMtexu6oYjKLzOU5J4/0+fSLFU3t2R83DAtoGlE27O015IVje3vG0fmc8e6UD35mk4US1vaCSDJLw+kjA65uTjFJZopNbChcQeMbqU6rSiXGltDUWCvxpDsjmRP3eko0jYIkaVNwW6l1h4rQ8aM1hto3EhDmnF5OkZ7L8MljQps9XmPsmJTPY/0e1pQ0sVEfiPBvjIG6DkRrKBczds5PWJgvubEbCccsszqwdnONB+41LC0tsL6xISj+lJjrWd5+2wpDHYNaKzpzH8XkaTCUue1GkdGKIqpReXkTAmVjWR1mvLAx4/x24o33GebmHA/e5fj41xv2zj7KtZ2Cf/6BL/Grv5fzt/7TiPimHpN8mfed3eLP/+nE5Svz/PaXp/yTrOIfHO8xzC3jEDU0TZ7VzFndR9nXZEcLCrtuZCfVjrEa9WCkpNWc7kOEvKCLXx23WhJVk7qluixl5cE8NF+wmBsuTjx1LRnYUwz9XsGjRxa47fAcD55Z5N7bj3PuyDIlhsn1LcYfe5wL1zaJTU0+zMgHBU7QsWyOatYmQQLBjEqYjWU+y3mySXz4esV2nTg6yJiMK8pBRjmC2d6UwWJfx9EBlxs1n8Llr77KudXDZP0eT33zQkeVqOuGLC+64tJ0Fbbwl+RgVARSXnRdTevqxqaumGyVfnmRkcWM67fW2dzaZVLNuOPsce675wyPv3Se4aEBxkEdEvOFZaUnHotb6xWFM/zou1f4dx/boDAzjvUd4xtjji85/sF/2+a+u45w16rj0rU9jqz06fVybm15Fudyjiz3idOaK9d3QKkcrfcsatw2HBiPs9xpKNqBGEn8a226a5uYqhglK2PplrrgjGHWNNjSah6HHuYhkus5JrHaUtR6FQl1pu02d13f/VIvPNpRFlFlxnJuQ9LRW6ZrhfYMp9utBN3xCalc9jwSGW7UH6L7CnFSS85Eq7LKnOtuW6sLRGutJKjNKrI8V2mnGvZ0eesbgaGJzV+lukba16izuEyR6hjl/+utKzwt2aLHlugbRXWBtlxt0lY3StNbIcXEdDIhyzOZK0dxlUo+7UEEr3UqHE5ohXiwIGrqhqIEh+Yf6++zvgm31vc4cnqFu29b5Q/++EVsnpOXBWVhyXNw/YzRfqCXW3JTc3UUePt9K1zcSby6NsNkllGTOLqUc2q5FhOYNaQmUeQlTVPJfiIlpbhCaGqSG5BJEKYGa3nFVwuYLYSIjw25E7NlbFlhOgZMMYqb2TrtVNTZ6hxN9AJlw5Mrgj+FGeQljcnBzyiKUuajMUAy5L2caxc3OP7IUUY7++TJsLA0YH97D59KDInJrMHUI04eP8KVK9folY6iSLzj9mV6mZOMDCs7llZGmDmnaBqd6rRGjNeQCnInL+CRpR7NK/t89XzN//pdJSlEXnd3wT//rT2e/kaPmZ/jh+5Z5a/96B5//u9d4OThF+Cud/ArT5zl537wPH/1rx7hxqUrfHmz4mcN/MzJAQtFZL8JxNgaUm3nF+ninE1SCJ7Tn1t5Qzq6aimyLSrCpUzGorp0NPqsThoZ02QSZYFJUBaWk8sFf703h10dwrDEDnqSNpgCZTTkTaCuamYvXeXaF14gjseYEMlzKPMMUwh+3KbI/tSzM2vYrwLRy19S5pbSWq40iU/sNzw3hSzBwMhz0+sXhDoSTCQrc2bTmrKXy2eSRE002Z1ybH6J+ZU5vvzF5zDOUqqr2mU5zhrqUMk7n4QFhZPnNcY2Qyh/TXibVscmdiPCGNBZv9OIA9/tEYos58VXrvDo6+7l5Ooq65sj+oOcJlpW53P6mZhjb+03LPQtq8sFH3hkyG98bspbj8ORQ45/8pFd5pbn+Zvfc5gPfWqN8QyK3jwRePXmlONH55ifK9nZbbixNqYsc6FbYDUMTjt6LTBQE6G8dlHfL7U6BBlbZXmBSSrY0UPdE/FN0vwkkYfHJIbmpgld8mLssFKa16F8tpSk4EfJIRLwafHJC5XgNYzBlhsoD6jp2IWJJKsGazqkiYybNVIjKGOwbhQnpYQDA5mIPw7CRxovi6HMieKqrmrKXqm4c5V/6RKn49qnpPNhyWTMslwX2uKszXOR/HpNuQo+YBRJEoIsh5zLFPjlDsZVRj0dusAxCjBrgWDChXGKbXcSLqNxoBYJbZLQnqYLnm9VNUI9VsNcCPp3my4juPaN/nOGXlkwGtds78uHdvLogOkssL42ZWlYUtWeySQnJ+BizXTSkM9nGGomU89bX3cM6wz/4Y+22a8SSxPPsYUMUpB44BhJ3nfZEFVdiUrDe8gGuGIeM1vDZb0O5SL+CdlHWK2QvQZ6SaaJjId8CuTOYRApoLG2W4w5Nb7lmSNE/Zw0J8PEGdiSFGR0hj4neZkxGdfMrfZZuW+ZW89u0OsXrK/vkoxjNGlYme8zmgbyVHHi+BGCD4Qycftyj0HhBCKYOcngsAczda+dZFLFHSRJqCPpbkLRCkmUSANneep8ImYDmtmY288MeP25TZ64cZH+kdfxy19c4Rd/bJv3f+IE33j+CotnN/nsq0f50Fc2+bPv2uCv/cV5ln9+zEcnFT9zPfJ3Vgv6JUwBAngjI4e2e21nxu3ope1KUopd4JokXltlGznBlmhEL7aVaBhm7VI2QWYTdTQsFIZ6UhNnDcX2SKi80VDp0jyEhrrxwnmz4ofJNLUyGkNmpJzaHNXsV4GZjxgbyRzkQOkcm8bym9sNX9uvaSKs9gtcmVE3nhBgbq5PrCNmUGCLDGaNZJhn8qy4wpGmgePHDvHMC5elO8gymkYLMTUmSzcGWSGodZE3C+m2jY7wXhzOUcc2xsoh54Oo9YhGF8tS8dpwILrJMsdoMiF3GWEayQ6X+P2ao3OOsnD4xnBxrebQ0BLyHvednXHuucTjG5bJjUgqevyTHzrJuEkUNhJdn6IsmE1qrm81vPENC7h+j+0ru2zvTCX6QAU8msUqv2+Ui6JuGkjqMCcJ/0wD92JIZFZ3zDGI0szYrjByTgo3iaGVnUa0MiVplawtk9BpTn1L7I5KI27J3Vme6YK+JXNY9dvJWZhC6GI2fHOQhd6CH9uJkcusRmk4pSn4Ds3jXxPfa9tDs9U2O+sUSigJVdlrIjaljc0kgKnN/VU1hFE8tmRZSMVmVX3gvaduGt0ViozMqjY5z3Iyl+mHIRVdG3dZ17Idc91i0mncZOguuwOJXOqWlWhV21bfBj0wM6emtUYUHO2Xp9JB+d3VzGgV+qiXZN0ELr26BdZyxx1HKft9NivLsZUekDGuIpOJp2c8/dKwMbaMfclcHpkFwyP3rTIsLdv7Dc7BqcMFw1IOy7woVDAgo5D2Ygi+Ih8ewfh9UX2ERk1sVpPOeA2ORI2Y6uVpgWhlljGpakazGbWPTKuGybSm9oFp7Q+McElw8TGKpwZfyw6g6IleXP08vgoUczmnHjtBqBOTvZostwyGJXWdmASHCZ7RqGays8HpY4cxDmLwHBoWNE2g1Cz3kLwmYOZMa989S7IARPEKqPlOLvym8Uwrz3DQ4/hyztMv7XP+8owQHUeWC77jsT7V9VeZ7zV849oyz94c8iPf2ac0DeXWZUob+a9PnOPV6wV3vmnAd9xt+HYDz1aef3B9zE6wDI3sB4zOaiWhTdhYB+bAqOGWciEnLaTyLJeXyx4s141xqmKzEk7UiMTcaredkizKe5mjzMB5iZNtUqSJ0o3KXsyxMDegl2f0ykIIqRjGdWR7Eri80/DKrRF7VUPtAzYF5pxhyTn2s4zfHiX+2Y0ZXxoFsmR58IhU3LNZraNNS1N5yl6BDYpby5KAS9WEYGLCephMqq4r9F1stdGluXQrxEhV14wrz7SSYrJuooIlI7nNukM0aWUstAJ5z0P0Mh6KAd0di9pIF8Z7e2OqumHnyh5hEkjGcsfhHGsio5lhfS9x5ugAUsPWzpT33iZesL3o+In3H6Ii48rFbQieOjhcglkT2Z0k7rl9GXLLhatbbO9NcRLn0bGu0IV2UlhmphdmS91od2a+7Tpb5RyG2ou3TQyQcr5WVYXNJPraudewv1pBtxZ+VS1nl8tsh2kiiZcoKwpm01lH0xZ3vjxjkl4eu58xxdRBHK2R8zglWc5HVbS1uxiJ6KUD0dJ5RyyWIC9riEGJvKlLGWwZKvWsxhpDUYpfoVDDGynha5F9tVplUUa1qWJRs72tToikbbNYdbXT4VBSEvNgVB195pT1oyqtoKiAlnbZJRlqlG4LLZP0rdhVhV5neJleDCaJr6SN343hINfYN6FTc0SVxyWVaIYUubY2hipwfLlHv7A8+coOq8sFlUclmYlpHSgyx1everYqy6BXYJuGgUkcWsi4tFbjrGF5YFjoZ9SzoFkAMm5z6pEhBvLBIWyWQzMhtcmOLbXYpG7kl/R/TzGSWdl/NI2MrfamNccXety+Osdb71jmHXetcuZQn5Vhzuowo/GRWRM602Qrr80yR26T7kqE6ukKGT2Z3JKcYTpuWDm9xMb1Pe3uHHuVocgMtyaW8dY6h5eHEC1lJtnuzmXUXqqsIstwxpAZI14I5Td5zUDHIJkxRhH7xlBm4q71MXFsPmdr1PDl52uW5iP/8fcnvLh3ituWtpnt7zLJF/jgl4/yPe89xlseXmB85TL5eJP9uMzPffQkRZlz+3ct8d3LGd81Z3mhivzUpT0uVJGVXkFqD05rJGJYfSrosxZSJKhj3OjFUjdiRLXJqdkQxcBL5VxkhkkdqYMhqHRSuhUjDmUd52FlttXLZK6dOXnpZ41n0ni2xg3XNiuu7DVc3Gm4MarYrxryzNC3hqUchmXGBe/41TH8Pzcq/mi7ZjJpONa3fNt9q7z1jkPMqUlRZJsi096fVsxlGeONPbKeTA7EDBopexnVrGZ3Z6S7RLq8C3nvAqNJxWg6pfGBQW6498iAO4/0CTGy0Hc0vlaXvemiWI1B/m+aCtlyw6TyzYWerYexc3KY7o3G3H3nad75xocYXdlnIc85Oy/jnp1ZYGcS6ePxG1tUa1vsbDXk3vM/vWPI4RPL7O3M2N0a4WNkOqvplxk7k8jufsNtpxYhBi5e3SVEqe7RkXd7SMheIKpSU7uvJPEMQt2Q30P+s3zX7TjK6h43xaDMKg2iaoIQDtR6kCv9oGm8XtixOxPbaIzgPT5Ekf6rV0/FXDLCT7KiaEOuWiq0jLK10LdWrQ5JwZBasBu6yUBTN9IhGQm0cxaypN4OgFwX0NYZjBedefBBkB8trlj/ApIoeCRDOnY8rUZ9DEYX5sFLV+Kj5HmjWuYsc+R5IfN8fZHqJnWcoFZbn2c5MZMDNWlFk+Uy3okagOWyg+Vca4Rxio+IsY11REdxyutvGoElZhmohvtADKC3sipOUpJF3+XrIzAZyysLnDm5zPmr+/zgO47hipyd8RjnehxZ6XPh1owvXkqcXoSdSWTejFm9/Sgnlgsu3JqBsSwMMlaHhitrDXNZTu6EC2Uzq0vLiv7K3aRmrwOkYQ1V09Av+vrli+KjLHJlBcGs8gQiwyKnzB23HZ7j/mNz4vFRffmZ5YFyngwXNkZc2JywOaoZTWfMDUrNm/cUuWMWIOrlVE8bUgwsHTlEM5Vuz5WW+YU59kZjYmrYHhsKl1jtBTbW1li829IbDCA29HMnByFRfRSRiY/MZjPm+jlzZU6KXi92Ud55DXzOnKHIoWkkmS/4wB1HFvjK+R0++tUZ63vL/NJHd7jtTMHdpxe5dPklend9C3986Rgvv/oqP/EDC/zEz90kbF6lOHaCr986y89/eJu//G17PLlieevVGg71+Oi252eujPnrp/q8eb5gqxZVVspb9ZXr3oPMZgo91BFCFGFHiLGjM1stgmIS2W8IgYl6QjJj8FHGjCZZNqYGprJM1RgqMuvxerhY55g0kV4m75GXu43MGhYyR4nIQG8Ew/OV5VlvOD8J7E1mzFnD2cWSEwvz3HlsDl1NsjLI2Zp4ikyo1lH9LWI+jdjcEiat8dOxtzEm2w8snVzk0o1NQJAsSeX3RWY5u9rjxNKAlYWSQZYxV8qY+KETS/QKy/a45vMvr7E3rXBWCongBZjZDvl88ORZThM8MTTktiCojDaEqKTohs9+6Ukeuv92jq8eYndzn+X+AtE3XLo+gdRw6Vbk1kZD5gPrY8e331tybKnE5Iab18acGmY8daPiyDDQKy0Xb1bkZc5tpxcAw0vn12U8t9DvvBRBLxFnM40mjiIC0vNSsjXkO4ox4Uo15jaR0Ih4yLbqRzVEGyMGwjwTUkckETXR01iR5NM54PWQ96ETeCS9cDOnfjY1/6Kx32VRqMxfwZ1oFIdSoY0xgolxRvD5zlI3tY4BEq4QT4hEfduOmJ7J4sRqtq/XHYWX5VYCm6mphRb8JcvlBIo2ybpUtnafkrtcxlgmkZXqwE108tw8L7quB2MxKarUV1qopHTUolDmkrrcs9ypXt9j44FJxnaZCLarhhL6BaksLcUkGdYaMiWa5wMVhTGaia4zznY3EoHMwGDQ4/zlHapxTbnY5/TxRV556RKHlvucOz7H1bVdVhcFlPiHL0Xed19OGRs+/PUZ3/Ow4Q6bc2Klz+MvTqmbSEbDmaMFX78wVmlx6PJIsszR2DnscIHq+gUilqTxvJnNqX1Na6MMOlqKEZlpDzPuO77I4YUeuRV8wqz2BPGnCTTQN+TOUWSOB47Pcc/xBTb3K755bYcLWxMKI3uQ0HhcXuBjQ6wj+aDg8L2rDFZ6xJmnnCu58c11js8ts+l3ccay2zjqEDg5D1tbm7xuGDh16ghXLlzUBbTMjC9sTHhpbSK7AP3Z3nhukdOLJVXdkFzqsNXOaOhQnlHkOU9e3efGridGUXVduemZzZ/k7JkZ+9tj5h44zuKz55lWD7JhD/H5S6v8+MMVP/iOef7jx1+mPHqaVWbYVy8z2k48/AMlf/jPAu9dDswXht+/bvnZl/b5n84N+K5DA/Z9lOApXYraZLvn17VBUiaRFK7YKWOU/Cudse52cExqOQDrIF4fmSe3oE/DLCaaKGjzzCB7DCx14+lnDqJnUGTgRDmzUwWuNXChjlxqDDdjYJZkbr3Sszxw9hDHl3IO9WVEOasbyePxgTOrc7yyMaOuPXkho5NeWTDZnzE4fYim9kI0CLJ/2rmyzV3HD3HxxgbOGpxr880FTfOOOw9xdKEn0QUpEEjsTyScqafpp0v9jG+79ygb+xUvrY2Y1J7FYc7+TM2A6WAXElsAZ6jJMsEboedIlpXUTcWV62scPXYYP9umZ/3/n6r3jLMsucv7v1V1wk3dfTuHyXl2ZzbnHKVdSSiAAgKBZILBRMOfaIONABtjwNgEE4wNAhEEynGVVqtN2rw7YWcnp56ZzvHmk6r+L6rO6eWjF9JHuz1z+9x7q37heb4P0hdcWEkJJPzgo1t49dAsW33oJII9YYJX8lmab7N4eYkbbhtgeiFi/9UTGAPn52IGB8qMj1bQzZjp2aYjRbhMGKkodD4uSyc/gO1OVWzAYj3p1GQaKZxggBzJYz87xp1ZWWpVoTrNFVPCTW1FkfmaJIml4Oo8WM4vxqIZeTKmJdtol+VhxSmeo5e7zt6JynPitr18HM/LdUNKKrdnpuhEtDF43kYIoFI++SzKVZ1uD+J2EnlwS5JYoJ7ne2RpQpKljoFvcQhBYCWwvsuPsJwsS8TN8l/WyxHqnnOV2y+ZdTsKcLe20Rsmo9xd6YrlDRe52EjMsg50By1z8zstLIlUeBa3YnImvrTIEwOIzCZ3GS1dl5E6to1FbRhjOxvpRhCep7g002Bhfpkt/WMc2DfCcy+cYm4l4oFrB/mjT8xx117DiYWURg++/+4BFucaHL4Y8c+vCbr6CvvGS3zx+YTVZsJgzWfnZBUp1pFCgWdhkImGJI0Ih7ehRIbQCVJ4hf/AtlqQYpepUtgQqYGSx/5NdSb7Q6SwHo9MbzB4hFPRCUPho5lrRCw2I+IkY9tImbcemODpsyscn28SSjtbVQjiCAZ31alv6kdHGVnXdplxL6WkAoSBbrtHOQhoxIb1HoQBxL0eUWeVzVPjnD99Bj/0EZnhxHyTly+3CHxFxffwlaAVGxrdFH8oJFGQpLaKV9hRaOh7KCN4dbbFCxeaNplOZyipEd1FwuYy+/eMs69+gVdOZezf4vPS4lmYGubJkyO8d+8FHryhwuGzixw7+h0++kFIki5nj1Q4eJfg4EMhJ77Y4i3vqPGef1/jP/+3Vf7wZIdLOzQ/NFKx2AqdQSbQ0o6c8sgDC/LcMFhpndvBKNDtnrQO67VujHaQvcCBF9PM4CsrfogzTd33iLOM1BgkUHJJfMJXtBPNmlCcaKXMa8NcAjNxxnqWEhub4VGWki2DPnvGamwZKlMOPbpxQrfnohSEh3Tu5T7fejrsgSgoh4rWSpu+bXXC/jKt9Q5KCcp9ZZJYE2SSTAlW19r40iuQP0mmuXn7AIPVgJVmRLns2QxwIZ0c2ricoQSDwhOCrQNltg1XSNKUauhzea3Lc+fWyNhgMinlIZRCOTqEVNLFBYOR4GlFEsesN3rsmaogtaHTyZheStk8Umbv1j4uXVrl9RNr7B4L8coh1WrIiVcWmKobVtYSWpFguCppdVKOnV9n17YpBgcrLCz1OHFmwZpq34RTzykZ+ehbOHJy8b478Y+TdtoE0NQUOz5tjP39XCgeAsJyyf6Msp8XoUQx9kI5kq5nke1WYiuL/ZBy5HDtUjCFi62I49RNfgxpWtAO3UXjUO8u8O9N0Uw2GwgXr+CCrfKuy0jl7gFrhvQK9UhmMcD50tCCwRwryr5eW6FnmsAPiKIYgcYPPOf3ME6aaa9Bi+ROC7KqkJI0tqlfucFCKdtR9HKMOzb9ilQUGmrtoh3zC8K+CfYhJWnOIaLQRgvH3fc8+8Eqlv0S+/f4tptRnkInTo6ZWY2+H1jTWBD67iKzrCJruNGstxKmLzfZcs02rjkwSaMd88bZde6/cZyPfX2GU3Ndzq8a7t2j6Kv4tKqS4SDlA3eM8qXDMTWvQ5IZLixljAxXmeiLqIaSKEncMszgKZ9YJ6igRrq+5BhkyuUFbEislfKQBlrdiOu21Ll2S90uLZPM7aCkw1/bD5+HIckMMueDGXju3AqXV7p4SnJ4psH1Wwe4Z9cYUao5M9+gWvIxwppHk3bshBaWHio9yzeTJY+J0RFmFpdJ45hOz7Dc0kz2W3Vd1F5mz44Jvv6NjKSX0kwyXr7cYNtgmX0TNeqhJE6glWqGyop212IUpDQ2Vz2nlhrbGS91ErbUAyb6A8aqAa045cnDLXrzF5kr7eXD92Q8/dI8uw9u4qWnLnBg3ONH91+03UNouHl3jbMLTV4753HXdWWuzES01vq5/Xth/VTM6cd7vPWROn/zZ5P86m/M8cnXelzqGX55a43BwGctNWSJPdxtUrKNRcYFdOVBQjlTTWC/+KnWeNoQaZDGcsCi2EpTPU/aWbiWBJ7iiVZMpqHsWbVghKGZwRopS7FgJU5oJ6l1QhuDZ6BaUmyuB2werDLRFzAUSmoVH6k8otQWd7GwRYVlUVmgX2igXvZYbkeQpGghUSS02z36haZUcqFG2pBEKSax4+o0SQjKNgdHA1P1Mvsn+un1UvrKbvkqBL04wwusqVfoxB1MNuddCPt5LPuSNNWMVkLesn+UKNP04owXL67SizLC0B3KOkNipee2FLRjnW43YWpbjRv2hvg+rHYTLsx1uWt/jVjA3l1DHD60ZA9PT7B89jJVE7Nl5winL/SIjce28RrLzYS55ZQHHpyESsD07Dwr6zGBK3rt7tRZGjIN0iCNfQ8QzuuWg1y1Mz/7itR562Su4NO5KionOlvlaeK8RTZZ8V+z18h5Z2Eu/pGF8AeMXcAbWxzHcU4OxoWm5Xh/61fJRSpgfV35uS8cfifTNuVVeY5g4Umn7rIS8hx8Kj3w8ohI5XuOveLGBdLeqlYQIfJ9UIEqyXHVWZoViOv8Ok0S+8A8T6GTrAjgyUdFYblE5hRAuaOzWJxpY0GOaepyRyR+KMnSXMqmCgqAcjb8HIqotVWM2CVyWigikjRBYm9lk9p5tCAnwJrCGORqRqtAyywBWPmq2K20OzEnzq9yV2rYvW0Iv1RiZjlCKI8Pv3WS3/3705R8w/tvD7k812Pr5kGOnmpSVYZf/p5tfPPwMk+dbHB6tsetV/UzXDWMDPhcXooIfEGSGYRM8b2ATATo3ixS2DGewY5EcjKnrzyanYjtY1Vu3Fa3KGtn2MoyUyh7pDaYQuFm41JFmhKEAVdN9LFntEwl9FnrJJxf6NBMFnlg1zDNbo+lVorEttrxWkRGhggEQgu8QOGXAjpei7mFFdIsJYpStFHMrWuunjAEStFYXGTr+CaUw9QstlJumOrjwFjFplP60DEZA6XA5Q5YaWeSpEhllR9JmoFOEUHA3Vv68D2PwBNEGQRKMT0/yPnFRUx1F7Org9x6zRLLUT+7piI2dw9xy74qS82Q/dsFr5yu8l3XZfzD0x32bQIPn+mzXXb3S27+8RLP/Oc2L/zhPHf/5hR//JtTXPV36/zpF9b5FdHm54YUB/tLrAhpv8xO3piHbWVZ6tQ4hiSLCVSAziAxVlmUZoZu7Ko9l0JobCyhzY0wmiCD6cRwqJPhu0zrVIgiITP0PAIlGQ0C6lWfiifpL3lMDQQMVQO7e9c2N/zSao8XzyzQjjNGqj537J2wIoTMiiO00ZRDn11jFebP9agEkrt2DrF9rI8zl9Y4dnSWiesm6HVSe3AEivUkcoWYR5SklPyATGfsHq0wUA6JegnzrZSFRkTgCTYNlFDOMCqdajAzlrqgk4zAs6KKxMUNKwxlT1ELfW7a0s8r02t044TAedLyJD7LFLOeKYDG2ioTAwNokzHfSFlrJRzYXgIp6cYxQ6WUoBzSabYISz71YY+40eXQpYwtkxVKnqQZS9bacGBXHZTP2YvrrK936O8PSBI7qciyBGNy6W5c2AqkFOBgnxa0GJBm1puFjxvpS8BmueexsCIPhyIvPJxs3Y2n7DloO9x8QZ5bJ8AxtZQqzNFa2wIf3D7XBV/pwutnlYP5aFUKu1u10x5n03CKWuk6GWN0cSf4Tr1aRA8Lh+7N0swuVnUe25knoNgxSJKktkrKGfGeKpADUtgX6edmIveCs9Rq33Pzm1Q2nTCJEke/3QgZ0llW8G+0M73lmQNxL97AmKBJ0rjwhViJmnA3vBvTyNztad8AT3mFwkAUaVpWVWC5NbY1TeKILEutsdJdmIl2UmK3E5meb4POGKwoqtWQ4xdWiaMet+0b4JarBtg24jFUkgWTq+QZTpxfJiPjBx7ZyX3XjnLkYptuK6XsSbaP+Nb4pYQ1cWYZWlpaL1nHFrcCSqHlfMWJzdVodnoMhIKbtw4SpdoSZZ3JKdMGKRzK3kWaeu49FY6RFacZ2wfL7B6tMdUfcuOWft55YJSdgx69JObg1ICVgnqKJNGoksKkdrYqfVCBo82WfS5emqHe31fEHK/2BIk2CDxmrszTX1HUB/tp9hL2j1a4ZryG8n2iJKXT00jfI9UZSepYUsK6z3NInKcUQnloBylM05RemtHr9ej2EnaND9DrpgR+ynOv13j3wzXOnJ3l4Yf38PmXPA7PKoYHFeurKQ/fktERAW+93uNPvhwz3O/RWfPorAuGtvpc/eEyrdM9Dv2fOVIUP/3BOn/y0+N0E49fuRjz5UZEyanheklscdhKOaOsInGfnzAokWQpRZQKhlTa71HoKUe0M0UnbbHdgq4x/MRkmY8frHP3SIjyFP3SUA899m4ewaQJW6qC9103zv3b69y2dYAbt/QxUA6InDRbG8F8o8dT55t86Ec+wp/96X/l7re/jS8dnaHTyxAOFa88RZKYIsnvvl3DREnGt44vsnWyn6GeptlICzxHKQgIBkpkiaZUsqa4Xi+mGki2DVV59cIyn39jma8fX+LQ5QYvXGzw2uV1JDhMiZ1QBJ5AZI6C7GYmkhxnbj/vvThhpBby0J5R7tkzSuCLwkyXZo7Hl9niJtXQWV+hJBKibsapKz2iJKMWSDqtmOZSi8GyQEddIlmmtm0zrbk2aWJ47UrMXVfX8Msh5xcilB+wZ1sdIsPRk0vO2KjwgtBONFw2h9bOHe6mI0mSWlqHE95obSnKUWQX0UJamkeqXaEHRFHkTMG2k/B9u0+y7nXbcQVB4Ap85Yp0U5ydeQx5TgpBCDzH+Usd/l0bY4Gmb6Y8G8u4U55HmqboNLVdUy5DdozDzIEe88+2caQGP/AKt7q1cRkQxuKHk8RW/tmbSY46tQRevWGksRiN3JRivw9pLjVz5Fb5Jk+HNUKlLvZRk6WWGqlzQJjIufnODOdwB9ohApLISchc5GiSxE6dY/PVsyyzaXPGGmRAkCZJkZOQ5ZRbbeVuxhikZ1+7yRzd1wEj3yyFU8Iu4OMkxvckx04uoFtt+kZr7N09wvRKCpmV4knlsXdM0GwljEzUWG1mmESzuJ7RXG/R7UTsGPM5M9tlsZHgKcHOMd/Ji60HxvcUiACRdNBpgnFvaBwnRQpZpxczNVji4avGrBEvM8ytRCy3Ittp5bJ9bfBDv2BJJU4+aDInW85jMbWm0YlJtGayFpBpTVkJQje7zjJNZaSKFyq7LDb270yihKCvxMTmccaG6wWyoxEb2ommXFYcv9wmjTNGh+usd7pWMOA+O6WSj+cpSp6i5LAy0jGhkgL9YBPeMm33InaUrFFoyoFP4AmGSh79lRKiO8+FuZD+WoXdo016sWDL/j38n0+u8/f/ssTpixH1MOPRuyTKlNg6Av/yTJupEN54NuXSsZhd94XseHeN+ScjTn5hjmbP8NZbK/zNLwxz444q/+1El9+/0qanJX2lkNjY55i5L6ySNoo1TVMC3yuiEqSQtKPUFSXG7iqNrSw9T5JkmijWeECt7DFWLbOaGVqtNlu3jvG5f/4DPv+pP+Yf/+a/0Cz1cWFhDU8YWr2IU7NNmt2ITAuSFFqdmCMza/z6r/wkb7v/VhqNJv/l136OH/+xH+G5M7MEvvVZRIl2aiLBzdvqdFM4dGkdaRJev9KkHkqiRpdKNaRc8mxueSfF8+zFAdDLMvaP93Fsvs2hmRZXjZZ4y75h3nrVKFePV9lUL1u/lxR4Qe7ehgyN0U66nKWgwFeCwJNOju5iWQ1M9QdFsZinMBqkkzwL4jjjmh191GuKTAhOznQhjXnt2ColPyVpR2RxTDgyytbbrmHu+Bzj/T5fP5uyeaLCtbuGyFC8frbBpsk+du0YxKSac9NrlHxnsnTo9Xy8o3GdhucVGSw2PdCen6nb/0opbaZGYv0rnm+LxbyQE844mDnqdg4pFG7hbZztIM/fSLONSOUcUwSiIB1YNJUpWFiWE6iLcznvHgy6kOtbEKjrZnIlnHGEafezWZo5QnTmorS1i951/ggrspV4gWe5S0oWyGXlKeIkwg98ctmv53kIoQi8wKaipe5QlqK4tfIQ9jiKXe60vfmkU6V4vufMR2zM1/NI0Bxf4ZQFOXVSKQcc8zybK6G1c1hbrYJODb7vjIH5m+oQzBvYdulgjQ6Z7LqsPCAmnxMa48YMWiNRBL7i7KV1FpfaQMb1+0Y5cX6d5bWILElYWulSLyvakV3SX7rSplIrUQvhL74yw+pyg9t29dGKMhZbGZVQsXvcpxRIklQX8aQon6jbtlWiawTzTiQzMFINuGN7nVLgEwjBifkWXzg0w7dPLRNlxuYX53rvJMN3zDDPYWsSbVUUGO34U6oABcbG5r8rIekr+S5lT9A/WrELbSFt1Zpqssg+p1LFZ36hQRrHlHzFWithvZmSJBnd2FBRKSNDA7SjhFJgP7BRbHj25BxfPjzNJ184w/Ryk5JvOw0pBSXfc0wpG32bptZAiDFcWenxqWfP8olnTvPs8WkGKh7DA300my0WmymvnvT5we+p8vy3TvPwu6/jCy8IXnhuhS9/J6bVSdm7OeLANSk7+qs0uppXL3SpS8FLj8ccfynmtn9TZ9N9fZz+TJv1cxHtXsa+rQF/8bP9/MTDAzwxn/ILp1c41IwZ9jw8aSW3wnlEPJdTk+nM5pH7EtA0OplTlll/i/1KWFqrhcQLqiWfipBcaEVc7GRII/i5n/lBXj92ih/7qd9gcGCA3/2tX+TIlXV6ieZixxDW+mlGhlara6OQO1327r8KpTMeePeP8pEf/3Xe+6Gf4iMffDeDk5tYaXcswj6FbpIx1l9ix2CFY1caHJiqsXOkjzSF0PNIWik6g1I1oLHWpur5GGHodhLiJKMWeoxWSwyWJG/fO8Q1UwNMDZSY6Au4ZesAO4dqzqOQL3udAEHb409JR6/VgunlNtMrHZffY8exNvtdMFkvk6bG4ZWE3Uc4Ay1ScGBrDd8I2u2YY9M93nnrEMvrgsWFmBJ2/FefGuHcsydQzTUuJCHPnMv44Uc2o1VAL0o4cbHF3m19VPp85hfXeeP0go1JsKMZe2C6YlUpie8ClvK9SOY69nwik58haWKl6Z7noRNDltizUCe5m9zuJbW2fCtj7JmWJilxElsPTOrOYpUnsLIROe58cCpPpUwpOg0bb+wVIiYEBTgvTVLXHbnRlRIkSVKsCoQUhS9OCBvBK4Q7tx2AV+Z8ec+Zc+xtlBUE23w3oVyQCU59ZSsuXaDehbBKnzTLEHmE45suiiDw7J/hzCsq8Gxr5b5QOe7AJmBldknslu8WYeI6FLfzyOW5hUPd2CWQ73sFpgOpCtSx7220XSo/lV32upSy+PNzV7HvFE7W/W3nnZ4SzM63uDjbBl9x7d4RVtcTFhsxCNfOlz1S4aPjlHajw8hojYGyx2Q54jc/cZr5lRZTQyFnZ9rINGLLqGKkHmKEzSCP4sRprxOHLbHZEL7nF2PGO3cNFSH3R2YaPH16mW1jFe7aNWyRINr+e3JD5GEz411Eq++5Dzc2+Q9stWyrQful9JVix2jVKkiMsA5i6S5xX4ISKEcpmJ1ZBCXtF8RYiGUjNiTakCYxtVLKpokhVlpdMvcePXHqCn3btvG7v/9r/Pwv/xiHViLOL61RCe34zmBhcEIIosjmxAS+5MJik1fmO/zoT/0wv/vf/yPx4CTPnJxhuL9Ctxcx1Gd46rWAm24aZjBY5tRsyuidN3B6OWLlUsI3j0CUSt7+UImBbYI9gyWeOhUT+RkHdkpOvJTw+isZN36oTmU8oDndQyrJ+nKPLDb8xx+s8f9+doxS6PG7Jxb4+HKMVB4DoY/0ZAEPFELZYCElHH/M0HM4FANEWYqvFBJBasBzks9KIKmXfU63ExbaEdcf3EWSpPzqf/pTnn32CD/zC7/DUP8Am7Zt5+xSC2EUb79hGwPVACOt1LQXJwyPTfDNp17GM4Ybdk9x4eQp/uTP/h+9zGO1mdoRIdawpl14UJbEzC51WWnGVAJLx1W+7aiS1B6GoYttkO5zOVL1qYaKesnmkqy1Y6I0oxenxKnt3APPdtZFdLSLoDUu3tWTHlGc8vJMiyfPrPHs+TXnd7Ej7zg17BvrI/TtLshOFBJrB1CSgUrIlkGPUtlnuSdYbWvuuLqfG/YEfPbrl+m0I1Tgcea5I4jmCt2BAf7X0wkfvGeCzVN9yECx3Io5P59w2w2TECrOX2oyt9iiXA4chcGhO9wBnsTpm7oGB0NUnhWvsJH+Z/IL0z0zO6p2I24BAuX2CroQ/+TCJv2mrqHIBHLPULudUS5WyDPTPc9zRksrOvEDR/bNi+8cReJ8cTpvCHyFdjDTHPFkeYZeMSHJOyfrQ3GRt6IYMWVFuLrWupCMSgeNy0m7uIWedvNpjCkS/grIWI51d0seIXFYEoESVlqrnaktcYa+PGTexsmqgvwqlXPlYkcvMucJaYNxPo1cTZDztqwRKd/NyCLIx7Lv7e/lu6WuyDsm1/VIKR0basMvYpzqSSlBu9vj1MV1UB4Hr56kWguZXUmsaTHNiHopnSil3YqIYxioCFIE12/1+YF7h3jsUJP19ZjplZgog/5AsHPcI+6lkBiHtw/QSRfpkDJ2OZahkYz3BZR9hRKC6ZUOz59b4/ZtA9y/Z5jxfhsTaqeM8k0LN2v6QTvXtEuHsheTAJOhpCH0FaGEagChLxgu22e0f6KGOLtGa6VLECqrHdcGzxOUyj6l4ZA0SanWykRxihGSXmbbc6E1uttiZKhOM9aUw4BTl5fYtHMH//Q3f8DbvusRfugj38vH/u8fcq6R0W53SWNb9Uo3QrNeEDuqOLXU4a/++Ld41zvuJ840f/G/fofy+Fb6fQfHbC0z1ww5cirjJ/7NGCefP8aO227guJiiZJo89libpZag0SzzPe/wGN7uM1EK+ObRNmFJcv1Vgi9/aYEf/pVl1lpd1GyPXjdBhTakKOpJHr6pxsf/8xS79x7gsZkWv3ZiiVeaMTUpKAWKzBhwKhphrFqzm1gchMgzG7AXvVCWB5ak9pCqlWzex2utDB2nTI0Pc/jwCfoCybW7JojXV/jjP/hf1DzNaNVjbrXF73zuVS4udRBC0Yo09b4apw6/SvPyKd5zyxb2jpS5Z/8kh5/6JvHaopODQuhLlIOmam2jg9d7trMY6w+YWe8iPEMWJWSxplQu0e7EDj1kfV/jfSG+svSH1I1ostTO8XOVkTaWFJtphwJ3i44otQvlzBh8T3Lz5n52jla4stLl6XOrxElKKbRm5UDaHZHNxdBkxkqSu7Fmsl+ydcgnjTJOz3Sphh5D1YBbDvbRyQRvTHcJPMVYX8hMFvC/X5Tcs7vKpkFFlGQok3H2Spt2z3Dd/kHwBIePz7Oy2irMunmEbOaS/nzfKwpkbdIicMnuQSS+g5xKt2xWnnIj+w3ulGsSnEHRJrQmaeZAhxvnrXarAkvotcozpZTl/rmuyDjsTh6qpR3FPIkdDiZLXbAfxQ7Hhl1Zv4oxELgYcZvFnhbRzn7gWym125nEcWSFScKmt9tqQvMmA4wqFoNxFDsmi4NtuT2GENKRMkXBA1Iu/ClNEmdYkS5b14IVdZbZOaiwIxDl5oc623CQ54qoLAd2udGX1rYVS7MNX4fy7LwuSV2GgLbjKCkcvdJokiSxkj9llT2eHxQkTT/wHQdG2UvMwSOTOHEzQFOMs+xIzUInj55YgG7MxHCZibE+Dp1eo1otUy0HrDd6JHGG9Dxkxd7wA2XN0mrMzi2D/OqHdvOWm4Y5cSUm0RLf97lhRwUxqOi7bpBgKKA6UHcuVvs7K+kyVYyh0UvpxCmQcWyuxY1bB7l55xC9rt3FYOySUUiLPEgcWsPzpTMibsT9Wt6idvkVmoVWxKmFNpdXbKdQr5YYrXncuXuIhw+MIs43SFNjTaEuMyaNU0b2jtDLeiwvrFsWmoHllk0iLJd91peW2DQ6SKIt/mO1l/KhD76HyuAYf/Yn/499Bx7i8cef5Z777+XSYqPIno7TjMxA6AcIDJfmlrn11pvYs3Mn7/3gz/KRH/5FPvaxf+Ytb3krge6wY7zOxelFyqUKn/lqxL13lvieW3qszjYZue1WHj8bsyXr8C/fhFoQkaRVvue7S+w/UKa5WOLVszG1IGNlaBensyqfbUlefinm7McXaHbBeCFZTxH1OpxcHiPY/2723vluVlWNPzyzxJ/PRqxryUig8D1J2uthEntgtGLbvafaekKElGQIG9lrDMqzu4+KlMzFhiONHvVqiYsXpokvnuCRA2PcuLnGffsnGNFrbAsj6rUKN02W2DuoGK8qyqGiHFh56NUjHuMly0gq+xCIlG0jfdy8uY/+UNLqxNa8Km0OizSwY3KAqaEqm0f7iJOUmXaX+lgNnRpMqqnWyviDId2u7bgrvmSyXqEbuzTGTDvng0vydGoiUSiM3HhVQhzHVk31pqX4jqEyd+8c5P69w/Qyw+GZJoFQeBgWmjGtXlZkhwdeQJakxInm6h1VSp4hTVJev9hh71TAxHCJVPZz1/WDXF5LUVrzxIzi7w9l/Mj9Q9y42WNhvYenrE/jjUttxkar7Ns5BFHKkTfmUMq58x2l26qdZBGyZ9hI+NPGFgW2aLb7Oz+w1bt23UXevVgJrCpYUxYlopyM347A8m7E7lOsA18qQRLFG+iYvCtxqwObx5S5SYs9W4S0u+BcdSrdBMlzqlu7m9auULXiivzik8IW8kmSuKX8RrSGcO+ntLG0FETIHJuRwwiV81VYoJdTm7iWyXNGPQFkiX1zBTZXIic7JklqnaOZfaCZw6dLudGOer4q2POWLaNdgImjBLulkTbaPgRHo5UOamcXUvYyELm5RtkHGARBYfjKX1OOM85hi7kh0e5WZCHp1Q59DhBHkeVq+Yojx+boNnuUB6vccHCcE5e7qDCgEipSnDY+NSRxRqed0Vf26fUyllfWyLyQB6/pZ2k14spyih+G7NnkMXGgjhgKC5qnyZIiCjjL0sLk0+gmGKAR23vths19pBqCQLoPtH1mScJGaiTWDZs6EJ2xbzyeW57FqeGJ06t86cgCT51e4akza3zt5BIzKy0m+0IeOzLLkYvrXDVSYf3iGqWKVZ9YBLRNlBs5MEptsIo0btYrbFZGnGqSpEd/NcQon/W2lT3v2rmDtbk5PvZ3nyHqaD7z6S8Ter7LhLeHSpJqPGGYWV5HSVuxbt22i5deO8WZC3Ns3rqNLz72OM89+wy1UqXI2m6tLXJmZoDnX2vywfvbsDzD2PYdLG+7nlTFrB6LeepoyvAAxInkPd+ruOPGGq++EvNXz/fxenYj26+5jnN+jX+KE86c6XHxcyu0GhGy4hGnPl88toMkWkf3TbDr/u9ncuf1PLXY5FePLfKpuQ4KGC4HeKFHlNgERiktNTfOtMvmsEZCiUBnUAsFAxWPY82I1cjKLTcHCUEa2wVqmhJFEcoL6EW2gykFkp2jJfrLkiwx9IUetVDSiDRz6zFGWy6dcBwlgaES+gyUPXzP7rN8H6Qn2FQvs3u4Rhgonn59kb7Ndft5dEWDMoK19R5nphcQUjI+UGag7JEhaEUZcaZJo5SBSomSF1iOletCLNlCk2rjkvGEw5WD79nuJ3aO/12jNR7YVefsUovFVs+OBh1WP4cKWgKFwVeG/ZsCwopHM8o4t5hy8xZFNRSUB2ps8hNSv8zfny9xqSX5/Q/v4LbdVZbWu1bSjsYowSsnm1y1c4jRsRrt9YSjJ+YohX4RwSycOtAa44xD+blL0QE0c2Cmya0AzkqgHOrFWhTsaNxScd+kpEI7JljO/KJQo/q+X7jHhctDsqZFG+NtdGoN025XnLis5MyN8qXyiomS53sWsgh2PCVcVIGkWKbnpmXlVHK5BSOJ4+ISEe58lJ5nX1zmCJHCSWFxyXBW9qqcp8JKZpMkIfB9MocDTtIUL/QdStwUM7z8UlIO1JUkbqxl8rAUadVdqXbtnS4uMVyqVz5OsxDHwN6URrvWz1boYf7/a3vBKGmpu7jbNadX5uqxNInd/NIU0MccyGi0y6c22uFUfIeLdxW8Epw4s8zMQgcqZW6/YTOnL7do9TK2bxpgvpHi+zZUqSQzwpJCeZIuPtFaE9FrsG/7MFtHSxyf7RK3GkyWBVtLiijRLnbXoHVSjK4yB5c0QMkT9Iceq62Ya7b02yS/yGrTLeTPYWPcF9fK/ijyU4x2VY5UZJkg9D0eP77IzGqH6zb38eBVw1y/pY+KgmogGa2FtGJNp5dSLSnEeo9MWy6SzMGGUYrvdkftKEUiWO9oUAqERxRF9JUF1WoZnWXUQ8ncwhJ4HpmRjI0MUKmUibttSr4VFChh8S2L7YT5yOPMUoep4RqHXn2Jb3/p0zx0YIIbpqrcOBqyP5sjSQ3Ti20Cobl45iIDQ/187DMJ+3Z43DK6xMpKwt6HHuDz8wNcM9bh4/8Ys7gWoeOIJPZ45/sDygMlvvhaQtBeoFSqsWXXtcyoMl9J4dyxdU5+dolodoWXLw1ybL5OxUvQaZd2N6ay92523/NevPokn1zs8Z/ONXl6LcbXdie4FKV008xe8MJmfmdo3KoOjaGkFFIbXm+lZMYwOeCzc6yPRi+zgVsZZCmkSUZQVmhhKPk+662EXpxR9g1zjZjnLnc4NJ8w3fN5dnqdi2sRb8x1efbsKq9cafGdi+s8dX6dlU6K7/ZkaS9GGs1anPLFFy6RDoVs2j9O1s1IelYxeO7IFYbihEdv3ETJV2yul1jvpXz5yAxfeX2Bz702yzdPL3Fqdh2TWepyHrWA0ZQCvwAp2iAlG6ssXLa3kjYwtRMnbK6XuWlrnZU4I4oTBkuKoZpVIGapHZFnBsYHPXZvqpC0I05catOJUg5Mlpg+Nk/7zHFmLi6RGsn92+AXHx5gdHyA6bNz6MwKPJQQLK3HXJiNuOPGcSh5nDm/wvnpVYJgYx/h+Z7jQCm7N3DOeBthm+JJr+goPE8VCtB8XKfTHJq4EQmQH+rWTqBt/LOyeHyEKGS9Gwmtuhg34UZVyvmkDMaiqMzG0jtN0yLiAeMK4lRbYKvzhaSOVIERhUjA7scESZTi+6pAtih31to1gxVAqXK58tGcW68cM8rOnB2F0UW/Ksfhz/lR+WIny70hylbeQkonTxNuJ2JNP55bvOsi3HOD0ivcQ/U8z/Ju3F5Fu3mdygPiHXreU/YBCwkmy4qo2iIUyqVnSekVvhXtUnssi98hqH3PhqQY8a+gYvmMcAMfovGD0F6cgcdqo8edN0xx1d5hRJzyN58+ym37B6lXBI+/NM/uEStK8KWhXlM8frLHi9OaRFtZ546xMnGqefZ4g4euq6OUz1qUcqiVYdY6JE0FJkJox/RSFl4Zp5qRvoAdwxZPUQ/ts0q1xggbfZsmGs+zo5F8pJflUTfGeis8334pykpwfLbNlUbMo1cPs2u8xmDFY6gasmMwoK/kI4Ezix0G6yETfSFLjZh0sOKAf3nXplifbbJ2donbh0Ome5rAE+wZFEwvJWweDti2YytPvjpNLZAYqTh29hJjI3WeeOo5Wp0eSRRR6y0zWPKJk4yyr1htR9Trg7ztxj08ffwSfYHHzqCLaq+ytV5m01CZ8ZpP4Hk8c6XFyHCV992znydfu0hfv0+zN8T2qR637Ut47MUAf2ScjinB4kX2jYd85ZWE9z7k004Ei6cTFh+PuWK6rEZdygPD1Kr9eLLM61cuEyvDpjilPRdz7T2a/r4WL1+sk2Yh5dDuEEW1j4FtB+nvrzO/NM+TCy3O9AxlrQmMtHJzkWEyW3EbYT/j9oBSbB4MMMrjE3Nt1pOMq8YrjFRDklTjB4rUpRamTtHUjVKOzLY4Otfl5FyLxcgw04XNE6Ncv38bQwM1UhSnF1tU6oNMTo5R7atRH6yzdeskz5+Zx3dFxcm5JqcWu5xZiqhtHWT7dZvIYmfuE4JOswezTa7fMcDcchdjMvaN1VjuxHTjjP0TfWweLKON4MiVJkLA5noZre18XmMTJbUzsW1EudrJq5T2EhXus55pwXAlYKBk95GZEZxdbBEltovyfY9WJ+PmPTXuvbYOAr7w4grn57o8sDugX3a5eHyWxabkwCRcta3M6L4tzF1ewo9Tzq9L+uplrt8/wsvHl3nsxWV+/kdvZvveCR5/4hyf/eobBKFncR7uPbJ+EM/hWVSRBKhdRIWLk7LqQXeeOpymk9S7aFkHgM3HVAi7CFdCFrkdmbEy27zSL5SpuTxX5Tkfokh8tSFnNgoXYUnnBoNJM4t/RxZNgTWw2p8JQouNyrS9oDIHl/VytqC7MLTOXIJhzvsCValWP/rmDHHbUtnKXXo2d8IL7KgqByZKzy5n8qq88EyoDU0x7obO3NhHuwAdTyk83+r/i4WaxEV92gwA+4KtTM7k3oX8FpTSpQwKl2Bn9xzS/b3WiOPbvYZr+axZ0rWETm1m57Febj63oYU53jh3uCOKjkg49YLnezRbPbZP9PHQvTuoBYrPfvUkRmvecdsQn312kU39KaEH2yarLLVTPvFij/dfI+j1El6bgaePrmFMxokrHW7cM0D/QJUgy3huxUn+1hVhtULUadllpLS7p8xArRywY6Ti+GDOdKesPt3zVCF/9TwnPigigG2X4Ac26UwJgTSG1+Y63Lq9n6FKQJTYL3nq0C4GqASSmUZEsxWzbbjKSjshHaxS7QvdfNZiOeYvrjKWZtzbX2Yhy1jRihtG4cKaoVot8a2XL/D6mXlOXlnn4kqbxYUVDj/zLPds62NLf5mdg2XqgSx8IJm2z3thvcPTp2dZbHTZMeBZWnAeOlQK6CuFPHdxndcurfBv33c3//au/XzrpZOcXWgwtWmK82fX+dEPVJi+0uN0Y4KJ3Vt4/LlL3Faf5+QZQ+TBnkl49a8lopMihOJMt207yIEBqvVRPL/EyZUFTAbX7g/YeZ3PDZNr3LBlkem1KtONGp5nUR1JnFEaGGFw10GmRIsLM7N8dd2wnmUMeYJBZTlQcZLaxaSQGKOolBTbh0ucb6d8brYL2rBjpEzVV3TjrAgN8jwsc6zb47lLbUrVPu64bhejw3XizOOWa3aSZYbnjp7h2NkrhJ7k4O7N1MohC0vrLK026CUp0sCuTWMcnVnj0lqPuWZEaayf3ffuZHRbnbSdEvcSvNDOy5curbG34jO91Ob8fINqoNg+VGGg5LFvvMZoNWDzUJndwxXq1YAXL6wzEAjqFY/MQWdzXp4Q2KxvF10gpaUWeNIKIZSSLmDMpee5MfW5lci5wO3YrxOnfPcdg+wcDVhrRvz9E0vcuL3E/EqP/SOaheWMatkwMlEhLVepVANmz65Rr2mePtXluv1DTA0GfO2VVWYbgv/wU3dRCQP+7GMv8srhKwz0l2x0gksnzTOTpLJnkK3cndeiyEL3XI6NLkbIVk3pzjh3KG+Il/Kdbh6yZfeTORgjl+6aPCvHeeVyl7lLNbMdipvA5Gdm/ucY53rPLyZjcmd7UCS5SuGQ+km2sVMxLiGyODNFYdK21AuJ6h/o/2geMyuENcJIaZVAxR/ucsILRQD2tnTP04WvK+I4KXwduQ7Yzs5smJTn+Q6G6G5Cd3vi7PxCSJf94UBkOW9ISptu6MxaltibORKy1ZTnOezCReAK92Yo94G1Pg9TRDMqqVyqmHVo23wQF77i1Ax59ZV3XXmgvUDi+T4feNtVlPrLvPLqLMdOLfBD79rFxbkupy81uHlHhdGa5FMvNRmpGO7d4SOSiPc+sptrrt7MyUtNjpxbZ+tomW2DARUfjq0lNIIQs9ClXB2ms75kMbXG2IWWS7Cb6vOt0cp5N5STYc+udbmwEnN+pcdaJ7YGPc9lSwtFnCWF2z9KbHtaDSW1wCakaWMvm0DlRkSolDzWuhmZztg/XuPiSo9ksIIXyMI5LFPJ/MU19qMZKYXMphkLkaYVGeYaGa+cb3PkzDy33biL//xL38Nt12xhZbnJG5caiCDkwGi5SFeTwrKCfNdFGaNptHvsHgwYqwUInVHyFZ7yCLXmcivlG2fWuGX/Zn7lQw+w/Llvcml5nVdWYmqhoRMPsG2r5vaDkq+9FrJS3swNmzQnX7vAw9tSPvUdxVvuK9Nd9bh4psXmQY9Oz3ChuYbnh8hyP/XBCdKox7HFBj/wtgqjdcHiqmZTf8Jb9s/R57c5tTpKMylT8SI80wNP8cjSEe4o2QLmUCvmO52YlSSj6gnqvsI39ndMtWG032OspHh6NeW51ZT6UJnhqsdgyXfVaZ6ppfAkvDTT4fYbDzAy1M/FxRaT40NcuDzHG+dnmF1dI45TAs+j2e5yYWaJS1eWWGt2abV7rDdbXJpfY73Z5q7r97F9ywh95YDl5Sbzl1ZorkdUB0K3KzQILWjMNdjdH7DSTen0YiaGakz22Vz7JDNEmQ1PiuKU0VqJctnj4kqP3aM1awB2xV+SZW7VTEGWjTU0uhmeFJSDPIvefoej1E4vEmM4PtewnbZjNNVrPt971xAlUs7MxXzhhWV+7f2b8ZTk7Nl1BpSlNvilMv3jQyzMt8i6PbvzuJDwvke2YxD8xRemuXr/JD/w3qtZWerwP//qBVrduBDPCCXdlEJsYFTERmqfzIVEbj9cXBJGFKMf7VzdKp+euKwOKXMSSFqMv1R+3jhZsHZCl3ySkmeuC+dCN3nkeGYLXXsBqA36bnGG2bNRSYl0nEErgKLYK+eFuXjT7yKdClY6O0WSpLYDyzSqWql8NFcTeDm212GAUyf7ytP+jLGtaOgHhT/EapjtzD0PgMpzOjyHDzHCOGCYLhQJ+c+ZNy2sLVNfFCYZcEbANCvSETOXF6JzbXRqH67nWyaNDVrJ+VYU0DFLpLX7GwtdFEWyHAiUbwGKnudUYy6XG6fyEu7CTNPUyh2bMd913zaGR0MWl3p88svHeOCmcQ5srfCVF1bZXI6p1xRPnTM8uk/iYyiXJIsdya7dE9x1sM7iaodXTje5d38Nk6TEseZoHKLaPeK1BI3trFJXUQSBohOllEOPsarvzJoQJRmvTK/z6qV15hs9Ep2x3k04udDGU4Khmr+RcSIMl9d6XG7EZMYw0Rc4cKbzhZgNjb4xBiUMrVijpWLXcIlDC22CiT677Eu1jVf1Fc3La8Sx5o0kY7odETVbXFqNaHUSEJp9uyb4H//+7exYXeCaqzfz7u+5DxHHfPnZkyxEsHckpFZSlAIPhV2wdqKMaqjYM1phsFrCV4p6ObB4egSXOpqvvLHAQCj5y9/+CNnhN1g9eoIlrXitJ9G9iPGpCU6dWedHf2SIM2+0qffW+cvvO4dfrvDlJxvsrGpMX4X3/oDHpQsZ8+cztpYN60ax0FxHGo/+yUn6h8YhWaXbWOKqyYyhwNBMQnSmuGVnm3u2XiExhrMLfXSDfnbOvMLelfPIUo1rSorb+gNqnuCldsTz7ZjZJKXiSQZRlAPJ1HiZasXj01daTK+lbL1hHG0k4wKEMC6lUpBlCS9dXmehrYmilCNnL7OwuMKJc5dJtTvMDAXlWghrqLWRrDY7xvMkSkqiNGNxeZ3TlxaolkJu2LOFft/Hj+HipSXGdg6hkHTjlM78OlvKIfONNq1exu7JAYaqgQWsOu5anoaZZJqx/hJjtZBekrHS7FELfTu2yUN/hWSlHfHMhRWOzrU5s9jmUqNHM8oYrQXFCN3zFdr5kOZaPbqR9Sx1E8PNuyo8cHWV0Bf807fn6cWGt982zM5tdTrrPUyvQ7XfZ2THJIMTIxx+4RL7dgQ8ebxHbbDG/TeOML0Q8bGvzfATH7mZ62/dzOHDC/z5x1+26zss9sjG6srCx6bdBZJfFNolEmrn01DKSpnt47DKKE8pjBCkqYWSSod/NyKPCXdpqK6I9/x89O85i0PuLxJF5LdNPEwL06J401I+TynM0gzhOaS7MBuEXxz1Ix/duy5HuLPAdjN+sXfB0YbzyGarEFWo/v7+j+bzfs+zPgq7OM6KHYAxEJZKxUIsjWP8ILDdgqNOCjcuyXPMEc5oJihwwbZTsRpqIS1ULZ9jWSQyKLfU99+0GLc/uzEmyxdPCOFmdtphMFxtk18cAoeiVwUSPt8y5WM2gbQVhhth4SoN45K/8rhb40KzlOcRhB7Lqx1uv36KA9dPUg3LfPyzR9mzucqDd+0k6nZ59tgaw/2Ky2sZ9+wpk2oIVcrcSkR9uJ92N2N0UPG5Z+a4/epBhodrlEh4ZUUSeZL103N4QWjh80YgZf7GwVjNY7QaYoxgsZVwerHJ8Zk1tg1XuGfPMAcm+9g7VmX7cJlK6BEqewgZJO0ooVYNeN9dB3jx7AIiSxmshU79lrPPNBphP5wI2rGhnRhEariIZHj7AN12akF20qAjzfqFNVZizXKzzehgjUfuv4af/qGHuOHarTz1ndd599tu4R7Z5djffY32uVn8QPDo974Dej2++sIJGsbj8sIaVxbWWGp0WVrv0ujaTmqh2WN6qc30couzqz3OtxXHliNOrkSknQ5//NsfYY+MOP2pbwKS082UV3sCHUWUqhUazQoLs/NcP5nygYNzNLqSW67r4/yy4ktPLjHe7/HQ7SGT+zVXLmlac5K9Fc3FLrQ767aQGJ3kh0vHMXMd/vyQx3At4uotZTLp0ehoKmHC3dvnuXpsmXMXWuw8cZJ+XyAdIij0BPsqPjfUPAJfcLiR8fxSyoU0YbjPZ/9oH0oFfHx6nUhrdtwyQaQk6ZU1+qsBAsViu8Ursz0yr8rVOzYxOjLI2lrLAvQcCgZs/DI54DFxHgJj95rajVeElIRBSBInDh2jyQykGCbHB1hfaNLoxvQPV8gMtBebbK4GtDoxyoPJehnPYTikZxeqmbGSVN93SiCteX2+hfE84lQTKvs6JJJWkvLKTIN6yWNzX8hgxSPSML3aZa4RkWWGwapPkmYkGjwhacYpK92McuCRJPD9D4ywbThkpRHzN99a4sDWCvdcP4yQgrXLSyzPNJjYNsjglhFe+vYZhvsyBkdqfOLpNX7g3XsYHSzx/MkG3znW4td+5i5GRmr882ff4MuPH6e/v2J5Z70eYRg4xIj1UgSBjxKOcusYUzZQTzrIicAYi/5Q0k4rtNt3GGMvgMzeNkXcrXV+Q+aYVpbHZ8GvbkliF9tplsepO+K2lecmUeKyzD1HL7dFu20OtPVuYRWdeXhePj1KU5uKmkc9KD9vJHRhZ8gx8oV9I7/IqtXqR3NwYO5c1jkvRVvdtnJwpSzTBEHgFigb2b9KWWQx7iHljkf0BiGyaLvc3P7NAVA4w5906OHMOd7zONqcHpnD1ITBdRDSqQ30m4JccEFL9gEFfuA8KVaenGkrSfR8z3Yjrr3L08Ryea90HYyncqxLrtiyi/Z2N2Z4qMbb79vKQODxya+coNPucNd1U+ycrHDkQoN/eW6Nm7ZKhioe9dE+lpbaTC8kNHuGnVv7mRqv8vUX5/F8j4ObQ6qh4uRCxHKlH9VLSFsZfmgvYT/wCyPldVM1KqEH0ibazaz3GKz43L93mIFS4BIEM/rCgFCCEfZyiKLMSqmTjFMza1xYarJ9pIInbPWYZBkVtzhPMk0YuqpKSU7MtbmyFjF89bhj7EgwGUEYsPDGHOuLbUYHS3zf227g1z9yDw+Nh4wuLLK76vGlo5e58/odbFtcpNWz3V168gxxJ+Ke730Hx46d5zuHz/HQu97J+K6ryKqDJJVBxMAYWXWY8thm6lt3cvD2O7jhjjt5/vDrtDoxpt3m9/7DB3hw/wTH//LTBSPoaCvleAIlKUl6KeWKx7CX8o7b7eJTCYgyxTsf6OPQmZinnmtww1WCyVHDtmsN81cEesFwoJJwiRLrjVWW5ma5Qc9z72SZydTw94fh9ErCjTs1Q8MBGp9GI2Pnlh6lo9OYi5qg6qGwDDKEopcm+LHHI++v894fGKWUehyfT/na5YhDqxHrxvDaYoRX9pk6OEq5FnDpcpN+YD1OOLFm2L9nG76SXJ5bZK3ZotXpFgZFIS0RWLmDA+dMNjpDCLd7cLk/QuCkpKLYW84srLKwss7CSpMtm0Y5fmiaKE6Y2DXC8nyTSd+nGnost3tM9Qd2jEGOCnIIjsxyvkq+x+nFDv31Gv/27bdw/sqiwwaZ4rXtHa6wZ6TGVD1k+3CFLQMlhMk4u9QmEx7jfQG+S/xTwnB6qUeSQhRlbBnx+MAdg4SBx2vn23zhhVUma3DXvippHHPx+BV6zRY7b9rJpVMreHGbW24c58+/vMTU5gHeevsYnu/zfz5/joGRQX7iwzdCYvjD//ciFy+vuDGScOmCFGF6uHNMu0W5MVZhhbDIn40zKc0joYqI7MyNrIQbj9sCPS3kwDirQR55Ycfzovjf0gmU8i4jnxxZ5qApdiJGGMfmku7zLjFOOizyM/dNxbhSnuPO2VXFBgndYeR15sZXyu1XRFGEq0ql+lEwxSGr3cbL932biS42dMHCeTdy3pV0+I/UtXh+GBLHsZOu4cx8dg8ipXJtlMWH5DsOHMMq/6W0U2ypNy/Bi0WPjcHdcFvbf98PfMt8CuwCK3/9Io9sdL9fkib4Lgw+cyA042IlczWZZe2oIhPFIo7Nxu/umDfKk/S6KR94y07Kg31cuLDO158+x3vuGieo9LF1IuTTz8xx62bN0ECZkdGQ4+c6RKlkcaVLR0s2DQa0uynfPrLKozcP4wnBWiPi6KpgYKSP7nzbLqwym06nPIVGMlj2qZd9FyVs1VnXbupHAW0H1fM9RZylNvpXYw2JnkX0z611OXpllb1jNabqFVInoQyU4MWz8xyebrBrrN8mTRroppoLKxGD+0aojFRslCbO19NNuXRkhmv3jvKXv/uD3FXSNL7yLAtHTtO5NM/p187zpaWI26/azOaVVdZW2gRKsRJrpl8/z9RAmYkb9vPJLz7FNQf28z9+/6M8/PDdPPrIg7z3u9/BO97xFt7zrkd59JF7eOjBe/n433+Sl146RCA9fusX3s07Dm7m2P/5FFGrY3d0iebplZizsWCiPySJNVtG4A9+YRPN1QiiLkHJEMXw9dPb+PHv8vj0V9c5djzlLXdp+qshOw+kTE9DsKa4vpbxags2Ta7z0oLB7yXsGQjYTEJ1W5kvHIakF7N7c0oQCKZXPN74imTQ1y4fR5Jpp6RxVNar3zvM9Qf6uf/OMm+/s8reqT5OzmgeO9VESfszYX+J2mAZv17m3JllGonkpuv2ce7SPMdOTdNLItqdnj1YlCg6R6W8YkYPVh4vhCyiD9IkJfSDIhkPt5uwYVhOcJEkVKol9mwfR0aa2elFKmMVli6usXfTAJ0kY2qoYserqXY/D770SIztZrQRJMbwxkyTlUaTitzwUGis+95zUQtISSeyF9C2oQqBJ5lrxoxXLLo+8BTTK13OrnSphB6d2PDQNTWu2xLg+4p/emaFpNdjMozBaPbt7OP0G4uMDYe0EwHr62zeUuXPv7XGxSb84g8ewAt8Op2Iv/rSJd711v08eO8uTp5e4vf+/FlLFRcbZ1hx6OamPTdWzrFPxgl1MKL4d3IqeJpmxZLaDwIHMM2KTsH3fZC4rJGsEPDoLEMJZXOOtM4bSvfP7fhe5+mrWiOFQ4t49swMAp80sWd2vsPNoy6k86yIotB2OCqXN5+5FNhcaavceZjHaNj4W4t4Uv191Y9ibOqeEt4GOEtY7LBxYxtZzOpE0R7nlb5w0rPEpVTlub75yMpTHmmWunYoK35cCOlMRRTUXPvisqIrMdq+NoRwRkWcIsIZs4zFmAilnLmHDYSH85XYNtC8KegnK7JLrATZK5b+ORIg/1kb2+s7XDLOEyPxFKw2ejxw52627J2iIj3++lOHeeDmEXZsHabbTXnmtSvsHYLA96mUJMdPr9Lf51MLDKuLa3zzWIdmBEcutLh6S43BiqLPM7wyG5NVKnieJFrukmFzMhKL1aWTarYNVgrA5MRAxVYyBetLYVyglC8lqbYLHe18Mf3VgF0jFfoDj76qjydgeqHBEyfm2XvTHVx74wHmLpxmsBaijWatnbCEYPzAOEnX7k6yJCWsBMTrMdPHZ/iJf/MQt28Z5fU//STNbgZhgKd8Opnkqys9NveFXEdMksZIA5c68O1GzNT5i2y7ehtPnZ3h8pUl3v2e76IbJ5aMUOAZNEMjQ/zRH/0Vv/d7f87BvVv5vV/5Hu4cCDnyl59CpzHawQmjTsyLieTEWofvfeRGto2U+dp3rjDc53HdXp9elJFpyR+9fBsfO7yTm7au88gNLX7/Yy2qfsDt12qSOGT/vR6z84Zqw2e7lxAPe1x/p8dXjsckzYSRcpXr9igeelTx7BuGrz0Po6OCxbMerWM9EpGghCR2iBAlMkrABeXzB2aQdC1lUKbs21LhtuuqvOveCvsmfc7OwvTlLr421Lf00TdcYWm+jWkL1lptzl2aoVwO8KQ1yeZ7x9QZdPM5dQ5AUi66IF+YCumhyQrTWb48TdKUwFdkKQgFCwvraDTXHNxK0ojoiJSuht5iG+3BYBg4Vpwukjy7cUJf2SPwLFWi7ElCYbg4v0aooOTC3KQTvmhjuWxRkmHrPE0mJIMVn6FAMDpQIXaiHqEk0yu9wj3/bx4eZWIoZHquzV9/fYH9kyEfvrXMN490uXrK4/ihefpqHjVf0xWKv3055cx8ws++dwebN9fxfcHh8y2++MwCv/rTd7Fl+yBf/vppPvmVY4SByxNy/A2J7ShkHlPhILNGa6d49BzXy1I7dJqROqOgdD66zFGb85FX/gyEsJ6MXMGaK6ryMZMFHWb5gWlH+s5IqWTOGlRFFhDOL2LYSB3MKSJZkrowKduK2vrdFQ9SOZ+LjcZIHIEERw02TtBSqMdcty9x1b8UNt8jM+4BOcmZkBtoX+V5Floo8tGTIMjd627MpaT1IthIRkGpFNpb2i3X8xFQPvsTTjIs8rAZ7bAnjvCr0c60Ih1/R2/wq/I/RzpFgiN75kErxrH68/jc/IHlprw8Icw4fbOSvovVtTncTqRvTUFvSpzLx2RrjS7PvjoLvR7X7h9h08QA3355Hk8Y4k6MMIZSycqjO50I3xcEKqMXC4aqHu+8ZYC+kkQJzeOHV0lRjPUFXD9i6KWacr2M1il37hnivbds4ZbtdTwJa60enSghijJSbeg51HLqRnNKGHyZoVyZ0VcOqJV8qiWP/mpIfzXE831aUcJzJ2b5zMuXueyN8uu//Z/424//CVNTEyRJ7C4eRTfVaGmTCFNtkMKiUVLnthdCUJ8YYeX0RVJfUOurEKSGRjtmLUno8xVXVhukZJA5Rk+SsJYYlnoJKy8e5padE3TTlMZak8mpzUxMTTI2OcXgwBCNRpPf+a9/wh/9z7/kx77vbv7m19/JjgsXOP63XyAMbGGhBHTjjEYimYs1IksY27qZD9x9DQOqy199Zp75hS5+tcJvvnAvL80NMjWwzm9+aQc7rxrnl39ygP/xDz2eOywZGC2T6AqDkaEnYvYMKu6Y1XQXPN7xjhqvDkg+fX6Vk6906J7I+Kn3K973No/Pf0fy1a9F9KXga4EMLWW37AukMqTdjInb+un5gv/4bfjw531+81uGp08mqFLAD7yvzqf/YJR/+746sxdXOf7EeXprMclayo4to6ysrRP4AcYIhMhHRvag8j3PERQMYd59vykhMchVN9J9b6RF8+DQQZ6SxEnm/jlUywELi2u8+MpphoYHyFoZO+/dzlItpKOtICXNDFoLAl9yeaXNU2dXefrYPHNLDQJpGKyWOLh5kIcPTFCvhnQjq3TUbsRVdhkXvtt/BqGFJ6apYajPctXsjtXQi2ynHacpt++usLmuiGPNKxdjOlFCzYeLizFXjRk+9rlLqLKAssfjZxL+4bDhzv0VHjzQR6WvRBrFoDOePrLM1s1DXL17BBNrnnzuYj7BJnB7C+Xl1GQXce0iXvM0VeWgpsWC2x36+c5WCkWaJsX7oYSy+13n7bAdDCRpQppao2IOMhTKRVIINhSmzpcRBr7dr3jKRWFI140Yl5JppzuWspFnetj0yCwzKOk5MZJ21N4U5ZIX0yQlLJcKSXKaZsV5LtxEJB/TebkJBofM8PPbNElB2GV2GtngHNz2X7zJH5E6+ZnWVnUlAK1sxkZmUnqR3jDe5DM4bciE7WzsJeIs82LjYsg1y9KNsdLMGojyVk0o3gRBzCVuulBdYSgqA+nmh8ZdlsbRebUTCggp8ZyUt8ARGAuIyzHInh/Yh+zAZJ6SVCslnnnpEj/XiakOV3nk3j089fQbdBNDbbBiM8X9Eu1uihBVhBRUqwFxlnBlMWLznowffOsko32Gf/rOGlopglLI3Zs0L5/VrE432DdW5Yatw5yZbTDVFzJ+9SjHZhpgMlAefs4Lk3YnYo1Ehk5k2VklX9GOInqJjbttdBNWGhEr3ZhIeOy7+ir+v595J+9++0NUqiHL84s88a1n2Fyt2HGir1hs9BB9ZQvcUwpkZllmxhCEHhLF7OwKohSRdhNe7Tb4xqUWs5lmLc3wSiXmlpo0KlU3LtTEGCJj6GWGxuUlpgi4PDvP93/fD/ODP/QDDA+PcObMBY6+foxL02cZKRn+x0/fz41Vxfw/f5vW4iqVvgCd2vHAWjNivWdY7mYsxxkiMwSVEtvrAQ/1C/55TfMXn22jrn2YNxbLVFSDOFFEmebXP7WHP/5Qg9dOpPz6H/X4l2sqLDyjWTwjUIOapZWUnbeUeM9bNP/4tOHOa/o5NdTlY0d6VMop0gu5/t46cSL4p+szjs4Kpp6QTDYMopzSTGK8zOCFkkffP8q9Afzzd9Z54kLAX78h+cR5j9unDA8PZdy0OeP3fnGIWhn+4GNLiOgUkyNjRELT6fZQKihGyL0oIQwCt2603w/PD4iSyObYCCwkVQtiF69qc+Sdt0jowjCbc4Kly9vWRhMEVgYcmYyKDun0IiZvn2Lu2xeIU0PZt9L5NLPilMCXPPpdD3BpvsGhc2cR3TWUTiiHHn2VEtVKSOjZKObAtwDJMPDoJBntXkqfUqQ6w/MEcZrhKUHofDOX17sMV312D9e4c38flYpPN8p48ug6IwMh/+5dU3zlmXlmFho8f7nHG7WA7JLmmu39/IcPjlNOu/zjc6t4JiXwfTqp4OnXlnjve25mYLiP6fPLPPfqZcolSwFPk5QgDEmSGE/6SOWTpU6Q45bcnq8K+wOu0M1lu/l+AhdH4Zdct+gummI34s4nIZwiyy2k0zhG+h7K8/MNibVFeNbQnbg4buMiv4PAJ+qljtIk3HkqEcrxdSUII1Ge/TvyiYrMDYnGFuNBGNhLPHYXl7BIHSUksUnASKcAc+d6vk03DvmdRLEzr9mbKo5tXrdwc7g8PD5zngrp2XwOG2MrHV7EbezdBl8IO4vNF93Sl3Zkhvv3iwdoCg+GMfZwT7MMz+mP87FYqjNL33VGR5t9oQs2f+jnaGvbNhrsjNAiqDfeQOdRchVEvoEyGC3d/2nptdrpqXPJnBLKAd0kR0/Oc/7KOnuurvDIvbv4py++xqWZFYYqAqUMrXZE6AJlEIpuTyK1/UhcPr/EyKYRbto3xOefW+TlY4s8dMMk+4YrXL3Y4onAYyz0OXpxibOza3SjhO3jdW7aMkC3l5DEMV4ptBGhOiOQCiXgxHyT12eajNX7MVJitP1glKoDjGwf4abxUa7av5frr93H1Xu3EwYha80WtYEBHvval1m+dJ6bD0yx2o2JU81KOyWcCNA6te+by79IkhRRlpTKPo994xVGdw7wt5faLPoe11y9lWs2jzI0XOGz3zjEkNH0hT4trek2EuIUWnGK8nw6zS6i22ZkuJ+K6fI3f/bHhKWALSNVrt0ywo+9ex/b0DTOXOH1M3NUah7lehlfQELGejOimRlSY2hKRUtYr8jkUJVu0uGewRovpCmfeXqN4bVphnduJ80MUa9HKDOOLPbxp09ezV/91ks88kMZf/g/uryrqoiqhlD7qKGMiZszVpYlH7on49uvJ5T9GsubAn7vGx0ebaV8pK/BE8sjDO9XzPdFrO0SzL4k2PmKT9iLUAOKOVXiqWNt7rg54Gfe3s+HViTHLmV8ezrlO3MBT52uMXaoxP1bMt51H3zuiTWakaDSF/L68XMuAc/6LQTShY8557OL2cq0JvR8qxBC2iz4fNFaSLkhdpBRKS07Ll+y6+K7Ygu7JNWcvTCLiiXRWc222zYTa00nSij7ocNcaOqlgGZrmepAH//rP/wK5y5c4sTxUxx+7SgnT53j/OIC7flVSBNCZagFUAoUlTCgk8HZdcO+kRJjVVsFW/OclaOmGLbUA/YqwYVGzPYxiecLZhYSjl/u8L231sBTvOX+KV59UdNt9KgNedyyv5+779jF4mKbucU2cZpS7Qso10IOHV9jtWt4y52bwVe8/MYCswsNymVHC3cmaKNBeDZAzZATLBSeZ5WgBvCFV5BwwYIipXLGQSVReI68a20Qyqgi79zzVUHAyBKNxhkWpSU5p86fliNRTGoK83VORDfGEMexDY/LNiCPvkOqWEmuwiQZmRPM5O+9drtgm/5px5zG+ffsCkPjSdttGJ2P4WymklICL8tseyKN3T3I/EOZ7xJybImbpXkuZlZi8Q2WNeiy0wvipEb5HmT2hLYyw8B1Ei5D2lgVVyY2bmpRyHpdvrkDARbqK+Fyw93pL6UoSJ6eZxd6vm87BO1aPS/wC+WA1trRPOw4yvcVqXAgRWHDroy7lRPXOkqpEFrb9lK4nYu2SyZfCeaXmjzx7AX27qhzy8FJRoaqfOvlOX7yu3czUPFpx3bnIPzAUjBFAr6gIwLOL2Tc2GgwMdjPXQdH+MZrK7ztlinaLc39A4YnREo7Ad3V7JvoY2aly5GLi3hmiOFasIHLx1Z/1qEv6S/7NNtdPvx97+bH/90PWWoAmlK5SqWvRuCHgKHT6dBoNIE29cEhpqdn+N9/+TGuGq3RixMUksV2ROQJxiZqG0hoLZCBVdp55YCp/SOcO7fAL5+6xAPXTvHb77+fqclJwkARnb/C5772CjdVPbIoJUlAAatRjFGKqAuen9IXBExUFP9hTx+lepVatUR/OaS3tEb0+kWmlxpowC97CCPQqWY9yVhrx3SSBGk8okzy7cyQJZpyOWBseIDk3AIVX/BQxeeTmSY5/zzrJUVleAxlMrQX0N8X8enDE1y7bS//8men+dx/7dJoZEhl6LYzrv4uiZIQ9TIuz0ruPQAjF3o8daLG972vxFefSvj6f+ux5dEWO/f5qJ6HX8pYfsgwc1Aw9aTHziOGbY/C008u8ulPwI13+Lz17mHu3Olx/36PE9MJT5zocqRT5huLFZ69UiPjPLVawOUrCxbP71i3nrTjqk6cQZI5w1pGxVcYI+loa7oNfBsIJh36Iok1JX9DUYgzoerMEiBkIZGPi0wJJSTdZsy+/ZuZPrPI0uAKsh5weaFFvRJYqoGS9JUUOybq/O3ff4k77ryfG266nq1bJnnkkXtJM8PaWpOVpRXmFhaZX1hkaXmZhYUlms02prVOtHiMEwsp/Zv7qTgzXZIlbuIAk/UaL5xdYXRYMj4QkMUJzx1bI1SCO/ZVkGWP1mqbYZXy0D6P/ol+apMjLK13OfTaNLt2DBIlXSplH4zkiZevsGv7CDcd3ISJM7721FmiJKV/oGL5c1FWjO3zoCclbUWvjbbZGcKKVVJH7BaeVT4p5wLPHO9KuSwUsCIFK/YwBcncYIh6kbM3SJe6Kooxu86FDi4jJCfz5oF7QkjrX0tTsswJoMgKz5uUkrgXF0h6iRXVxFHspi2a1L2Pxmxkgli1qpOHa0vhEG5cZ1zhoGr9fR+VDmgm3a1oJWmqkMXaJZlyFCu3PHL3bW6osQtp6ZRTecpYnlwoCgOL1umbLge3QC+IpHZslBNKAfxAOfS6R5ak+Mq3y3llFUXChbWQKwgyXVRFFqdsOyzfxWnmWe+ZQ8hnTmUhhbAsf/cQDcW1X+w8rB7c2bQxRYqYNoL3veN6ygMlLlxc5dsvXuF7H9nOG2fXWVlpMFRWhIFP0osolSRSKl6cjjmxKjk/ExMGPjftrfGFF5bZv7nE0HCNyZri9ELMSsOwtaq4uNimm6VEmcZHMDpYxVMUMkPLR1KkWlMLPdqR4flDbyCymKmpCUbGxhDSI45j1tebdLsROknxgoBSfZCXX36Dn/ypXyZameOWXaOWuyQFRy41YbRK/+QAWWIzmnOvjZ0LS0zPsHhpje9/9Fp++8P3Io6eYfHbzzP/1Mtcef4ET7VS3lrPJcUGpSUvdWISbbjWl5QrIUmW4ZGxL4nptrt0Z9dYm1khbrRI0pSw4uMJg9QZvcyw2Mk4v9wlE5I+T7EUC/45kZxu9vC1oVwJ+P73PUh2+A2SxSUq0uNKTzPdjqG5SlIZxi+XEMKgtaRUMjx7foKbxlepvL5IK4WoZRjZK9h1k0GnHij7xY5Tw5HHBbW5hC2jATffXGau5/H8tzKWL0cM7TOUSj7xgsGEhuhWwflxQdoMePCaHrdtb3PmuOYL32pz6HQPbWJ2jMAduwT3bU25ZzJi62DKGyciwqDK8vKqJSxntjuwdh3N5nqJwUAx3Bcw1lfi5u11dgyWQcDmeonrN/VTL/us92KqvqReUax37ThaSGuktWNbBxrV9jOUL0+tpNQeeH39Va7av43p47NUtlZpXFxnx0QVMhs/64eKPl+wkgb8wIc/iCBjda1BkmVESUa5UmFoeJAdu3Zw8MA+rj24jxtuvI7v/753srDc5uhz3+GOXWNFnCoO6mrRNhlpBsdn2nz/WybYNOqz2k75f1+dp172uXpUUe+XzJ6eZ3lmHVUKCMdHqA/XOHFsnqGBkNNX2uAHPHT7NhYaKX/8yXO8/53X8OD9e5m+sMp/+7NnMNp+j9BgFbtOeSnsqC5neNnPvilUop4b62c6K4RGefdh3Hg/dWFwphhIufPD7SCko+bmJuwikM+JVXOUjXLkjjfpmGxGSP59dCFljoTixlwU+PZ8BbARPe4VZ/S/XuLLgv+Vtw+eZy0AeeSHAFStr++jONWGcGTafA+BAU/57rbTRRqgAbsISq1BRTnVRcHF92ShNdcuftbkFEoj8Hy/6EZ8LyhiFpXbN+QJgWCK28463d3MNUkL97vOMpfwleep20Q7KaWDP9o3Hg1xmm5g6R3hMndkCmENUbl6S70Z9CisEgOR/4wzIQqJ5wlm55q848E9DE/1U0LyV594jXtvGmOyDt96cYH9kx5RKhke9MkSw1o35VvnDXsGNXfsr/DtY21OXGwzs7ROnHncf/0oWioUmq8dXWN3PQAhubLSppNoxgcr+NrgSQh9q3DzpHQRn/bljdZ85tspTzz9Kl/8wmNsnprg4MGDdOOEcrWCCisE/UN0ehGnv/gZ/vF3/xvPnpvlgavGqZWtmGC5lXB8qcPYNWNOB6+I4tgFbG204OvTTfR6h9/66A+TvX6GmS8+hZSCVqIphSFvdFLWkowwhUakieKMxzop1yrJtorCq4R4Screso8Kg2LcWesvEwTW4Iq2n4X5RsJsJ2Ut0oTSwwjJ4VTyyWbCxdU2fcqjk6Rs3zTK973nXlaffRlWG3SFoBRnvB6nmCwmi7uUxzcTODRDUAqRAr5zcZJ6ukDfSgu/L2T3Qwnl0C9UgUFJcflUyuXXYKWVUdOwLZTM1RO87Rmvf0OwNK0JhmNqo4q0B2WtYBR2TGasnIp45aLi/psV77lNEmQJX326y2NP9zh2NiLp9rhmk+aaqYwvPJmw3oLG2jqlclCgdfpCyc3bh7lm8wCT9TKTtRKT/SG+Eigp2DxYYrQaUAp96iWPXSNVdgxX2D5cY/NASH/ZZ70b04lSN9Ky/qgwCMhcRHQu7ZXucpmbXSEsl6hVQ/zxEmsLbfqRjPSH+J4ky6CvHNJYW+fbr53gXW9/iNHxEdvJuiCiKLIxrWtr60xfvkQ1rPAP//Il/uLP/5pbtgxQCjbc81Zy7xE4OnOnHTMy6PHwTQMEnuDZY6s89vIaP/7WYda7gmp3nfmzy4S+pjoywOjWEWYvrXP6XIPrrxrgE4/P8cg9U+zfN8Vjz17im682+O3/7z7Gpvr44tdO8A+feY1KJbAKJyd91loThiW3R9IWyeKQ7DaO1+460jRzZ2SOHxGFD804/5lwdggppIvtzo2Jb5LJOnVTDmq0WehOPkvum9vgX2lDcYEZbVzchXXOO/Kii8QQxdmsdUYQhhgX22HNpgnudivO64LoTVZccGmSIj1ZeAEBayTM1QNFte1wxYXkzmxEo8KG2zyXAFqtuXCqJ5cO6CnLi3GdjNHa2e1FoTu2agT9ryS15s1QMaegsq/LgfZMUTDZgCrfZqO7pssFsshchFCAzXLmjFUzOJ22w5ngFvieZy+WJE03HPIOyWyTFk1Bq8TY1lIYw/Jal51b+rn92hFGBip8/ltnWG/FfOChCb78zBUGVEZmBCP9ISUleOpsj+EwYzRU3LSnj/fcP8XkeJ1WF1480+CBa+oov8SWQZ/nzjeYWe6xY7BCJhStbsxQOcBXVv2S40fSzHYDSWJnloHvsWO4RieDoNbHg/ffxcjwMGXPg0aLZH6W2We/xbGP/yVf/+KXeKFjuGHHMFsGSw5kaHjh/CqlHXXKwxXSKMOY1MaICovdVp7EK3m0lrqUEnjfd99FdvI0zUvzhGFI6El6EZjU8HdzXV5pZbzeNbzW1QwqwZ2+RPlQQqOFpKk1FSkJS74LOIM4M7SijOVmxFovI0oNvrAKolOp4fNrXb4xt04aZ/zoh97O3v1beea5o1x3YAfveuAGVp58nrjZIfCFzR5PNYfaMX6nCcqjPLKJpJeQZYKSrxF+jXsfbNBYahD09zh4s6LXEUgPstQWI6eelPS6Bulrlpc156Y17bsUt75VkHmCU99JmTuhGNvuMzRpkEbRbWa8fdzjA3eW2bFd8A+fbvLKBZ/RoYD33lfmges8ei3N0y93eepVeOV4zMVlH4RkaXGFclhCa8Pdu0e4ftMANd/KXzMXCyulcd2w9QzZw8LhaJQt3DSGqq8Y7/fZMVxjtD9kpRvT7ESUA88mJLoCK18Kp2lCpjOCwPKnOq0IUVMQerSvtNgxUUUnGULYXJmJwRqnTpzkC996gT179rJj+xYq5SqlwKdU8ilXygwP1hkcGuaTX3iMv/7z/8PNUzX6qwFxamGmQggC5YGbVigh6cUp99w4yOSQj6cMf/nYLJWS4pc/OMWrJ5osXVxmvGZQpZDxPZMkCTz38ix3HOzj8UOrrKeK73tkJ3EK//vz55jYNMRPfuRmRKb53T97lpPnFghDn9SRK5RSRZVufW6qiM7OM4Jy0oWfB9T5vlWJ5qRqhzmRuY9NWo+ZcYevQGzAQ4220xGzsQS3awRZeD6sgs0W2hY6a02AVlFnNgp/F+wnhFVdGZ0hXJdk3C7b7n4ded3zCq+bMWbDYOrAtcrzrFXCFd2Z+zOQ4Mk3XR6+b7MbPN+OZoRHwZ6SxmY347hWcS+mVAoLdVQhX3P7CeFaT2m8Agmi3YIHF5qSJxrmlNwscS2gU4X5yneqBF3sXtIkQmhsfG1i1QNBWEJoCv27dn+HUt6bjD0bOSVCKYQzduUuy1yWhpBF22byMBj3H2v2SYsPS+Jie6vVkC8/cZaf+MjNVMdqvO3enfzdpw7x/31gGw/eNMZjT1/hkWs0M8spU8OShUbK+27tp9lIee34KoNDJXbtvY5/v3WAH/2d5/nmoVW+774K7VaX77l5kN/+1GX6RINNQ/0MlSwyIgh8ZpoxE1XwHXolSYXl9ZvUac0ND+0d4tRCi//8G79Ltb/O9mrAWNLGyzSLUY+zsUT2DXDbjoB6ySPKNN1ewqErbdKBMhM7h4nasc2Rl/aDa0d3VvrnpQpPCdZbEc1Gi5oziD67lvBaQ3MxAa1TtnmSzBOsZyldBF0heUwJ3o4kTDM84NtdwzW9HiNxQKw1WicE0qKuPQElY5UqryWGb632mI5Tsl7M2+66gY985J2842338PO/+kekqWbz5AgqjumuNDBKECqPIEi5taK4FCsOxYLg8mmagxP0j2/GlymNTsDbt17kKv8ij93mcYV+wvmE22sxHS3wSpqZ4ynrc4pIpJTLJfpI+Hoj4dWnBrizlXLToynDE/Di50Je+FjKnkczrrsf+lXIUKxprkKYKG7rKcb2dlgdzvjDz/SoV0Ju2+fzC+8NCQP47PMRK22fqkophQFrrQ67RqsMlXziVCNMihJ2nBr6FnEjXYVssT4eqU5cNeu6ak8glCJ2KI3RWsh9u0f51ok5Gt2YWqVk1U+eZc2B9QTgWE5RJ8akGd5KyOQ1Y5y6uMbl5R7bhsv0nE+pm0TctHuKszPT/Lsf/gluvP1Obr/9JrZs20TgK1bXGhw/eZoXnv0OevkKj149RpYa2r0YoSSB5xOniR0tGyus6UQJ9YGAHVMB1UrAifNLHD7b5mfftw0TDnHNvoiX5lYwWcrI5ADawFPfucJNe/uYWY557JU1fuPfHqBUKXHq8jrPv77M//7tmwhritePLPDcq1col3yLXPesoS+OI3zPd/ERthPI8eloy4ArqOPu0s1cwl+iMyQb3YEttC2zL0ttDLdjkdvPtrLP2GjLz9Ku2NYmQ2TCRep6bnScwxI9pG/Pql7UcwBcZwZ3EboWhGvHVMqZGpXngwvms6/XYkw8F0Xuh6FLccWZwJ1a1VNkSUoQBvR6kfv3DV6SJAUy3AILLUDLuEUxjhilPIkw1seRJYlLybKdAcouoJVSpCbBuK6BN8H5sjS1Ma+u+jfGWB2x5zm5oSly0aWwH3Tn+Ed6ijiO7T/XwnYdjjCrtXVFCwe7klLa8BacOTLfz7huwsLQbFBUHuGYX6D55kVKWVj1ExcHKZUkSiOkZyXJaZIUGuxQGV57Y44XX5vhzrt38J6Hd/GX/3iYbx9e4wOP7OLLzy1xfiFh/6aQmdWYvoqkpgxjU2W+83rMaiPFnD+HT8ZHHh7nr742wztv6EMGipu3KG64bpSjr83SjjK2j/cRSMGV5Q4nFtrsHKtw/ZY+el2NNmmh3CiH1g/SSzX7JvvZMhByZTViud1jSXmgQqr1CrfWfcZDSTcTZFJjUsMrMy1Oz6yzaf8Y3W5CUPZcTrkpQJdWwGBpyZXBMhfaXb7znWM8MjLAf72Ush54XCUyHqz47Oy3vpAg8OgmCesGHl9PeLWXsTeE2wJBZgTnopg4gluEjzQZAwGUPY+lTkLkKS6l8ORKxIlmj6nxQSYyzU/8/Lv5/g++lTg1XL6ywKlTF5BCMjk2SLyyRtxoE1ZLxFnKaOjTCGK+a7DE7HpKN02ILxxFj04SGUG91OUHbzyB34OTT3fZeZ/kiKqgOgm3VgVtkzJ3EhJtSIVAJCmdrqayp0pJV/jm4ymXrvS47WbB9/ys5iv/kPDK5yTdNcPDD2bsmPTpeRlnX+2RrGuWTwe88701Hro+5pUjXc7Nxfzj0x5DQz7HL2UomRKnMXGcMTLYx+27J6zB1RdIvGKcYUyGNnbR7/uKNNVkpHhegDGZ48JlzgtiEMZqtpI0pex7vPXAJM+eW2VhrYvn2emCxmLWhauIe1FKtRxw2417OXT0PGv1Va6+exfPffUNQm+MoYpPmiV4nqTdjdk6Psi2CcOVUy/xySPPEzt9DFlK2YOtY3VGd4wRJymZxo1rDO1ez7m6Nb5SXFxp89rlJj/8ls2M1APidptvH2sjPcmNu/poNCLGByReFhNWQi4sZMydmOXAZp+5lYg/+eoiH3n7Vg7uGcFow9dfWaK/r8LDd+2CDL7y7XMsLKwxOlIjiqz508r+7fPV+XQlRy8VDnXrX8vHTJnWFnsk3Cg/zfBCN6pP7UJeC42L/CjiLoy2naEldHhFVIVyExOjsBLaJMXuCyiK4SRJyKQuQI/aXQhK+ZjMdg5pkiJ8ZXcX7qzzHK1ASIGxf11BOtfGqraUkgWuhfz3cSF9frGC0BamaJyL0aLWLXhPuKVyodTQGTp7E+rDVTG5xMwL/AJjki/e832CVMqqodLMqRkcjMvYPYJEWPWTw7znb4xGb3C0hJXqSk8VHUOms2J3s6GvNsW4Kl/wZ47HZdVXb/ow5Kov1y7mrahVM9iHKAqcM8VFZ8dezoRo7J+51ugw0F/mrffuZWyszjMvXODFQzO85+4xJgcUf/fNea7f5LPclbQ6Pa7ZJKlPDKIxfOfIGjuGDCWhqYeCb77eQAY+126rgsmQFY+LQZ1uFHN5pcfMWsy6JxneM8ClC6v0Ys1w1Sf0FEiJrwSZCwGLEosy0QjqtYCtQ2U2VQK21z0mSoJaKIiB0IderHny1ALTS21+5COPortw9NmTzF1aJKyFVOoVlPBcSqIdH5JB2B/SXWgzP7PGiWbMhfk1/suOkNtrMOVL+kKPsOSRZhnVQDAWeFw3UObxlS7DaK4drKJLAV9ajahKxUTJYzlJmc/gpV7K062MLy92ONTTXHfLAfZMDfMffvYDjAzVSXXGNVdtYW2tSaPR4ROf/hadVpvve9/9TMURq68dJRPWn+RLha8U3ShmR+DxategOk1UOSDt385/vOs1rp5s4FcNA6bDP30x4rrrNGc7JaTUDLY1x75je9ySJ6GX0ioH/Ogv1TkwDm/MSc7M+sxf8ShVBPe8VVKqejz3WI+VQxHXb8/YsaPEya+lxOsxpUHD+MGAbkswt1DmnfeUefhWDyMi/uI7ZYKhEp2FBnGccv/VkxgBL59f5uRMg5lmRJwkjPWVbKCYSZ2IxY4YUm0PJc/hhow2TK9FTK/0aMQJvpL0VwJ7MBnYOlJlvhOTZBsLVG2sZ0c6blMviplZWGHb+BAXTi7xyz/5fkbGh/jbzz7DcH+JoVoZgSTwXZaHEIzWa2we6WPLUIXJ/hIT9SoTQ/0WABhn6MxOKjyHHxdSECUZgVKcnmtyZLZJrVbmQ/cMMzJcZWG9zf/95iKhMVw3kjHVr5k5d5krczEDAxXK5ZBNAykvnI34+HNNHr19ko+8YzuZNsSZ4Lf/9gTvffsBvvvte2k1Yv7T7z/JymoH6c68PA/c/rflWmW5A9toayQ0YiMCPBcTud2rEBKd6mLJbowoCly7J85DtHSBcLIHeOaGWmxE9rrzR2s7htQOGmtFRcG/WisY7boOl0+ii5WBdKFTzj8nKPa4ue8tczga3J6lGNu7zwHauVDyuEPnzFdS2B2ILIBhpkg888PAOVspeCjK9zYiZ3OsR6EjFwWfHgR+YHHjSIekwOrXwzC0uuUkxpOW7Cg9WeQd5y5H+1Bs65xL14SU1mAj3NLcGIIgcDewLH5OeaK4RXPNNMbO7HOoYx4bqbXZkCobCyHUWhMGwb/qyPKlWhD4BZkzb2CFsMKA+YUm73/bNVRHavhpjz/5+0M8dMsY9904yrnZBqcuR0z2w1wLbtkWopWCTNNZbXBqLiXODPu3D9JfC/nEU8s8eKAPv+Qz4iW80fHwxwfo31SjOlmhb7yMVwuo9Jc4f3qF1W5Gf19ILbR5EcJyEuzvlmkCXxFHCb1UE7toYeWBMBpfwErP8MVXp5ncvo1bbtjHLTdfxa/8yg9x2y378TLJ6y+eZXlpnYGpfrvfAktgTVM8T1HuL3HxyAyvnZ3h391/kJtGyrSiGL8S2nFhklhVTxBQMoqT7Yhn1mMmPcnVPjzfSnmlrYkEvLje4eVOxivNhGUVUh2p8zM/+k6+/91381M/9DYuL6ywd/cUo4MDfPKzz/DAvddSqVQ5c2GWT/zL1/Gl4Ec//CicOEXr9EWCSslVcTaLo5tlbA59MuBcamjNrfKB+xJ+4NYl1td6dHsR+/Z5NBcMX3/RcNctivORonlJkL2eISsZEkEgNK/ogPEDdb7rDsWeoTZdU+JcNMAbZ2F1TnDdwZT7bs7IzsT8yWeazC9qtnRSvFjSf32ZLQcDWq2EtRXBxFiAMPCnXzdM1yeoDZZozzYZFILxeomvHZ1luRXRjDSNdsyFtR4jVY+xWonUTm3JtIswFtIut7UGLWilGd88scDsWo8r612ml9vUQp/BSggSKoFivtljYa2Hr3JDGwSeXxxwvufRiyKGhuuMDvVx123X8cHveQuZzvjs118FDJNDZWsU1oZq2c770ywjyayRVQpJktkldeZwKqXAjmeSTFuyhYY3Lq9yvh0xuGcnD187xF07Q6QUfPbZBY5fbPPhByZZW4u4bafHubPrxJFh/44yJ2fafOFIhPY9Bvt83vvQFoaHSigMT77R4DNPzfLff/Uetuwc4omnL/PnH3+RSlmhlG8L18Ar1KDCyV0FuREZl8DpESexU2dZAq41d+aCUFPkj+d+C89V9PnvbnexG2BDqWRRVOeJrzgWVapTt3sWGxkdbk9i80MUXqA2mIJGu92IdN9VUcAQi7wkt8+QQhaJhsKdoRsMQLeXcTiTnMJhLzA7gZIGQ5aaIhLWuMO1WCZJ5eyIrjVzPo7cvJK6xReOU+U57EjOxTfaRifqNCsO4xxHUrgxUzs7TJO0eHhZltkZozNK5QCx3H6fJwembg+RR9zmIVjCmRQzByWTLl5SkCdqZYUkL9deW1e9dGx/syFHNiYP/yNNU5Iktd2Xu6GzTBP4gvNX1vnWs2cgi3jLvTvZPDnAZ56apWsCfvL9V9ND8rWjDXwlaXchieybWK8HXDWluLzY5a++fImSMEgDX3mtSSlU1KshD22S9KIMUo3wnKM1MVTGBjj46FV0PcnXD1/hlek1Vrp2lukrSeh7BIGtwAPfI1SSWugRSKucmW9mPH5qlb9//iJb9uzkf//hz/OB9z3C177xAnHU4+p9O/it3/4x/vkffot7r7qKk4+ftBdS2e6nhKdsRvdIidp4HyrT1EeHycKQONUo38cTAl9K+kIJccapdsrfXGnSascYKVhVkq+0MkSaUq3Xue7Ga3jfO+7lt37hQ/zp7/wkaZygjeYtD95KJ8o4e2GeWrXGVfu3k6SauYV1huoVLl28wtp6l+H6AKP1EtHsHEHZjkaSNEEqQ2Iydg5VaeiUR/oUO8oecbeHunyMOOpi4phqSbK6FvIT7y+zvyz49hM9ttQyzo6APijp0z7Sh26c8oq/iZ/72yn+9xfLHNxf5Xe/D37y6jl2TxlOtmv83RM11ppVfu6hCj90c5VPP97kl15p81wv5qrbQ2olwaEjGk9khCXNpfmIFy559Pf7RK0uSaYZrIYcX+yQZpqtQxXu3D3M/fvGeGjvGAMl33YRWuB7doxl+Us2rMpSLQy+FNy2rc7e0TI7hitIIXnm7CK9NKOkBCSayVpgDYTGBsQZYw2HOZ3bZosoOt0e2kCz3SOKYx59623s27eZS+2EL7x0hmYvYaAaILINgYwwBiEVSaYRKKI4RQoDbspQ8iWVQNGLUl6faXB6vcfYgR0MDlW5b7v1X1yZb/KZp+a5/aohHrl5mLVWxKuvL/DSuYhX5zP+79Ntnj6X8uD1dX7unZsYHSjhuUA2LT3++fEr3HLdFNfuHUF3NZ/6yiniyI5u8kAn48zKgdtx6CxDKZfEl8t5Mxd94IrYLLPTGe1ECHkERJ6hZMPvbNXtK4uCt2APOxbaQKVbL5xx8FeDcMVzPvHJimiLzGQbWUkuo8QWxvbMlq5QziPKLTE5dVG3HkmcOhSKVY1JoVx+U1aQgXmTmkxKdy66eGKdWe+Lqg8OftQOvywynAL1a3/JNEnxAt95LuwDzqFdVq5rCkVBvpjJDYjWBW6jY33fsygQzysMNQW3xd3w+bhI6w3FQfF6HGs/l/nlioJcuZAjA2z87Ua2eT7esrA56bhcFBdW3obmeBbjBAE4irAqcPKFLqJoAe1ITOEHHp6UdLsxnXaP73l4B5UKLM6u8fdfOcd7H97ESL3CwW0VPv/CMiN9sKVf0sss/78dQa3qM1KT7N3ez6nZhIvzHU7NJ7z1pkH8IGQ8hFcXezS61llvM5VtUNDS5TW6y11uvG4vEYKj5+Y5NbPGcrNHpA2x0fRiTTfWdDPDfDPm+FyLIzMtXp/vIqo1fuQHHuXXfvEHKZV8atUqX3rsea45sI1aX43GWpPhgSrvfs89xJ2Mr3/2WUZ3DjniQE5wtkXE3Pkl7jm4mT0lQ9ptU1KClbWII2sJz7czPjvf5gsrMdVywC37h5ld7fJ6bDi10OTuW/Zz4MAufvM3foK3Pngb196wn6WVBpcuzzAzv87jzxxBeh4vvvIG73jwNmoVn0NHT7O21uaBB27iHz/xdV589TQ3HNzJd919FXOPPYUSBuOyG/KIASGMJR93UrYqwaE049VjmoO7Bbu2Sbo9kKbH/CW4fTDliaM+ugrbxxXxLoG/ZPBmUrr9htmJTSxlYzx/ssarxzx2TcQ8cLXmhsGMwbTNeimEdY+9KmZ8sMJb91QgzfjaTMZzx2MGShlJ4vHgvf1AmQuXIr50CGRfCZXByuk5JvvK+L5i/1gf+8YHmKhXqYWS/sDHU5Bqje9Zn4jG4CtJonXBnRJOtTZYDdhcr7BpsMTeiT66ccbZpS47h0uk2jBcDVloxax3E3toAkaKwhEupTUZDw/W8YKQd779dsZH+vjCF58iLCl+89d+mOPnF3n6yEVanYjQV/RVAnSSEQSqKA59JaiVQ3xhfQ3aGBbWu5xb6nJ4psFC0mHXLTvRQY3tep0H95UJy5J/eXKB5040+d67Rti1vc6w1+X4sQXwFDsmA/ZPGL73wS3s2VRjudnlhVNd7rlhkMG+KseupPzFZ8/yqz9+C3fcsZWTZxr8xh88jlCmiMO2pAxdVOG2aLSqzCzdyFwx6EKSa8nd+YTkTdMSYc+XnMRrJxumQCTpzO4iPE+RJAlBELiOw7jxmMIYF5shlPWqKZfBIZ3zXRvne7PCB+NUtEVQX4F6dyRkkbvrct+JTWdNs7TAxBfwRKdMzb/jdp3hEPNpnrcOqlKufDR1um9b5WdOvpU7MRVZ7HAg7gNkKbz2Ngx8nzR3XCpJFMVFzCO5yknaFxHHiav2dTEeszegPdiTJC2ckP+Ka5XzqRzSXWc2N1h5ftFG5vO9fIyVpygqKYmTxOWnu6tCCldp2FvHGgftOCbXvZOP0QwWaObaRasjNu62ti1dkiQ2eyRQXJxp8NCdu5nc3M/W8T7+/BPHmBjyuOWqIQZH+qiVfN54Y4E9w+BVQgb6a1y+vMpATbKwmjCxaYCbdwb0i5iXzid4wnDD7n585RMnGa8u9CgJ8CshUsOpZy8wEdT51V/6N/z8T32Qa67ayanT53jXdz/EYrvHxcUG5+fbnJ5tcnahxZmlNnPrCcFAH6VqSKUU8DM/9i5+8kfewfJygyTTDA0N8OwLrzM0WGP3jk32C5RldFo93vFd97C+2uGbX36BkW2DpImxCo3UEJR91i6v025H3DtWJswy4rUOf3WxxZfamjfaCT2h2TxY5r/89N3ceedBPv3kaU7PrnHH9dv47d/4Ub76+EscvGY3ygsxMuDipSWmL07zR//933Nldom/+vhX+Dff+yjXHtjJ2kqDKE546rnX2T41xic+9xTTl+Z5+L5rePjAZma+8RxGKLers1W5dHTiauDRTFNqUhAaweux5tXDGW+92yMgK+bgz/0zvPtG+PIRj5HxlNGKh9miab9qmLimxM98BNZWNGeXq1xYrfPNw330epo792vu3p5xq1znoGrilUtQ8ShVfG4aSnjbtR6iFvKxr8Y8+0ZCpxuxd1IyPqT5l29oKpsHMXFKutDkhk2DbB+yl0aSGTJj6MUpGhvrrKSymfeeQknLkZLY4ilOMve9NSRJRpJqFLYj2TFUI1R2T5Im9vt4fL5BO0rwXOSzQBD4vu3clf1ODw30Mzo6xHc9cjOeEnzrycOUQsnD99/AfXddy9X7d3Li4iLHppc5eXmVZpzRTlJanQgtBL04Y7nVZW69y6mZJocurjDbyZhZ6/JffuPf8vADt/PCi0dZmF3kwzcPMT5U5vxsm//xyXOgMz744BjDo2UWzlyk7qXsmAqZqmVs2TFGaWyS6XNzpH7A0emIR++cQsiAP/7MWTLp8Qf/8T7K1TJ//U9H+cq3TlKrhkRxbBWW+R7V+ZuSOC0KUc9TDvnk9rbOtCvd6F47+kV+QKdpVmSgIxwK5U2xHepNvL3c4pCPspIkLkyLtoNwloQ8T8lRuOM4sXsuzysUrrIwbTsToYvV5c0AW6eQteFSnsOWZAWryzjpsdbaesyUW1Vk1hMShL7b2YCq1fo+KmQOx3ImO/cHKeU5f4Y7dAtjiizyB7LMXjL2ArKjqVw+m2cOqFyqK+WGG9MpCXKFQw51M2+CiwksxjnJUirlStFaaePwx8YgPFn8OTrLyJJ0A30iKVRDedcgiwx1uxQzTo3gbB+FGkw7gw/aol5F/po96Qi+xipWXPWRpilh4LHejKhVSrzlvt0MDZY5fXaZLzx5ie++ZxOegpW24cVDM1w9puhpn74+j9XVHqNDIeuNhEYE45NVdmweZqDm8fGvzfDQjUNIL2CqFHCqldHxfOLlLqefmea977iP3//9n2Lvrk0sLa6Tttc4fPQUj77lLn7p33+Iu265mrc9cgd33nEtj7zlTt759rv4oQ+9lbc9fCtrzSa/9FPv4ff+56epVEtctX8HcWpJn2+cuMB6q8n1B3dY4q6yEsBOJ+W+B27lpedOcOLUeUa2DxG3EqSv8Ms+Kko5dXqJJa3ZPlBltRFxITKc6Wl2DYV89/4A4SluuHUXm7Zs49vPHOe2Ww/wf//4ZxgdGuSzX3iWm26+lqnJcZrNLkvLqxw+egIvqOBJQV9JsbC0yucfe57HnniFY8cvoNMeL714jFPnZohjzQffez/bsg6LLx7FK/kokUvUHa5DG6JM0x8EzHZjdoUeTSV5eSZhadbnHQ+FdJuC/lHLYjvy+Yx33Sw5PBOyezLDDzOiUFHpCFTN5/0Pa6bCJqcWKzSyCkfOVnjxhGJqOObgTgGpR9rV+F5KnKV8Ym4vY2OK997W432PBvz/XP13mGXXWaYP33vtdHKonDtUdc5qSa2cg6PkbGNjj8EMAwyY+c0wMMAAYgwMw4ABG/BgDNjYcsLZkixLlmSFVuecc1dXznXyOTt+f6y1d7U/XZcvS5dK1RXO2Wu97/M89xO6Ad/4cYtvHwhoBh6nz1Zp2hmspI0/XWG4mFTVriG2urGbuiknjGhPTSgnKyGnD13X1Jo3xPPla9TUdZK6RippxHbTQsJAaDpNx5cHkA5T5Ra2sq5ahnRX+YGc4IUmmF8sc++dW9i1fS2CkBd/epSOtgzbN6+hXKmxfmSIJ956N/uOnGHtujXMVhrMlh2WnJCZssNU1WOxBSRTDK0b5GM//xjFQopsOsVv/up7WDu4ivvuup1EbYptHU2KxQz/+INrXJ1scPeWAltH8qT8BucPXpEtoUKj5cOqWzZy7OAYQ50mPzldobMjw727eri24PF/nz7Hb358Nw89MMzivMfv/p+XWVquYhiyME4oikXUPup50iQShfrk1ODHiCfZCx4FmIO4dC4CrkY9HZqQayArYcsKBNu6CdRqxOV8saMV2TAoWy/CFdhsdNBoGprQVUmYFiNnXNeNQ9ZSsNfiTZGu9BaiyUQl7DVVBihJ6XLCsExbWo9VfCFyo0q2oTwNoo1TGIYqSEioOnmFqkMMMC1L1dSGEmFhmzd9MrGCVQ9DOWpFQT5lFZTthFHTXxDz7KNVlRC6au2SOzmJDNBi5V9oAmHIQ8kQAs/xYs0BhVP2XDcub4nEqpjXFZW63HRISc1Dfk9CBXQM04q70SVNM1A6iR7xBOJ+lDByNYDSVJRzTf3sQqT7a2yyzHsfXUcmb9OVT/LZL59g14Y2tq/rpF53ePngJNuGLGaWfLrbbMqlGtm0RQuTy5MOezbmaOoZ1g9l+MnhGZYqLndtKxIEkBaw71KNsUM3+M+/9n5+67/+HM1qnVqljm0nCf2QpOHxl3//bbZvHWb98GqSSZN1I0OsXjPA4GAffT3t/PDH+3nmR28ShAGnz13nhRcOcvnaFLfuWk9/bwdnzl1nZnaBB+67lWq1jtuo0Wp5WIkMIgzZsWOEr/3ri/imTyKflIcyIZm+HM2ZOtcWWrw0WeH1ZY8xJ8ALA2wDNvUYzFc9csUsG9b1c+DEOIlsmmrV5ZkXDnLy7DUuXRpl394jPPfMyxw/epzQbXLy5GkuXb6C16zhtmokDJ/uQoKh7iyD3Vm68glOXJjCcT1+5eNvJ3ntBuWr1wlNm7Dlyru0LLqWHfdIO3hCwHzLZRiNqYTB4SsBfsvn7ts0yjWDrg0GU5c8rh7yeesdBm2dJk4NfE8nScDSjEepBrdsC7hv7TxTMxWmy0kWmllevZikVGky3NYgcDwsA2rC5mvzd3KsNcKpiSTJFNw/6HCnW6LT0Pje0YCLE02ELrDbc9RH5xhpS2IZJq4fcuDyNPvPT3F9bpn2XIpcwsBxfWzLwFX7a1OYKw875GvXNAQz5SbXZit4gUbOlpyslkL0a8qsUkiajJVa+KECmBIg0OWmwZcYi1K5yqaNQ9y6cz2JhM1P956gvT3D1o3DBIF8AH7mH7/La3tP88n//CHGx6e4ddcG/urP/wv33r2dxx6+nfc9+SDveNtdfPSDD5EwBH/1d9/jv/36e8hkU7RaDYrFHL3+dXS/wbXJMn/57Rt89JFetg2lCBotzPICMzNNsklBSEDfllXMzMP8tXH6+lN8Y1+Fjzw+RHt7ln997jrXZx3+9o8eIVvI8q1nLvCvXz9MsZCSnLhIC9Wii6uIyeDSsrrSPRQJ2ZLVJXlT0aZEZvhWEOmGYUhWnzJ7BlGAT9PkzT9ukZQ3/ugiGoQyCyc0EV9Yo4uqUAeS0ASmriaPIEA3dbnq0kQMgdRvQqyEQdRFEq3/QzVVBQpoK9TFWlfOVRmeJJTiuiSth3Gfu1zd+ejpVPqpUAtjjLCmXnZBSCzSRIwpoYsYviVb+Yyf6fgIglChkF25Uwt/FmUiRSYRJzkj4qOu8MkrB1KIMFSKXdkMVyKCYZy61HV1giJ7KXTV7xuq1VmEI/E8XyXZpZ04UK1gQlkHUUKjrl4kpmnKcI8WyuBPlPBUAn1k5laZIuleUMl0U9eYmavR39PGnp0D9HXnOHB0imNnZnnnPQMkDXjl4Dh5w6PV9OhoT2FqAX7LZboWsv9KC7PuUKo2wLApJuA7+xfYszFHLpOgJ63z6hujPP72t/Dffv39zM8uqV+6TA9ni0XyKRtDa/Hnn/k2w6u7WTvYTbPlUqs28DzZGPfFp3/E9k1raS90sG7jIH29RYq5LH/56W/Q2VUgm0ty+sw1Hrx3J41WCzSDZqNGKpOjXq3R01vETlr84FuvMbSjT64iAStpku3PYJnSgYdtITIWnuPhez4jHQauB6n2NnbvWM2ViQo/eukIfqPM4vw8q7uzZPQ6tmjQkbfoytl05S26C0lWd6XptGFV3mYgo9Nrh3SnNdoMqFbrvHxyksGuIr/8oYdYeukNWkslNHSKd2+jMTNP4Icx+yxEWiltXUcYGjUnYF0iwUnP5cgZjZ6ekJGiTss1WLvL4szrTS6farJmi0Znn05tXKNV0zDMFtUZh0pTY3BNksd2hvQlZlmueJRbNs+fMumwq9y7QdBsukw1U7wxPYBGwFi9i1enNzB3cpH2sSkKIsk9wyaPPqxx40qdyaZNXviMFFMsVR2ePX2DwZFhPvKht4DQefHAGQaKGfJpS+6nQ9kJ43grFaR+GJK2DPZfmefkZI1Mdw+nx+aYWiixqj1LywvwFLUhYRksN1wuzdWUdT/KV8mLnBf6KrCokUxYPPLQrWiaYP+hMyRskwfu28Xc3BJ/9L+/xLM/3k+xbw3PvXSKK6PzLC/O8N4n76enM08qYZJJ23S15zh05AKf+I2/5hc+cD8P3LcLz/VJ5wpMnHiD2o0TdPS08dffucrUQos//LkeqqU689MVRsdKLFUCsobLmm1DVH2b04dHeezuTv7f8/OsWdXGW+5fy2zJ51NfPMvH37eDdz46TKsB/+PPX2ZiZjlu55PmGfncM5RzUz7zZC+S6zryYikLiQgIVApcV04r2UYoV+Hy+eTLBqZ4fRiJ0rIETiJSYtiikgZEJKQLQajo4p7rYBhCuUvVpdVX/T+qoMpXNA9NXWBjRp0XquKolYqMqEAqzuHpepz1kHw0ob6WiCsYqGddhDLRY4wUmoaey+eeikRkmTBW9rX4BRSFx8I4MBGEAbZt4bZcZWET6ga+soOLcyQR8CtUI5BY+feKlqJGQHkqu46rBHiZlgzU6SitwqFKiouV5kF1oBHDwFCsLYUECMA0LQzdwPXcuBdYlqOIlUIrLRLQiNk0EVJA01CoAdUrrTj8Eb9GF1rsFtNUJeTUdIUPvesWEoU8HVmbv/qXw9x9Sw+bhrKMTS1z/toiO1YlmSu5dBalOPfKhSY5v0nS9yiXmhy/VGWmJphZrlCuOty7s5N6y2VDf5YnPvxRAm9lXSe94xJbYCQz9LTb5LM2f/R/vkaz1WDb5iFy2TSWZXDu7BX2HTzDn/3hJ1i7qptVQ/08++MD/Okff4I1qzv58tde4qWfHqa3u8B99+zCcaWdsVmrEXoN7IRF6La4667tHN5/iYtXx+ke7gSkmJttS2F0JCiu7SDXlyO7KoNXdSnP1hhqM8knBfMtgx0b+5lddLh4cZy7t3TRnrIwvCY9GZ0OO6Q9AcUEtNlgOXUSnouuaZSWaswtNxhbaHDs4jTnJyscuLjIxFyZu27ZwBN7Rpj88U/xnRaZ1X3kto+wePgcoSF9+Lr6WelCTpgZ06DkOKSCkJTQOe/DiXMawdkGrSM+q+4z6d9gcObVJpePufSt11izLaQ6DwEWdiKkuuDIN3QmS4/d5MF1y2wpztGRqHNrn0MxA0lLcLpUZN/CGtxmg4ytYyUt2mYmaFuYoeyazMwb7LrDZucal+++WEU3YKQ9x8vnp/jFX/wgn/7j32DdmkE+/okP02w4/PjV/WzsK+K6fpwL0BW1od70SNkmZ8cX8Yv9fOkf/4Jf/dWP8/B9d/DsK/uZm5mnO59SWqaI10Fjy005VQtBEHgYpiH5WLqBbek0Gw65bJp3P3EvmVyKo8cvUS5X0TTBr/9/n2XXzvVkshlmqimGN2+hb/Vazh0/xrZNffR2F2nUG+SzCb773D7+7K+/xi+9/14evv9WAiFIJlM49Qrjb3yHrs40x68s8vc/GOeX3j7IPTvbCRslXjiwyMRSyNp2neG1BSr1gOuXZtk9YvHKiSonZjU++f4RUtki33xljGOXSnzmjx6l0JXh5ddv8OnPv0E6Zam8hMIliRWmlOwSlxfeIPDlNKAul5FpxDItVSkrL8mmIXUoQ3WCRM2glikhsPKSt1I4JYVpVz7YVVmTqz7OUOwwTUj9CaV9aOrhjxZKpBDKcRppxTIfqrIoflxBgVphasqeHRf8qYNHN43YzCB0ue0h0oR1SWQOAlUSqCYoqaGE6JlM6qm4bUuFZVBp7kiPFpGHNYgaulbQIELX5UPe0PBVgEYmIqUVLa63NeSax1XEW5TIEx1W0VgVf1MKkxz9s6HcW6b6//iNolxhkejDTS2FhNK4G4k/clW8skoTmhaPb0IYcYrdUAgDX60EhCHiEz4IwtjTHSoSplB0ymi/blsGU3NVtm/qZuOmHlb1FXjxlctcH1vkLfePkM9YPP/mFLesNrl6o0ZvZ5KSq3H2RpPd3YLOtM7avMfGwRSr+lIslAPevFrnljVJCgkdK6hCGJAb2EjgOghhEFvLtBDbstGtNB1pwa5tA/zz0y/x79/bS8rWuXXnOp578RCZlMn99+6g5kB3Tyf7DpwBEfLQw7dz712bGOht5847ttBWyGMbgmwqidus02y53Bib4+SpKxw9fpGpmQVuXBhHb4R4Sy1Kk1Wq5SbJlE1S1wlcaTFsLTWpLzQY6rQYzBtcX/TYsWUAYST4yeunuHskTSJwgJClSoux+QZnxyocvTTP8dEyR66V2XtxgSPXq5yYrHNmqs6VeYf5ekgNm4Vyk0qpxvuevIudnSnm9h4hcEPaH7wVZ6lM9dIYRtKKsR5CV5cPzyMINXKWwbIfMGgKaprgXNllwtFpX2hRuqCx4WFB1yqbsz9tMHE6IN0n2HCXRmNBZ3lOQ+g+wg/o6tJpugk8R6OnYLCjr0XW8Kk3oSMXcGSpn9OVAUzNw3d9Wn7IcH2crtY8wrAJbY9rowF77tQ5ez1gpmxStH26htfz8Q+8jU/+1p/z15/9Cm7L41f/08f53rMvUTA8EqYkOmiBvHzJS5ns6Rn3bJ7+4t8zuGqQUrnGyKZtrO7r4ulvPctQewoRRIaTkFzCJgQmlpvoysYZbQ5CNY27nsOv/tK72bF1hEbLReiCf/va8+w7cI7//MtP8v733Mc/ffklst3raSukyOUKjF0+TzFt8I633cH42Az/+2/+nTffPMXv/MrbuP22zWhWGttOgm4xvf/7GO4yGgF//u+jVKsev/yWLpXj8Th9rUEqDNi5yqTmWywvVFg3mOTIqMd3T7T4/V/YSPdANw0v4I8+d5T3vWMr73tiC15L4w//8jUuXZsjYVt4jqve+3psnvGjZ5dQF9EIoq9WQNIhZcbGnVARcgNPhvIsxfeLGIAolAih3GZEMN5ogxIo3TWaZKLLeARnXPkcavvjS+dVRNAVqvJW0wWBF8ZlT0KBX4Wi30aO1WhFFkqav9RcVH/MzeYm13MVvt6INz6her5Ez3aNED2VSD4VFZ2E6tYdiTPR6HIzoj2aMHx1Qvqeh2aosichYo6WYZiYhoHreti2Hf8gArW/E2IFaxzRHn0/wE5YsVsgULH5iDSiq5uA02qBkE4GX01M0SEjWf0rrWsR6DHCIIeqZTASnfT4MAjiJKymRk35fQY34ZLDWMMJQ0WtVLY32cmgPkbXcN2ApaU67314FaYRkrLgr//lGA/fOcTWoQSXbixy/tICm4cSLFZhvOSjOQ67Bi2ujzfI2T6hadDe0872XoMT16tcmGjywOYCrUCntTBBsmsEPZlVVk0Rj8rye7QwExlSts8j924iDEO++s1XePbFQxw7dZGPvP9hOjrbaLY88vksTqvFufPXeOzhW/E8j5F1vWgIrl2fZP/h87z8+gl+/PIRXnr1KG/sP8Opc6OMjc+RtU1u3TnMmq52Bgo5dq3p4NzhMeo6pAopWjUHM2Xh1B2Wpqp0ZHQ2dVtcn2vR29/J0GAX33/pDPvOL3ByssWJ8QYnblQZXw5YdAxKrqAV6Ag7hZVMk0qnSdg26XSK7o4C6XQSocHiUgnP8fjlj7+F9PVRFs9cwkok6Xj0DkrHL1CfW0TTDWlH15UuFw3VgXxDpIVGyfVYZ8C0sLje9KkkNIpTDgtjglvfqZPtshg97nLjpE8yZ7PtTtB8C79icH7e46VRk/vWCSxDx/FCAk1AqGOZgrE5n+9eW0NNFAm9JtlMgmorYFX5Bu31eRwtwe3vNLl4qslQr46jGxy6ImjPQqF3Fc//5A0OHT0DCI4cO8lHP/Jerl69wcL4NdrzaXw/wPVDab0NAyxdZ3JumTsffZSerjaefPdHef75F7l91zbWr1/HMz96mVTQoJhOYOjyFq5rBmPLDWbKDSxDxA9RyzJU5kojlbR421vuJJdOkbRN1o0M8paHbuP977qH7dtGOHTkHK/su8z6LTtoNVo0HJi+dobS4iLlSovPfuGHrBss8sn/8CjJVBYz30Y+X8BKZZg9d5jyxb20d+R47ugS3zmwxDt25tk5ksYtL5HJJZmYaNJlenQWU3iej2kbfPNQnZ9cDvjjX9rI8NpOUpk0X3r2OkculfjMU4/R1pZk74EJ/uSzr5BOWfHqSNdNxZ2SNu+o0iKuqFWO0ijIHAQBoRauEDLUwzlCuMubuUKea+rCbOho6mJu3HRR124yJElgo4ZhmMrVuUIi13VDaQ5Rv4ja1KgDLUJOydBgqOzC0jXqeS6ow0g2q+oxODZ63Qe+7Kb3AylNRD1MEeJKj923WqwdaqqLQ89ms09JhV36gkOVdwhjkqR8UEZtfGEoS5p834tPYiKvszq9dKURGLrEV4RBgK+QyEKNW4GvHF26tMgJJcy7rqvcHnrM0UIxrCK7bTQ9BIRxsYv8n4hdVWgyYyHxAsT7wmjVJEV+Q7b6mVa884tK6F3HVeOdroCRK/0AamGovgfQ9YjCKXEonudjJwyuXF/kjp0DrFnbxvDqbp556QpXr83w+J0DDHYl+eorU6xq07ANjcPXGqzqTHDH+gw1J2Bs2mV4KIUXBrjlMuu7Lf517yIjgwWG+9I06y18p052zS1ofgtNGGpl50saJyGpdBLDzuA2q2weaeNtj+6i6Xj4gcYvfuxJkukchVyawPXwvBbf/cFrTE0v8vxLh/nGt17nxy8c5I29p5gYm6RemsdvVCmmLfrzFlv7k2zqsVhT0OhNegx3GOQNhw1rupivupy9ME2t3KLhyrY73dKpTFYwCBhpMyg1PaqBwbaRXl45cIWGo9HV3UkyZVPI5Si0ZSnmciSTCUzdotVySCQTCp2tK4y2SvUCswtLdLbl+dWfe5jKa2/QmF8kOdRP8dZNzP1kn2rLROlaEs4ZrWn1uPtAkBBQcwN2ZE0uNEKu1gNEJqBzNmTqRsi97xbkOgQTF3xunArwXJNt9+ikc4IXR3X+7783GZ1psrFXMNxnU2vKB7pta9TCNN94o0J1cZFEoRsTaAUB61vjdDSXCJMmO+7XGb/hkUpotHdZPHvMpy2TojE9Cc0aA4UM3YUUphFiprLUajWqk6Nous7lmRKm0LDU7TBhmdRaDrffdz9f/tp3uHbxOkuzc8wuLvHE2x/ma9/4HtmwRWc2IfMyQlqeG67HdMW5CamhHj6GQdNxGBzo4vqNGZ574U32HzrD8RPnWS5VqTYaWJbN5OQ0+04v0NY5iOe0cDyfnFjm3j2b+Oq3X+LXP/ogd+7egIdBrqeLjp5eEskkzVqJqb3fIp1JsFxr8KlvTrF7OMlH7s9x6cIiAykPz4MrV0qsynuUWz7HJz2+fsKhHJr88c+vZsNwBz6CxarP7372EL/4c7t5z5Nb8Zoev/8XL3PhSkTdlR3gbqslUUUqFGio3vNItA7Vyj4MQsmQkgJrjByRDqaViUMG7qSGKoQW86dC5ZSLOH2GacQZOkPpCjJSIG/5nuvFJXywos+Gat0uNz+eCvqFceha16IqC+dnSgGjZ2KorMaBmjqjf3bVOi1KnEsIpMr7qUu2rwxDUTBRhQsN2X8rVvp4pXgU+Yf1eFCIdAdfdeiGwU2HDT5hKP9A1/MwLTO2Bnu+p1DGWmxFi/gsEMRjY4Rjj7IiIvqlmjpo0SgVrNjd1ArMc33sREL+gtTXHXV8xEyZYKXNS+4N1ferJi0hBCIQccuhEBGKPpC6iWHIX45q/tOFBrpQrJsgRsR7vo9pSJiZH4R84Zsnuf/etSTTFr/64R38yh88yy88sYadawt84G3D/Ms3TvPrD+dwPZiYqRIaaUY2tbNQF1wab7JxwKdabjAy2Ms7d2n8y08m2DlkYSeTNGauUhm/SH5gHX6zRqhJumfgyFuP06hRLGTp7buDxelpGvUKv/CxjUzNlnlt32mmZpeYnl1mdnKB5UqZVELnxLFTFLJJ1hZ0kp1pLKFRSOrkszq15aoaeSFtQaXWIDQEk8stKmNL9GcFRm8bmgmrdEF/AOMTFRrjZYKUKR+ivqDcCunJG5y7Ok8qZTHQ087s4jiWZRIiu+fTdoJqvY4WamSzSYYG+7hy5SohGraZoFqrYiaTQEil0aRWbbFn13qSzSrj16bwmpDZNIyzuISzXEWkkzIoFq2vNBVOVcJ6xEjLpW1C3WGp0eJ9OfiyL9i77JMpaGw75/HSvwnufV9IrjPH619vcvzFGgsTaW75mMHgziH+oNDkL/91lJPXl/mNdzg8cXsKQxjUXZ2RjiYP3b2Br768hHf2NRpdm7D71qIFGr4Tku71sWwNKxMyM9tiy4BBynCYW67zyMY2Wo6PrsmHjxemeeEbXyOZsElYSc7P1hhqy7FQqmPflEi3DZ1qpUIYgG3puJpJq9Wk2WjRqtfo6LSl3V5I4J7QdXpyCUyzRkiI4/qEnryFJg15B71/pMiurWtZcqFca7Fcb3Hg4GlqDVceQI06zVKFs8feIJ9vI5HOUa9V+cCTd/MLP/cgqWwKp+lT6OgkNCzq1QZNzWTp7Bt49RrJfIGvvD7H2GyD335XF2k7YGreZXufxfGLDUaXQ87Pw4UFn8E2nc19Nu2dWdavyVH3oC1r8fdfv4JhGvzSB7dDELLv8BQvvnGZdNqS5AjbxvflIRK5rRyV59J1HU2srM8lZ0oW2Wm+fI9pig8YhmDZK7UUkfYaEDGkNIJQ5XLQZGZKTRYyCuEToqMbAs8JECIquxOxayoMtbjMLxbFVWLd0I04t+a5LkI5Xy3LkhdslXfzPYmn8l2Jag/inIoMSkalgEEsQ8gtjh8GEMgyLA05Oema1I4D38fwPE8VMmloQUgo5OgVhgGWpatynzB+mApNQ4TyRaxB/DCXJVOy+EfokubrhSuthZrCFBuGQeD6qm8kjNEpjtOSrgVl3RWq7CQModVyYpRAVAIfxC18Bhia6m4XsZAepdAJiXvXQzUqWqaJ50o0QHRzkKA0Pc6KSMFKivuE6ocbSL+0ZGopR1ZUO6WcCWEIrtMiDCGTsfjxTy+xf/9V7rxjNe95+1b+5t8O84/fvsTf/n+7ePL+VZy7ssT3j82QSRhMzPtcurKImckxPGhx+pKLPxGQRsOfWOb924scvTHH9w4u8YuP9rFY9lk+/QqpriGJhVDVpkLIl6vnaew7dIGzF8eYm5lnbnqWcqWCF/jomiCdSJKwdfoKSbb3FcgYPvkEWEKO9K7j44SCqdkyYxNNvBAm5qqUXcFitUWlJRsRPTREIonVqPA/h/vJ2QJBwOquLB0pKeT5YcC+kkbNCyi3fAY7bKYuLjE1X2Hzuh7eOHoJx3XRhYGmhZQrZe67/04OHDxOe77AuvXDnD5zBsuymV9YxLZMwhB2bdvMiz/dh+O4bN+6BvfKdZqVGlbCIr1hgNLJixIIapm4jqoQUD33vuvFDZNCOQEbjkfGNJivOKTDkA8VDL7kWbzZgkwigEM+jSWNJ37D5Mlf1njp6wYeAVdKPquzJp/8jRwDmRpfeNbhfz5d4dVTLX79bTk2DlpYukAYKbLrh2HqLI2ps9RDh6DpQguK3QG68EkWdMozLnrgkUvK1UHLCWk4HomESb3WxLQthnuy6IHGmckStqlz+/perkzM0Wi2MA2dRtOjo5Bl78uvMLVUxU4kcLwamUSCl17dz0AasimbRsvDxZW7ej9goergtKTFfqCYYEtfjtD3OTlZwQ0MvPkpTr4+STqTIplLsb69i3vvX4sndMx0Dk+YLJbqnLs4xpUbM4xNXafRaPKff+8fyacSDA12Mzyyhm2bR1i7uo++gT7GL5xg8cJBOvq6uTy9zPcPLPHg5hQja4qYNLG8Kb57yOX4tIttaKwuWPzyRostAwleueJzecFFIyCd0rk86/KV5y7xB79xP4MDbfjNkM995TCNlkfOttECEdv2/SjF77QkusPzEMJUa/QwFrXlc0JuXWT8QK2MtRXd1rStuPhMlktJR5TQdQwRieAAXryKIo4bSLKHXGrIuECr1YrXV0RgWqFgsjdVjUdUDsMwcF1HRjJu4mUFKonoqSB3NJEI5TQzLVOuApXt2ECo/nMp6kehRl0PY8dXBGDUM9n0UyjgXqgIjRHV0dBlkCpGikQIEaQ4Z6g2vGjfF7V26aq9KsJ/hOoPtJQ+gqbF454sYfFlNsMwY7xIGEr7LSoEFrNqQvnDMpXeEcYMr0DtKP04yS51FT+2xmmhFh9Ohmmo0c9X4MZorxmulGWhClsiVMpNNMsVEq9+U4m9QRASA9EMQ6facKnXfZ58fAuJXJZi2uQv/ukge7YW6G1PsGNtlvMTTd44s8jWwQSr2iwWll0621LgtOhoSzBRgomynAZTusd3jla4Z0uefCqBW18k1AJS/ZsIWnWFJ2iSSSb463/8Fl/50neYvXEDp7JEIekxkNPZMpBjY3eCHUNpBlIe3QmfvBHQLNeYW25y/sYyZ8Yq7L24yBtn5zl0dZnjN6qcmawzVQpohCYIk4SdpL29iOa7eJ7H3HKLNT1ZDD3k8tV5OhI6auFHwjSYqTnM130G8oKOpOD8VJ3O3k562vO8uPccbcUClm3FfoC3Pf4Ily+P09HZTqNZJ53OsHnjOmr1JsVigfVrV9Hb3cH+IycwhMYvfeB+9INHacwukezppvuxW5l58SBOraGcKWod4wUrdaFBiK8cgpEe4gYh+XSSiuuRC32GEwbHqx5XWyFr8zrJJZ2lccHQTp9VOyzW7Gzxg/MdHJ5eyy3tCzy4O8Nd25PMTDV47XTAC4eXuDAZMNQWMBqs4tJyEUNPEFoFGnPnELUZ1gSCNXsEff2Ca2M+Tc9k3bDN0z+tk7FNhjtTyv0nMx3gK7EzxDbg/FSFS5OLBL5PPmWhhRqGAX4YYoc+3Umdka4MG/vaCMrzHN/7Bn3FNIGmcOF+iO9LXtpMpcX1pSbb+/PsWZVnsdZCFzqb+/JMlxtcK7tMl12uL1SZXa4zMTnN7NgYizfGCJamCSolsiJgTUeajas7efi2ddy9ax0bRgbpLGYpL5U4eewcz7+wl69//3UuX77MRi7QcnzaO5P8+b9fx9B8PnRLhnqpxvL4HHqjRcIW7OgK2NPl89hDwxQzNm6rxYsXW6wdSLF7YxvCSvCnXzpHowWf/dTjpDM2r+0f5U//7jUyGYvAi9bzirHn+6qkTubOhMKkS9MOqhRPyJu4JiexIAjU80VNCL6sKgj9MA5BRxqvRKJEwMIwdnMFgU/oc5NuKT/WNI1Ye45ChdFfkcgdMXEN04wRTTIYKN83pirZi/Jt0X+wUlOuNj2K3EwEWFTab+B7sRNNNxSdFwhcL6b1Gsoirmcz2aeiuHvkrNLUGBeReaPYvWGYNxWkhPGJGDVqGbpB6MuWNFTtbYRCl0wVgee5hMiPlY4AaaMLIxtshBtW9YyRaO358oAwTTPWYyIgGQphHDkFpCVYahqhmkjkTTNQDiyVKldTSgRmVOUHK7WOYiW96QeSMxS9gCRwzVrpIdEFnuvIE9001QHokbQNzl2e54E7VtPfn2f9qgLPvHSBExeXeOL+LtAM7tzaxivH5unO6ty+Ps/kfIPBoSJT48v092axjZC1/SkG1g/S11Pg/ESNczdqPLarDRed+uw4VqEHO1fEbdbRNIFl6Tz9rZfJOUs8tLWNoXaT7owgY4S0pQTlSp2rExVOXF3m4KVF9l5ZYu/lZQ5fWebyvMN42aXU0tB0O34hF4t5MqkEuUwKV/Wj3//gXTiexsLCPL4XkEnoBCLkxKVlhoq20r/kQ3u55TDfDEhqIWvaLEIB8y2TRx/YyfM/PYll2di2JW9BhkVPTxdzC8v09HZQyGeZmprl0uXr5LMZdm7fQj6T4tS5C5w7f5XevjY+fMcIS6/sw3E8eh/dg56ymXnhAJpl4PshhqnFGB5DBUFDpPVa00JCL4hJqWEAKUPW8nbYgkFD57gLY27A1qIOyyHT1wV92128Nvin5/p4vbya164b9FpVtg20eOyuHAmhcfBMnbob8qOjLgeOTGGaHnY2g5EpkOka4vrEFF2ZMg88mCKZ1pkrQdAMGeg3+cprLTqyNms6s3Glgh+GBD4qwCqdiX05k6xtkDQk9A4BSVuGDIWuoxtyV29YGklLUMikqDYdKfpGphchL09NZYjZ2G5zcLzCbLXFTLlJwrLImDDtpwjtHEu+xXjZZ7IKY5WAK3NNLk+VGJ1a4NypC8xMjFFZmKa5MEPObNBpugxkNG5dV2TPhi7u2dzPcE+GfnuWrN6iszPPC0cm+PJLM/zZL6yiM+lz/eQ4QbVB09PpzIakLUHv5lX4yRSXD1+no7/Avx8q88EHOhjq7+DQ1Tp/+s8n+JPfeoAH7h3BaQX89p+9xKVrs1imoZAccrMQBAGGZcZ6bUT59pUW4fsSABsZeuLyPS2qn1iJOkSpbpTo7vmyr0gCDMO419yPW1aVIfRn6rVXHKzx80lRQOKHu6LtRgl3TeVKNHVgGapdUEITddVSqP4+ytipfqYo86aF4IeyzM8QQvEIA0VvcOIWWV/x0AzTwG1JN62eTKWeEkLWIMZBD41YaAlVlD7iYEnvsuTvBF4gS0ZUAYofBuoHHaFOpEAjKby+staCIXQVwJPJTV2XJTgRfsH3fQkuI8RXVYtCCKVraHFgJurplmUoK0EgGXo01C8TudsV0huvR1NJZIuLMO4RzV6lRCPbr/wYX4VniC10ppqSoniMrily3U1fWxQOankhlWqLdz26BsPW6O9K8+efO8KOdTm2DOexEjb7TkwzM1thR6/J1JLLmqE2FuaXMXQpbM3MVli3sYvugW6GOy0+9+xVBjsshjttWk6IWxons2p7TPXUdcEb+8/SWp5n26ocB67UmC4Ljk20eONCiQPXa5wcrzJZ8VlohHjoJBIpstkMhtBJJW3S6SRBGJBIJwl86bxzHY9qrU5PTycPPbiH0nKFRCbJhYvXSOkBbq1K0PRI2YLjM1VyliBl6gS+h4PG6FITQ4O17SadeZuzEzUee/Q29h26iBfIw3nD+tWsGuplYXGZof5ekqkkY+PTpBI2+VyKyelZZmbm8DyXq6MTjI9PceedO7gzqLN87TpGKsOq9z7I4uHzlC9cx0rZMX4mDCQKIlBOGfk7D1W9qAo/ql2vZVlYAhabDv0JndWGzv6az6xusjmtURrX8Jo6N5I277nPIh3W+dHpFHsXumnWXfasbvHATuhvTzI5Dx98R5FNqx2uHD7FxaPn8Zwa2UIaWrC6Y4lHbvXIZgUzc5AwTNB8vnXAY7CYpDNlyIuZ0DBuylJFe2nD1MhYBsmUoWpM5SGDcldKR2SArcsLYtNxcVwfTchJ3NBVaZGmkbR01rQluLbUZLbaYl1bkmIuQakiD5yKb5BJ21iGoJjLYpkGZiJNS7NZdAymGgbXywGXF1wuzTaZWG4xMVNiemaRytIic9MzLM0vUJmbp7ctYOPqNH4IdcfhU98YZ8+GDO+9u4vFkkNrsUQmaZJICXzXpbCqi+61PZx8+TzDa7O8ermFr5t88OHVNEODP/zCOTLZJJ/+g0dJWjo/enWUv/nCPrIZO06SCwWKlK2CgWxCVZdCSaIQscNJ0xQHS93+Y4aVcisJXSBYKa1baT5V2JAoGIis5pWVs6jMnbLFqmlG5nCEZE0Fsjcp+hqjaSFQfUrRnyNFf/0m4kagmgR1dbBEEw7xhBJVWURrMD8IVDme2hepwywOwWkSlhoJ/1FWJvQD9Gw28xRKy/B8H9M0JMsqYs9rKwAuXzlgoqlB8laIVz5xkjLKYqh/FyE/wphGGfU0E8PJNCVMRQyaaNWlkouxXgEyRBPt5TRN1pxaptw/GpaB760EeVzHVVgWVWavyfWWYcregUAtIGXK3lC7UXlCR26JqOJRCB3TNhXiWY/T+dEkFIv0KjAZlWnZhuDcpTnu2NXLqh6b4f4se49M8pP9kzz5QB8iDJlaanL84hL3bUgxt+zQlrdYLHkU8xYdOZ1L1ypkqNGqt1i3rp9qpcmXXpnmbbfmSSQt3GYdv1kl0bMRA49UJsOPXzlGbXaGtZ0ZvnOyxpJnUXVkJW7CtsgmbZIJC8s0JA1UkUc1TcOyEtTrNUxdV0n+ENs0MC2T3p4uTGGxfccGrlwbZ+/eI+SEy0PDedb3ZulJm2zoTpPSNU7NNemyNayEhROETFYcfD9gpMOiv83i4lSF3bdtY2xingtXp0klk7S35dmyYT0nT56hu6uDGzcmuHD+KrfdtpOLl64wOz9Py2mxtFRlfm6eVhDw2Eg3GxcmqNZb5EbW0L5nK2M/+Cm+21LuOi1mnEn3oHTwRWb4MILxaBqeuwLytISGbQpKrsegoTNs6rzW8BiremzK6VgCvn0xwzsfaeNt2z3yosbZ6Swv3Wjn+A2btLfMk3cJtq4KOXDCZ9O6dn75F3rZsU7nyolrnD50BpIm9YH72De9idKyR8pbYCBTZ3rR4ztHAtb1JCimpJ19se5ydGyZRq2u6hJ8bFMnl7IxhMC2DMJAI2EZ2KbANAxStoEB2IZBqdai1PI5OLrMTMWhL2eiGxLv76mKA9sUNNyAqwsN8AISlqDS8OktplhquCw2Q9KplLwsCmX/9D3aCgVMXSObstGFgZVI42Cz7FhMlDXOz7mcnfY5N+twYbpB0/O5c3OWpaUaXe0p/uGZCV4/XeLjD7YxvKaX8lKZqetLtKehVndoG+qif2MvR35ygWxakOtI8fmfLPBrTwyxZrif18+U+MxXT/LXf/gIt+3so1L1+K//6wWujy9gmXqcqpYWfl1lMzxFsg1kelthlKQ2YuIHnpyIha7cWlI8jyIEYRDE7s2oAz0SpaM8iKfEeZTTT1c67s0leZ4XYNmmzHfEluAwnniiGEK0eTF0I37WClX653vydW6oZ7jrurFVN7IbR6s51BrNME0ZFhRS64lcXTLbZqiNC+iGhudKmGToh/E6X08l008JZU81DKk1BKo0fQUipkrdIxywul17nlwf2bYdc+Ylc0TpA4Ye3+Z0deOPe0Aie5m2ko5cGbVWOkGEKokPYovtzRh2+WLwfV/lTLQVeKOkkcUnte+r/Z1YaTBEBWlc11GFUX48PkZTRJRNCRUDzHM8dVJr6pQP8XwPT2kvui6BkbHDQqVXm67H5HSF9711PYYh2Li2k7/4l+MM9CTZsSZLLpvhp4cn2dTh02xAe0eKhhsiCPFbDpMzdWZKHvmcgRAWdwzn2HuhxPnxGo/d2kXTNXGWpxCpDKniALrw2H/0IhOXr7NtVYGzcy6e65JOJ5RtMMROJXFcD8dxMAwDyzSV7TlA1yGTyWAnbZYWl7j3/jtIpVK0Gk1237KDN/Yd5OTxc0zPLFBtthhMQU82SdWFuuvSckKylmB8uYmlC0wB+AFTVQcPjc4UbOixKdUdzGIfuazNwWNX6e7ppFZp0Gq1GB+f5eTZs0yMT1GpVZldmKPZaNFWKLBhZC1u4HHt2ijJdJoPdOgkGxXqNZdV77wfv+Ey8+phNENHGAaB42MIQ+WZXDT0FdOE+j1KPUxISGYYgCcvAoYhqNV8AhEyaGusReNHDYlp0Q2fb0+nOT7ZTV+bxwcfheHUHGOzJsdLAzx/NsXYZJlbV7cYKGgsT5WYvuGx67YOPvGh1dy/KUlzfp4zRy4xVUtzJnyQs+4dzAR9nJyEsak5BgtJ2lImPhrTyzVKGNxy912MVVxuLFS5Pl9mbL7GxEKVG7MlppeqTC5WuTFbYWy+yvhijWvzNS7MVJhzdY6OLlBImKxqS9KWSeCo0idDN+VbOwgoNT2afgh+SHvOIqEbWELj8nyNZCEvA8KhRstpYlsmpmFSr9dB06jXm9iWjWHIA83UIZNKUSjkMe0URiLHREXwlt1JNq/OYFomb55b5vM/nuUduwt0ZzTWdetcOD3G1HiJ3g6bzpFBij1Fzr1xCY2Qnbu6+F9Pj/Hg7T08fMcqGprFf//MITav6+J//c4DmIbON354gc9/9QiFfEL1C62I0NEK3lQOKgk2VD0XIgrf+SqkrClBW1eacBA/y3zVIR+1SoTq9o+aan3fRzdlR4eGdKiGocxzhMgO9rhUSqzUI6jNuFyjqcsqKtAYPQd9NTVpShsOIhsw4AVebOHVhRFbcaPDIaJ3gLaCLIlrLXRFUg/iIYSbALu6QryDhiFv5NL9pGlarDHIlj1DnaYqOKgrG68XoKNj2bJ4KvA9CQSzdFWRKD+fr/zRjutCKE+7qK7W8zzZBaxGKj8OqghCvLg4JYr1R3tDoQkwwljoEppAtwy5a1Q3gGgiikqjZNc6cZWk0IiFtGhlJcnDgcSdhBqo1VkYBPHkEv2l63K95nluDDjTDZNAVfbapilbxDQN07JwPZdMyuSnB8Z57vUJnnx0NbffPsQnPrCDT3/5FPdv76A3pzHSn2Gp1sS0DRxfo5g2qDYCegbzrJ6vcX0p4LnDZapHr9Kez/DItjb++SfjvHpykXu3tbOwrFG/vI9EcZBCoZ98yiKRssinBZaAVugTEGBbBk0CqvU69WqDQi6jqMoGQShrRqvVOrfs2o4QBpcuXWPT+hHGx2eolKuUK2XsZALd0ElYJmghlxaWmW6U8MOA3qzBls4UCB1T16g4PsWkjik0bMsk8GGyplFvBnSnNC5dusHtuzcod4lGuVxm/6FpOjs7GO5eTbG9SEdbgVK5JpcFWkg6mebi1Wv4Qmd11qLXaVALQpJdbWTXr2H0ey8RBj66MCGIGuUkBltD3tgig4iupvUIZqebah9uhFiG4MZiXYrtrsaScFmbSvBbAj5X8Ti66NDRrzFZgT/49wyjS00+dC9sX7PM73+twY9KBb55rotbctNkbYNsu42ueRzdO8HhnM3bt3Tzt3/cx/jUAv/w9HFeO3WapYE7ebVtM4bWha2fotb0pNirGbTlEhw+PcHD9+3hV/7jx1moNmi2WiwsLjE/v0S9WqNer+O7Lp4fYto2uVyGfKHA7p1beH3vAf7+bz7Dves7abk+XhBiCgNCH0LpxPEVLy1raugpE9fT6M7ZjM9XaQibwY42FpeWIJQPQNdxVyCqWkAqlcJ15brT91wsy6babJBUWPJy3eXhjYKHNyVxQ51So8VffOcGtw4n+MjDnbx+cJra5CiXLs6zqtOia7iPesPl7GvnybRlWDec5k+/coPO3gIffscIyWIHX3rmMsfPz/OTrz5G0jJYKnn8/b8dwtCJ8xa2bcuOjiiv5brS+BNprIpSEWUvNHWJ9FwfyzRVcdRKojtKhwfIwLOjDpkoNG2pEHQYSHdkBCUMwpXnies4IFaC1fLKqFoJ1WERZTGk3diKV2W+H8YBZyKQbajQRqEAXZmVTKH6TCzV0aSetdpKADta8eqGxNgLXY+t7RFkV1P6r2macssjwDAMmZaUig5xOEtXbJcI32taFq7rYpomlinXXfKTaDieh2XZ8Rcv8egiTncLKQaodr8gdkaJCCuiCUJNrqk0NRbH3eSatLt5ngfBSghGJtNNQr9F4PtYhqVWTfLE15SnWY6lGoYlHV6amn6iIHxUfas6czB0U/H7JTAt+rpjF4WuxciQiMkVqLIYOVoKXNeNacXy7z0M08ayQv7qH/fz6N2DJFPw3/7jbXzjmXP87Tev8Ne/vpmRviTjk3XuHEkxeqPCHVuK+JkCtaUlkgmdzoKJWGiwc3uWacdnsSkY7k3xmR+Ms6HXJpdJUKuWWb74Kk74EDdGJ5ireDx9YJGJxQamlcBouXiej+e5DK9Zw8jatczNzXL8xEmazSYJ5RV3Wg7LyyV2bN9KNpuj1fJIpmy6ezop5LM0ak20pIWvVnU5W7CuqNOez5KxNOo1Fz0p6CgkmZytoheTJJOCTNkjaAZUmh4LJZe+tgwHZ+YoFG/B1KHZbPHAg/fELhddwOTMHFevjXJjfIrlxTJChHR1djAxNgmGwe0FnaQIWF5sMPTQTkLXp3rhOsI0VGhUQ5i6ulhooIipmtAkERVpO/Y9+WZRDchYhsZstUW16WMJg1YQoAfQwqOmW3ysT+fVuYCzswv0ZFK4yXb+8eUOLs02+M3Hlvm7/1DmT77V5Pp4he5kSF23qbekj94rasxqVf7pVJkfZtJ0lndz/1vW00xeZt/hI6SqFxG5QXTbpuy4+JrAcVySuqA3a/G3//xtnv7CZxgcztNwJCMssoUSyiIjx20S+gH59nYQNq+89DJff/pr3DPSjeNIb7/rS8Zbypb9FhFhopAyyYQajaZHJmFTb7Q4MbXEuk0bMSwLp+mTzaakscQwqFSrIDSSlkmz0SSZTjI02Mvc3KKk/pZdapUqumHQlQn5wO1ZKo5B1gz52x+M0fLg19+1Gq9RRwt9bkyFLFUCHr2twMJMmfpyifYOGyup8amvTWIU8/zeL23FN1JMLTb57DdO85H3bOO+u9agCcGXv32C42cnaS+m8CKLq+PIB6Mu6d9+IBkckStJZh+0mAweIZFEVLAkBIbQ1eXVljZZXdpdXdeNL6sikOt9oQt0FdTTtAgou1L0FAQBQv3OosNXR+rJIVoshodRS67anGjRyt2/qRY3DCXCRDGvJPQxjGkckfknrtPVFXtLKPE+DOIJFCGHh0jkjUCzsiPHo9VsYicT0umVSKaegijaflMhuzr6IqxxyMosE6UvhS6LmWJBOrgpJen70gVFGKNEPPWLicJ80YQShNKCFtluUTpFTJaMqLnw/9cMKA8IqWVE7gUgVOI4IabCjaB0llCRdUPVHR2JT1JYkgl5TYlOceNecDP2QgnqQq4/JN8myp/IHWSgRasPI+6SDv0Qwwi5NDpPX28nt27rpK2YptVw+fy3z/L4nYN0ZzWe3TfNPSM2N2YapNM2ne1JbpwbJwx8XGFwbdqju6CxZV2RVQW4pc/kRyeWuDDV5PHbe6g3IWc0+cenX+SFvdcZ6MmRT5s8dksP2wbSzFVcSlWHt7/9Ubo72zlx4hTXr9/AMi1cr4VtmNx73x10dbUzOTnNXXffHqMbsuk0ASGFQpGz5y6ye/cOdKFTqjfoTQWsLdrUWz51JyCRNNEJmFqsEpoWg8UEgeOx1AqoOD7dWRPhuWzuzXB+skLv8BrOnR+lUm2i6zoTY+PMLy5w7PgZJianqZVrNBoO6UyKxx65j5brce78JVIJi4/02Oiuh+eHjHzgLdRHp5g7dBLNsOILh6b0jajAR2h6vM+NLgy6IakDQr3UW27AdMWRtllDwyPE1g2uuD4TYcjb2yzu7U1xtVTj/EyZbDFHLmdx8prNwRsptvY3+eCdVVbrJWZmAgzDpel4mJbOcdehYZkUbEG5GvLTl3Zz6HSRVLKXD+24znBbg+FcheszDuW6S1/OllNvCP1tOa5cGeXN05fYuesWent7CIKQZtOh6TjU6nW5OVAw0Wefe5G//KvP8tUv/AvbOlIUU5ZaK8u2wTCUjCbbNLmyUGey0qInm5TraUswtVDj0PgSG7dt5oF7b+XggeM8eP8eZmYX8YKAeqPG1i0b6OnpZGZmHk3XuO3WHezevZ2+/i7uvWc327ZsYMOGddwYn+CDOzyGh9rIpG2+/foM//LjcR7bZPDAhjRBGHL+8iKj16tsXZOhmBZoThPf0Lkwp/H5V8uMbO7l9z+2Cd9Kkk6a/NVXznL66hL/+lfvorstw+hknd986kfK06KpwroVrTUIgpUKB9OI1zehRlzNEKjLarRKj557NzuZNEJVkheVZaOK9fSY2hEXMWkrQexAdXlomuyw19S2JWJyBcpeHKFDuKlcSho+iOspojWV0FB6bDRO3xT0VoHCSESPpg/D0GXsQK2movpb86YzIDInRZ0hYSjNB9HH6oVC4alQFSupoxLPlyMbaLRcR2UdVipiwzBqoZOiu6csYRHgK/INR/wYLVb0iRPu0jlixvF4qVmEcQ5AVzZLKXp5cQ+IUIEcVGZFi6i6aiyNQIlRyRWxe0HueePymNhFIaGJvoJHRux9w9Tjw8xQP1zDlA1f8vtQOBZFyYxw9pZyTUhEvAJFRusT5co6dX6O9711E6lskm2bOvnWM+c4ebXMr35gCxevl5ieqbBjdZKFBuQTgtnJJRK2TmnZY6kmA2Z+rU694dLbnWNLf5ovvLpIPhmyZ1MBLxSs602wY02BR7a28eC2HvZsH2RzX4J63eHg1TL33rGDH73wCgKBnbCwkxYf+dB72bNnF3v3HmZ0dJx8IU8+n2Pzlg2cOH6G3r5uWq0W8/PLvOWxe2m2HK6PjiIImVssM1vzuL7YYLrSoCNlkksnaPkBx68v0p7SaE8ncIXORNVlqC2J6/uMdCSYWW6Q6ullYbHM5GwZ13XZuX0jnZ2dZLJZZmYWGFrdh+u65LNphgb7eWPvAebLNXa2p3iiy6JSa9A2PEj/Q3cy9oNXaC6VMZK2GsulgyeyqEemkTDw5KUHeXsPPB/TklkeUxfcWGgQCgiQtk18jXk34IQX8t62BC1Nrgce7sxybqnK1fkytmmSbUsyUzZ5/kSatKhwx+YQEQQsV32EoYMOZwwoV+pk8iYLkx04C/0U8zrtwuE/3X6YNe1N7t/lcXlScPB8g+6sTtKSPfReCOv6CkxevcxXv/UMc4vL2HaSTC6HZSbQdJ1Gy2VsbIZnnn2Bf/j0XxPMT3LbcI9aR/uSxKBK0ixT2kmXGy7HJ8vUXI3FWpOZcovR5RYTdZ9Mvo1bd26h2JYjDEN2bFvP6TOXCTyXnTu3UGxrp7eng0TSYmJimjvv2M3JU+c5ePgkh4+d4+LFq1y8Osm7tnvsHNKxrCSnrlX4X1+5yD1bs+waSDCUC0mbGvvPVig3HAJD58ioy/4bHm9eD9g/BT//jhF+6clhKs2QpG1x7EqZP/78Sf7nJ+/jXW/ZDKHgqb9+jRdfv0AmbckLqC5wPTfeIkgCtxc7nAJlCbWjRLl61kT8vcjQY6jucj+QN3pPuaOCqCIg1DBNPdZn45xYdAhpqCIogWlZ8cfYCTvm/8l1mqccsCGRudM0jRi5Hj3/onW9ptLkgdLx5OXcj/FOstLCwzQNNH3l+RwdPtGBqmlhTDCPgobR+k9W2irdWtdjMKSeTqefiiGBN1XCBlFA0FDgLCJmkDqVXV89ZCXmV7+pHztyCojY3STkCBWuFKPouvSr+54np4Q4oOjHnJlIuDEUDVdE6jTE3mhC2aPuut4Ku0dpHn4Y7StddQBqMapdCHkbjUdVxZaXFmFdjn+RA43I7aDEI2UHjCBjYbjC42KlxidGqeiqeSwMNdIpm+mZMrqu8/C9g6SSCbK2wae/cIBNIwWe2NPFN1+ZxAocUimTdNqm5Xjk0yaVkoOVMHDcgEpLoz2loxkGq3ty2HrA3z03yX1bChQzOkbCprMtydyyx6UZhzdPz/ONV8fZe34Zw9IZGV7LnttuIZvLMDM9C1pIW7GdM2cucvnKdar1JsW2HL3dnbz22kHGJqbIZtKEQcArL7/BxOQUR4+dpNlyqVYqeJ7PuvYE6zoyrG1PkbIlInq52aJ702ZcM8X4jRv4QchCU6M9baKHAR0pSNo6M02DbC7DkeNXWbWqn507t7Dv4FEG+rootOfYtWsHzUaTeqOJoQsOHDpGYNo80Z1ge7tJuVSn/9H7EIbBxI9eQ9i6zHVE1QHq0NBVb7V8PRsqd7Ti2pMFZhqTZYdSS7ptAuTrOgg0XnUD3pJLkNA0QgNsU2c+MHil4uO4LUrzS3iOSyKXwXE1Xjljcd+qZfq7LLKFDK26x3TT4Yot8B1wwwqlyUGcpXa80KA3XeLedTcQ+QR23qTpaDy3v0XC9OnIJPAVGdUPAvra8xT0gFNHDvHcD5/lxR+/yI+f+zHPP/sjnvvOd3jxu9/mxumjbBvsoLuYw/Fc1fOjeiECuYK2DUHD8Xnj2iI+FkPb7iNM5Gi6DYxUmq7ODkzbpq2YZ+vmYVw/oNiW49iJC9i2wSMP38N3v/MszWaLW3fv5NTJs9xyy2ZOnj5Pq+mwerCPto4u1ojLPL5Rx0xlWa62+G//dIm+No0P3plluRwylA946dAixyddLpU1JqsBtm2QT+s8sClF2U/y1jv7sWyBlTAwLYM//MJZCoUsf/cn78C2QvYdmeZ3/+JFcplE3EYawVojsriGpuy40uhimdKl6bZcNC1Ujaph/AyKmVfqYaorC2tk8Y1eY4ZpoIsVnLuk6xJPOGqNE7/OIstsdKn2vQDdMmKtTihar+xCl24tVFBZU/y2SDmRk4IqmrLMWMyPu5MMge+HsalHaLoyToibarxViVVcPMVNB2xUkxGt3uTXohoJV5qxQrWa0tQuP9IM/FCO+7qpx55iP5D2tng/psiTUu0P4lVRSLjy5oxqICOKbmz7jTDFOgKFRxeyvyEeGxVJFbWSUK8Itb9UnmtlCw6VRU9E4l78C5QhRfkm0lewKSHopojheoEXxPW68fqKQB2scgTxQz+uwI3hakG40tao1mxCzY5RjsS0dE6cmeXxu4fozJts2TjA4VNTPPfGOD//1hFG+tP84MAsmuvSlTNYXGgyPJxjca5OtebTkTM5MutT9wWua3BpssGGvjQXphu8dKbMO+9oB18D32Xv2RJPvz7DlfEyve0Z7tjWzcXxCmvWDrF//wFee20vOhLXcunSVcYnpkgkLGlQC0OKxTxHj57GdVymp2e4cuU6QRBSKlXp6upg3frVrFo1yNTcPB1WiKFpVBTEMm3qXJpc5oOf+DC//3u/TtPOcnVyllq5RtOBQkIQuC4b+rJcmW2QKLRz9sIYxbYiq1cNMjY6QWd7B/liHtOwuHr1OlbSZnJ6irmFEoWUyYe7LXQ/ILSSrH7Po0z8ZD/V65MIw1DBQOL6Tk0REiQLS4vLxqJispbjYhmCmhswXXExTREz2tKWwctNn60JnbVpEycMEUDStvmba0ucXKyzLZdgUGtxdnYZ020QJDq4pb/F48M1KjUwLY2hrjzjjsvJRp1CJo0X+MzfGMQUbTQDnW2D89zaPsvkxSZtay0SKcFPDoZMLVQZbEvJN7x6iPkh2JbBYGeenqxFUbik/QpZv0bR8OkrpmjPZSWWxpfYek2tqoLQxzAFBlB1fPaPLUmO1NA2iqs306rMYlOjrVigVGsS+D4b1w1RKpX5/jMvs3bNEO3t7XR1d+G0mjIHlE6xfv0aZufnGOzr5sixszRrNTL5Nj760BC3dUzT0pOI0OePnr7C6FyD33uynaQIuThWZ2m2gh9qbOqEjYWAJ3YXeGx3J3ndp9CR4cB1n7ff2YmPRrGQ5HuvT/Fvz17lX//ynaxfncf1dT75R89zZXQBQ8jAaJSRuHkaiC+OQlfPGS0Wq31W+Fe+umh6nq9It7IrPfBDQi2Mrd5R5sLQDQL130fFexFgNggCbMsmVDURmpIBorVq1MPuu15MuYgmnujZE7lShUqRu47zMzGKUOkgcZ+6oqajnK6GalmNti+6mjBiOUGh4OXUJrUTwzSkHqzovoEfYJq6ymRr6Llc9inflztuyzJVj7kKtahgjdKspVofq/bhTW2Crio0kQ/eIJAWtah8RFNR+cgWFukiUVF9EKh2w5sQ6fFUoE5HQ5fJcKflKOZUqHbYgRT+dRHfImMYpJpUdFXTGPmeo6BV6MsdaHRaezfx+qOC+9CTPegEWpwl0IVyhilODso8EIYStIZgBasSWYhVZaimSwdQpdpiaq7Oe9++GSORYtu6Hj7zpUPUmwEffnQVTrPFt1+fZlOvSdPV6O5MYmkekzMN2jttphoaVyab4HnMOgaXFj22ri1y9FqdS5N1Hr+1g1JLZ2RVhttW2bz9tl7eftcAfZ1pXjwyzaqhAbo629m4foS77trDwtIylXIV3dAxLVPt1JtcvnoDoWskkwna2nL0DfSxffsmtm0dYdP6NRTyWTzfZ3xikutzJRabHrM1l6TQWNWR4epSgz333cu6NQPcunsrg4P9vPjMCyw0oK9o4no+azuT9LdZ5Hr7OXxylHQyyab1IwwO9TM1Pc3Vq6PYtk2ptMTicoWr18dpBiF3tCd5oCBYXKrRfcdOcutWcf3fX0ALQ9nprOgI0egdvZyiET0MQ0yhK89fqPAMcH2xiS40TF3DDcHWdA7WPbKW4La0AbZAGBrthQzfmq7xwlSZfCrBnpzOOwsWGUtwamGR6vwMg9kWb9lpous6lbJDOmFytNFgyncxdFlINXN9mISZAC3B7oHrDJszvPY1h3S/Rme/4MBJaOltTM7NM1RMY5qGhCMa8kBsOvIBF2gCSzexLAPUA6TWkn0xfiBDtIaa1jWhkTAFFcdn/9gyLWFRHNhMx5ptOLVlzh94nmtXJ3B9j66uDnwvoLMzDxpMTszT1d3B8lKZ0WtjdHd2cvLUeVat6iWbTjJ2Y5o1qwe4dmOcRx99iJkLR7gjfwWHNB1tKf7p+Ws8c3iR//qeAR6+cwCExtx4ibWdOp1ZnQQN+rozDG5Zw/6Tywy1wQ9PNWlrS/LWuwbwMFiuwyc/fYQPv3sb//mjuxHAN5+5wKc//zqFQkI9m0QcCF5Z+6iipcCX9Fu1mRAqER6EKxdNTUT4ET3uClnZXoRxU2lIiK8cTJpKkEfv/QiNHkRg16gBUUh8fpQXkTgpX2kbWtwr4rpuTDs3DBNXIZ08VT6laepSrS6nMjpB3IpomGacbYtyXvHCJepHjzpHtJXCq8AP4i6YSFKQ2biAIJCBcl0X6JmMZGFFfuIgiOpidWXh1eVJGPmCfcWhUu4lCWPUVcWtHjOgIm3Cj2/sWoxKJ9RU4M9fEbiFpsqZ1FAWruQ1IotZtFvUFMcI9bCOVmRhGEjfv1jBMxMG8RQVpc2jYqzwpvEwCLWbQGUynCgnHi0uWTFNU665NDl1BMrdoGmatEELEWdItDjBL5TTIlDhHzk1Wabg9IV5tm7oYeNwkZ6+NtyGz99+5Rj37x7g/lu6OD9WY3q+wYY+k9HJOps3FqnUHWwtoNSES7MuazIBt6wyufOWdu7as44tqwr8wzOjZBOCOzcVcJyAhKWxXAv46ck5vvf6GFVPR9N11g2vplarc+zYGcbGJmi0HBzHwTR0stk0gwO9bNo0wsZNw6wdWU13Vwe+H7CwuMiFi9c4deYSx0+cY2F+gabnkxU+96zKs6m3wEBHlqVShRstg//0yz8HmofnOJw5e5krJ44w1xAk8NBCn6wVkk2EdK8e4tD5OZqNBp7b4troOHOzMzgtn9HrN+SbMAyolGsIQ+ejfSbdSR3fC1j9/rdSOnWJ0snzCNuSrpRYi5MXojDutBY/wxeK3sRCaIwuNnE9OUmhaQhfMO6FLIXwYEqnFmoYIWSTFpNC8A+X5kigYSQTPJgKsEKPDQmdLXmb89Uqb15eYHTa5ZZVgo5Clobn8WKtihuAHzi4js3CtXWYIiT0DR7fMUG+NcuVAy3SPQnW357iwPESXf23UejM8cbhc6zuzJCMeh10ETt6DKHRdGVnh6d27gnLUlA81fwphCT66jrXFhscGl1EJDIkC6soDmwhk8tz4eBLbBzK8Lv/9WNcvTrB+UtjDPZ3s7BYZXx8Gk0T9PT0cP78JWZm5ujt72LXrm0Mr+nnjb2H6e5sp1jMsO/gae65dT0f2LhMIpUhU0jz4tFp/vGHo6ztTvLWPe0YzRoTV+aZmK5TMFw0NBKZDF3D3Rw4Mk+eOi3N4OlDNf7w45vQDINMOsWnvniGqfkm//bXT5JMCOZLHr/2Bz+iWm8owHbEcjKUsyyWYeXqSU0NEaZEV5uG6N8LdcEQ0ZQSr55CpVFrscEn6jiKJp1Ih404gXKN7sd/nm7IXppAZUh0ocsKbyHdXhFxXEIUNZWdWzEwGapsT1cXiLiGNjoAFQ9K+gvVZdYPYzlCaEJpGPJ7joi+YRgosGKgHF2qlE9E6rsWY+ENpX3rqXTmKaFHNbErHxw9WCN3UiSqoOy1aEEs0hBKm6+0vvrqdA7ing83BiRGe2ctdlWFEQ8/COMCdxEJOOqUjNZKsS7iy5MzGv+E2mkG6oSQ+7noGkE8UvqewhGYRsxBkjgWI6ZdgtRUoq8zypDIiWclzBhzwxRTx1e9J6aiFHuePGi1UFPubsmukT3v8pcvNI2TZ2d4/1vWY+mwe0c/P375Eq8eGuddjwyzeXWar//4GtsGDUo1n0IxT8IM0AlwrASHRl0KRkCzEVBIaVRbTdav6WNNV4JPffk821cn6cmE+LqFMBN8+eVpXHRM02CpVOPY0VPcGJtA0zXWrh5gy+YNbNm8nuGRNaxeNUAqmaDlOFy+PMqpMxe5fGGUK1euUa5UaLkOmzYOs+e2bVy+MoZpWcwuLFNzZIJ/odpk//VFPvqLP8edezaztDiPbVocOnqWsRNHqbiCRtOhyw6oBoI1HSbFri5GlwLGJ+ZIJmyWlyuEmkaohSRVh/TiconlaoMteZMP9SUoleoU1gzRdf9ubnztObymq6ZWlW1SCVwUVUAGpHz1hhfKrQKWJZgqO1SbAUlbxwtCglCjEQoueD4PpQ08TUMPQxJJEzud5tMXZ5kuNdCFwDZM7tRDujIhrhbSicYOW8ezdb53vsT+cw67V9s0ihr7Wi66G6DbHo1yG82F9ZiJkFa1wTtvHaXN9rm2z6V9rWBgm8aNcY8DZ+Fzn/kvnDg/zssHL5BNGrTlkljq4RSEIa4fVTtDEApMQzkVhYal0tYtz2eu6nBkfImxUpOuvh6MdBvF/s3Y6TzzE9dpTp/gM3/537h190be8ba7cVour71+lLa2PMIQFPI5QgLmZudJJhLMzMyQTqW4em2U8+cv88B9t+K4Hnv3neDBrnH62i30dIbxhSr/4/PneeSWNh6+tRdT1+mhxN69E+D79LbZtPV3kOnu4OiJJfTqHKv6bP7wBwv84jvXsnO4iGFZvHR0jk8/fZbP/e+3cdvOPnRN41OffZMfvXKeTNokDCMjjYIKmiu02YhmoYFEtEfOT+WojF4bodrF+CrvESOLDJk5izh5QsEGhULE6DcJ6NLtJA00gedJW3kYoml6/KyLJojIzRrpDuqB9DMahKZAiTLLJmKTEmGAF9WCq8s+8TNRDQSqWjdUn1s3okClXBRFiXdfHXRxIv8mjTuazCKnmVph5Z4KVLgkKumR/PgAw1r5YWiKKR8nzMOVhsKY5RJqcdcCrEwymvKZyXIWiRLXFREXTWknphHv8SQdd4WAG1fdqrExqsaNAoFCrbcMXcT6h6atHD6asrRFYcnooRIEQTyJBKryMR5RgwChg+cGcedJ3GoYJTRj6q8WQ88C1ZUSAdZuroCMe9bV+s8QMDlTAzQeuncA2zQZHmrjzz63n3whzdvv7KHa9Hnh8CKPb0vx5tkKve0Jiu0FFisuR67WeeKWDDMLDrWFMgMFQaDB5jVdCNPgz755lft2tNHTlkQELm+9s597b1vH4UslwlBj147N3HvfrQwNDBASMjk5y/jYFCdPnefIkdNcujzK5OQsc/NLqoRJ0NFZlDdcz2PPbbvwAzh69BTJVIrQD0gLj8myQ9XK89v//Zd5zxP3sVxaRgQ+mUyOI6euMn7qFBUnYKkVsrHLYLoWMtKhk8llqOlFDh69QHtbHt/3aTZbKtEb4IYhS4tVnDDkwwM2w0ZItdxg4J0P4cyXmXv9GHrSjB1xqDWqbsmdbcQ70kK5h9aFwPVksHKm3GSx6mHoOm4YEiC7oMf8gM1Jg6yaRPN5k65Cju/OVnlhfJliMoGeSJDLhLz3rjTNSXDrLvXQw0wY7DZ1epMaBxdbvLC/xpWWj71aU/gQn5nxfkrTvRhGSM72eHLXGO2DIYUBQb4PCu0WC2WP1w6HvOdd9/LkO+4h317kJwcucGV8nqbjKfFWIMIAU02/liFImvIgbbQ85qoO1+frnJ4oMVZq0t6Rx7QTmOkCWrKLdHsfdsLm3P4f8x/efzcP3H8ri8tlXNfnbW+9i2wmyTe//QJdnZ3Uqg1mZ2fUmlhqjhMTU9QqFexUguWlEscvTPCRPRaP7ipS83UcP+S3P3cKfJ9P/fJaqi1YmlxEW5hncilgqNOid10Xrqexf/8N+nMeG1an+JMfLrNjSzf/8Z1rcUKdUjPg1/7iIO94cITf/8270cKQn+yd4Pf/4iV1eKxc8DzXU42mfqxTRM7NCOfhKY0gqqWO4gORyG3oK9sVX1UFRysgwzSk3VbFDmT9gx5HIDR5hZeBRUVPDBSxPIoJCPUxEX8wqt2OGgul5iw5WrLESY/7yaPnHDeh3XVDl/W6qjtEV2xAlTAggvqtdL8H8QWLyEwihEq/K8KvH0o2odJYdEOPdRE9k808JePsKgzo+xBqWAk7zoEQajcVva9MAfKXE8YsIdNaEVziKH4QxJY41BdtmhZe4MdeZtQeMoxK4rUw/mYjkTM6lOWnWaHqapqQog7him02qsy9qTU3erFLUq76vgQIFEBNgGFaK7Y9ocXBRTlhSVuz0ISixepxJW7kxgqUpqKp3IRhyWCObsofODdNLqZhyj20bXD41CwP3j5Eb2eSNWu7qVYc/u4rR9mzo5snHl7PG0cnmZiqsn3Q5PSYT0+HRTJhcvBShfs3ZenrNLg271OpBSRNKLVC7t7WydhcnS+9MsNjW5MYoU86n6G03ODlo7MIXWehVGX0xhQHDpxgamoOzw+Ynp6mVmvR1dHGk+95O1OTs9RqFUzT5Ml3PYamaSwuLjHQ18OO7Rsp5NPUm03GxqbQTIO3b+nACl223XMXv/aJJ1lYKHH69BXa2nKk0ynOXZ7i0vFj2JbJ9cUWI50WFj5u6LN9XSd6x1r2HjgNwqDmBFQbDn4oX9T1epPFcoXejMknVqdo1ZtY7e30vu0+Jr71E9xqDdOWukp0ydFUxoMwgsHJ3a5uSreLITTKLY+psouhdAO1ZcT1Q7pESMYQoAt6ulN0d+c5ueTw2fPTmJ5PfrAPJ4AO2+HXflUjP2hQmRe4JQ3XcagJn92JBP0bM4x3WFx8vk5lQiPZE5Iqtpi/voN6KUEYBrSnHIbGx5g+5VDcICgOWviBS90JeP4N+PBH3k4um2HX9o08/vgdNAOdU1cmOT2+xLWlFrPlBvO1FovNgLmqy9X5GleXmlxbcrm+WGe23GDNyCC//VsfxQ8Clms+TqCTbl9LKpNjbuIalM/z2//1Y5i2KV1BhkGz0eTW3RvJZdJ8+1sv0N5ZVFDSANf14pu9H8if58xClYfW+7z7tiylhsDQ4akvnuHSRJPf+fkRNg13UK65lKcXcKstdMtm48YiLUdn6vI4q7sFjqbz58+V2biph//+wWEqjYBcLsn//KdTTC22ePrv300hbVKq+nzit37I4nJNPjfUezR2RWrcxIzS4iY9oabP+N8hO4VMy8T3op6YlY2MGxuEIkRTXPEU0yjCIJDWXrW5iPNyIgK0rmgRYRBiReypQB5UQRDiOE6s3Rqmie96sYYcxRWk1deP9QoZvFa8QmUTFoaO57sI1VAoL/aR+86/aZ23cphI4KYbr/Ijfpema/iur7pAxErEAg09kUw+FUfyvRWBKC5gVxa02B6rJomoUAqNeDrx1QGjxSK32i1qqvZW1wgDxcbSRWyT1bUVsUuovVwYuQXisIwWu6U0IT3ZkQOL6NuJujqQDwzP8+IOYpAP7UjAF0IosrCI12Wo3uCItWzohmz5U44yudMk/jqEIVZuJMqNESgysBB67MeWNGAFQ1MBNqEKsgxDp9F0OX1xgfc/PoyOxu6d/XznR+c5cmGZt93Tz13b2/nB3mmyRsjqDoPXTpcY7LI5N9lk50iOjoJJb4fBxFSVesvj2nSVC+NVeospDlyscPhKnXfd00GIRdb26CtavHG+SrPu4HgBqZTFQw/ew9at62Ojwx1330at0qRer1Is5Ni4eT3Xr43htFoslipkM2luv207L7+6j/GxaTSh0XQchooJWi0XUezioft3MjaxwP5DZ7ll5wZcN6DZ8nn5hdfIW3B9ySUhfG7tNTgz67FtdZ7O4R28+PoJsrrHlp4k/TmNrR0GnSnByesL1L2Ax3pS3JUTLC/U6H3sbsIgYOaFNxG2ie9LfTKqLdZibLuEcsodrgTDWbaJ44eMLjSxdEPufAWIaCWhArCNVkB/R5KOzjR1T+cvz01RcnyS6QSZ/h6WZ5fp6zB4110WZtKjY6tOcbVFcwkohejNgMqGNOkP5CjmdW4crjF7xMXzDVr1zaSSnbi+Rn+hySZ/lPOvz3HjpM/kKNhZHT+p8f2XXNatHyaXy2AYOn097Tz8+KM49SqV82fpSNjU3ZDADZitOczWHALDptZw8QKN7VtH+NjH3spv/ubP093Tyf/7wnfRjCQuSbIdA1i2zaUjr/LwnlW85a33Uq81sBMSxWGYslL4rtu3MDk9y+t7jzMw0KeMNfJjPN8llUpQa/rcOqjz0TtSVByDtoLN558f59l9M2xZXeDBnRkSaRNqFRYuz+E7Hqv7kngtmJtdpuwLXr8R8o3jAW+9q49PvnctgWmTz9p8Z+88f/f1s3z+L97OPbf3o6Hxqb/ZyzMvXyCXsSWXStfxlQtS0252gIax6CyDcRFplrj7Qmq3UXWFoULOQZy10KIkXiizEFqELIgqY5VLVBI8IqepFgvjYShDyisEizDWXDz13DANQ/aWo0kTgALRRjluXcEWCeX6POJXobptpE6srUwTUShQWzloUMl7EXW6xxd5uf4PPJlXijWdMMQwTWULJt6ooIVyhUUQTQ3KL69LPEdEhIz9xKHq3oh2gpqmJhawLUvZV4WiPkZfLGjhSjdwhBzWlCjtKWeDpkWpdnUqws80xvmBPHEN08Bz/fi0jB7o0S8q6syW9k0vPqSi/ElwU+955Mn2PQ9Nrd98hT8BuVpDi0SziMWlx5/TdV1JylSmAW5eVUWhm5sCQrqh47oOBMSusSAIMQ04f2mOZDLFvXcOkrJ01q/u4E//35tkkgZvu2eIzoTDv/10FhuHXatSnL9R5/Ksw8ZeAwufzqFOcBsUUhoJPWD92jbWb+zh/m1t/GDfHKdv1Hnb7hxNDzrTAT1Fi5NjHo1qmY2b1jG/sMSbbx6kUq2zes0gXT1dnDlzgVqlSl9/N57rcOcdO8hmc8zNzTM7P082naZZa1GtNvBcl0bDoS8VELg+Yb6Hd7zlTvYfOs/84jK337KRcrXOQE8bR8+PceDIKXxhkjB1bhsyqfg21dDksUf28KPXL6A3l3hyU5qOBHTacGSsyvUFh2zK5BMDScxmC81O0PfEA0z+4DWcpQqhkM2aurpJeMpmLlerIi4PQuiYpqDluIyXHIW0CFceApomOVJCo9706e3K0N+fwxYG/+/cLPvm6qSERnGwCz2ZYHFqjq68ydt2CurzMsnbNqjRvdUkv9ZATAjGchpj/TptPRZ9dyRpLDhceLlCY/EquYSDY7Yz0hkyUj7PfLWOnfJZmmgxfcbl4pU6RxcEs+PjvPTSPl54+SBvvnGMqbEbvHnkHB0L0/zChi7u6c/x+HA7D/bkKBaSnJxvcvfudfzHX3yCD3/oMdatHyKdyvDyK4c4eOwCLQeMdC+ZtnYCt8HE2Tf5lV98B93dnTgtD8MQlEtV/vYfvs5Adxed7QVu2b2Fg4dOcunyOB1dbfieJ99zYUi51mJ9j8mvPJzB0yzacjb//sYkf/H0eX77g2vpbk8hhEaPXuPk3iucnw6Y9iwuzPnsvVzlXMXgwrLGxSXBxx7r55efWEe55ZO0dC7MevynP3uTD797G7/3G3cjdMFLb07yO3/2AslkhBrRleEmUAYeX/GhwPdc6b6LqhbU7dtT9On4th35faRVL15da5q8/MmAsOpH9/34cJGd4TLEFwnscl1GHD40TVPmTTwvTndHHR4R8UIm171YjA8CX+moyI2G3LdL95TrqU2LFscJBCJ+VkbrdpnF09VKP4owaKqrBPW8C2O0U9TKiirTWolFqLA1K7EKPZvNPhWnJg3ZmeE6HnYiscLCUpZZIYz49IyIuJoSq6LpI0qpB74sO9ENSzoeWLHHGroUt8L49BYqRxGVpKispLZC61WIlthSFkOBtZWSoMhuFxF/daEmH8OMLcfxl6GyJYG3YocLgjBmzoSBL5EkCAwVJosw7Zp+E4/fMlkJYor4RhPRgSNXWXTLED+j04jY3ZZOW7xxaJx7dg8yOFBkeE07y0t1PvPVE+wYKXLL5h7WDmX54k+mqNY9Ht+R4fCNJoNFk5TmUOjLsTjXYHmxRjOAWs1jeKSbYi7H43et4l+eu8H5ayUe3ZGmVGqweW2OYtLj1VMlenu68IKAliPFxLb2Ijt2bKJebeD6HqdPn6PZaLJcqvDKK/tIpW2q1Sa1WpW3v/UhSqUKc/OLEIasKlhovkNQ6OYd73yAn/x4H9VymT23bcb1WhA43H3HdjoGBjh75jKtZpOt/UnW92fYe6nCQ3dv5upkmZ8eOs9SucGJ68ucmXG5vOxRc332dNq8s8ticaFC9723YyZtJp57Az1pKxFTk9W1hjKGBLIwynN8NF0WMAWexOyPLbdwPLBMXXZsxJeiANOQK6xs2mLtUJ6UYfHseIUvX18ib5oElk66owihYHlhCVv4vG2d4NgPAzwR0D1sYRgmVrvGwC7B2cBkPA2tyTpWwWZhtMVQGLBljcHFk5epzZ+lXJ5henYJEUBHysIPfAoJj0OzAXNWjoduW082naaQNjD9KudOnydsNVgwLM5Xm5xaqnF0fIkw8Hn+6hwf+Pkn+YPf/kW2blmDnbSwLIukbfOt773E1NQSobDJdK0mmUwyO36NTLjAf/jo27FNRVwVUK20+OlrRzl85DRDgwMM9HVzy+7NfPf7L5FMpTB0S9IiAo2eTMB/ut/CsBLkkhrPHZji/37rGg/s6uDX3pJlsRqyNFEmXZmn6UAma9HXl2Ttqjx37GxnpMPkvq1tzJU13vdgP7plYhs6DV/nk39znGw6wdc/924SlmBxyeejn/w25WoTyzKk60w5iYSKGUgtE4LAUzZdOUEYylIvsyFyQxH48uIp1Drdi1ZdYqV3iFDqaYRhrBtExiJYYePphiEZfAoJH7H7QJM8PgVZDXw/XoVJfTlQQ5E08XjKbi0PB3mp96LAtHJVoZxgWkRHF2AIQzXLRhdrCOKHPze1skp9J7gpqxddljVNlU75N2GcwkARHOTPUAiBns6knxJRND0MMS0z7u+QP5AgHn2CwI8PDQkK9DEM+bCMEr7RqR+tmqIRS6gAoK+0D11Z1G4WoKM9auTu8pRob1pWTMpcwa3rcffwyr7Rj8usZP1jGLspohWTpukxdCxCnkSdwuKmaUI3dHTDUGjjMKZlanHPsHwR6er7Rx1sqEN6eAAAzaBJREFUgSIAC6G4/jcxxKJ0q65G5ZhNo1aGjabLiXPzvO8dWxDC5f67h/nhj8/z3JujvPO+AboLNnfv6OL4jSb7z5aZrbgMdptkjZBsW5LS/DKB52CnLabGa3TlfdLtHXR0tHHnhjx/9a0r3Jhtct+GFKWaw1CnwZq+NIeu+ey5fRd33LaZpUqD+cVl+ro7cZ0WwtBYmi1hJ00WFpdJJixAIusr1SqHjp7h0rUxfLUn7s1odBdSHBst8bZH72D/vlP4gcvtu7fiOC3pJQ88HrznFi5enWbfkdOMFA229qcYnWuSymfp6cjw3CuneNeH3sPdjz7EfN3jypVxEpkUv9BvkQ9chJmg/72PMPn9V6BVB1YoCSIGzAl5i1SWxTCaVoXG2EKdVihDeX4ob3lExT9CyBVXwmLLxi6ySYvTCw3+z7l5NDck25dDtw10UwrolcUyWSPg8Y0hY+dcMl06xX6DhYstUgUdF5NM6DEbuIy1NBL5NFP7S+zps/iHP+ljqD9PbanO+UtT7LvR4JITcHaxScWHcqhxoJ7gN94yzEgqoLRUplKuytAvOrZt0dmWpWOgl3T/EFpbB7OWTSqbZHFuie888zovvHKEE2fHmJhapFKr88xzP8X1NZLtg/Su2YTQNK6cPsDWkQ7e8thd+J6Pq8CTmXSKo8cu0tvbzvETl3nkodtoLxbJ57P88LlX6evtptFyaU+G/NrDNrapUcymeP3kLP/722MIDX7hkU42b2xnYbbM7OUZulIugS5Y22ezcW2G4bV5EjoYocZMAy5MB7znvk7qTbATNn/61Uv89NAY3/zc+9g40oYwTH7rUy/z0huXyWZtWTqnqBkoHIevtgCGZan3pogdWJGwGeUfDF3qYVEY2Fe23Ajzzk1UC11hlzRNqFpubUXH1fSfIWSIm0CsqOxRZOWP6nRNw1LPEV0ecjdNJtGzEogPoujQkJiRQLllfdVdZKpsinRrCU0ZlVQQ0fMCFUVQ05bSM6JDUk7fYTydRY+uKI6g67o8+FRRlmVZstLWV4wowzRii6nruGp8kTvFKKCn6zftzaJ9n0bMzY+cSr7vKYCXhqlEJ02hRjRVPuUr8m8EIYv6M8IgJFC0Wy2C3gVh3PBlmlacDTEtQ/r9A9lCKAyZ5jRMM27xkr3GBq7jqva1CEFC/IOJKimjjHuUQI982ZEFOFp7RTx8qd1IId3zPLWPVDtNdbITair5qix8QQSO9PBcD8uyZB1sJsGla3PUGi3e+vhGTD1k97ZuPvvFw1ydqPD2u/vI5zPcf0sP69YWOXKpStr06csKPFUoU634pGyNifEKS2WHpVKT42cXGR8vs2FVlm++ucDoXIsHt6bxfFjVodPXpnHyeoPRsWmuXR+j1XS4dPEKo9dusLS4jGWbtFouLdehVK5RqtQpl8sI3yWr+6xps9jYZXPbYIa1nRl0DX586BoXr4xTLlVJpWzuuH0ztWpNOV3kdDY9V2Lf4dP0pAWDRYtcxub8TJN7bt/It350lN/9H7/KIw/dzRe/+G0mF8rs6kzyoYEUs1NLDDx6F5qmMf3ifkLLXsHtR8wyTxojtFCTNNVQmjOEJpiuupQaHrrS7yLBUBYHSeqCF8KW4Q6StknJ9fnDE1OU6z7ZXJKOHQWapRZBSyfdlmVpao72tM47dxnMXTXo3miRzgsOfbZK+4jGck3DGfN464hGghS1dI5Lh0sMp0L6RoboaXd49A7BnbcWGewMOXTGpWVlmPR1rjQt0rZF4DrMlxukEzq9xSS9eZvBYpKiHWD7Dk6tzvz8EkboYoU+mUwCW0BPPoHWqlCZm+Laucu8+vJ+6tUq1UqVuuOxtLCIZZlMXjrJWx7Zxa27t9JsOAhd4DQcOjpynDw7iud6jE/Mcu/dO3Fcly1bNnDowDEuXR5joDPFL+0R5DIy33Ds8gK/9+Xr7FybZ31vis3rCmSEy9SFGQy3RTKVIldMYGcNEm15Ag1Onlhk/doc//bSAndsLbJpVQE7keDp12f4u6dP8n9/71F+7r3b0DSNp79zgT/57Eu0FZI3rab9uJPIczyFLTLi6cCLrKcoDp4SlKMCJy0O9IVyq+D78uoZyGedqSy4vmJOCUMoeKG64auLL3FlrH9TMZM8HGTh0wpGJ2IH+n4Y22KFChiG6jIsDJlDkRRdEU8t0QbjZsdSBGP0fS8Wuz3XVSs3Mxb6I4SNhNSiPqd8LxhRl5IKZUfmn0jbJgLrRrbmdDrzVKgsY5HV1f+ZVkAl8ijBmZsQwbJ8xEQXQpYS6aZyLqhVkxmtsuR4pQGIFWdCdOpFp78RCc2GTuD6cTZDCugruZEo/CWUh19TzBnfCwh9H8OSSd2YBxNEk5EeEyojl1c0SURCe9w9rJAkQjFn5C3BUIfaCvMrSrJK8rNYaWGMetYjBpch8Lwg/nMlo0fC1zz1M/L9gGTSYP+RCTasKrB5Yxf9gx10F23+6gsHKRbT7FidptrU2LS5l+n5JsfPTrKxU8NIZci3J6ku1tE1+eIbLWnMzC6xti1g69Ze7t8zwqahFP/4zCgXJho8siOP03TozBsM5B2+9+pVlisuQeCytFzBcTyWy1VC38ESPm02DOYMdg5l2NWf4461RW5fm2dN3iKNhF5emy1x+PoyqWSSq9emOHPpBj//ocfo6yngOC15yw80Eskk9YbL8z95ExOfdlunvyvLoUvL3H3XVl54/TTrN27k8sVR/uWL3yKRy/DzXTpdnkOYTDPwxANMfvsneE0HTY9oACve9+iyQiin1EARpyeXGpQaPpZpKGu5Fq8vhbI3On7A5g1ddHVm8L2A/3VimnPzTRKaoGNbHsOQjrvaXAszlaSxXCZvejx5m8H1ow6d601SacHFV+r0bhJUKh7j07BqvcWQ63N3xmPvgTr5TMgh91G++maRgWKNO7c3yOZ8vvG8w8aNIwz29tBezKKbNjM1OD9dZ9+5OQ5fmmX/uUmuzVaoNCTifX1fjm39Gbb2pWgzPFqlEo1mnYnZBaYX6yzVpF5XyNiM9LexYaDIraszDCaXuHjyMLVakx3bN6IbOn29HdL1iMR/tLXl+b9/+1Vu272ZO2/fRr3RoLOrSCabY99rr/K77x4glxJ4jsuV2Ra//c8XKKRNPvXhIQ5dWGbPoECfXWBuZpn2QpJa3cUyIDRtiu05zpxepq8gmK7Cm5ddPvmBESw7weFrDf7Lp/fz/rdt4v/8/uMI4NSFZT7x378nLwMiKn3zYk6dFzuo1KoGTW41wqjCNUDT5HMpximpwHDoK4urL+3/foxHUkwoT+odiKgCgFgD1XWB70YTg6rDUA61FeuvCgG6vsx+BDdpIMokFIntkdtLHjpyRS+fS4p3ZVqKsivXZHKS0OJ+9CAIVyQCddGNut6FctCFMQJQUcR/puo7siJr6tm9Ek4M/CB+v+jpdPopQ92oow8mXBGzXdfFNHSEJhEX8UmoHuyEckm0wptXrVuEscsoCFSVYxQmNHRCf8XVJfd1evwLi9ZlMThR/QLDKMWpCbzAUweRPHGDQLKJoqklJuuqdVXg+ypxGSruv0SnRCX2IdI3HkSwMSHU51o5tXVdtth50fcX2+rC+CAkjEKPqq1MN2RJvSeJl4HvywCTLpTlTosBdxFqQ2jw031jPP7wBtqLaXZv72N+tspf/fMx7r19gOGBDE1XxxQhz74xyq0jNrVKwPrNfTjVKs1anfZCgulyyEwFkl4dHQ9hG6zuzvPI7g6+9OIER682uX9LjmQySWfB5pYhkxPXKoTNFgNtNgMFi239We5Ym2d7X4rNfRk6UwJ8j8VSnYsTSxy7MsuZyTpn5hscvLbAxEKTpGGyuSdDJmlxabpEd0eWu/ZspdVyMC1brsB0g2QqxQ+e3wueQyFhMNSd5MZMnZ7VPczPV3jt4AXOnT3P7GKF4ZzJJ0bylKcXWfW+x2mMzTC37xRWNi0PX08GqUI/WnWq1K8pf+eCkMWGx2LLV2sr5UaJ2ilDmc52wpDhVUW6u3JYocbfnJ7jpbESBcOkbVsRO2PiuT6JrEV1rkGr7uJ6Ht1ZeGijxvjpgLV3JhGB4MYhlzV7LMan5Kph/YYkpeWAYhr2HqmRbkugd27h6pzNkevdzC6FNMtL7D3mkc5m5esm0BCGQdK2Sdk22UwGYZhopk3ZEVyfrXJibJmDV5bYf2GOk6OLzFVcUgmDkf4Cw+1p1nWmGMob5G2o1+vML9aYXqxyY65GR1ueW0baOXmtxKaN6zh77gr33r2dRqNJMpmk3mjS3t7G3Xft4L57duC5LdozSS5dnWPqxkWe2OSQMuQ648Zsnd/54kVStsXvf6CfnBVy8EKZrckyE+MlLNvAaTgUsyYd/UW61/Rx6twyXstny0iGT31zlo+9czWbR7q5PNPik39ziMHeHF/57HtJJU1qDY2P/Pp3uD42RzJhxWv2QF3uZNWqWr2rUFy8wokzE5pqENTi4qYo++DH+35VoR1Vyyq9hMihKqKHvX5TRkyPzJvSjKMMQNFDX5oAA8kTtEz1dcvVk1AfF4Y/u4mJ9I1YawlC5RILYgeZoRoCPc+NJ6GoJEpT35OuHFYREUS7yXocqAuWUBiWSO+RG5SVjIlpGcqspNa8ujyU9Fw+/1RU8BSlF4W6xcXdHeo09z1PCTPxGm1ljNNU94ZYGbWi1YBQ7W8hEEQPe5Usl50gIr79yzh/qPIdRqxPEHu7icm8kR4hNOKvUzcU+0gT+IEXW9p0xcoPohVWVLVr6Oha1GMufiaUE/eiRIn3+HQOlbc6CgEp73mowjvK7RVGnQJRtuCm/WhkrZNcfkO53ORtwbQE5ZrLkaPjvP/t69F1gwfv28Dr+6/x1ecu8cgt3RRyafr78uw7PguBS9EKqDkhvX15KqUmeuCx7JkcHXMZLOiUSnXyaVgotfCaAas7bb722hTHbjS5e30KC5+krXPP5iS92QyFEDrSglK1wcxyk+OX5zh4eZ5joyWuljWcVDt96zZy90P3cXFhgUZex/d11qcTrO5IU2k6ZCyDpGXww1eOY+gG9969m1azqQp6dAr5LC++dpTG0hIJEdCVMcmlTdr7uyl2tvOv//ZjGo06nmnzwdUZNgYOYvUgbbs3c+3fnkWkEnJVgCbXVtHEjIamh/HBYidMFuou8xUXUwg0Q3XA6OC6QXxhqLUc+nvyrF7dTt4y+NLFef79eom0r5EZzpDuTOI13VjoFIZGba6B4wUMdeu8454EUzcCRvYkcEoBF14tMXi7zZUx6Qbs7dZJpAzMhM33X6/Q3eth5nq5sVDE91ucnh7g6IkmNCfo6h1Aw6PZahGGAY1GC88LcH1JVrWsBKah09FWJJNJY9oJAiNBqQVjSy6nJiocvrrE8evLXJqpMF93CBF0phNsGCiysT9PfzHJT4+PkjA0Jpbhox95O6++fpj77t6ubuHyPeE7AT3dHRTzNpVSiS9+43We+c4P2Wxdo5jWMQ2D0dka//MrV/EDwX95spvbt/UyNlHi6MVF0hqUGyG3bM1T6C6Q7ytgZRKcPrXI5GiFe3fm+Nvvz9Ldl+c/PLGeqm/ze58/yaUri3zncx9g3dp2hGnxW3/8Is+9fJZCIRm38UUUXD8qhlMNp0F4M0U3iEVfz/PimgZd5UHCMJAPZ4XKibFFCrEktJueIUIQBuC5jnrd6GpNrgw//srryTBkiZ2hVuQRuSJeYwn5HIoaYEXsalVtsK4MiZqqf0OytbSbskpCTlci0nBF3DESFUVxE8ZpxUmqxf9NpD9HZiWNm1Ln6uMi1E9M5Y0BtaCnEsmnhC4DNtEaRRPypi1trJoCcoVxkDGMJ4IVIQpVRhUJ49JipmCJRDAufwURHzu/ZKAuGv9iW67qA44KxsMwkLx+P4i7iEGKZL6n8hyqxlau4nT1AI8ELZlGNtQ+Mb5BhL66YUhvdxj3DEvUSTS1RBbgiAMWqod/5FgIVGozHok1oXzk4QopM5C1mZouVMpVTk+RyE8gR8gAiQm/cn2J8lKTtz46jK5p3H/XCP/89WOcuV7infcOEGoG/d1J/u0Hl7l7U5proxUGBor09rYjghbXZppMNXTuHLG4MeNydaLOzGyFkmvR3Z7hQ48O8syBaZ45NMt9W4oUMzbCMOlqC7g41+CZExVqehK92MXI1s088PCdvO/9j/HzH3qYD773Id7yyC7uufcW9h05w8XRG6TSFovTVdoytrxNBR7tSZNMNsG3nz9CLpPg7jt3EiAnxJRpcujERU6dPE8+YZK3BT05jYHhQdr6V/HKG+fANCla8EtDGdyWz6oPv5Oxb79EY6ksUTmqLzr6PQehbLUMQgg1DdMQLFZbzNV9xE31n5qGJC6rN5gbBrS1p1m/ppOUpvH960v8/Zl5Uk5IbjhLfiCN33Tl5UetvOyMhVf1qS473LLe5okHkvRuFqRzFtPXPMZOh3RuFJR8ncaMy+y+JkI3yHcaPHOgQXuHwUcfmYdWmaRe4/b1FbrTE1y4UCOdyWGZctqJjCie6yr0RBjXomqahmWZJGwbCEjYBm2FPL2dbSTsBAE6ZVdjquRybrzEyfESr52bBGBTb4GlpoOtayyFaX7ho+/gBz94iV07RujpbMe0DDLpDKlEguujN3j6Gy/wd194Fn/hGj9/V46uvIUhdJY9g//xr+cpV1rcv63Io7va8MslRq8scmTcJ+O57ByAQtGmHlg0Kk1mR0uIVpNdq22+vW+J0abN73x8I5lCkf/z9Dn+/fmL/PNfvJPHHhoBQ/DPT5/gz/7uFYqFBC3Hu0nwDtVtPlypmVVYkciZpSukUuzgjGgXCr0k75hanCyX/R1BbM+N3J+x7VfItVUY3BRojjXQUDkr5dcl11SK6acpgKturNh4pX9XPj9NybiKOFy6EfUWyUlI9rivhLsjS66rwIZRcDG6eAuVKI/K1DxPPqNcNR1JCrB83msaN32tflx8FeGtfPVsEpp+kyNLQy8UC09FH2iaRswRigtDVNFOZKkNtZuY9l6AaVmYhqyBdZ1WPCJFb0xJvxSxDTdQglQkpAshqzp1XeG2FX9LogjcOE8SJSYjkJiIPfvE5VRBVFQfNWYFUQmVFNAjYFmcvlR2vQh14LkuCOl4kB5piU4OA/m16BFLyVe/VPUgieBnxJ83QiGopq9wxeER9U/okVssSsyq1Z7nS0eQ73lkc0n2Hhqns5Dk1tuHaCtkuHVLD3/292/S9D0e3dNPV3ueheU6Pz4wxVt3F3jjyDzFvE1vT54bcw7HRhvcsTpBZ95gcq5K0m/w4J29DK7qpitn88DOTt48W+IbexfY2m8z2JXCFyZbVtncfssIv/Brv8YH3n0f9925ic0bB+nqzGOZBo7rUi5XEFrIsVNXOXbuPF29BWZnKuBDT1tCWmO1gJ5sgmIuzXdePEIQBNx3725aDYd0JsnE9AL7Dp5gIGeTTBj0ZAI616ziRwfGOXDoIoFl8O7eJLuCOsUnH6V25QbXXjtJKpfG1AWeH8qwapTpUTWhQsg3ZsuX3R7S4isQOirMiiICCFquTzJlsXldN+1JixduLPOXZ+axfUj3pymszuA1VWmaKVSAVCp3ybTB5OUl7txi8fgeA6+l4TcgnfUZ2GOQ7zE4e8HBbHiYZY3law7zZ5s8c6FG96DNg7eZrO8vs2dknjs2LpCySvxkn04qkyPwfDrasjKF32ghhE4ymSAMwPddkokkAdLaGQYBzVaLbC5Ly/WUG1Hu4ZMJm3TSpvD/I+u/4yS96jN9+DrnSZU7h5nuyVnSKAtlgQAJMBiMbQyGhXX22muvd727zl5jr9c5+2fjdcbG5IxBSCAkhEBCKMwozmhyDp27q7qqnnDO+8cJT2lffz62Qerp6a6q55xvuO/rHhpidGSEgpDxBHZvaHJuuUu/22cui/nxH/keHnjgUcbGRti7dyvPv3CUBx56gn/+t3/nyw8dJE01331dxOv3VFhf6TE2XOHwuTa/+I+HSKTmV9+9i5MX17lqNqZ77jKPHM3JCHjfXeMEgWbx8gr1akSscypk9EXIv357jbN5hf/6g3vYunmCjz98gd/5u6f4pZ++lf/6E7eDFHzjsbP8xC98mjCyhWBhmXm2q+hnGYGDHgbShh9BFIT2cjAObyFsIWvHOsJNWnz9pnziHgNppW4iIaXwewNhmXnK2xEERZ6Z88J+kSdUOJafKLORHC7JXALGPmF2OZY3aJf+GkwYlp3qKFuEhrHJovfLcIRXwzqlFs5758Z0YUiW5wRSeBS9Az+aNYPwP69w9BBX3Vpll+volJUcB81G4/3eiahACGVvcvOQGWSHqdac3Vtra7+3kl+jcsjL8Yz9S4LAtLda45n6TqLr4IRODWFye6V3UOKouEr75DA/sgqkx5k4s6Cwsjt3W7vFuZSl491LbbXynYxZNmHNfcKYduwuwsEaA2EXs6/42Rng3mC5l6pcyNrLBGtqki7C0n74lb9UArs/sZWA/aCGYWizQ0IeevwUN105ydZNI2zdNsaGsQa/9effZHos4Zq94+zcMsrBY8s8c2ieu69ucfDlNnlekGYZ33ipw3VbYmq1hFqsOTOv6K6t0u/2EVmfTTPDvP22WU5cWONvHjjPtsmE/ZuqrHYUDZbI1+fJa7NkSrG2tkI3NVyfvMiQMqRaq3LkxHm+8fjz1IfqxI2Ys+dWWVvpE9k2XEjJ1rEqG0cafPKrzzA3v8htN+2n3qzT7/X4zH3fZNtwhNSKnRti+tUh/ubjT6O0YKoa8xMjkuFbrqcyNMTRTz7IWhhTk5JACBQaoaXXs0szz0QKE5J0brlHFAQoYVE3wsyepS10+nlBtZ6wf99GxpKIR88u878PXoReQXOqytieYbJuRhgZv0CRaQ+rIzfj0PZSn81KskVJUlFQaYIuMkYmKiSx5vlnU0ZbgnwuoDUUkXUVX7mYM7sp5NXXV+gToXQFghpnz/d47GlJpd4k63e5cs8WkkqF0+cvccP1V3H58gKFKqhV6yidW1WhqcS7vZRAxIRRzPLKij3wNGEgSDOD5Y6jkH6WsbEZMFWPOb3So5JUmO8Kfvjdb+S5Z4/w+S9/i2dfPMILL55gbLjBnXe9ind+z6u5eWqJSvcCWaaYGKrw0HOL/OZHjyC15g9++hqGWwnPnljjxtmA54+s8uyFlHdd3yAOJd2+6fTXleSlczn3H9Z87qU+2/du5H/+yLVs3TTGo88u8DO//zhvu3c3f/6/34zUcPb8Ou/66Y+zum7fR+zYOTDREtj9h5PMugK1yAov+zfjpND6f4SPpkjt2WHgiXbkJQYoJfbzooX2ORwolyVSYtijKC7jajUIYbI3AjsF8RFzdgRuYmnNBaYL7XEortN0HEA3hpPCkHudA9z8OSPJLTJDzXX2CddhaR/9je923Bjd0MeVl7u7PBSXZ+Sc6sqmxcZJbAPV3PMlLCMsJ6jVau93el9pDXJSmhFTFIfWBW4jD9080GuVzS2X20PCOctNWx3bF0/4Ct/N2Zxb0+i17S9mwV5uAW3UFE7Spj1ewN0EppOQZea67SyCILRgs4Gltw040loZMJr9XcwiPfe3s9PtBk7S61QM9u+RwvBiPNFSGvyFM1WaD6H2ajXh3zhzmeRF7pd6ZpluX19bbTgApXCOeSEII0mRa7768FHe9NptjLRq3HDDLL21Pr/1109w1c4Ge7cMcfNVE3z7xQW++fwSd+2r0VntoZTg2QuafTMRRT8jIqfTzUkCzdJih6LX4dz5JZ49vsrG4ZhuDh/6xmWKfsZd+0foq4Du0mXU0jFkc4pK04zGstwcRjovSJKQM+cXuO/BbzM00iBIAmojFdp9xYW1HnOdgkvrGSeXe4wkIftnWtz/6Is8/OSLbJkZ58rdW/n0fY9T6a+i8pzr9ozyhSeXeO7YMnkA/2kqZv+ujVT37+Pcx++jj2RlvWC4GiB0mflgENPCS7OzQnFhLbNgOkEYSqQw3qBCaZP1kWuCOOSKPdNMNKs8dXGNXztwHpEL6sMVhncPUfTKPVrhyM32YRZWYbd4eZ2NmaJyUvDiMwUXDyuyrkSEBdVRzcp8wIYNEfMvGb9QBjx4OWX/lXWuHAs5ejBldl+NztI6q+2C+x7NGBkZpdPp8KobruLAcy8ThhF33nYdhw4f954kpTSbZ6e58oqdvHzkJPuv2sX4+BAXL84zNtpiaXmZ17/mZrZsmuboyfOEgSRQiuXVNptaIVsnGpxe6KCRHDh6mXtfexNDw3UeeugJfuvXfpi9O3dw9MQcoV6jdvGbXDx6jDiuMFIVfO6pRf70C6cIhOB7bmhww/YGnfU+R08uIZfmODIPWV5w/HyPR19a4uC5Ps/NCR450uPpuZBrr97Ej71tF2+8ZQOVSpUzl9b5yd//Drs2j/CxD3w/1SQizeC9P/dpDr50jmYj8UWXK1TDMDAeHjsuFkJaCGE5snZGPxespAcvByktZUm7wYvpYNLcnjXKLplLjJLbcThMu9LYHBA8xRz797v9bmCX0somADq/khuzucvPF6yB6aQQ2o7UII4Ts2dxnYt9LeIkthMZ6XcTQSD9pEhasJtSynQsosxICmwMh08ftL9rGDv0i02jVVjhgBgwYVt7xNDQ8PsHlynOFOdu7xJnrK1vQ3ind2nqk5Y4aV6Ukjhb+DGTZxP5hZSBlZVZJHaE5Zzh1qwnHKQRkydc5KUUzY3HnG3fJXthc8vDMPJvgEknC+0NLryqwi2nXOqhGIj2dbct9vcKwhBlcfVuLFdYB6u1sPvlknC/l9WEh7EVGGi3DAv8fNWp36SFQNr3zMx284IoFCx1Cr79xEne9oY9xGHE3Xft4OSJRX7nb5/itqvH2TJZ59b901xYyfnn+y8yORJxxWzC40d7DMWaWBY0hxLWM0GjHjLcijkz12d5pcON129h376tvPGGSVpxwf/96iWePdHhpp01mtWYleVV2meeIxcxtYlNSF04Wz5xGHDm/AL3f+NJRiabKFVQUFAbT2hMNoiGY5KpGqIRcnKlR2815c4dI5w9e5mPfv5Rjp66yKWlZaJ8nQ2tiLYKue/pZTqF4u7pOj+6qUI+Pcnlbx8kW+sRxiGdfk4lDEiiEGMpEQShIO0bFlKmFKcWezaTxbbpGGd5YTNgcmVS167Yu4HZsSYHzi/xv569xPpaTr0ZMXblKDpXtpAoLA8otHwlI6AQUiJjydrCOokquGEkhASW5/ssnwlYOiJYPSnYc1cFIQXzBzMyVdDX8NClLvt3x1w7XeHCsZxdNw7x7OPLHL/Q5TsvB+QKWvUat9yyj0qlytZtG8nzguefP8LwUBMpBd1+n6uu2sUtt17HA1/5Bve+/jaOHDlFt9vl2qt2sWf3FmY2TiEEHD5yiigMWVi4jE67rHT6HJ9f4cJ8h0tLXSIpeOzgEW686SqeevJFvvrwE/zZ//0Ml84e5fWb16DbI4wSGhX456+e4x8evMj33zTC7GSd62Y0jaLNiVMLXJzLKHKoNKvUWzWiVp2dO4Zo1iOu2TVMpV5hZLjJz3/fdqIoIqrGdDPFT/7ekygl+OIH38nEeI0gjPlvv/FlPveVF5kYb9Lvm9A6GYZesZTnWXl+aMM7D4LQKoisotPKVf1OIDI5RC5IKc/VK/JhChs6J+2C3NF5tduKWqaW0mVgnqvkgzAg76fGjzFwjghrU3AFty45jObcsPRd4wuR5Han6yKygwHnvNs9oIXvmpzVQljDpPvZHKsrTfs2skJ7tJLLT/J5Oc5ZL6RPShRoO6o1QVy5pfwK+zsjMCgTPTCecS5uNys0apXQZmFYz4ZtddzB69yKxi2pfBslbXytAK/Qcvh0NMgo8G1ioXIT4ORUTlaVYJZH0vs4lG3jXCyjsrdvHMX+zTS4EOs9kcIv7M0ITFlTos04tl1GGIR+b+E+oGbeao1ItnV1M0Bh5XdmnopNbTSMHJfD7DTYWugBT4oqgY1e+Rb4D5RrgwNpdjIyCClUQSUJOHZqhcMvX+Z7v3sfIHnDa3bw5DPn+btPvMxt+4eZmW5y05XjzExW+MKTq3zrUJulrmJqNGZiKKQ+1KJZFSwt9ZgcrvDsuYxmmCLzHnk/pdqqsXvDMLfsG+FLT87zxScX2TcTs3W6Tq+r6F54ibS7Rn1iF3EcovOMKEm4MLfEV77xJLVWxT6Yuc9WNuMGQRhLGqMJc52cC/N9Ng3FTNUkLxw8ytpal1o9Yd8ofPKZNiu9nFZF8Ks7Wwhg7eQ5dJYTJBECzWo3Jw5CatWAzOYgoCFMzFji3EqKVoKs0GgpCCW+NY/CAF0UREnA7l1TTDdrPHtumV87eIH2ekZzqMLY1WN+ti1MGIyvOAv7fURgRwxhQKwF5y50WFCCbmYUMNWKJkYjViRXvHWMuUNd5p5rk9QT1nL45nzKLVe32FErWL6csenaJi+8uM6L5zq8fDogDkMmxoYZb4U898IR9l2xixOnL1KtJNx40zWMjgzT7XWZnhphy+ZJDh06wfSGEZ599jDLKx0brAbfePQAFy9dpttOSXurXDNV4erNkwwlIZUwZLxZYWOrxlVbRvnWgWM8/eJprt2/i5W2Zve05Lfeu5tGoImjCAL4k8+e4AtPzPMTd7d407U1Xj7X4Z4razx9vMfxhYAbtwVcvbfJ3u01rtjV5JZrxti/ucL6WsZNe4f50lNrvOWOKSbGa6S5SUH92T89yJEza3zhn97B3t3jiCjm9/78m/z5P3+LoWZE6joCDXEckqaZ9akp++wIP+qR/rESJZHXwgdd6qDvHq1ASCJ8eJ3LPMemo/o4bWnwI8574RRVRWYKyCAMbVBdWAZAWZOxtudMZotxJzPG0seDIPApk2Umuc1mF2UmiQM6eiKwPVydsdmNydx+V1jxkRuh5Xavq625Jffjfxue5jhfllptlKyO0BH6UZc7o0CYEZbLh/ZAMTvfC4PAyMZ0GWCifViK+cHN/4YeJua6B58Rbsm9Ra6Mo90qGXw4lZUQu2CXIJBeDea+l+s43AfCdUJOs/0KRYQQr7Dih1bxJexcvBhAfZcOy5J1IwbmlUIaJ6iwYVCuCxFC+8VZIEtZnZcX2w+kcyEFUg7Sn23XZDPepd1/DLyu5su1GQM6DpkIaNQjnju8wKWzi7zlTfsIpOTN9+7ivq8c4R8/d4xXXzNCVSq2bRnjDa+aJI5jLi70yNKM3RsC1lPBpokqC4tthqshBy+Zi3ppaZ2VxRUmmlAfaTJUS9gzJnns0BIff2KVRqS5YUeVfipIV87SmTtL1JwmaQ5TTQJeOnKOrz32NKMTTbI8M6NBoZD2dSoKs5/odwuaowm9EC7kmnYoqI5VSVoV+u2MkVrES5cK2lrxvo01rq0ErKxnBIGJZDXGLcnKek4SBUgUQSgocnNo54XgzErPKGCCACUEUSAJbMBXEkiyzHCydm4bZ2Z8iOfm1/hfz11ktV1Qb8QM7x0yZlRHK7CjRimE+fxaXb7KC7St6KJaTBRLTiz0ebmreKmrOJwFPL9ecDQvOHJZceBAh4XlDGTEYi/jOwt9brgyZMdwxMpqyJabKnz7qRXm+5rnDhdMTo4zMtRgZsMYDz3yLGMTo2RpSrWaMDMzSa/bYX5+kekNEywvLdJqtVC64OCzJkvlfe99M0uLKxw5ch5VKPr9nJlazrVbp+j1c+LQBFJFAupxSCyh1azw7NGzXJhb4XX7a7zvrmEiYZbRy+t9fvcTx3jy5RV+9rsmefNNI1zsCC7NdZifT1laXOPGzYIQRbedISNIGglxIPnG04uMNQNePt/lxXN93nPPFErENIeq/OIHnucbT17mkx94O3fcvAkhAj78qRf4n79/P8PNxI6BHOMuML40a3IOPdtOe4IFTinlRCt2t+vwJcISsd0zb/ak+OdV2iJPWAe5MxUHUpL1Uy/1N7xFIx9XNh5bSGFx7u48UVa2arqdMLK0X3t5FQNdUq7cKE56A7VSZQqgQ9AHYWizzM2/D61giAGYpOFm6RIuO3B2Ohe5YEBQ4HZKNmckT5XvsF0XUzYZdhfkRu/NZuv9QphZnBiQcA3GwDp9s5d8eXZViFYFWZYNzIQNhEtbh6aRAht5JQPucP+CBIF3aoYDeejCpySWjkrlx2B4VLF2aBHnIrfSXSHxsr7MAsiKvLCsLz+BsVI7E+Go/YfVBA45bXie58RRaC9SgR9SWeNOEEWWc2XDXwbQCsq5Su3i3I+43LhKuAwUbGaAMTRKa8DECxIMNqZej/jmMxdR3R6ved0uoiDgTa/fwefvP8JHvnyM1980SbNmCLTX7B3mmp1NPvngeW7cmbC8rpkebzI52aIaaE4uaUQI18xWubiccvLsCucvLPHCiUXiJOY99+4hjAR//vlTXFzKeNWuBs2hJr21FXqXDlGIiLHZ3Tz8jWc4cOgIcSUx0j9tRnNeby8EuSoIQwO3i6sB9eGIqBnAaExQC1k80+H4QkanUFw/GvPjszU6aUESBBSFuSBUYdA0ncx0kY1YWs28JEdzeqFnsDZCkNqIVyftxXmOQsGefRuZbtZ5/Nwy73/uImtrOY16yOiVIwRuhOjdty4e2YxGXTCzHPAgCSFIRuuQQ7bcozU+jk5arIkqC2GdZ49qjnRDDqYBT60rDqeaznqfO29o8up7JkgmA6o1wfxCzoXVgCOnIEliGvUa/azP/MIye3ZuptdLee65w6ytrrPWWWNpZZWr923l6488zWvuvpHLc8ucPn2R6akxHn30IC8fOYsMAuqNmrnA11YROkdjVF1RFNO12I9KHFIJBKeXOrzz9mG+/5YxMiWoCMXRCx3+98ePcvxch199zx5u3z9Eu5NyeU3xladX2DUmmU1SFhZT+p0ecaiY3TbG1FSTA89fppuGXLm9wZ989gI/9t3b2LxxlDCKeP8/HOIzD53mg3/0Jt7+pr2IKOIrD5/kJ375cySVEG2RNx6Hbp88pR05whZ5YLO+S2+HUzeVUxS7tKbEnKf91IyU7TMfeqsBnmHlF81ID3X13Y+V/br/Xqhy8hBY+b/7mQqLB3G7BmGRIY6m63hdjuKrbfyF22m43ayTbzu/ntu9IKQHPGr7u4vAZsGrAq1MppGbLrmiWvvX0nyYCxvp4bhYqlDIKHQPQOmtcYqvRqP5fos48bRaQZkDHoehD1OJwtjPCrMs93hgp7POi6JsxaTw7l5saqEQJvpW2BGAFNInBpoqPzcSOFGy/A1B3u4Z0D7wSQYD/CnPvbfjBfshcbK8IAiR9oIKQuckpyRN2g4oiiIKrb0pzczPpZ9VuoAdgzEQHgofSAOBzPPMHyy5vazcnsflKJsRXuFntWbGGlpNuHyF0q2URGt/QQNUqxEPfusUrUhwy63baNYS7n31Fj762UM88J3L3H3NKKNDFVZ7BVPjdY6eWeHl40tcuSnmzCJcs2+CoeEm5xb7nDi/yp37hmm1IlbaKYuL61y5pcoVO6eIkoS7r59hz8Y6n35sjkdeWmP7eMjWyToqiGifexH6bT7wmSc5t7BMsxlBIOyHNfCva6YLotB0q/00Jc9MR5v1M6JIIBQsn+mSFZokgF/eNkRcZICgnxv/BsIqUBB0UvPPhmsJMpSs9jLOrKQUStFXcKij2JAYaaREmaAopanEAdu3TzDdqvG1M4v83qFLpF1ojVQY3jvkYwNcrDHWaIodAziwZpEbwq9nHklQuSYQAaoraI6ME4dVatUKzUaDWi2iWqnSrDVojUwgajVW5xd57U0Nbr55jKAiSDvrXH11hYe/0+E7z7SZnhqjvdZmbmGNMI7p9zPyQtNPc/I85+L5RRr1hMnpMb7xzQPs3rmFpaU25y7M8frX3sKmTdNs27aBLZsmuXhpASEF/fYa7U6XpX7KuYU2l1bXWUtzWlFgDzjJO++Z5u6rhllbLxhqxHzj2Tn+14cOMV4N2DIiuWVfnXRtnUZc8NmnOvSXVrl3X0IaBmzb0WDbjlG2XLmZpFHlmeeWuLSoeevto/zl586zbabOG2+dRoRVfu+jx/nwfUf4m99+Az/83uuBgCeeuch7fvZjZEVBFBnzcRhYVpXATw+k7SCMMrRUMfkJRWHyLLL+wOgHBszBGu3CppDloW6NeL7jGSgUjQ/C/BxZPyurcE/RNcIh3K60KCxWXVvBg/maQpfBctrCYKX9LOVZ7kneoS1KDdA1JMszD7EdTFT06ik/qTDPn7KCoTAISG3XrWyiYhRHKK38eTXgsfD5ToEbV0kBhV1HKO0nOm63HTTq9fc7DDoW/uX4LGahrn0Eo1ZO3WSNgGKAfAoDgU+BP/AMljjzEl630PRLAKG9AsHPF2VJqtXCuIndUp2BoBZhs4udi96hRGQgPXvL8aycekxpkxfgLjYQXiZckmPsTavwLlRtScUOn2LkxJbFJc0LHlsaZhBID/BzTDEhhamE0T7bXfplemH/s0AVIIX2WQQGmRJYs1IpHawkEV9+5AQbR2KuvXYjo62Ee+7ezt/927M8+twit185xFAtIi8E+7a3+NTDF9kyHjDaDPnO8/Ns2zRMruGlU2vs3RghkCShotdPWVrtc+rUAqHqEFUrXLF1lL0bKzzw5GX+/UCbJNJcOVOhWh9m4fwprp6JWMsFZ9YzogCq1YRumvkoWgPOtA+LligKSzbQJNWI/rKidzmjFwf86HSFVzUEXaUp7OhB2kPCUQ0W11NCKRhvJqykBWeX+mZGHQZ8p52xoxbQsPjrQJntXT0J2b13A2PVhE8dm+dPXp5DdRSVRsTolcPG4Gl3eSa/xmIFtEaEEuEgcoAIS5SPttVXKCXpeka6pqk3TGZ4lmfkWUaRFShhPDGVSoW022bp4mXecNsIe2ZD+u02SSKJ4phPfXmNkxc1mzZNoZXxAMRxQjfNydPUvP9KU2hFr5fxwovHSeKYo8fOcOLUBb7nLXfywgtHefKZY2RpysqKCQ5bWVvlxk0trtw8zWQjZkOrxmSzzlgtIksLGnXNG25vsXOmjoqqNBLBv33lBH/4qWPctKXKj949ysETa1w3K5Fpl4vdgAefbfNjd1XYtnOYqU0NJjeNEtZqnD7b5uBLXZZXunzva6b48FfOML8OP/OuPSSNIf7go0f5l88f4o9//R5+9sdvBgQvHF7mnT/1cda6PZLY7jqF9UvYMVKSVPw4vSgKRPjKREDcbtESagNp8OxhGPhZvuPfRVHk1ZquQHTyfwc1dLQMN0nQCA90NRMEBpAmeEihxx/psjhmwKvk1F7OI+ZGvA4B78ZPhaVIS9s5qYG9c5m6KL3R2uFKhFd1mVWAg4warp8CKYnj0t3uyMXaJrG6CZTQpdeuDBW0MmG7mwmGR4bfb9oWs/0Po8gbCY1rO/CxsFoY+WyRm+7Dyc0CEcDA4WagYVGJNraHu6s6HZtfa2Vz1e0HwrZzzrDiqna3tHb+EgcuFEJ796SjBytrEtIWIe/eDJco6CSzeuANK5Q5JIo8NxhkCzEzs0TlpcMOTOJDqezNbBZdemB/VPgdSmFDr1wAS6Hd7LaU+LoxnfQXp/Z+E3N5BKgitzJF7dVaURLw2S+/zPaZIa7aP8PEWIXX37qZv/rgQb7+7BxvvnMLoRDU6xU2TcT83RdOcc22KjOtgOeOXKZRETx3ts+VMyGddo84CunmmnqoWEs1i6sZKu3xtYNLnJjr8713bOTanUN84IuneOb4GjvGJJPDNVq1mDt3NNg0GnN0PmOuk9J02eQDhqy8yG0qo+kO3chy+UiH5W7OHcMRPzFbY7mfgh0ZBHZcmheODQTtVDFUi1jpGp9HLCWRFnyjnTJbCZgNJUoaVIlSmpGhhF07pojDkH84Msc/n1wi6kJjssbI3iF0DkIrC7U0QgyzHJUlUFSWwWTYIsFlLqjCFBMq17Qvdqk1m4QyNCMvP3c2S8p+b524fYbVTHPqrOLAU+s892yXTidgYTnl37/ZIVUxw0PDTE6OsXl2A0uLSwwPN+n3+zabB7r9HlJKqpUKvV5KmpqQtX6vQCOZn1+in+asrXYRwHpnnaDoMV6NCKwENhKatJ+za3uNu29u0qjY9FAkH/j8Cf71oXO8/aZh3ndbk8tLbV46l3HH7ohKo87f37fAW/cLrr9pljyqo/KMbqfP+VOrLC+mTE8E3LR3iE8+NMeBkz1+7l07GZkc548+dpK///RhfvsX7uYXf+YOFHD8VIfv//FPcHFumVrFRBK7CAhpF7nGHKz8wRhEoT3MBqp5W2wpZXwYLufH4G60Bw5KmxOE1l75VBSaKLJFpdAlfdZW9Uo7zp4xkSbVxKiShF11WAGO85ppYf8FgiiKfVqq2d8UXqzkoIpmPyJ8l2EW6NojnooB64Mr8gcvFPd65WnmM5UKl2VuP8PBAIUj66d2l2J+JlRpk3DjO7drcWmGTkFb2DMxjGKCaq3+fmFjPBGlCqDwO4+SeGsY+rlHnyv7wrlD3emeFZigJjs68DZ7T7QVAy+GsEh4l7VbbpuFnZ8rVficDqdECG2wu3LjncCMlcIoIi/ygRhIQWCzRYIg8tiD8kbHpyGGUWBdvRqF8u1gGEfWVGMUU1EYehWXrwbc9tvG8po43sAgofVAVePoli7ASuB9Ka7dNZWK9O52VZShV0EYWL22ueSrtYTPP/Ayu2aH2Ld3kqnJFq+9Ywf/9PHneODxs7z2hgkqUciGqQZbZpv805fPMZTAzfuG6Kc5z5zJ2DRi8j2mp5osrvRQGmanm1xaTnnxxCrTlT733rqBzTMTXLtzlBt2D/H5xy/z8ccu06yG7NtYBRmwqRVw165RCq05vpKznhYkUtqLw3wWVKGMoct2XKune3RWcsYakl+eraGzlFw7abORYCpd5rDkSjNRr9LuF1zqZNRjgdCar6302VYL2F2NyYVGKBMlPDXZYsfMGAup4rcPnucrF9ZoEZFMVBna1UAU9v2xEcVebSlDqwQxlZgzpZq5mNdgGtSIPegQkK720VoiAwMeze1OMbSyyfZ6n92tNvdeM8755ZxHXlzm+Lzgied7fPHRNsvrOWPjo1SShJ1bZ9i9axvDIy2Wl5fpdc2uMctSGrW6Qa5nOWhJrhW1ehWBZL2XIoSg3+9TrVTp9lNUlrGwtMK5pXVWO10KpWhUcm6+rskVuyvEScJIo8LhU6v8xgePcP+T53nj/hr37olI+wXnLndYzmNefUWT//vAAjdvkfzAd22hHzdptap0FpbJ11OmxqqMjgQsrhR86CuXOT3f55d+/EqmZ8b5w4+d5a8/9gK/+4uv5Vf+26tRuebCpT7f+6Mf4czFJZrNSglzDWVZ0ftIT1PNZ3luLwAni3V54eZwF4NcKluUOhNckedGhBOYZNE0S5EWiePUodqde34JbkZYrllxQhlzRpbeFDMGNwF6aa/vYyKc9N/LbaUg7ad2YoJdYgeetWfGYuaZN0h6POtqUE3q1bD2dSnynCAKbTFrLllzwVgklDUROne+i9swk5SS/zc4DdJOHSpdEJv0v2+e5wT1Wu39LkVP24dJFWZ+H4ShmZNJ6SvmMDSyWrfQ1to4fIMgLPcmcWRAYHHoO4coinw3YcJ9yoWNsi+Cy/ww6oGMIDARsMou3nMbBVn4aFxt9cz2ApEuSMhGUmpzMWobwyil8D9DEFpFmbIXiT07ojgy/Cv7wjnlTRiFnhemBtAGIMyMXRfeZINHvptOpPwZ7Rs40GqqwcpJKwSB/b0zP5NV9mdQajAdTPt/D5p//+pRdm0bYe++aTZMN7n7tq3800ef5ZFnLvOm2zaSVEK2zA6zf1udBw4uc+TsOldsrvD86TaNSsy2qSoyiFBFxnq/oBXBel/x/KWCqkrpLLVJwpysEIyONHn9teOgFf/84DleOt/j+h11mo2IMK5wy5YGV44nnF/uc2a5RxxaX1CuSJIQrSQE0F3OWT29Ri8o+G9bRthTFXQG8qcdwkYKI6QJpSCRMN/LudDJaYWShbTg62sZ17RiticBndwUPnEo2Tg9wrYNIzw+3+X9z13g6GqfWi6ob2rR3F6n6OWejeVNI16dZ8a25rMn/FLThQq4ECAKjfunUSUk7xbka33qzSZp2i+Ve4FAFRopFMcv9+j3cgKpqdXqVGoRy/2UxlCD0bFJxsen6PV7zGycIAgFcRTzwotHqNdr5NZtnWcpcVxBCkESx7QadRtpUJClGXEYoAoDL60mCb20z407JplsNllsL3PDlXVee9MIE+N1CiVJZM5nHj7Pr37wKCNNyZYNVW6fyaGf0UwEl9MKy13NU8d6zC90uXJa8OypPkeOL3Pk8EWWljLOzBd86dvLPPhsyvOn+1xz1QQ//o591JsNfutDJ/nbT77MH/7q3fzCz92FyhVzSxnf9+Mf4+UTc1QrRiGnLRyRgaJBW9mQIUaUI2CXqOcBgeVg3Owq7OjcTyaCgDAMfSHnVKS5B7pqz8Fzh6hbWBsZbmbk/hZnHoZmfCUsY81YCOw+aSACwilM3c7MXUBRGHmJrLadraeVOxy7xsfIKlxkt/QXlrbucq01uSp8zEShcl/UOEqIc5MH9ncLggBpcfNKl6o2MKmGYRwZP4nd/Rmfn/B8vyAMCJqt1vulhX2V6VUarcWAE1vYkZX04U2Bw3PYmz8YoNS6RZK374uBPHQ/EjISOLdz0MpUGGijhnilwz0k62cG8iiNvn0QLaABofAHv5CUN6dNIAyk8AFPnjwcOCWUNF4CS+bUKMvDUj7/3XG1jG9FvCLz3XhjbDvtJXVGbx1FJm5TWrWZkx8HoWmjtRpQbFj1mNOGK63AKSvsks9d5ggXAVwQx+b3/vR9h9i3ZYQ9u0fZsGGUu2/bxd9/+Ck+/42TvOnWGZJAMTxU5TXXT7FWBDx1ZJ355R7drGDHeMjyWp8tM2OcvbBEEgZMTw7x9DnNWAXy9R5L8yvotCCMI4ZGqtxx7Wbu3D/JgWPzfPCh84w2q+wYC0hzTSOEu3e2GG8kHL7UZ2EtJTSoUAgEWa+gf77PSqH4rqkm3zOWsJYVRtoc2JAdaebMue/UJBc6OZdWU2phwKFOzov9gtubEdORJJWSRBp0x47ZUYabVT5ybo2/PjpPmkNSCIZ2DdOcqaK6Runmokhd/oFfaBa5KUqERdRYz5EpNNwFN6BYsXTTpBXTW+3S66RUqnUD+8wKg6+QkHU7BEHOggxY6RTcs2eKEMHFtYJNm2YZGR1hdHSIahLxmrtu4v6vPcYTTzzP/qv3snv3Dhr1BpcvLwCarJcyM7OBK/btYNOmDaysrLJ12ww7d2zhhReOEoaSSqViDr+0x8EjF7j7ukne97oR9m5uUWjzegH80SeP89Gvz/Pmm4Z5y61TzC/l7J9U5N0+RT/lwHnFkyd6rLZTWq0qT55MOTafc/D4OkUh+dbLa3zntObm6zbxuru28s637eXGazdRabX4Px86zj9+9hh/+uuv5ed/8mbyTLGwlPGun/oEzx4+T6MeGbl/EHp6gzuQoZROSwGBT+aTlsAbENqxuvAoJImy46LC8aPs72nQOiBF4M8yV4zJgStI+PPDoExcro+wPrNBUrcTvwhKxpTDtDMwAHd0bodx9/sHL+vVFrUivB8Daxp0RWMclXHg5e4akzEiMNOdQJoOK5B+heA9MWiiKC73O7Y4F5goc0cML+yC34iT7I4mMBddZC94pbQNlLLME6WVR2kIi2B3MjPHlve36UBbqAYCmwY/AO7QU5Y6KYTwSgMJvoUyjuxS8+z+HpeRUeQKLRwATfjQGFXkr/CIyAEPSZ7lZokWlPI/BuIn3RgsGwimx34Pt7wKotBWAngUvPZXHn5xpgbmh47AKT0B1IXA2NwPGRDacCnhU8jsfN25W/2IxHwftPLzXjOT1z5CWPihn3nfvvCVl9m5aZi9V0wyPd7k3rt28JHPv8CnHjrFHftH2TjRYnk1Y//uYW64agNbNjT4zMOn2bWhgkLQqteYmmxyeTllw1iVEwspUhe87uoW8ysFR063aa+sQNpHxpKUGFXkDDciPvzIeR4/tMSOqQobR2PSQrBzPOHmmSpCS84s9Vlez4kjydrZHqtrKTtbCT+/dZhullonfykKyeweSgLdVHF6NaXT12RInu5k5AJeNxTTiCQdZTI/mtWIXbPDrEUJv/P8Rb50boWqCokiyci+IaqjISrVJkRH2bGhHQ/IgVAg4d1MisGzAnvoCIvF0Hb0VRQFugAZCaKhmO5Kh95qFykhy1J63VXay3No0WNo+xAjG1v0L69x7ZYxlnt9ziz1aLUa1GsVzpw5z8hwk1tuvJrNs9NcuLxAtVohCSOSxGSQ33LL9XTXewwND7He7XD06AkWF1eY2TjFlk0b+c7TL1CvVlBFRjeFKJTcvS/knmsqJPUqUaVOLVR8+8UFfu2DRzh6dp333tnk2h11jp/vs7TcY++44NLlNqkKGaoW3HVFneu3hezf2eB119TZv6nC1TuHecOtE3zqW8v83A/u5t1v2cPkxBBBVKGXBfz6B57lo/cd4W9+90389I/ciCoU8wsZ7/7pT/L0S+cYblUNGsd2dcp6yZQuTFKnVWK5+FgsKbckXJvizZv3rITWSWdDa/Iz1bcyC3VbHDhun+tunFzYp426ptSSLpxNIc9zv+/UXmWlfYw2uhQD5ValZci4oc1HwucOuXNWFRbLIwznSunCYqBsKJotRrPMmCi1nbJ4EZMwPjVH83CwyCxLvRLNXBaRxZNov4dx3KvAnr3+LLPKV7cXtJt1+1TYC6TZbLxfeJAXfomklHGfmoW5gX8Z7oppC/tpRhTZGZ9v4wVJpeLzhp0CIPKY5HLJZAX0nnulGVy0Ymm1hqgqRBnYUhQFcRLZpDFFkiTlgz04Z/RmG+VzNnRhQWE2sbBQhYeKmQV9YSIcXadkczvcKM6Notyozs0sy1Ge5Qe4Zb+X61kSrP3AFVbGa15zq0ILyoQz43kwIxIjVbReF12O4rycVLolvh2poPnMlw6zZWOTq/ZvYHJiiLe+fg+fu/8w//S5I1y7a5QdGysstVOCULB9dphuL+Urj1/klt1Nnnx5hdmpBhunhgglrK1nnF3R7BkVjCSKOFEstzPmzy5x9tIKJ+a73HTVFG+7cys37mjy9NEl/uWhS5y+3GXXVEwSCop+xo2baty6bYjJkQYvHV5lZa5HUo34pe0jNEVh91flXFkpg7SXQjDfybjUTpEaelpzuJ+zpxZy43iV9VxToKlHksnhClPDdR5pK37rxUucbmdUMkU0kjCyp0UYC7JugQyM6EFImx+jQYQCldl5shR2FCAs3dfEEDuEv5sjm52UfWADYQ8l01knQxWQOf10nbzoETegMpbQmKkTxpKsm5MurHP1bIvF1T4XO4pKtUK/3/MP6CPffJLxiVEuXb7E4cMnSLM+L770MtdcfRWbZjbSzfq8+MIxtm2d5brrrmJlpc3KapvX330Dh186wtJaDykF128J+aHXNrlhk0SLgCQQXJ5f4QOfP84/fOUiu6YSfuUHZtDAxOQQt187xsHnLjEiezQaIePDEfv2zTCzZZSZ6RozG1rESYUXT67z+psn+bNPnaXebPCz79rHxfl1KpWIdi/j5//4aR596gL/9lfv4F3fcyUozamza7zjP32UZ146z3CravZcKOIoKmOJrbrIoDO0pX3Yg9oG0Xkzcxj44lZYoq2wox6P+hiwJpiU1BI2GEahVx0Z0zAlDFW680d4d7jzlhl/h9k7GIhh6P1cbhzqOtQwiOz3kGZn676/26xTfpbKkDzlD31tp0Nu7yLAdmCaQBqlHnbaE4bm91PaXAx6IGBKuWymPPMXqKGB2AhymwpbWIKHO2fEwNmKTa3Ns4IgFAZl4gJGHMxLijKNT1h2S2ClYSY1y2rrrfnGjV/CyGLQtVEsRXaGFkgLPUO6LdWAOcjmxfnWCP93uohRqybz5FpnqgHIbIsqKbHEWEKpm9UJN9e2OQFeASXL1k6Al+FKK+WU0mZ3FBYL7/NCQq+sUgMcscFMbmHTy1wIlmNdlSorYfHzVsM+gFbO0szOXYWfzbquyeFPsDkQ9uUw40AbehOEgs9+6RBTIxWuvX6W4XrIO958Fd85eIk//uBBxkar7N9SQyhNN5NcvXuMbz47x+kLa9x59QiPPbeAQLNn+wS1SsIXn7jIDVvN1482Q9K+YLIFQb/LxqjP+FCMrFWYGq5y762b2LmxwUPPzvOJb10Erdm3qUGuBaPNmFBr7vvWPN1A8GObWtzUEKzn2szrLXwzssvslX7GpdU+nUyRBMJmmgu2VUJqYUA7z6gkARP1hO3TDRaDkL84scbHTi8hMqggaWxrMbS1ic7tQlQagqu0UnRcIE+hCSLz/4V8pSTcLXPd1wppnUBu1GY/I86gmmcGBFoZrZK0YmrjFSrDFaLIKFtUbkxneqXHtTPDXFhoc6kvGB1uWZaY6zAFL770snGOJxGrax1kEDE/t8CBZ1/gxIkzvPENd3P8+ClePnwMtGBsfIhma5ivP/ocV2+t8IO3NnjNDkWzHqIJqFVjvnZgiT/41GmePrbCnfta/Nq7NtEtAtp5wB3XDvHg109Qz9qMD4Usr2b0C8nIeI1+r02cJKhMc9+3LnHlpoSHDy7z6cfm+cOfvZpmq4rUOQvr8NO/9zRHzqzw8b/5Xu557XYQihdfXuCdP/UJjp5eZGSkSj/NfHFkqnKjkAxsURraA10OzuaD0O8/EUZA4ZRCDNBqHdLcoZXCuAxtCwaC4fxo2CmbrKXglT4Q7ach7pmWYuD5xiyiBxl3wuNRCm8wxnLYCqsqLVRBEidm/G0NgC5kTri9nAvJ87lCykdquw5GyvAV4hqDgVc+b8kV0HJgfxsNXkjOy6jtqFaXmSWuaPXm7cgoVN2kJRgdGXl/XhTkuSHW5jZQJQpDw/1xiA0HBrNGG9eWuWV6EBnXeRRHfmamcQosP6qzxEizt4jCEJQ7xLVfsDtMgasGhJSE0qSOmTY0tCq8wmRae6lnbu39JnfYoZ29HM5J3bLcJI4NKKnCOCRLU6/7HrT8C0S5U7FzRW0VCc7L4uiY5g0z4EmPvVAOT4LXXttBDQMakzJD3aqtzFw097h4E96lvfnHqTjMdyqR81op4iTi8/e/TBIobr1lM5VY8P1vvYoLl9r8/t8/STvT3HrlBAEFOgi47cpxPv/YBc5c7PDG60dYXFzn5RNLTA2HPHFinamxOmOJJkBBlrHeydgwHnNpKeXMyXlqUZ/KqGFXJQHUEsnyao+vvbDCA08vMj4UM9IK+d//dorj8x1eN1XnfTNNlrOM0HaixvQJvSznwkqf9cwsSCNb7QUWgpjbz+TWiQa7JmvIOOaj59f5yxNrHF1Zp55DpZXQ2NUgGYoo+nbUGZhtfJE5TI7jGhnVj0oVMnKFQDkzRmtUrpGxzZpws+HcUkktOSEvCn+gIDVZTyFc0l2uPJK7yK0arZ1x1WSTTgFnVzOGmnXiJKafpmS5Ig5DoiQyHawM/eh4fX0dIYzH4eKly2zbMsv1N+zn+cOnaNYko+Eit25a463XNxiqaCAkDgSHzq7zx584ygcfPM/dV49w9bYGI42AG3ZWeOH4Ols2VvnW42c4+fJlZlqgZUC1EbNxc4tY9kl0QRjFfP1gh+lGzsXFHr/36fO8/8f2c/2eEYQqePlizk/+7jO0uylf/Jcf5KbrZwDFkwcu8Y6f/Bjn59aoVUMLMRQDYphSAacKRRDbRL4gKvcH3v0vBsyDJSFAWCaakIEFr+Z2IV28ImLWfS/HpgtsZEWRK3KvhDSTAb+rsNnkHm0SBX4xHSeRgTbawzn0uw9dymQtwqjIC+8GD4PIUypcbHcgDe7EoVfc7tSTLyziKRhQgmI1Q04Y4ApqP35y8bnW5yTtOYj12jg8lYPJBtLsldCm2yizUaTPQ3IKs6Baa7wfOyJSCvugBiUk0R1WXm+tbWUfeKCgA525HUSW5USVxEPEZGCraos6zgubN2wDwB2zxXFlzA9eVvnmzyhfBRjDiycF+PGWdHJhax6MLWJAWfOjwwBoUSJVVFHY9hhvNiyKnDiJ/cLUxFgWAyZGbV3O5SLe3zj2dQotfEzKwLTNtupwpibzAQq8n8DhDNwHFF0GYjlJYhBJ7w8RThHkEM+68A8HViJbqUR86WvH6K73ee1tW1Cq4LvfuJd6vcYf/d13eOFEm5v2jTAzVkVpuOvacZ450earzyywe6bKxhFJIwmZWxe8dKbNvgkT4zrahKW1lKBaY2wkYamvOHe5TWdxlaeeu8Qjzy2xYarGO14zw5tvmeDCUsonv7nAN15Y4dJaxqbRGv9ztknW6/sMZqlhrZ9xaT1jqWPCf3QBgRBIaZR+ShgfzXAt4trNIwzXEx683OMPjq3wtYtrREoTSUlza5PapooZCRbmYJECkxEhpX0olFk3hQb7rZRChtJ0KkojpHmNi1xDYDATeV5YyadAW7VYYHETxiwvvZjL5TQIKcjTwpIVrDlLmTl+ttjjus3jrPb6vHBhjUKZjjqJIrS2M38Z2rgAc1iGUYhtiEFCe7XNsVPnOHPuDPvGurx5V4/dw2uM1QUySKgEmlMX1vjg1y7woYcvcnGpz/vuHudH3jzDM8dWGNIrXL2xwiMvrvDSoSVYXeWWq0dI6iGbdk0xu6XF8EiTar3JymqXR59ZZseGKrVqxO985jI/+LoZ/uNb9gGC+55c5qf+8CBbZ1p84YPvYPvmIWQU8cDXT/Le//JpVjo96rXYRj5YOawtrKIotOohbXdRufUe4A82vxuV0j9PRV745z3PC5QqzOFnu3LzPEi/g83y1BrzrBTePtcgfca6U2eVADvhkSVC49NOtZcGFxY2q30MQ57lFscU2IW0srvi3EblGhGA0mb/EVrirusECreXVsobucWAqVu5JFhbfPmzsfyRnTzJd07m0MdHeLtJhvQ4J3OWSjGQSmhHfU6260QnPv52qNV6v6PfegelJVJqSwk0Tu3C+xAEBtvgLPzSgwOl1wybNYHy39cclCX6JAojCttyOad2kRcIjPRNF9aKH4UE0sogrezX/TJxnBjDoHQ5I8ZwhAemuVhGBnYSpj1z/gs3f3TQZTc+G1SUmRli5LsbPbCAktbnYLoRaaSFVrLnIiSd+kpbF6zx0ZT56+4yFkHZuv//jdiEqZ5De7mXHZKwSGhsqmPgkxG10NSrMV/75mnOXWpz72u3I7Xijptm2bdzir/76PPc9/gldm8ZYvuGCmGlyq1XjtKoxnz9hVUOnUtRCsYakm+fWOfmXU16acHEZIuiyFlZ7REVikq9wmOnNbK7TtbtsLmZs32yRrWaMD5W5bU3TLFrU4OvfGeRPBT8+KYWV9UC1tOcyJquzq70uLyWGhm4gDgKSJUZRQohyLUmCQVXzAyzaajBY0spv3dsiU9d7LDW6VEDkvEKI7tbJM0QoUyug7bFnGluHa7Bdgk5A3RWQGqfA+108kFsWnb8w22x/Vp7NH9RaH+gW0PwANNN+KJCICgyRRAHaAnFSs7eyQbDFclQVbC20qbdWWNhaYX1bkqaFXR7Bg4pMLLJ0AYpZbmimxa06iF37Ep4xw0Vbt4kqISaqFanVgm5ONfm7++/yN9/5QKnLq4hhObnv2cj1+1u0mlnPHxgntt2RCSy4IGDKdeMdbnrplEaky3iuiKqxAiZsDjX5vzFHueWMm7c3QSt+ZV/PcXrb5zil/7jtSyvZfz9A5f41b96lrfdu4OP/H9vZmw4IapU+fCnXuAnfuFzpColiQKytPCZOWEYkKWZKbbs6ETYw8zle7hOo6zOeYUCytG7/ejLjnh8AezMv7bokrbzwIpOApvj48b3wifyGfac8aGY5z22SYX+LLEVrLQXlIuKMMy90I6ZBiYM9ucN7EUUhg5DL3yhaEQ2RiGm3Eh8EB9ix1juNXNFvSeHBOW4ymrH7BkhLHtQemWX8AWxuZDCMPK7ETybr7DeFSdvtoW663hqter7nTPaVfHGh1AwgMn3N7myhppyWWTyhBk47IIw8G5dT5+lTAYTA/uUQQ23UxeEYWgP08KaeJT9fla6K2wKmDAyVueK1xapLoPAHOpecWVuYrMwky7guESoW7CYYe4LDy8zM9ayDgnDwEg57Rvgll5uPjVICFZKeaa/uaBBi7JbUQOlQhBYAKXj+RelZNflALhwLf/6yxL0JnxKI+XYRZRJjK1mhcefOsvjT57lNbdvo1UpuGr/Bl5323YeePgE//S5w7RaNa7Z2qC3nrNv5zh3Xj3B5FiNYxd6zC0VXG5n1KshG2qKbhawefMkK5eXSHspY42IoyuCdpAwOxxw5FyP+eU2y3PLoPo0R5o8+MwyTxxeYv9ojfdNVFjtpyBgra84v9wlVQUVu8fBJrAJAamCJArZOdFg60SDQ334q3M9/vXcivmZNFQaCY2tDRpTVVSuUamp7Aw6W1iRB4SRHRXkGhHbfZRdpotAoFIjMTZfJAaKFukZQQ7p4wQLflYvTNcR2D9vvtZ1k/h9ogzwmO/+cp8NzYSJ4SqzwxWu3jrCrvE620ZiNo/VIG0TZH2EzuisrdButzl7cREJTA0Jbt8m+IGbqtywQZAEIVE1IRKKi8t9PvL1y/x/XzjL/FrBz7x1mjAM2TAc8NprhpG55MSRS7x0oc9brpngo4+3mW4UvPdtO0hHN9CanmBoqMrSQo/+ekqawsRIws4NIU8eS/mtj53ljbfP8Avv3sVqGvEb/3aSf/j8UX7xp2/iL37z1dSqAYiYP/37Z/il37mPajUiCgMQAUKa/U5gUxTNrL1EbphOT5X7BOf1cogRJ0l1EwyXje4mFM417UCBVq3osi6kLJ9ZYe0L7nuVOxKbaqqVnwx4IoQtVH0mklKW3gF5mvmgJjf2938f+FwPt0tTxUDhJwR5mtkRe2H3wcqPpZyRVVgCRmjVeIUVHWll1g5uNKsVAwF8wi/ki0JbubNTjNkGTZYyZuc7Q/AKz517Hf3vLARiw8YZ7cip0jJbnGw1y3OTn2APdZc5bi5go9MOnDXfISpc4MmAllo50GBg0MpKa5IksSElptV37mwtrAKL8nu5O7wocqI4sTI35e8Bl1Aog4A0TV/RfmqljHIjEj6hzBt0HJ3TzrTzLPeRltpG0cZxbHY6bvZo/SmagbAY12r6rBTlDKx+CW4gZtYTIEtOv3Oxl1Jn/CHl0AFSlBRO7ELedGYD6hM7XvEmQ6SHXcrALKYXl7ps2zTCX//Ovdx64wykgvOLPf7n+x/gk186xHffsZH//o5tTE/UWG2nxJHGeIhCHjlwlj/9txf4zXfMsnCpTWOowuzGYZYX2+g058RCj396qsu7b2qSdvvMrfRZW9eMNEO2bxniX58oOLWc8z92j3BPQ3ButUdbBaynBVmuiCwKW6uCQprR0UQjZnakSdRqcGC5z+dPz3OwbbqiJCsgDog3VGiM1dBFbg4EYWq3IjcZz2WanLbcNIFS9p8JZS7lVBFEpkrTuefGIzQooXz1VmRFWTQgTLVnd/MikOYysiMGX3wIO8LtK6LIeAqUKojjiPX5Lmun20y1EoarIRPDNSoyZHy4SrMS0qrFBLIglDm5gtVun7Tos3E8ZseGGlWpKLSAXFFvVjm90OdT37zAY4eXCKOYN17X4tpZycxUnQ/cd4l79kQE3Q7zF9c5tBxwbD3kji0xjxzr8+a9AVfubBG1GgQSQl0w1Kqz3m7Tz+DQmT5f/vYSi7nk3a+b5n1v2cETR3P+5wde4ux8lz/8lTv5jz+wD/KM1W7IL/72I/zrZ59mpBl7uvVgVoVDA4VhCV8t7DhZ29xvz35zFFtVEIWxJ9a6SYQLRIqj2O5wlf9n7iIw5FztpbmOfeaUUEopfz44ZZVyNGn7XKmipG4I7QCuOVqbkZjxbZkYBlPYitIT5nxkQvgIbtftKLtTUEVhu1jz/CqrAHQbz0JpBvCq1jsTWBKwu0DMRSQc9aJQRGFgYbrKd1eIsjtSNh7XccJcx+cc/igDnxR+laHsJQ1iw8aNujT64T0dURz7+Eaf+m1la86LEcWJaYUCc+MJ60o3HZ30en4nfyvsSAJhKsPczgOdrM20TNacJUyGg3NfC5swmOWZTRnTPjbX54nYql/lue1glGfpq0INECkLixyxail7ILhRmE9MFOIVJki3GMvzUqNeStuMpyQIBEKYF1fb7ysCm7RIueTSzsDmlEAe92I9JVZWLIS0pGKbm+Lnkn7lUu5zXFZ9aDD7xf+DkYmikE6nj5QBf/Lrd/Put18FKqAQAX/9j4/x/j/5BlOjCb/0nu3ctLNOqg3jLK5USOKIP/ngExx96Rw/+bpR5pcLUgLGR2oMT40wVE/4xQ8eYnNLcfuOKmmuWOtkdNKCdl/w5SMSRMRf3LyFiXSd+bWU+V5OL81tl6BYzzKSJGQkSWhVEnrVhCdW+zw83+PI6jq9bp9KaHZTrQ11GhMxkVVO9XsZcSUmK0pDlZFCh4YjhosXKLySyi8OhUQGhqfmqsksS2237NASBUEYm4fQft88K5BCW2KDstk3Zhlv8PXGC5QOxAEYoqxESzNCTdsZa5c7rC/1TZyAbYBCKUgqCZVQcuuuFm+9eZxmK6FatNFa0k011Yok1wFHL3b50pMLfPvIGlfNVjh3eZW33jnLG26cYHm1Q6YUH3vgHNc12lxYzOimOWMTTV68nNNTAUMxhAJE0adeC2h3ckQYcnGhINcheRgxMdli/+4m3/3q7eycbfEPXzzJb/3TUbZtbvE3v/8GbrxmElTB0VPr/Kdf/BKPPXWKkZGK6cpcnHRRPjfGKBeVOHWl7OFvqRY2lExZI6F2TvAwQnv4YOjhqH6OP9jxO7S5EMgw8peKy+8pbPFqFtklbUJZiKywRl2hpU/zc7J5pwQ1oW+Swu1xrUGvsBQKn00kIE/NmalUPmA4lpbobVYGzmQoXZKrZU+FQeiFSEWRlSZCZz7WCikC73XRAyF6rpB2a4/AjekLx/ky+8DQ4ZOwcmVLMTfZRZJcGWGQQnnDpti4caN2ln1sXnQYhGYpOxCc5w2EUpCluU2bC+2ip/BfgzbrJB/QLi3Sw3pXDHSuIIoM7sS9mcYNKX2VKKWwM2y7sEQSxRFpmpaZyLb7UYXy6ih3o2qlTD6FNeWYGzoitwdDXii7fAqs98NqxQtFGLnLQCIJbFVUqhS01VjnqqBMiLA7EwUyNIwk6fJMLC0zDAxi3BN9i1LKrAYEB4V98+wiySrT1P+jgtNe6YayChQ0QgRepODGaQ4JopQijkN6vT6d9YKffu/1/MZ/u4NKZE6Qx545z8//2lc4ePgiP/LmTfzIGzaSxCE5AnRIEMOf/utTXD41xztvG0GlmovzXQoRsnGizhcPtTm9rHjfjTVkGFJNBGlW0A8iPvBgj1QLfu/qDexrhCytdugVin5maLtJEBr5bpLwYk/x5bMrPLO8zmqamYzyaoVNY4JXXzdKXIt5Yb7HuXbG4mofooBQQ4AmSkK7pNaIcABJgbY6eEcA0DZLxuwlZCR9lSatrjF35lgpbICUGWkFoXFIaxXYMDCLoim0PbACP8Z00nJTnBi8T9o1lW2SGFKzEoL+SkrRV6S9gmw9Y6wRsm2ywpWbYvZNBWwarxLX6tBPaUSa+bWcbx7p8vjLaxw83WbzkOadd06yoSX42wfm+IV370KhKBYv852TcOH0ZYZFDy0SNk5I6o2EqfEm692MTEuGhhJQivNzHcJA8vHnBPfcMM3oWJONU1U2jVfZPDvKqQsd/ugjh/no/af4/jft4o9+424mRxNEEPDIY+f4qV++j/OXlmg1K3ZZbjA/wo153bHnJPFSDoywtIWcBj6W2sBRzZI4y3KfX+79DdZkiGNC2VFUUeR2chJ4/pSw0FanIEWbw7ewAWJODCGFIEtTYyS2OCR3viibUZLlmbWylZ2Eq+j84W5hrFqYZ9TJcV3Aldba/L12UpFluSeVO3l/7goOhe8IvPnP7pD92MLuYrQdcwd2rK0t5EMpG+ec5TbdFQILeC0bhVKZGicxKjcXexxH3gPj37NcITZu2KBzZYxtoZ21OaaTz0IvcoQIPOPKycNK2q4dCQXG5WlIq4EJaI8CQxFVft7k2ytVGER61u8PLG+EVzYopcqMcj/GsjhnB6kLBsm22C5De7yRC6DXFlHv21gbnuUTyKy233F3jNTXoRUYyIHHjymUvUCiODZvplVguJwQp0wrCnfBBp5k6QQLriWV1gCEMwpiVEMiDGwVoHximhsZah8uE1h1STDw8bMdVV6G4GR57tHVoZTML6xz182b+evffgNbt46BVswv9vitP36Ev//oAa7Z2eDXfmg3u2drrLULikAQCsEn7z/KN586y/7piJv3NRFpymo35FJP8rEDbd53Q4wEtm0eobveQ4aSP7u/w3xXc/V4jd+4apJWmhI0m4YsqwOOtvs8vZrx6PlFnp/vkOcFo7UIpeCabQ3eckODfTM1QpXT72kyLVkr4MRiyqGFlDMdzeX1nLW0sJ4iabhbaUEQlegJIQOU0IR2jNDPMiKngLGdSZ6mpjixqNUit6wk2/pLpM2Zjsx7mtsuSpkDIU0Lv+sqCkUUW1SNHSfalR69NCdNzbNSC2A8EuwYidk7HbNltMJoMyYKA6RQRFIjoojj59Z5+MACjxxaZb1X0Egk33frEFdMFIwP1/ncd+ZYSyU/9/bNHDxwnn//ThuRr7OrZbqridGQjRtG6eQFtXqTLMtoDTfRSJ567gI37Rnm77+5zvhYg194zy46KQZCKhQvnu7y2/9yhKOnF/idX3wdP/GeaxHSOPr/7qPP8+t/9CCFyoltzAHSpf7lVi7qMOSOHWWist2+MAwC0jwjEAIhAzuKMc94nhv0jNlJCvIiMyo3LW1uSFBmudhdgQzcJSF9N18GtJnCUOXFK8gRhkdnI7EtSkkVuc8h0cZMYQkUyuYLZQNeDYMQSXt9f84ZtZWycdfluWGXtmZ5DaRZbtNbcz9GiqLQe8hcNrmQAdKqWaVNTS2yvCyg7YgptLtg531xPHHnXSpUYWLINV5l6s5RR5COXOdsR1pSBkSBOUu0BjEzM6OVvUAM6TPzcY7YZCs3cnF46rzI/aKpXBINBCINtEvS2fEFrxgFyYFlcVHkNsxJ+SWZM/mUMrrcviHSVhfmzXIUy7TfN6iSgSW5W3a6bsN7METpNjWVobL/PfBpZ8IryiiR0ZSXR2hfL3crl4WAY9RoW8HaHZHFCBh9ufLhNB7/rLQFoNkFnKNt2pvIL/+tWsu40MtqxkXyug+BW7oF1lBZ5ObDUijtPTcygPV1xeRolT/+tdfypnt2UfQgqFf43H2H+YXfvI/5+Tb/7V1beefrNqEUrKz1GR6Jee7lJb7w8CnmLq0yGiuu2Npk62SN379/nqum4YqpChs3DhPmHaam6nzuqQ4ffazL+HDCVD1me1USxAlz6wXzWc78ep9uARGKTaMJ86s9VtaNum2iHnDFbMS1W+psbQmalQgZh3T7BZGEIFJ01hVLKZxpK86ta852FPM9RTvX5ICw8DuhIQpL1LYWA0YxgVcFuc+rtktHN37VlAeQN39pq6mXZizgixlM95ErKAL8ojcEWrFk00jMTA22DkumKoJWEhIUmiACpBnVVSsx5+fWeflCj0dfWuHZk22mRiq89po6RRHQXlzkzTeM0OkWtJKCf36kzb3X1BB9xZeeanPN9pBrtja4vNgn7ayjdUpSrTE8OUatVafVqvH0wdMsrmTccUWDL72Qcv8L6/zDz19FnITEcUSYVPjoQxf4sw8f4eorpviT/30PN+6fACVYWs74ld9/iA9+6mlGR6rGBV4UphgbcIvjCNXCxSzgD7i8yIniuNSgaoVHz+mSkKEstyqO7CgbkwFTFAWudnM4Djfydhw/bblWwk5ZnDpG212qy/txP4PWmizN7MViirkwjjxby8luTeJpZvJF7DnplvPeDmE9JqWbfGC6Y4PupJSkNu/FkbfdmCyw0Ei/ZrDSdN9bF8pGj2t/YOkBzJEf8VHufp2jPEszQvf3DYTkuUW6ayCkDHxQlovjFUIgZmdntTsEhY0eFW6+l+eEcWzDfzLCMLYtnB3lWIOXk7YJf0ib5ZR2F4uvCoTHmXsIqpQ+d9pFhKrCari19nuFLM0Iwsi/OMplVdtkMMOhsW2bvS2xYTLa//1mV6OKwibnuZ1BRJqZObWyqAARBJY5FdoLS/lOxL25wjrLXXXjLj0pAk/BjOxrplFWNWJm8Vma2Tz30GeYKPumuuWf13rbttp1hoVdLAayTEwT9nV1h582nzLzYtqeUVr5qRu3qcIkBXb7GWlf8Z9/6EZ+8b/cQb2aAJKz51f5lf/zVT7y2YPcc8s0//UdW9k6VWMtVciiIIw0R06v8ezLSxw/06bfSbm0ltFH8MO31jl+bp3X3jRJP9VUq5LPPdXhay/0aPc1SgoajYT11BQUIxVNLBTXXTVJtR6zPL/GzTsr5Boee2GR5092WEthsiG5Yirm2tmYvZvqVAKFjGIWl3sEQUBSMQFJMoyZW+kxnykudgsudBSX+gWrOYanZeOVwyhE5ZkdhwR2HFWqTnRu0uFkIMucBVX4Nj7LCht7a1hnRWbinKXNmkkiSS0QDFclW5sRm4YDNg3FTFYEtUBRrSSsdTK0VlSrpuNqr6es9hXPn1jm+XMZL5/rsbTWY2ykwdtuafC660YZaiZ84NMnuWeXRKUpI8MNLi/l/PuBNtdOCV48r/m+O5sMj9RRUpjxXihZWe4QNZqkmeTipRWOn+8Sh4JX7arwxWc7/PuzPf7iZ65l7+YKSmmOXtb8wYeP8p2XlvnJ91zN+//na2g1Y0Bz8IV5/vOv3Mdzhy8wPFQhTTPf+Zad94BhCxMtWxTKQ0bdNADhFuNGBZpnmccH5cqa72wmiNltOi5ZYYU8pfLRVO85YRj5oksPBJI5FaS0ERShfe6dMAZrazCCKeXd6UFougN3TkrbKaGV2c0IKLLMFONov3cIrSTZZR9poe2FYG0BHqxXRiWbHVtuf+/AB00FViqc57kdtZYdhbJ/3iK8fYCfI3+4hb7jdkVx7MdmQkAxQDP2o3EhPEXcqWoLC4UEENMbNmhpkdlucaLLDHq/JC60UyFJH9IibYav4JU7Am+os9gNM9OXllYZWOmZcx+7uMrcZ4u7A9RX08rO9JH2NpTeJxEEofeoGE6MDWNCe7ewWzrJMEDZttR5QqQM/I3q89mlq5ZsC25d+s78U+QFcRSRu1ljaLs0pbwsUdk5q+PoezyLe0OtskFKYSODJY6KrOwl5N44ZxLy+wy7IFZKEyZmEamsCqtwD6U1PLkPv7A7F5eI6B4450gVEpaW+9x0zTR/9L9ez7VXb4QMdBTxwX97il/+nQfIi4z/8gO7+a4bR6jXIpZX+whpDiYBrK92WVnP+Pm/fol794Xs21BnoV3wqr1N+n2oNwVnFzMOnMq5sAyHLubctbfO5pGIrLfK0nKPuW5IRMqbbt3I0MQIgaXr9lLN6Ytr/NUXz9PTAZ1+wWhdsn865NrNDbaMVkBlRIEgTCJ6mSZKYtL1dSQ25S0M6RMw185YVQHz64rlfkFHQSfVFELT18ZjkSLIlEFGoFz3IPwsHwqSQBIFEqk0gYZaJGhVBHUhGK+HjFcFE82IRBXUIkk9MhHRWho+WBQbkUIYmNyQw2faPHeux/HLGUcv9ijynC0TCbddMcTnnlzh3XeM8tZ799Dp5rx46DyPfec09+4KuLjQZ61bcL6TcOhSRqwVb7uhxg1XTaAaI6i8w+kj541QIYV+HrDYLuhnsGM6pjVU4x8fWuTwouY337OTW6+ZZm65x8cfXeCvP3OcidE6v/8rr+btb91HkWkC4J8/9gK/8WcPsdbuMtSqGpJrUC5x+/2+D2tyKsNwkLtkl+hal+onoY2a0FkjpB0fu4LH7Dps3LAtyAqtTd5PUdiuQngDZ5GXQVQeQWIFLeXO0iYUhpEh2tqxm7I8KKwx2UdXF7kf1zusv7TmOw3oXFnEjSotC7mbngQW527C+wLpqONWvCEA28E6lZXwarXA5y25DBDXGZSSYFOQOmlv4czZFvHu4jKcwTFN+37y5DolA3O0SYhIzxl0Oe5Y062vC6Y2bNCBFIbzjpW8WjutcsHycvDmsUYqK0XUNuZQYWbFXsaIw3FAmmZlCIowePY8y+2s0S6mhOkElAt0sRWAsBV3oRRB6JaYof/FwjCkyDOrOCoNeNrK7YIgtJdd4OWu/oW0qPm8KKz/Qnuqp5Tl4S0GsCHKKqZcxrKD64XuErHOam0JmqYDMOmJbjHocCUODhnYbksNjNrCMDRKnsB8cBBG5RHIwBqXdNlVWIaPEwIIayQUUryinS0JsgaHbV4iO+6yJqv2ekYlCviFn7mNn/mhm5BIiENePjLHL//Wg3z6y89z+7VT/Ow7tnPVpoTOeko/M69vFEU0qhEPP3mBP/jg8/zcG0cZqccsdWDTeMT4eIN6qAmEJmpU+MsvznFurs9/vr3K4mqfS4tdermi18tBCrZvGWGoGRGEMZmU/Mu32mybafD21+3m0mrK0y+c4+FnLrHYgUTk7B6PGatLdk7VmWxJWs2IamTGfVEc01nvkucCKawa0KCeSTNFLiRhGFGgONuGSiKpVSP6PdMdiDAisq54rTVCKWpxSCgVUkHez6hVJMJSTQMpjO8nDMmVopdqolDQqMaEUcDyap+5dsbFlYIXT3c4MZfRywtaFc16Jrlxe8KNOyvMjNeYX9X83/su8FvvmSHLodtTPHZsnbULc2xtmLzrnqxyfjFlvqu5uFxAlrFrUrJvU8iWDQ2GhxrUqwFLCx0KFdBqVeh2ehw80eazz/YZGa3za/9hDxumGnzrUJs/+sjLHDyyxnvfvpff+PlXsXXzOAQJ5y50+I0/+Bqf/PJLNGohYWAMd25k4xSCxruh7GVgCNTSui1NCFZhd3oWRhkac52yz2ZgfRdpltpoCDUQaxz4otGbfW3+R2gVUS42wnGoSlCjkwwPsM4KN39WbnttR9F235LbEVVe+FTUMIq9gAItDbnAjZ/t7+0SS7UuPWHCBqTlWeH9H24/ZwzUhR87BZHhUeGAsIHZx0g7UvUjMOd2t4gV16E5DY4bgesBBqEzV5sC2fyuhVI2LK9UcyqU/T72wnIGUEtQFtPT09rENNj8BetAL2xKoFNWYcPqDY5E+Nl/6BVBBk+utTYLHReJaMdYMjAoYmkzP4yz0+i2i9x8X7+wFiVIzOWuu8oksId3mmYkSWIOdXt5RVHo0QAyCHzbZjI8rEpLQr+fEkdG1pdmJiTGdVCOoe98Fdpilt3P5SieDrrolmSuQ3AVqhyI+JWhyZKQQWAUbHYu7JRUgTdCmfCaLC/R1do65kv3K3bcpTzgD4eKL02vtiMr/OzTLfhL8KP26We+a7Sm9kIVLC/3uOeuHfzuL76GPXunQAuUEPzrxw7wG7/3NS7NrfD9r5vlnXdPs32mwXq3oJebXVotCfnSo2f40BeO8t7bm1y3Y4hTcxmFDIhVShJLNmzbQL+f8nN/d4Tv2iG5aTZmrqNI0xylFYvrmmoS0AwlQQQPnJbcsHuEd7xmEx3RIO33qdKjmwUUUYXvPHuaf/7KOUZbFVbXM4o8pxoFbBhJ2DIRM90KiYWglkhmxysU/R5aGD6VQCOShLyXk+uAb57s8YZ9DWQUEQlNHAkUAd311JgRhR3F5Mq+fiZ0J+v1qNQisEVYGIWsrhes9gVzaxmLPcXZuT5zHc3CyjorvZxqHHPj3hE2DSluv2qMrx9YZLwhufPGac6cnaNejfjst+ZZX+3yvTfW+ey327QzaCSSWKVM1DSbJiWbtk4TRBFREnFpqcf5uZT5lT7feXmNQisma5pWPWZmKKReSzhxoc2JuZS8UuN1N03zA6/dwem5lA/ef5Z/uf8ku7YM89u/cDtvu3e7ddRX+Mx9x/jNP3mYo6fmGB2tG6VhEPiTygXABWFkTbDKChoilA8jshJZp1TTpmh1o5zQfv6L3HD1DLw1txL2YCCDQ9jRt/SjrMJ2Gobjpzz9242RtMJijcrYB2d6zu1/lr7mx595ZodZZpMXVqk5mC/k4xbsGDsMQyPxRdhuxYEd3c/rk6Ks5y7zGJUsz73KalBBZjhqgzh3PPTTTXI8TNEptJADOyAjSHABUm5U61YPDrjr9ilFbqgJFNr748x0R1mlrrYyXq09asTkeBtUr/bObKf80R4bLqW08MSYNM/M3BM7srEfHJ8k7sdZ2hu0cDTH0Cgl9ECUYgmkM8unKHTLHHsr2+V12WG4/2zZVhrvQHeytSiK/HJLYUY+oZ2dhgMZ6O7Do7QmjmLS1Fx6QpWyPFcJhVaHzoDCR8qAPEvtoSzKOaAYAJvZFlF4Z6lxxqpCEcaB/X2EjdQ140OXtZLbrHrX5jrFmduveNoxBiYlbJaxtvsfh54prPRa2VxmYXcphTLdSBxHLK+s06xG/PefvIX/9L7rTdVeqXLuzDK/+2ff4J8+8jRDDcl77t3M2+6Yph5plJTkmWBktM4jT57lg587zJ4xwXff2KKZCBbairNrsNwVVELBt491uLBW8HOvbrK6XlD0uwgtyKWkmynGEsGZNjx6vMdN0wHDlYjZ2RFGRioEOqc5VOfssuL+Z5a4dt8wt143w2q7x/HTSxy92OXMXJeXTi5zeaVPJKGSJERCM90MacSKjeMVNjQkQ/WY4VrEU6d65Fpz764aqRamC8m0VwJlSISyuOsgJO0XdNOMPAjpd3v0Cs1CW3N5reDCSs5qL2etl9NJBbVYM1qF3bN1Flf7nLyc8qvv3s3sxlFi1ePi/Bqf+8YF3nFzTCcLQCpatZg///dlrpyUHD3fZ7yu2NjK6HcL+oVmvCloNKqoQjI0VGVy6yTVSszUWJ3nXrjIQ0/PsXtDzCcO9JgaqtBNNbroMzFaZd/OCd5y2yaq1Qof+spZ/vZzR1jrFvz0e67m5370OiYn6iBDLlzs8n/+4jE+/NmDJJXASLsLXR5WqrBemJzAemrclLyw7DgwxZGDtTpiRZ6bwkkifEyzUV5qn/+du6rcpgpqt+gtFKG9QIQUfpk9mJHuiLNOClvuFKX3ffnxk9K2yOtbAUroIYRoQaGtCECY/egg1t+Z8Fw8tVNe+srfSW0tfDHLMgtNFN6LYn4eq66yuRthHJtlvgeolt4P92eMhLncmTovh7+07CWJKEkKRj6sSyWWjQeXLo9JlIbysoOShKFZ9rsVg9i4cVZjvRd5/v+a1UylID0t1bxxpsMQ9gWTntToql63aM7T1EPM3ELHKFOso93atcMg8G1bFEc+jEVIQZHlHqIWRUFZ7ds3KbByNed0ddgAgyHAy3ZdO6scEkCWmeTaUylNK4etVJzz0rg2pV26GtKvwAZOaWVetwAC5MCCzsRN6oG/w2PDrYROisDPWIWQA+2mJQIU5e7EsbccVbhkZZmHNPC5JBZBb30lg7gW98BnWe5HesouhIWQXs7sWluJJs1ylpd73H3rFt7/3+/khhtnIRcgA771xCl+908f5qtfP8HOrUO8554Zvuu2WWqVkPnFPs16yOLSKp95+AQvHF5g84jgui0Vtm2osLiwxtD4CFpU+dVPnuO7r6qypamBgryfkauCNEuJg5ipiRoff2adExd63DITUquZ8UtYq3KpndPL4O5dIddeMUlRa1KvRQRxjXoj5vDxi7x4vM1Ne8f42NfPcPRyzjU7Wpyf67LcLUjznJfPdagnke8ampWQii4IpaKbS0KpieIAnRopaZZn9AtTFKS5Ii8UldgYC/tKEseS4WbEhuGYpbWc5dUe775jnJFGQL0SMjwS88efusDtuyu87pYZ1tcFvfYSj7ywTn/xMjsnQtI0pVKrcHlVcv/LgmaQctvOhBt3hFxazNAUDI806fRyKkGMLvpU6jGtkRYr7T4vnGizulbw6ivqfPb5jLBW46fftpO1FIIoYGyoSqtV58EnLvAHH3qRF48t88ZXb+LXf/7VXHfVFKiMQgV84gtH+d2//gYnzywx1Er8GVAupE3rqz1hOzIBSPYscKeeEIIoMl4CJ2FXaKIg8ngQR4nIs9x6zNy3l7YLkd4M6JzZCPN3am0KLW3Hzk7x6aJslfbmC9P9FyUk0HciFgefppkVqMiBfYT2XZYMQ4osM8BDm05qombNKsAhUfwy8v8ZKzuah4uWdYWfcaTbnCDrDxN2ZyQCx40SVmilvajGsQOVxoz5lSbA7NUCm13iznYnWHDeNzepKQrrzbPLefd1Lu7bxAs4N7s1mAuBmJ2Z1W60FIaRVwfJMPB8FhFIAhGSF6ln5zv4mFtEG+pkYRfx5heNo8ikbdk3s1AmOMbpq7VFk3tBmvvBB4OVlEJrQZalRNYZ7FRd2kXCUWKNpRT28BfeiS7sPsdchpSKDzuDdbsOPZhJXAyENVnpmggCv2ASNte43FMU9mgXfjZbKItCEJJ+PyWQwvtdHHBSBu4icx8wyg5lIFbT7F6slNplDRTKv27SqYOC4BXmRj2QIe3wD8JGhZY4hFJhpJQx4xl1kZE7RmHAWjslDAQ/9q6r+bkfv5nxySHIDIbkY194iT/7q8d46cQCd1w9yY+9ZRtXbDEZ3ZkqqFclR08v8shTlzhyqk2/36ciFfVEMlYPefxEl0BofuTWFmfn21y/Z5zVtT5rSysQRkzPDNFLBX/wpUWuGFfcvT1mYaVLGlU4fTFlcyOlVYmIKjFJYi7vidlx1kXM48e6vPlVUxw/1+YTjy3wU9+9nenJFiLro1Sfv/vKRRaW17l17xBZP+drL6xy+xUjjDc1DxxY49pdIyQig7BCrRqwsLDKgSOr3LRrmEpFUE1CQhQqT7n/mTXee8cISRIS2djlv/zqIq/fFTDdDMgtgfnoXMaDB5b5k5/cwsVlwdOHlnn2WIexJOfKLTHdvqZOiqwEdHXMC2dzXjqbsbVVcPf+GhsnW5ye6xJECaro06hEVCox690ec4tdUhXSCHO2Tdf48BNrLKYRv/4f9yGTOrVKgqTg5dNrfORrF3jg8XPsv2KKX/ypV/Gm1+0ksIXeS4cX+a0//yb//tWXqFQD6rXE0G49Cbjk15kQMDMijWxcQ8nPs/nbQpBZyKDZlVofgjXSeT6UneVLO04yYV1yAA3itepEcUhmFZ2BFGRFCW/1ZGpp5NSGnBEbf0eaIYLAZBpZkY6wHiuzRzGLZxEILwvGPq/KFqoOjJjbkZ22XelgMrsfcQWBj394RfxFXtidqTmfXACU67B81LIjOZeZEAMjK+V9IYXdcUShGVNpd1EroyQscptg6LJCXFqs88wMwhUdnXzAwuBAjI5abEZYGzZqfKys5VmhPU7Yu0ODEtXhYiQLDyw0b6Cwkldlq+sgDLzkURUWUSxkWVk4E5aF0bks4JLoqy0CweIQlA2SEuWOwewCSoKmq7qN+7KEkBV5bm/YzF8cxUD2ifSYC6NQyLK0BBXaN8S03LnfhYRhaD980u4UyuWeB5DZy9OoHYpXmB7N0jywH06LX7AfqCzLDdk3z32cpJPp6oH3xOVDG826HYmpUjnplHHmw6s9dK3IjRLGO1Oj0L6nhTdhuQ9znuXEkQFsLq322bdjnP/xU7fyjjftIajGACws9PnYJ5/lr/75O5y/uMLrXzXFe9+4lb2bG/TSnKQaIchZXu5y9OQcZ+d7LK5oVtYLpEp5/MUFvve6JlPNCBVIdm9p0u8pVpfXkCpl976tfPOlRf7wc6f40ZuHuGpjQlKLyLo91tp9O2Yyuc5KQ1CL+PY5zbWbq1QbDf7l4QV+9HUbmJkeYn6pQ7u9zgOH+sQi5503VVDJEGfnejz83Dw/dNcQX3wuZedMnTfdsZ1uZmWVvS4feeAwr9oUMTU5TKeXoYOE4ZrkH+47xRWbWnzXtUOstFNWV9b58OMr1GTBFdMBy2s96tWAZiPiEwdz9m+tUAsU335xnWqoiEXBSFWybTpkdDgmCRWjowm10WHIc46e7/GZx1Y4v5hz62bB3deNsdxRXFrssbqegYwYqUlqOmVitMqZZc1nD3TYtXWE//K9O2iODaNzzYmLKR/5ynG++M3zTE/W+ZkfvZl3v2M/zUoAgWB1KeVvP/QMf/XBJ1lc6TA8VEHbBWoQhN7E6ySkUgQUKn8FkkfaItOol7QnxrpkPBOzGngFZxCE3p+Q56kXeoRB6OMOnKnYSeWNOS6wnblGRuaZcXQHb9S1UEbnAVMDnqrAZgOVVgMTziQwKX+BRXtoO9EwqCs33rFFqhUGOemxlCVlO5SBHT9bYoZ2YzI3LcAXjW6s5kjeyvq1XFqqI3kEg92VD7QqfRtORSrd+WJTTUPrhzP7XF0ikorM6HyEtnJjyyV0Z78uc1iMr86IkjwxfGpqSvugeKevtjgB86aVCyCnEnIyWmUvBM9xKnIiuzdwOu9AunjKvMQOCDcTtWMjp4Cy80RdaAuhC1C58WC4FMAszcx4xY2wgsCw/LXzUpShPllucj0otDc2Onmds/O7gD+/uHYSZnsJFPafCyFI84w4jn2oCqLMEHFVkfuQu+/hP0hKURS5ReMr1/lbHHxBIEPyIvNqDEN3tfnxXuobmEM+K/wu0bXLZo9jOVDW8aptrnK5MDQXjFO9aJt2JRzx0z6kLkwGIFcFeZrbLArzenV7GevdnDtv2sT/+MlX8ZrbtgMBVBMunFrhHz/yBP/woSdZ66S89oYpfuDuDVy5YxQRSFaXO+iiR62RkPYVusipVUK+9K0z/OsXj/Oz94wjpWRuJWPPjglGqpr23BxxbYjZTaP809cv88XvXOKX7hmiEYYMD0vSVNHrZVycW6VaCUiSmI4O+MLz64zGmm4e0e91mR1P6GcGtbO8bsywm4cr5EXBcCvg0GLErikjJZ/rhvzcO/eyuNqnyAvW2j0+9egckyyzdzykm2qSRDM52WSxK/jiwT7/+d4hzlzs8+KxJc6lNc5c7nDNuEJIRahyMhVyOQ85eBmmhwJGgoLb98Rs29JiudPn8uUuKs0YaUmiUNFqJeSFIKk1GRprkYTw2a+f5xOPd5ltCF61XbJ5LCGSijip0O1rTi2mHJ4TXGoXvPv1m/ium2fJCs2FxT6f+/YcH/3qOapJwH9817X82A9ew8aZEUgztIIHvn6K3/mrRzl46ALViiRJIqNWcgqcAZp2nmVelOFGT8b/pdG5faa12/tZ5Iu9GMyiOya38bNFXngJv7ZqqyJXA1lE2o+23fhEDITZyUBaxAel58Fix4vMhkI5SKAuR9oMFEoO+R5Iw39KM1OQubPI0bed7DeIAntZSe9fcSZnT+nQRuEZRebc08oVc8KP0g1IUlLowk9UBMJL8h06xOxMBFEcmhGbNR8qm0ViLkz3O5kJinK7OsfZ0uV+V9hxXDEw0tdWjOQVYWrA+1Zy2e2+2IJkZ2ZmNO5QsQezMakEnjxpqg7leUJJEvvDSg8ohdyIxFcIQYjSpiowEZPC65jzPCvTyGwmBKpUEJjs85g8L8iybEAhZnMawqAcP9nFtLv9c6uScKE0fp9jOwUPNLSzv8CO67RS5HbBpAtj6tNoQykWTkEiSjKxdSuHQWQNfoX/nYRTLAjtx2FeGeXDX9w+pAy2F4EoA448z8qOsexIqVSr4Su/IAxReW4w1QN5IeYHkWXeilEweMhlKe1V3thpVGmFd8LnWVa25XbEJ9EsrfYIJbz9jXv5rz/6Kq7YMwkyAik4dnyBf/zwU3zw4wfodFJuvWacd7xmM1dtqtKoKjrdlG6aEYURMoxI4pCPf/UY9339JP/h9hF2bGgwt5TS7/WoVBKyrGBstMnUWI3f/NQJhisBP3lbwtxij82zDXIRs97uML+wRqYCNk42WS3gL7+6yI7RkF0bEoSAWkWzvNIjSSLW+5rVtqIaKTpFyIuXBBWpyIl5+zUVOhlcuNQhyzMudTSJEGxrZVTjiGoIlVijQ/jqCcloPSbWmoXVDCkLAq2YaAjiOCRBMz0KYaNOQcATpzKeONbjmomCO64aYWwkImwO0WhF5J0VLpxfplqtkIQwUo8oij46rNFLBTLP6BYBH/r6Mk+f6tMKcxqJQUsUApIk4LYrx3jfm3YxNlzl0NkeX3rsPB9/6Cw6iHjP267kx37oBvZtHzOXPpqnD1zkz/7+Cb708MsEIdQqIWme+efFdeXGjxWY5wDhiyA3IXAO7dyqjYwr36HT7fjVjW5lGVLn0v/MAV9GGkhp94iFHZXbgq/IB8fI0sc5FHkBjiBr+XJRHFl2HyXJwY6azIJeEQbmoHe7Gpfz4b7O73XtYeqMe+aiszuDMBjwSNiJhB646ChDnqT1dznAqxQD/q2gPMuKgTNAe85WmQFSGgXtmEtpz78Snpzhnm3zTAcDijAz2DBdmAzLxbubbLi9j1PJFXZvEtpJi0IjpqamdTDwxjhddBCZeaDTWQuLYhc2HU7aGaeDA7qFjBdbSSe9Mr9YUq165ImwHBfsReMUT2hNoQuPGQkC2+6GEl1Qpgk6rX0YkaY5Umqv93bJhsLFP1q2jNBlV+GWfu6yc4YZhybAjrykbXmlKNtGZyz0lctAQpmgDCjKs4woSgzSxHYSxoGe2gpBG0ZTbjMoBj4QpfxY45E53mGq/FK9NGkVSCy/X1u4nzMi2kvUAxkl/kMR2JZVW8JwEEgrH7SSYLu3yvOcUJbuWmEhb3FkPoyLSx1GGlV+8O37+Mn33sDWzZMQSJBw9MgCH/70c3z4M89w4swKV2xt8pbbN3LHlaPMThrn9Xo/J6k3qFUivvj1Y/zrvx9i51jEq/c0mKxrlA7oIsmJqIQRzVbM737mNHfujHnTnoS5lR57rphkbTWj2+6ytrRGr5uzYdMEz11K+csH53j17jqv2l4xgVWRZGUtRYSwuJwSh4LpqRbdHD7/7Q4Ly2vsGk3YMSTYMJQzPBSSZ5r2SoqII2KpGWlGJM2ErqjwyKF1vnm4w85heP3+KkEoiaVgva9pr/VY7eTIJGa0FTIzVWViLOapUzn/9s02Vd3nli0h+3eNUG0kpFlOJZbUanXaax36/YJuPyfvtZmeqFGpVHngYMqBkz3e+foNJHHE0dMr9NOCjWMxt+7fwEizwncOr3L/05e5/4mLBHHMe95+JT/6ruvYt3sYEVchyzl5cpkP/Osz/NunD7Ce5tRqoQUHSVvIFPazJ/0O0yXWFdY7VdhCzGRqCLQyMcxRHNlqtvRYBCLwmdzus+1UlE4Q4gyuURiaLsDhixwdNhD2AnFR02YME8nQCGusv8JNG9xnnMFnyMV0O9KFDGz6ZMmTcqoyBzSUlslV2Lhub2a2RbG7gMxBjCdaSHexFLmJ2nUZMtGAejQKKNLy2TTREoXHKMkBsrglgprJgSvAw9CSpEML/cRbJRxHL7Bue3euOFqvR7Docgfqi9SBCHK3p3VueWUX/cHIyPD73YvpGCyDkY8uvAW7dnXsJsd5chOcIAjtm2UXxY5FFYR+Sa7tvNO9Ic5051RT2oXF2Io5V7lZUFnarpP7ujAaaXM13AvsFtbuchD2wxbYNs5UTYX3SRiESOBbcS2ckanwP5dnTNk3yodPKeVzRLQz8YkS6+K+1piOAt+V2HfUdlMu6MX5l4yj1OFZRBB41YVDPHtGviMeC6Ov97G+wrnztU9VdK9PaPdIIG1wmBzIVLEKE6+IM3kw2F2P76zs+xqGIbowlWOzUUUhePSJ83zmy4dZXumybWOdoUaF0eGEV796J+94015mNrR47vAin3jwBA8+M8+l5YJaEjIzViMRmtWVda7cMcr1e4Y5uZjy4MElDpzqsdxT5L0ew1VBpHpUYsnu6YhPPrHEpvGIzeMVXjw0T6MZMb5hkiQOqTdCUiXYs3mc0YbkY08sMlYVjA8nKBHQrIX004KhlnGtt9cVk0MVbt3XYKkneeZ0l9mxmFuvGEVryehoi/GxGpHq0qgKOms5QRzTHG5y874mcQhfP7xOLdDs2jREpw/DozU2TlRIGhF5ppAqpbfep7tesHtbi1ddMcShOcWjh3u0F9cYjXvUkoDL8ynLKymrq33SNGNspMrG8YTT84q/fWCFlZ7il//DNq7eP8P0xAj7tw9x2zUz1CoxDz23yl987hR//6WjLPcFP/TOa/iz33wt7/7+K5gYryOCiHOnl/jAvxzgv//v+3n4iZNUqyGVSoTLspP+MMcnA7rO0+w4S9WSOwBKhLnwhY2b6zs5rqLE/pSZE2XyZikaFP6z6xzU0mX8lKKlUnRj9yHmgtCWXKt9VzKYceEjKgTliC0opbZR4ALv1EAOD74IcztF9zy518bbCHJlxl0W5+ICydzvLIVJLdW24wjt7jWwFb5LHvVjQYdbtf/H+UjwEmG8M92N0wMnpLGCHGeC9CF0Ns/EYaWcAlS7T4A0ZkFvB1DWna+sSnVg/CU2btio3dJbBJZQK6QfoRT25hm8fUttnvSLX5fq5QiRzmDo9gUK7VUGhSpsgAlkeWqqEMuDcnRcI9sTJVjMbpGUfbEGw3+1smFKyswpM0umDGwWu7nVzbzRy1TtDRtYfb/n24f20glCLy30i2UHmJQBhVvW2QW+k9i6d9tF2XqGlVVBOfOhjyQe+AC4cRRWrujMPVqUrCyfoOiIxFhgoq8W8SE65qEUvnNwlZlSLrAK77x3qhrvVXEUZj1gSrQPpEuwNDui3M+7QynopwVra32mJ2u8+3uu5L3ft5/t2ybMZl8K1lb7fPXRk3zk08/y4GMn6XdT7rp2itddN8aNu4fZMNVCWKjnmUtrPHd8hePn17g8t4ZOc4IipxILchHSTgvjZL9njA11zbn5AiVyFFBttVha7lKvVdkzlfD146v8zYOXeO/tw4zVQqr1mOmRhHYnpUCz3iso0pwN4w1Gp0b52y+d4fEXlnjLVQlv2d9gfrnP9m3DhEnM2lKbtZU2mYiQaJpDDRqNGl88sMC/PbLA9+yvcv3WGp0ckjik0Yyox6bIWu/0QCsyBK1WHSkDHjzc4+DxVZIiZcuIZN9MheFmSCUAZMDphZzHXlrnTDvk1v1D/PBbtprDNaqy1s44dr7DV5+6zCMHLnN2ocNVu6f4sXddw/e8YQ/Ts8PmWQkiTh1f5KOfe5F/+fSzXJxbI6lIksi47eMwIi9yT1TA+rqULvzh7kjEhfV9YMmygZQoRCn6kEbW6sjSUtrCRRiTsR8P6XLcjDZj9DAKfdPidgulV8MBVvEHZbnoNQd0nqaIYJCsW7LL3GTD7EsKS8iNrMJL++mJ9uj/kjWHj5ywjDT9SmK5q9DdeBxc1obj1ClcuEMQSis3DjxvyxVnyhVsdtdgaL3W3+YI3JY+EdhdpVOMmZ9dQWHNy3ZUVKaqCu+NcS+yv6QkZNbHI9yS31JGhDQKOu9EERAFBiwpZmdmtU/+cxGxhYGDOZWS+QOm63AHjZfo2Tm+9i++Nao5b4SQLo2hdKD6hZayphnpiZFBFPrDLc+c09Qu4UPpUQOGJRPaeWnhI2KlHXO5nzFwizCXeh5YnpZzploZrlss+XvJVhqDuHrnMsXhSixwzIDzrATYgHhI077BAlgDlIvbdGNC7PJQhqEdX5nxUxInHv1ilBO5vTDsh1GX3kSNJokqyFCy3mlbhZawyjab6liUOSLmQXJ7LuXT9VyHEkWxvVRN/kWhSpSLEE6x4h5uK5EcwNI7b00YSnq9jLVOn6mJOt9zzx7e9/1Xc/WVkxBFZgyXKQ4cmufT//4Cn/nS8xw7vcSWqQavuX6au64aYct0g4mRGsKOJTNVsLretxWy4umX5tkyM8JLJxb5h8++zJv3V3nVtirdTs9A7QgZmximu7bG+HCN7Tum+Pg3zvJvj17g+28aYXK4ytxKj02TCaPDDbJeRr0mSZKEuFqnUkv4wGeP8ZWnLvOOq2O+65ph5pd71Box4yMV6kN10n5Od3WdII7Joiqj4w0+940z/MNXL3H37gr37q9zYS5lpZtTqcXkKQw3Y1KlIU0ZqWumZ8ZoVkI+8dgiX3mxz9bxgOWlHjE5cSyJ4pBmPWH7TIt7XrWBbTMNeirm0mKbRw9c5stPzvPS6RUq1ZB7bt/OO9+6l9tummV4uA5BAH3F0TOrfOhTz/OxL7zAmYsrDA9VB5SQ1kCa5wRBhAjMAtuhfgJntrOFlEstdZ8rp6rCLsuVyo2Sycn8VYEMzAI8cB2C3Ye4syO3z4QTeIShkfaq3O4npdub2qRAd6EUisJA9wkDk+ddZJm/ULw1wE4bvEhGBia10u1o7KWiHFTWdhhO/up9ZQMRCj6aQWn//GgnqtHKdG8uUlu6oL3Ajp6cMEGjdEHiQrYKhbSXaVG4VERlC0ybbGql/4DtbgRhFFlxkdl/uphx4WIgHHXcFuBuqe/iB5Q1H2Nz4118gWkI8ldwCZ2E24kJxKbZWW1ItMZEYpbclvwYBd5z4A4PXz1YpVahCs9k0pTLXwfnMsbBMndXaMiKDJQiDGOvMcZJW0tNk1+wBVbyp6wjNbJeEqcIUcVA/KKNITVpW9KiUCwKwGcRmwsnsDubwlbnSlsBNBqlcyvPs5eAPXxdS1jYVERn6nPmS10owiS2+ciWuTPAozHdV+gDaNyb5eCLRlFmqb42AU1rQ3wNw7CU+dl8hDCIfHaBV2PYGayXCzqZovWThKGhDpgZcWl8iuPIdyzownRZtkU3iz9KBUngVC+q1IzbPUtRaLMCCUzlsrbWp9Wo8MbX7ODd33c1t10zS2U4NvjPIODy5Q7fePwkn/3S83ztsROsrayzc6bJHVdPcvO+MXZtHKJVD4CCdnud2miT8+fmaLSGmBit8uDjx/nrT5+g6PXYNRaydTxivCWphQHDTVOQpCQIGfLPD5/j7GrBe+4YZftEhaX1guXVLhAQJyFDjYQiz5gYqbJ54zBff36V+5+8yKs2SL7rmgZLa31WuppQQiWQEEV0UkG7lzExVWfHTItnD8/xlw8uMlkLeOP+ChtGq6RpRlZIAgEjrYReL6PTT7mwonnkSEY/qvHDb5rlur1jnD7fZqXdJ5YB462ETVN1IiG5tFLw7LEFvvLkBQ68vEI3V+zfN8X3vmknb3jtLrZvHjXRglKQdhRPvniJj3zqee5/+CiXF9rUGglJFPncjEAa9aPJjw/K5y2IyuhWO3byiA07E5eyjIZGUo5IBqpXV1Eb4Kj5vJiO34oyXD6HLcpM92Pd1QOsOAEe1CiEVSRagoobJwf24C+0pshMrLVBANnxknWWm+kEtjvBk7EDi5Z34zhpD1LzzCmPBnFQQyd1L5RhXzEYUzvAohoc/ZqwuOIVXC5PnRUOJSI8adh9jbK7Vsf/c6dkIAOyIjeyfEtJt1bBUhxDaT42Zbwtpl3Mrl9TWK/cwM/tCcHutbcdjJuqCCEQMzMGZSKdUkcZPlCcmAdJ2B2AC0Lx8jGrzTaHWfHKw7Fwox8zR4xsmJHfhSjHaLJZHDYvHKWQzoik8QcVwty2GkPN1Nan4tAlUkqyNCdJYnIbF+kMisqSWKWLolQWFQ+ozFwCbhfgHN1l7q8s8el2f+MvyIGJFfbylD4et5Qmuw2ea7f9bFWVGAThs1KkHxN5k6VVT7nENLw82Hyfaq1KUSjSNLXOfIxuTNi/W+HVYT7JMDedZGCVHmaebTsluxh3SzmXD+9DsHzwkpE0FrZK1DZThYGxYGDBbNKy1tY6KYHQ3HDVRr7vu6/kTa/eyezGIYgCkCGkiuMnF3jom0f54teO8O1nztDt9dk+XeP2/RNcv3uI7eMxYyM1sixFqYJMCUbHapy9uMbDB+Z49uVF1joKioyInLFaQLMRstpR5DLmmt0jyEDwj/edZseo4DVXNpgejtAo0lRRa1TJs9QwsKo18gw+8a3LHFtQzFZyvveGJrPjMZ2eptNTxLWIfqop0pSil5LKkEPnc1ZFQj+IePqlRYYjxcaRiFpiXrv1bsF6EZIqSa1V57pdNb7rzi1MTI7RXuvRqgbEcczqSsqhk8s8/fI8Tx9e5Nj5Llmh2b51hFffspV7Xrub/fsmqQ0lYA+/S+fX+Nrj5/nUv7/AN548QZop6rUYgcswcfN9U6WaatS6vM3gxIRuFdo/7y7rG7vfyNK+r8y9d0HaKAi7+Q1sfINzaJcBQcLO/4WXjrrDEKW9u9tEt+Z+H+j2GK7gc3taCtDSKCUD253kqTlUlcON2NFRKM3zawyIuVEliQDsyNelIbrdn9tPuj2M84AEYWBzk8RACJ5BOxm1qREbae1eGztFDAIvHihseqXbHyONaiu3DvnBy8qp3bAki8EsITPWM8ouFAMxzS5PvAyv8/8uMCIjP563U5sBRKwvDlxha35+c7EXmTnnc1VYFpaLg7SUXZD+RUdj3gxdWFil42aZ8ZBTLgh/iEqPNRh0eDtnqbK/hNIud8Rg1g1EMPekX2zSntKm6wiltMFQJW7YfPCEJ/gOzj09qpzSaDPo13DVTqFdLK2wmR1YqiYIoXz2slNJGdSBrXDsIskF2GCx0AUuwCmwCY6izCYvzL7GM2d86mPhl2zKCQNy5bNVcJ2D24sUijiJ7a4Jev2+zxdxD7RWDr2Sl7wrSi2+dKPQgSRDv5tReHet1qIcg1nvTVEU5ntqM092eGfnrtc+sEd79ZvzvXS6Gb1uxobJBve+ejvf94a93HDlJPWhBsgYpKTf7nH87CJff+wk9z90lCcOnGZlpc2myQZXbWtx/Z4xtm2sM9oIGG8GNBoVpNQQhiwtdWl3MwqDaCWKFEm1QpwEhmaA4PSFVf7pSyd57tAio4lmdiykFiqiyIxPOn3N3JJGJBXuvGGUN94yyRe+dZmvfus8E2HKzumEaiyIQ/MZnFvNuLAGPRFzxb5J3vqaGYabMQcOzfP484ucu9yjvdYjiQPGmiHbN1TZM9tkz44JxkebtLs5633NhYUezx9f5emXLvPSqVXOzXWoJAFX7Jnm9a/exV23bOWaK2eptWIzotKavNPjucOX+eSXXuK+Bw9z9PQCSRxRrUbmbdFl9euk7sKJYvye7f9xPmu8r8HgNZT3MQwuXK2ex8/VxUCOhhTS58cXnoFlCkMjXpE+0E0Kk9QYRJKsn3rXtvE0KY8dsg/EAC1b+OWzUSP1/Zhca4ii0gNhumkDtTMeCaxopYzDLjsGPMW7TFyV5KnbE9kGSpWu8DIJVJaZMcoFkhnRkQkss/vm3CJc3N5Jlw7/shg1P5ObDhlLhTaswkJ7lVaWFYSh9LsYY08wqjgnLlDOxAj+InLdUVEYNmDhorgHuyA75mMgt8QF54nZTbM6z4py6ZznhHaJ6uz0zmU9CAhzy2gvp7MpYrhMDGEUT2Xyliixwna+6dphg8uIyJUqE+FwC29THZWzrZI6KcOQIs1M1G5eeH4LFhDob1yHXtAOU28OdRGUAVjCYt2jIPTOcKXK8VIURvT6fUIZ+s7DwSeRJvHOB9hrh1AQfvHlHjdH0Q0iYwiS1hFqKgmzFDcVYF5yxJx7feADLr35T9uxkcFCCCmtizawFY7w8jtctWLdqQjtL8Isy+zSTFBoZXZXDrniEPs4zpby+yJpd0puqWnGidrGM8tXLNmFZSW5jVSa56y1u4RCcvUV07zprq284TV72L1tnEqzZqrqIEClilPnV3j8wGkeeewEjz99hmOnF1BpzuRwha3TFfZuGeLq7aNsmqhTj3MmRutE1ZBcQ2qlnVpouut98rygWjNY9RPnVnnh2BrnL6+z0u6ji5yNkw2mRqpMDsdsm2kyMlSjX0CjEnD4xDyPPT/P6QtrrK31KbRmeiRh2+ZhhlpV9m4dZrwVkxWKNM+ohiG1mvEiqMzmWRQaGcesdRXn59Y5dr7NodNrvHRyjeMX1lFSMj4ccdN1m7n9plnuunU7WzePUm3VbQVrWExHX77Io0+c5Ytfe5mnnj/P2npGIBX1euK9ElJIX1m7Q9wlzTlwqNkRaH9oK5u14fYTrns1yHCj0Hb0bpMi6sbAoeW9qVdgSZwxNs8LW4za3YO1AThShQX0epBoZs3H0o297KXgLkOTEFk+G8qZFF0ErlN74Xh35WfWF7leWeQy2AuCILKXlmHzuXx1c4C6Yb32SnuHZTddd0SWpt58F0ahleQKb4xGDnRkfjur/VjRxGYoP052l770FG5pQtAs986fv0qX0w67dnDwSZ/Sam0JrrswVF3lOw5PI7eFbpHaqA13uQRy4HITiI0bZ7SwnH4XdeoOnMIeGl4dpc182+0RTDC79heU03H7hXFmW2PnsbAacjOSKvzOwsn4Sviag4mVh68zvjn1j8ENmINBYipe1/XEocFHuwPUG2fs8llZ/EEcmQvMZYmUjBldSnttwJMbXwk7djN+gtirlaRt2V0Gu9Cm6xqEIEohywCZ0Dq/3QNhQXBm5yK8usU9/LlVt0hp91H2wQutXt2NE0pKqPDIBePJsaoMGwzmurIwignsh7/4fxDxpRO41O1Lm1TpMDcuK8XlLrsxhu37vBxUSkdMNahqFwYkrV9ofT2jvd5neCjhhqs28obbt3LnzZvZt3uKsB6bikiEUMDqUpvDx5f41ndO8e2nTvH8S5c4fWGJLFM0E8mWjVW2TtXZtanJ1tlhxocqjDQTkgiqsURoM8JL04xaPTadlgjJMuWzIZxcsd/LTDUZSnrrBdVGhXotot/PbAKhycWpRCZ6NE0ze++Z2Xi/l9NNBZ20YHkl5fxizjNHFjl2vsPcSsbCSp9caabH61yxa4zbX7WFW26YZs+OEcbHWxBHEMWQQbGec+5yj68/fpQvP3SYbz99joXlNmEoqNcr9tAtSu+Ww3xISVZk5QjF9dyBLAPH7FJaW8VikTu5q/QIiyyzuTFe0m9EMe6zJB2oz4o3EOaAyrIM7IXkOl0XY6B9Wp6wRlbhDzdXkPgK2k49VDEQrGZ/X63VgIjDIIZcIJW7PJWCMCx3oa5T9125/Zz6BEOlB7oq03IEAxMWv6B3lgJVKsSEdNMZvALMHe65MoF0hZUVm7xyUCq33ZrtDMUr/SZ+/yIcDl6XHjFrH1CW3K1ttkoghZncBNIwwQpD03Bnhh+jew6WOYvcRS9sNrozGrtmQNhzQMxsnHFrbD/GYcAbYA6wCFUY3o13UXseU+jpug4B7O3ztj328ZL2f8wBZEZPLhLXLXGkCTYuXyT7gc7T3M8a3RjIYY5d1KRbBkc2u9jd2gOB5eZrVeHbQzO7NQ5ad3C6P+suTrccxwLFjEwx90oyV0dozBvm5bd2FuwZYIHVXVvprKvUHPsGG94SWpGAlG7xZZVOFgDnjJvukilUSTNV/jIXtiDQVrVhE9+saUgDSZJYoJrykZ4leyyybXzmNeWWye/HFbp80swYzJKBhcDjz52Xxh0Egxp0nxvtA3XMn1lr98hyxUizwjVXTHP37Vu57boNXLF7muZIzRiohYSsACVZWO5y9tIahw5f4JnnzvHiy3McPTHP/OI6vW5KrRLQrIZMj7fYNFljdqLBaKtCsxZQTWB8pEI1lEhVUG+YXVwQms+LUHYUgCDLNFFicuVzi3PIuik6jFjvZvQLWF7tsbDcY6mdc3mpx8kLbS7MG15Vrw8Zkno9ZMtsiz3bR7jp+i3s2znG7u0TjA1XiVp1m+ksIctZX+1z+PgSjz95lke+fYpnX77M5flVtFLU6onPttAMmF8HApTczq7Ec4sS2+GEEBI7birKRbCvsLUfL/kxgCjl/lEY2UApu+y2u0aPMVelM5wBGrQjXAhRxtk6H0gYmUwQR9l1+wYhBuMQpI8mUNaQbDqUAikjMyJ3xGy3r7UKSW0LRSNyKVDa7EK8TNbne5iC2n1mjdvejKzd6FjYPNhyYV54pl/ghDVAFNksIJ8uaotq5yvzERhi4Pk1ijjnCfHB1PY1y92FPkAmL5QmtKMrV+gK93640bo1D7rnTzhopc048e+7HV+XLvcyl8jtusTmTZt1nmdGsheE5Db/u3RuD1AgbdMrbQdSzkHxN/Mga8Yff9K52C0aQNuqweZgFHlOrVohzXP7z0oTor/chMsztyY2XeYReKSANfu4JZQeSBhz7a9zjbrKwGUVC4EXAHi2jSXwFkWBznMfmCWsXFh6B6r2H2jnvDeUXbsPsDd7XpjxnEGnGLhinmVoZSs3KxF02AizRIc8V/7Bd3sShDbdoFWCuRbQZcw7GbKb9boVkKvCjG8jNB9Oocl62UDM8CAEMvDGTA9qsyBNT+e0SZRYyTZ6QBXnKzgxkGbgImKtVDCzC3srQQ5t56aUptMxI6dGNWb75mGu2TfO7Tdv5rorNrBpY51qowph5OiZUGiKTDC31OPc2WWOnVrgxKlFjh6/zOlzy1y83GZxpUc/Lcj65nOfxCGB1DRqCUkkiKQhFCdRgBSarLCFTgBFYQ87NFlWkBfQKxTdbkG/bz0CCOIkYqgRMTFeY2aqwc5to2zdMsaubSPMbmiwcapFpRqZ9y2OrUpFkndzzl9Y5bkXL/CtZ87z7WfO8NLRS6z3CmQoiCPJ/6+sc2uy7Diuc2ZV7dMDAjMDgjIlkqBCipBe+GA/+A/iLzrCDoblW4TCpCnQvAAECEzvqko/5For68ChB5EEZrr79N5VeVnrW29ergz4Agk72VTLunktPa32GIGDIndhXWj/2EvPNJVWVOpxDxczTW45EllHWJFBMloTAYOZMFZApAFv2ExURgf7acdWzKtSBocfRN4TNtoORRK+FqTdmaKJr4kL6b6nvVwPS21g2J6ciJQUn925m9u9ECft9FrE/8f6YyFNakNr4xgJ194jJa9ZaKXTfOvPMCKBReO40lAtKCPVUDBsmhngrvle3zfODowV+TU3PCwDCaNlCm+V896eFVVzFgwzz0hXCFe9m64IXnr8tmTXUzsc/8Xnvwi3Opi38j6agtT3ztzrG4obxsDSxUnrPpfr4s1EyFjIiujxeNj9+lpjD9zqcGnkf5ZDPbBXOiTCZOuDDsmDs+MgahcCsZD1a37sTqzUg1woBsNZ5lJeO0OYUh12CfPCh5SjrHF1jcDIAZJhkFGVOLwv0DAfj4cq/eYZzsJqjBGgWQ2GOoisDOzJ6OTIHiD+PQLO+qtBem222B3gJSCuhdkfjCt2o0wbijH8XFoAzvLKyPzIq32HtYtquI6gsJ2RpXDak10WkTLCTlIxfC1cOuogQpHhlgbWhhTJ15lKqb2X/eTTj+yf/uEz+4///mf2H371d/arf/5b++XP3tmn79+YvblywQyKsVk3e400OX71jf3p6+/sT3/+1n7/+6/t629u+/LLb+zPX31vX/35O/vmr7f95dvv7b63fffdq61wjEozOOnlMcy72ZuXYS+j2ft3b+ztJw/78acf2Y/ff2Sf/eRj+8lnn9hPP/vEPvvsjb3/ZNiPPv6R2ZtuNqA93dtsLrMP2776+oP99ndf26//2x/sP//6d/affv1v9l/+5Xf2+z/+FSgJtzdvHqJQPyXZUY1HWN7RdS8VAA3YmoZkPdNhT29EseL86e+UIirKcLpjYW+X4olwx5hkaUyrKGwcmJzBO7IpmCtObwLVkWl+2/I50TwoHxL2nFkAVXe08HNxKc8ReOa05yI7EFFB6rWR/O0lcvWjMGY6aUqED3kxEU5cmJMuwmJP6CC8u/g7RNVFSN4YAyTslbvYlV6PwLnEfI6JQx3VgGS3hmmCH6MzOu/1O/XCtziqx001VWDsD9d7A16qgyDMM3itaQMWhzimCHkuNfOf/fzn4ViocdSQH1RTkEv+79OCGcRIGSTSRC5KjiQOdyNJkr11qyI0dDhnvCJuZaCCGWA/79ustcyrwIPIBSBVVlz8Uka8ebGBb3O/zswg2Nu6tZzvwiRkysMpt+nG2IlLopeXF3v98EEZJFxkZSttZr4l4SVWgBMfM0dVAo/3SKnddT0EOfvwIfch950dwGDimsBopZThQpsjrC6SqMudn+O9ViE2K82KYXABW1Y/e+VosHW3+8PUZ0uwneFgYMhN7FJp5e/HAFer+WgShekt2cJmt14CDFbFZqXUSfN+KJaY0DbiqfNST7pBB7Z7zWn3Mvvuw22XNXv39sV+8bO39s//+O/sV7/6W/unzz+1f/z7n9hPPh32408/th+9+xh6fdykLw/METi33mYfPpjZgDz5zsOt+0GSfrUxHtno7LDxeEEHMZAFSkqgmU0cjHOa+bD7m1f701ff2Jd/+Nb+9V//bL/+l/9r//V//NH+52++sv/+v/7Nvvn2tntNG93to49ebGA0wd3AvJcOXML+Gh60LEa6ABEDLKW9t43HqJHnrKJhzjt3au4ymAqeB7JDZuoQ7b1VdOWkokZWIfAegtn82BuQC4Wij7wmClWkBLJivXFUPRE5MNAFs5rmRXldDxEYFkx1fL9Jm037gNl4wGyHYrbQSaaxDW0WHCmnuXYIgdJRpe+5bDzgxKcCk34wETjqsKU9QGKBwzPlOHfZmUnx5SWzTwELCjNcYv0gC4e8XQmAnXNas6ad0pzTXq4LQXJN+50nB78oADnmGtfQyIvY/Cd/Hgps//zzz6NkeHkrbfDgObowkizhpH48ksab1XVl6ua4C7iLVvM/zsvIt6G3QbTNvaUAkOYbSq+8XZMfc7++SjNOnLrYjVB/2AFMI++/AwUwxjBDGMzrh1fN3ROE6LYdHyokauVxCAHIxpX56Uw6m2zFoaOee+qzZHa5JK4Y5zwegCy+3qoUG2B1JHLGTkCiMNlRnTyrSinV8P0xethbHrIXEPh7B3IQqAQr8xBnnYYlYgX2uLIDWFTE8TPwoptz55m5z9CbKEw/woNKgenH0vWQg0Zexm49MxasChKjCuwwja2VoxVGfc657MPrtLUCiXJuH3/0sL/57GP7u59+bL/8xXv7/G8+sZ/99Ef2y59nJO77tz+yd+8+QkyrZ/LgNax55OLa/cD2wDGXSo08POa2aYnveL23/fW77+3bv3xnX/3l1f741av95rd/tt98+Vf73//nK/vd77613375jX35h2/s2+9mHvC282tfl1FJ3dtla09UzTR/bhVuc1HJBDl6b3ZTXeVNSpuTKRVI/VSOBzqJ2mZB/cMANi6YjyXx6Wa2I58isLTdqzDqZF8RPdIV41BrFOZcK/GTSkSOrRqW6uGIYOi6jOLIqWi95WiNeCAxu+owNi7Z0XVciD4wjmZwMMfRqbEQpvucYEaOealA7e7p3XAQuCmFhoqrY6FdadPcKeVhL0giJ/XwmEVkftLVe110oD+suezx5oJXh6osmHut5PO8LMv4uY/i/giT4kQH4yyOsTl+nnd6ZxhxmwUf6MQ///kv4lRB+cHJr9FUl0x3jIHErmUTjlYiiu/XKSmtnJ0AgC1JP0ML935dUjnJKQ1lVYLZmpzQJjDhqKhXPNhsFR0Vv2F5Nl9frblZH5et2Fr8pXojK4r8vm9lk3BJ1KzZtpWVeq8LkBiTCoziXoL02mkdM9LJyo4wF9xvmqsSOwDl2hm803H4M8LXj9ECFR38s6xiNA6A54L59bGPDIRTKUc9NxanGisdAsMiF4SyYmqRWngIhmB1tO9o8I8LiW7iJcQML8E4woKaIwvhnhW0RZgn58UbGBa+DPjcr/FiYam+ecUCf69cwr6+LjPLr38Ntzcvwz56NHv78Yu9fftibx7d3r19Y29emr1czd68uexHH73YXNNeHonv995tzW3f39PmDPvw/W3f39u+/vo7+3CHffWX7+2bb5d989fv7Z6p4LKWlXwfbtejW8eujYfTwuW8qPkHeYAeCBYHzQtYSDUeVYdOfaqVzH1jv9hoPJUzfEF0ssusii7EERDFooVBamUoc+yp1hMqZ28uZuNpCsFKds3MwYnDdR62ldHNsTTjXjOe9Rh77jImrmPRz3ht7kHNnwu1XEBX5o66ex3kTQ97qY3Mti11Z/O+M8kUQoQlAQvc8xfPjVyg8/M8vRkLWUeF5ChmUqBoMoRH5Web58fOUHm54hPVMlPIMZE1T1AlztA0Rg/QRLYQ8CcAc2HcmGcDqSPD9j3TA+d1DuUlZNr1qqDD//nf//0vg+Azoka2OhK0xI+U4jVr9vjoja1562CfGL0UZrkehrmSV6W2Owxt0oYbtpuj6g/6COh6tcoh10Jop0Ij28CGkVhIwhqRO4S10hBG3MkS+pjxrUsvG+fDp9xW8bhoQSkQMK/2jVA59+rftqVssfYgnIUuZK5X5TV6uudZ2ZU08DwEEOyFQ4RcfpG9dsZWngY/UxJby3nrXTnL4WB3gVlEzK9buWMHnf8Y4aVnZGpflBVlU+dVe6XKcFiHP0gyXq8ALIsyNwVstr7taCeP6NJD/y5TIvYyRDhQ285Rzj4iOzl0S8n1QpdXoorX1yw45l6pXMm5ZGZDYzQqgN8RStQI7Ft37smwp+2927jyPSK5uXZiCBzbW2MjByGa+wF+BpsXK4oBvj9CTqBqJO7muZKNmg5Yxbs6DgPGMOR4kUY9oHowr+e75EH1kxVbrdViWdIa7b4MqYRRu8GRVOmFhfxWUFEWauc4SYdex7sDMx8vTosaj9WOpeSldKm71ftszaXKoqs7jmrf8a6vqEkILy1OHvbemYMRWxikQKIfI3ADl0B+3lSjuopnCocCo6WUjFcHTmRQmvya3usGA/OcUzvhxM3feiZYhMl1biFiBmOzlbCqC3XIl5avbAhJxEX63gwFA85fjDGcF5OjqGXWzS2aW0A1sUBgzM6CpNp53GKhtiaX500SQVthLQoe2K2ZdY6ZWPmUOojkzvm6pJLYh5opb8xUbHF+zgVVWKRChh2Gmb3eE215cexz8VSJWipa6FLFzHc8LiE8IrbNGXKSMwxnxxKbSgqRQGdzdSOgkuOgvReCdRzO8y0gIs2DWdXPGsYijKp1t21m49G05MtWOmft0Y4ClByikZMWh+RyXGXaZGtKinAfXVWbjJc8GvADE/Tmxm6yFf/LzBpfNNsyn6ZHKGOSA/C6duBQ2EH6zp9lL7jU9haokx0XD9A+GvYBMLiC0WVuQrLYBGIDM+XRh70iujQ/p62l7ssjL/w3XiY4XmA0kvLiPHeEDlWS+3Wo/tAVzG0toSDWA/iPeVdiprn1fmEXONGJLOtIBr03MiraIcUdKTMvbxNe6rkQJ7OtR4fn49XaJk069xfeSWzIS4EPf+AAoXdnjNxjNU/ybDA0zmAYDbeO70kG4ajs7PQqRErxrQynxYcK+Q1aUE1Ixze7rI0RUEW2dgTDMdUwIhEmyuGgORA3BEdY5YnJ/eNA/EJOXPDP98bv2A4KbXktePjOVc/dOlSoIkrAu2FSIobl41TZQbmgxn6DH45HddibxcLWnmKHWXhllezY5muZk1vH4mwvyZRzpN5tMWVQwYBpVr5Gio5yckyJsklZ2rtLKen0gcS2GwmL2TSY9c8+++yLvaOyPGj0WSEyJR9aGpF4S5UYpOmHaKqo0W3gQy7gYa+F9ekY56h5H6jkgWoNiWdrlSGJFYY6Avw9HR/gALMrgEUR+A3jlgaeU6qVyprPr8dKs8KcqjXewuKeLnPgSg6pqmFUNkaHhyYxzVuekgFFWy6x1VGFP6s+zJ6k0ZI27jAfrqVcOw5dfl/q6iKMSCKqyJhAKHVP2BGo03WwcTHIf4cv5rb690sdtOU6pkx675SPkurZ4STufchLwx3bBkc79gbYUq1L5bAwz2Htp69PmbdzETr3Uz59866lInlemxA6N/3vc+66kMFzsrAE8SE35kbUsvOCY37FIE116xm779su0BU4fvQDbpnmPAdGPMUMDoVHaB9RKXvqNNwOGnYeNNyBOTp4egBo3KWnpx07DZc5bj89d0Wozi/G/At2fxrleFa5LdzauIjClmGQklkevqNf+LOzdnESB1hVvXjPqXMI8eOKfmGtqVNw7DeZPeK9VW4G/hJ2YgujOW8u1VU//GtZZDHFtCSw/Jk5Mud427Coz+dpi6ah3aCyoGpvw30TxSTjwu+pt4r7NhOq/ekyxJnkfKa8QIe0QSwF8/UKg0JHbHo+XNxCFkKU8BNAu4CaIRfsjBrub99+8gV9HmylM2OiKyCmQbrLamVzaYSXx9lzq9BnuEn9YBaBCyUkDRZkLUzqEJn/vOz+/OBeHo/KBsG/x3aVi2X+WV5apR3I7+fxApJu5CK4WuJ8cSU/DNN4zVABGUBqDVrgSwuyPDipySffRpepPpc6IFYsXHZI9pLEsGHZ7zb6sDcfvVQGe8UT6hLXAaI0nqg9C350SnDJRDK4xnlK8sqrzqlJy89O8fy8tSE50qgqTz5f4IYORT6QXQoOel0ovNiQKCrMTEFYHX4EP+I8S8XlRx59h2rIdZ0jvwVVWn4PyLHnxjqOvcNBZM6RQUUzX2PAFdwVeUrETCnL6pL15qqYnxfRTctsqo74s7AIau25WDFFBDgOHxhe8XMOwA712W12eX7ESydAb4N3JE8DoXsyxGEBbk1eEqrhWm9pnjRNABX0xFEl4XylbKvncCLBkN9P/h78GD+ZWFf8eideY/Qm1eXa267rwjiqaRwmiiwnGItz/IqErmfouQB2+bd6Sd3xLjSE5DHnnEt0hWCRAo64CxVbXsFv9Y64Rp1xdEi5d8+/W/kdR9CTxu1wrTNrpB1KRwJb6TlpPzh3mpiFUXklGBlSJl7ma2aTdMjuTVlNT6P8wKEky/1aQA9gYd66MokXiLpdTnGDxBSRluguqENm5dJbz2UaKu+tbqDpZbFdrk+pTog1b1n9BJ72jIt0tb/JzLuVQaBZKg8OKIIa1BccYc05VYnw5V1zqwKhrr33Uv/MeyZ4jFJJGgk5t6eiay2M57ILkQkQB+3oaXZiaMt8xYvXzOa6bQnTvG1cl12PRxmCPEBDXvDOhA7GcSxoGUK18Hk1zl2R365MalQ18y5+EA1mNS5rklwHCcW4UOLgpoUd+nQpgejvXHpGrOW4i/soRyRwSje3kNutNxiqTGBNKsrOTJZ5zxQAUKbIbrMBvod/n/LmJsnrlURWXEZEknM27GY2MbPv3tX57UVIXr5kVJlRs0/F09zotLprGdl7R3UM31Njte3q8uIssAwVKXI7zJqor9kdgd1GpWHkyLD3oja5hY3HkIJpozJOujaX2U0z/I2lM7vaJVNf8aq8Nxwu22JuseEaoxHoudphj5eHFEpc0q5Vla5bK1oFlYeHQvG+J34uF5qEtAr+jtOXtQrX7ix6UgQ0+pAilNOPJ3UJWFiUDVNGThmvyBM8YHmYMrIBSszkdzW5t7mozvFwHFeCw2TZwaAKiYwMP8M5UF78no+upZSkaUpOD1vT76ija2zodm/kjLQ+0FWmOdGOCQ1/xYlJytyQjO3Yh0gnP/v+/v27L/jBYFeb2vFdBwX5/NmVQDGzTBI9O+bn5F1xtp4H6da4aCPRj5LglLWOrBaOIBiL0jDnHHzauicqs2atpTa+gIxNFY7tfEjzl8K5a4gGTJ+GH7ynvFRoxoHTku7RAwmR53eA+QS/SXtWMJm5WmXi5gP6+sFFJUCOVKUk+pw7IWKrl821bK5bnZBrGeuSzHIhvKMk0caUNC4U45mfKtf5hroKlyCBeI3BYYdSrI+GBbiVIkWCMH+CZtay0qUy2we4UiMVSogDrbqUSBCjslokHqIxIAzwuFajDo4/4nQQ4xDbc0NM0WSmi5ZgzTxceqnbzOXObi1J0aY9zJaKjWqorUQ7jkpKX284jPlZUvEo3pvMbMcYUMogF+Np7VA6JEdDLF1b96cQJXZXAdp14AIMhACNfqgASWowPUoQ1BRfyY+4ZnHQ8OU3yc8iJbjGl/zrc8y/Va0rqRO7xy2E0D6e5135GTRHri1Zun5MKrtajcGJ9+EeIg/vrd+rnjlmm++oZ3Z0mfrSPFxqpd4KmuhPn1nIo8WvSbNvXYINEtkmZWXzVjTkI19I8mtGVuD5ejwSccRwOWmh2JG2UlbyonNLQzdFFzTjbZxfG5/5VhYJzj1cYIIAYC3RnegUjGzfvv3kC3kcUNXYEawkg5EehqXRlmFuv2fNpbM9bkXH3Gbj5YIkLvED5Vso5pWgia0dhjy3C0ym0BjJChSGtjxlb1vy1Zr05MPM+M02uu17q6WlvLG5Hall1SLaIffLDxeS2FaRmlSf+cF4qjzpJtopxzk0WhGbwIOGqgeOVVTZBxArNxRfREf7gVQWkZRu36awHj88+BxlUOfNKon4i+t6KcoVDrlYh548KpNa6hwUDIXDtoreBUiTEsk2eh1so0lm672YRdkdlDNaRi+yk5ofiqJCW6SGn//dc38WYYYFvfcmtRuVc+Y1QqBsau8aVZWT/kBqHEopLkBTNdML+3N4gNT6a8TochVrrLdSIk81mZsVQt8K7WFhhyzba/cnEnSNXDlOXZDFG5zS/ciV8WPccvqQBhBCXqspnQNMIVTXOyCXxSGrgb/7E46kRr0hEYKOwMOzwN8vo2eJNtkwqPL30PC1xfoS847PS9OlRxRH0zj9+ffMZ6BfvfaTzGO3uqgDrDoafBuUji5O3PNYUh6ujXTRWbssfi0WFOxQ7cgZOiMySNYIjAXbD13mGgUmnywZfXEkx9pB9nVhhChkGqAtk/KtDv7YdzTsbMm7Q8hezts7Meecm4/+5Loe16WKtfVeB/3a9nhcikv0H8ymxyOdp/eHV8XTEqncekcrkBW4XjrGPkaZ2kgryM+j5GrWXAafgMElf+U5r31cF/6enZyZK78HjieM1F8u7dzFxOpj2HXl5fd4ecCh24xLjZwtbs1us4VGbOSjHsa9FlAPObbxyAf66hwf5KXGVDWzwji7m72+3raNSW0H+gGLTfOs3oklCThcNRfHCM8BZ6RIgXjmjnbWUWE3c3GDnrD9UTwz74wKXYrNNFwoox+jNFTGBuhkyU29cq0xmlFRorFXeUwC6pN9FC8UfWTuCegBOxQxbFLxHVnaeRrqZU1/h8Nt3fA55f5rPLpiip3GNtvioMmH4RklKsURNfiHUzsLhXHYeXHwcIzUGyrcg/15LKgJ0Zsfbu2LCNtjZU6JtJlZNIOXopWIZJcBj11IHpQNh1SDKXTJQJzZNf3JBHy/vmo/lp17N1+HeNDpnfLqZtHFjWvI5+A9RRb5DDvFSzIRxkJFHlXc3vB26Z97g1BhafsV3gQJzEhcMaOtt6azrTG6+qjeuZvbGI27QQhzXZKfBzwqnDwU84/Sepqqc9zZR7O5w1q3RPS4IRu+phRuzWKhUwFJoyHIiqFUVGvxWOR52yGa4YW39tRYzyEy6D3VkOyILbgKaIpk6NiPnjJ8frE+Mo53xqxL15v1d2/ffVHfVenvJ1zYRbFMU01HixeHIojI8aBvQW1v03ijX+XsZqoVXZRuIfXBXju9DpSgQr/tUpwMcZwMznd3zqT9CXHCai8O1IAAZzsjXInt4BhkzVVsGY5xMMJL/DLghMCOpKO+KyiHOxdVnJD29WNhT2nmXqWE8KihsZaf24AFyYqdGch2VLU1yvDDMV+gyTlRFESplSSPdBf8kKWPM9tgzcMhbkLE06HOLpQZIax0ydTiYpcIlPrdVDqaw/0acnJBDbRS0RKoGmCpEoyPOS5UuOhA9lLOlWTXRSo5cTCmzBI/qNFHfvzpL1A2DWWiTV6bfqBnuLNQN8bcBGdGW4i0wDFnY1dNTtWBOD93UJjLouKUYl7fr1vTIdyImbDaA+VIaMBcV+MijpSk3mPmPS69vWokKlICyw3y1EjxPgQXDD7jgWta2tcuTcII6IvjwKIzRvV8Z/xg6mnUdl5YB6qe+8XRRl2wu8jQNLUaiiCOveinys++kChEH5kVJp9nCRfUPGiFPQGFgb4JO6gd6Z2BuIWnaUtDLAszdso0mta+0so/JomXVWrj4T5n8VdOetM/xxyt3kOO1ZrpWW+IneAomTL3jQlG//GPf/wFxxgZY7uk9qChqFRatZQeCAYioXEfypK9kRWCB6AxPYsHKzHCaL8mmPgaS3E+/iRLDewFpqBoO8osw9Y8f1EV39oOXAcdr03O+grdWTtFAgOhTAvmgoWl/VrpK2HMq8vz5uJV5WLTtUhtwhDkL7gT+YAZ+rgecLHOIoBCdST576oD03tVU/wsqYTiwxEAS5LaOTgC5C6GOyhGUp8iAnzeuZgNUZfVHvd6UN1KuUdDG7MnOMZi/gMvrOaneoqXpYHdVfudfhWyXj4AfJaa+TbTjoHQR/6+ecHSQ9JUMVlx3EbXASfFUasxEDlefNnyUBvl3j70fSf91Lz+e462rOSix5xZJIET092YGpgjjlJJleS2S9q663A2msW2zIC91UiF0crEp6+dfKPcPyyp6+bEpQlTmaNy5a3Rx5XPOWWf2J1Jfn4URwOqOH0emP2PXkWpaA7oShWtHFEMLOOyv8sgF/KFmD4fjXO5yCnlRoUpxQF1xNyfHg1eRPSTXCM70vA4YmL9KU+8BAiujiu7JkYtdI3zSypmOvdIg2goohfOzUWIrdISqbJEwYJiaOtipnqt6bJXN2jPQX5ExXBceK4qWkME+GIolT193i7DeO2O+tu3b7/w7gh+mXqZmUIntIbmwlZVkNVcU5XqXsl6kkSsIE4Mup+s1AP5EHC6xor6wVouEgb05mtve3l5eSLZ0s3JCkvzT5js9jF29iNWlol63hMz3xHnSwUK8Rh86VprydqJmstDTlbxnYjl5X7EW5rudLNzDnxA8HM+XM84oZIe1G4zZyQ/o+JwwX2NB3jvQ44bxyJW+JhD4okdiGnE5Yfh0aRua1JyNLnGjXj6Vf4KHqXtyHxwUFGZUcBL1D1+EB0K9z2MSh6V42xRF3EcrmfXoXr6HEqeWIIBmBgHYgP6YZJ004FSmJfaU8XBYuOBw26LuxiDqdQ89y3rPkgAP2AKZXd77G/MrCV2Gc9TLZzZYcpQ3ptiCPj7yaVt1yWnXA/toLZiSbMabkVHuPKZ33vr98L/PFB0MZCItInATo8oHXev1MJWWzalVuKZIuwz0d9LFGBHMqb2OpLO5thuPK5aarc6XzjO0iLcygfjVn4i9m0kCfdryOyoXBEcqFIxRVXZ5YPruqi11zgqe8YvRPxAzNXakwCEMERvlWAY7NKxVzFdKK6gKJo4OwUZyDHZsELwMu9PkdQoZkHEKLk5C/vxpLaip0Y5SGEaZ3I6M+DnIs6ktZYZdWkIz0N4y2jDPUM+7I+Xl0yfWlNcLCmftKMoE093Is1zBHI9HiUNZPTqdZm3dEM2sHAC8YynVprKpTUJLlw63GXeQjseWP710SGFrZzyBSXYvFdWHWPkrmTlB70m5+XtaO+KbmsKa1qSv/LTfzwu65jjRkR2DJAy0k+V5xUInTCbrYWqepm4NGNcyp23lsa18eg1FpgzTUpXxnw2K6hj3eWFtMi9Tpcy5DTf1WHqWipmV7bLNQyIoHschyNYYK0qbRriwthNhOTHsUMy8JCfpbAmfujPN4KomCPSgWR3yJ9L3u0yryYdPTQuoJ8oURKR4gN2B8drxha+weQZtnBAdr3sNHtWo1HjnBDCw5PiQKzLQY2Npz+3bc9pgYPXcTGwuqU/x5FBI68GqqD7XvV54QMjeXqv/Lw3qvlrXKAHhPK23fNdZFXFfUsVdgikYpaIBwy5l/kGPqORZL2L4rqyYiXepV9dB3sSpgE4xDNrx4gxuxp2QaYR7bpneYRWPF2g5sVnSpNcjqhyoVsenUYfDj0aO8fnDn/PIj0YuSV7BhSYO3PEkZMTymK3p0mHslNWinAuCAvkU4NUfUkgsSX00E7W3cx7ess6gYgoVGJb9/KqMXBr9IYpTXZ591x2z1shXr0NxSMYyBjsrLWIj5IR059EBSYf/LVWPi8i/i7tOaks89asv3v/7gseAJm8FuXExINshPhZUR2NFdoxpeVLRcidfhDM+1hBMT5TV7IdgD4FSoVkhHRFG6R9E3PcrCQ6DqshMGMuqdHBtqaR1sYiK/CZnu1ztvC1D7GVtFSOGfYxdlNFD6nnFsSqKv29ciR2vQwt1rm/IaDQhPUAdsQSzrYwfxy96zCmwmLh0LAjzGcJSkjejgalWrSVsguHFpEdIq1WRkgF6CDe10ovHziU4kC5B5IlI45AoCPEqw2vqYKUVRtadDtIwvHUQVLmzTlyZTGH2fnY2OE92mVE5ayZeA0qUDiOIM+nHZJMjjLIOWOFqaX8ceGl/6LkoRaZGtk9YZCSWjtjmzu2BVsRwPQ//XD8IDgmx6UeT4owIU0ijrS+/OZTZlrkes7oJ6MSuD9bJEFs7QUc+88xehZWRxBNIAOE8lruWyRIIfmZEQun0ovjQIkgTF4KVskcjTXFWeeX7qM9JXmeoXWmEW6rXJNg59gTbnoUsLm7ObJ18Dlor4e/a80sKMaD4ML8PuIQCkmSfozLlLljp0Ksy5NyfzgYa5hAxNpCG7HLalR4CiXfy3+FfRsLijP9VV1iVIGzVbSVDSELhS0neu6B/InsUGKolA0PdOLczUW49U8//fQL7S/A8vFm8itQIcUKtfWRI4eDoxNwmTNkKsLteoy8GcMl+QrmF/AQ7QdRtbuqjXbI7lxcpco85mzbn2JRTdGLWsz1hPTxcBmDYDuTHE4HENRDtazNGSere+6FFkmp7VjycXYtfT2hckC4oNVsgBBS9jpX7p1q8W7K9eA3xs9jH/gLLdZwwLTmQiRQXns6jTkWzLUOlStQySHekhJULgbbsffaqgAzlY/L6YbuwFuNo9jJicDbatzZe6sDNfxQx5h8ER4uoJyd465dh3bv9f0buh66kPmZa8R24Hk0liKgMCiDXWatG4v7EFCv4I4dqXaNggRJkcvdb8fIy5ofyA1TMdR6/Zk4si9Elb5KLsvvM7vWoXk1nxE/Ot9tDJDyClDjsvMY9bVjqetWS+4Sd4SWpXRfB9WKEskUkZCYce9+oGPakbPhB7I/zYySvVtJbPMCqKxx+qT8iGao7B8/9qIlgBADuh07r2M3wyyhMbrGXCnFrx2r2/H1+Bq2EsZ0ImWinOvBCp2ZQSoya0TYeo6f26iMD3YtDd9DOfL9SdTByF8HPp74p42s+s0kxGOPuAMqQnxN+qf2Qd5Iavp4EmywyBsjaeq1gyQJ+ii2mtcOhNTQAWYRVQB8AMy2WlBmmqcRxrQwZKtqUDPRnU4tci5qy5gyYUJyOEXDCsuRKp7QgcfqpMBqW4vAsBOz3moJSvkZF2ataUHkxg9vpwcBF9LAIW9E3M/ERQu5HWHX46GMdKqhGqSqshS0ZAMpH95KJXHG7r5588Yej8te71tjBb70vXVdYPugjB6EdEmrT0yDHeY0M4x/uJDDqIJa/mcYZI0IlCrJ0CGp10Iyb4WB4SGjVLi6ntBnxOqsuoqCKfZWhiaO32S44zMmfSvBeKFxJ7MQ2JWMa+Di2mqNuWQO4s1PthBnvrPMV9wB0ABLk6BrK26HQiikOmTl6d4KU05b1UF5diiVlCQneXPhd5Rvg5m8eyCQqd6VhSJmXED0AICnDBUscSdTFQEC5CiM73jHtGH0Y690qHkMYg4u8Jv8RPnZ2q7iJ/97YUHk49Hz6WLNcQG+Uz+qRfqGw/9otp/QO2XMbCUHt3Oqvqs4DYRv4efhCCo7qVn7AEhT2VW0dqjqOBE4/Dje8HvrJVAgOJK7SRbJPJ15UWiPopTRvGSSWmGSoqvTQQFLwrJFjaQIjOU75ZZ7uXmzq6lnNke0oa4ipKxzIV1OkoAptmI+ec54rvb37z/94mQZbcRYmhe4Tg5GWfVrlkmFkEmRlH9mXCMP4PuWT8MwL2aHkIeYiVvDb05IZAwSmUvAH4xV5LjSNc4HjHI9zucIZeO+JI4x1H3PzFvYZP+D3opDPz98LgTtkC13jX+kNJPUD4vFQ2IbR0AWq5IGBVNrw14eD3u9byQJVo6Cw1j0+jpVyTg6ATrvhenG6MyjVSBP7FRLuJ94x6Lh4odqECnw5ybwMqRmqUOhHbyqjXwCykf5vPGSdy+/iVzxXgYmjZG40G8HBwrU1X0yiwB0c40wmIfBdt6edmOxFqSoNfHS7/JwMYcCebY+40LG2PHClqyRIoINAYfteDLn0e3L38+OXUiJAwjJQ4d4E6MRz0q5le99K2HAiooQOEQtuSfEi79LIKELq9WSl05yLr1PMYH+fzxn98Thgi9dtGks1Y4DnMRcGQdHexLfUNVJskGSmEu4wYtKI9koLLvAnlDN8bIWm4uYjYMlpYuHhR7BjEcHT9wIzYabUbIcT3LMrRHyD1z6dgReHsiXWEW9IPGi9SrozstLDDjE/nZkqLQ2FAQWwPVoRBt5jvJ51D7UKrwqdPH6QReBAsxbJbDiZ1krngtJGqm5WMcoG8UxsnIh4euYh+3NcJT2pNfeVgudfrRS+gXhvtt7ZWU4chG6Y6NSNKmItuS7QLw3B7U0jkVoVn4Tc8JOdARwDvugzop2iiUVcSpuDtImk7i2IhuzQkKWtCEVD6owIgT4gkUUxjl22P3hfpop0+251kT4VT9UJGGtJcWVLS3HTq1RUofFrrlyDq7RxZ8KN/v+wysumn0QTnkxJHI9Z91u16NQ4xwJzNdZ8EUKY7GnYjstWaX07xV1mulsmevAkYrF1pJ7XJBntyFkTRtFeqa+Pl3x4EPdSz6Axbk8DwLsPDSgAGwOGkPzngcsTVJpyHSL1qxf3dw6uoimhTvVgNTus3uigmdN7jVchtaNHBXD5dFb8qwmL/uGfHqMJwzqMi71F0gCypA5zGADxrqMPe4aX9KJTHknhQNbl/USRdZsS77rrZhxiSGfeD6zQ+Hfuw/5s0ChnX6jmndvfS5Te8iIkrSaUeoMn4cVVqT2Wmd7kN0BCbAam2Uwh625gCcxJZeSAixCNkU7I/ccdHJ7c1sMo9oVHyuxxkHz5Y4qZDis/KFcKJvGpQxo6hDIJCqpZMiqyqW6coytoTRVaBY6lO7acQn2uYqTlUXBrX++j0A+xx4jnDLkhWcWo7VI43Ka/46kYIzImO6YZwqMigMcu7msU56N/QtnlUoXbU3Pen/79u0XrfdqlanlN5j0gA3goUPiI6Mjiebgz92vS5UMaaisTwcc3lIs7bALUlRsFUsDTbUIYGR0q641RcjcAtkd8acnOnnQW+JKVatFHt2XXdgIbh49qk1XBecVi0k0iWthXjgUVslCXQPQp9lqa/Z4POQ8b63Zh+8/lM+FBk5xtUrCGNjnNMtOoz4nPx6UDXDl0sy3ITAmXc1dcamsdjSGimJfMTuFS0C6ZZ+X7KEXWnG6VhXPSVY+K9JzVhuY48oXYvY04w6r3HoVotiLaSd0sH5cksZDbhu1W+NB6hAneHNRSQOoGXZBjWoeZE5rpNBdLCFVudzzKGSr8PP08NCMScZVpfFhPNZrX5AdYQEl95E5QcNi8tMKisgAsrNqFrXhlCYfJl/hOIwMKz9ECBthX64dXkDe6r0mee5xIEOOhEE/qnErnpbc/d7FueK7qtHigSdn1ZwqRcRM7xrrEjVkR0GQn1d/ou/2IytjINOCysPa6ZnN1xtkh9PM2LQU1/4AJtBz7znGgIu9IqNPXArHwixWTpw6x4gBxA7N0XrwlY9z+GisdnbtIBlwIhIn/U4U3Wc/iveKj/DWDpVvAVk5Cm1C50Mi/e79+y9YjXYvZAQvDHo2eKDuex/VxcIuoqGyr5yPjnm3HgQwo3rLCplL9LUnkBlNeQwxl5bhi271Xg9CWILFqFCIoFIrNFdtownaxy7BvdcysaW+Wf4VZBivNQsChzl0h8eDlxJDXTQ7bK6XIpU/XumCx/y9ISfjuiD5wzrku++/t8fjkVUZFq2uIB9gUkYh7hciRIXGbsAUnAchlnvMNTdPPbxzROAEyuUBLqcuiQJtqBPLaNyr8iMg6eSuiQvH8mPwW82XweF54MvsUICMcZVbVml56EB2XTQMApKLHBf0mktzaZnumh2O/AOtP9phojJJmF26+l5GRximNp6FDmQ68yf8BxdGLX5BJhjsejDGwwF9jaHDjqMzs/2DhMPkJbFCHX3gEjz4ZSvHGPscUWyzdkGyveFDUciSqWjT3J9IcSiweisacOsVgdBhXhXO5li2n7kQhs4wAYR4JiCy4YyexWdd6JmMuZC0Z+C+nfG8JRnNd3DupT3g44F9kFfIlGFJ7ig8K+mzQYIfUg9mMTu0O+gI7SJ2ReZEqqYsUSQdBVVvTWZMRUhA6UUiR1BVBpKHObuEnSijozggOdgRqe3OS68oBpIoa5zpOaXYtZd7gTCHEn9eflk8ACvkB2GbwWkUHEES3VtFcmREOVRqY4A80e3/AWbQhcFn5w2+AAAAAElFTkSuQmCC"
def get_b64():
    return _LOGO_B64

def show_logo(country=None):
    b=get_b64()
    flag=FLAGS.get(country,"") if country else ""
    # Flags hidden by default via opacity:0. Canvas test reveals them only if they
    # render as actual graphics (wide) — not as fallback text like "LR" on Windows.
    detect_js=('<script>(function(){'+
        'var spans=document.querySelectorAll(".tp-flag");'+
        'spans.forEach(function(s){'+
            'var cv=document.createElement("canvas"),ctx=cv.getContext("2d");'+
            'ctx.font="32px serif";'+
            'var fw=ctx.measureText(s.dataset.flag).width;'+
            'var lw=ctx.measureText("LR").width;'+
            'if(fw>lw*1.5){s.style.opacity="1";}'+  # real flag renders much wider than two Latin letters
        '});'+
    '})();</script>') if flag else ""
    flag_span=(f'<span class="tp-flag" style="font-size:3rem;opacity:0;transition:opacity .3s" '+
               f'data-flag="{flag}">{flag}</span>') if flag else ""
    if b: st.markdown(
        f'<div style="text-align:center;padding:0 0 .2rem;margin-top:-1.5rem;display:flex;align-items:center;justify-content:center;gap:16px">'+
        f'{flag_span}<img src="data:image/png;base64,{b}" style="max-height:280px;filter:drop-shadow(0 4px 12px rgba(212,168,67,.3))">{flag_span}</div>'+
        detect_js,unsafe_allow_html=True)
    else: st.markdown(f'<div style="text-align:center"><h1 style="color:{C_GOLD}">Teacher Pehpeh by IBT</h1></div>',unsafe_allow_html=True)

def ico(s=20):
    b=get_b64()
    return f'<img src="data:image/png;base64,{b}" style="height:{s}px;width:{s}px;vertical-align:middle;border-radius:50%">' if b else ""

def pprog(stg,tot,msg):
    b=get_b64(); lit=min(stg,tot); pct=int(lit/tot*100)
    # Animated pepper icons filling up with pulsing "Generating" text
    pepfill="".join(f'<img src="data:image/png;base64,{b}" style="height:40px;width:40px;margin:0 4px;opacity:{"1" if i<lit else ".15"};{"animation:peppop .4s ease;" if i==lit-1 else ""}{"filter:grayscale(100%);" if i>=lit else ""}border-radius:50%;transition:all .3s">' for i in range(tot)) if b else ""
    pepfill_fallback="".join(f'<img src="data:image/png;base64,{b}" style="height:36px;width:36px;margin:0 4px;opacity:{"1" if i<lit else ".15"};{"animation:peppop .4s ease;" if i==lit-1 else ""}border-radius:50%;transition:all .3s">' for i in range(tot)) if b else ""
    icons=pepfill if b else pepfill_fallback
    # Status messages with model icons
    status_icon={"⏳ Teacher Pehpeh is cooking...":ico(22),"Asking ChatGPT...":"🟢","Asking Claude...":"🟣","Asking Gemini...":"🔵","Combining the best...":"🔀","🎨 Creating illustration...":"🎨","✅ Done! Content is ready!":"✅"}.get(msg,ico(22))
    done_class="animation:none;color:#10B981" if "Done" in msg else "animation:genpulse 1.2s ease-in-out infinite"
    return f'''<style>
@keyframes genpulse {{ 0%,100%{{opacity:.7;transform:scale(1)}} 50%{{opacity:1;transform:scale(1.03)}} }}
@keyframes peppop {{ 0%{{transform:scale(0.3);opacity:0}} 50%{{transform:scale(1.2)}} 100%{{transform:scale(1);opacity:1}} }}
@keyframes barshine {{ 0%{{background-position:-200% 0}} 100%{{background-position:200% 0}} }}
</style>
<div style="background:{C_NAVY_L};border:2px solid {C_GOLD};border-radius:16px;padding:20px 24px;margin:12px 0;text-align:center">
<div style="display:flex;justify-content:center;align-items:center;margin-bottom:14px;flex-wrap:wrap">{icons}</div>
<div style="background:#1a2744;border-radius:10px;height:14px;overflow:hidden;margin-bottom:12px;border:1px solid #2d3748">
<div style="background:linear-gradient(90deg,{C_RED},{C_GOLD},{C_RED});background-size:200% 100%;animation:barshine 2s linear infinite;height:100%;width:{pct}%;border-radius:10px;transition:width .5s ease"></div></div>
<div style="{done_class};font-size:1.1rem;font-weight:700;color:{C_GOLD};letter-spacing:.5px">{status_icon} {msg}</div>
</div>'''

# === PROMPTS ===
def _p():
    return "You are Teacher Pehpeh — AI teaching assistant by IBT. Warm, wise, practical. Speak like a trusted teacher/village elder. African idioms, proverbs, analogies (fetching water, clearing farm, cassava). Liberian terms. Cheerful, encouraging."
def _g():
    return "Groups: 1(0-4 siblings,most time), 2(5-8,limited), 3(8+,very little). Consider socioeconomic, computer access, parents' ed."
def _r():
    return "Rules: stated resources only, max 3 problems/group, self-contained tips, WAEC format for exams, African context, Socratic method, print-ready. Use **bold** for key terms, objectives, important concepts, and answer highlights."

def build_sys(reg,cty,grd,subj,task,cls,res,lng,abl,tm,top,sch="",mano_ctx=""):
    s_tag=f",School:{sch}" if sch else ""
    mano_block=f"\n\n{mano_ctx}" if mano_ctx else ""
    return f"{_p()}\nCLASS: {cty},{reg},{grd},{subj},{task},{cls},{res},{lng},{abl},Time:{tm},Topic:{top}{s_tag}\n{_g()}\n{_kb()}\nPhysics=extra scaffolding. Group 3=SHORT. 58% no computer=paper first.{mano_block}\n{_r()}"

def build_chat(reg,cty,grd,subj,cls,res,lng,abl,sch="",curr_ctx="",mano_ctx=""):
    s_tag=f",School:{sch}" if sch else ""
    curr_block=f"\nMOE CURRICULUM CONTEXT:\n{curr_ctx}\nUse this when answering questions about curriculum, topics, or lesson planning." if curr_ctx else ""
    mano_block=f"\n\n{mano_ctx}" if mano_ctx else ""
    return f"{_p()}\nWhen greeted: 'Hello, how can Teacher Pehpeh help you today!'\nCLASS: {cty},{reg},{grd},{subj},{cls},{res},{lng},{abl}{s_tag}\n{_g()}\n{_kb()}\nIntervention WORKS. Teacher's work matters.{curr_block}{mano_block}\n{_r()}"

def build_free_chat():
    """Free-flowing chat system prompt — independent of classroom config."""
    return (
        f"{_p()}\n"
        "You are Teacher Pehpeh, a friendly and knowledgeable educational assistant for African students and teachers. "
        "You answer questions on any subject — Mathematics, Sciences, Languages, Literature, History, and more. "
        "You can generate practice questions, explain concepts, give study tips, and discuss any topic. "
        "When generating MCQs, always format them as:\n"
        "1. Question text\nA) option\nB) option\nC) option\nD) option\n\nAnswer: X\n\n"
        "Be encouraging, clear, and culturally relevant to West African education contexts. "
        f"{_kb()}\n{_r()}"
    )

def build_stu(reg,cty,grd,subj,cls,res,lng,abl,info,sch=""):
    s_tag=f",School:{sch}" if sch else ""
    return f"{_p()}\nCLASS: {cty},{reg},{grd},{subj},{cls},{res},{lng},{abl}{s_tag}\nSTUDENT: {info}\n{_kb()}\nTargeted advice. Compare to data. Risk factors. Interventions.\n{_r()}"

# === AI ===
def _sub_blocked():
    """Return upgrade message if subscription is not active, else None."""
    if not SUBSCRIPTION_ACTIVE:
        tier_msg = "expired" if _SUB_TIER == "expired" else "not yet active"
        return (f"🔒 **AI generation is paused** — your school subscription is {tier_msg}.\n\n"
                f"Your student data and past content are safe. "
                f"Contact IBT to renew: [institutebasictechnology.org](https://www.institutebasictechnology.org)")
    return None

def ask_gpt(sp,q,h=None):
    blocked = _sub_blocked()
    if blocked: return blocked
    if not OAI or not OPENAI_API_KEY: return None
    try:
        m=[{"role":"system","content":sp}]+(h or [])+[{"role":"user","content":q}]
        return openai.OpenAI(api_key=OPENAI_API_KEY).chat.completions.create(model="gpt-4o-mini",messages=m,max_tokens=3000,temperature=.7).choices[0].message.content
    except Exception as e: return f"⚠️ {e}"

def ask_cl(sp,q,h=None):
    blocked = _sub_blocked()
    if blocked: return blocked
    if not ANT or not ANTHROPIC_API_KEY: return None
    try:
        m=list(h or [])+[{"role":"user","content":q}]
        return anthropic.Anthropic(api_key=ANTHROPIC_API_KEY).messages.create(model="claude-haiku-4-5-20251001",max_tokens=3000,system=sp,messages=m).content[0].text
    except Exception as e: return f"⚠️ {e}"

def ask_gem(sp,q):
    blocked = _sub_blocked()
    if blocked: return blocked
    if not GEM or not GOOGLE_API_KEY: return None
    try:
        genai.configure(api_key=GOOGLE_API_KEY)
        return genai.GenerativeModel("gemini-2.5-flash",system_instruction=sp).generate_content(q).text
    except Exception as e: return f"⚠️ Gemini: {e}"

def best(sp,q,h=None):
    """Try text models"""
    for fn,nm in [(ask_cl,"Claude"),(ask_gpt,"ChatGPT")]:
        r=fn(sp,q,h)
        if r and not str(r).startswith("⚠️"): return r,nm
    r=ask_gem(sp,q)
    if r and not str(r).startswith("⚠️"): return r,"Gemini"
    return "⚠️ No models responded.",None

# === VISION AI — Grade photos of student work ===
def ask_gpt_vision(sp, prompt, img_b64, mime="image/jpeg"):
    if not OAI or not OPENAI_API_KEY: return None
    try:
        m=[{"role":"system","content":sp},{"role":"user","content":[
            {"type":"image_url","image_url":{"url":f"data:{mime};base64,{img_b64}"}},
            {"type":"text","text":prompt}]}]
        return openai.OpenAI(api_key=OPENAI_API_KEY).chat.completions.create(model="gpt-4o",messages=m,max_tokens=3000,temperature=.4).choices[0].message.content
    except Exception as e: return f"⚠️ GPT Vision: {e}"

def ask_cl_vision(sp, prompt, img_b64, mime="image/jpeg"):
    if not ANT or not ANTHROPIC_API_KEY: return None
    try:
        m=[{"role":"user","content":[
            {"type":"image","source":{"type":"base64","media_type":mime,"data":img_b64}},
            {"type":"text","text":prompt}]}]
        return anthropic.Anthropic(api_key=ANTHROPIC_API_KEY).messages.create(model="claude-haiku-4-5-20251001",max_tokens=3000,system=sp,messages=m).content[0].text
    except Exception as e: return f"⚠️ Claude Vision: {e}"

def ask_gem_vision(sp, prompt, img_b64, mime="image/jpeg"):
    if not GEM or not GOOGLE_API_KEY: return None
    try:
        genai.configure(api_key=GOOGLE_API_KEY)
        import io; img_data=base64.b64decode(img_b64)
        try:
            from PIL import Image; img=Image.open(io.BytesIO(img_data))
            return genai.GenerativeModel("gemini-2.5-flash",system_instruction=sp).generate_content([prompt,img]).text
        except ImportError:
            return genai.GenerativeModel("gemini-2.5-flash",system_instruction=sp).generate_content([prompt,{"mime_type":mime,"data":img_b64}]).text
    except Exception as e: return f"⚠️ Gemini Vision: {e}"

def best_vision(sp, prompt, img_b64, mime="image/jpeg"):
    for fn,nm in [(ask_cl_vision,"Claude"),(ask_gpt_vision,"ChatGPT"),(ask_gem_vision,"Gemini")]:
        r=fn(sp,prompt,img_b64,mime)
        if r and not str(r).startswith("⚠️"): return r,nm
    return "⚠️ No vision models responded.",None

def best_all(sp,q,h=None):
    """Query all models and return (best_response, best_model, all_responses_dict)"""
    rs={}
    for k,fn,nm in [(OPENAI_API_KEY,ask_gpt,"ChatGPT"),(ANTHROPIC_API_KEY,ask_cl,"Claude"),(GOOGLE_API_KEY,ask_gem,"Gemini")]:
        if k:
            r=fn(sp,q,h) if nm!="Gemini" else fn(sp,q)
            if r and not str(r).startswith("⚠️"): rs[nm]=r
    if not rs: return "⚠️ No models responded.",None,{}
    # Pick best: prefer Claude, then ChatGPT, then Gemini
    for pref in ["Claude","ChatGPT","Gemini"]:
        if pref in rs: return rs[pref],pref,rs
    first=list(rs.items())[0]
    return first[1],first[0],rs

# === HANDWRITING TO DOCX ===
TRANSCRIBE_PROMPT = (
    "You are a handwriting transcription expert. Your ONLY job is to read the handwritten text "
    "in this image and transcribe it EXACTLY as written — word for word, line by line. "
    "Preserve paragraph breaks, numbering, bullet points, and any formatting the writer used. "
    "Do NOT add commentary, do NOT analyze, do NOT summarize. "
    "If you cannot read a word, put [illegible] in its place. "
    "If the image contains diagrams or drawings, briefly note them as [diagram: description]. "
    "Output ONLY the transcribed text."
)

def generate_docx_from_text(text, title="Transcribed Notes", school="", teacher="", subject="", grade=""):
    """Generate a Word document from transcribed text. Returns bytes."""
    try:
        from docx import Document as DocxDocument
        from docx.shared import Pt, Inches, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import io
        doc = DocxDocument()
        # Set default font
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)
        # Title
        t_para = doc.add_heading(title, level=1)
        t_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in t_para.runs:
            run.font.color.rgb = RGBColor(0x8B, 0x1A, 0x1A)  # IBT red
        # Metadata line
        meta_parts = []
        if school: meta_parts.append(f"School: {school}")
        if teacher: meta_parts.append(f"Teacher: {teacher}")
        if subject: meta_parts.append(subject)
        if grade: meta_parts.append(grade)
        if meta_parts:
            meta = doc.add_paragraph(" | ".join(meta_parts))
            meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in meta.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            doc.add_paragraph("")  # spacer
        # Add a thin line
        border_para = doc.add_paragraph()
        border_para.paragraph_format.space_after = Pt(12)
        pPr = border_para._p.get_or_add_pPr()
        from docx.oxml.ns import qn
        from lxml import etree
        pBdr = etree.SubElement(pPr, qn('w:pBdr'))
        bottom = etree.SubElement(pBdr, qn('w:bottom'))
        bottom.set(qn('w:val'), 'single')
        bottom.set(qn('w:sz'), '6')
        bottom.set(qn('w:color'), 'D4A843')  # IBT gold
        bottom.set(qn('w:space'), '1')
        # Body text — split by paragraphs
        paragraphs = text.split('\n')
        for para_text in paragraphs:
            stripped = para_text.strip()
            if not stripped:
                doc.add_paragraph("")
                continue
            # Detect headings (lines that are all caps or very short + bold-looking)
            if stripped.startswith('#'):
                level = min(stripped.count('#'), 3)
                doc.add_heading(stripped.lstrip('#').strip(), level=level)
            elif stripped.startswith(('- ', '• ', '* ')):
                doc.add_paragraph(stripped[2:].strip(), style='List Bullet')
            elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1] in '.):':
                doc.add_paragraph(stripped[2:].strip(), style='List Number')
            elif len(stripped) > 3 and stripped[:2].isdigit() and stripped[2] in '.):':
                doc.add_paragraph(stripped[3:].strip(), style='List Number')
            else:
                doc.add_paragraph(stripped)
        # Footer
        doc.add_paragraph("")
        footer_p = doc.add_paragraph("Transcribed by Teacher Pehpeh — Institute of Basic Technology (IBT)")
        footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in footer_p.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            run.font.italic = True
        # Save to bytes
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()
    except Exception as e:
        return None


def _make_mcq_sheet_html(qs_json, title="", n=None):
    """
    Render MCQs with a sticky answer-sheet column.
    Questions left · Answer sheet right (sticky) · Answer key + explanations at bottom after submit.
    """
    n = n or "?"
    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{background:#0a0e1a;font-family:Georgia,serif;padding:14px 16px;color:#D0D8E8;font-size:15px;min-height:100vh}}

/* ── top bar ── */
.topbar{{display:flex;align-items:center;justify-content:space-between;margin-bottom:10px;flex-wrap:wrap;gap:6px}}
.topbar-title{{color:#D4A843;font-weight:700;font-size:14.5px;letter-spacing:1px;text-transform:uppercase}}
.prog-wrap{{flex:1;min-width:120px;max-width:260px;height:4px;background:#1a2a3a;border-radius:2px;overflow:hidden}}
.prog-fill{{height:100%;background:linear-gradient(90deg,#8B1A1A,#D4A843);border-radius:2px;transition:width .3s;width:0%}}

/* ── main grid ── */
.main{{display:grid;grid-template-columns:1fr 220px;gap:10px;align-items:start}}
@media(max-width:520px){{.main{{grid-template-columns:1fr}}}}

/* ── questions ── */
.qlist{{display:flex;flex-direction:column;gap:7px}}
.qcard{{background:#0d1e3a;border:1px solid #1c3060;border-radius:8px;padding:14px 16px;
  border-left:3px solid #2a4070;transition:border-color .2s,background .2s}}
.qcard.active{{border-left-color:#D4A843;background:#0f2347}}
.qcard.correct{{border-left-color:#81C784!important;background:rgba(129,199,132,.06)!important}}
.qcard.wrong{{border-left-color:#EF5350!important;background:rgba(239,83,80,.06)!important}}
.qnum{{color:#D4A843;font-size:12px;font-weight:700;margin-bottom:3px;text-transform:uppercase;letter-spacing:.5px}}
.qtext{{color:#D8E8F8;font-size:14px;line-height:1.5;margin-bottom:7px}}
.opts{{display:grid;grid-template-columns:1fr 1fr;gap:3px}}
.opt{{display:flex;align-items:center;gap:6px;padding:4px 8px;border-radius:5px;
  border:1px solid #1c3060;cursor:pointer;transition:all .14s;font-size:13.5px;color:#9aaecc;user-select:none}}
.opt:hover:not(.locked){{border-color:#D4A843;color:#D4A843;background:rgba(212,168,67,.05)}}
.opt.sel{{border-color:#D4A843;background:rgba(212,168,67,.13);color:#F5D98E;font-weight:700}}
.opt.correct-ans{{border-color:#81C784!important;background:rgba(129,199,132,.16)!important;color:#81C784!important;font-weight:700}}
.opt.wrong-ans{{border-color:#EF5350!important;background:rgba(239,83,80,.12)!important;color:#EF9A9A!important}}
.opt.locked{{cursor:default}}
.ob{{width:17px;height:17px;border-radius:50%;border:1.5px solid currentColor;
  display:flex;align-items:center;justify-content:center;font-size:12.5px;font-weight:700;flex-shrink:0}}
.opt.sel .ob,.opt.correct-ans .ob{{background:currentColor}}

/* ── answer sheet panel ── */
.sheet-panel{{background:#0d1e3a;border:1px solid rgba(212,168,67,.25);border-radius:10px;
  padding:10px;position:sticky;top:8px}}
.sh-head{{color:#D4A843;font-size:12px;font-weight:700;text-align:center;
  letter-spacing:1.5px;text-transform:uppercase;margin-bottom:8px}}
.sh-note{{color:#8899aa;font-size:12.5px;text-align:center;margin-bottom:8px;line-height:1.4}}
.sh-grid{{display:grid;grid-template-columns:repeat(2,1fr);gap:2px 4px}}
.sh-row{{display:flex;align-items:center;gap:2px;padding:1.5px 2px;border-radius:3px;
  transition:background .15s;cursor:default}}
.sh-row.current{{background:rgba(212,168,67,.08)}}
.sn{{font-size:12.5px;color:#445566;font-weight:700;width:15px;text-align:right;flex-shrink:0}}
.sb{{width:17px;height:17px;border-radius:50%;border:1.5px solid #2a3a5a;
  display:flex;align-items:center;justify-content:center;font-size:9px;color:#334455;
  cursor:pointer;transition:all .14s;flex-shrink:0;user-select:none}}
.sb:hover:not(.locked){{border-color:#D4A843;color:#D4A843;transform:scale(1.18)}}
.sb.shaded{{background:#1a1a2e;border-color:#D4A843;color:transparent!important}}
.sb.correct-b{{background:#81C784!important;border-color:#81C784!important;color:transparent!important}}
.sb.wrong-b{{background:#EF5350!important;border-color:#EF5350!important;color:transparent!important}}
.sb.key-b{{box-shadow:0 0 0 2px #81C784}}
.sb.locked{{cursor:default}}

/* score + buttons */
.sh-score{{margin-top:8px;background:rgba(212,168,67,.07);border-radius:6px;padding:6px;text-align:center}}
.sh-score-num{{color:#81C784;font-size:21px;font-weight:700;min-height:22px}}
.sh-score-lbl{{color:#D4A843;font-size:12.5px;text-transform:uppercase;letter-spacing:.5px;margin-top:1px}}
.sh-btns{{display:flex;flex-direction:column;gap:3px;margin-top:7px}}
.btn{{width:100%;padding:6px;border-radius:5px;font-size:12.5px;font-weight:700;
  cursor:pointer;font-family:inherit;border:1px solid;transition:all .15s}}
.btn-submit{{background:#8B1A1A;color:#F5D98E;border-color:#D4A843}}
.btn-submit:hover:not(:disabled){{background:#b02020}}
.btn-submit:disabled{{opacity:.3;cursor:not-allowed}}
.btn-reset{{background:transparent;color:#8899aa;border-color:#2a3a5a}}
.btn-reset:hover{{border-color:#D4A843;color:#D4A843}}

/* ── answer key section (hidden until submit) ── */
.answer-key{{display:none;margin-top:14px;border-top:1px solid #1e3060;padding-top:12px}}
.answer-key.show{{display:block}}
.ak-title{{color:#D4A843;font-size:13.5px;font-weight:700;letter-spacing:1px;
  text-transform:uppercase;margin-bottom:10px;display:flex;align-items:center;gap:6px}}
.ak-card{{background:#0d1e3a;border:1px solid #1c3060;border-radius:7px;padding:9px 12px;
  margin-bottom:6px;border-left:3px solid #2a4070}}
.ak-card.ak-correct{{border-left-color:#81C784}}
.ak-card.ak-wrong{{border-left-color:#EF5350}}
.ak-q{{color:#D8E8F8;font-size:13.5px;line-height:1.4;margin-bottom:5px}}
.ak-row{{display:flex;align-items:center;gap:8px;font-size:13px;margin-top:4px}}
.ak-badge{{display:inline-flex;align-items:center;gap:4px;padding:2px 8px;border-radius:10px;font-size:12px;font-weight:700}}
.ak-badge.user-right{{background:rgba(129,199,132,.15);color:#81C784;border:1px solid #81C78444}}
.ak-badge.user-wrong{{background:rgba(239,83,80,.12);color:#EF9A9A;border:1px solid #EF535044}}
.ak-badge.correct-lbl{{background:rgba(212,168,67,.12);color:#D4A843;border:1px solid #D4A84344}}
.ak-expl{{margin-top:5px;padding:5px 8px;background:rgba(212,168,67,.06);border-radius:4px;
  font-size:12.5px;color:#9abecc;line-height:1.5;border-left:2px solid #D4A84366}}

/* result banner */
.result-banner{{display:none;padding:8px 12px;border-radius:7px;text-align:center;
  margin-bottom:10px;font-size:13.5px;font-weight:700}}
.result-banner.show{{display:block}}
</style></head>
<body>
<div class="topbar">
  <span class="topbar-title">📝 {title} &nbsp;·&nbsp; {n} Questions</span>
  <div class="prog-wrap"><div class="prog-fill" id="pf"></div></div>
  <span id="prog-txt" style="color:#8899aa;font-size:12px;white-space:nowrap">0 / {n}</span>
</div>

<div id="result-banner" class="result-banner"></div>

<div class="main">
  <!-- Questions -->
  <div class="qlist" id="ql"></div>

  <!-- Answer sheet -->
  <div class="sheet-panel">
    <div class="sh-head">📋 Answer Sheet</div>
    <div class="sh-note">Shade one bubble per question.<br>All bubbles must be filled to submit.</div>
    <div class="sh-grid" id="sg"></div>
    <div class="sh-score">
      <div class="sh-score-num" id="score-num">—</div>
      <div class="sh-score-lbl">Score</div>
    </div>
    <div class="sh-btns">
      <button class="btn btn-submit" id="sub-btn" onclick="submitAll()" disabled>✅ Submit & Mark</button>
      <button class="btn btn-reset" onclick="resetAll()">🔄 Reset</button>
    </div>
  </div>
</div>

<!-- Answer key — rendered after submit -->
<div class="answer-key" id="ak">
  <div class="ak-title">📖 Answer Key &amp; Explanations</div>
  <div id="ak-list"></div>
</div>

<script>
const QS = {qs_json};
const N  = QS.length;
const OP = ['A','B','C','D','E'];
const sel = new Array(N).fill(null);
let submitted = false;

// ── build question cards ──────────────────────────────────────────────────
const ql = document.getElementById('ql');
QS.forEach((q, i) => {{
  const card = document.createElement('div');
  card.className = 'qcard'; card.id = 'qc' + i;
  const opts = q.o.map((o, j) => {{
    const L = OP[j] || String.fromCharCode(65+j);
    return `<div class="opt" id="op${{i}}_${{j}}" onclick="pick(${{i}},${{j}})">
              <div class="ob">${{L}}</div><span>${{o}}</span></div>`;
  }}).join('');
  card.innerHTML = `<div class="qnum">Question ${{i+1}}</div>
                    <div class="qtext">${{q.q}}</div>
                    <div class="opts">${{opts}}</div>`;
  ql.appendChild(card);
}});

// ── build answer sheet ────────────────────────────────────────────────────
const sg = document.getElementById('sg');
QS.forEach((q, i) => {{
  const row = document.createElement('div');
  row.className = 'sh-row'; row.id = 'sr' + i;
  const bubbles = q.o.map((_, j) => {{
    const L = OP[j] || String.fromCharCode(65+j);
    return `<div class="sb" id="sb${{i}}_${{j}}" onclick="pick(${{i}},${{j}})">${{L}}</div>`;
  }}).join('');
  row.innerHTML = `<span class="sn">${{i+1}}.</span>${{bubbles}}`;
  sg.appendChild(row);
}});

// ── pick ──────────────────────────────────────────────────────────────────
function pick(qi, oi) {{
  if (submitted) return;
  // Clear previous selection for this question
  if (sel[qi] !== null) {{
    document.getElementById('op' + qi + '_' + sel[qi]).classList.remove('sel');
    const old = document.getElementById('sb' + qi + '_' + sel[qi]);
    old.classList.remove('shaded');
    old.textContent = OP[sel[qi]] || String.fromCharCode(65 + sel[qi]);
  }}
  sel[qi] = oi;
  // Mark option
  document.getElementById('op' + qi + '_' + oi).classList.add('sel');
  // Shade bubble
  const b = document.getElementById('sb' + qi + '_' + oi);
  b.classList.add('shaded'); b.textContent = '';
  // Mark card as active
  document.getElementById('qc' + qi).classList.add('active');
  // Highlight sheet row
  document.querySelectorAll('.sh-row').forEach(r => r.classList.remove('current'));
  document.getElementById('sr' + qi).classList.add('current');
  // Scroll question into view if triggered from sheet
  document.getElementById('qc' + qi).scrollIntoView({{behavior:'smooth', block:'nearest'}});
  // Update progress
  const done = sel.filter(x => x !== null).length;
  document.getElementById('pf').style.width = (done / N * 100) + '%';
  document.getElementById('prog-txt').textContent = done + ' / ' + N;
  document.getElementById('sub-btn').disabled = (done < N);
}}

// ── submit ────────────────────────────────────────────────────────────────
function submitAll() {{
  if (submitted) return;
  submitted = true;
  let score = 0;
  // Lock all options and bubbles
  document.querySelectorAll('.opt').forEach(el => el.classList.add('locked'));
  document.querySelectorAll('.sb').forEach(el => el.classList.add('locked'));

  QS.forEach((q, i) => {{
    const correct = q.a;
    const user    = sel[i];
    const card    = document.getElementById('qc' + i);

    // Colour answer sheet
    if (user !== null) {{
      const ub = document.getElementById('sb' + i + '_' + user);
      ub.classList.remove('shaded');
      if (user === correct) {{ ub.classList.add('correct-b'); score++; card.className='qcard correct'; }}
      else                  {{ ub.classList.add('wrong-b');   card.className='qcard wrong'; }}
    }}
    // Always mark the correct bubble with a green ring
    if (correct !== null && correct !== undefined) {{
      document.getElementById('sb' + i + '_' + correct).classList.add('key-b');
      // Colour options
      document.getElementById('op' + i + '_' + correct).classList.add('correct-ans');
      if (user !== null && user !== correct)
        document.getElementById('op' + i + '_' + user).classList.add('wrong-ans');
    }}
  }});

  // Score display
  const pct = Math.round(score / N * 100);
  document.getElementById('score-num').textContent = score + '/' + N + ' (' + pct + '%)';

  // Result banner
  const rb = document.getElementById('result-banner');
  if      (pct >= 75) {{ rb.style.cssText='display:block;background:rgba(129,199,132,.14);border:1px solid #81C784'; rb.innerHTML='<span style="color:#81C784">🎉 ' + pct + '% — Excellent work!</span>'; }}
  else if (pct >= 50) {{ rb.style.cssText='display:block;background:rgba(212,168,67,.11);border:1px solid #D4A843'; rb.innerHTML='<span style="color:#D4A843">📚 ' + pct + '% — Good effort — review the key below</span>'; }}
  else                {{ rb.style.cssText='display:block;background:rgba(139,26,26,.18);border:1px solid #EF5350';  rb.innerHTML='<span style="color:#EF9A9A">💪 ' + pct + '% — Keep practising — answers and tips are below</span>'; }}
  rb.classList.add('show');

  // ── build answer key section at bottom ───────────────────────────────
  const akList = document.getElementById('ak-list');
  QS.forEach((q, i) => {{
    const correct = q.a;
    const user    = sel[i];
    const isRight = (user === correct);
    const correctLetter = correct !== null ? (OP[correct] || String.fromCharCode(65+correct)) : '?';
    const userLetter    = user   !== null  ? (OP[user]    || String.fromCharCode(65+user))    : '—';
    const userOpt  = user    !== null ? q.o[user]    : '(not answered)';
    const corrOpt  = correct !== null ? q.o[correct] : '?';

    const card = document.createElement('div');
    card.className = 'ak-card ' + (isRight ? 'ak-correct' : 'ak-wrong');

    let html = `<div class="ak-q"><strong style="color:#D4A843">Q${{i+1}}.</strong> ${{q.q}}</div>
      <div class="ak-row">
        <span class="ak-badge ${{isRight ? 'user-right' : 'user-wrong'}}">${{isRight ? '✅' : '❌'}} Your answer: ${{userLetter}}) ${{userOpt}}</span>`;
    if (!isRight && correct !== null)
      html += `<span class="ak-badge correct-lbl">✓ Correct: ${{correctLetter}}) ${{corrOpt}}</span>`;
    html += `</div>`;
    if (q.e) html += `<div class="ak-expl">📖 ${{q.e}}</div>`;
    card.innerHTML = html;
    akList.appendChild(card);
  }});

  document.getElementById('ak').classList.add('show');
  document.getElementById('ak').scrollIntoView({{behavior:'smooth', block:'start'}});
}}

// ── reset ─────────────────────────────────────────────────────────────────
function resetAll() {{
  submitted = false;
  sel.fill(null);
  document.getElementById('pf').style.width = '0%';
  document.getElementById('prog-txt').textContent = '0 / ' + N;
  document.getElementById('sub-btn').disabled = true;
  document.getElementById('score-num').textContent = '—';
  document.getElementById('result-banner').className = 'result-banner';
  document.getElementById('result-banner').style.cssText = '';
  document.getElementById('ak').className = 'answer-key';
  document.getElementById('ak-list').innerHTML = '';

  QS.forEach((q, i) => {{
    document.getElementById('qc' + i).className = 'qcard';
    q.o.forEach((_, j) => {{
      document.getElementById('op' + i + '_' + j).className = 'opt';
      const sb = document.getElementById('sb' + i + '_' + j);
      sb.className = 'sb';
      sb.textContent = OP[j] || String.fromCharCode(65+j);
    }});
    document.getElementById('sr' + i).classList.remove('current');
  }});
}}
</script>
</body></html>"""

def parse_mcq_for_sheet(text):
    """
    Parse AI-generated MCQ text into list of {q, o, a} dicts.
    Handles formats:
      - Standard:  1. Question\nA) ...\nAnswer: B
      - Titled:    1. Subtitle\nActual question text\nA) ...\nAnswer: B
      - Sectioned: ALGEBRA (4 Questions)\n1. ...
    """
    import re
    questions = []
    lines = text.split('\n')
    blocks = []
    current = []
    for line in lines:
        stripped = line.strip()
        # Skip dividers and blanks (but keep blank as block separator)
        if stripped in ('---', '***', ''):
            if current:
                current.append('')
            continue
        # New block: line starts with a question number
        if re.match(r'^\d{1,2}[\.\.\)]\s+\S', stripped):
            if current:
                blocks.append('\n'.join(current))
            current = [stripped]
        else:
            current.append(stripped)
    if current:
        blocks.append('\n'.join(current))

    for block in blocks:
        blines = [l.strip() for l in block.split('\n') if l.strip()]
        if not blines:
            continue
        # Skip blocks that are section headers (ALL CAPS, no options follow pattern)
        if blines[0].isupper() and len(blines) <= 2:
            continue

        q_lines, opts, answer_idx = [], [], None
        i = 0
        # Strip leading question number from first line
        first = re.sub(r'^\d{1,2}[\.\.\)]\s*', '', blines[0]).strip()
        i = 1

        # Decide if first line is a subtitle (short, no maths symbols, no '?')
        next_is_option = (i < len(blines) and bool(re.match(r'^[A-Ea-e][\)\.]\s', blines[i])))
        is_subtitle = (
            not next_is_option
            and len(first) < 70
            and '?' not in first
            and not re.search(r'[=\+\-\*\/÷×]', first)
            and not re.search(r'\d', first)
        )

        if is_subtitle:
            # Skip subtitle line, collect actual question from subsequent lines until options
            while i < len(blines) and not re.match(r'^[A-Ea-e][\)\.]\s', blines[i]) and not re.match(r'^[Aa]nswer', blines[i]):
                q_lines.append(blines[i])
                i += 1
        else:
            q_lines.append(first)
            while i < len(blines) and not re.match(r'^[A-Ea-e][\)\.]\s', blines[i]) and not re.match(r'^[Aa]nswer', blines[i]):
                q_lines.append(blines[i])
                i += 1

        # Collect options and answer
        while i < len(blines):
            line = blines[i]
            # Answer line
            if re.match(r'^[Aa]nswer', line):
                m = re.search(r'[Aa]nswer\s*[:\.]?\s*\(?([A-Ea-e])\)?', line)
                if m:
                    answer_idx = ord(m.group(1).upper()) - ord('A')
                i += 1
                continue
            # Option line: A) or A. or (A)
            m = re.match(r'^([A-Ea-e])[\)\.]\s*(.*)', line)
            if m:
                opts.append(m.group(2).strip())
                i += 1
                continue
            m = re.match(r'^\(([A-Ea-e])\)\s*(.*)', line)
            if m:
                opts.append(m.group(2).strip())
                i += 1
                continue
            i += 1

        q_text = ' '.join(q_lines).strip()
        if q_text and len(opts) >= 2:
            questions.append({'q': q_text, 'o': opts[:5], 'a': answer_idx})

    return questions


def generate_result_docx(text, task, topic, grade, subject, school="", teacher=""):
    """Generate a polished Word .docx for any AI result."""
    try:
        from docx import Document as D
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from lxml import etree
        import io, re
        doc = D()
        style = doc.styles["Normal"]; style.font.name = "Calibri"; style.font.size = Pt(11)
        # Title
        h = doc.add_heading(f"{task}: {topic}", level=1)
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in h.runs: r.font.color.rgb = RGBColor(0x8B,0x1A,0x1A)
        # Meta
        meta_parts = [p for p in [school, teacher, subject, grade] if p]
        if meta_parts:
            mp = doc.add_paragraph(" | ".join(meta_parts))
            mp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for r in mp.runs: r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x66,0x66,0x66)
        # Divider
        bp = doc.add_paragraph()
        bp.paragraph_format.space_after = Pt(6)
        pPr = bp._p.get_or_add_pPr()
        pBdr = etree.SubElement(pPr, qn("w:pBdr"))
        bot = etree.SubElement(pBdr, qn("w:bottom"))
        bot.set(qn("w:val"),"single"); bot.set(qn("w:sz"),"6")
        bot.set(qn("w:color"),"D4A843"); bot.set(qn("w:space"),"1")
        # Body
        for line in text.split("\n"):
            s = line.strip()
            if not s: continue
            if re.match(r"^#{1,3}\s", s):
                lvl = min(s.count("#"), 3)
                doc.add_heading(s.lstrip("#").strip(), level=lvl)
            elif s.startswith(("- ","• ","* ")):
                doc.add_paragraph(s[2:].strip(), style="List Bullet")
            elif re.match(r"^\d+[.):] ", s):
                doc.add_paragraph(re.sub(r"^\d+[.): ]+","",s).strip(), style="List Number")
            else:
                p = doc.add_paragraph(s)
                p.paragraph_format.space_after = Pt(2)
        # Footer
        doc.add_paragraph("")
        fp = doc.add_paragraph("Generated by Teacher Pehpeh — Institute of Basic Technology (IBT)")
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in fp.runs: r.font.size = Pt(8); r.font.italic = True; r.font.color.rgb = RGBColor(0x99,0x99,0x99)
        buf = io.BytesIO(); doc.save(buf); buf.seek(0); return buf.getvalue()
    except Exception: return None

def generate_result_xlsx(text, task, topic, grade, subject):
    """Generate an Excel workbook — best for quizzes and MCQs."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        import io, re
        wb = openpyxl.Workbook(); ws = wb.active
        ws.title = topic[:30] if topic else "Quiz"
        # Header row
        hdr_fill = PatternFill("solid", fgColor="8B1A1A")
        gold_fill = PatternFill("solid", fgColor="D4A843")
        navy_fill = PatternFill("solid", fgColor="0F2247")
        thin = Side(style="thin", color="CCCCCC")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)
        ws.merge_cells("A1:D1")
        ws["A1"] = f"Teacher Pehpeh — {task}: {topic}"
        ws["A1"].font = Font(bold=True, color="FFFFFF", size=13)
        ws["A1"].fill = hdr_fill
        ws["A1"].alignment = Alignment(horizontal="center")
        ws.merge_cells("A2:D2")
        ws["A2"] = f"{subject} | {grade}"
        ws["A2"].font = Font(bold=True, color="D4A843", size=10)
        ws["A2"].fill = navy_fill
        ws["A2"].alignment = Alignment(horizontal="center")
        ws.row_dimensions[1].height = 22; ws.row_dimensions[2].height = 16
        # Parse questions from text
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        q_num = 0; row = 4
        ws["A3"] = "#"; ws["B3"] = "Question / Content"; ws["C3"] = "Answer / Key"; ws["D3"] = "Notes"
        for cell in [ws["A3"],ws["B3"],ws["C3"],ws["D3"]]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="1A2744")
            cell.alignment = Alignment(horizontal="center")
            cell.border = border
        ws.column_dimensions["A"].width = 5
        ws.column_dimensions["B"].width = 55
        ws.column_dimensions["C"].width = 25
        ws.column_dimensions["D"].width = 25
        for line in lines:
            # Detect question lines
            m = re.match(r"^(\d+)[.):] (.+)", line)
            if m:
                q_num += 1
                ws.cell(row,1,q_num).alignment = Alignment(horizontal="center")
                ws.cell(row,2,m.group(2)).alignment = Alignment(wrap_text=True)
                ws.row_dimensions[row].height = 28
                fill = PatternFill("solid", fgColor="1C2340") if q_num%2==0 else PatternFill("solid", fgColor="161D35")
                for c in range(1,5):
                    ws.cell(row,c).fill = fill
                    ws.cell(row,c).border = border
                    ws.cell(row,c).font = Font(color="D0D8E8", size=10)
                row += 1
            elif line.lower().startswith("answer:") or line.lower().startswith("ans:"):
                # Write answer into previous row col C
                if row > 4:
                    ws.cell(row-1, 3, line.split(":",1)[-1].strip()).font = Font(color="81C784", bold=True, size=10)
            elif line.startswith(("A)","B)","C)","D)","a)","b)","c)","d)")):
                # MCQ option — append to col B notes
                pass  # already embedded in question usually
        # IBT footer
        ws.cell(row+1, 1, "Teacher Pehpeh by IBT — institutebasictechnology.org").font = Font(italic=True, color="888888", size=8)
        ws.merge_cells(f"A{row+1}:D{row+1}")
        buf = io.BytesIO(); wb.save(buf); buf.seek(0); return buf.getvalue()
    except Exception: return None

def generate_result_pptx(text, task, topic, grade, subject):
    """Generate a PowerPoint — best for lesson plans and schemes."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io, re
        prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
        C_NAVY = RGBColor(0x0F,0x22,0x47); C_RED = RGBColor(0x8B,0x1A,0x1A)
        C_GOLD = RGBColor(0xD4,0xA8,0x43); C_WHITE = RGBColor(0xFF,0xFF,0xFF)
        C_LIGHT = RGBColor(0xD0,0xD8,0xE8)
        blank = prs.slide_layouts[6]  # blank layout
        def add_slide(title_text, body_lines):
            sl = prs.slides.add_slide(blank)
            bg = sl.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = C_NAVY
            # Title bar
            tb = sl.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, Inches(1.0))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = title_text
            run.font.bold = True; run.font.size = Pt(24); run.font.color.rgb = C_GOLD
            tb.fill.solid(); tb.fill.fore_color.rgb = C_RED
            # Body
            bx = sl.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(12.5), Inches(6.0))
            btf = bx.text_frame; btf.word_wrap = True
            first = True
            for line in body_lines:
                if not line.strip(): continue
                if first:
                    bp = btf.paragraphs[0]; first = False
                else:
                    bp = btf.add_paragraph()
                bp.space_before = Pt(2); bp.space_after = Pt(2)
                run = bp.add_run(); run.text = line.strip("- •*").strip()
                if re.match(r"^#{1,3} ", line) or line.isupper():
                    run.font.bold = True; run.font.size = Pt(16); run.font.color.rgb = C_GOLD
                else:
                    run.font.size = Pt(13); run.font.color.rgb = C_LIGHT
        # Title slide
        sl0 = prs.slides.add_slide(blank)
        sl0.background.fill.solid(); sl0.background.fill.fore_color.rgb = C_RED
        ttb = sl0.shapes.add_textbox(Inches(1),Inches(2),Inches(11),Inches(2))
        tf = ttb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"{task}: {topic}"; r.font.bold=True; r.font.size=Pt(36); r.font.color.rgb=C_WHITE
        stb = sl0.shapes.add_textbox(Inches(1),Inches(4),Inches(11),Inches(1))
        sf = stb.text_frame; sp = sf.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run(); sr.text = f"{subject}  |  {grade}"; sr.font.size=Pt(18); sr.font.color.rgb=C_GOLD
        # Split content into sections and create slides
        sections = re.split(r"(?=^#{1,3} )", text, flags=re.MULTILINE)
        for sec in sections:
            lines = [l for l in sec.strip().split("\n") if l.strip()]
            if not lines: continue
            title_line = lines[0].lstrip("#").strip()
            body = lines[1:] if len(lines)>1 else [""]
            # Max ~12 lines per slide
            for i in range(0, max(1,len(body)), 12):
                chunk = body[i:i+12]
                add_slide(title_line, chunk)
        # Final slide
        sl_end = prs.slides.add_slide(blank)
        sl_end.background.fill.solid(); sl_end.background.fill.fore_color.rgb = C_NAVY
        etb = sl_end.shapes.add_textbox(Inches(2),Inches(3),Inches(9),Inches(1.5))
        ef = etb.text_frame; ep = ef.paragraphs[0]; ep.alignment = PP_ALIGN.CENTER
        er = ep.add_run(); er.text = "Teacher Pehpeh by IBT"; er.font.size=Pt(24); er.font.bold=True; er.font.color.rgb=C_GOLD
        buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf.getvalue()
    except Exception: return None

def synth(sp,q,resps):
    v={k:v for k,v in resps.items() if v and not str(v).startswith("⚠️")}
    if not v: return "⚠️ No models responded."
    if len(v)==1: return list(v.values())[0]
    p=f"Combine into ONE print-ready document:\n{sp}\n{q}\n\n"+"".join(f"=== {k} ===\n{r}\n\n" for k,r in v.items())
    for fn in [ask_cl,ask_gpt]:
        r=fn("Expert editor.",p)
        if r and not str(r).startswith("⚠️"): return r
    return max(v.values(),key=len)

# === MAIN ===
def main():
    # Load Teacher Pehpeh logo as PIL Image for taskbar/tab icon
    try:
        import base64 as _b64, io as _io
        from PIL import Image as _PILImg
        _icon_b64 = "iVBORw0KGgoAAAANSUhEUgAAAEAAAABACAYAAACqaXHeAAAf90lEQVR4nIWbeZwd11Xnv/fW9tZ+/XpfpJa6tVqbrUi2vMZyErKYOIkT42yACYTNGQKffAgDBAaYYQsfBoYwQCZDPoEQCCEkdnC8xHZiW4liObJk7busllrqfX39tnq13PmjllevJTOlT+u9V3Xr1ll/59xzbone3l7FGx0CUCCECH4rhUpeABACIcIzK2cSovVkfFs4H+o/Hd4cI1aME+EV1frcaAIFCBXeKVAquBSwoBLzgBRvMHmSPqVU8BdeFyKgSUSk3UCEIsFNNJ6QEPEGNymVuLc50w3GBfSgVsyiVHCvaD0vRCCsFvrC/2V8ISRKKYVAhIxe//CIQEHMCUKIhBauJ10p4jmv41txI7XH90opW+i4jqIVt0YaFgiEEqGSgueqSECACieSEQOa1IjUJKSIJ4sEEf1F0hcJ5jVNxoStsKfrBSlEYLrROUHTPhNMqMT35jDRMnalgmTiWYHgmh4XXYueHSgkFEAr1QJf+fHApvmrloeq8J+UommKkVSlbLln5SGFREp5w2sqnDswsOj5fnOmxJRqBb74NJ/p+35oXKG7RMS1kCTQg4+mNqOHIprokGRESoEfa0WiAM/3QkxI3J8gMqntYIzED+8JxoCMQAUV2KdqgowQMjbhiE6FCl2rKRiFQgqJEhEQhnilwMePNZ+kT0bS8n0/Fo8QAhlaAoSYELqH7weTCyWQAnzPj9QWPNz3W5kOiYv8P5pHEJilJPiLCRICKSVStCol+p7UXgRk8RklYnCMJCsSUSwpwOjQEzMGWKQUUpMBkSoxeUi2ItQWCn+FhXu+1+LOItSqEBFWaHiehxfSqHwwpEQQmKsf3qUF3ARCRaD8JB5FtEUajmMRylcofCLsiawC5Qe8iOuFoMdSRsS+p1TzM+nzTf2tkKQieTYOj0KT+Mqn4fn4jovvOGRzKTo0iSElOjBbb9DQAgBOC/Crdcqeh0xZ6IYWsCNCQTQfhpABk57vNV1PBO4jEtpu0hUoNCBZxNmFLkKTC1ygCTxvgF8xNqgwWgRSj0BLxlO7nkejbpNNW6wv5hhZ3c327cP8YN9RHn74LazbPISha3z9y8/wT88fwQP27t7A+x68m7Nnx3jhhdc4NbGA13CoItBNA0PKOBlTKJQf+HxgvU36hRABxoTqEIASKkycRIBjfgCOMkb5FaFNChmDTfCABCRGSYjX9HcpBLqUuJ5PuVylty3Dh9/2JgbSBp/+tYf43T/6ODt3bSSlSXrXDVIt1Wi4LrfufROe4yFcj1t2b2LdtrXcdfcOPvnpj3DbSB9/+lsf5T23bqTX0nE8H8dXaJoM3DAERl/5ofuGv30/xq2VKYZSCs8L8S4Kg5G5RGFDCBHkAqoJgJFZxQlNaA3Rh69gfqnMYFuWR99/NzsGOvj0b/8Ut2wb4dKVKYTUWFqqslyqsDSzxOrNqxhc28exV04iNEnG0hlZP4hTd7DSJspx2XvfTm77sd38wqMf4MPvvoORnMlIe5q64+L6Ks4oV2azSfOVMuGmqnWEiMJg7FcrJNUEmWbqGINLGOeFALvh0p8y2btnM37d5hO//iGefvwH1GyHd73nLr739CsUilk6Ott4+OG34ZUqHHzhCPWqzdBQP3atTt7I0d1Z4PVTl7k8PsfoyUt85NEHmZua5+ih07z0wmH+x5/8MpdOnOdL//AMJSTzrh9bQoxSYVaqlGomPzfINCMQCAUgmkKIzF4lB0bMBg/zlY/QArN3XI/1xSy/85sfYd2mIX710b/kqSd+yHs+9Bae/PcXeft77uTIS0d4/lv7qdYbXJ2cYWZqDj1jgfJJ59M88MBuauU6//avz3LyzDWePXSOn3jLLvoHuzh7cpRL56+wftswt96zjRef2s9n/9cn2ffcQf7nP36HVD6L1HXwabGEgA/wQ5BsXogWSKHbt6wGRRQRwkitVIi24HtJiQChvzvVOr/z8/fz/kfewcMf+D0mlyo8dO/NfPoPPsZv/crn6BnownEb+IZicDBPJiXpaM9AvUy5Wse1G+SyaZAaVxdsrGwapaeZvjyPtBW779jBhaMXefdPv535mUUy+RyF9iy/9PE/47Zdm5gbm+KbL58i1ZZDxFgmY2wI7Ta0guSyMHCPWABRJBBSNCUmgtiajKkBooJwPW4b6ePY6xPcctMQf/v3v8Ff/Nm/ctNNQ2Ql/OClY+T6C1gpGF6Vw6gvMD81R31pkbQBQirqDajXbDRdIpTD6LxGX3c7ncUUmWIRJ13g0qVFlmaqrFu3lrPHL/Lf/vrX+NwffZmJyQX+5K8+ycy1af76s//CU0deR+o6EnB9L4gQUsYWm9R8ZP5CiFYLEAKk1JqhDa5jHilwbIf33rqJz/zxxzl//ipf+Nw32Hv3Dj748+/mle8f45lnXqZnsI2RXgN76irTYxNIKenIGNSlTs0HK59BaJKq7eE0XGyZodKl4wuX1IVJNuo2fi6PXuiA7iHOnp+nvOiy/aZhzp65wm989pd55ekfousa7YM9PPzQ70EhDzJQWly7SCysboQHsQDecOmbzNERSKC8VOb9e3fwW//9Y1y4MMEd99zCd772PF//2ncZ3DLAji29qGvnacxPUccilzGZ9C0WfMlk3We+oag4Hh7g+B7tWo477uxidPEyM+MO660864RLHpesXWKiJulZv4ZG+xBPPn6Iu+/cxbaNq9CzFuu2j/C3f/hP5IttHD16nmdPXCaXTeN7XgzkCIEK85zr+Ovt7VVR8iCFwAtXUXGCGS0vNUm1ZnP72l5u37OJv/vHZ/ncn/8yN9+2hWcf+z6jVyZoH8yRc5eoXjpPsT2FbllMyQzHFlwul21KDQchFZZlYGoalq5Tq3vcflM/E3PjXBgVyLSiUfbozJuszjhstHwKbpWM77BAhtz6zVybcvnG11+ht6/IjrX93PXWXex5521cPHKOP/j9f+DcYhVNgB+G7ahS1LI6jaCs+bs5IBk6Ixl4vqJH1/iZj72DW+/Yyjvv3cG3H/s+yrb5xn98n0xvim5nGm3iAvmMiZ0psK9q8uXTMxyZXgQp2Tw4xJ6NO9nWv5GBVA+VqotQEqnXmV0S6BnYOtDDvTt68aUkP3w7zqa3kLv5zdRy3SzPz+OOnmHjap37P3IHU6Ua7/7o27n1rbs4cfAUCsG73raL6uJyWOuKfF41XXhFPiD6+vpUFDtb6mZKoUS4dtc0qqUKP3v/rfzq7z3C0UPnWJxeQJOSA/uPM7S1D7Ewhjt6ge41/Ywpk6dGSzRUive++a3ommTn1ltYOzDEoRMvcXbsR5hZm9pSmn2vjDEyBJdmGiyXBX/8yEbwK3x1n8lf/M4fUq3WmF0qI5bHObv/m0yePUjWqcParczYeYbXrKN/dTcXTl/mqcf2YZgGRtrkheOjCKkl8pUwgfcj0w4+9RgUkgsb1cypBQLP9djcV6Q6v8j+p1/h3nffwdWxGc6fuESmO03encOqXmOuu4eTNY3nx8vkMwU+/cGf47Zdu/F8F00z+PK/fIXzUy/hegLPcli7qo31azTGpxeRfpoPv3UjOVXjn5+eYvimexgbmySVTpExBfX8ILc/9EkOPPZ3XDu2n+7ZUbq6hnj11SPcm7mDc69dYHjLML/0ax/k9MvH2f+jz1PLZtGjxU/MJ7FLCEQTA5JA6Cu/mfJKgbO4zJ/85ofRLYNf/43P818efR/33ruTbz25j+F+gX/5HGZ7gVFl8Ny4S7W2yE/c9z5+5ZGPMz49QcrMcPzkUY5eeh4vs4jvglnvYH5xloX5OnrBo0vm2V7MMo7PeLWMYQqymQzru/ew46a7aS+kWC5X6Sq0c+jJv+H8gecRCLpv2c2+717l4UceYM99Ozlx+AxPfPk72Erx9RePoHQ9sGbV1DpRjiNEKwgmaz9BHBU4ns+23nY+/6XfpFqp8qUvfJvnXjzCQG+BD3zgFkrHD5HLmkxbOZ64aqNJ2LluB488+EE6Oop4ruLkyRMcu/YtFux53IYklXNYLjWoVVPUZrMM7a5x4rUZqBQY2KwwcRF1k1xRQ9XT9Frb+eB7f4aXXzvAodMn+P1PPMo3//YzVMZOkMq34fdsZGLMpbZc44O/+F5Wrenj1ed/xK//7heZR6KpFvhrUbbo6+uLr0UlpWA56SOkxK3U+NRH7mNk4yrmphb5yU+8h6/8/dOcvjbGajVLB2UWcx18+1qDmZrNHZt28plHP0XdqQNw6dJlfnTmG1ycOUWlZtOou2g5E2Myz5rZOloHnPUa9C1LCprGZFkx3euTGxQMDuSYn6+R0lKsbbuTB9/5U/zVP/xvZpfL3LdpDaXjzyBrNfo3r+fIOYf73vFWNt80xD9/4XEef/Yg0w0/Lt8l0V8IgdCCpXyLAAIklMhwnR9UfRSDGQvHcZldqvDmHcMMrO1maMiAsbMY+RwHKgbHS4qCafKZj3+KtWtWoxQcPnyCb3/vq2S7HDJmN2mjQKXWoHFtgV/Yvgbv9FmWXp+kLiQGCk/BRMnBKWY5vH6AY+OXuHXrVjaOrOfatWnefs8DdHUVeew732bPztupXniRiz/8Nr4v6LrpFs6eqfLa0YscvjSJzKTRW4o5Cc2H4O77XrAYimsBIlH3S9xwpWIjBBjtOZ547RI/t7mHxswEWU0x7eucr0hM3WHD6vWsHx6mUquiScnWrRsYGfmveB4Ucnk6O9p55dAR9FeeQMxNMHV1GmmALhWOAksKcp7G4swstw9vgC330F1s52Pvf5SJmUlsp44mJY+8/0MITXJF1jl58CXMhVm08izn5yp8/+IUXR15PMdtWcAFnp8onckAF2RcvhZNhltSBqUwNIEhBW7dYWh1J2t6TAqWTcNIc2oZlJSYhkHZrlKvB0QqQNM1CoU8HR15lHBYrCwxcfUCnY0Si4fOoRouCIHv+WQFuL7CMgRTSmfyzEnefuvt7Du0n8eeexyEh1I+juNguw1qlQp/+dV/Y96wMDM5Ji6Ns+fmAdpyKTzHbVacW4J+82tUAZORZFSsddV0B0BqWvBVCOoNlw1rO3CrSyzMVXAMk6t1hWVqeAqMsLkS+50Cz/XwXC8kQGJJyTOnpjju6XS0Z9B1nUJngW8uSx6ftvGEoFfXEIsluts7+NC7H+Ivvvh57IZDLpOlkG/jythVfv9v/pzzE68zVpPYhoHpVulJuawaKGI33BX1ynBFq/zr1gJ6M12OqqZNGQShIloKBKur1f3Byq7mSa6UFQ0lSWuSuu3R095Ne6FAqVIKmRbN7owCTQg8BbVKhaerHjUjS7Hu8tTUEr6hcVtPjnqjQT1tUelu4+lvfoWPfeJTzMzO8fKhg6xft5knvvskL766D8MyaM/nmKg2ED1pjMoifmWZDWuLnL84RTpSbFxBDF0gWthFKX5Y3gkyJl+1SEhoMpYgSNJpk76iSUp6dHRmmawrRLjkNE2TM1eu8OqhV5mcnkMzLGTT/QBBo9Hg+OFXGd69haEdmznX8Dhbq7HxzbexYc8tzLe1cVATTKRNMg2HbKGX5fIyP/3gT3D03Ou89n/+lOyF02zfsgtLN9A1iYNkzpVUlUF5fpFVPZmwrN8S1FtSYBl1SyIXaCq8FfyUHxQQhRQ4rkt7WwpTemjKQzMEywTlbR+FJSy2iAov/uNnOfC1v+K1l1/EFzJoaymFaRqkLItUNo02OcfeTWvZuLqHtZpPZ22evVuGWVutY1XqDM+VEBWbe975LjqLRaaWl9mpNVj9+ii3730fm4dGqHsNdEMDIZmxFbYCqRx62k0s02gB8ibAAUKFLTSaAmgy3fwerQcCt/BxPZ98zsQUHmPTNWbrAldo6Iak7iq2iQZLlQUOpbs4tgSHn/5nzpw4QrG9A0PTuXT5CjPz86xeM8J8ZZGrzz5Hpr7MxsEuFl87xo++/FWWvSoDuQw+CjefQ6HQTIszJ45R+NFLnOtazWJbmpcO7iOXz+OH2F3zFIU2E0t6ZFM6pqnh0+rrzWQvWeZd0RlqLYsStpDDfpsPactA+B4GHrYn8MKuRcYTNCplzrW1k25LMWb71Ms66f3f48DFi1SXFzlz9QJKmrxj7UbuKhaZnb+GM1PisuNxX1sOx/MxpGTZh2uOS6qrl4FVqzhz+hTO/qe5ZuvsNCsc+dqXsFMmmqkBXgC0CExDp1yuksHHNDQadgOpNcviSvlhV6lZARVSNNvjNxSCagKHEAJdE1g6tGVNDE0GaK8JNOWi6RIpNdyZGqmqSzWXZXHqGvuPH+D0xAVSmRS4da5Mvk5DS5MSPq6mkerv4oyvUfLB8DxKyw2kAKe9nVqtxuXjx7BOnKErn+G16QrrylX6hcQJ/ToIUAq34eHZLkr5QSKnYg3GfMXl/+jTV83WWHQhKotB2GIOw4RSCs/3sRseolZHT2UQvoko+dQzJrpw6ViqceLsPCkpSPfmsYs+HeRYmKwxs7jAUKcgrducLlncYpk4SlGeWWS9KVFSYiuB7/qY2TTfOnCEL/3w59m6tpsHU2mqV2YoS51nfMGMBEsJPKnwlI+JQvk+ZlYHIXA9Pwa666rcCR0LEfYFfN9vsQJf+bHWA0NQSE3geqBZJjUrhSYEasphZmaZfG+Wg76HW63gSQPba1CbWsbtakPNVJi9ukAxa7I457P7VofZdIkTJclmKUlZQbFFk4JKw8f2FWO6Rg3FnZ0aH0hVubxQ57vd7ahMirF5B3ViGikhV0xD2kf3fdKWj5VLM2kr7LoT4JkUcTMnWfBJlsljDIgaoslwEVkEBOWy+cUqVjpDAx9TKlLKw/aB6TIKRSFjUDR95ss+qYJF1YPVHRbtWoGhzgw12+H4nMdMqo/i+iXUyWnWmgbSCPID2/ExUhprUpLf7rSQjsuLl2b5nq2R39RF7fw8Ym4ZzZCU6g7lhSrp9gw9O1OUqj69nRmWKw412yWT0sKETISlseuLokqpBAiuqIlG6wMVZjG6Jlkq2dQ8DZlKkZKKzpzkwqxBxnLZMdDGYCFF2XY5PrHM8mCB+eUaazokTiqDp0AIjVfP+WzbWmFbcR5rj8bEhKR41aYrozHtKaoNjzNmG7X5CrKjSGakk8GxKezlOo1ynZuH2ylmTZaqDmemyszXG2R9i4xwQZqMz9phEqYHfYKwFKZpEqWa2WAiEWpKILnloGVvEME+oHLFZnzOJp+x8Oo1NgxnUdJldZtFe0rHbjgoz2EwI+kxJQuuz2zZoXvdCMulMjVH0duWZrFhUPHgxNUUiz15nh/eyvfKKZZLdWpCxy5YpLvy2K5LZblGsS2Dcn36czqWBHxFShcM5C162nW6TYdaA8i1MTq+hK5J4u1RMb41wU+E+4daXCCKF0IIpCbxos4vQaNUEXRdz16aZ8ueHipzcwz0uqzqt/AbgZnVXBcjl2cwA1OOzaZijvWuQ3Z6lNzue1CLM3ScP8OpZQ+xaQ1tV0YZPenylvcIzl8VLG3IMWalmZmpcHW6RqnuoAG5QoaCEgzmTYRmUKrUQQsWXMNdBl6lTHdnhmVP5/XLC1im3rpHSdDc6xS6Q7R9piURity/ZXuACIqjCIFlaZy+MIufylNuQG2hzFs3ZtGVR6Naov+mm3nbjz/IeNXmmmowVEjR4diUTo+x5txJto2PMqV0VjmSr+1rw1szwkd3LXH01Rku76gwcrfF5ltNri3U6M3q3NyfZU1nilqpzOJinclqjY07d/O2jz6CbqboTPus7zOo1R1ERxcXJ23m5stILdjeEzd1ohQ/qgfSFEZLGGwBxDD8SakFS0cBlqEzMV3i1GiJ9cOrcSauMVBocGZDAcdOsTxxhe9cvsgdb7mTHq/K4f2H6a05jHcUUOk0+cEhZL6N7ccPIUaPkr+rwFTZQOrdqKyGMTGOMl2yKYutPSlSpsZi1SZjSI6Pl9n85rsZu3yOS+dO4WuStsEieW8K0hputsD+5y6ih+sAKWVzM1dc85RoUjTDO+EOkZaKaaIqDOG+n1AeSgp0TfLCgSvs+dhOLo1N4F2e4AN33cVtD/wUxw4fp1r3yKV1eqpprrXnebLh8qZVHUxiMnplEq9ymb1Fkzs6BAeemmdmWLBnpELdSPHc7O30Thxm22CRumtQXS7RtX4L9XNnkWqZkVVdjKXz4Do8/OYNvPr8f1A6P8eares4Owdnzs2QyQTrACEEic1iAf2ejy8EiGaZPMaAuFqysjgShhLCTks2bXLh0hyvnltg3fBqvNfPs3DmJMvbj3Lw5WO8fmkC2wPbk3SlFG2Ly0xr4OQyLNop1vk18ppkVsEO3eLVw1WW8jn0NXmUM4nobUeOLZPrWo3Z2U5paoKsdNncm+WJJ39IKmWypjfPuVeuwbVLtLXnsQs9PPm10/Gyvrl/QwJBCiwSOEDcA0mmwqKJ/NH3qNUc7Q8ILEJhmpLHnj2DKPSgtxWoz87xwuOPk3LmWPCyaNkO+odHmF0oM7hmNZqUZAyF26ixPe1SXygjgIbyWJWSHDtQ4eoFl+1vKrHUOUgq65MtdjM4sh63ssSEY6IXuzE7+hH5HixZYez0KWoN6N22iQOnFzl9YYps1ozBW8qAbsSKhqhIbr9VCRBUwYIBQbhtNhwvW0MiQCplMDlT5snvj5HbcBO2YSKrCwy0e2zsN7BdGB+7ipVJ88BPPkTvQB9uZY5CZY6bMjq15Sp1BQvLDpOaRSVTYbd3kWfPGRhj42zbMsDCzDhH9z2H07OGnpv38LM/ez+OXWdD2zy92jxtwqE4soqrNYtvfecMuayFF5m+EES7SyPL9nwv3i/Q3BcpEgKIGwZhvSyxYEguIYUU+D4Ucmle3H+R/edKrLvzVnTpMXFhlJHMLDt6auzdPUw6JfjBcz9g1+6dCFz2dmeQ8yUqjs9sqYGbS/EfaR2KBjdv0nhvrsKCofjK90qMl20abd2U56ZY0+bT31XkbVtM6lOXSZVmaRvohZ4h/u+/HqHhuGhhYSZewodKjeBdCoEmZZOPSKl9fX0q3j1Bs0usSFSLE2WkYAnZrCDX6i6/+NFdbO/TOH/oKKZUZDq7uPud7+QL3z3FiwePM9Q3yLqZOR5qM5gv29h1j3nL5BumxmJBp6vqsM0V9BZtUhuzvHRyiHwWFssN5Mw4m9d30d2ucGcv05nyqBlttK/fzOf/7Tinzk2Ry1l4brQJIsCqyIojS5BCIKTE9/w4RMYCiErhUsgg6YlTSJHAAnGdKwgBvq+wbZefefgW7t6U4/KJUyzNLCDNDG0DvRwsubx+ZpHPZCzchst4zeOSofNS1sDJalhAw1MsTaXpMcr0FTI4Vj+a8DBw2ZhdpCfXYGl6mmzOZPXafkrpfr742ClOnJmgLZ/Gdf0EAAo8z4vBMErlI1tQ4bafuFja39/fzHlEk/nYfGJIDVBVRu4iEltplKJSd3jgxzZx/229lK5eoT5+FUdpTKo21k8IUq7LBaU4ZAjOWhppS0NH4UaZGWA1FB01g66ONKu70/RlayxMTdORVvQWLZbMbmqZbr70jWNcubpALmvhhhXnCNVjSAtdIZkMkcSzsD4g+gcGVNwoDENdsDtbxRNHQpRChsxHMmlaghCCUtlm68YePvTjm+g1q4ydG6WyWKG6aLLku6R7suhpnWXbZ77uUVeApqH5ipSEjrRGp/RISYXlVWnYLjndx811URwZ4fDFMv/+1EnshkcmY+I5UXVPxdWe1nY4LTnNykMIrreA5Gcs0KgyJBKSjqxKNK9rmqBSa5CyDO7ZM8TenX10yCoT4xMwP0+t5tOdhQXPYKHmo6TOqoJG3XYpVxxcKWnPKuyGR8PRKPZ2klu1mivzPs/uH+XkuWmyGRNNk3hec92uCDdPx2VwWgSQCPzXC2FgYKBFAOJGkyQkFkWEwMSawBmc89G0YB1eqTRoL6R5045+bt3STWcKdKeCvTSPU61RcT1cR1HQPDQBC45E1yRWe5G2Qp6yzHB5psaRM3McPz2F5/lkMwa+Cjc5xKw3aW+J96Hmm0WQhDySyh4cGFStiU84vFk6j90gULhsMbHWdwPCPkG4n9dxPWq2i2XqrOovMLKmSH9Plp5iipwpkL6H1CWu6+MhmS05LJYdRq+VGB1bYGqmDL4ikzGCTNSL8EIkWG8t3MSUNnHvDQ8BiMFVq9SNqiWxVG9QRkriQrJsFgFOlGWBiAuUjYaL43pITWIaOpapY5oyztTsuovd8LAbDgIwTQ3TCDL1ZKstelbMruKGNMYDEvetvCgE4Q6ShBTVihua5h7hQOBPQl03Xct9kQKiBoVpalhWwJDyFY2Gg21Hr70IpARDF1iGFcCar+JsLZ5zZa+/hc8bqPvGw8PRAYV6tDE2CWYr7458PCCC+DUUqUs8z2sdHi46krtNo/tWdp6CndyRuQbC8KJiRFMzLSw209toczQo1cp8ZIAq8VupJMY1BSujmB69DiOkQEQtwWQBMVFVjb77vtfcehpSGIXEVvGrG2ojehUmegEqYlhFBYtoFtFcpTZZT8hn5bzJx4nkeYXjOi1tsxWtMRH38kSYOkYMR9eB+JW3mGaCJAPR1LRK5AtErrNSCiutNizCCBlYR+u7CTGRwdA4JCfmEqECk48NLwT9SZO77roLy7JiHmXr+wARiU2ED361/lYJUIqJjCOGSNyTZLhFIjc+RDP6RK+0NGt7iaQjMWcCmpFSomt6oiHUzGt832PXrl3cf//97Nq1K3bdFb3BlUja+uCWvCCuIsWqiVdirUnI/5/v5BGsy5qFy/hkYspkCG4Z5isczw0Luxqe5xKFTE3TOXzoMACHDh1G07SAp6GhIbXS31ceieh6w4yqKfAVAnqD7OsNWG/O1CJcFdq7ivHwulfgEjNEhxSyWRgN5w7eF/KCV33D1FmPJ/jPiI14EQIRMaYSsf9GR+J80nCjNzmIWlZRME+aSZzUhJ/h7whg3zBTbSE5ekcYovRJSoGURsJyBP8PlFNFLjHnE5QAAAAASUVORK5CYII="
        _icon_img = _PILImg.open(_io.BytesIO(_b64.b64decode(_icon_b64)))
        st.set_page_config(page_title="Teacher Pehpeh by IBT",page_icon=_icon_img,layout="wide",initial_sidebar_state="collapsed")
    except Exception:
        st.set_page_config(page_title="Teacher Pehpeh by IBT",page_icon="🌶️",layout="wide",initial_sidebar_state="collapsed")
    # Inject favicon via JS — bypasses Streamlit head control
    _fav_js = (
        '<script>(function(){var l=document.querySelector("link[rel~=\'icon\']")'
        '||document.createElement("link");'
        'l.rel="shortcut icon";l.type="image/png";'
        'l.href="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAB0eUlEQVR42o29d5hk11nn/zk3Va7qnKdnuidnhVHOkiVnbMsJZHudMGlh0w92YdldzMIuLLuEBYwBg22cAWfJkixLVrDiKMxocuyZns65crrh/P64oe6tbpnV8+iZDtVV995zzhu+7/f9vqK/v1/yL/wnkShCQUr/pQIIfw1CEPo9CCEi30ffr/UOQoDAfa0QovUKoUDo7yXS/6Tgde6v13+Gey2tL6Lv3X794Sva+N4FyrrPkRIUReA4TnCv4c8If+/+3rtX77Xh3/vv4X4aCKF4z8T9ueM47tWFPsf9l+CpEHq/4PnKN1ra1v0rYt2vRNvDFO7FICMPGP9/JAjp/dy9cOFeWeQ9ECL4O4FEURQUIZCOe+GqqroXpCjui6SMPHTR+kDvQbX93n9/f2OI1iW679n6V4j1GzJyb95rFcW9D/dnIrzfgw0vgsWObnpFcRcQ3MX2n5G/+ErwWuF+Ttvucv8WnGCh2z5XiGCh/fvyr8PfDMF9+n8vQnfsXa/iLoxoXdy6EyXXLahEEnn+0t9t7t+7N+wvimjtOdlayOBBKCpSSmzHCf5eUZXg4bQ2UOhECiWyTRWhBAsfnELp349/w62T1Fowf6OIwNr4h8a27eDUutcpUBQF6cj11kwIhGchw9Yp2My0rJy/QRxvgd2FFpEDFLYQivferWcpQpbJ+9q3IAIcx92UjuNErJoI7XRFURDevSlSeo9dthbMX0T3IjxTLOX6IxO2nMHNRXe6/xopJULxFta30Z7ZC6yBBMe2vRsiak69Iy282/YvR3gPc725k2EH1XoUgtD1ET1a7e8Q2oTuAyWwVN4O8yzkG7hOz0UFz7LtdxLpWQoZ2gwyWAf/M8OWJThcG3ye4m0UGbFUrc3t36vjOEjvx0rogIQuNGoVhKIEDyg4Qb4D8T/NP4XBArVWSOCZeG/R/cX0N43jOMHJck8SgZUImzf/zp11Dzm8hhKhCM9vB9sxZHZbO/eNNk2waWGdX5fByYXAeoYWMbAk/jWHrFL40xTFPe2OI1tuAtHmjtxN4Nh2JK5wrYT/tFvb27fg/mv8axHe2oQ3hf83SntwE/YdQrim0LEdkKAqCsIzhYhWANGy0K7JkoFZXh/wiCDwEMFOF0KgqIr3vi0TLoR38oT3sIOnKCOuQQlOjOdThdKKAYQIFkFGDWLkdEhHhqyBt3C0/LgiFFRFDT7fdyGuH3WvPbxZlOC+1NDPo4covLH9exGytThBXNF+ooN/RciSKNHnI+WGsU1ghf21CH4pJaqqejeiRP2YAKEILwoN+9bWaW/tvPAub31tWVZgpmzPZwUBjbepFKUV+QtFRAy469ds7zXRgNWRMrSzfXMtW5vGj7yJxiXhUylU7/684x+Np92TaDt2yLG4Gy74rNABagWb7QG1jOQfvusTbZu+tdhuINo69bL1/sFzcp+V9KOCILARoQ0WMm6hOEj6G8DdhEpgnv0PbE8lwkEM0rUOInKzMlgMKaR3Ioksjh/Euabfi2YjZjB6QsLRtaKowev9WCUS7QJSEHmKMnToRSskDB66IkL3KcLmUrRlMeFrEpF01HYc1415D9WNSZyWv/W+RoYX39tsouX7gwDb2xx+nBK2lP69K6rqPTNwbMdLJVsuE0FglfzrUlXVtQDec1MVxbcA3sV45lYicaTjRZ9KEKBEH0A0IJKEfL63gZSwCQjlr34KEU7PhFBai9GWmAYBmKK47khEfbK/MdrNb3t6654oBVUoEdsYhRPkugAj2PgibIJZZ/ZD0V0oK8C9vg2CQindAyMj8bTrdoIsIfR619rZrluWTivOCNycQujxBtbaNwiO43jroARuVfPNnm8qhXCjrPAJk9I/GeuBniA4k+6pR0r3FHofrqqKawmEZ6pFdGF9DMDxTpF/0pRgkdy7UVXNu+nWjSveDduOHfH37uvkOhwjfDqwZWCxwiY4kueHUi5/Z7anya3P8eNhEQn0pJRIx/GerW+mW1mI77udkIsKMpTgNIhQ3KMEMYsQwn3moYDR/1kQv4QQtxZGoSAdB1vaKOEoknAm4Jscz0cLIXAkQWSuesGNaM9twm7EC/YUtfX6cPolhEDTNBzHwbLMwH0oShR1VDxL44TSrlaARjvKESCXkdA35Mb8HB8v744kiqH8fV3AtQ5tlOvBJO++/NRUSnfL+ie2FXS0rJATwl9kKD32F9CPA1q2VqJqqpftiEhcE6TMbQFj2Ho73oYE0Fr5+HoUMDDzIgCPEP6O9X2exD3l0g0w2lHVYFd6X/tm07ad4CHatt0K2PyIPAShut7JDq1nK6qO3HhgzYQbWDpKEPkqMvRZ4YOqiPXgjg/OiFCwKaIpXStgE54bcjeT7dihDEKG3I/wAlsvdvIXQ0ZTOdrTR8+i+bCmkD7g03rWwj+0iuK5D9cNKqqCbdnr0MvwAikyFMW3AhnZQgC9BZGh0+/I1muFENihr4PUKhTw+L462Il+SimE+zACMw3us3ctgx/thndsOyQbCUyDky+9a/buJYKf/xQcODDdoU3vRd2q0pYhBViCE11kxHoX6VtWx/1fyI2gdxmJnzaEl2Qky4ugsUJR2t5RuvFSxB2tx7zC0ZAHydrBBziO45lLGUSY4cAuvKCiLT5Yj0G2Tqrj2G3AWxQ1FF664th2xDT6mUb4QUdMvwwFZn5+HViQ9UGs/0eKcLEDJVScac92/EA0nK34Gy7IBrzX+fGJfy2i/XmFoPVWwiKi6a3YeHNG7zpUEwii/VbcFnFRzsaFIaW9zIJnPv0qnRt9q7hlg1a1RVGVEEYvAkypdUHhndp6H/+B2cFDlhtCsLZtt7JwEYKTfd8pW8dBBPm3wLLsAKXzASc/q1H8DeTj7h7O7oI1bdVMKSLIqB/M+cFX8HofrWwz24FlCwNQ3heqb2G861YVtfU8ZfsGDwNfRLOr0HrYnosUG+AR4RpPuOIqAa1lzIigcK2UhZZf8/F6x/F8UQtl27DAGsBPUei1Pc2LVA5pRdMRk92ekgGqpobgWYl0iIAwimhBrSIId1tuxX/Qjm1HAszAnfnpDDLY0P694EfzPnAoQ4UlQcRl+vGOH0o4XswQhqj9YNkv5oSzDSnDBkGG1kK21Vb8f2QUuo+kwi0sQbgbIASgtNewvfMjvEDHfYhBYTtkslqpkm8q2yN5gRsZE7WAwQ34CxepjcvWa3wE0JFONDYJZ9B+xc57MLZsQaISgaK4wRGOB0d7mY0PWMk2DMHxU0E/UBPRCqkMFVcipVZvpaXnyiIFn1AJ171O1204oZMvVAUsGWQTwqu5BIGuI4PNJttcUsRB+JmHaGVo4aBcCNHaAIpQgijVB1uk4wTwsOU4kQfkR8rBRQZmMZqThpEoaOWtgfnzMQDvlETIIUKujzW8zeI4LiZAyF9KD5WTyADilEBc1+hBojmuG1A1gaaoqNKh4EhmAcs/Zd59qYpCn5AI26YuBDUEpgDbBTRaRbNw4ScChfuwtKTltbyysqoiPRhZhu7XT21t28FBeiXb1mlRNQ3Hg9Tbgvk2Ak4IR/ArqKpXyg7xCBwpXSDIP+WK4qZn4cWVQRTunjgb24voWyfN3aESR7Yg4TBqBqCo7qZQhRJcgKqqXr5MUOUTInxSQosbCsp8/F14G0pK19/HhESxHeJA1tDoz8VZMB0s0+I9t1/N8M7NCMUhmUoQTyWolCo8/9BzHD4zxWWzFTjaQK8iue+a7WRyaSYvzbCyVqJYrlEwbVZVlartHghHEQhHhrZCy/QKIRCqfxhA4GZQ2HYLvZQe+8kralmOuw7CcSKm3n+tbCNjtQe2MvQaESnweddih2B31wKEfWsrSvWRMOGhRkGaZjuRhRXBGokwbuaVOp3gVPjVNsc7tf5N+aXOdshGhCHq8Gbw0TtVIL2CTEZXGUsadKeT9A92M7Cpm807Rtixd5wvf+Y7PP3iSUav2sHbf+5uqoUiUk+gCdAUm7E9Yyz/xl8xPVfAEgJVEahAyrHZf8+13Pue21meXaJaqDA9Oc/rTx/ltcOnOVm30CSkpU1RU6lIieN4eTXRhQmsoxeAOh7vwc9AVM21gpZth/y4Z96dVvblF9R8oxDmSYRxkHA5O7AItgzijnCkprkQbpjXIUMplfdhoVhOeKlN2JeESX5hGoYf2bo+zvHIJU4Q2bhpXmhLh4IZFyGUSMeOXLQUYAhBDuiMqSwLHaNp8fa33chd772ddDpJT38O27FoNgW9/V0Ix2F+bolGqcLiUo0H/+KLdGzq5e0ffQv9gz1kBnuwriyjxGJICRqSXCpOT28Wp15B0VRS/V3csm2Y6+65li//r68x+dAL7BsfZMfmPk4cucBsqcoaUEZgOSHGHKHnqbqHyQerfPPvp9r+g3SLZC3mVSsFbYcu/OcfwjmIboIWdhDFbvxHrYW2nIfVt4Ie/5qkIz1EUESAuAi5MwRSBIG7I10T7RDNzcPpUTjSb5mT4KG0Klmuy4mpClf15jh0wy527R3jK1/8IWdnl6lZNoNDPVTLNf7pSz/i/NFzVIsV6qsFpGFg1hoIRUWP6dRKZWKVtGuGpUOjaeF4FBKfNtbZmaFvqBdFN3j8Gz9g9rljvOW//QIH940yvGcM5ZtPce3N+/ngL72bk6+c5uRzJzjz6hmmZpeZNh2WLfwIBSE8q+UtuhKO2hEhNlWIgyjD8LSMAGZCEQGw1E7KDdcQWptBtpWGW1Za87dIwPLxSYh+OhH6Y0VE/Zz/Ropfu3dChSFafsv39YKWAUARkfcWbYWPgL8HZFSFHlUw0JlmoVhl94ExPvWfHkARGq++co7zs0tMnJthbaVAIhHDiMW4/var2b5nM3/+P79Mw5zDatqYlkNCNfnE//olFBQ0HF557hizVxbRdM09nUKgS4fewW46urI06haNpoWaSZFMxnBsh0a1SVcyzoEb9iDrNa46tJsbbr+Wmck5Tr50iu9/5RF+MrNKTEqaQtBwwtCwl1G1VS+CANxDCluv94NL2eIBSiJQuR9Iy7ARFW18QNGWiXjfa0Eg5X56hGYlhBLBs/2I0gkKRl5VT0YhehGharmnXRUqtt0ifkpbRtMTKQNYUkoJisACUkJw6/ggt953HXuv3c5f/uFXOX3sElMX5tk0NsBt9xziheePMTW7xOpykZHRPt72zhtBUVldLpLSNbJSUq81sS0LKeH8iUnGxwYxcnGOv3yWermKVFTwahIxYHR8iFjcwKrWeOBX303VlCSkjSkli+cm2X9gnLHtQ1iWw9/+76/RlYrx5vfdw46D28n/3fcZyya4dusgl89PcaHSZMlxwS8R8CVaoI///Fx2kWz5dw/bD0Ji0bLI/oPWNdV1lSEcIED+fNpdiPkk2thFaiqd/nSokL8Bx74NkGljmqiKGnDxVCXKzpVtTGTaclIZIpIIoQSnP6GrZJEkhMRQBNfu3cyH/t37GBjqplSs8/JjL9I51MuuvWN093Vx7uQEF85PceDqnezcs4VjRy/wZ7/zeY48c5TBTX10J2I0aybX3XMtmUyWz/2vr1Jr1Dl04z627tnGlZMXmJhdoemRYrp0wTvffQfj2weYmV7i9NEJUoaOZUuee+wVXnrwGW59373s3DuGbbmFnVOvX+Clx17i+NNHOTuzxKGrdvKvf/+XGNjUgzmzyGK+6qap3qlVxHoiaiu/pw25C0PpBPFBmPDhZ2IyzCoK+Jrrq6LCg7+1MMwYWad2pk2IqROlJhEpaxJN4lqvEa2oVYZQsWBXOO5N9GgqmxMx9l6zjd7uLD/48RFePXaRM0fOs/eqHVx32wGyXSmuvWEPtm0Ti6nc9dabOH74NFcuzdFsmvT2ZPjIL7+b4U19jO7czI+/9yw/+Jvv0KjX0KRFd8bg8EMvsHPnGPGkweYdm+g4OUlDChQp6U7H6OhKEUslmbk4z+d/74vYnTkSMZ2VtSLd3Tn2HdqJZZrk14rcevfV/MwDb+HPf/uzfPUbj5POZbj3vXeSSWhce8chHv7mU2zLJejNJDl/ZYFJIajaoAgZIsR4+L3jgJBtUPDG9K4wI9onioqwFY68hWzjukgcabeAoHWtNeFo3i9qCNqICh4yF0pFws5NVVun2g8wI2XbMMqrCAY1havG+nnXJ97OtTfvJZlOMzm3yrPPn+Do86fYdXA73V1J3vLOm5meXmZhYo6xLb1cfd1O9u3ezOSpCRwBvf1ddHV3MndlnqMvnODMkQsUV0v81i/8Kc1ag3LDpGw5nPv1z1IzLRwhmDdtKg0TRwou1U3+9Hf/ge27NhMzGywlEswUKwAkHMk7bjtI/0CWRq3JV//uIRr5IjfddS2vv3aWmpS87bYD7L96nFq9ybe/8DATr1/gLR99O2957108/dVHefzHr3B4pYLjB30R5o6InP4WkUVg21awPr6VCFBMJwqayTATy0sj3VJ4lPWiRajePvcuZOZbqJuM1LiDMi2tgC/YDMHFtXL+SJ1eKKGd615Sv6ZwaPsgn/jtj7Bz1xamLs/xzKMvYtUbNBNxnn/2GHe+8waGhvs4e26GP/mtv6GzJ8d//uNfoX+gk9vfeTM//uIPePQbTzA3Oc+xY5eYX8xTt0w6ejP07xhkvDNFOhujsytDKhVH0yXxmI6iqNTrJs2mSbFYZXG5wupikbNTc8wtV5guNqjXGwhdZSht0N+XQyoK8YTBJ3/t/bz6/HFefuoVzGKZbcO9vONn34RjNmmaNqv5Ev3jI9z9jpsZHMyx/77reOzJVxGaQtJxsBE0bOnV9ENtZIqCIsAOytpOUAiSHlnEcUJ0MCUAB6JtbWEwr51DgUQMDAzIjXrIWiBVtDLVikRbOLxPGvFpXesrWSJyQWF404dHDyZ1PvWfHuCud9zIay+c5G//+Jucnlqgrqg0NI3OWp1f/ZV38fafvYdyocT//b0vcfrYBP/+v36UVDrBI994gid+cpKGkGza0s3OvZvYtmOAocEO0qkYzbpFrVYHx6ZSqoLVxJFQrDbRVIVcyqDpuO6puzOFEY+7KJ+0WFmrc2k6z6lzi0xdWERfqzA+Psg977iRa26/ivHtQ0zNFfidn/uv3Preu3jPR95MpVgmmUoigJnZJbo7s6CqfPa/f4EHXzpFKmZw3Wgfi9OLnCjVyVsSVay3wgFdXfgcTOk9ZxnhLLTHB9GGEY/R5GUK4dhCkxsFfOE2qxDvPEzM8DebqopI50s7IYJ2dFyGaE9+migEiqaQ68ri2JBfrXB+co5CMoHTtNka0xgfHOTwoy9x81uuY8uWQe568w0cOXqe//0/vkqj1mBwrJc3f+gmdmzvZaAnR6PeoLi8wvT5y5SKNbrSBj2dcdaWV2hWK1iKztpqkaTSpCw1VrQk3UloNi0WnBgrZZPjF5a49podbB7tZO94F9cf2AQKXJ5c4pXnL/J3n/kO2c8/wl1vu560odE7OsB9770N4dicOX6Jl598jbd/6G309+dQdYUXnzjKc0fOoygKdx7cxqd++6M8+U8/RnzjMU4pDqumg+pbgFBPhaZqLkzvUdNbbiAUTIdb20Q013cI598ej0FRkdJBzaRSnw6HGn5TpGwjcCoh/y1C3T3SL6DIDYKVjZgM3jdZVZAVDo6uUZcOCdNi144RNm8foW+wi+WpZSbOTtKnqXzwvbfy87/9EU69eo7qSompK4s88k9PcnmtzJarRnn/h27izvv20ZOLU8+XuHh6gurCDAlhYhUWiVUXSVJibS1PY22RhLlCQjSZXrGoFEps7bIRjsml2QJjuQbFYplKvcG+zQnsehVZWmVttYLSqHH+4hKJVIxDN23npjftId4R5wc/eIWXnjrG5k297Dy4HT1h8A9//A2OvnCc3TfsY/PYAGsrRf7+L77JxFKe67dt4hd/68N0ZmLsvHYnImaQP3GRkm1TlaGqgpRomoZQhMePaBXERKgWLUIZXOvn61lFoo1KJyWo6Uzm0xskfxEGadBs4AEYLb6+0sYCcosyYv16h4ingg4huWXrANftHaOxsErJdqgoKmJ5jatvOUAqmWR812YaSwXe/PabedeH7wEk549e4AffeoYfPv064weG+cQv3su1+0ZYW1xh6uwEKaeKVlnBKM6zurDK0twq+bUyadEkJk0ScYMrFYPOjgSqqiKNFGoigW1aCMBIZ5hda7K500EoklMzNp0ZjW61jllZ5cLkGrtTNYprFRpNWFmu0DXQwTvuP0Tv9mGef/4Mj33zWVYnZjlx7CK3v/cu3nH/rTiOw4++9ROefOYo40M9/OpvfZjNmwdYLdYoVSpcf+sBUj1drJ28wFLDwnTcxVe8SqzPYwwTAwQCTVNDaxUlVQR9gm09hbSBQmo6lfp0uFVKtqWFPg9OhEqMYVKG37jh48A/rdMeBP0xnYN9OT7yGw/wlg/fR99gN9paibn5FeYX1+hIJDlww25iCtx4z9UcuG4nxWKdz//R1/n6t56h/9px/t1/eBs7dg4we/4y1fk56gvzlJcWmZ9ZYmV+mbQmySoWnUmd5aaBiWS2qDCZV+nPavR0dWI6Bjg2Ix06VSdOzVLJJmMY8RglmebwtEEt0Uly8xiLepx0uUSzUkVvVtjb1WCpVKNLbzB5epKJyTU2j3Tx1nddjZ6L8d0fHse0HO65Yz87r97O5IU5Pv9//5lkOsG/+a1/xbbtwxQrdX78pUd45TtP0X9gF4du2sPaapmjr56jqmioXqm8VScg6NMI+gJD7GMfRonCyG6WoEQwl+jJ1GS0hTJiNnwzoUS1GiKki1Z/u1sOjrCB2nABVVEYVOGDv/weDl63k/Jakft+5hauunE3Nz7+Gk8+/BKPf/1HxBM6d77ndtR6k9ePT/Clv/o+l+aW+Nf/+X3s3dbHa6+dxyosMdQdY3FukbRZoFs4aCmN2VqciZrOYCqBica2EQPFaaDGDFAFpi1Zq9WwVQUzaXCpbqMmdJS4jmPA5Ixgxejhug/Eaco8K3NzNJ0m5ytJOqs6WodNpVzBMPNcnFngrj05jq6sUJzSmL60wJ6DW7j9M3v4hy8+x2f+5FtcOT9LvVxHSMnHfu197Nw3BpZFo2kyd3maRDKJogqa1Sqm6WAJhQ4VmlJQsyyUdnqX43gouusWaOMYEurQUgRBq32YmxFp/Y8ohISjylDeHmX3yAhEHAgvSBkCgsIyHVGe4FYFPvKxt/CuT7yVuellzp27wg3X7yaezpJfK/LsYy/yyiMvkUzFKTqSw6+eZ/+NW/nox+9ieb7AxVPn0crLZJQGhgFzKzVUw2CsU/LKNORyGQazglwmRh2VxUqTqapkpmyxUncoN8F0JNLb1C5zVqDogtJsmZ99+23suNbk+NxRbFuyNiOpFR02ZdNc46RZLlU4P7HEwW05HKkwYi9TJcbFWozNvWlOzZvkNo1y+537OHpshj//0x+Ssiw++MA9fOI3fo7CaoFzF2bZv38rxXKZWr3JcE8Hi0sFfu8//CVzi2sc2r2FK1fmeXm1TNN2aykCgWVbQVTvONFuLEUoHkMqtLheRdaybd6oiV1Ne1BwhATqNWJIxwmKFGFViojfD7kFRQivlTyazsQUhaQCTaHQRMLSKvtu3MfQpj6+8bkHmZ4tcuCa7eia5NDN++kZH+KZ7z/LD05O8vO/cjcPfOBmDj97HK22Qqy2zKZ4GdVukDAUymqGYlMhnkyxuTdJpsPgXFXwxHSdJ6brHF4wuVS2yVsSWwg0Q2DoAkMTxGMK8bhKKqXSKNe596qDDO2qcOryZSZOWpx62aFcUCjlBfPLAiOj0p/TGR/MULMUVosmZq1GX8ZgdqlGs1xk3wBUTMnFS2ts3tLH+z5wDc+dXODMi6fpyiWplOo8+D++wFrNZOfBbfR2ZSnXLb75V99h4vhF3vr+u/nIb36Y8sIqC2cusaaqOLYTZGKtNnkZIX/6KXvU50epZO3SNd4GyHw6TFyQXvXOr/uHscWAse4jVkHLU1SOxAcsHCnRhGBfR5KDAznq1RqLCCjX6e9Ms/vqbWzeMcaefVswm3UOP/06jbrJF/7313h9Kc8f/+8PM9yX5S//8iE6KVFaLdAsFSk3IZnSeXlOYXQgy7aBOHmh8eN5k4cu1Ti53KTqgBFTSCY04nGFWExD1z2qturR34SC0ASmcEg7Ge6+c4zXp45x8fUmC4sGie4kjqKDYSBjGk+9OMeRVYdlYdCd0dnWn6Ai40zXVKqWpNOwyDhNNBqUqiYnj1xksWjz67/6Jk7PFfjuV55i5tRFzpWaXHj9PJePT7C0VOTxr/+Ic6+d5T2/9gHe+cCbiKkOg7tGWTh1mcLiGgWhtNDCiAbS+rpNpPrnVxAVsUFtxysGZTJuGuj21Hv9aiLMS29xzcN9+jJEDhEhnYigMcOreAwpcOjgOD//3z9J2rK5cH6aRdvBWCty1S0H6e1K4zgmf/NHX+fx7/6EF586wlStyu///v1UCzWeevhZrspVENUyXbrJ1rEu5sw0FUdn16Y0VVXh25fr/PByheW6QzLmnmhD11A1FVXzaN+K15kTlq4RLo5RqTXZ2TdMaqjApYk805cUssMxN+tRJEJzOXpdeopMwmDVcTg61+RMwWKkN8FoLkEqEaOGiqMYKE2Tqfk8h8ZjXLi4xLNHZvi1X7oLtTvL3z18nEpcoyAULs+tcuLFk6QTMR74zx/jxtv3USlVmbi0wOimPob3jDH78mnWilXKTkvNI7z4ilCCU95iGrtWXNKq2UTooiHKfssFKC3tHhEUJ1zEyTc9Ua2gcLu3aylUoQZb0wG6NYUd3Wke+A8fYP/BrZw5fom1s5dZsh3yhSpjA91s2zvKwuwyj37nJ7y8VETtTPDnf/ZRXnjhAueOnmSTUSWGhaEp5MtNzpbibB1I0Ner88iVJt86XaFkOeTSOrGYiqqroELTcaHdptmk3qxjWZbLR7QFilRcPQCPjlctmWzt76ZzS4Ujz5dQM0k38LXBsi2skiCGw9hgN0XTQmKiCotqU/LiVIlp02brYIK0AnNVwXLZIW0IllZrDHTqDPSmOXFxjdtuGuO6qzfzgx+ewohpZGIx7r7tID//2x9lx44RrKaJ5Tj85e99HrthcefbbqR70wBTL59irlp3+1kjghxs0O0kA4a2CGkltaukyaAcnEp/OqKWEeKvhdNDRVU27C5pb0D0CxFpVWWrLnjvp97BPW+7ge98+Yc896OX+eXf/SRiqcDc9AKlmUVuf9st2LbksYdeRM3E+cwffYjvfvdF1k4f49pem3RMUhYx1Gwniw2Va7ZmmGxI/v54hamqJJdRiBkaNhKhC6q1GkpdZSg3xNbBcbYOjnPd7kPsGd5D3EyxpWOE2qrFldU5UrmEawHyJpt6UmT7TI6/ViHeEQcpsRybWM3gj//9e+hJWayUl5hZqpHNpvml93+S99z9Vt56823s3X2Iy9UMuY4OdmSbzMwtka9BJmmQioGCSWfM4fuPn+PWW3Zzz607+PaDR7lz5wi//ge/REcmxsyVec6emWLHzlG2bB3mu196hN7BPsxSlZdeOMl0pY4tJQotAYhWS4+MUOlFW6/FRkp4fiuZ1l4Klh65sT1l8Bc/InYkWoVfH4vwgaNBVXLbm6/n7nfcxOTELA9/5zl+4f/7IDfdfhVPfu853nrv9aR7O/jGFx9m7twUK806/+ePHuDZ507QmL/C3pyNWaqS6c0wX9TppsktB3v46skih2ca9HUn6FAFupOgbtZIxZI0qxb33/Rm7rz5Fro7O0inEqhKjGeefZ7zk0fYd42Dpayy8+p+Bl/RODYxg53xxBpEk0rexLYFiuIehPJak3dfvZWBnssc3NGkJ7WJIydPsGfbHu6+/lbK9aobL0nJjXv2UzElql1hbNejfPtb36MidA5syfLKkQm22SXeNJric3/1IJ/41H38xR++nz/53Qd59uHnecv770SoOl/6v/+EWa3x1g/czc692/j6732BecviymqJrTGNBVuSb1oo+E0rYUEjGVEekT7bup2n10Yhb3ECFdefuGRPJaApBemF9MQFwizUULdJmEGk4tCZinPTW24kHouhKQr/8X9+ik0jPXzxz7/NydfO8Xt/+xukc2n+80f/gNMz8/zJX/8yj/3wVQpXLjAUtzg277B1KMnJWcm+0ThK2uB/PrfKmiUZ7k2RL5XoifXxu//+P2HaJol4DCQMDfbTMJtomsbiUp5Hn/geTu40eo/O8YkittpkZCTOzfeOcm9jJ1/81nOsmXXKtSpaseluYgE2DoalcGB7nUuXJ0kkEuhCoVJr0JPtpGG6xSWXjmWjqRpXJicxjBgj+36WB/p38eKDf8+Pj82xsz9LubBCrya4dTzBw997ng9++D4++qt382f/6+v0DvVwy5sO8fO/8XN882++z9nXzvPoU0coWRZxy+bNt1/F7j1beOjLj3DYcmjaoc5QIULNL+GeQY+AKsOye0QZxFKiptOtWkCUHCIDpa1wQ6TfHEJIv84nJPilyowQJByH6TOXyfZ1Mb5zFF1ITMvi8UcO844P3su+g2N85wsP89jjr/Lvf+d+Lp5foDxxnDdtVVGlhYgZTBYF12/PsSJU/s/zy0hVoSMTo1JrsKVjnF/5yMfYunWUdDpFIhEnkYpTqdYAyezsMt9+6GsUxSmmFla4cmWVldUSuQFoWJLJ6VXsusW+A2mefWqWmAHxjM3yAmgZlXKhzntvu5E9YyZgo+Q1fvBylSu1Ch955/sZHR7Gcixsu9Xyns5kWF5aI2EoyI7NjG7fTbY+weNHF8nmsjj1Bp2dCXpikgcfP8W73n0dVaHync8/zt6rx7n1nuvQY3H++XPfZ6lpcnDLIB//5ft55wNvYs8Ne1m+ssCV81MUpQj58lZE77N8fA0mRSjRCqyIxgQCEcYBRMS8t0iKTiTyUNz+qnWKoD5tsVcR3LN/lJHxQaYnZjnx3HEsy2Z831ZUReXq63YwtnWYU69P8Ce/92Xe8ZGb2L65l1efPcL1w2AYCoVCg8tllWu3dXCpKfjMSyt0pA3iCZ1a02T3wG7+4Dd/i77+bqrVmivWJB0cx4VNu7q7+KvPfomXjrxEQu8gJjvIaD3U8wbLSw6ry1UynSorpSKa6mA2HAr5Jpu2GyxOS2wLbtq+k4+/fxvY81w5VuavvzLNTy7Mcff117Br8y5On72EIjT6+rsRQNM0MXSd3v5e95SZJaqk6Bs7yHh8liePTjMy3IOmq6j1Ev05lcdfnOJTn7yLZ165xMThC3QN9/Kjr/0IIQTv+eCb+MivvJvt+8bQFJ1CscLFk5e4dG6SZVv6yh4R8UhC3cFhudQWZYwIm0sI0UICI4og7f7Ce8CtNK8tRvBMkKEp7I+pfOI3P8wd77yJv/5vn+c7jx4mrSrcfet+3vWL72ZgoINCqcof/Ie/gQ6FD334Zr79Tz9hUKtAvUoiaVCwYftQBzNS4W9eXqUvF3NVRmIqzYrFf/uF/8TevTsolUpomhbpVEqnkhx9/RQXJ19nZKugqiwyv7qEoghyyQzx+iZWFkx+8sKrKF1rbNmWYOqMTX7FYf/dKsd+CP/hY3fTk5hleWaW89MJTs1B19YYg0OCQt5hda2BJlSy+jA0u7j9tlsZGx+hkC+4sLgAHBOz0cAkTode58qPP8MzRy6xOW1TXFymqzPJmXICkRngAx+8nV/7N19ni6ay57aDvOkDdzMyOkCjWsNxLI6+eoGHv/AQc5fnSHRkuFyoMFkz0UKqZmwg6OVbarczWmmT9XEPswsEBZJjbuNGVBolRNtu4wSGFQccoAeHqw+Mc/+n3snMpXl+8OVHOVeus6CozE/MMfnKaQa3DPDsDw7z5EvH+Y3f/hm+/50X2WzPsSvn0NMV40JFY7QnTSVh8Fcvr9DXGUdRFTRDAxVkXfCOu+4llUoEoJWmq1imRTaT5vs/+DGnJ57DNpZ55dRxZtamqTQrLCyWuDy1wPzSNI5d5sY7RlieSDA/W2Npzmbnzk3UtAXi8STnXl1k4tgMr04k0Ea76dluYitrnDqzzNzyKnqmRnrAJNVfoatT4ZXnz1OrS3buGqfRaAZs6UQiTToRoyGSdG/eh1g+wYl5k/6Bbk6cnOOqfb0UVtdokOD2u/fyrQeP8qu/9QDbd45SLFXIF2v09GS4cGqC1546wts+8TPc+/67WH7tFIuVBrVwQC6JcgJFSFnV6/mIKKz6Ql/pdOrTAZrX1uToEged1qlv1yj1qctIdEVhSFd47yffyfadm1heWmNlYQ1ruUijVmdZ15gvVrj8wkleOjrBx37tHhrlKoXJiyTsBvNrFku2RiIRp7s/xR8/v0Rn2nAXP6ZTa9RpFCw+/PYPcO01B2g0G+iaQaFQYGZqmrHxUf75W4+wkj/FpoMW08UJQKVeiNEoCZIJHRUDpODKlTxzawt0dqksX1HYNDDKjkNx5vNTxJOCC5Mljpxx2HtnF1XlCjPza+SLFqsLJh1dBo2KyuIlk9krFWSyyoGrBnnqiZdZmKtx7aH91Go1jFiSP/7cX/Psq4e59ZprsPQOBjZtYnXiRY5cLNLXFSebhK0jOZ4/fJFrrttJrWlx9LmzbNs7xjc/9yA//OfH2XXtXnbuGWPfrQc5dMNutuzaTLna4Pxrp1lF8RBh0daY2kYKaRPsCOsGqJlMFAr2++rbBQUUVQnRulohpyIEUgi6pMO1+8Z4z8ffikDQ15Pj6lv3s3v/VtKOgzW/QlUIXl+pcO1NY7zljl386OHD5KiiCUlqqIeqrXBga5Y/fXkFoWkYuoam61QrNW7ZcRO/8qGPc9P119JsNgPNHduW9PZ2c+bMJC+/8iP23RHj9fOnWLrioCjQ3R/DcWyGtxoouklxzWJoPIZh6Jw7WaJSq5HKOlTjV9AMlYWZGtWyRd+WGKV6mYXpGukuk74hjc6eGCuzJsNbdXI9KrnuOAtTFj988BRvfedV5FfXmJspsXfvbgQKLxx+mX989Ft0dnZw4949WIl+NvUmmDp3jI60AbUSpYZE1WPMzhd5x3tv4vvfeYVnH3yOCycv0zXYy7b92+jqStORS+EAJ1+/yOtPvsrEzDKrphPoAxES6g06rpQIBWRD4Qg15TGC2jVx30hav9W67UWZiiCmCMZjKh/61fvZuXcL3/riD1hdKjK4ZZDh4S72Xr+HA4d2YU0tMrtc4jf/7d189R+f5+p0nu1ph+6BDs7kBfsHYnz7cp1LBYtMUkfRNRqNBtu6tvHb//bf0tPb5QI9XvG7aTbp7MhQrdr84z9/g+vuSfLUC0eYPm/SM6jTt0nn7JECQjPJl8ssLtRAWJx6vcTrz9TYMhKndwusigLLKybLKyZLs5AWChnVoSuh06hqHDncYHWlRDLtgK3TrEDngIrtmHT2KSQScSanpxgb62fhSplUsoPe3g4OHbyaSrXGNx/7HpPzs4wNDFFRBzGq58iX6mDkWJxbJhlXqNebxLNZBrb08vRT5/nF//hzfPBT76K/J4PQFBYX8zzy1cf47t99l1fOTDJthURSNpCiXRfPSf8QB/qgXmdQW8DXrlAtPNGkQLGK6O/9RsauZBJZr/PKc8d57p+eQiA49cIJbrv/Dnbv28z49kFmlkv87HuuYnm1DqtLdHTpTOcla406/bEYE6bCizM1+nIxT71GYthxfvVjn8DBplZuuj2CHqetv7eXy5em+dZ3v8P+G2McPzuBaCbZssuEeJMTRwrUmxJDSOolqK6CVe5iv+zgQ92TdE41cC4Lzkl4TZMkNJXba4JhW2A3Ya5WYX9WcCgb57nZBC8vr3LoTpt6Nc7MFRsj4RBPWCQyKjlDY8E+TmZoMz945CE2b/5FdF3jkx/4EIV8icPHjvLj53/CNdsPsbm7m325Jc5O1kgacQzFoTsrePHpo7z/w/fy1OYOUuk43V1p5ufWeO2pI/z4209w5vI8CyiUpQi6uILyvYhK9oblasKCWsKTp7U9AakgDWxtABk0Yob5ZT712PEULwI8WUosBFXL4dKLpzj64knOlJtcNiWzF2e4+OxxpONw7OWzPH/4NB/9+dt55aVzdKtVzk1X6BnJgKLS35fi8ydKJGK6G9gZOuVyhVt338S999xOtVZHUxVPn0ilVmvw9W98lwcf+2c6R9eYWZ1ET9RpOHXOn1tkcbpBrQC21cSuxFm+JBCVET59393cZc4xsGUI8mU64ho7Ukn2NwTXWIJ+RSERU+lPa2gqrNQc+hsObx3rIBHv5fBFi3jc4vK5KsWCw+JChfnJKsUFg1rRYGZxgdePnaYvt5nxrZtQVZV0PM2pc2fp7Ozi1z7ycYZGd3PhwhlyaoVs2sAqFZCqjuI4NBWD0bFBHv76c3R0pfnmZ77NY99/hhNrFealoCk9jf+QPIUIcTijDSJh2TwlIp8TVBYHBwflRrLpYY666hWD/BEpgZChI73WcrfrVfPYp1Yobsgg2aQJynWbN73/KkY39XDk6cO8ba+OJRUOzzjsHEjxfEnw7HSdrkwMFAVHSGKNOH/4G79NR1cWy7Qi5epyqcpn/uaL5GtzNK0yQnWIxZJkYgNcvfcqhgcHqVarOLYkm07zlW9+j39z/Qi9CxOoN15NY3aFmYeecU+SoqB6kbLpgKoIdEWgCcFcocli1VVLu+/GUT63orHnvvfwje9/ibn8Eh940/1cd/XVdHbk0FSdaq3OmfNniBkGd999C416AxCUK1VymYx7eBwL05bMPPu3PPWTl4kLgSyvYHR0g5Hi1nfcwZ/94aOsTq+yZFnkhUrTkSgyUqEJDmS4RBwl8IigshvWaHRVYGyXEhbu8AlzzCPze5ARNYqwwkTQUYxLuMDjDwrhKmcUkJRMQSwe485bd/K1bzzH7UOSSs2mJiFhqFQ1nRdmynSmY6Co6IbK2kqB973lZxgc7qNQKLrKWt6NW5ZNOpPkd3/n1ykWKlSrdRzbQTd0OnJpdF3DcVw93e6uDv7+S//MzeNZepcn0W48SP6V06y8eByha56kucTx+tV0D/y0HVfhszOpUTWbmKrC5XML7FIt1tYWeM/b3smffvWznJg+zr964H1BAU1Vu9ixcxTTNN1mEq/Cms2ksWwLaUkc2ySWzDJ2zX288vKr5KsaaaGh1iqUqzYL0/Psv20bf/gXT9HVncSxHO/UE9IIaolH+HqL7TFcWMjSB/UCso+v0h7IlMoouyScGTieEpcvd+pswEbwFURFiF8oAUNRKRZr3HPLVqo1i5xSZ6kqeelinYnlJkOdMV5YbCIU1dUjUn1FTYVatdbqLmrrW7Rsm3y+gFAk2VyCjq4UiaROrV4lXyhQKlepVGpMTc1x/MwZro5LagLKpyaZ//ErKPHYevHJti8loCmClKFgScnraw22CJXXH3+Et9z7Vt571zt4/cwJJqdnMIwY9VqDeq1BoVik5p18vyfX8mXvhKv1YzWq/MU3nyLbmaYnq5NLx0lKky3dKi+/cIZDe3vp6kjiWHZUTC1a18F2WsLR7XpKvoJaWGAi0ED0WVxhSU9/uJEWyK957KBQD19LizaKAsqgICEjiKF/aXfcuIkzpy+TSir0DSYZ7TcAlZphcGypSTapuxCzr2EbVynVq1EtQELDqoQIAkLLsrEsX8NP8XR0IRYzWFpZI6HY9MQF1bUKy6+cQs+lXDZT6P1kWMTa4zMonoyBpilYEs42bRqmoKeS5/zlK9x2/U0IAV/7/j8xOzdHOp30rIC6TgjTJ1eZlkk6k+TS5St877mnmZRxkokYtXiWWsPBrjeIWxXspsmtN2yhVGmGhVM85a/wNBYvwg/VZNwpJm5vpq/zLCOBuxO0myvt9X3HltFZOdKTMAv3B0TF4QKJ+YCmGBJyqtdNdo730NuZ4OjRC9ywLU6PYSPjcXo7ExxbaWL7zGIvt5cINF1jpbiKZdob0JkjasHr5/bI1gCMfLFIXAVd00h255DV2jqas3TFEN0+etsmrUhyCjQdSdF20FXQFUGPqrBiSXqxmbwyxd4du9g+soOfHH2R3/nT/8XlqWk0XQ0VW6I8fCklPT3dXLkyy2e/+kV6elMcnqmSTKvYQmOVBMcuFlA0lfJqnttu2OwWm0K9fO5kFTXUjNvCbBSh+FrqnqRetJO4nThCqNwX8fuuELET8P6iooQyomejqmrgn4UfdATtR4JKtcmt129hZq5Ad8xibbnExEydfMEinjY4vmySjms4QoTValyZWtNuMZU21BQV676NKL8rAtu03FOznOfiT47h6AZaIP0qkaaNFosR60yjNhrEMimeqMKXlk1O1ByKXlAogAFNoWk6GLZDMZ+noyPLp//Nf2Tb4DgzhRn+9itfQlP1iE6SK+9ioRs6qqbx5DPP84d//RmmS/MkDYPlksV01UaTJrGYIBZX6VCaXDg/z/bRHCMDOZqm3YrqIxKwIipM7QNzgUhESOLWl8tta/5R1o1b832/FxwpQl03Ii48/KAlKxtq9ZYudqAIBU1T2Lujl2phFUURPPxaETtuMNxlMF2RFOsSTVXaTKbAbjiMDY+2GK+IqAt4g94TGSKmSkeiajqGhO8s1fl3Vxr8n0ULs7uDoW2DOAhs0yY53Itm2yQ3DfEFK83fTld5tebwQsX2gC43qtaFIK0qGKZJo+n63qGhAW679mYqpTr5WoFisYym6p7ql4OqaWQzaZq1Bl/86jf4n5//M9bsNTKZDLaU6KrgTMFy4xepsX/IYCQraeZXaNYbHLpqE5VqMyIG4YQ0CsNBX6DNqGwwWLN9gbwHpqbT6U+LsKBz22iXiPCgjJ4ysW4GT7QEaZo2vV0p7n/zTg6/fJ7+pGTfgMqFFYdcUud4RWGx7hCP6QhVaQkvK2BVbD7wlp+hp6cL2zKDcrRo06ON7Ie2jaEKQa1e59z58xxoVNlcWON4zeEnZcnNW3sRxQrnShZqvY6tafz5ksOJC3N8oMfg3RmVW1IaKdW1PqWGgzBthJRUejqYp4GqGvT09bNtbAtdqW4e/ckTdHd2cc2e/UjHIZfNUCyV+dYPHuIvv/L3nJk7R1dvBzhKIA+jCCjUbK7pj2GZDr1Kg0LdoanopLJZUh05nnj2AqmEEQzcECHmT0tZtaWPHOhuh6yBaBtk0eoMilCNW9oA7YMOI4oTb9T/FXqBogrKVYubt/dRb5hY5QpX749z7nKDbEynjsqVkklcVyNDJtyARsHCZmpmjv37dmPZNmpIKIGwZO8Gay9D6hnxmEFhLU9BVOnMqLyvS+d7yzV+84nz3N+t0ZnQ+If5Oit2g/MLa7ztwDDpoS4OVxrojRpjtSajzSaFuoU10MXSaC8r9QaNhSl+8r0vMbBpC/29HTxw//309nTx0OM/5D33vJXLM9M8/KMnOXLmOHl7hWQmTiKewG56VHqv8qqrCoWqTd4GQ9e4XIiRXyuS602yOL3E5n3byaRjnuxtWFm8JbvvBEM4lUAvSIRnQKwbbiGjcvGR0WqR4kJreEFEG85bAEVRWjr1oV2neHGBbUt2jXextlpkayecmKhwpaJQq1g4hsZazTWBtEmqI8FI6kxOT5FIGGSzabKZDIZuuGyXyIyWqPlvBaES3YjzzOM/Ym32Iqdtwdyu7cyObmK4wyDXCRXT4RpdZSAGVjLOL3/4rXTv2M7j5+doJhPEEjovaZLntRizezYx15tGm5kne2oCu38/+ZJJfnUFUFheWeGWa6/j2l0H+IX/8uv8zh/9EdekS/zabTtRTBVDxlBN1aPe+Qwdrz9BFcxUHToSKqlcnHQyxkCHxszsCqkYDA/maDYtb6ZSa6sHk0/YYFZz+GHItvlGsm1egAzpvsgwOTSsFhLyJTLUNIojPUaO9HJc1Wtdcqtxw0M5FKvBcl3Qt6mTd949xNBAktmyjU1rZg7BWBb3c/S4zhMvvcBnP/cVvvWdh3jw0SeZml0gm8147siJ7Em5Ph5ESkF+ZZGSo3Hftm28/cAODtxwI8+eWGbL0Ai1boMnMnFWnRq7hjv59//2Z+nu6eHixSWGxnexe61MzrG4bKgoVxboe/0yV1frdBgKe/bv5F//199n06ZhTNNEUxUsx2b/vt0szhf41DVbuGN+HvHEi1yzZR9//pt/wObeTTQt031GqvBG1SkYqsp8VVKzYXq5QVa3WVouEVdsrKbF2GgXzYYdyvFbo/ZkmAweUSdtrVl7w29EKCpI/wgpTLUxffwJkUGOr6pIb6ik7Y16VTx1T+k4WLaFIyGXSTDQk+TshQV2bs6yOQfnziyQSqgsNiWqpqy7KlVTqNZqxEsGd+/aTGz6aeonH6Zy5nGeePCb/OM/P0QsZrjixyH+e9tc5GBOXy6XQVcc6vUqTz/0NKPjfXz5y/+DXfv2YC/k2VdsMKwbWHNznDo1wXvefRNf/8b/JFEvYxbLmELl+lKFA6ZFPKaTdwRFRyGTzbF505DbgOHxJpKJGNW6w/17NnNgdoaZc5f40dQa973lbSRjMWaW5ojH460GG2/qqK4Klqo2ttvNyUwZZlYbdGV0rFqDkcEMlhMi5eAO3fBBABk65RvL9YuQInkIs3EPuGi1CkRcgdM2Q8ebQyOEC7PaTqjnvCUVb3sXZlo2vV1JDA0Wl0qkhMXpC2ssFhyaCIqOwPA2gJ9JKCrUmybJapL7948wbp2nurTA8mKe2uwEN3WvkK6c5W//7itoWmvIg2ybM2w7DpZjIZEMbRqlaTvk0mn6l5c49vffIn3iOPOvPo+iG6Rtiz1pg7Te5Md/9kVeefQFpn7wJEPPPk8ikWS63CCjKNiqQsVxkEJSQCGbTVOrNlokC02lWK7wxKPf4067gDm7ynP5Bp3vvJ+7br6Ohx57gnyjSCxmhNrp3PtWhaBqSmpSoTOjYiQUxkfSJA2oVWr096QD8Y6wQEcwR4kN5gtHuJ3raX5BINnSDfDHjImI1l/QVyZ8vfnWWNVwM4njyGCypiIEpmnTlYthmRb9GZXXzuc5X9CIJ3RQVMomXvoXEi1SobFi8XMHO6hPvcJ3jpd43OrlZaObyewBvvz0AsnCebpqx/nHbz5EMpVyJ4Q4jje80caI6aTTaXLZLIoQjGzdiQ48eeI0qcFOBs0Cjz3+GGuVKj/3MzeQ3d7DbiGJJWOcyC/w8rce5OLLh6l0Jnhah5RlEwMaUhJT3FpBRdPpyGVwpB2od+QyGR596nn2k6djrchDyxWeHtrLxz7+IWqlMs+/fphURzJ4ropP1vCUzk1L4gjJWtUhp1qM5tzGFLNp0ZlLBlI8gXNXRIuoK5QwF8xFAAPx1pZ+sIsIRmcZaeGYUHhCQpH5dIE7CI1fDWsJhOThXfRNcwtDtiSXjSOk43UVqewbijFTkMQRNEwHVWsFkUIVlItNRjUNvXCZRy86FMc205F1u2OXKJPcvIMvPvMyH7ld5eT0WSavzLFtbIR63fSUxy2uTE1zeWaGar3Orddfj2kLxkY3k9XqnBUZnHKNmKOTXqkyd2mebM3keKHJ7niMzZkUeVWhYEOxWmdTMsbPjPYyf3oGWwoMb+gEmRzdHTlMywVa0qkEk1PzTLz2Ar9omCwWKzRTKe4fj/Htr32ZZNcgs2vzpLpTHrlWRANrATZQtyWDPXFqpSwXJ9dQ+/rpUGwsQ0VXFW8WQlgsJDqb2S/PS5+w449HDMH94cZRZDAxREZ1Z6UTFSJum6TLBkzUaCex63NymRi25bC8VuXmrQZzCyVWzRjZ7gSmA7HWhGf3byoOW3vTHLu0QL6zn54OHbNmIiVUrAoVUaXa08mZSwXG9m/liaeeYWl5J6Zpo6saz756mKPnTtKQTWzV4uWTR9jaO84dN+zlxsVprhy7QL1UQ9M18l0Gh186B0IlZejsNwQ5BUzToWmD6jjksgmaUjJluvcbUwUNy0Ib7Kczm6FZr2HoKsvLa3z1W//EHWmLtaNXeMqKkdMV7pyd5ozi8BsP/5iurZ3uY3PcfkRHRsmc7nxBQUxV6OjUODansHVAUC3XSXRp6D7ETGSIUuR0EwaEvOnemqcyEp5a4q+pqir+BhC41eCIxEd06LCHrQdTJ3xQQoYnBAgc2wmg9kRMo96wSGsOKys1hB5jKJekZjpIKSKAkiMlSVUjS4VTDZXYUBKnaXsOVqLprlKWGMnx+tk82woLHFlc5LHjT6E63oBrXZDqTJJQ3EBrrrzEiy8e47f/1ftYmcszNbVEIpPCaTYZyBi8Y6iDYrWJLiWKI3FUBVOFumURi2kszefdMquiImxJUlNYrDXIDI4QN3RqlTK6ofPVr/0j78qsUbYVno51MFeucmdGoVJvsHZ6kvuv2sHFWpXpZo24oXszl9qyWI9wU6tbLK6U2T5gIG2LQtkkp7kFKduR0dlBoUkNkSkutCTmbOwAX2+n/QnhzQ4OUDbv5Pq+xp+q5Rca/P6AiGSp/5qQQrg788+lkSUMV61joSqwFY1Sse6BF05IIdyFjm2ziV0rYukJNK8ZVSpuqU4qAtuRxHXBWixNcS3PlmwcI5Ei25Uh05Ehk0q3dHMsGweV8S4Bzgra0DC2aaEpgs6kwaWawzOFJqc7chzt6uJpS+XHqyarNnQkdC+rUZGqjuXNTczoCvOmZHznNuKxGJl0itmFNZTVWVInJuk+M8Hd23vZvXeYcsOibsPTs2XeWl1gh6ZRa7YKO6GaTctqWjbxmMpgTiVhmRRKjSBQVBSx4axDGVYKeUOMXKwr8klcCXolPNhZVdV1lTd/V/nzeFupRhRjboFJ/nb2cAQcLEVh+3CS68dTZHNxbMctFPkDOn3x5KaEeDzGcEeWSqnmVgaXm5iXKoj5JmLVQloOei7JaqVO1qpSLzU9H+oEfEVsEKpKrVjnYL/g+WMXqHZn6etIk9Rd+XbpSDbHNEYycXb3Z7l+UyeKKjhWtlytPi+YbVgykNhThUMxmQbp8NorrzE3t4CUDqlkjLm6xZHFOvaZSXYuLHKy0OQfZ2vsyKjEheDFuRUyqVhomGN44dzFiWuCWsPmyppkarVJE5VsUqFp2aEZxy2+v6txsV77J3xAA3hd+HwAKzR3Wba0gsOdQRERKE+Lzh9bsk56NKRL4//O9ubrIR0qNZNauU53V4zZiws0yJJKJl0TKGQwjUwIgZFLcH5hmeu3Jnj+oknNgvJEgWqzga6p6IpKZ7MDLWswu2ayLVNFmO7EJ63o4NRs6pU6UjroHRmG1CYjCYt/PJZgKD7Fzwx3Mb+whqZqDBkqy6bDxWNXaDoOe7sS3NsVY6HhUA3FPw3LzZ01CaYqmKhIXvncF9GR6IbB+OYR5tbW2JWNESvU+FHexpQWaybkNbhVkXxmskx+IEOXrmCaVvDMRAjAEUh0BRTpkEmqXCkZXNsfw3RcQMuyHE+5xdMM9Lp/w2N5A/1GxQ3a3YYQ543nCIfHxoX1gH3+WGuekYtaiZDc6DpR4tBMHP+mLNtBj8VoNh2mF2vUTIhnLGK6Cz3ZZYHWVCEuEZ2Q6kxwYlZjZ6HE/TtH+MITF8E00XQd07apmybW5AqZgS7mmw57FJuMpWFerrC2UsS0rACYUpcqaJ3QO97kPeOLNM4UOaGpjMV08g03TunSBX0DaWwJlnQomjZpxY3IASxbuvEKkFSgIhQuF6toatMd8JQvc3F6kZv2jlJGZ3dnjHi+gaIovC8X55HFOv+9AKl9HeTiOlbFQirSrQUJf5qKK8KhArqUNCxJwmlw3ZiGJR1MNCzLxjQdVEUgpYhK7ysKIjS7QTphvebwfKaQMISPIjq+BRBR7HhdD5mqedMpndZ8Pl8pLLAGnryMqrqkSl2j1nBIJHRUXWO2LunrTlCt2+i1BumaTr0IS8U1YrqOtqSRHMuSyHXy0JkF9o4ukkqkqWngWHWENNFUlXqziSiWceIKpm3Rl0xz+OIlDE3BcbxJ3FJQazSZXnVA13nzrlUuAxOWSfWiwo5ujZrlYAOm7bQSEfzFd2PrmiUx3ZHpdCc1TkhX6EqxHEqmTSoV4xN37OKBbp2ZI+d4ctWkKuGkDYbdpJ5NkOvMEis4NGfLNGsmqAI9pSNUgZl2sBV3TJwuwFAg37QZigviCYXLJZvciEGp6gpPJ+Iatj9TIHDFTksSxtP0C5Rbw6N224Aif/i0FhZ9kDI6sTNi0j0lcDffdpDCCWb1EpoQ6kgHVahomspaoYGiKvQP5EgIkx29GseWdbAbxE2YWi0S0xUq9RqiLlBRsEyTcl1w5dga6YSBpqjYWgzbUlCcOgiBEdMgE6NpNhjqSdK0HGK6RkqHXEInEXMLUfm6yeSyQiKtY3Y12Nmr8cPlBpmiYDinUPYEmteRTbwaSM0bQ5t3TQRPFRqgKvR3pzgwkOZN/Sm21IvMvbLIigmLKLzU10F2KENRgfiaTWOywGKpEsQnElBWBDFdp7s3RzMhqEmbroSCIgXNWoPu4ThTJZVEwiCViXNxqYZlOaHhEeHRPH5A54Ssb3ukKCPDN1tu3mmlgcFbOQQCxMEsGz9a82bauPq0SmvGfTDmwv1w07JRFVhcLqHqOgiVGHXWlkrUCzrJ7hTZmInt+EpjKpqmUCxVsGwLTVVJxFTqTRMpTTS1gVB0HFSQFuBgOTqrpTrjozEShs6Noxl60zFvPI0r/jS9ZpLrqlFVY0wam3jy1QVuetc2nvvOcW6rJRhOQtHyZVeigbQE6pZDQhfUO+LMGDFu71X5ZHeK8bhKplxj+eIUr+UbnLUFzzdgSTPYNJBGJFSMuRprF1Yp1htIj1irKgLNi69qzQazs8vkkklEStDZrRFTwFFVXrlQJd3dyUivhqUazC8shkq6/mSxlgaA7Y/ObWN2heXi29nCfgynhae4uGakJf4cEX+I6AeI0Ow6b9ijJGjZ8n3OwnIFUNHiSWorayxWHPoHkyRTGgNZB8tywFBQ9CS2WXF1cR13bk5PJkZPysByJCtVk5VyA6G67KTYYIJmGVaKFrtlg/GeNL1pg4WySS6uoqugSkFCF7x4TrJtxOKphX7uGi7yw1cz/PJ9A3zj76e5tTfJNVnNm/knQ3wEMG3PH+sKBxSICxvFtKheqXC51mTZgvO24Iyt4aRiLNfqJIWCbTkoDQc73yBfrZGK6/RnY+iqYLncZLXS9FA6t46/ViphyAxbEjr5ap1+o0FNJtjSo7FQtBgYTXBltoCmhiaP+8h8u2RPm95rkJm1DYoK68Fq68aS4KxjkKwjZXo6QpGOosCNuA/S0BVK5QYrhQYD/VnOzC+h6Q5YFnMrKgMdisu/1TMudOw4NGzJnsEMmzsTJA2FuKZSadhcWq2STWhMLFbIbMqR6EhirdRYrkDMLtDTmeL8QpGaDQk9ju6JIazWTM7MmCzVdTLqFW47GMd54jWEpfHzb03woyWFxSsWd8ZVNEOlbjmonjJK1bSxHEnTdHh2ucFys0rFtChoLpPZSBp09CTZbkBnTKNaNanZNkZMwWrYVPI1kjGNG7bkSMc0ErpK03KoNC0KNZsLi2XWahaGplKr1unTUziWQFENdqUtlhaKlI0uhKoxOZ0nFtM92fcwp2+jEq8SNPL4s4uD0DwU/YuWMhARsxAZDiFlBFUKC0KHEUm/QOSDDf50atOGS9N5srkUvd1Jtg/odCcktqOQNqArqWFZEuwGCIV9g2kGMjF0bxxrtekSOjflYlw7kiFjKJDUEKpAMxSWy5Ll1QrZdJyemGDvSAfZmIoALNNhKGtwcCSN4phcv0Nw/MgC77ypybFjJbJbFB7YJrjmA518X7WYKZbJqR6rBkm1aftVcHqTCqf6Bki85V5im/rZOpJla1+a0Z0HSA5to1BtsqkryUhKJV5uomqKKwGbM7AcmFqtUW9aqAI6Yjq5uMbt27oY7YrTsCVJXWIoNovLVfqzCvNrJnUlQa6nk1LNYX6pQkxXW/o/7aruIdWPcMWQULVWRIaCtgZzKxEssm0KZRikkCFRIhFqDQvTsv0LUBQfSJGcOr9EKpemWGpgNiVWo4FKg4QO27ensISDZbnK3j1Jg3ylGQRmAomhCcBhdq3G5q4kXY5baDKyGkULijWTXBz6t+9k65vewZpjoAINy8Ec2IkxsJm5GZPBIYNvXRzi9Pk699yT4eWzkn84YjA8V+IXfmGALyR28+CaA2aNmHD9vw2kdIW6UMnFFFLLV+iNq1iNBg3LZvbUUfJXLiEl9CV1VEvSaNqu8oaA/kyMnqTGjt6kCy45kqZtM1uoMZuvM5KLU2s6jPcnSOqCoa44py4UmShq9PZniKWTTM2uUa42W7MZpYg0iPgFOUe6pXgfyW1Thw+tW9SiB+3hPssnIjVGS/pNDTUWBmwgX0/Ap4b5c/BC5INKrck737STCxeWiJkl6paDHk9i2YKOLWleu2Si1OrEdIVtfSmGO+JUG7aXTbib4PJqHV1VSBsKTtnEziVQ4yrOVAFVmGDEWCPFxWMnSDRKSMdmpWLRbZUoL8ygJjI0qzWuuWczX3zY4O37V5g7r/PF2a28frrKwVyZ4uZDfK+0m7MNi+LKPLUmJBzBpi6DJ8s2hu7QrTTdgVNSo2AplJoOlaZNRWrUmxZrCR36sziLZbqkTcJQmc3XMVSFuKZgewO2hzviJHSFfN3k/EKVt12bwzSbdFVXiWdijI9kWCia9GwZ5anDk5y7tEIiaSCdKCSsBhlYK8IPU/tlSMHFL0EHZX4RsQAiOhCqLTUMav+BeqgT1ASk44019QdKem1iEkk8rjE1W2R2vsjoll4aWppSpclw2sLWdEaESW+Pimm5J85QFc4ulJlaq1GzXACmZMKWDp3+uKQjriE0BUdVyC9X6OtIIhPdxCx4YItOR9wdC1+u22SSCWbXytDdw28c7OG1FxYRy/N8+pdTPPaMxc5Bk7vKJ5np3syxiTr/eutj/P67LrNrOE1WCuZVh7O6xYk1m4lYnGRngsW64GJN53ze4txCkYnFEucWSpyZK3CpLLExsGs2SrlOTlfY3JVgd38KSyrYwu0umirUOTpddLMCBJ0ZjcGswHYEXSmF0U6VQrkJiTRazODIiTnSScPrA3TrLopQggKaP33dnwsQbvEPBwjScSIDo1TF1YMOLMC6OkKIWeqecBFtD5MiKleuKG6Z0xs1L3AncRfLTTYN5di/o49608KuVRGNOvm6RDgCtSfF7JJDThPkEhppXWFLdxxNWqwVyqysFkjlOklu2kphfp7ZRBwnp5NuCH7r4AAdS6s4QrBneZmGrqNefx+f/I+/TXxwiMnLl/jEUJzqqYs8L9MMp0zuuDbF8cpW4vHL/MwBFfvUAvMNm1pRcnDXGvpPpunLQ9/BDrI/O4SVFcwUmqwaSeYW6izkSzTqTW/T+xmQQ73RpFKoolsSkhr9hgtMlRsWtUqZWq2OYpsMZHTiMY2m5bBQbLJtxKA3ozE1uUKGJhPLNn1b+oh3d7NYdPj2I6dIpwxsP76STmhET0gBzNsQgato60rydYLaGUOtDRAe0x7ZPbIlVR6mHAsiRQaUECfNaY2DVRWFctXknffu5fKledR6CavepOGAo+hsHUiwWhBkpEUyptKVNphcKaEObuXme9/MdW96M2/90MdYnpvl0qmTrHRmMFIqzXydTK1OIZXGKBSoFRsszaxQWZnHnjqL8fpPuMcuwpV5vnJ+hdEb9jCSr/HrT3axktxLp1FHO7nCddsUdt6lcP6y4HghxundOqplIftiWH1J0ptydHWoHH1ulXyljm3bJHSVmEd+taWkaXvFLUXSrDYxejPEFNAqVbr2XccHPvoRlESa7qFBluqSSn6N/lycYt3i1v0p5vMNRtMWqw2VLVs6kEaMzPAoj/1kgtMXlkjE9FaI5geBigims4f9vWxvlhDRmUEtZpB7mrXWzHnRmjYdqjYLoawbQxZWEndjAYIqV3i8PAhSSYOT5xeZWSyT7e1kfn6WlGqwWrAZymlkHMnWXsFLV+p0xBTshEpnwmBtaYn+TVsY7ergkb/7a159/llmYjFk2hVGasY0dm/pI3/0HD9abXLScbg1EWNvLU/utRdwpOSx2TITmU7m9+7jUMLglobCtyZeYC6V5dX49Qx2zbH/ugqNeUHvrOQxxWG2aFN+U5wPx+osn5iknkvSIwzqTZO+pM7Vo71Ix0ETElu68fVKtcmxmRKWdKHk5mKN+Q6FPl3BXpikqcZ498c+ztzsNNOf/zwJQ6FaM9k1FEOoAh3JsGGRTSVZKTaJJbvJajo/OTxJNh3zsiwZGQTlu9/IrBchIqV8J2jdF+sGR4mWWHTq0yI0ITQ8ByjSTSKUFvVY+u1gitusqLaqh0rbbGFFuCQHTVO4/YatLM4vsbBcY9+IRkdMYWqpwc5tKY7XNFYVwexcnk5dJSOrvPT4j3jxR49QWJ1H9gwwdnAvKytr1BUTQ9GYubjGiDS5Oqmx2JTM1hq8UoMJEkx09dG8eh+3vv9etoyN4ExeYUe9yGbV5KXpVT72dpstY3M8V1BovGQytajSe0cPdUcjf9pky2qDPg2uvt7g8KTCK+dsBmIKEkF3SgdHkDBUlqoWAsG2vhRX1upIBLVyjd27x+no7+PysZNcfOU5Xnz8h7z82MNo5RVw4NXZMrcd6mRqrspYxmZursLFNYdNo13Eevo4fXGNx566QCYda2k+eny/QLTLLQPSWj+fCKoGYl+KaE2Bi86AXucC1peEw6NHAo0gTyJGSulWCEOzBZG0jTF1y5ExQ2Vico233bXVpWqnNaqrJTTZpGKrZGM6mc1ZXl1qcOeNN6Am02iZLLnePooYZLeM86GP389VO7dw/923cPT1Cyw5FSp5k/vSOj+aKWBv7mEtk2XzrTdy7yc+SD1usGPrMEqtxiNPvcpNaUFsfoHuTJxsucJ0c46P3Kpz/CUHZ0hgJFJMxbYxlNC4Y7BO8bJkwQBdtfnhVAq7qjEoLDQhScY0pAfNVkolLNuhL51gpdKk1HBoNpvcdmg/d99ziGRfD8urReq1OkYiSV3Rsbv62X31CM7aCo5tMWSvEu9KMzycI2/HGBwb4a+/9DL1pt2S5fXMvogcxFB6J6J6joEQpKpEZze20/DTqfSn2/l94cUM95f5uoC0CQ/L0Mz7VmuY6vWyu+zf1XyNro4kVx0Yo5ov0SwVsGyJKgWTKw7bc4JbDt3Jr3zyA/T1dWACb3v3m+no60HTVKqVGjawZ/cWRroG+PHzryN6Yjxbk5wsVLnn6n4Gt49jri6yeOEs5YUZFi5fZuLsRRbWaryjL4XdqFOpmYxldY6cgJeL8OYOBbMmGdtisTVVYbmWY1gpoQwM8WDpdnKvnqNZLFAxDXbt2kWsqxczv4yuqlSbJlsP3YIqwCwVaDgKK1WLer3Bnr3jbNs6hCLgiWNrbNu/kw/87JtooPDhB97KTmOJ549cZFivk7TroBtUtTTZkRGuzFX5xoMn6MjG8WLqUADnBPhLANeH5OLx8Aa/bN8aLL3R8AgRbg5tGzAYqIWEcf5QA6kQkU7R0LjD1qyakJR8OhXn9Pll3v/2fSws11CsGhNXCvQqJlnd5vJilXfdvo1qchPdHSmScZXHH3+WoaFBtm4ZYWlxDrta5vAzL3HhxAlGNMGgrrPqivvQbTUwhc4Lr00yu1rj0nyV6dUGF5eb7OzLcWe3QU1KzJUiUlPZ36HwzFGbI4sW13RBvSooIVEpcfNWeHaulx9d6OWm5hQ3bqtyNq9TKku6e7upVMo0q2XiXb3QPUL34DD5iTOousHUWh1VwIFtA4wMdDG/tMaFxRi1UoG+tMMdd99JqnCOv/3Cd9kxlMAqVjh2pUnXSC9kuhjduok/+qtnMU07tFgy0tCpKGp0LJPiDZkWod5JoQS9A1K2sjrbsdG8dn7pdwcTaQ9fzxxpvbnXaBCohUeHFbcUqVsK1tLLBHRVZXmtSkxXuOHQONVKkbVSg66UIKc6VEgyPz/D9rFBzHgPP3niOWpLc1w+e57XXniNucvTHDtxkSMnLrGYb5JvqNjVMl3LZag2qVl1urp7mFh12Ll7Dz39A65uUDrFoViTHd1JCgsr0LC8fFhhXwe8NmcxmVe5c7fgbO9mZsoZzp1dZLrWwSsXJTf1T3HdgTqvnpfkFIXm0hXypRrxmIpmN7l04jhavYghHEoNFz4e7U6xtrLC1IUJ1GaZoXiJLZ0Smk20+grnXj/M8mqFruoqwjBIdiRJdGXIDA7z2slFvvvYGXLZeHDqNU2PCGT5axGWgfd/54RHxK6zEGFMxy0oqal2CyAIjYyRkaBfhOb9hr8OSxS2KO8iYKEoqopl2SSTOq+dnOP+N2+nULXoTDo0ixUyKZ3ujM6xBUmyPkX/2A6+/92neP7IZS4uN5lcrjOdt2lqSRIdPWiJLKVKnVSjgGYIak2Hcr3Gri3drNQ14oaBZZk4jmRxeoa3d0hyPR0Uzk+h6JpL/HAkqlDZnRHM5x2enNaIbx8gN9RH1VLZk5llbLDClMwxURllabGOLhxiw9sYGe7HlgqJTCfd/f1cunSFVDxGxhBUTQvLtJlpJliswfnpIlMLRaaXqpTXVtiRmuXpMwVuHnBIY6LGVBKdnSw5KTaN9fM7f/IkiYQeZVu1teg7Xvk3IHwAqq6FLLEvn0NbcC8i+s+BTmBYRpTQ5I8wmLCRTItCaAiT0hpZEhlvFvQluLlyo2ExM1/kZ++/lhdfuUJvzKQnG+Pi2XmGu2KcmCrTo64wvdxECJVY1zC5bJJcJsXY5hH27dlBpVph9soU3TFB3FDJl2s0pWBLt8FcXUdVNGrVCiXToWthkvdcO0ZhuYi5WkDRNK8KJjAdyWLJJKkJsrrGgxcEWneO6wYq9ChVZocrrJS2E29alCpLOOUmA9t38fZ/9SmuvfNN5AaHuXTpEolmiZgGc2UL0wajp5/+kUF0RaGru5O+kRGaIs5b9gtevViiX28wP7VKwRRo3b3U4xl2HtzBl775OsfOLJJKGgEJtMXOah1Fd1q7CGIsPyNwHAfNm+ROaJK4mwkoAdwfngEVCEWGWUAg12UF/sQtf4S5aPsbJxwYRkxOqNooJcmEwYlzS+zd3s2OXZu5dGkJ2ajQlVK4MlcGI8lMocbd+7Is5SV2vANVVSiWqpiWxeLSCtKxUXSVRn6ZjmyS+3/5k5w4dZlOUcFU4yiJTjKZJOcmJnlrzmb3zi0sHb8Q0Kb8nvmFUpOZuoNhaHRmNS7OLHL2eJktmwWH9ml8e7bO5fN7+Y2D51DiJV48q6EvnefEC89y9uUXOPvC09Tmp9E1hQtlh+nYOKV6hWxXN0ODvaiqSn9fN4VSjQcONjg/OU9M1uk0KxCLUTNSbNszQjPRSbFi88efe46e7iSW7bhWWBERRNaf6ehPBAtPbQnmNQm/P9IOYjIpnZByuIfReLiCQpu8uAyNIXe86pKqKl5KoQYkEBlWnhbR0THhFmVEdLK1I6GnK8kffvZ5RgYyZDZtpqzncAyDuhGnR6vSETN47tQyN+5NMj+/gmVDPOY2VFTKFWIxnUwuy1rD4pq77+SuO69DQXBlPs/OwQRnT5/h+OlzdJZWufPAFlaml8G0WumUgLlSk3xDktEVxrti/MGFMmPbdf5w5yon/vkSf/rYKtVyB7l4Csess22zwft/4QMMXnMDy2tl5mZnWS7VKMcyTNox5mI7uXbPIF/4i19l02AHzz//GvNz8xw/fZk3jxW5fGUeUygMOVXiGli2w87tvZydtxgf7+f3/+Jpctk4thNNuaOaQEQGY8gQ8OZjLtKr0xDpPnNpfG43t+WxsLy4IZ1qScVGxsMLF+UTighx/wjm0OA3OIZiAlcsSka6U/yxMoon8yo9alSt4XDm/AI//6HrOXlxlYQuWFqpMD6soVYLNBI91BsmewccTswpZDs7SSaTNE0r2IQZUefmu24lruuU600unz7Jzp2bkY0GxWKJ29OCq7cOs3TqEoqmBkHRQtmk2HDRtQOjHXx1qcFRO84dux3uvc9gx/YEDwuF+mt1OpoVCienmJ6qQLaPgeF+ijPzOE1JcuceHvjkezk5XWW1ovML795GJpulK5ckX6xw/OwcP3tIo1ktcn6hwlWZGlVTcqGosPfgFqYKkkO37OH/fPY5Tl9YJp3SA2aS/xQDYQfRUj9zQpwL2jZIJOdvoXpBAQlEqx8zcAFtZttdXyWiQ9NSCgm9uRTetI6WC3CrVf6HiGByZThFkRIScY0zF1bQdME9d+1ibqHE+JZe1lZWSapN6oUKVZGhUCzx5qvSHJ8oMHdlDru4gFVYIS2rpBXJ04fPcO0Ne7jnnht48scv0aPVcEb2kp+Y4aP7+ygvFDGLFVRVC05+oeEycfeOdfPd1SZPqw6peJzhZpVNCGxDcminSsPsZMtmnUdPL3Oh1E9hcY4Lpy/Q2Zsj05HAcGqcOXKCejGPape445arSCfjTE6vcNONBxmsncBpVjg2WeLO/gZxp8nAaAedPV1MFQXbr9rJE89O8NXvHqe3O+Wa/sBKiZYGUKQML0KlGhE09AR0MR8xVNRWldfXFmgD+tzJof7o2Aha1Boc5Q6OVFucQC//F+GuYOGJOHvBhR1SHPddgQjtSFd+DrIZg6demOTA7n7Gtg6Rz1eRlkNCg7PTJXJOCSXZxdRymbu22oyPDnLojrt4z7vuZE1VefrVswxYDU6fm+aWO27ixLkZsvU5Xrm8xg2Gxd6uDCuX5tB0HSEki2WTQt01hTs3d/F42eZ79SbDY2nyKw3u2y+4+RqDk6fraGctiglJ744B5uZsTJFltWqzWhfMrjVZrkoWVqtMLhQZTlj0dqVJ9/SRzmQp51eon3qchZUi08t1NmkN9HoZLW4wk9dYcJJs2buN2fkSv/t/n2KgN41ly+i8vxBrR3qmX4g2hcTgWcqW9GtIy9HdCErbgMkQRwBQ0xk/C/CjeBEIL/idJoE4QTA3WAneSFEUN5AIS9dFAsrQbo0MLvDMeDrGD544z123jNHZ08naahXbgQuzNbIxi3izQCqd4XJRY7jL4ZZrR9m87youz67x4+PHUboyLKwu8+qzx8CRnL40TbJs8+HNHVw6N0s8pqIosFA2WavZCGmzY7SLRwsW36rUGNqUwjEdREyFs1ViqzbDW2OceKpGUTR41/UFfviqzeZN2+ntyZJJJxgfG6JWb9I/PEyxVENp1nBQubBUZ/XSMUadyzz58iUq5QrbsxYmgsWGhsj1oHV10z26CYnCv/30w3R3JluZk4fetWR5oodIyighxxf1kBEWV6uuoyqqpybiCVOFO7m9tVTT6fSn26eDtEQHW26gFWiogfCAIoSnWtliJ0b6B0W7kqdY36UKJOI63330NPfesZXOng6seoNtY92s5askMwnyi0sMpASX1xQOv/gKvXKZ2dUqT568hJbVUDviLNULXFlboTBn84u9OlqlSs2WxDSVy8UmS1WLLkMwONjBP642+ZFlMjSSwjbdQMuIK1ypC549H+f1xRSvFGDHdpUDe1J87Uc1LFvHMFz6ejzhqpOUK3XSzTwdhoaoN9kz2GC8o8GPX55k+3CMDqps7RWMjuVQYwnqIkZuZIRYIsYv/db3yaQNb9gVkaJbRPYtdIIJFs4fDCmjk968WEvxSCN+A4gIYzhttDAxMDAgI0BDqBbgK0uHJ4epquKNMydCPpRtUyqEIgI16yjPoF1poIVVF4t1/uR37qMrG2fxwiRZo8HExCLCqpMVNUpWAr17kIrp0JuUNFMJfnipzMRSjWzaoLRk8e8MlYMZQdl2Gb0ThSaFpsOunE4lEefrq00m0oK+njiW6WsMuVxwu+5gr6SQWppmtcBYaonRIZULq90Mj4wxPz9PuVxBUzU6OjIsrVVJ1EvcuSPO/p0pLi/XuXBhjoxqEoupHNiVpVmzqJkal4sq267dj5Twr//LD0gk3OHWPlMnTMAVodxZeA0ctmW3LLBfj/EIuS3I3l0H1aPPBwrioqUOAtHycGhgRDQNbLWFC5x17WIyEiSGS42BJo2v/ysJpV9tZBOf1OhZk1hc5zuPnOGqfYPs2DvG5ckVrt7Tx+xKnXRMpWSCXlphqCPOKglSmspVPQabOgxeOVvibbbkrX0x8k2HYt3iQqFJUlcYz8U4bKt8odCk3GfQnTPcnoSQT5SORDEEttOARgORaKArKucnNXpHhuntTIFQ6e/vpaO7G0WLoZsF3nIwznC/yskreRaXy2wSBXKdCfREjNNXGmzaOsqCleHAjftZXKnw7z79CMmkjqopBB1cQTbemgQansy+EUXfd89ukO32B8rQsOhAE8DHbzy3IR0ZHOxgYIS/oH6030r5ZEQO3rbtYKCUYINZtKLFGwhGm0dSEsVTlPCriiEdQuH7N1hYKvGLHz7E+962hxeeOcVI1mZtcYWVlSpmpcK2VJ2mnqSc6mGtoWApGtPHCvyrnjhLNRO7aaJKIGYwheCRUpMrCZWunhiKdGdECCXcpt8atiQUgZAOpgOjDUFa6lyqgu0o1EzoyCS4bkcH14wKnPIKlxbKzC6V2dmrMNqlElNtkkmD45MNLCOD0dXD/mt38sOnzvPHn3ue3u5Ui0ofiGuuH9gRsHbCgh2h2oxoP1SSqLBnKHWM9HK2DfoWg0ODMtpsYLspk/dC27ZamjZto0eDExzRohMh0YcNIlDRBjvTpoCJKx+3uFzhtutH+PVP3cTlSytUF+bYOWRw4uIK6fISEpVjszUOjmf4/hmNX0gmiFsmc1WLghQsaDrHHYdXTRsjp5NLKAhHogYCTXK9vKgnP+hIMKVEVG16m6BK2DqYYHNfjKTh8vhnV2vULYfa8iqdCYW6DbvH0wgjTqUmuLTisGXvVsa39/HZf3iVh588T39PCumRZ4MeCxkdxBGpwnqWKQoJO63N4KXcYRKvE/DDwsLdUY5AxGIPDQ3J9mZCHzcWQd+5N0lYCSNU0WCv/UFGYWURoZK1fizecMa5pinkCw2yaZ3/8mu3smmkk1OvXWAk26SSL5BfKvHqRIVbtxu8NqeSzMfpSRqcq5tMmhadHYLrNsfQYwpzZZuZqmSlKSnbULcktvewbP8Buy2PaEBcE3QnFfpSKmnHps9Q6Erp6LrCzFIVpVHFsiU7BmM019ZIZjPUHJW5EoyM9rNiJ9h39TjL+Rq/92dPMT1fpKcrSUgbIqiPhKt5QhGRyN0/ta0o3wlpObcgYJ+vKdp0gMLplvC0fluzhF3dATE4NCRpQ5ICClFokkg4fZBhtRC/0CNEdKCEWD+wMIo2bMAxFESmW6qqgmU5rOWrvPPe7XzigUPMz64xP3GFLR0uXv7669PsG1A5U1S5vNBkUK8z2JOkbOtYRpq46tComWSTKg3TpG5JHEVStQV1qRDX3BOnAopjk9QEuqpiqAJFSjRdZblUZ63UJBHT6FDqZGgwV4/TkZSk4wJHqlxeAz3Xxej2ITaND/LdR8/wpW++TjKpk4jrgXZ/kC3LaDNnIOQQUoMWsiW45bN9AuvgawX6WUAABIkIwTc69ZV1+sGRDSA8KZhI6hFOMfzeAMcJ4MQweyjsy9u/l0gUlNbrIwrj0ViilVm4O1VTBWuFOumkxid+7mpuODTKpXPzOIU8HYaNWSmgmA1evlBhU7zJ5m6Fs6sCXVHJJtxh10ZnlvNzDXRVw2w2aTYtcmmDob4sq2sVDFViWja2I5hbrZNOCKRu0JdSycgCNRFjeq7OSM6Vv0tmEpydM+nq66XQEHQO9zO+rY9T55b4u68d4cpsgZ7uVFvTlYw03AYFHSEiNK9A6UOELEXUS25odoUQP1XLe93sR3CRwDDU7I+A9/N3X+xJ0Eo/ZJvipAipfUfGy4Y2g99OHpZ299Oa9QYhBBm5hB/SKXfGz+PPTPDa67Ps2tnPjn1j1DAo1AU2gmt291K0VUzLYaQrzsVli80Zi61dkmKhwpWVBmNdNo5Vp25JTFOypUvQWF2lJ2FzZdUmI5p0iDqpGCwWLboTFh0ZhZEeHUczmCrFyFfhwpxJ92A/3VuG2XFwG9WG5LNfPMxXvn0chPAKO3I9H0+uU9ZvqaV7k8DC+sxhmk4LiGu9xm8Pi1iO9sFHQrRdRyg4Hxoakut2ZxjF897YsZ2WauhGBMP1rj76HhFrERaYUiImy8e/12n/q27OrKiCet2iWKqzdbSDd7x5F1cfGKJRb7I2twK1CquLebKqRTatsZRvkHUqpGMK00XJ5ZUae0ZipFMxzs07KMU8N22No8VjnFqCS3NlxnsNckmFpqpzZhE2dWpYDQtLMege7KYqYvSO9NDVl+P8pTwPPXaWl4/MEI9rZNIxLNtZdwzbO6/CMlGibVBnUFsJ8JWocEH4e7nBHJXoJpAREe4wz1MIEEPDQxL5xgsavjBEFLxZJ1HOOpffVp0Sb3jBQbtSOPPwrVJEs1AGfW71hkWpXKenM8GtN27hxkOjDPVlEaZFfrVAo1jArFWZns2TVh22DCWpC53ZuQKq2aCrI+HyBxpVrGqNeMLATqZZzddZWaoQj+v0DfchjBhqLMamLQOoiRiFcp2Xj07z9POTXLqyRjyuk/ZIHEGb/BtMNCFCpGoLvyJTQJ31J4r1at/R8bChCNsP+ELHK+RVWrO3hoeH3X0Wys3XE0Q3qDbKn+JkRHiK7Xote3eosS92HC5uEJmB7rNZ3GDQjqaOgS6uwLIdyuUmjiMZHkhz1Z5Brto/yMhghpimoOJQqzYol+uolkm90aRaMdFwFbmqpqRat4jHDdIpA0XTELpBIpMgHo9jA4srVY6dXebI8WkuXlqlVrdIpw0SMVe3xx+hK9v5eP6Uz1AKHSq9bNiK31IBb1fjC0XJol2bIUIDCmEG4SxetPcIuRsgMlo8wJnXm/L2NC6CYQR/GwyPi5BEZHh4gTfOzLYs1u8uGSGXIlzh44AKHbqAsJitorp4WtO0qdSa2JZDOqUzMpRjbFMXI8M5+rvTZFIGqaThKZWB5UhUxb1205FUayblisVKvsrUdJ5LU2tMzeZZXashpSSRMEjEtWDMLm2Kn+HZPO3jWiIBc8RHhHoq2k1+++mXGz6yDQO8N063QltgeGRY+pGmqqioioJlu/KiMoQktQcVIlw/EMHQ8hYlPGSuwxO9wpiHf0EueugEHTBBgBglHSM2FMsU0TY4b3KW8JRFm6ZNo2F54+fcHgXDUDF0FVV1N447jdTBNG2aTRvTclzSqCIwDJVYTHPnFgtPBOONLKCIHM7g2fkHauMm3NBcBtFamnZih9wgZY4+DPHGZjlaBQ7J+Ag0EVKeDNfyW90oEVnaSDFBhoiHrfw2+kGRryUb+g+3cBTdnW0iV2ghHlwUio4OUpYSb76O+waGrhI31KCELUOQqWn7ctqur9Q1DUPXIpp6/jwe25Zv/IzFxh4ymGcg2xYhrO4RcZkbH/JWoc1ZNxpFSOEdVPlT0r7o0LDwDEGt/aYCxcnQ9On2u43gzREcNbpDFaEEY01FSI1KyOh727Yd/UwZ2W8tcWxHhoYlRk+aFBvbSem4qkcSO+RC/OCybeaQt1hywyBHghPC5MUGxzlsrSIimiKqvxhquZMbQKlig9MfYVmHPKRc54P4l6YrRiqwipAicjMtECgazIg2OpiIDhBt/T40t97x6GL+cVC8bhXaKlrr0ssNvnHlUUOaN7IdEGnjxXnS7G+4QjIUK0l+ynDKDaZSsn5gVeR0y+gpDmI+KdeNPRQRCf71pn995E/Ixf6URRcbdAS3FZAEoMjIUVu/U8JiEaLdGYuwkEQLuFGCmbVEFCrCUyw1TUPTtLbNH1JtlG88HHI90PFTlLLlunArLKLZ5kp+Gny20aq1rYRs9fC14HBCmj2CaH60wRdyAyXxCEV/o+3Rmh6yTjA6dPFStrkxb85g6wOCwC00CijSSRJSn0JEiKJB3VpKHO+0Shkt+kjZek9VVVqn+qf41Yi18X29lO3H51/8T4Z19iPrJv+F5RcRVxORSJMbFMHaJle1sj6x4fXIUIdveNOE94MIzQducavanr1vFUPPTvPayGSE7RUNpkNzA72mgpB584MxV6VaiW5S35eHGEG+b46OGZStoQahh9O0zJCsuYyMfW23eRIRihHkxqjiGx178UYhsVy3pryB74+qpLWkV4NFdJyIsAahyVxhzoF8I9MmCY3jEeuvIDwOPsxdEK3fyw0yCyLuVW4wQJrW1DDCqpIBf0wG5UNfW1ZKufEiefwzX08ofLAF0Y0Rrn8H4lM+7t22qAHVSco2Ey6i40Nl+LX/kmVojVxZ5+6idjS6qOGFDcU6QS8lraaa1jQcEX3wgjd2tW34/xsHc1Gr5bsW2bYmlm1hW1aAmEixfscHjEQZgRGJUI+bpollWS3euYfZI6JUcn9xHcftCFaVKISLaJ0Gn1coI4vbUsKW7UOPA7mTDSBHEY0DxAY4eXtJdCP9nNaltihVEW52QHRpDdeQRGsnkVgxHN7JMFLY6tBFrA8qpQjPUnqDZCMcYSrr97uqqtTrdcbGxxkf30qtXm/1DYbCBQ3f94ccSIQA5CtPhbtQFHdCqG3ZIQHj6N4MdxY7SAxdR0ongHQ1Q8Nsmi2igpAbV83WLXgLEZRe4Nh69nL9ppRttXDBOmQpQsgMK2vhXleY3xD04YW0EINrERttJ9Eq97ad+JY3b222lqsVLYk+xR35Ek4aw1PS23BUEJJqtYphGNxxx+1I4OKFC1QqVXdwZYiJJkZHR+W6KlWoJs8GAUiYEhbMCgqZILEBYNLue9ybcqJTr8LgbvhzZRsS9kbEEvn/HhSGsxjWzdSVHvVatvrwnZb2bpgF3UbLayGlbYdKhnZhmOPfzrlof9aGYaAoglqt1mIDyfbnGgVgpJSMj49x++130NPTg5SSlZUVnn7maS5NXHLJO97LtQ2LNSHTFsazfc6ZoI3wgYhSveRPXx0pJdKWkaHV672cbLl4n9YUqmU73vyCdQl3mFQpfzpC2nLu7Zs1SnVz7FYK2xqkJaPAmI+mCr9yKtcFnJENLMJZVogb2NaPIYFGvRHy8zJaW3jDXS9Yn7Wvx9JdC8BGUGJ0wcLpYCsg3OBkh8DZcJAWRrTaf9fu0mUbU7b9ZLRk7cIDLPmppdN1UG07wN7Gpt1oCFPEPYRr8jLqPtutkWyrXATxVlBmjyKjwYFSRASdjZiPKBYf5pAhBNRqdXTd4OOf+DgC+PznP49pmoEL8O9Vi9xYpMS3nosuQ4CQm5q1etIiuvRsAG4oIeClrSARwbXd4aaRHLidUCEiJ/2nADYbHnrZCukjiVqoDLuR+wqi//X196gVlm2gTFSkiVBs1RYfr1PvCDqAaYk/xOJx6vV6KBuRkWfqX1YqlaJSqfCTp58GITBNk1QqFTSLBMjk6ObNkg0kRPl/5JS9Ib9TiI3NUzhH2gByCdKgDd5fBgBI2+v+H3z/RoZBvCGKICPBZTg1bb+38ABNwnFTWxlTsJ4XGCHuOmHXSxRoE9Ho3s+gIlzAkFUJC3c50g7wAOmsB02VN2KaRODddnTupxw72Q7GvDFbZD17OVImWf/u7SQT2kQP/1+R3DdE/0KcBREpfsh1M4Ui6+v75/AihDmRYeq3UCKW1I1nWhY43PMfGRNLqybSDuzIkGuKCER7wWyYAhZVAhbu8Og3jNTb6tMyXAWTch1qJcL144jUTOiZKOINrI0IBU9RLtu6tFBsTI3a6ExHhCwFQTlVsEEZVbT58TZbEYk327X3RVRnWYTYOYEmTxj88mcvhM1xm5+XrMcUNgCo111TlHixQeUgqE56EIIQG5L5Qh+iRPLgVrFDRGVlpWirGPrYvYgqVBHCv0PvGbFPYUKpbL91ESmstIis66G2FgVLhHVLItf406G3dnk1EcUa2soDUS5Ea57CessXfeiRYhvrpXrfoAQUacCJtOGHZeE22NT+Y/n/ASS14tdFjUYXAAAAAElFTkSuQmCC";'
        'document.head.appendChild(l);})();</script>'
    )
    st.markdown(_fav_js, unsafe_allow_html=True)
    # Sidebar starts expanded so users immediately see Configure Your Classroom
    # (throb animation on the toggle button draws attention when collapsed)
    # ── LOGIN GATE ──────────────────────────────────────────────────
    if _login_required() and not _is_logged_in():
        _show_login_screen()
        st.stop()

    for k in ["chat_messages","students","conn_checked","conn_info","gen_result","grade_history"]:
        if k not in st.session_state: st.session_state[k]=[] if k in ("chat_messages","students","grade_history") else (False if k=="conn_checked" else None)
    # Always sync quiz state for ALL subjects (catches new subjects added after first run)
    for _qsk in QUIZ:
        _qk = f"qz_{_qsk}"
        if _qk not in st.session_state:
            st.session_state[_qk] = {"lv":"easy","qi":0,"sc":0,"tot":0,"stk":0,"done":False,"sel":None,"hist":[],"manual_lv":False}
    # === PENDING PROFILE LOAD — must run BEFORE any widgets render ===
    if st.session_state.get("_pending_load"):
        _lp=st.session_state.pop("_pending_load")
        # Restore text fields
        st.session_state["_school_confirmed"]=_lp.get("school","")
        st.session_state["_teacher_confirmed"]=_lp.get("teacher","")
        st.session_state["_phone_confirmed"]=_lp.get("phone","")
        # Restore dropdowns
        st.session_state["country_sel"]=_lp.get("country","Liberia")
        st.session_state["lang_sel"]=_lp.get("lang","English")
        if _lp.get("region"): st.session_state["cfg_region"]=_lp.get("region")
        if _lp.get("grade"): st.session_state["cfg_grade"]=_lp.get("grade")
        if _lp.get("subject"): st.session_state["cfg_subject"]=_lp.get("subject")
        if _lp.get("class_size"): st.session_state["cfg_clsz"]=_lp.get("class_size")
        if _lp.get("ability"): st.session_state["cfg_abl"]=_lp.get("ability")
        # Restore toggles
        st.session_state["moe_toggle"]=bool(_lp.get("moe_on",False))
        st.session_state["mano_toggle"]=bool(_lp.get("mano_on",False))
        # Restore main-area settings
        if _lp.get("task_cat"): st.session_state["task_cat"]=_lp.get("task_cat")
        if _lp.get("task"): st.session_state["task_sel"]=_lp.get("task")
        if _lp.get("agent"): st.session_state["agent_pick"]=_lp.get("agent")
        st.session_state.profile_set=True
    for sk in QUIZ: 
        k=f"qz_{sk}"
        if k not in st.session_state: st.session_state[k]={"lv":"easy","qi":0,"sc":0,"tot":0,"stk":0,"done":False,"sel":None,"hist":[],"manual_lv":False}

    if not st.session_state.conn_checked:
        _checking={"en":"⏳ Checking connection...","fr":"⏳ Vérification de la connexion...","sw":"⏳ Kuangalia muunganisho..."}.get(_lang_key(),"⏳ Checking connection...")
        with st.spinner(_checking):
            st.session_state.conn_info=check_conn(); st.session_state.conn_checked=True
    conn=st.session_state.conn_info; online=conn["online"] if conn else False

    # CSS with dark/light mode detection
    st.markdown(f"""<style>
    /* Collapse Streamlit header and tighten top padding */
    [data-testid="stHeader"] {{height:0 !important;min-height:0 !important;overflow:hidden !important;opacity:0 !important;padding:0 !important}}
    .block-container {{padding-top:0.5rem !important}}
    @import url('https://fonts.googleapis.com/css2?family=Source+Sans+Pro:wght@400;600;700&family=Playfair+Display:wght@600;700&display=swap');
    /* Dark mode (default) */
    :root {{
        --bg-main: {C_NAVY};
        --bg-card: {C_NAVY_L};
        --text-primary: #D0D8E8;
        --text-secondary: #9CA3AF;
        --text-muted: #667788;
        --border-color: #2a3a5a;
        --chat-user-bg: rgba(43,125,233,.08);
        --chat-user-border: rgba(43,125,233,.3);
        --chat-ai-bg: rgba(139,26,26,.1);
        --chat-ai-border: rgba(178,34,52,.3);
        --tip-color: #8899AA;
    }}
    /* Light mode detection */
    @media (prefers-color-scheme: light) {{
        :root {{
            --bg-main: #F8F9FA;
            --bg-card: #FFFFFF;
            --text-primary: #1a1a2e;
            --text-secondary: #4a5568;
            --text-muted: #6b7280;
            --border-color: #d1d5db;
            --chat-user-bg: rgba(43,125,233,.06);
            --chat-user-border: rgba(43,125,233,.25);
            --chat-ai-bg: rgba(139,26,26,.06);
            --chat-ai-border: rgba(178,34,52,.2);
            --tip-color: #555;
        }}
    }}
    /* Also detect Streamlit's own light theme */
    [data-testid="stAppViewContainer"][style*="background-color: rgb(255"] {{
        --bg-main: #F8F9FA;
        --bg-card: #FFFFFF;
        --text-primary: #1a1a2e;
        --text-secondary: #4a5568;
        --text-muted: #6b7280;
        --border-color: #d1d5db;
        --chat-user-bg: rgba(43,125,233,.06);
        --chat-user-border: rgba(43,125,233,.25);
        --chat-ai-bg: rgba(139,26,26,.06);
        --chat-ai-border: rgba(178,34,52,.2);
        --tip-color: #555;
    }}
    .stApp {{font-family:'Source Sans Pro',sans-serif}}

    section[data-testid="stSidebar"] {{background:linear-gradient(180deg,#4A0E0E 0%,{C_RED} 40%,#7B2020 100%) !important}}
    section[data-testid="stSidebar"] .stMarkdown h2 {{color:{C_GOLD_L} !important;font-family:'Playfair Display',serif}}
    section[data-testid="stSidebar"] .stMarkdown p,section[data-testid="stSidebar"] .stMarkdown li {{color:#F0D5D5}}
    section[data-testid="stSidebar"] label {{color:#F0D5D5 !important}}
    section[data-testid="stSidebar"] .stSelectbox > div > div {{background:#3D0C0C !important;color:#F0D5D5 !important;border-color:#8B3030 !important}}
    section[data-testid="stSidebar"] hr {{border-color:#8B3030 !important}}
    .stStatusWidget {{display:none !important}}
    .stTabs [data-baseweb="tab-list"] {{background:{C_NAVY_L};border-radius:8px;padding:4px}}
    .stTabs [aria-selected="true"] {{color:white !important;background:{C_BLUE} !important;border-radius:6px}}
    .stButton > button[kind="primary"] {{background:linear-gradient(135deg,{C_BLUE_D},{C_BLUE}) !important;color:white !important;font-weight:700 !important;border:none !important;border-radius:8px !important;box-shadow:0 3px 10px rgba(43,125,233,.35) !important;padding:8px 20px !important;transition:all .2s !important}}
    .stButton > button[kind="primary"]:hover {{box-shadow:0 5px 20px rgba(43,125,233,.5) !important;transform:translateY(-1px) !important}}
    .stButton > button[kind="secondary"], .stButton > button:not([kind="primary"]) {{background:var(--bg-card) !important;color:var(--text-primary) !important;font-weight:600 !important;border:2px solid var(--border-color) !important;border-radius:8px !important;box-shadow:0 2px 6px rgba(0,0,0,.12) !important;transition:all .2s !important}}
    .stButton > button[kind="secondary"]:hover, .stButton > button:not([kind="primary"]):hover {{border-color:{C_BLUE} !important;box-shadow:0 3px 12px rgba(43,125,233,.2) !important;transform:translateY(-1px) !important}}
    /* Status bar: clearly non-interactive */
    .status-bar {{background:transparent !important;border:1px dashed var(--border-color) !important;border-radius:20px;padding:5px 14px;font-size:.78rem;display:flex;align-items:center;flex-wrap:wrap;gap:4px;margin-bottom:.8rem;opacity:.75;pointer-events:none;user-select:none}}
    .rh {{background:linear-gradient(135deg,{C_RED},{C_RED_L});color:white;padding:1rem;border-radius:10px 10px 0 0;margin-top:1rem}}
    .rh h3 {{margin:0;color:white;font-family:'Playfair Display',serif;font-size:1.15rem}}
    .rb {{border:1px solid var(--border-color);border-top:none;border-radius:0 0 10px 10px;padding:1.2rem;background:var(--bg-card);color:var(--text-primary);line-height:1.7}}
    .ct {{background:var(--chat-user-bg);border:1px solid var(--chat-user-border);border-radius:12px;padding:12px 16px;margin:6px 0;color:var(--text-primary)}}
    .cp {{background:var(--chat-ai-bg);border:1px solid var(--chat-ai-border);border-radius:12px;padding:12px 16px;margin:6px 0;color:var(--text-primary)}}
    .qbox {{background:var(--bg-card);border:2px solid {C_BLUE};border-radius:14px;padding:18px;margin:10px 0;color:var(--text-primary)}}
    .qok {{background:rgba(76,175,80,.12);border:2px solid #4CAF50;border-radius:12px;padding:14px;margin:8px 0;color:var(--text-primary)}}
    .qno {{background:rgba(239,83,80,.1);border:2px solid #EF5350;border-radius:12px;padding:14px;margin:8px 0;color:var(--text-primary)}}
    .qsc {{background:rgba(212,168,67,.1);border:1px solid {C_GOLD};border-radius:10px;padding:10px 16px;display:inline-block;color:{C_GOLD}}}
    .qtip {{background:rgba(43,125,233,.08);border-left:4px solid {C_BLUE};border-radius:0 8px 8px 0;padding:10px 14px;margin:8px 0;font-size:.88rem;color:var(--text-secondary)}}
    .sc {{background:var(--bg-card);border:1px solid var(--border-color);border-radius:10px;padding:12px 16px;margin:6px 0;color:var(--text-primary)}}
    .ft {{text-align:center;color:var(--text-muted);font-size:.8rem;padding:1.5rem 0 1rem;border-top:1px solid var(--border-color);margin-top:2rem}}
    .ft a {{color:{C_GOLD};text-decoration:none}}

    /* Solid blue Configure expander */
    /* Mic recording styles — visible on dark backgrounds */
    .mic-wrapper {{ text-align:center }}
    .mic-wrapper [data-testid="stAudioInput"] {{ border:2px solid rgba(239,83,80,.4) !important;border-radius:12px !important;background:rgba(239,83,80,.06) !important }}
    .mic-wrapper [data-testid="stAudioInput"]:hover {{ border-color:rgba(239,83,80,.7) !important;background:rgba(239,83,80,.12) !important }}
    .mic-label {{ font-size:.72rem;color:var(--text-muted);margin-top:2px;text-align:center;font-weight:600 }}
    /* Photo upload area in chat */
    .photo-upload-area {{ background:rgba(43,125,233,.06);border:2px dashed rgba(43,125,233,.3);border-radius:12px;padding:12px;margin-bottom:8px }}
    .photo-upload-area:hover {{ border-color:rgba(43,125,233,.5);background:rgba(43,125,233,.1) }}
    [data-testid="stSidebar"] .stExpander {{ background:linear-gradient(135deg,#1a3a6b,#2B7DE9);border-radius:10px;border:none !important }}
    [data-testid="stSidebar"] .stExpander summary {{ color:white !important;font-weight:700 }}
    [data-testid="stSidebar"] .stExpander [data-testid="stExpanderDetails"] {{ background:rgba(0,0,0,.15);border-radius:0 0 10px 10px }}
    /* Sidebar text inputs — bright text, clear borders */
    section[data-testid="stSidebar"] .stTextInput > div > div > input {{color:#FFFFFF !important;background:#3D0C0C !important;border:1.5px solid #D4A843 !important;border-radius:8px !important;caret-color:#F5D998 !important}}
    section[data-testid="stSidebar"] .stTextInput > div > div > input::placeholder {{color:#C0A070 !important;opacity:1 !important}}
    section[data-testid="stSidebar"] .stTextArea > div > div > textarea {{color:#FFFFFF !important;background:#3D0C0C !important;border:1.5px solid #D4A843 !important;border-radius:8px !important}}
    @keyframes throb {{ 0%,100% {{ opacity:.3;transform:scale(1) }} 50% {{ opacity:1;transform:scale(1.2) }} }}
    [data-testid="collapsedControl"] {{ animation:throb 2s ease-in-out infinite }}
    [data-testid="collapsedControl"]:hover {{ animation:none;opacity:1 !important;transform:scale(1.1) }}
    /* Red pepper cursor on main content (non-red background) */
    .stSelectbox,.stCheckbox,.stButton>button,.stTextInput,.stExpander,[role="tab"],
    .stSelectbox > div,.stSelectbox [data-baseweb="select"],.stSelectbox [data-baseweb="select"] *,
    .stSelectbox svg,.stCheckbox label,.stCheckbox input,.stCheckbox span,
    .stButton>button *,.stExpander summary,.stExpander summary * {{
        cursor: pointer !important;
    }}
    /* Blue pepper cursor on red sidebar */
    [data-testid="stSidebar"] *,
    [data-testid="stSidebar"] .stSelectbox [data-baseweb="select"],
    [data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] *,
    [data-testid="stSidebar"] .stSelectbox svg,
    [data-testid="stSidebar"] .stCheckbox label,
    [data-testid="stSidebar"] .stButton>button,
    [data-testid="stSidebar"] .stExpander summary {{
        cursor: pointer !important;
    }}
    /* Reduce sidebar gaps */
    section[data-testid="stSidebar"] [data-testid="stVerticalBlockBorderWrapper"] {{ padding: 0 !important; margin: 0 !important }}
    section[data-testid="stSidebar"] hr {{ margin: 4px 0 !important }}
    section[data-testid="stSidebar"] .stSelectbox {{ margin-bottom: -6px !important }}
    section[data-testid="stSidebar"] .stCheckbox {{ margin-bottom: -8px !important }}
    section[data-testid="stSidebar"] [data-testid="stMarkdownContainer"] > p {{ margin-bottom: 0 !important }}
    section[data-testid="stSidebar"] .stExpander {{ margin-bottom: 2px !important }}

    [data-testid="stSidebar"] .stExpander {{ background:linear-gradient(135deg,#2B7DE9,#1D5CBF) !important;border-radius:10px !important;border:none !important }}
    [data-testid="stSidebar"] .stExpander summary {{ color:white !important;font-weight:700 !important }}
    [data-testid="stSidebar"] .stExpander [data-testid="stExpanderDetails"] {{ background:rgba(0,0,0,.15) !important;border-radius:0 0 10px 10px !important }}
    /* Hide Streamlit chevron on Configure expander — show Teacher Pehpeh logo */
    [data-testid="stSidebar"] .stExpander summary svg,
    [data-testid="stSidebar"] .stExpander summary > svg,
    [data-testid="stSidebar"] .stExpander details > summary svg {{
        display:none !important;
        visibility:hidden !important;
    }}
    [data-testid="stSidebar"] .stExpander summary::before {{
        content:"";
        display:inline-block;
        width:30px;
        height:30px;
        min-width:30px;
        background-image:url("data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAACgAAAAoCAYAAACM/rhtAAANJUlEQVR42o2YaYxd51nHf+97lnvP3e/cubPbM/HYnokbO3ZimTZN2oTQlip0U0Qr1BKIRCUQKl8QiA8IRUA+UT6gSkilEgIEQY0KFSFt1CpNs9o0bRYc24ln83j2O3fm7ttZXz7cdcYT1CvN3Kujd/mfZ/k/z/MXo6OjiiM+ovNfCUAphBDth53VQgiUUgPrBYojj7rz4COWKRSic6tCgWrfIQ/sFWJgAyAUUogjwQ3uEd1DD2PprBtc/2Hv0AV3+LdEgJSyt1sMAqJvJU1qB1+ga9U2AqSUB4F01rTPUXcCPcIo3bO6OABkF62Qoo29s6D3pqLvToXqrOsCaJ8npUCTEk0KuvdJKdug6Z/bB3zoJQ4+OGAYvX1xe5+UkiAIOgcOukQRKHXgbE1KAqXwAlAqwHMdPM/DssLomqBpu0ghME29vdYP+uBEPxb7AAdiaACzfthtR2dLP3yFAIWgYTuYmiRthYhFQhwbG2J8OMmV/11io1jn4okxai2ble0CdcdDNw2skIlSQS8FEf3k6D3vuEW0L0KLxWJPHY6HwYxCHICH5wcI1+PC6UnmprL81V/8NlqzxSOP3s8TT3yaW6s7XLtxm2/+5e/wm198CLdS56P3ztKstdjcLaKbRpsVgKAbSqL71WGNrkmUaru4l41daw7GXwee6/uYQnB2ZpTHHr3APedPcf36LaKpOKNjw1SLdVZWdtjPV0jFo8TiUa5fu8XDD59ndn4a6XjcN3+Mly5fp+gpAgFSiQ7ItkGEEKhA9Y0kRN+CPesNcJ7oxGLT84gbOk//yZcpFSucOjPD6HiGQq7Ewvu3CSfCbG5usbCwTDRpMjGRYm1lk3/7rytcOH+ava09lCb5/OOfZG56hMWbaxQrDVQnse7InE42CwX6UbzWoUAQikDBVDJGVBdkx4ewUjEKW3v8aGkDz68zHPVpbVY5ERP4nottu3hJwa5t8vUvzZJf/4CmGyYeiVPYK1NxPL72W4+ibIenv/08gTSQQnFHrnRjsVtJhJC9WtC1nBsEiJbDv3zrG9Rcn+LmHqmRFDcXlkiyT9Sr0mw4VH3IuQoVCtFoeKSGoxQsh7HdGmdiGi0jxrafolzTmRgZ4TOPXWJ/v8zffvNZLn+wjjQMUEEf1wDta7FY7Kk2yfZDTwqBGwRMJKPcOztGy/d56KHz/Pyt91lZfp9xlSPk1Xi76PFy0eOtpmRdRbhZ8WhFosiEwXu3GmzaEfKuImyXGXLz6BGT3Ybi2WdfwzJ0vvK1TxER8MbbC5gho1Pe6CXKHTTTc7MmsSsNvv4Hn2P+zHFef+ldfvj869SqW9wTqbBlw+tVkyA2yoX5E0wNj7C2dZXdUoFmw+bWRpkz4xlmMgI59Ekunj1N6b3XaPz8h0wP2+wMhzlx5gS53SJevcnpyQyr+1UMKbu02+fbWDR6IEmklNi2x7njWUYTFmfum8OKRdhcW2QuXGSpofGfOzaf+sQX+PPf+0PunT3L9YVXKDRXGctaVIs1TGnyxYujFMot5k9f5J7Zc4zO3c/I9Alyb/+UYxnJVhHyGyVOXzjFaDTMm+8uIwz9ADghBFq04+I+cYLmefzNXz9J3fP5zj88j2G4DLc2yPmS53Y9zszewx999Xdp2i4/eP2fWSq+iRGTbG06BJaOGbVZ27PZM2vsFN5meXWZY5lZVssO2emTbF99Fc0IENYw5y7Mcdf0GIs3b3Nrp4hpaN1qh0CgxePxHlFLKXE8j5mRJLrjks0kGBlPExSWSYQV/70PY6MnePobf0q5Uuf5n/w7H2xdoVYVuBWdSxWd4Z0i8ZxNYrVOyxGUQoJmq0i94hALpbmytMpkJo23u0BDGTTrkudeuMLlqytITbbBKXpd1EAWi15PpwKFYzu4js+Tj5/n3kSBX9RDvLLv8sdffpJPf+IR1je3Wd++xVh2nGrdg8s/YHxthf3beaJhg4btseabvHn3ac5/7GM8eOE+/CAgEomxu7PKT7/9Z2RHx/jOizmuL+6QSsfaDDPYJUG/HxRCoAjaDzWBGTYZnRpmNOFRcTwWHR0rpKMEuI5LIhnj4r2XOD41Q9TSyC8tk9vMYyZC3HB8NoSkls/x8Owprq0usF8qMJwdpVQu8Pff/w82RQRZKzF/KoMVCaPJfi+gUO3uSak2UXcrSLd90qRGq2UzNRIhRYvFKjSERBMS3YhgWBbCsanVqpjhMK1KhTfLDrHAINOS/MJMcn/QwEgnuL22wmce+zzvLS6xuniTZ15+gVytwlzCgqDETFpDN3UCfyA7BvoJ2W23BjsZ1UGaTYVIW4qqbuIGipQV4dqL3+fF7/4jTqtFNBIhnUhRLVeYnsyQuGuKhiF48IFz1DRJvdEik85y7v5LxHObOK+9wvzcOaywQUXp5BqKISsgEQvjH27HOumsd/rmHmSlFIEKQAiSEclOoU5NRUBKJhybJeXx/uWXKdZqnP/Vz7H4szdwGk2y5RJ6s4nmO1RefZWTpknR9Zg5e5a9rdu4b75OaPo0m3vbmIZBy1NIU5IICaywTqmk0BggadXuBw4QtUD0TSskYR0sDXwhMDWJcF3qKoKtG7xz4xo/Xl6i1GowEYny+4ZFM7fHXjTGlLAxbJf9sMVay2fte88i98s0S+9izIwiTBO8FsoH13YxdNkdguDQIKYPZk2/nW+3914AugRZdamVoRATFDd3Kdc9Rk7HCfYdajstxj7SYCMkGREmo76L1CX1mseO1Pnu332Lr2TD1ITOVRlie7uOtCRGOCBkBmhhHdvxEaJT5tTBcUAfbLHa7T5omsD3A2xfYkRChFouxdUmRlJjKBqi5gXUlWDKhJkpixABN9JRPLvETMFFR5JXinTE4ImYyVK5xeWhFFrRJlKts1tzGZuKED2h0fJ1qjUbrcuBh8ZS/ciRVQmEUmwXbLRpnWPDPtEtg0vH48SkTyppkAOiiRimprFdaDGRCZDjLm9ZwwytNWhWy6yPJ6n4AhImad9n2AyYGYmyWPCQwiaqC96vQKPpEtLbM46UsjegHeDBQRcrwDAkq5sVKsoiaznMT4apVxskxmcIRyIMh3ROVqscC8f4qCUIKjrrrXESVpncTEDz4wlWXIfFrTLX1orUyy1GhlIYmUkiuuLEkMKIWCzutvBcDyHbFgqCoMfHdwzuff8rDENjN19juSyJ65KzU4LEyWniQynOPXCB1HiGvXoLu1bGvWuSuyqKqFdjJO3iZgyyI4qoLnhgJsm5iTjNRkB65jihZIqpmQyTEZeSFufqzX1MQxKovmul1DrxqJC9EqcUqhODQaA6aAP+51oBkcySbe7x1ccfpKkLvGINrR5wwxQsx1K8tVHlnrjL7GKR/KpkJjPEzZU00xEdI5lhPJNgPAT5UgMMmM8qwiGdd9Z9crkKpqn3pI5+LoiBmaQznIuOOiBEe5gxDZ3tXIWR42PMRGxuXXuP5X2fny1WES2buYk0O+UmU4HLR7w6ftNnZRM2LEFoyiDqGcSSI+RLVYyhLIWGT8a9jb5/G3f8JP/0/CJS9AemwdGk1/4dliG607/sBIUVNnjupRXyiRkipscxPc+lU3HS2TiT8/OMxCUXgip2ucmmq3grqXF/NE+oUmOrBgu3buHHh/jClx7h7mSJcX+H+Kk5nvnJOvVaE13T2pUMhVCiN831JJd4PP5Uew6WPaLuykFCCAxdx265LGzWOHv+JOFmjqxWJzGW4l+vvMt9tRZz1QY3GgEvJMKsJzT2N3Qs4bLmjVBXYYa8Xfzta4yJHOHJu3jmjQLvXtsiFgkNjBkShUIKiRwQncT4+PgA/fRlr54VBeiaRrPpkE5HeeI3TjHm71Ld22O9FOVkPmAlrPFaWKcekuhK0XIlIzacGEoykXBJODnGhkxq8Um+90aed97bJGoZbTmlOwt3DHPQg7QBHimTDQ7zqt2Cua4PQvJrH5/hV2YNzEqZ7Z0SMiyo+7DfgGhYolyPdESSkD5hU0dPZfmgHOKFyxvs5qtEI2a7OThCG+x5sNtVT0xMHGhYDwfqILFrHRJttDyymRgXz2a5ezzEsG4jfQfbC0jpARVH4QgTR7dY2Q+4ulJlYSWPrglChobvq74m05XmDoubXaNNTk4pIQaljp6a1xahBnST7j4pJa7rYTsBYStEJhNhOGURtzSUAMdVlCo2+UKDSqWJQBEOt6kkUH0F67DlDitbQgjE1NRUz4I9pj7k8jvk3Y4bpBAEKsDzAny/zf+iS7aaQNcEutYeJVXAAWup/0cWHgw1vSdYDqibB1RURJuMgjsPCzq2NXQN0xAHLux26IHqlyjxYbr1wPUqUPi+D0KgaRpy8FbVHe05KDT2irfohwGHZOdAqf5foA60TXTVAiF6393nUtMOSM5+EPDrn/0s8/PzOI7dLnX99lBwyJMfIioJfhlBvzv4tJmja3oxYFHRLq9iMMYFq6urFAqFdk0+fvy4OkozvoNyBmLzl/6oD/drtwc9PGYCOI6DlBLDMPg/IfE62g0z2WQAAAAASUVORK5CYII=");
        background-size:cover;
        background-repeat:no-repeat;
        background-position:center;
        margin-right:8px;
        vertical-align:middle;
        border-radius:50%;
        flex-shrink:0;
    }}
    </style>""",unsafe_allow_html=True)

    # Sidebar (defined first so country is available for logo flag)
    with st.sidebar:
        # Country first - drives auto language
        country=st.selectbox(T("country"),COUNTRIES,key="country_sel")
        # Auto-detect language from country (user can still override)
        if "lang_auto_done" not in st.session_state: st.session_state.lang_auto_done=set()
        if country not in st.session_state.lang_auto_done:
            st.session_state.lang_auto_done.add(country)
            if country in FRANCOPHONE: st.session_state.lang_sel="Français"
            elif country in SWAHILI_COUNTRIES: st.session_state.lang_sel="Kiswahili"
            else: st.session_state.lang_sel="English"
        lang=st.selectbox("🌍 Language / Langue / Lugha",list(LANGS.keys()),key="lang_sel")
        st.markdown("---")
        _cls_word={"en":"Classroom","fr":"Classe","sw":"Darasa"}.get(_lang_key(),"Classroom")
        # Confirmed values + version counters for key-swapping (forces fresh empty widgets)
        for _ck in ["_school_confirmed","_teacher_confirmed","_phone_confirmed"]:
            if _ck not in st.session_state: st.session_state[_ck]=""
        for _vk in ["_school_v","_teacher_v","_phone_v"]:
            if _vk not in st.session_state: st.session_state[_vk]=0

        _sn = st.session_state["_school_confirmed"]
        classroom_label=f"{_sn} {_cls_word}" if _sn.strip() else T("my_classroom")
        _logo_b64=get_b64()
        if _logo_b64:
            _logo_html=f'<img src="data:image/png;base64,{_logo_b64}" style="height:36px;width:36px;vertical-align:middle;border-radius:50%;margin-right:8px;filter:drop-shadow(0 2px 6px rgba(212,168,67,.4))">'
        else:
            _logo_html=""
        st.markdown(f'<div style="display:flex;align-items:center;margin:8px 0 4px">{_logo_html}<span style="font-family:Playfair Display,serif;font-size:1.3rem;font-weight:700;color:#F5D998">{classroom_label}</span></div>',unsafe_allow_html=True)
        # School name — key changes each submit so widget resets blank
        _sv=st.session_state["_school_v"]
        _school_raw=st.text_input(T("school_name"),value="",placeholder=T("school_placeholder") if not _sn.strip() else f"✅ {_sn} (type new name to change)",key=f"_si_{_sv}",label_visibility="collapsed")
        if _school_raw.strip():
            st.session_state["_school_confirmed"]=_school_raw.strip()
            st.session_state["_school_v"]=_sv+1
            st.rerun()
        school_name=st.session_state["_school_confirmed"]
        _tn_confirmed=st.session_state["_teacher_confirmed"]
        _tp_confirmed=st.session_state["_phone_confirmed"]
        _tn_col,_tp_col=st.columns([3,2])
        with _tn_col:
            _tv=st.session_state["_teacher_v"]
            _teacher_raw=st.text_input("👤 Teacher",value="",placeholder="e.g., Mr. Kollie" if not _tn_confirmed else f"✅ {_tn_confirmed}",key=f"_ti_{_tv}",label_visibility="collapsed")
            if _teacher_raw.strip():
                st.session_state["_teacher_confirmed"]=_teacher_raw.strip()
                st.session_state["_teacher_v"]=_tv+1
                st.rerun()
            teacher_name=st.session_state["_teacher_confirmed"]
        with _tp_col:
            _pv=st.session_state["_phone_v"]
            _phone_raw=st.text_input("📞 Phone",value="",placeholder="e.g., 0886-XXX-XXX" if not _tp_confirmed else f"✅ {_tp_confirmed}",key=f"_pi_{_pv}",label_visibility="collapsed")
            if _phone_raw.strip():
                st.session_state["_phone_confirmed"]=_phone_raw.strip()
                st.session_state["_phone_v"]=_pv+1
                st.rerun()
            teacher_phone=st.session_state["_phone_confirmed"]
        if not school_name.strip() and not teacher_name.strip():
            st.caption("✏️ Enter school & teacher name to personalize")
        elif not teacher_name.strip():
            st.caption("✏️ Add your name — parents will see it on letters")
        else:
            _t_display=f'<span style="color:#FFFFFF;font-size:1.05rem;font-weight:700">{teacher_name}</span>'
            if teacher_phone.strip(): _t_display+=f'<span style="color:#FFFFFF;font-size:1rem;font-weight:700;margin-left:8px">• {teacher_phone.strip()}</span>'
            st.markdown(f'<div style="background:linear-gradient(135deg,#3D0C0C,#5A1515);border:2px solid #D4A843;border-radius:10px;padding:10px 16px;margin:6px 0;box-shadow:0 2px 8px rgba(212,168,67,.2)">👤 {_t_display}</div>',unsafe_allow_html=True)
        if "profile_set" not in st.session_state: st.session_state.profile_set=False
        # All classroom settings inside collapsible Configure block
        with st.expander("Configure Your Classroom", expanded=False):
            if not st.session_state.profile_set:
                st.markdown('<div style="font-size:.85rem;color:#F0D5D5;margin-bottom:8px">Select your school setting, grade, subject, and preferences below. Save your profile to reuse later!</div>',unsafe_allow_html=True)
            region=st.selectbox(T("setting"),list(_regions().keys()),label_visibility="collapsed",format_func=lambda x: f"📍 Setting: {x}", help="Urban, rural, or remote — shapes the type of content generated",key="cfg_region")
            grade=st.selectbox(T("grade"),_grades(),label_visibility="collapsed",format_func=lambda x: f"🎓 Grade: {x}", help="The class level you are teaching",key="cfg_grade")
            subject=st.selectbox(T("subject"),_subjects(),label_visibility="collapsed",format_func=lambda x: f"📚 Subject: {x}", help="Choose the subject you want content for",key="cfg_subject")
            clsz=st.selectbox(T("class_size"),list(_sizes().keys()),label_visibility="collapsed",format_func=lambda x: f"👥 Class Size: {x}", help="Helps Teacher Pehpeh suggest realistic group sizes and activities",key="cfg_clsz")
            abl=st.selectbox(T("student_level"),list(_ability().keys()),label_visibility="collapsed",format_func=lambda x: f"📊 Student Level: {x}", help="Mixed, advanced, or struggling — adjusts difficulty and scaffolding",key="cfg_abl")
        # Map French display values back to English for AI
        _region_val=_regions()[region]
        _grade_en=_to_en_grade(grade)
        _subj_en=_to_en_subj(subject)
        _size_val=_sizes()[clsz]
        _res_val="standard resources"   # fixed default — resources config removed
        _abl_val=_ability()[abl]
        st.markdown("---"); 
        # MOE Curriculum Toggle — Liberia only
        if CURRICULA and CURRICULUM_AVAILABLE and country == "Liberia":
            moe_on = st.checkbox("🇱🇷 Align to MOE Curriculum", value=False, key="moe_toggle",
                                 help="When ON, topics and prompts align to Liberia Ministry of Education standards")
            curr_subjects = get_available_subjects(CURRICULA)
            if moe_on:
                st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:8px;'
                            f'padding:8px 12px;margin:4px 0;font-size:.8rem;color:#F0D5D5">'
                            f'📘 <strong>Aligned:</strong> {", ".join(curr_subjects)}<br>'
                            f'<span style="font-size:.72rem;opacity:.8">Lessons will match MOE learning objectives, content outlines, and assessment strategies</span></div>',
                            unsafe_allow_html=True)
        else:
            moe_on = False
        # Mano Language Toggle — Liberia + Rural only (hidden silently if data not deployed)
        mano_on = False
        if country == "Liberia" and _regions()[region] == "rural" and MANO_AVAILABLE:
            mano_on = st.checkbox("🗣️ Mano Language (Bilingual)", value=False, key="mano_toggle",
                                  help="When ON, lessons include Mano vocabulary, grammar, and cultural context for Nimba County")
            _mano_stats = get_mano_stats()
            if _mano_stats and mano_on:
                st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:8px;'
                            f'padding:8px 12px;margin:4px 0;font-size:.8rem;color:#F0D5D5">'
                            f'🗣️ <strong>Mano Library:</strong> {_mano_stats["total"]}+ words<br>'
                            f'<span style="font-size:.72rem;opacity:.8">{", ".join(_mano_stats["categories"][:5])}...</span></div>',
                            unsafe_allow_html=True)
        st.markdown("---")
        # Save / Load Classroom Profile
        # ── Save configuration ────────────────────────────────────────────
        _pf1,_pf2=st.columns(2)
        with _pf1:
            if st.button("💾 Save Configuration",use_container_width=True,key="sv_prof"):
                st.session_state.saved_profile={"school":school_name,"teacher":teacher_name,"phone":teacher_phone,"country":country,"lang":lang,"region":region,"grade":grade,"subject":subject,"class_size":clsz,"ability":abl,"moe_on":moe_on,"mano_on":mano_on,"task_cat":st.session_state.get("task_cat","📋 Planning"),"task":st.session_state.get("task_sel",""),"agent":st.session_state.get("agent_pick","")}
                st.session_state.profile_set=True
                st.session_state["_show_save_opts"]=True
                st.rerun()
        with _pf2:
            if st.button("📂 Load Configuration",use_container_width=True,key="ld_prof"):
                st.session_state["_show_load_opts"]=not st.session_state.get("_show_load_opts",False)
                st.session_state["_show_save_opts"]=False
                st.rerun()
        # Save panel — name the file then download
        if st.session_state.get("_show_save_opts") and "saved_profile" in st.session_state:
            st.success("✅ Configuration saved to this session!")
            st.markdown("<small style='color:#aaa'>Give your profile a name and download it to your drive so you can reload it anytime.</small>",unsafe_allow_html=True)
            _default_name=(school_name.strip().replace(" ","_") or "my_classroom")
            _prof_name=st.text_input("File name (no spaces):",value=_default_name,key="prof_filename").strip().replace(" ","_")
            if not _prof_name: _prof_name="teacher_pehpeh_profile"
            _pj=json.dumps(st.session_state.saved_profile,indent=2)
            st.download_button("📥 Download to my drive",data=_pj,
                file_name=f"{_prof_name}.json",mime="application/json",
                key="dl_prof",use_container_width=True)
            if st.button("✖ Close",key="close_save_opts",use_container_width=True):
                st.session_state["_show_save_opts"]=False
                st.rerun()
        # Load panel
        if st.session_state.get("_show_load_opts"):
            if "saved_profile" in st.session_state:
                if st.button("↩ Restore last saved profile",use_container_width=True,key="ld_restore"):
                    st.session_state["_pending_load"]=st.session_state.saved_profile
                    st.session_state["_show_load_opts"]=False
                    st.rerun()
            _up_prof=st.file_uploader("Or load from a saved file:",type=["json"],key="up_prof",label_visibility="visible")
            if _up_prof:
                try:
                    _loaded=json.loads(_up_prof.read().decode("utf-8"))
                    st.session_state.saved_profile=_loaded
                    st.session_state["_pending_load"]=_loaded
                    st.session_state["_show_load_opts"]=False
                    st.rerun()
                except Exception as e:
                    st.error(f"Invalid file: {e}")
        # ── Subscription status + logout ──────────────────────────
        if SUBSCRIPTION_ACTIVE:
            _badge_color = "#27AE60" if _SUB_DAYS_LEFT > 30 else "#F39C12"
            _badge_text = f"Active · {_SUB_DAYS_LEFT}d left" if _SUB_DAYS_LEFT <= 60 else "Active"
            st.markdown(f'<div style="text-align:center;padding:4px 0;font-size:.75rem;color:{_badge_color}">✅ Subscription: {_badge_text}</div>',unsafe_allow_html=True)
        elif _SUB_TIER == "expired":
            st.markdown('<div style="text-align:center;padding:4px 0;font-size:.75rem;color:#E74C3C">⚠️ Subscription expired — AI generation paused</div>',unsafe_allow_html=True)
        else:
            st.markdown('<div style="text-align:center;padding:4px 0;font-size:.75rem;color:#8899BB">🆓 Free account — AI generation not active</div>',unsafe_allow_html=True)
        if _login_required() and _is_logged_in():
            _login_lbl = st.session_state.get("_login_label","")
            _lcol1, _lcol2 = st.columns([2,1])
            with _lcol1:
                st.markdown(f'<p style="font-size:.75rem;color:#8899BB;margin:4px 0">👤 {_login_lbl}</p>',unsafe_allow_html=True)
            with _lcol2:
                if st.button("Sign out", key="logout_btn", use_container_width=True):
                    st.session_state["_logged_in"] = False
                    st.session_state["_login_label"] = ""
                    st.rerun()
        st.caption("© 2026 Institute of Basic Technology")
        st.markdown("[🌐 Visit our website](https://www.institutebasictechnology.org/index.php)")

    show_logo(country)
    _subtitle={"en":"Curating Personalized Content to Support Underresourced Teachers","fr":"Création de contenu personnalisé pour soutenir les enseignants sous-dotés","sw":"Kuunda Maudhui ya Kibinafsi Kusaidia Walimu Wasio na Rasilimali za Kutosha"}.get(_lang_key(),"Curating Personalized Content to Support Underresourced Teachers")
    st.markdown(f'<p style="text-align:center;color:#8899BB;font-size:.95rem;margin-bottom:.6rem">{_subtitle}<br>ChatGPT &bull; Claude &bull; Gemini</p>',unsafe_allow_html=True)

    # Unified status bar
    if conn:
        keys=sum([bool(OPENAI_API_KEY),bool(ANTHROPIC_API_KEY),bool(GOOGLE_API_KEY)])
        act=[n for k,n in [(OPENAI_API_KEY,"ChatGPT"),(ANTHROPIC_API_KEY,"Claude"),(GOOGLE_API_KEY,"Gemini")] if k]
        # Network: only show when there's a problem — device chip handles the "all good" state
        if conn["quality"]=="none":
            net_html=f'<span style="color:#EF5350" title="Teacher Pehpeh cannot reach the AI services right now">🔴 <strong>No Internet</strong></span>'
        elif conn["quality"] not in ("high","medium"):
            net_html=f'<span style="color:#FFB74D" title="Connection is slow — AI may take longer to respond">🟠 <strong>Slow Connection</strong></span>'
        else:
            net_html=""   # device chip already shows connected — no need to repeat
        # AI Agents section — friendly count + names
        if act:
            models_html=f'<span style="color:#7BB8F5" title="The AI assistants that will answer your requests">🤖 <strong>{len(act)} AI Agent{"s" if len(act)!=1 else ""}</strong>: {" · ".join(act)}</span>'
        else:
            models_html=f'<span style="color:#EF9A9A" title="No AI keys configured — check your settings">⚠️ No AI connected</span>'
        # Image section — simplified
        _has_img = bool(OPENAI_API_KEY or GOOGLE_API_KEY)
        img_html=f'<span style="color:#C9A0DC" title="Can generate diagrams and visual aids for your lessons">🎨 Image generation on</span>' if _has_img else ""
        # Divider
        div='<span style="color:#3a4a6a;margin:0 10px;font-size:1.1rem">│</span>'
        bar_parts=[m for m in [net_html, models_html] if m]
        if img_html: bar_parts.append(img_html)
        # Voice section — no key details
        if ELEVENLABS_API_KEY:
            bar_parts.append(f'<span style="color:#81D4A8" title="Text-to-speech — Teacher Pehpeh can read responses aloud">🔊 Voice on</span>')
        else:
            bar_parts.append('<span style="color:#EF9A9A" title="No voice key — audio playback unavailable">🔇 Voice off</span>')
        # MOE Curriculum section
        if moe_on and CURRICULA:
            curr_subjects = get_available_subjects(CURRICULA)
            moe_html=f'<span style="color:#81C784">📘 MOE: {", ".join(curr_subjects)}</span>'
            bar_parts.append(moe_html)
        # Mano language status
        if mano_on and MANO_AVAILABLE:
            _ms = get_mano_stats()
            mano_html=f'<span style="color:#FFB74D">🗣️ Mano: {_ms["total"]}+ words</span>' if _ms else ""
            if mano_html: bar_parts.append(mano_html)
        # Unified status bar — Re-check is the last segment inside the bar itself
        import streamlit.components.v1 as _conn_comp
        _status_bar_html = div.join(bar_parts)
        _conn_comp.html(f"""
<style>
.tp-bar {{
  font-family: "Source Sans Pro", sans-serif;
  font-size: .78rem;
  display: flex; align-items: center; flex-wrap: wrap; gap: 4px;
  padding: 6px 16px; border-radius: 20px;
  border: 1px dashed #2a3a5a; opacity: .85;
  margin-bottom: .3rem;
}}
.tp-bar span {{ white-space: nowrap; }}
.tp-div {{ color: #3a4a6a; margin: 0 8px; font-size: 1.1rem; }}
#dev-chip {{
  display:inline-flex; align-items:center; gap:5px;
  padding:2px 8px; border-radius:12px;
  background:rgba(0,0,0,.2); border:1px solid #2a3a5a;
  transition: all .3s;
}}
#recheck-btn {{
  display:inline-flex; align-items:center; gap:5px;
  padding:3px 10px; border-radius:12px; cursor:pointer;
  background:rgba(43,125,233,.15); border:1px solid rgba(43,125,233,.4);
  color:#7BB8F5; font-size:.78rem; font-family:inherit;
  transition:all .2s;
}}
#recheck-btn:hover {{ background:rgba(43,125,233,.3); border-color:rgba(43,125,233,.7); }}
</style>
<div class="tp-bar">
  <div id="dev-chip" title="Shows whether your device is connected to the internet">
    <span id="dev-icon">📱</span>
    <span id="dev-lbl" style="color:#8899aa">Checking device…</span>
  </div>
  <span class="tp-div">│</span>
  {_status_bar_html}
  <span class="tp-div">│</span>
  <button id="recheck-btn" title="Re-test connection to AI services" onclick="window.parent.postMessage({{type:'streamlit:setComponentValue',value:true}},'*')">🔄 Re-check</button>
</div>
<script>
(function(){{
  var icon=document.getElementById('dev-icon');
  var lbl=document.getElementById('dev-lbl');
  var chip=document.getElementById('dev-chip');
  function setDev(ok){{
    if(ok){{
      icon.textContent='💻';
      lbl.style.color='#81C784';
      lbl.textContent='Device connected';
      chip.style.border='1px solid #81C78444';
      chip.style.background='rgba(46,125,50,.15)';
    }}else{{
      icon.textContent='📵';
      lbl.style.color='#EF9A9A';
      lbl.textContent='Device offline';
      chip.style.border='1px solid #EF535044';
      chip.style.background='rgba(139,26,26,.2)';
    }}
  }}
  if(!navigator.onLine){{setDev(false);return;}}
  var t=Date.now();
  fetch('https://www.google.com/favicon.ico',{{mode:'no-cors',cache:'no-store'}})
    .then(function(){{setDev(true);}})
    .catch(function(){{setDev(false);}});
  window.addEventListener('offline',function(){{setDev(false);}});
  window.addEventListener('online',function(){{
    fetch('https://www.google.com/favicon.ico',{{mode:'no-cors',cache:'no-store'}})
      .then(function(){{setDev(true);}})
      .catch(function(){{setDev(false);}});
  }});
}})();
</script>
""", height=46, scrolling=False)
        if st.sidebar.button(T("recheck"), key="recheck_sidebar"):
            st.session_state.conn_checked = False
            st.rerun()

    # Compute general MOE curriculum context for chat (available across tabs)
    _chat_curr_ctx = ""
    if moe_on and CURRICULA:
        _grade_num_chat = int(''.join(c for c in _grade_en if c.isdigit()) or "10")
        _chat_curr_ctx = get_curriculum_summary(CURRICULA) if CURRICULUM_AVAILABLE else ""
        # Also get all topics for this subject/grade to include in chat context
        _chat_topics = get_grade_topics(CURRICULA, _subj_en, _grade_num_chat) if CURRICULUM_AVAILABLE else []
        if _chat_topics:
            _chat_curr_ctx += f"\n\nAvailable MOE topics for {_subj_en} Grade {_grade_num_chat}: {', '.join(_chat_topics)}"
    # Compute Mano language context for chat
    _chat_mano_ctx = ""
    if mano_on and MANO_AVAILABLE:
        _chat_mano_ctx = build_mano_prompt_context(_subj_en, _subj_en)

    if not conn:
        keys=sum([bool(OPENAI_API_KEY),bool(ANTHROPIC_API_KEY),bool(GOOGLE_API_KEY)])

    # Tabs
    if online and keys:
        # CSS tooltip layer for tabs (st.tabs has no native help= param)
        st.markdown("""
<style>
/* Tab tooltip system */
.stTabs [data-baseweb="tab-list"] { position: relative; }
.stTabs [data-baseweb="tab"] { position: relative; }

/* Generate tab */
.stTabs [data-baseweb="tab"]:nth-child(1)::after {
  content: "Create lesson plans, quizzes, activities & more";
  position:absolute; bottom:-32px; left:50%; transform:translateX(-50%);
  background:#1a2a4a; color:#D0D8E8; font-size:.72rem; white-space:nowrap;
  padding:3px 8px; border-radius:6px; border:1px solid #2a3a5a;
  pointer-events:none; opacity:0; transition:opacity .2s; z-index:999;
}
.stTabs [data-baseweb="tab"]:nth-child(1):hover::after { opacity:1; }

/* Chat tab */
.stTabs [data-baseweb="tab"]:nth-child(2)::after {
  content: "Ask Teacher Pehpeh anything about your subject";
  position:absolute; bottom:-32px; left:50%; transform:translateX(-50%);
  background:#1a2a4a; color:#D0D8E8; font-size:.72rem; white-space:nowrap;
  padding:3px 8px; border-radius:6px; border:1px solid #2a3a5a;
  pointer-events:none; opacity:0; transition:opacity .2s; z-index:999;
}
.stTabs [data-baseweb="tab"]:nth-child(2):hover::after { opacity:1; }

/* Quiz tab */
.stTabs [data-baseweb="tab"]:nth-child(3)::after {
  content: "Practice quizzes & WASSCE exam simulation — works offline";
  position:absolute; bottom:-32px; left:50%; transform:translateX(-50%);
  background:#1a2a4a; color:#D0D8E8; font-size:.72rem; white-space:nowrap;
  padding:3px 8px; border-radius:6px; border:1px solid #2a3a5a;
  pointer-events:none; opacity:0; transition:opacity .2s; z-index:999;
}
.stTabs [data-baseweb="tab"]:nth-child(3):hover::after { opacity:1; }

/* Students tab */
.stTabs [data-baseweb="tab"]:nth-child(4)::after {
  content: "Manage student profiles, risk analysis & parent letters";
  position:absolute; bottom:-32px; left:50%; transform:translateX(-50%);
  background:#1a2a4a; color:#D0D8E8; font-size:.72rem; white-space:nowrap;
  padding:3px 8px; border-radius:6px; border:1px solid #2a3a5a;
  pointer-events:none; opacity:0; transition:opacity .2s; z-index:999;
}
.stTabs [data-baseweb="tab"]:nth-child(4):hover::after { opacity:1; }
</style>
""", unsafe_allow_html=True)
        t1,t3,t4,t2=st.tabs([T("generate"),T("chat"),T("quiz"),T("students")])
    else: t1=t2=t3=None; t4=st.container()

    # TAB 1: GENERATE
    if t1:
     with t1:
        # ── Task category selector → filters task list ──
        _TASK_CATEGORIES = {
            "📋 Planning":     ["Lesson Plan","Weekly Scheme","Term Scheme","Strategy Guide"],
            "📝 Assessment":   ["Quiz (10 Q)","Quiz (20 Q)","WASSCE MCQ (50)","WASSCE Theory","BECE Exam","Rubric"],
            "🎯 Activities":   ["Homework","Group Activity","Reading Comprehension","No-Lab Practical","Educational Game","Illustrated Lesson (AI image)"],
            "📚 Study Support":["Study Notes","Remedial Material"],
        }
        # Tasks that need topic selector
        _NO_TOPIC_TASKS={"teaching strategies","term plan"}
        # Tasks that need time selector
        _NEEDS_TIME={"detailed lesson plan","homework with minimal resources","group activity","reading passage with questions","hands-on zero-cost activity","zero-cost teaching game","5-day scheme of work","lesson with AI-generated visual","catch-up material"}
        # Tasks that show options/extras checkboxes
        _NEEDS_OPTIONS={"detailed lesson plan","homework with minimal resources","group activity","hands-on zero-cost activity","5-day scheme of work","term plan","catch-up material","zero-cost teaching game","lesson with AI-generated visual","reading passage with questions","revision guide"}
        _NEEDS_MOE={"detailed lesson plan","homework with minimal resources","group activity","hands-on zero-cost activity","5-day scheme of work","term plan","catch-up material","zero-cost teaching game","lesson with AI-generated visual","reading passage with questions","revision guide","10-question quiz with answer key","20-question quiz","50 WASSCE-style MCQs","WASSCE theory questions","BECE-style exam"}
        _NEEDS_IMG={"detailed lesson plan","homework with minimal resources","group activity","hands-on zero-cost activity","lesson with AI-generated visual","reading passage with questions"}
        _WASSCE_TASKS={"50 WASSCE-style MCQs","WASSCE theory questions","BECE-style exam","10-question quiz with answer key","20-question quiz"}

        # ── Category pill buttons ──
        if "task_cat" not in st.session_state: st.session_state.task_cat="📋 Planning"
        _CAT_TIPS = {
            "📋 Planning":     "Lesson plans, weekly & term schemes, teaching strategies",
            "📝 Assessment":   "Quizzes, WASSCE/BECE exam prep, grading rubrics",
            "🎯 Activities":   "Homework, group work, reading comprehension, educational games",
            "📚 Study Support":"Revision notes and remedial catch-up material",
        }
        _cat_cols = st.columns(len(_TASK_CATEGORIES))
        for _ci, _cname in enumerate(_TASK_CATEGORIES):
            with _cat_cols[_ci]:
                _is_cat = st.session_state.task_cat == _cname
                if st.button(_cname, key=f"cat_{_ci}", use_container_width=True,
                             type="primary" if _is_cat else "secondary",
                             help=_CAT_TIPS.get(_cname,"")):
                    st.session_state.task_cat = _cname
                    st.session_state.pop("task_sel", None)
                    st.rerun()

        # Build filtered task list by matching English TASKS values to current language
        _cat_en_keys = _TASK_CATEGORIES[st.session_state.task_cat]   # English task names
        _cat_en_vals = {TASKS[k] for k in _cat_en_keys if k in TASKS}  # English task values
        _all_tasks = _tasks()   # display-language dict
        _cat_display_keys = [k for k, v in _all_tasks.items() if v in _cat_en_vals]

        c1, c2 = st.columns(2)
        with c1:
            task = st.selectbox(T("task"), _cat_display_keys, key="task_sel",
                                label_visibility="collapsed",
                                format_func=lambda x: f"\U0001f4dd Task: {x}", help="What you want Teacher Pehpeh to create for you")
        _task_val = _all_tasks.get(task, "detailed lesson plan")
        _IS_PARENT_LETTER = False   # Parent Letter moved to Students tab
        _show_time    = _task_val in _NEEDS_TIME
        _show_options = _task_val in _NEEDS_OPTIONS
        _show_moe     = _task_val in _NEEDS_MOE
        _show_img     = _task_val in _NEEDS_IMG
        _show_topic   = _task_val not in _NO_TOPIC_TASKS
        _show_wassce_guide = _task_val in _WASSCE_TASKS

        if _show_wassce_guide:
            st.markdown(f'<div style="background:linear-gradient(135deg,rgba(139,26,26,.3),rgba(212,168,67,.1));border:1px solid {C_GOLD}66;border-radius:10px;padding:10px 16px;margin:6px 0;display:flex;align-items:center;gap:10px"><span style="font-size:1.1rem">📋</span><span style="color:{C_GOLD};font-weight:700;font-size:.88rem">WASSCE Prep</span><span style="color:#D0D8E8;font-size:.82rem"> — Review answer sheet shading technique before generating</span></div>', unsafe_allow_html=True)
            if st.button("📋 Open Answer Sheet Guide", key="wassce_guide_btn", use_container_width=False):
                wassce_shading_modal()
        with c2:
            if _show_time:
                tm = st.selectbox(T("time"), _times(), label_visibility="collapsed",
                                  format_func=lambda x: f"\u23f1\ufe0f Time: {x}", help="How long your lesson or activity will run")
            else:
                tm = "N/A"
                st.markdown(f'<div style="background:var(--bg-card);border:1px solid var(--border-color);border-radius:8px;padding:10px 14px;color:var(--text-muted);font-size:.85rem;opacity:.5">\u23f1\ufe0f Time: not needed for this task</div>', unsafe_allow_html=True)

        _parent_letter_override=None
        _pl_delivery=None
        _pl_mom_edu=None
        if False:  # Parent Letter block retained here as dead code — actual UI is in Students tab
            st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:10px;padding:12px 16px;margin:8px 0"><strong style="color:{C_GOLD}">\u2709\ufe0f Parent Letter</strong><br><span style="font-size:.85rem;color:var(--text-secondary)">Adapts to each parent \u2014 audio for those who find reading difficult, text or email for others</span></div>',unsafe_allow_html=True)
            _pl1,_pl2=st.columns(2)
            with _pl1:
                if st.session_state.students:
                    _pl_student=st.selectbox("Student:",["--Select--"]+[s["name"] for s in st.session_state.students],key="pl_stu")
                else:
                    _pl_student=st.text_input("Student Name:",key="pl_stu_txt",placeholder="e.g., Janjay Kollie")
            with _pl2:
                _pl_concern=st.selectbox("Concern:",["Struggling academically","Frequent absences","Behavior issues","Needs extra support","Outstanding progress","Parent meeting request"],key="pl_concern")
            _pl_details=st.text_area("Specific details (optional):",key="pl_details",height=60,placeholder="e.g., Failing math quizzes...")
            topic=_pl_concern; _topic_en=_pl_concern
            _pl_name=_pl_student if (st.session_state.students and _pl_student!="--Select--") else (st.session_state.get("pl_stu_txt","") or "the student")
            _pl_info=""
            _pl_mom_edu="Unknown"
            if st.session_state.students and _pl_student!="--Select--":
                _sel=next((s for s in st.session_state.students if s["name"]==_pl_student),None)
                if _sel:
                    _pl_info=f'Profile: {_sel["sib"]}sib, Mom:{_sel["mom"]}, SM:{_sel["sm"]}, Works:{_sel["wk"]}'
                    _pl_mom_edu=_sel.get("mom","Unknown")
            _auto_idx=0 if _pl_mom_edu=="No HS" else 1
            # Build sensitivity context from student profile
            _sensitivity=""
            if st.session_state.students and _pl_student!="--Select--":
                _sel=next((s for s in st.session_state.students if s["name"]==_pl_student),None)
                if _sel:
                    _sib=_sel.get("sib","0-4")
                    _sm=_sel.get("sm","No")
                    _wk=_sel.get("wk","No")
                    _mom=_sel.get("mom","Unknown")
                    _factors=[]
                    if _sib in ["5-8","8+"]:
                        _factors.append(f"This child has {_sib} siblings — they likely help care for younger ones, cook, fetch water, or do household chores.")
                    if _sm=="Yes":
                        _factors.append("This is a single-mother household. The mother is carrying the full burden alone. Be extra respectful of her effort.")
                    if _wk=="Yes":
                        _factors.append("This child works after school (possibly selling on the street, helping at market, or doing farm work). They may have very little free time to study.")
                    if _mom=="No HS":
                        _factors.append("The mother has no high school education. She may feel embarrassed or powerless about academics. Never make her feel blamed.")
                    if _factors:
                        _sensitivity=("CRITICAL SENSITIVITY CONTEXT — READ BEFORE WRITING:\n"
                            + "\n".join(f"• {f}" for f in _factors)
                            + "\n\nBecause of this home situation:\n"
                            "- Do NOT assume the child is lazy or not trying. They may be exhausted from responsibilities.\n"
                            "- Do NOT say 'make sure they study every night' as if it's simple. Acknowledge the parent's hard work first.\n"
                            "- DO frame it as: 'I know things are not easy at home. Even small-small time for books — 15 or 20 minutes — can help.'\n"
                            "- DO suggest REALISTIC actions: 'If the child can sit with their book even while minding the smaller ones, that helps.'\n"
                            "- DO position the teacher as an ALLY, not an authority: 'Let us work together. I am here to help, not to add to your load.'\n"
                            "- DO NOT put the child in trouble. The goal is partnership, not punishment.\n")
            _delivery_opts=["\U0001f50a Voice Message (audio for phone)","\U0001f4f1 Text Message (short SMS)","\U0001f4e7 Email (detailed letter)"]
            _pl_delivery=st.radio("How to reach this parent:",_delivery_opts,index=_auto_idx,key="pl_delivery",horizontal=True)
            if _pl_mom_edu=="No HS" and "\U0001f50a" in _pl_delivery:
                st.markdown(f'<div style="background:rgba(46,125,50,.1);border-left:4px solid #4CAF50;padding:8px 12px;border-radius:6px;font-size:.82rem;margin:4px 0">\U0001f4a1 <strong>Voice message recommended</strong> \u2014 this parent may find reading difficult. A warm audio message they can play on their phone is more accessible.</div>',unsafe_allow_html=True)
            _close_loop=("CRITICAL: End with a specific, simple way for the parent to confirm they received this. "
                         "E.g. 'Please send [child] back with this note signed' or 'Please call the school' or "
                         "'Tell [student] to let me know you heard this.' One simple action.")
            _t_name=teacher_name.strip() if teacher_name.strip() else "the teacher"
            _t_phone=teacher_phone.strip() if teacher_phone.strip() else ""
            _t_sig=_t_name
            if _t_phone: _t_sig+=f" ({_t_phone})"
            _t_phone_line=f"Include teacher phone number {_t_phone} so the parent can call." if _t_phone else "Invite parent to come to the school."
            if "\U0001f50a" in _pl_delivery:
                _parent_letter_override=(f"Create a SHORT VOICE MESSAGE (under 45 seconds spoken) from {_t_name} (teacher) to parent of {_pl_name}.\n"
                    f"CONCERN: {_pl_concern}\nDETAILS: {_pl_details or 'General concern'}\nSTUDENT: {_pl_info or 'N/A'}\n{_sensitivity}"
                    f"SCHOOL: {school_name or 'Our school'} | {_subj_en} | {_grade_en}\n"
                    "RULES:\n"
                    "1. Write ONLY the spoken message. NO meta-commentary, NO key elements list, NO analysis, NO duration notes. Just the words to speak.\n"
                    "2. Natural, warm, conversational. Max 6-8 sentences.\n"
                    f"3. Introduce yourself as {_t_name} from {school_name or 'the school'}. Say something POSITIVE about the child first.\n"
                    "4. State concern simply. No education jargon. Use Liberian English expressions.\n"
                    "5. NEVER use the word 'spirit'. Say 'sharp', 'bright', 'hardworking' instead of 'bright spirit' or 'good spirit'.\n"
                    "6. This parent may NOT be able to help with schoolwork directly. Do NOT ask them to teach or tutor.\n"
                    "7. Focus on PARENTAL ACCOUNTABILITY: make sure child sits down to do homework each evening, goes to school on time, takes studies seriously.\n"
                    "8. If classmates or older children live nearby, encourage the parent to let the child study with them.\n"
                    f"9. {_t_phone_line}\n"
                    f"10. {_close_loop}\n"
                    "11. NO letter formatting. NO section headers. Just natural speech.\n"
                    "OUTPUT: Only the spoken message text. Nothing else.")
            elif "\U0001f4f1" in _pl_delivery:
                _parent_letter_override=(f"Create a SHORT TEXT MESSAGE (SMS, max 300 chars) from {_t_name} to parent of {_pl_name}.\n"
                    f"CONCERN: {_pl_concern}\nDETAILS: {_pl_details or 'General concern'}\nSTUDENT: {_pl_info or 'N/A'}\n{_sensitivity}"
                    f"SCHOOL: {school_name or 'Our school'} | {_subj_en} | {_grade_en}\n"
                    "RULES:\n"
                    "1. OUTPUT ONLY the text message itself. NO headers, NO titles, NO meta-commentary, NO format notes, NO character counts.\n"
                    "2. Max 300 characters total. Friendly but brief.\n"
                    "3. Greet, state concern in one sentence, one action parent can take, invite to call/visit.\n"
                    "4. NEVER use the word 'spirit'. Use Liberian English: 'your child doing fine-o', 'sharp', 'bright'.\n"
                    f"5. {_close_loop}\n"
                    f"6. Sign as {_t_sig}, {school_name or 'the school'}. Simple words only.\n"
                    "OUTPUT: Only the SMS text. Nothing else.")
            else:
                _parent_letter_override=(f"Write a letter from {_t_name} (teacher) to parent/guardian of {_pl_name}.\n"
                    f"CONCERN: {_pl_concern}\nDETAILS: {_pl_details or 'General concern'}\nSTUDENT: {_pl_info or 'N/A'}\n{_sensitivity}"
                    f"SCHOOL: {school_name or 'Our school'} | {_subj_en} | {_grade_en}\n"
                    "RULES:\n"
                    "1. OUTPUT ONLY the letter itself. NO analysis, NO key features table, NO 'why this works' section, NO commentary about the letter.\n"
                    "2. Simple reading level. Short sentences. No jargon. Warm and respectful. Use Liberian English expressions.\n"
                    "3. NEVER use the word 'spirit'. Say 'sharp', 'bright', 'hardworking', 'doing well' instead.\n"
                    "4. Include: concern, what teacher is doing, what parent can do at home, invitation to meet.\n"
                    f"5. ONE page max. Proper letter format with date, greeting, body, closing. Sign as {_t_sig}, {school_name or 'the school'}.\n"
                    f"6. {'Include phone number '+_t_phone+' in the closing so the parent can call. ' if _t_phone else ''}{_close_loop}\n"
                    "OUTPUT: Only the letter. Nothing before it, nothing after it.")

        else:
            if _show_topic:
                topic=st.selectbox(T("topic"),_get_topics(_subj_en),label_visibility="collapsed",format_func=lambda x: f"\U0001f4d6 Topic: {x}", help="The specific topic inside this subject")
                _topic_en=_to_en_topic(topic)
            else:
                topic=_subj_en; _topic_en=_subj_en
                st.markdown(f'<div style="background:var(--bg-card);border:1px solid var(--border-color);border-radius:8px;padding:10px 14px;color:var(--text-muted);font-size:.85rem;opacity:.5">\U0001f4d6 Topic: not needed for this task</div>', unsafe_allow_html=True)

        # === MOE Curriculum: override topics + show badge when aligned ===
        is_curriculum_aligned = False
        curriculum_details = None
        curriculum_context = ""
        if not _IS_PARENT_LETTER and _show_moe and _show_topic and moe_on and CURRICULA:
            _grade_num = int(''.join(c for c in _grade_en if c.isdigit()) or "10")
            curr_topics = get_grade_topics(CURRICULA, _subj_en, _grade_num)
            if curr_topics:
                is_curriculum_aligned = True
                topic = st.selectbox(f"📘 {T('topic')} (MOE Curriculum)", curr_topics, key="moe_topic")
                _topic_en = topic  # Curriculum topics are already in English
                st.markdown(f'<div style="display:inline-block;background:linear-gradient(135deg,#1B5E20,#388E3C);'
                            f'color:white;padding:3px 10px;border-radius:16px;font-size:.72rem;font-weight:600;margin:2px 0">'
                            f'✅ MOE Curriculum-Aligned</div>', unsafe_allow_html=True)
                # Get details for this topic
                curriculum_details = get_topic_details(CURRICULA, _subj_en, _grade_num, topic)
                curriculum_context = build_curriculum_context(CURRICULA, _subj_en, _grade_num, topic)
        # === MOE Curriculum Preview (expandable) ===
        if is_curriculum_aligned and curriculum_details:
            with st.expander("📋 View Ministry of Education Requirements for This Topic"):
                objectives = curriculum_details.get("specific_objectives", [])
                if objectives:
                    st.markdown(f"**Learning Objectives:**")
                    for obj in objectives:
                        st.markdown(f"- {obj}")
                content = curriculum_details.get("content_outline", [])
                if content:
                    st.markdown(f"**Content Outline:**")
                    for item in content:
                        st.markdown(f"- {item}")
                activities = curriculum_details.get("suggested_activities", [])
                if activities:
                    st.markdown(f"**Suggested Activities:**")
                    for act in activities:
                        st.markdown(f"- {act}")
                local_notes = curriculum_details.get("local_contextualization_notes", "")
                if local_notes:
                    st.markdown(f"**Liberian Contextualization:** {local_notes}")
                materials = curriculum_details.get("lab_materials", [])
                if materials:
                    st.markdown(f"**Lab Materials:** {', '.join(materials)}")

        # === Mano Language: vocabulary preview for selected topic ===
        _mano_prompt_ctx = ""
        if not _IS_PARENT_LETTER and _show_moe and mano_on and MANO_AVAILABLE:
            try:
                _mano_prompt_ctx = build_mano_prompt_context(_topic_en, _subj_en) or ""
            except Exception:
                _mano_prompt_ctx = ""
            try:
                _mano_preview = get_mano_preview(_topic_en, _subj_en)
            except TypeError:
                try:
                    _mano_preview = get_mano_preview(_topic_en)
                except Exception:
                    _mano_preview = []
            except Exception:
                _mano_preview = []
            if _mano_preview:
                st.markdown(f'<div style="display:inline-block;background:linear-gradient(135deg,#E65100,#F57C00);'
                            f'color:white;padding:3px 10px;border-radius:16px;font-size:.72rem;font-weight:600;margin:2px 0">'
                            f'🗣️ Mano Vocabulary Matched</div>', unsafe_allow_html=True)
                with st.expander(f"🗣️ Mano Vocabulary for This Topic ({len(_mano_preview)} words)"):
                    _mv_cols = st.columns(2)
                    for idx, entry in enumerate(_mano_preview):
                        with _mv_cols[idx % 2]:
                            if isinstance(entry, (list, tuple)) and len(entry) >= 2:
                                st.markdown(f"**{entry[0]}** → {entry[1]}")
                            else:
                                st.markdown(f"• {entry}")
                    st.markdown(f'<div style="font-size:.75rem;color:var(--text-muted);margin-top:8px">'
                                f'💡 These Mano words will be woven into the generated lesson for bilingual teaching.</div>',
                                unsafe_allow_html=True)

        # Reading Comprehension: show literature selector
        lit_book=None; lit_info=None; rc_mode=None
        if not _IS_PARENT_LETTER and _task_val=="reading passage with questions":
            st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:10px;padding:12px 16px;margin:8px 0"><strong style="color:{C_GOLD}">{T("lit_library")}</strong><br><span style="font-size:.85rem;color:var(--text-secondary)">{T("lit_desc")}</span></div>',unsafe_allow_html=True)
            lit_book=st.selectbox(T("select_book"),list(LITERATURE.keys()),key="lit_book")
            lit_info=LITERATURE.get(lit_book,{})
            if lit_book and lit_book!="Teacher's Own Selection":
                st.markdown(f'<div style="font-size:.85rem;color:var(--text-secondary);margin:4px 0">✍️ <strong>{lit_info.get("author","")}</strong> ({lit_info.get("origin","")}) · {lit_info.get("genre","")} · {lit_info.get("wassce","")}<br>Themes: {lit_info.get("themes","")}</div>',unsafe_allow_html=True)
            rc_mode=st.selectbox(T("comp_type"),[T("pass_short"),T("pass_fill"),T("pass_essay"),T("pass_mcq"),T("pass_vocab"),T("full_comp")],key="rc_mode")
        # WASSCE mode — lock out irrelevant options
        _IS_WASSCE = _task_val in {"50 WASSCE-style MCQs","WASSCE theory questions","BECE-style exam","10-question quiz with answer key","20-question quiz"}
        _WASSCE_VALID_EXTRAS = {"WASSCE alignment","Local examples"}  # only these make sense for WASSCE
        exs=[]; add_img=False
        if _IS_WASSCE:
            st.markdown(
                f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD}66;border-radius:8px;padding:8px 14px;margin:4px 0;font-size:.82rem">' +
                f'<strong style="color:{C_GOLD}">🎯 WASSCE Mode</strong> ' +
                f'<span style="color:#D0D8E8"> — Options limited to exam-relevant choices.</span></div>',
                unsafe_allow_html=True
            )
            with st.expander(f"⚙️ Options for {task} — add extras"):
                _wassce_extras_en = list(_WASSCE_VALID_EXTRAS)
                exs = [e for e in _wassce_extras_en if st.checkbox(e, key=f"wx_{e}")]
        elif _show_options:
            # Options label and available checkboxes adapt to the selected task
            _TASK_OPTIONS_MAP = {
                "detailed lesson plan":         ["differentiation","formative","takehome","wassce_align","local_ex","literacy","large_class","cross_curr","ai_visual"],
                "homework with minimal resources": ["differentiation","local_ex","literacy","large_class"],
                "group activity":               ["differentiation","formative","local_ex","large_class","cross_curr","ai_visual"],
                "hands-on zero-cost activity":  ["differentiation","local_ex","large_class"],
                "5-day scheme of work":         ["differentiation","formative","wassce_align","local_ex","literacy","cross_curr"],
                "term plan":                    ["wassce_align","local_ex","cross_curr"],
                "catch-up material":            ["differentiation","formative","local_ex","literacy"],
                "zero-cost teaching game":      ["differentiation","local_ex","large_class"],
                "lesson with AI-generated visual": ["differentiation","local_ex","ai_visual"],
                "reading passage with questions":  ["differentiation","formative","wassce_align","local_ex","literacy"],
                "revision guide":               ["differentiation","wassce_align","local_ex","literacy"],
            }
            _active_opt_keys = _TASK_OPTIONS_MAP.get(_task_val, ["differentiation","local_ex","wassce_align"])
            _ALL_EXTRAS_KEYS = ["differentiation","formative","takehome","wassce_align","local_ex","literacy","large_class","cross_curr","ai_visual"]
            with st.expander(f"⚙️ Options for {task} — add extras"):
                _extras_labels = [T(k) for k in _active_opt_keys]
                exs = [EXTRAS[_ALL_EXTRAS_KEYS.index(k)] for k, lbl in zip(_active_opt_keys, _extras_labels)
                       if st.checkbox(lbl, key=f"x_{k}")]
                if _show_img:
                    add_img = st.checkbox(T("include_img"), key="add_img", help=T("img_help"))
        # === AI Agent Selection ===
        _avail_agents=[]
        if OPENAI_API_KEY: _avail_agents.append("ChatGPT")
        if ANTHROPIC_API_KEY: _avail_agents.append("Claude")
        if GOOGLE_API_KEY: _avail_agents.append("Gemini")
        if _avail_agents:
            _n=len(_avail_agents)
            _all_label=f"Best answer — all {_n} agents" if _n>1 else f"{_avail_agents[0]}"
            _agent_opts=[_all_label]+[f"{a} only" for a in _avail_agents]
            _agent_sel=st.selectbox("AI:",_agent_opts,key="agent_pick",label_visibility="collapsed",
                                    help="All agents gives the best result but may take a little longer and use more data. If you're on a slow connection, picking one agent can be faster.")
            if "all" in _agent_sel or _agent_sel==_all_label and _n>1:
                _agent_pick=_avail_agents
            else:
                _agent_pick=[_agent_sel.replace(" only","")]
            # Soft advisory when all agents selected and more than one available
            if len(_agent_pick)>1:
                st.markdown(
                    f'<div style="background:rgba(212,168,67,.07);border-left:3px solid {C_GOLD}88;' +
                    f'border-radius:0 6px 6px 0;padding:6px 12px;margin:4px 0 2px;font-size:.78rem;color:#C8B06A">' +
                    f'<strong>Heads up:</strong> Using all {len(_agent_pick)} agents gives richer results — though it may take a little longer ' +
                    f'and could use more mobile data on slower connections. Just something to keep in mind!</div>',
                    unsafe_allow_html=True
                )
        else: _agent_pick=[]
        gen_col, clr_col = st.columns([3,1])
        with gen_col:
            gen_btn=st.button(T("gen_btn"),type="primary",use_container_width=True,key="gen")
        with clr_col:
            if st.button(T("clear"),use_container_width=True,key="gen_clr"):
                st.session_state.gen_result=None; st.rerun()
        if gen_btn:
            if not SUBSCRIPTION_ACTIVE:
                _show_upgrade_prompt("generate")
                st.stop()
            # Build prompt
            if _IS_PARENT_LETTER and _parent_letter_override:
                _tn=teacher_name.strip() if teacher_name.strip() else "the teacher"
                _is_rural=_region_val=="rural"
                _locale_guide=("Use Liberian English expressions the parent will recognize: "
                    "'your child doing fine-o', 'small-small', 'I beg you', 'let us put our heads together', 'the child can be somebody'. "
                    "Avoid the word 'spirit' entirely — do not say 'bright spirit', 'good spirit', 'has spirit', etc. "
                    "Say instead: 'your child is sharp', 'bright', 'hardworking', 'doing well', 'trying hard'. ")
                if _is_rural:
                    _locale_guide+=("This is a rural family — farm analogies are appropriate: "
                        "'education is like planting rice — you must tend it every day', "
                        "'the child who studies is like the farmer who clears the bush early'. "
                        "Reference walking to school, community support, studying by lamplight.")
                else:
                    _locale_guide+=("This is an urban family — use town-appropriate references: "
                        "school transport, study space at home, after-school time management, "
                        "neighborhood study groups. Avoid farm analogies.")
                sp=(f"You are {_tn}, a teacher at {school_name or 'a school'} in {country}. "
                    f"You are writing to a parent. Be warm, respectful, and use simple language. "
                    f"Do NOT refer to yourself as Teacher Pehpeh or any AI. You are a real teacher communicating with a real parent. "
                    f"{_locale_guide}")
                q=_parent_letter_override
            elif _task_val=="reading passage with questions" and lit_book and lit_book!="Teacher's Own Selection":
                rc_prompt=f"""Create a READING COMPREHENSION exercise using the novel "{lit_book}" by {lit_info.get('author','')}.

INSTRUCTIONS:
1. Write an ORIGINAL passage (250-400 words) that captures the style, themes, and characters of this novel. The passage should feel authentic to the book's setting and voice. Include dialogue where appropriate.
2. After the passage, create comprehension exercises in this format: {rc_mode}

For "Fill in the Blanks": Remove 8-12 key words from the passage and list them as a word bank. Students must fill in the correct word.
For "Short Answer Questions": Write 6-8 questions requiring 2-3 sentence answers about the passage's meaning, character motivation, literary devices, and themes.
For "Essay Prompt": Give 2 essay questions (150-250 words each) requiring analysis of themes: {lit_info.get('themes','')}.
For "MCQ": Write 10 multiple-choice questions with 4 options each.
For "Vocabulary Exercise": Identify 10 challenging words from the passage with context clues exercises.
For "Full Comprehension (All Types)": Include fill-in-the-blanks, 5 short answer questions, 1 essay prompt, and 5 MCQs.

3. Provide a COMPLETE ANSWER KEY at the end.
4. Add TEACHER'S GUIDE with: key themes to discuss, discussion starters, and how this connects to WASSCE Literature requirements.

Subject:{_subj_en}
Grade:{_grade_en}
Topic:{_topic_en}
Book context: {lit_info.get('genre','')} from {lit_info.get('origin','')}. Themes: {lit_info.get('themes','')}. {lit_info.get('wassce','')}."""
                sp=build_sys(_region_val,country,_grade_en,_subj_en,_task_val,_size_val,_res_val,LANGS[lang],_abl_val,tm,_topic_en,school_name,_mano_prompt_ctx)
                q=rc_prompt
            else:
                sp=build_sys(_region_val,country,_grade_en,_subj_en,_task_val,_size_val,_res_val,LANGS[lang],_abl_val,tm,_topic_en,school_name,_mano_prompt_ctx)
                q=f"Create {_task_val}.\nSubject:{_subj_en}\nGrade:{_grade_en}\nTopic:{_topic_en}\nIMMEDIATELY USABLE."
            if exs: q+="\n"+"; ".join(exs)
            # === MOE Curriculum: inject context into prompts ===
            if is_curriculum_aligned and curriculum_context:
                sp += f"\n\nCRITICAL — MINISTRY OF EDUCATION CURRICULUM DATA:\n{curriculum_context}\nYou MUST align content with these objectives, content outline, and activities."
                if curriculum_details:
                    objectives = curriculum_details.get("specific_objectives", [])
                    if objectives:
                        q += "\n\nMOE Learning Objectives (address ALL of these):\n" + "\n".join(f"  {i+1}. {o}" for i, o in enumerate(objectives))
                    local_notes = curriculum_details.get("local_contextualization_notes", "")
                    if local_notes:
                        q += f"\nContextualization: {local_notes}"
            want_img=add_img or "image" in task.lower() or "AI visual" in str(exs)
            keys=len(_agent_pick)
            rs={}; ph=st.empty(); s=0; tot=keys+(2 if want_img else 1)
            ph.markdown(pprog(0,tot,T("generating_content")),unsafe_allow_html=True)
            for k,fn,nm in [(OPENAI_API_KEY,ask_gpt,"ChatGPT"),(ANTHROPIC_API_KEY,ask_cl,"Claude"),(GOOGLE_API_KEY,ask_gem,"Gemini")]:
                if k and nm in _agent_pick:
                    s+=1; ph.markdown(pprog(s,tot,f"Asking {nm}..."),unsafe_allow_html=True)
                    rs[nm]=fn(sp,q) if nm!="Gemini" else fn(sp,q)
            s+=1; ph.markdown(pprog(s,tot,T("combining")),unsafe_allow_html=True)
            result=synth(sp,q,rs)
            if _IS_PARENT_LETTER: result=clean_parent_output(result)
            # Also clean individual model responses for parent letters
            if _IS_PARENT_LETTER:
                rs={k:clean_parent_output(v) if v and not str(v).startswith("⚠️") else v for k,v in rs.items()}
            img=None; img_src=None
            if want_img:
                s+=1; ph.markdown(pprog(s,tot,T("creating_img")),unsafe_allow_html=True)
                img,img_src=gen_image(f"{_subj_en}: {_topic_en} for {_grade_en} in {country}")
            ph.markdown(pprog(tot,tot,T("done")),unsafe_allow_html=True); time.sleep(.5); ph.empty()
            # Store in session state
            # Clear cached voice audio from previous generation
            for _vck in ["tts_audio_pl_voice","tts_audio_gen_pl_voice","tts_audio_gen_pl_sms","tts_audio_gen_pl_email","tts_audio_gen"]:
                st.session_state.pop(_vck,None)
            st.session_state.gen_result={"result":result,"task":task,"topic":topic,"img":img,"img_src":img_src,"rs":rs,"grade":grade,"subject":subject,"lit_book":lit_book,"moe_aligned":is_curriculum_aligned and not _IS_PARENT_LETTER,"mano_aligned":mano_on and MANO_AVAILABLE and not _IS_PARENT_LETTER,"pl_delivery":_pl_delivery if _IS_PARENT_LETTER else None,"is_parent_letter":_IS_PARENT_LETTER}
        # Display results from session state
        gr=st.session_state.gen_result
        if gr:
            lit_tag=f" — 📖 {gr.get('lit_book','')}" if gr.get('lit_book') and gr['lit_book']!="Teacher's Own Selection" else ""
            moe_tag=" · ✅ MOE Aligned" if gr.get('moe_aligned') else ""
            mano_tag=" · 🗣️ Mano" if gr.get('mano_aligned') else ""
            st.markdown(f'<div class="rh"><h3>{ico(20)} {gr["task"]} — {gr["topic"]}{lit_tag}{moe_tag}{mano_tag}</h3></div>',unsafe_allow_html=True)
            if gr.get("img"): st.image(gr["img"],caption=f"{gr['topic']} — Generated by {gr.get('img_src','')}",use_container_width=True)
            valid_rs={k:v for k,v in gr["rs"].items() if v and not str(v).startswith("⚠️")}
            if len(valid_rs)>1:
                _synth={"en":"Synthesized Response","fr":"Réponse synthétisée","sw":"Jibu lililochanganywa"}.get(_lang_key(),"Synthesized Response")
                _from={"en":"Combined from","fr":"Combiné de","sw":"Imechanganywa kutoka"}.get(_lang_key(),"Combined from")
                st.markdown(f'<div style="background:rgba(212,168,67,.1);border-left:4px solid {C_GOLD};padding:6px 12px;border-radius:6px;font-size:.82rem;color:{C_GOLD};margin-bottom:8px">🔀 <strong>{_synth}</strong> — {_from} {" + ".join(valid_rs.keys())}</div>',unsafe_allow_html=True)
            st.markdown(f'<div class="rb">{highlight_result(gr["result"])}</div>',unsafe_allow_html=True)

            # === INTERACTIVE ANSWER SHEET for MCQ tasks ===
            _MCQ_TASKS = {"50 WASSCE-style MCQs","10-question quiz with answer key","20-question quiz","BECE-style exam"}
            if gr.get("task") in _MCQ_TASKS:
                _parsed_qs = parse_mcq_for_sheet(gr["result"])
                if len(_parsed_qs) >= 2:
                    import streamlit.components.v1 as _comp
                    import json as _json
                    _qs_data = _json.dumps(_parsed_qs)
                    _n_qs = len(_parsed_qs)
                    _has_answers = any(q.get("a") is not None for q in _parsed_qs)
                    _sheet_html = _make_mcq_sheet_html(_qs_data, title=gr.get("topic",""), n=_n_qs)
                    with st.expander(f"📋 Interactive Answer Sheet ({_n_qs} questions)", expanded=True):
                        _comp.html(_sheet_html, height=max(560, 220 + _n_qs * 48), scrolling=True)

            if len(valid_rs)>1:
                _indiv={"en":"Individual Model Responses","fr":"Réponses individuelles des modèles","sw":"Majibu ya modeli binafsi"}.get(_lang_key(),"Individual Model Responses")
                st.markdown(f'<div style="font-size:.85rem;color:{C_GOLD};margin:12px 0 6px">📋 <strong>{_indiv}</strong></div>',unsafe_allow_html=True)
                for mname,mresp in valid_rs.items():
                    with st.expander(f"{'🟣' if mname=='Claude' else '🟢' if mname=='ChatGPT' else '🔵'} {mname}'s Response"):
                        st.markdown(mresp)
            elif len(valid_rs)==1:
                nm=list(valid_rs.keys())[0]
                st.markdown(f'<div style="font-size:.75rem;color:#667;margin-top:4px">Response by {nm} (only model available)</div>',unsafe_allow_html=True)
            # === DELIVERY HANDLING ===
            if gr.get("is_parent_letter") and gr.get("pl_delivery"):
                # Let teacher pick which response to send
                _send_opts={"Combined (Best)":gr["result"]}
                for _mn,_mr in valid_rs.items():
                    if _mr: _send_opts[_mn]=_mr
                if len(_send_opts)>1:
                    _send_pick=st.selectbox("Which version to send:",list(_send_opts.keys()),key="pl_pick_model")
                else:
                    _send_pick=list(_send_opts.keys())[0]
                _send_text=_send_opts[_send_pick]
                _pdl=gr["pl_delivery"]
                if "Voice" in _pdl:
                    _vk="tts_audio_pl_voice"
                    if _vk in st.session_state and st.session_state[_vk]:
                        _aud=st.session_state[_vk]
                        st.markdown(f'<div style="background:linear-gradient(135deg,#2E7D32,#1B5E20);border-radius:10px;padding:12px 16px;margin:6px 0;color:white"><strong>Voice Message</strong> &mdash; Send this audio to the parent via WhatsApp or play on their phone</div>',unsafe_allow_html=True)
                        _aud_data=base64.b64decode(_aud["b64"])
                        st.audio(_aud_data, format="audio/mp3", autoplay=True)
                        st.download_button("Download MP3 to Send",data=_aud_data,file_name="voice_message_for_parent.mp3",mime="audio/mp3",key="gen_pl_dl_mp3",use_container_width=True)
                    else:
                        st.markdown('<div style="background:linear-gradient(135deg,#2E7D32,#1B5E20);border-radius:10px;padding:12px 16px;margin:6px 0;color:white"><strong>Voice Message</strong> &mdash; Generate audio to send to the parent</div>',unsafe_allow_html=True)
                        if st.button("Generate Voice Message",type="primary",use_container_width=True,key="gen_pl_voice_btn"):
                            with st.spinner("Creating voice message..."):
                                b64,src=speak_elevenlabs(_send_text)
                            if b64:
                                st.session_state[_vk]={"b64":b64,"src":src}
                                st.rerun()
                            else:
                                st.error(f"Audio failed: {src}")
                elif "Text" in _pdl:
                    _sms_text=_send_text.strip()
                    st.markdown('<div style="background:linear-gradient(135deg,#1565C0,#0D47A1);border-radius:10px;padding:12px 16px;margin:8px 0;color:white"><strong>Text Message Ready</strong> &mdash; Copy and send via SMS or WhatsApp</div>',unsafe_allow_html=True)
                    st.code(_sms_text, language=None)
                    
                    email_result(_send_text, f"Parent Message: {gr['topic']} ({gr['grade']})", "gen_pl_sms")
                else:
                    st.markdown('<div style="background:linear-gradient(135deg,#4A148C,#6A1B9A);border-radius:10px;padding:12px 16px;margin:8px 0;color:white"><strong>Parent Letter Ready</strong> &mdash; Download, print, or email</div>',unsafe_allow_html=True)
                    st.download_button("Download Letter",data=_send_text,file_name="parent_letter.txt",key="gen_dl_pl")
                    email_result(_send_text, f"Parent Letter: {gr['topic']} ({gr['grade']}, {gr['subject']})", "gen_pl_em")
                    tts_player(_send_text, "gen_pl_email")
                st.markdown('<div style="background:rgba(255,152,0,.1);border:1px solid #FF9800;border-radius:8px;padding:10px 14px;margin:8px 0;font-size:.85rem"><strong>Close the Loop</strong> &mdash; Follow up in 3-5 days to confirm parent received your message. If no response, try a different delivery method or send message home with the student.</div>',unsafe_allow_html=True)
            else:
              # === Download format selector ===
              _QUIZ_TASKS = {"10-question quiz with answer key","20-question quiz","50 WASSCE-style MCQs","WASSCE theory questions","BECE-style exam"}
              _PPT_TASKS  = {"detailed lesson plan","5-day scheme of work","term plan","lesson with AI-generated visual"}
              _default_fmt = "📊 Excel" if gr["task"] in _QUIZ_TASKS else "📽️ PowerPoint" if gr["task"] in _PPT_TASKS else "📝 Word"
              _fmt_opts = ["📝 Word (.docx)", "📊 Excel (.xlsx)", "📽️ PowerPoint (.pptx)"]
              _fmt_default_idx = {"📝 Word (.docx)":0,"📊 Excel (.xlsx)":1,"📽️ PowerPoint (.pptx)":2}.get(
                  {"📝 Word":"📝 Word (.docx)","📊 Excel":"📊 Excel (.xlsx)","📽️ PowerPoint":"📽️ PowerPoint (.pptx)"}.get(_default_fmt,"📝 Word (.docx)"),0)
              _sch = st.session_state.get("school_name",""); _tch = st.session_state.get("teacher_name","")
              _fmt_col, _dl_col = st.columns([2,2])
              with _fmt_col:
                  _dl_fmt = st.selectbox("Download as:", _fmt_opts, index=_fmt_default_idx, key="gen_dl_fmt", label_visibility="collapsed")
              with _dl_col:
                  if "Word" in _dl_fmt:
                      _bytes = generate_result_docx(gr["result"], gr["task"], gr["topic"], gr["grade"], gr["subject"], _sch, _tch)
                      if _bytes:
                          st.download_button("📥 Download Word", data=_bytes,
                              file_name=f"{gr['task']}_{gr['topic']}.docx".replace(" ","_")[:60],
                              mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                              key="gen_dl_docx", use_container_width=True)
                      else:
                          st.download_button("📥 Download Text", data=gr["result"],
                              file_name=f"{gr['task']}_{gr['topic']}.txt".replace(" ","_")[:60], key="gen_dl_txt", use_container_width=True)
                  elif "Excel" in _dl_fmt:
                      _bytes = generate_result_xlsx(gr["result"], gr["task"], gr["topic"], gr["grade"], gr["subject"])
                      if _bytes:
                          st.download_button("📥 Download Excel", data=_bytes,
                              file_name=f"{gr['task']}_{gr['topic']}.xlsx".replace(" ","_")[:60],
                              mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                              key="gen_dl_xlsx", use_container_width=True)
                      else:
                          st.download_button("📥 Download Text", data=gr["result"],
                              file_name=f"{gr['task']}_{gr['topic']}.txt".replace(" ","_")[:60], key="gen_dl_txt2", use_container_width=True)
                  else:
                      _bytes = generate_result_pptx(gr["result"], gr["task"], gr["topic"], gr["grade"], gr["subject"])
                      if _bytes:
                          st.download_button("📥 Download PPT", data=_bytes,
                              file_name=f"{gr['task']}_{gr['topic']}.pptx".replace(" ","_")[:60],
                              mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                              key="gen_dl_pptx", use_container_width=True)
                      else:
                          st.download_button("📥 Download Text", data=gr["result"],
                              file_name=f"{gr['task']}_{gr['topic']}.txt".replace(" ","_")[:60], key="gen_dl_txt3", use_container_width=True)
              email_result(gr["result"], f"Teacher Pehpeh — {gr['task']}: {gr['topic']} ({gr['grade']}, {gr['subject']})", "gen")
              tts_player(gr["result"], "gen")

    # TAB 2: STUDENTS
    if t2:
     with t2:
        _stu_word={"en":"Students","fr":"Élèves","sw":"Wanafunzi"}.get(_lang_key(),"Students")
        stu_label=f"{school_name} {_stu_word}" if school_name.strip() else T("my_students")
        st.markdown(f'<div style="background:var(--bg-card);border:1px solid {C_BLUE};border-radius:12px;padding:14px 18px;margin-bottom:10px">{ico(20)} <strong style="color:{C_BLUE}">{stu_label}</strong></div>',unsafe_allow_html=True)

        # ── Parent Letter (moved here from Generate tab) ──
        with st.expander("✉️ Parent Letter / Communication", expanded=False):
            st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:10px;padding:10px 14px;margin-bottom:10px"><strong style="color:{C_GOLD}">✉️ Parent Communication</strong> <span style="font-size:.83rem;color:var(--text-secondary)"> — Adapts to each parent: audio for low-literacy, SMS or email for others</span></div>',unsafe_allow_html=True)
            _plx1,_plx2=st.columns(2)
            with _plx1:
                if st.session_state.students:
                    _plx_student=st.selectbox("Student:",["--Select--"]+[s["name"] for s in st.session_state.students],key="plx_stu")
                else:
                    _plx_student=st.text_input("Student Name:",key="plx_stu_txt",placeholder="e.g., Janjay Kollie")
            with _plx2:
                _plx_concern=st.selectbox("Concern:",["Struggling academically","Frequent absences","Behavior issues","Needs extra support","Outstanding progress","Parent meeting request"],key="plx_concern")
            _plx_details=st.text_area("Specific details (optional):",key="plx_details",height=60,placeholder="e.g., Failing math quizzes...")
            _plx_name=_plx_student if (st.session_state.students and _plx_student!="--Select--") else (st.session_state.get("plx_stu_txt","") or "the student")
            _plx_info=""; _plx_mom_edu="Unknown"
            if st.session_state.students and _plx_student!="--Select--":
                _plxsel=next((s for s in st.session_state.students if s["name"]==_plx_student),None)
                if _plxsel:
                    _plx_info=f'Profile: {_plxsel["sib"]}sib, Mom:{_plxsel["mom"]}, SM:{_plxsel["sm"]}, Works:{_plxsel["wk"]}'
                    _plx_mom_edu=_plxsel.get("mom","Unknown")
            _plx_auto_idx=0 if _plx_mom_edu=="No HS" else 1
            _plx_sensitivity=""
            if st.session_state.students and _plx_student!="--Select--":
                _plxsel2=next((s for s in st.session_state.students if s["name"]==_plx_student),None)
                if _plxsel2:
                    _plxfactors=[]
                    if _plxsel2.get("sib","0-4") in ["5-8","8+"]: _plxfactors.append(f'This child has {_plxsel2["sib"]} siblings — they likely help care for younger ones, cook, fetch water.')
                    if _plxsel2.get("sm","No")=="Yes": _plxfactors.append("This is a single-mother household. The mother is carrying the full burden alone. Be extra respectful of her effort.")
                    if _plxsel2.get("wk","No")=="Yes": _plxfactors.append("This child works after school. They may have very little free time to study.")
                    if _plxsel2.get("mom","Unknown")=="No HS": _plxfactors.append("The mother has no high school education. Never make her feel blamed.")
                    if _plxfactors:
                        _plx_sensitivity=("CRITICAL SENSITIVITY CONTEXT:\n"+
                            "\n".join(f"• {f}" for f in _plxfactors)+
                            "\n\nBecause of this home situation: Do NOT assume the child is lazy. Frame as partnership not authority. Suggest realistic actions.")
            _plx_delivery_opts=["🔊 Voice Message (audio for phone)","📱 Text Message (short SMS)","📧 Email (detailed letter)"]
            _plx_delivery=st.radio("How to reach this parent:",_plx_delivery_opts,index=_plx_auto_idx,key="plx_delivery",horizontal=True)
            if _plx_mom_edu=="No HS" and "🔊" in _plx_delivery:
                st.markdown(f'<div style="background:rgba(46,125,50,.1);border-left:4px solid #4CAF50;padding:8px 12px;border-radius:6px;font-size:.82rem;margin:4px 0">💡 <strong>Voice message recommended</strong> — this parent may find reading difficult.</div>',unsafe_allow_html=True)
            _plx_t_name=teacher_name.strip() if teacher_name.strip() else "the teacher"
            _plx_t_phone=teacher_phone.strip() if teacher_phone.strip() else ""
            _plx_t_sig=_plx_t_name+( f" ({_plx_t_phone})" if _plx_t_phone else "")
            _plx_close_loop=("CRITICAL: End with a specific simple way for the parent to confirm they received this.")
            if "🔊" in _plx_delivery:
                _plx_prompt=(f"Create a SHORT VOICE MESSAGE (under 45 seconds spoken) from {_plx_t_name} to parent of {_plx_name}.\n"
                    f"CONCERN: {_plx_concern}\nDETAILS: {_plx_details or 'General concern'}\nSTUDENT: {_plx_info or 'N/A'}\n{_plx_sensitivity}"
                    f"SCHOOL: {school_name or 'Our school'} | {_subj_en} | {_grade_en}\n"
                    f"Sign as {_plx_t_sig}. Warm, caring Liberian English tone. OUTPUT: Only the spoken message text.")
            elif "📱" in _plx_delivery:
                _plx_prompt=(f"Create a SHORT TEXT MESSAGE (max 300 chars) from {_plx_t_name} to parent of {_plx_name}.\n"
                    f"CONCERN: {_plx_concern}\nDETAILS: {_plx_details or 'General concern'}\nSTUDENT: {_plx_info or 'N/A'}\n{_plx_sensitivity}"
                    f"SCHOOL: {school_name or 'Our school'} | {_subj_en} | {_grade_en}\n"
                    f"Sign as {_plx_t_sig}. Max 300 chars. OUTPUT: Only the SMS text.")
            else:
                _plx_prompt=(f"Write a letter from {_plx_t_name} to parent/guardian of {_plx_name}.\n"
                    f"CONCERN: {_plx_concern}\nDETAILS: {_plx_details or 'General concern'}\nSTUDENT: {_plx_info or 'N/A'}\n{_plx_sensitivity}"
                    f"SCHOOL: {school_name or 'Our school'} | {_subj_en} | {_grade_en}\n"
                    f"Simple reading level. Warm, respectful Liberian English. Sign as {_plx_t_sig}. OUTPUT: Only the letter.")
            if st.button("✉️ Generate Parent Communication", key="plx_gen", type="primary", use_container_width=True):
                with st.status("Drafting parent communication...", expanded=True) as _plx_status:
                    st.write("Preparing sensitive, culturally appropriate message...")
                    _plx_r,_plx_m,_plx_allr=best_all(build_free_chat(),_plx_prompt,[])
                    _plx_status.update(label="✉️ Communication ready!",state="complete",expanded=False)
                st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:10px;padding:14px 18px;margin:8px 0;white-space:pre-wrap;color:#D0D8E8;line-height:1.7">{_plx_r}</div>',unsafe_allow_html=True)

        with st.expander({"en":"➕ Add Profile","fr":"➕ Ajouter un profil","sw":"➕ Ongeza Wasifu"}.get(_lang_key(),"➕ Add Profile"),expanded=not st.session_state.students):
            c1,c2=st.columns(2)
            with c1: sn=st.text_input(T("name"),key="sn"); sib=st.selectbox({"en":"Siblings","fr":"Frères et sœurs","sw":"Ndugu"}.get(_lang_key(),"Siblings"),["0-4","5-8","8+"],key="sb"); me=st.selectbox({"en":"Mom Edu","fr":"Éducation mère","sw":"Elimu ya mama"}.get(_lang_key(),"Mom Edu"),["HS Grad","No HS","Unknown"],key="me")
            with c2: sm=st.selectbox({"en":"Single Mom?","fr":"Mère seule ?","sw":"Mama peke yake?"}.get(_lang_key(),"Single Mom?"),["No","Yes","Unknown"],key="sm"); wk=st.selectbox({"en":"Works?","fr":"Travaille ?","sw":"Anafanya kazi?"}.get(_lang_key(),"Works?"),["No","Yes","Unknown"],key="wk"); cp=st.selectbox({"en":"Computer?","fr":"Ordinateur ?","sw":"Kompyuta?"}.get(_lang_key(),"Computer?"),["Never","Rarely","Sometimes","Often"],key="cp")
            nt=st.text_area("Notes",key="nt",height=50)
            if st.button({"en":"✅ Save","fr":"✅ Enregistrer","sw":"✅ Hifadhi"}.get(_lang_key(),"✅ Save"),key="sv") and sn.strip():
                st.session_state.students.append(dict(name=sn.strip(),sib=sib,mom=me,sm=sm,wk=wk,cp=cp,nt=nt.strip())); st.rerun()
        # Excel bulk upload
        with st.expander(T("upload_excel")):
            st.markdown(f"""<div style="font-size:.85rem;color:var(--text-secondary);margin-bottom:8px">
            Upload a spreadsheet with these columns:<br>
            <strong>Name</strong> (required), <strong>Siblings</strong> (0-4, 5-8, or 8+), <strong>Mom_Edu</strong> (HS Grad, No HS, or Unknown),
            <strong>Single_Mom</strong> (Yes, No, or Unknown), <strong>Works</strong> (Yes, No, or Unknown),
            <strong>Computer</strong> (Never, Rarely, Sometimes, or Often), <strong>Notes</strong> (optional)
            </div>""",unsafe_allow_html=True)
            dl_cols=["Name","Siblings","Mom_Edu","Single_Mom","Works","Computer","Notes"]
            dl_example=[["Janjay Kollie","5-8","No HS","Yes","No","Never","Quiet, needs encouragement"],["Hawa Sirleaf","0-4","HS Grad","No","No","Sometimes",""]]
            if PD:
                template_df=pd.DataFrame(dl_example,columns=dl_cols)
                csv_data=template_df.to_csv(index=False)
                st.download_button("📥 Download Template",data=csv_data,file_name="student_template.csv",mime="text/csv",key="dl_tmpl")
            uploaded=st.file_uploader("Choose file",type=["csv","xlsx","xls"],key="stu_upload",label_visibility="collapsed")
            if uploaded and PD:
                try:
                    if uploaded.name.endswith(".csv"):
                        df=pd.read_csv(uploaded)
                    else:
                        df=pd.read_excel(uploaded)
                    # Normalize column names
                    df.columns=[c.strip().replace(" ","_") for c in df.columns]
                    # Map common variations
                    col_map={"name":"Name","student":"Name","student_name":"Name","siblings":"Siblings","sib":"Siblings",
                             "mom_edu":"Mom_Edu","mother_education":"Mom_Edu","mom":"Mom_Edu","mother":"Mom_Edu",
                             "single_mom":"Single_Mom","single_mother":"Single_Mom","sm":"Single_Mom",
                             "works":"Works","works_after_school":"Works","working":"Works",
                             "computer":"Computer","computer_access":"Computer","comp":"Computer",
                             "notes":"Notes","note":"Notes","comments":"Notes"}
                    df.columns=[col_map.get(c.lower(),c) for c in df.columns]
                    if "Name" not in df.columns:
                        st.error("❌ Spreadsheet must have a 'Name' column.")
                    else:
                        st.dataframe(df,use_container_width=True,height=200)
                        st.markdown(f'<div style="color:var(--text-secondary);font-size:.85rem">📊 Found <strong>{len(df)}</strong> students</div>',unsafe_allow_html=True)
                        if st.button(f"✅ Import {len(df)} Students",type="primary",key="bulk_import"):
                            imported=0
                            for _,row in df.iterrows():
                                name=str(row.get("Name","")).strip()
                                if not name: continue
                                sib_v=str(row.get("Siblings","0-4")).strip()
                                if sib_v not in ["0-4","5-8","8+"]: sib_v="0-4"
                                mom_v=str(row.get("Mom_Edu","Unknown")).strip()
                                if mom_v not in ["HS Grad","No HS","Unknown"]: mom_v="Unknown"
                                sm_v=str(row.get("Single_Mom","Unknown")).strip()
                                if sm_v not in ["Yes","No","Unknown"]: sm_v="Unknown"
                                wk_v=str(row.get("Works","Unknown")).strip()
                                if wk_v not in ["Yes","No","Unknown"]: wk_v="Unknown"
                                cp_v=str(row.get("Computer","Never")).strip()
                                if cp_v not in ["Never","Rarely","Sometimes","Often"]: cp_v="Never"
                                nt_v=str(row.get("Notes","")).strip()
                                if nt_v=="nan": nt_v=""
                                st.session_state.students.append(dict(name=name,sib=sib_v,mom=mom_v,sm=sm_v,wk=wk_v,cp=cp_v,nt=nt_v))
                                imported+=1
                            st.success(f"✅ Imported {imported} students!"); time.sleep(1); st.rerun()
                except Exception as e:
                    st.error(f"❌ Error reading file: {e}")
        for i,s in enumerate(st.session_state.students):
            rsk=[]
            if s["mom"]=="No HS": rsk.append("🔴 No HS Mom")
            if s["sm"]=="Yes": rsk.append("🔴 Single Mom")
            if s["sib"]=="8+": rsk.append("🟠 8+ siblings")
            if s["wk"]=="Yes": rsk.append("🟠 Works")
            if s["cp"]=="Never": rsk.append("🟡 No computer")
            st.markdown(f'<div class="sc"><strong style="color:{C_BLUE}">{s["name"]}</strong> — {s["sib"]} sib, Mom:{s["mom"]}<br><span style="font-size:.82rem">{" · ".join(rsk) or "🟢 Lower risk"}</span></div>',unsafe_allow_html=True)
            info=f'{s["name"]},{s["sib"]}sib,Mom:{s["mom"]},SM:{s["sm"]},Works:{s["wk"]},Comp:{s["cp"]},{s["nt"]}'
            b1,b2,b3,b4=st.columns([3,3,3,1])
            with b1:
                if st.button(T("assignment"),key=f"a{i}"):
                    with st.spinner(T("creating")):
                        r,m,allr=best_all(build_stu(_region_val,country,_grade_en,_subj_en,_size_val,_res_val,LANGS[lang],_abl_val,info,school_name),f"Tailored {_subj_en} assignment. Max 3 problems.")
                    _by={"en":"by","fr":"par","sw":"na"}.get(_lang_key(),"by"); st.markdown(f'<div class="rb">{highlight_result(r)}<div style="font-size:.65rem;color:#556;margin-top:4px">{_by} {m}</div></div>',unsafe_allow_html=True)
                    if len(allr)>1:
                        with st.expander(f"{T('see_all')} {len(allr)} {T('model_responses')}"):
                            for mn,mr in allr.items():
                                mico={"Claude":"🟣","ChatGPT":"🟢","Gemini":"🔵"}.get(mn,"⚪")
                                st.markdown(f"**{mico} {mn}{'  ✅' if mn==m else ''}**"); st.markdown(mr); st.markdown("---")
                    email_result(r, f"Teacher Pehpeh — Assignment for {s['name']} ({subject}, {grade})", f"asn_{i}")
                    tts_player(r, f"asn_{i}")
            with b2:
                if st.button(T("risk"),key=f"r{i}"):
                    with st.spinner("Analyzing..."):
                        r,m,allr=best_all(build_stu(_region_val,country,_grade_en,_subj_en,_size_val,_res_val,LANGS[lang],_abl_val,info,school_name),"Risk analysis using IBT data. Compare to 183-student dataset.")
                    _by={"en":"by","fr":"par","sw":"na"}.get(_lang_key(),"by"); st.markdown(f'<div class="rb">{highlight_result(r)}<div style="font-size:.65rem;color:#556;margin-top:4px">{_by} {m}</div></div>',unsafe_allow_html=True)
                    if len(allr)>1:
                        with st.expander(f"{T('see_all')} {len(allr)} {T('model_responses')}"):
                            for mn,mr in allr.items():
                                mico={"Claude":"🟣","ChatGPT":"🟢","Gemini":"🔵"}.get(mn,"⚪")
                                st.markdown(f"**{mico} {mn}{'  ✅' if mn==m else ''}**"); st.markdown(mr); st.markdown("---")
                    email_result(r, f"Teacher Pehpeh — Risk Analysis for {s['name']} ({subject}, {grade})", f"rsk_{i}")
                    tts_player(r, f"rsk_{i}")
            with b3:
                if st.button("🪪 ID Card",key=f"card{i}"):
                    st.session_state[f"_show_card_{i}"]=not st.session_state.get(f"_show_card_{i}",False)
                    st.rerun()
            with b4:
                if st.button("🗑️",key=f"d{i}"): st.session_state.students.pop(i); st.rerun()
            # Student ID Card display with save options
            if st.session_state.get(f"_show_card_{i}"):
                with st.container():
                    card_bytes, card_fname = generate_student_card(s, school_name, _grade_en, _subj_en, country)
                    if card_bytes:
                        st.image(card_bytes, caption=f"Student ID — {s['name']}", use_container_width=True)
                        _save_label={"en":"Save card to...","fr":"Enregistrer la carte...","sw":"Hifadhi kadi..."}.get(_lang_key(),"Save card to...")
                        _save_dest=st.selectbox(_save_label,["📥 Download to Device","📧 Email Card"],key=f"card_save_{i}",label_visibility="visible")
                        if _save_dest.startswith("📥"):
                            st.download_button(
                                "📥 Download Student Card",
                                data=card_bytes,
                                file_name=card_fname,
                                mime="image/png",
                                key=f"dl_card_{i}",
                                use_container_width=True
                            )
                        elif _save_dest.startswith("📧"):
                            _card_b64=base64.b64encode(card_bytes).decode()
                            email_result(f"Student ID Card for {s['name']}\nSchool: {school_name}\nGrade: {_grade_en}\nSubject: {_subj_en}", f"Teacher Pehpeh — Student Card: {s['name']}", f"card_em_{i}")
                            st.download_button(
                                "📥 Also Download Copy",
                                data=card_bytes,
                                file_name=card_fname,
                                mime="image/png",
                                key=f"dl_card_em_{i}",
                                use_container_width=True
                            )
                    else:
                        st.warning(f"Could not generate card: {card_fname}")
        if st.session_state.students:
            st.markdown("---"); st.markdown(f"#### {ico(16)} {T('grade_work')}",unsafe_allow_html=True)
            gs=st.selectbox("Student:",[s["name"] for s in st.session_state.students],key="gs")
            gw_col, gw_mic = st.columns([5,1])
            with gw_col:
                gw=st.text_area(T("students_work"),height=100,key="gw")
            with gw_mic:
                st.markdown("<div style='height:24px'></div>", unsafe_allow_html=True)
                grade_audio=st.audio_input("🎤",key="grade_mic",label_visibility="collapsed")
            if grade_audio:
                with st.spinner(T("transcribing")):
                    grade_text=transcribe_audio(grade_audio.read())
                if grade_text:
                    st.session_state["gw"]=grade_text; st.rerun()
            gsub=st.selectbox(f"{T('subject')}:",_subjects(),key="gsub"); gt=st.text_input(f"{T('topic')}:",key="gt")
            _gsub_en=_to_en_subj(gsub)
            if st.button(T("grade_btn"),type="primary",key="gb") and gw.strip():
                sel=next((s for s in st.session_state.students if s["name"]==gs),None)
                if sel:
                    info=f'{sel["name"]},{sel["sib"]}sib,Mom:{sel["mom"]},SM:{sel["sm"]},Works:{sel["wk"]},Comp:{sel["cp"]},{sel["nt"]}'
                    with st.spinner("Grading..."):
                        r,m,allr=best_all(build_stu(_region_val,country,_grade_en,_gsub_en,_size_val,_res_val,LANGS[lang],_abl_val,info,school_name),f"Grade:\nSTUDENT:{info}\n{_gsub_en} {gt}\n\nWORK:\n{gw}\n\nGive: grade, praise, corrections, tips, next step.")
                    _fb={"en":"Feedback","fr":"Commentaires","sw":"Maoni"}.get(_lang_key(),"Feedback")
                    st.markdown(f'<div class="rh"><h3>{ico(16)} {_fb}: {gs}</h3></div><div class="rb">{highlight_result(r)}</div>',unsafe_allow_html=True)
                    if len(allr)>1:
                        with st.expander(f"{T('see_all')} {len(allr)} {T('model_responses')}"):
                            for mn,mr in allr.items():
                                mico={"Claude":"🟣","ChatGPT":"🟢","Gemini":"🔵"}.get(mn,"⚪")
                                st.markdown(f"**{mico} {mn}{'  ✅' if mn==m else ''}**"); st.markdown(mr); st.markdown("---")
                    email_result(r, f"Teacher Pehpeh — Grade Feedback for {gs} ({gsub}, {grade})", "grade")
                    tts_player(r, "grade")

        # === PHOTO GRADING — Batch scan student work ===
        if st.session_state.students:
            st.markdown("---")
            st.markdown(f'<div style="background:linear-gradient(135deg,#1565C0,#0D47A1);border-radius:12px;padding:14px 18px;margin-bottom:10px;color:white"><strong>📸 Photo Grading</strong> — Snap or upload photos of handwritten student work</div>',unsafe_allow_html=True)
            _max_batch = min(10, len(st.session_state.students))
            pg_cols = st.columns([2,2,2])
            with pg_cols[0]:
                pg_count = st.number_input("Students to grade", min_value=1, max_value=_max_batch, value=min(3, _max_batch), key="pg_count")
            with pg_cols[1]:
                pg_subj = st.selectbox("Subject:", _subjects(), key="pg_subj")
            with pg_cols[2]:
                pg_topic = st.text_input("Topic/Assignment:", key="pg_topic", placeholder="e.g., Algebra Quiz #3")
            _pg_subj_en = _to_en_subj(pg_subj)

            # Create upload slots
            pg_uploads = []
            pg_students = []
            for slot in range(int(pg_count)):
                sc1, sc2 = st.columns([1, 2])
                with sc1:
                    stu_name = st.selectbox(f"Student {slot+1}:", [s["name"] for s in st.session_state.students], key=f"pg_stu_{slot}")
                    pg_students.append(stu_name)
                with sc2:
                    photo = st.file_uploader(f"Upload work", type=["jpg","jpeg","png","heic","pdf"], key=f"pg_img_{slot}", label_visibility="collapsed")
                    if not photo:
                        photo = st.camera_input(f"📷 Or take photo", key=f"pg_cam_{slot}", label_visibility="collapsed")
                    pg_uploads.append(photo)

            _has_any = any(pg_uploads)
            if st.button("Grade All Photos", type="primary", use_container_width=True, key="pg_go", disabled=not _has_any):
                _graded = 0
                for slot in range(int(pg_count)):
                    photo = pg_uploads[slot]
                    sname = pg_students[slot]
                    if not photo:
                        continue
                    sel = next((s for s in st.session_state.students if s["name"] == sname), None)
                    if not sel:
                        continue
                    info = f'{sel["name"]},{sel["sib"]}sib,Mom:{sel["mom"]},SM:{sel["sm"]},Works:{sel["wk"]},Comp:{sel["cp"]},{sel["nt"]}'
                    img_bytes = photo.read()
                    img_b64 = base64.b64encode(img_bytes).decode()
                    mime = photo.type if hasattr(photo, 'type') and photo.type else "image/jpeg"
                    if "pdf" in mime:
                        st.warning(f"⚠️ PDF not supported for {sname} — use JPG/PNG photo instead.")
                        continue

                    sys_prompt = f"""You are Teacher Pehpeh, an expert Liberian education AI grading handwritten student work.
STUDENT PROFILE: {info}
SCHOOL: {school_name or 'Not specified'} | COUNTRY: {country} | GRADE: {_grade_en} | SUBJECT: {_pg_subj_en}
TOPIC: {pg_topic or 'General'}

INSTRUCTIONS:
1. Read the handwritten work carefully. Transcribe what you see.
2. Grade on a scale of 0-100.
3. Note specific errors with corrections.
4. TAILOR feedback to this student's profile (risk factors, background).
5. Give 2-3 specific next steps personalized for this student.
6. Be encouraging — celebrate what they got right FIRST.
7. Format: SCORE | TRANSCRIPTION | STRENGTHS | ERRORS | PERSONALIZED FEEDBACK | NEXT STEPS

IMPORTANT: Extract a numeric score (0-100) on the FIRST line as: SCORE: XX/100"""

                    with st.spinner(f"📸 Grading {sname}'s work..."):
                        r, m = best_vision(sys_prompt, f"Grade this student's handwritten {_pg_subj_en} work on {pg_topic}. Student: {sname}.", img_b64, mime)

                    if r and not str(r).startswith("⚠️"):
                        _graded += 1
                        # Extract score for history
                        import re as _re
                        _score_match = _re.search(r'SCORE:\s*(\d+)', r, _re.IGNORECASE)
                        _score = int(_score_match.group(1)) if _score_match else None
                        if not _score:
                            _score_match = _re.search(r'(\d+)\s*/\s*100', r)
                            _score = int(_score_match.group(1)) if _score_match else None

                        # Save to grade history
                        if _score is not None:
                            import datetime
                            st.session_state.grade_history.append({
                                "student": sname, "subject": _pg_subj_en,
                                "topic": pg_topic or "General", "score": _score,
                                "date": datetime.datetime.now().isoformat(),
                                "grade_level": _grade_en, "model": m
                            })

                        _fb_label = {"en":"Feedback","fr":"Commentaires","sw":"Maoni"}.get(_lang_key(),"Feedback")
                        _score_display = f" — {_score}/100" if _score else ""
                        st.markdown(f'<div class="rh"><h3>📸 {_fb_label}: {sname}{_score_display}</h3></div><div class="rb">{highlight_result(r)}<div style="font-size:.65rem;color:#556;margin-top:4px">by {m}</div></div>',unsafe_allow_html=True)
                        email_result(r, f"Teacher Pehpeh — Photo Grade for {sname} ({pg_subj}, {grade})", f"pg_{slot}")
                    else:
                        st.error(f"⚠️ Could not grade {sname}'s work: {r}")

                if _graded:
                    st.success(f"✅ Graded {_graded} student{'s' if _graded>1 else ''}'s work!")

        # === TREND ANALYTICS — Performance over time ===
        if st.session_state.grade_history:
            st.markdown("---")
            st.markdown(f'<div style="background:linear-gradient(135deg,#2E7D32,#1B5E20);border-radius:12px;padding:14px 18px;margin-bottom:10px;color:white"><strong>📈 Performance Trends</strong> — Track student progress over time</div>',unsafe_allow_html=True)

            import datetime
            _gh = st.session_state.grade_history

            # Filter controls
            _tc1, _tc2 = st.columns(2)
            with _tc1:
                _all_stu = sorted(set(g["student"] for g in _gh))
                _sel_stu = st.multiselect("Filter by student:", _all_stu, default=_all_stu, key="trend_stu")
            with _tc2:
                _all_subj = sorted(set(g["subject"] for g in _gh))
                _sel_subj = st.multiselect("Filter by subject:", _all_subj, default=_all_subj, key="trend_subj")

            _filtered = [g for g in _gh if g["student"] in _sel_stu and g["subject"] in _sel_subj]

            if _filtered and PD:
                _df = pd.DataFrame(_filtered)
                _df["date"] = pd.to_datetime(_df["date"])
                _df["date_str"] = _df["date"].dt.strftime("%b %d")

                # Summary stats
                _avg = _df["score"].mean()
                _latest_avg = _df.sort_values("date").tail(max(1,len(_df)//3))["score"].mean()
                _earliest_avg = _df.sort_values("date").head(max(1,len(_df)//3))["score"].mean()
                _trend_dir = "📈" if _latest_avg > _earliest_avg else ("📉" if _latest_avg < _earliest_avg else "➡️")
                _improvement = _latest_avg - _earliest_avg

                _m1, _m2, _m3, _m4 = st.columns(4)
                with _m1:
                    st.metric("Avg Score", f"{_avg:.0f}/100")
                with _m2:
                    st.metric("Assignments Graded", len(_filtered))
                with _m3:
                    st.metric("Trend", f"{_trend_dir} {'+' if _improvement>0 else ''}{_improvement:.1f}pts")
                with _m4:
                    _pass_rate = len(_df[_df["score"]>=50])/len(_df)*100
                    st.metric("Pass Rate (≥50)", f"{_pass_rate:.0f}%")

                # Per-student scores over time
                st.markdown("**Student Performance Over Time**")
                _pivot = _df.groupby(["date_str","student"])["score"].mean().reset_index()
                _chart_data = _pivot.pivot(index="date_str", columns="student", values="score")

                # Add projected improvement line (target: 5% improvement per assessment)
                if len(_chart_data) > 1:
                    _baseline = _df.sort_values("date").head(max(1,len(_df)//3))["score"].mean()
                    _proj = []
                    for idx, dt in enumerate(_chart_data.index):
                        _proj.append(_baseline + (idx * 5))  # 5pts projected improvement per session
                    _chart_data["🎯 Projected Target"] = [min(100, p) for p in _proj]

                st.line_chart(_chart_data, use_container_width=True)

                # Subject breakdown
                if len(_all_subj) > 1:
                    st.markdown("**Average by Subject**")
                    _subj_avg = _df.groupby("subject")["score"].agg(["mean","count"]).round(1)
                    _subj_avg.columns = ["Avg Score", "Assignments"]
                    st.dataframe(_subj_avg, use_container_width=True)

                # At-risk students (below 50 average)
                _stu_avg = _df.groupby("student")["score"].mean()
                _at_risk = _stu_avg[_stu_avg < 50].sort_values()
                if len(_at_risk):
                    st.markdown(f'<div style="background:rgba(239,83,80,.1);border:2px solid #EF5350;border-radius:10px;padding:12px;margin:8px 0">'
                                f'<strong style="color:#EF5350">⚠️ Students Needing Intervention ({len(_at_risk)})</strong><br>'
                                + "<br>".join(f'<span style="color:#F0D5D5">• <strong>{name}</strong>: {score:.0f}/100 avg</span>' for name,score in _at_risk.items())
                                + '</div>', unsafe_allow_html=True)

                # Export grades
                _exp1, _exp2 = st.columns(2)
                with _exp1:
                    _csv = _df[["student","subject","topic","score","date_str","grade_level"]].to_csv(index=False)
                    st.download_button("📥 Export Grades (CSV)", data=_csv, file_name=f"teacher_pehpeh_grades_{school_name or 'class'}.csv", mime="text/csv", key="exp_grades")
                with _exp2:
                    if st.button("🗑️ Clear Grade History", key="clr_grades"):
                        st.session_state.grade_history = []; st.rerun()

    # TAB 3: CHAT
    if t3:
     with t3:
        # Free-flowing header — independent of classroom config
        st.markdown(
            f'<div style="background:rgba(139,26,26,.12);border:1px solid rgba(178,34,52,.3);border-radius:12px;padding:12px 18px;margin-bottom:10px">' +
            f'{ico(20)} <strong style="color:{C_GOLD}">Teacher Pehpeh Chat</strong> ' +
            f'<span style="color:#A09080;font-size:.82rem">— Ask anything · Any subject · Any level</span></div>',
            unsafe_allow_html=True
        )

        # ── Africa Discovery buttons ──
        # Each tuple: (button label, hook line, full prompt for AI)
        _AFRICA_SPARKS = [
            (
                "🏠 Huts & Aerodynamics",
                "African huts are round for a reason physics teachers love",
                "A student just asked: why are traditional African huts built round? "
                "Answer in 4-6 sentences max. Be vivid and surprising. Cover: aerodynamics (wind passes around, not against), "
                "structural strength (no corners to crack), heat circulation (air rises evenly), and how this is the same principle "
                "used in modern wind-resistant architecture. End with one punchy sentence that connects this ancient wisdom to modern engineering. "
                "Use bold for 2-3 key phrases. No bullet points — flowing prose only. Emotive and proud tone."
            ),
            (
                "🌍 Liberia: Pepper Coast",
                "Before it was Liberia, Europeans had a spicier name for it",
                "A student just discovered that Liberia was once called the Pepper Coast by European traders. "
                "Answer in 4-6 sentences. Be vivid: explain that the Grain of Selim (Grains of Paradise) — a peppery spice — grew wild along this coast "
                "and was so valuable it drove European trade routes. The coast was so associated with this spice that Portuguese and Dutch traders named entire "
                "stretches after it. Connect this to why Liberia's geography made it a crossroads of Atlantic trade before colonisation. "
                "End with something that makes the student feel proud of this heritage. Bold 2-3 key facts. Prose only, no lists."
            ),
            (
                "👑 Mansa Musa's Wealth",
                "He was so rich his vacation crashed an entire economy",
                "A student asks: was Mansa Musa really the richest person in history? "
                "Answer in 4-6 vivid sentences. Cover: his 1324 pilgrimage to Mecca with 60,000 people and 100 camel-loads of gold; "
                "how he gave away so much gold in Cairo that he caused inflation and crashed the Egyptian gold market for a decade; "
                "estimates of his wealth at $400 billion in today's money; and how the Mali Empire sat on over half the world's gold supply. "
                "Tone: jaw-dropping, proud, cinematic. Bold 2-3 key facts. Prose only."
            ),
            (
                "🌍 Oldest University",
                "The world's oldest university is in Africa — and still running",
                "A student just heard that the world's oldest university is in Africa. Answer in 4-6 sentences. "
                "Cover: the University of al-Qarawiyyin in Fez, Morocco, founded in 859 CE by Fatima al-Fihri — a woman — "
                "over 200 years before Oxford; what was taught there (theology, grammar, rhetoric, logic, astronomy, music); "
                "that it still operates today; and that Timbuktu's universities held over a million manuscripts. "
                "Tone: astonishing, empowering, proud. Bold 2-3 key facts. Prose only, no bullet points."
            ),
            (
                "🧬 We Are All African",
                "Every human alive today traces their DNA to one place",
                "A student asks about the Out of Africa theory. Answer in 4-6 vivid sentences. "
                "Cover: that modern Homo sapiens evolved in Africa around 300,000 years ago; that all non-African humans descend from "
                "a small group that left Africa roughly 60,000-70,000 years ago; that Africa contains more genetic diversity than the rest "
                "of the world combined; and what this means — that 'race' as a biological concept has no scientific basis, because we are all, "
                "at root, African. End with something that reframes how the student sees human identity. Bold 2-3 key facts. Emotive, precise prose."
            ),
            (
                "📐 Egypt & Pi",
                "Ancient Egyptians calculated Pi 4,000 years before your calculator",
                "A student asks how advanced ancient Egyptian mathematics really was. Answer in 4-6 sentences. "
                "Cover: the Rhind Papyrus (1650 BCE) showing multiplication, fractions, and a value of Pi accurate to within 1%; "
                "the Pythagorean theorem used in pyramid construction centuries before Pythagoras was born; "
                "that the Great Pyramid of Giza aligns to true north with 0.05 degree accuracy — better than the Greenwich Observatory; "
                "and that the Ishango bone from Congo (25,000 BCE) is the world's oldest known mathematical object. "
                "Tone: proud, mind-blowing, precise. Bold 2-3 key facts. Prose only."
            ),
            (
                "🌊 The Niger's Secret",
                "The Niger River flows backwards — and baffled explorers for centuries",
                "A student is puzzled about the Niger River's unusual geography. Answer in 4-6 vivid sentences. "
                "Cover: that the Niger is one of the only major rivers that flows AWAY from the sea before turning back toward it — "
                "creating a massive inland delta in Mali (the Inner Niger Delta); that European explorers spent 300 years trying to trace it "
                "and repeatedly got the direction wrong; that the inland delta floods 30,000 sq km and supports millions of people, birds, "
                "and fish in the middle of the Sahel; and how the ancient Mali and Songhai empires were built around controlling this waterway. "
                "Tone: mysterious, epic, proud. Bold 2-3 facts. Prose only."
            ),
            (
                "🦁 Ubuntu Philosophy",
                "One African word that philosophers say changes everything about how you see yourself",
                "A student asks what Ubuntu means beyond just a Linux operating system. Answer in 4-6 vivid sentences. "
                "Cover: the Nguni Bantu phrase 'Umuntu ngumuntu ngabantu' — a person is a person through other persons; "
                "how this philosophy underpins traditional governance, conflict resolution, and identity across Sub-Saharan Africa; "
                "Nelson Mandela's use of Ubuntu as the moral foundation of post-apartheid South Africa; "
                "and how it challenges Western individualism — you are not a self-made person, you are a we-made person. "
                "End with something that makes the student rethink their own identity. Bold 2-3 phrases. Emotive prose only."
            ),
        ]
        # Rotate buttons based on day of month so they feel fresh
        import random as _rnd
        if "spark_order" not in st.session_state:
            st.session_state.spark_order = _rnd.sample(range(len(_AFRICA_SPARKS)), min(6, len(_AFRICA_SPARKS)))
        _todays_sparks = [_AFRICA_SPARKS[i] for i in st.session_state.spark_order]

        _hdr_col, _shuf_col = st.columns([5, 1])
        with _hdr_col:
            st.markdown(
                f'<div style="color:{C_GOLD};font-size:.82rem;font-weight:700;margin-bottom:4px">' +
                f'🌍 Did You Know? &nbsp;<span style="color:#8899aa;font-size:.76rem;font-weight:400">— Click to discover something amazing about Africa</span></div>',
                unsafe_allow_html=True
            )
        with _shuf_col:
            if st.button("🔀", key="spark_shuffle", help="Shuffle topics", use_container_width=True):
                import random as _rnd2
                st.session_state.spark_order = _rnd2.sample(range(len(_AFRICA_SPARKS)), min(6, len(_AFRICA_SPARKS)))
                st.rerun()
        _spark_cols = st.columns(6)
        for _sbi, (_slabel, _shook, _sprompt) in enumerate(_todays_sparks):
            with _spark_cols[_sbi]:
                if st.button(_slabel, key=f"spark_{_sbi}", use_container_width=True, help=_shook):
                    _spark_user_msg = f"🌍 {_shook}"
                    st.session_state.chat_messages.append({"role":"user","content":_spark_user_msg})
                    with st.status("🌍 Discovering...", expanded=False) as _ss:
                        _sr, _sm, _sallr = best_all(build_free_chat(), _sprompt,
                                                     [{"role":x["role"],"content":x["content"]}
                                                      for x in st.session_state.chat_messages[:-1]])
                        _ss.update(label="✨ Here it is!", state="complete", expanded=False)
                    st.session_state.chat_messages.append({
                        "role":"assistant","content":_sr,"model":_sm,"all_responses":_sallr,
                        "is_africa_spark": True
                    })
                    st.rerun()

        st.markdown(f'<div style="font-size:.78rem;color:#6677AA;margin:6px 0 4px">💡 Ask anything — explain a concept, get study tips, generate questions, start with <em>draw</em> for images, or 📸 upload a photo and say <em>transcribe</em></div>', unsafe_allow_html=True)
        st.markdown("---")
        for mi,msg in enumerate(st.session_state.chat_messages):
            _you_label={"en":"🧑‍🏫 You","fr":"🧑‍🏫 Vous","sw":"🧑‍🏫 Wewe"}.get(_lang_key(),"🧑‍🏫 You")
            _by_label={"en":"by","fr":"par","sw":"na"}.get(_lang_key(),"by")
            if msg["role"]=="user":
                st.markdown(f'<div class="ct"><div style="font-size:.75rem;font-weight:700;color:{C_BLUE};margin-bottom:4px">{_you_label}</div>{msg["content"]}</div>',unsafe_allow_html=True)
                if msg.get("photo_b64"):
                    try:
                        _ph_bytes = base64.b64decode(msg["photo_b64"])
                        st.image(_ph_bytes, caption="📸 Attached photo", width=300)
                    except: pass
            else:
                allr=msg.get("all_responses",{})
                if msg.get("is_mock_test"):
                    # Show a collapsed version of the raw text (so answers aren't immediately visible)
                    with st.expander("📄 View raw question text (answers visible)", expanded=False):
                        st.markdown(f'<div class="cp">{highlight_result(msg["content"])}<div style="font-size:.65rem;color:#556;margin-top:4px">{_by_label} {msg.get("model","AI")}</div></div>',unsafe_allow_html=True)
                else:
                    st.markdown(f'<div class="cp"><div style="font-size:.75rem;font-weight:700;color:{C_GOLD};margin-bottom:4px">{ico(16)} Teacher Pehpeh</div>{highlight_result(msg["content"])}<div style="font-size:.65rem;color:#556;margin-top:4px">{_by_label} {msg.get("model","AI")}</div></div>',unsafe_allow_html=True)
                # Auto-render interactive answer sheet for mock test responses
                if msg.get("is_mock_test"):
                    _chat_parsed = parse_mcq_for_sheet(msg["content"])
                    if len(_chat_parsed) >= 2:
                        import streamlit.components.v1 as _comp2
                        import json as _json2
                        _cqs = _json2.dumps(_chat_parsed)
                        _cn = len(_chat_parsed)
                        _csubj = msg.get("mock_subject","")
                        _chat_sheet_html = _make_mcq_sheet_html(_cqs, title=_csubj, n=_cn)
                        with st.expander(f"📋 Interactive Answer Sheet — {_csubj} ({_cn} questions)", expanded=True):
                            _comp2.html(_chat_sheet_html, height=max(560, 220 + _cn * 48), scrolling=True)
                if msg.get("image"):
                    st.image(msg["image"],caption=f"🎨 Generated by {msg.get('image_src','AI')}",use_container_width=True)
                if msg.get("docx_bytes"):
                    try:
                        _docx_data = base64.b64decode(msg["docx_bytes"])
                        _docx_fname = msg.get("docx_name", "transcribed_notes.docx")
                        st.markdown(f'<div style="background:rgba(76,175,80,.1);border:2px solid #4CAF50;border-radius:10px;padding:12px 16px;margin:8px 0"><strong style="color:#4CAF50">📝 {T("transcribe_done")}</strong></div>', unsafe_allow_html=True)
                        st.download_button(
                            T("download_docx"), data=_docx_data,
                            file_name=_docx_fname, mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key=f"docx_dl_{mi}", type="primary", use_container_width=True
                        )
                    except: pass
                if len(allr)>1:
                    with st.expander(f"{T('see_all')} {len(allr)} {T('model_responses')}",expanded=False):
                        for mname,mresp in allr.items():
                            mico={"Claude":"🟣","ChatGPT":"🟢","Gemini":"🔵"}.get(mname,"⚪")
                            is_primary=" ✅ (selected)" if mname==msg.get("model") else ""
                            st.markdown(f"**{mico} {mname}{is_primary}**")
                            st.markdown(mresp)
                            st.markdown("---")
                # Find the user question this responds to
                user_q = st.session_state.chat_messages[mi-1]["content"] if mi>0 else "Chat"
                email_result(msg["content"], f"Teacher Pehpeh — {user_q[:50]} ({grade}, {subject})", f"chat_{mi}")
                tts_player(msg["content"], f"chat_{mi}")
        # Voice input for chat — with red recording indicator
        voice_col, photo_col, label_col = st.columns([1, 1, 3])
        with voice_col:
            st.markdown('<div class="mic-wrapper">', unsafe_allow_html=True)
            chat_audio = st.audio_input("🎤", key="chat_mic", label_visibility="collapsed")
            st.markdown(f'<div style="font-size:.75rem;color:#EF5350;font-weight:700;text-align:center;margin-top:4px">🎤 {T("start_recording")}</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)
        with photo_col:
            st.markdown('<div style="text-align:center;padding-top:6px">', unsafe_allow_html=True)
            _show_cam = st.toggle("📸", key="chat_cam_toggle", help=T("photo_hint"))
            _cam_color = "#4CAF50" if _show_cam else "#2B7DE9"
            _cam_label = "✅ Photo ON" if _show_cam else "📸 Photo"
            st.markdown(f'<div style="font-size:.75rem;color:{_cam_color};font-weight:700;text-align:center;margin-top:2px">{_cam_label}</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)
        with label_col:
            st.markdown(f'<div style="background:rgba(43,125,233,.08);border:1px solid rgba(43,125,233,.2);border-radius:10px;padding:10px 14px;margin-top:4px"><span style="font-size:.82rem;color:#7BB8F5;line-height:1.6">🎤 {T("mic_hint")}<br>📸 {T("photo_hint")}</span></div>', unsafe_allow_html=True)

        # Photo upload / camera area (shown when toggled on)
        chat_photo_b64 = None
        chat_photo_mime = None
        if _show_cam:
            st.markdown(f'<div style="background:rgba(76,175,80,.08);border:2px dashed rgba(76,175,80,.4);border-radius:12px;padding:4px 12px;margin-bottom:8px">', unsafe_allow_html=True)
            cam_tab1, cam_tab2 = st.tabs([f"📁 {T('photo_upload')}", f"📷 {T('photo_camera')}"])
            with cam_tab1:
                chat_upload = st.file_uploader(T("photo_upload"), type=["jpg","jpeg","png","heic","webp"], key="chat_photo_up", label_visibility="collapsed")
                if chat_upload:
                    st.image(chat_upload, caption="📸 Uploaded", use_container_width=True)
                    _img_bytes = chat_upload.read()
                    chat_photo_b64 = base64.b64encode(_img_bytes).decode()
                    chat_photo_mime = chat_upload.type if hasattr(chat_upload, 'type') and chat_upload.type else "image/jpeg"
                    chat_upload.seek(0)
            with cam_tab2:
                chat_cam = st.camera_input(T("photo_camera"), key="chat_cam_input", label_visibility="collapsed")
                if chat_cam:
                    _cam_bytes = chat_cam.read()
                    chat_photo_b64 = base64.b64encode(_cam_bytes).decode()
                    chat_photo_mime = "image/jpeg"
                    chat_cam.seek(0)
                    st.image(chat_cam, caption="📸 Captured", use_container_width=True)
            st.markdown('</div>', unsafe_allow_html=True)
            if chat_photo_b64:
                st.markdown(f'<div style="background:rgba(76,175,80,.1);border-radius:8px;padding:8px 12px;margin-top:6px"><span style="color:#66BB6A;font-size:.85rem;font-weight:600">✅ Photo ready — type a question or send with default prompt</span></div>', unsafe_allow_html=True)
                _photo_q = st.text_input(T("photo_ask"), value="", key="chat_photo_question", placeholder=T("photo_default"))

        voice_text = None
        if chat_audio:
            with st.spinner(T("transcribing")):
                voice_text = transcribe_audio(chat_audio.read())
            if voice_text:
                st.success(f'{T("heard")}: "{voice_text[:80]}..."' if len(voice_text) > 80 else f'{T("heard")}: "{voice_text}"')

        # Inject CSS to make mic red while recording
        st.markdown("""<style>
        /* Make ALL mic buttons red while recording (Streamlit changes aria-label to Stop) */
        [data-testid="stAudioInput"] > div > button[aria-label*="Stop"],
        [data-testid="stAudioInput"] > div > button[aria-label*="stop"],
        [data-testid="baseButton-minimal"][aria-label*="Stop"],
        [data-testid="baseButton-minimal"][aria-label*="stop"] {
            background-color: #D32F2F !important;
            color: white !important;
            border-color: #D32F2F !important;
            box-shadow: 0 0 12px rgba(211,47,47,.5) !important;
            animation: pulse-red 1.2s ease-in-out infinite !important;
        }
        @keyframes pulse-red {
            0%, 100% { box-shadow: 0 0 8px rgba(211,47,47,.4); }
            50% { box-shadow: 0 0 20px rgba(211,47,47,.7); }
        }
        </style>""", unsafe_allow_html=True)

        uq = st.chat_input(f"{T('ask_about')} {subject}... {T('draw_hint')}")
        # Voice input takes priority
        if voice_text:
            uq = voice_text
        # Photo submission — either with typed question or default prompt
        if chat_photo_b64 and not uq:
            _pq = st.session_state.get("chat_photo_question", "").strip()
            if _show_cam and st.button("📸 Ask Teacher Pehpeh about this photo", type="primary", key="photo_send"):
                uq = _pq if _pq else T("photo_default")
        if uq:
            # Check if there's a photo attached
            _attached_photo = chat_photo_b64 if (_show_cam and chat_photo_b64) else None
            _attached_mime = chat_photo_mime if _attached_photo else None
            _display_content = uq
            if _attached_photo:
                _display_content = f"📸 [Photo attached] {uq}"
            st.session_state.chat_messages.append({"role":"user","content":_display_content,"photo_b64":_attached_photo,"photo_mime":_attached_mime})
            img_keywords=["draw","illustrate","sketch","diagram","picture","image","visual","create an image","make an image","generate an image","show me","dessiner","dessine","diagramme","illustrer","chora","picha","mchoro"]
            want_chat_img=any(uq.lower().startswith(k) or k in uq.lower() for k in img_keywords) and not _attached_photo
            # Detect transcription request
            transcribe_keywords=["transcribe","transcription","handwriting","convert to doc","convert to word","notes to doc","write out","type up","type out","handwritten","transcrire","écriture","manuscrit","nakili","maandishi"]
            want_transcribe = _attached_photo and any(k in uq.lower() for k in transcribe_keywords)
            with st.status(T("thinking"),expanded=True) as status:
                if _attached_photo and want_transcribe:
                    st.write(T("transcribing_handwriting"))
                    r, m = best_vision("You are a precise handwriting transcription assistant.", TRANSCRIBE_PROMPT, _attached_photo, _attached_mime or "image/jpeg")
                    allr = {m: r} if m else {"AI": r}
                    # Generate .docx
                    _school_name = st.session_state.get("school_name", "")
                    _teacher_name = st.session_state.get("teacher_name", "")
                    _docx_bytes = generate_docx_from_text(
                        r, title="Transcribed Handwritten Notes",
                        school=_school_name, teacher=_teacher_name,
                        subject=_subj_en, grade=_grade_en
                    )
                    if _docx_bytes:
                        msg_data={"role":"assistant","content":f"📝 **Transcription Complete**\n\n{r}","model":m,"all_responses":allr,"docx_bytes":base64.b64encode(_docx_bytes).decode(),"docx_name":f"transcribed_notes_{_subj_en.lower().replace(' ','_')}.docx"}
                    else:
                        msg_data={"role":"assistant","content":f"📝 **Transcription** (Word doc generation failed — here is the text):\n\n{r}","model":m,"all_responses":allr}
                elif _attached_photo:
                    st.write(T("analyzing_photo"))
                    _vision_sp = build_free_chat()
                    r, m = best_vision(_vision_sp, uq, _attached_photo, _attached_mime or "image/jpeg")
                    allr = {m: r} if m else {"AI": r}
                    msg_data={"role":"assistant","content":r,"model":m,"all_responses":allr}
                else:
                    st.write(f"{T('asking_claude')} {T('asking_chatgpt')} {T('asking_gemini')}")
                    r,m,allr=best_all(build_free_chat(),uq,[{"role":x["role"],"content":x["content"]} for x in st.session_state.chat_messages[-11:-1]])
                    msg_data={"role":"assistant","content":r,"model":m,"all_responses":allr}
                if want_chat_img:
                    st.write(T("creating_img"))
                    img_url,img_model=gen_image(f"{_subj_en}: {uq} for {_grade_en} in {country}")
                    if img_url:
                        msg_data["image"]=img_url
                        msg_data["image_src"]=img_model
                status.update(label=T("response_ready"),state="complete",expanded=False)
            st.session_state.chat_messages.append(msg_data); st.rerun()
        if st.session_state.chat_messages and st.button(T("clear"),key="cc"): st.session_state.chat_messages=[]; st.rerun()

    # TAB 4: QUIZ (works offline — no internet or API keys needed)
    with t4:
        # Safe fallback: keys may not be in scope if user jumped directly to quiz tab
        try:
            _quiz_keys = keys
        except NameError:
            _quiz_keys = sum([bool(OPENAI_API_KEY), bool(ANTHROPIC_API_KEY), bool(GOOGLE_API_KEY)])
        # Mode toggle: adaptive quiz OR full WASSCE practice test
        # Show a live connectivity note — both modes work offline
        _quiz_online = (st.session_state.conn_info or {}).get("online", False)
        if not _quiz_online:
            st.markdown(
                '<div style="background:rgba(255,152,0,.1);border:1px solid #FF9800;border-radius:8px;' +
                'padding:7px 14px;margin-bottom:8px;font-size:.85rem;color:#FFB74D">' +
                '📴 <strong>Offline mode</strong> — Adaptive Quiz and WASSCE Practice Test both work fully offline. ' +
                'AI generation and chat need internet.</div>',
                unsafe_allow_html=True
            )
        _quiz_mode = st.radio("", ["📝 Adaptive Quiz", "🎯 WASSCE Practice Test"],
                              horizontal=True, key="quiz_mode_sel", label_visibility="collapsed")

        if _quiz_mode == "🎯 WASSCE Practice Test":
            # ── FULL WASSCE PRACTICE TEST (100% offline, no API needed) ──
            st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD}66;border-radius:10px;padding:10px 16px;margin-bottom:10px">' +
                        f'<strong style="color:{C_GOLD}">🎯 WASSCE Practice Test</strong> ' +
                        f'<span style="color:#D0D8E8;font-size:.85rem"> — Full exam simulation with answer sheet shading. Works offline. ✅</span></div>',
                        unsafe_allow_html=True)
            _pt_subj_display = st.selectbox("Subject:", _quiz_subjects(), key="pt_subj")
            _pt_subj = _quiz_subj_en(_pt_subj_display)
            _pt_bank_all = QUIZ.get(_pt_subj, {})
            # Gather all questions across levels
            _pt_all_qs = []
            for _lv in ["easy","medium","hard"]:
                _pt_all_qs.extend(_pt_bank_all.get(_lv, []))
            # Pick up to 20 questions for the practice test
            _pt_qs = _pt_all_qs[:20]
            _pt_n = len(_pt_qs)
            if _pt_n == 0:
                st.warning("No questions available for this subject yet.")
            else:
                import streamlit.components.v1 as _c1
                # Build question+options JSON for JS
                import json as _json
                _qs_json = _json.dumps([
                    {"q": q["q"], "o": q["o"], "a": q["a"],
                     "e": q.get("e",""), "t": q.get("t","")}
                    for q in _pt_qs
                ])
                _practice_html = f"""<!DOCTYPE html>
<html>
<head>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{background:#0a0e1a;font-family:Georgia,serif;padding:10px;color:#D0D8E8}}
.container{{display:grid;grid-template-columns:1fr 260px;gap:12px;align-items:start}}
@media(max-width:600px){{.container{{grid-template-columns:1fr}}}}
/* Questions panel */
.q-panel{{background:#0F2247;border-radius:10px;padding:14px;border:1px solid #D4A84333}}
.q-header{{color:#D4A843;font-weight:700;font-size:13px;margin-bottom:12px;letter-spacing:1px}}
.question{{margin-bottom:14px;padding:10px;background:rgba(255,255,255,.03);border-radius:8px;border-left:3px solid #333;transition:border-color .2s}}
.question.answered{{border-left-color:#D4A843}}
.question.correct{{border-left-color:#81C784;background:rgba(129,199,132,.07)}}
.question.wrong{{border-left-color:#EF5350;background:rgba(239,83,80,.07)}}
.qnum{{color:#D4A843;font-size:11px;font-weight:700;margin-bottom:4px}}
.qtext{{color:#E8EEF8;font-size:12px;line-height:1.5;margin-bottom:8px}}
.options{{display:flex;flex-direction:column;gap:4px}}
.opt{{display:flex;align-items:center;gap:8px;padding:5px 8px;border-radius:6px;
  cursor:pointer;border:1px solid #2a3a5a;transition:all .15s;font-size:11px;color:#B0BEC5}}
.opt:hover{{border-color:#D4A843;color:#D4A843;background:rgba(212,168,67,.06)}}
.opt.selected{{border-color:#D4A843;background:rgba(212,168,67,.12);color:#D4A843;font-weight:700}}
.opt.correct-ans{{border-color:#81C784;background:rgba(129,199,132,.15);color:#81C784;font-weight:700}}
.opt.wrong-ans{{border-color:#EF5350;background:rgba(239,83,80,.12);color:#EF9A9A}}
.opt-bubble{{width:16px;height:16px;border-radius:50%;border:1.5px solid currentColor;
  display:flex;align-items:center;justify-content:center;font-size:7px;font-weight:700;flex-shrink:0}}
.opt.selected .opt-bubble,.opt.correct-ans .opt-bubble{{background:currentColor}}
.explanation{{display:none;margin-top:8px;padding:7px 10px;background:rgba(212,168,67,.07);
  border-radius:6px;font-size:10.5px;color:#B0C8E8;line-height:1.5;border-left:2px solid #D4A843}}
.explanation.show{{display:block}}
.tip{{display:none;margin-top:4px;padding:6px 10px;background:rgba(43,125,233,.08);
  border-radius:6px;font-size:10px;color:#90CAF9;line-height:1.5}}
.tip.show{{display:block}}
/* Answer sheet panel */
.sheet-panel{{background:#0F2247;border-radius:10px;padding:12px;border:1px solid #D4A84333;position:sticky;top:10px}}
.sheet-header{{color:#D4A843;font-weight:700;font-size:12px;margin-bottom:10px;text-align:center;letter-spacing:1px}}
.sheet-grid{{display:grid;grid-template-columns:repeat(2,1fr);gap:2px 8px}}
.sheet-row{{display:flex;align-items:center;gap:3px;padding:1px 2px}}
.sq{{font-size:9px;color:#666;font-weight:700;width:16px;text-align:right;flex-shrink:0}}
.sb{{width:16px;height:16px;border-radius:50%;border:1.5px solid #555;
  display:flex;align-items:center;justify-content:center;font-size:6px;color:#777;
  cursor:pointer;transition:all .15s;flex-shrink:0}}
.sb:hover{{border-color:#D4A843;color:#D4A843}}
.sb.shaded{{background:#1a1a2e;border-color:#D4A843;color:transparent}}
.sb.correct-shade{{background:#81C784;border-color:#81C784}}
.sb.wrong-shade{{background:#EF5350;border-color:#EF5350}}
.score-bar{{margin-top:10px;padding:8px;background:rgba(212,168,67,.07);border-radius:6px;text-align:center}}
.score-num{{color:#81C784;font-size:16px;font-weight:700}}
.score-label{{color:#D4A843;font-size:10px;margin-top:2px}}
.submit-btn{{width:100%;margin-top:8px;padding:8px;background:#8B1A1A;color:#F5D98E;
  border:1px solid #D4A843;border-radius:6px;font-size:12px;font-weight:700;cursor:pointer;font-family:inherit}}
.submit-btn:hover{{background:#B22234}}
.submit-btn:disabled{{opacity:.4;cursor:not-allowed}}
.reset-btn{{width:100%;margin-top:6px;padding:6px;background:transparent;color:#D0D8E8;
  border:1px solid #3a4a6a;border-radius:6px;font-size:11px;cursor:pointer;font-family:inherit}}
.progress-bar{{height:4px;background:#1a2a3a;border-radius:2px;margin-bottom:10px;overflow:hidden}}
.progress-fill{{height:100%;background:linear-gradient(90deg,#8B1A1A,#D4A843);border-radius:2px;transition:width .3s}}
.result-banner{{display:none;padding:12px;border-radius:8px;text-align:center;margin-bottom:10px}}
.result-banner.show{{display:block}}
</style>
</head>
<body>
<div id="result-banner" class="result-banner"></div>
<div class="container">
  <!-- Questions -->
  <div class="q-panel">
    <div class="q-header">📝 {_pt_subj_display} — {_pt_n} Questions</div>
    <div class="progress-bar"><div class="progress-fill" id="prog" style="width:0%"></div></div>
    <div id="questions"></div>
  </div>
  <!-- Answer Sheet -->
  <div class="sheet-panel">
    <div class="sheet-header">📋 ANSWER SHEET</div>
    <div class="sheet-grid" id="sheet-grid"></div>
    <div class="score-bar">
      <div class="score-num" id="score-display">—</div>
      <div class="score-label">Score (after submit)</div>
    </div>
    <button class="submit-btn" id="submit-btn" onclick="submitTest()">✅ Submit & Mark</button>
    <button class="reset-btn" onclick="resetTest()">🔄 Reset</button>
  </div>
</div>

<script>
const QS = {_qs_json};
const N = QS.length;
const OPTS = ['A','B','C','D','E'];
const selected = new Array(N).fill(null); // selected option index per question
let submitted = false;

// Build question panel
const qContainer = document.getElementById('questions');
QS.forEach((q,i) => {{
  const div = document.createElement('div');
  div.className = 'question'; div.id = 'q_'+i;
  const opts = q.o.map((opt,j) => {{
    const letter = OPTS[j] || String.fromCharCode(65+j);
    return `<div class="opt" id="opt_${{i}}_${{j}}" onclick="pick(${{i}},${{j}})">
      <div class="opt-bubble">${{letter}}</div>
      <span>${{opt}}</span>
    </div>`;
  }}).join('');
  div.innerHTML = `
    <div class="qnum">Q${{i+1}}</div>
    <div class="qtext">${{q.q}}</div>
    <div class="options">${{opts}}</div>
    <div class="explanation" id="exp_${{i}}">📖 ${{q.e}}</div>
    <div class="tip" id="tip_${{i}}">🧑‍🏫 ${{q.t}}</div>
  `;
  qContainer.appendChild(div);
}});

// Build answer sheet
const sheetGrid = document.getElementById('sheet-grid');
QS.forEach((q,i) => {{
  const row = document.createElement('div');
  row.className = 'sheet-row';
  const bubbles = q.o.map((_,j) => {{
    const letter = OPTS[j] || String.fromCharCode(65+j);
    return `<div class="sb" id="sb_${{i}}_${{j}}" onclick="sheetPick(${{i}},${{j}})">${{letter}}</div>`;
  }}).join('');
  row.innerHTML = `<span class="sq">${{i+1}}.</span>${{bubbles}}`;
  sheetGrid.appendChild(row);
}});

function pick(qi, oi) {{
  if (submitted) return;
  // Deselect previous
  if (selected[qi] !== null) {{
    document.getElementById('opt_'+qi+'_'+selected[qi]).classList.remove('selected');
    document.getElementById('sb_'+qi+'_'+selected[qi]).classList.remove('shaded');
    document.getElementById('sb_'+qi+'_'+selected[qi]).textContent = OPTS[selected[qi]] || String.fromCharCode(65+selected[qi]);
  }}
  selected[qi] = oi;
  document.getElementById('opt_'+qi+'_'+oi).classList.add('selected');
  // Shade answer sheet bubble
  const sb = document.getElementById('sb_'+qi+'_'+oi);
  sb.classList.add('shaded'); sb.textContent = '';
  document.getElementById('q_'+qi).classList.add('answered');
  // Update progress
  const answered = selected.filter(s => s !== null).length;
  document.getElementById('prog').style.width = (answered/N*100)+'%';
  document.getElementById('submit-btn').disabled = answered < N;
}}

function sheetPick(qi, oi) {{
  if (submitted) return;
  pick(qi, oi);
  // Scroll question into view
  document.getElementById('q_'+qi).scrollIntoView({{behavior:'smooth',block:'nearest'}});
}}

function submitTest() {{
  if (submitted) return;
  submitted = true;
  let score = 0;
  QS.forEach((q,i) => {{
    const correct = q.a;
    const userPick = selected[i];
    const qDiv = document.getElementById('q_'+i);
    // Colour answer sheet bubbles
    if (userPick !== null) {{
      const sb = document.getElementById('sb_'+i+'_'+userPick);
      sb.classList.remove('shaded');
      if (userPick === correct) {{ sb.classList.add('correct-shade'); score++; }}
      else {{ sb.classList.add('wrong-shade'); }}
    }}
    // Colour correct answer sheet bubble green
    document.getElementById('sb_'+i+'_'+correct).classList.add('correct-shade');
    // Colour options
    if (userPick === correct) {{
      document.getElementById('opt_'+i+'_'+userPick).classList.add('correct-ans');
      qDiv.classList.remove('answered'); qDiv.classList.add('correct');
    }} else {{
      if (userPick !== null) document.getElementById('opt_'+i+'_'+userPick).classList.add('wrong-ans');
      document.getElementById('opt_'+i+'_'+correct).classList.add('correct-ans');
      qDiv.classList.remove('answered'); qDiv.classList.add('wrong');
    }}
    // Show explanation
    document.getElementById('exp_'+i).classList.add('show');
    document.getElementById('tip_'+i).classList.add('show');
  }});
  const pct = Math.round(score/N*100);
  document.getElementById('score-display').textContent = score+'/'+N+' ('+pct+'%)';
  const banner = document.getElementById('result-banner');
  if (pct >= 75) {{
    banner.style.background='rgba(129,199,132,.15)'; banner.style.border='1px solid #81C784';
    banner.innerHTML = '<strong style="color:#81C784">🎉 '+pct+'% — Excellent!</strong><br><span style="font-size:11px;color:#B0C8E8">Great performance. Review explanations below.</span>';
  }} else if (pct >= 50) {{
    banner.style.background='rgba(212,168,67,.12)'; banner.style.border='1px solid #D4A843';
    banner.innerHTML = '<strong style="color:#D4A843">📚 '+pct+'% — Keep Studying</strong><br><span style="font-size:11px;color:#B0C8E8">Review the explanations and teacher tips.</span>';
  }} else {{
    banner.style.background='rgba(139,26,26,.2)'; banner.style.border='1px solid #EF5350';
    banner.innerHTML = '<strong style="color:#EF9A9A">💪 '+pct+'% — More Practice Needed</strong><br><span style="font-size:11px;color:#B0C8E8">Focus on the teacher tips — they show exam technique.</span>';
  }}
  banner.classList.add('show');
}}

function resetTest() {{
  submitted = false;
  selected.fill(null);
  document.getElementById('prog').style.width = '0%';
  document.getElementById('score-display').textContent = '—';
  document.getElementById('result-banner').className = 'result-banner';
  document.getElementById('submit-btn').disabled = true;
  QS.forEach((q,i) => {{
    document.getElementById('q_'+i).className = 'question';
    document.getElementById('exp_'+i).className = 'explanation';
    document.getElementById('tip_'+i).className = 'tip';
    q.o.forEach((_,j) => {{
      const opt = document.getElementById('opt_'+i+'_'+j);
      opt.className = 'opt';
      const sb = document.getElementById('sb_'+i+'_'+j);
      sb.className = 'sb';
      sb.textContent = OPTS[j] || String.fromCharCode(65+j);
    }});
  }});
}}

// Disable submit until all answered
document.getElementById('submit-btn').disabled = true;
</script>
</body>
</html>"""
                _c1.html(_practice_html, height=700, scrolling=True)

        else:
            # ── ADAPTIVE QUIZ — always available, no internet needed ──
            st.markdown(
                f'<div style="background:rgba(43,125,233,.08);border:1px solid {C_BLUE};border-radius:12px;padding:10px 18px;margin-bottom:8px">' +
                f'{ico(16)} <strong style="color:{C_BLUE}">{T("practice_quiz")}</strong> ' +
                f'<span style="color:#7BB8F5;font-size:.85rem">— Adaptive · ✅ Works offline</span></div>',
                unsafe_allow_html=True
            )

            # ── Subject selector + manual level toggle ──
            _sub_col, _lv_col = st.columns([2, 3])
            with _sub_col:
                qsub_display=st.selectbox(f"{T('subject')}:",_quiz_subjects(),key="qs",label_visibility="collapsed",
                                          format_func=lambda x: f"📚 {x}")
            qsub=_quiz_subj_en(qsub_display)
            qs=st.session_state[f"qz_{qsub}"]

            with _lv_col:
                # Three level buttons — active one highlighted gold
                _LV_LABELS = {"easy": "🟢 Easy", "medium": "🟡 Medium", "hard": "🔴 Hard"}
                _lv_b1, _lv_b2, _lv_b3 = st.columns(3)
                for _lv_btn, _lv_col_obj in [("easy", _lv_b1), ("medium", _lv_b2), ("hard", _lv_b3)]:
                    _is_active = qs["lv"] == _lv_btn
                    with _lv_col_obj:
                        if st.button(
                            _LV_LABELS[_lv_btn],
                            key=f"lvbtn_{qsub}_{_lv_btn}",
                            use_container_width=True,
                            type="primary" if _is_active else "secondary",
                            help=f"Jump to {_lv_btn} questions (adaptive engine continues from here)"
                        ):
                            if not _is_active:
                                qs["lv"] = _lv_btn
                                qs["manual_lv"] = True   # flag: user set this manually
                                qs["qi"] = 0
                                qs["done"] = False
                                qs["sel"] = None
                                st.toast(f"Switched to {_lv_btn.upper()} — adaptive engine active from here 🎯")
                                st.rerun()

            # ── Adaptive mode badge ──
            bank=QUIZ[qsub]; lv=qs["lv"]
            questions=bank.get(lv,bank["easy"]); qi=qs["qi"]%len(questions); cur=questions[qi]
            pct=f"{round(qs['sc']/qs['tot']*100)}%" if qs["tot"] else "—"
            stk=f"🔥 {qs['stk']} {T('streak')}!" if qs["stk"]>=3 else ""
            _manual_flag = qs.get("manual_lv", False)
            _mode_badge = (
                f'<span style="background:rgba(212,168,67,.15);border:1px solid {C_GOLD}55;border-radius:10px;padding:1px 8px;font-size:.72rem;color:{C_GOLD}">🎯 Auto-adaptive</span>'
                if not _manual_flag else
                f'<span style="background:rgba(43,125,233,.12);border:1px solid {C_BLUE}55;border-radius:10px;padding:1px 8px;font-size:.72rem;color:#7BB8F5">🔧 Manual: {lv.upper()}</span>'
            )
            st.markdown(
                f'<div class="qsc" style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:6px">' +
                f'<span>{ico(16)} {T("score")}: <strong>{qs["sc"]}/{qs["tot"]}</strong> ({pct}) · {T("level")}: <strong>{lv.upper()}</strong> {stk}</span>' +
                f'{_mode_badge}</div>',
                unsafe_allow_html=True
            )
            st.markdown(f'<div class="qbox"><strong style="color:white">Q{qs["tot"]+1}:</strong><br><span style="color:#D0D8E8;line-height:1.6">{cur["q"]}</span></div>',unsafe_allow_html=True)

            if not qs["done"]:
                cols=st.columns(2)
                for j,opt in enumerate(cur["o"]):
                    with cols[j%2]:
                        if st.button(f"{'ABCD'[j]}) {opt}",key=f"qo_{qsub}_{qs['tot']}_{j}",use_container_width=True):
                            qs["sel"]=j; qs["done"]=True; qs["tot"]+=1
                            if j==cur["a"]: qs["sc"]+=1; qs["stk"]+=1
                            else: qs["stk"]=0
                            qs["hist"].append({"c":j==cur["a"],"lv":lv}); st.rerun()
            else:
                ok=qs["sel"]==cur["a"]
                if ok: st.markdown(f'<div class="qok"><strong>{random.choice(PRAISE)}</strong><br>✅ {cur["o"][cur["a"]]} is correct!</div>',unsafe_allow_html=True)
                else: st.markdown(f'<div class="qno"><strong>{random.choice(ENCOURAGE)}</strong><br>Answer: <strong>{cur["o"][cur["a"]]}</strong></div>',unsafe_allow_html=True)
                st.markdown(f'<div style="background:rgba(212,168,67,.08);border:1px solid {C_GOLD};border-radius:10px;padding:12px 16px;margin:8px 0"><strong style="color:{C_GOLD}">📖 Explanation:</strong><br><span style="color:#D0D8E8">{cur["e"]}</span></div>',unsafe_allow_html=True)
                st.markdown(f'<div class="qtip"><strong>🧑‍🏫 Teacher Tip:</strong> {cur["t"]}</div>',unsafe_allow_html=True)
                recent=[h for h in qs["hist"][-5:]]; rc=sum(1 for h in recent if h["c"])
                if st.button(T("next"),type="primary",key=f"nx_{qsub}_{qs['tot']}",use_container_width=True):
                    if len(recent)>=3:
                        if rc>=4 and lv!="hard":
                            qs["lv"]="medium" if lv=="easy" else "hard"
                            qs["manual_lv"]=False   # adaptive took over
                            st.toast(f"🏆 Level UP → {qs['lv'].upper()}")
                        elif rc<=1 and lv!="easy":
                            qs["lv"]="easy" if lv=="medium" else "medium"
                            qs["manual_lv"]=False
                            st.toast(f"Adjusting → {qs['lv'].upper()}")
                    qs["qi"]=(qi+1)%len(bank.get(qs["lv"],bank["easy"])); qs["done"]=False; qs["sel"]=None; st.rerun()

            st.markdown("---")
            r1,r2=st.columns(2)
            with r1:
                if st.button(T("reset"),key=f"rst_{qsub}"):
                    st.session_state[f"qz_{qsub}"]={"lv":"easy","qi":0,"sc":0,"tot":0,"stk":0,"done":False,"sel":None,"hist":[],"manual_lv":False}
                    st.rerun()
            with r2:
                if st.button(T("wassce_tips"),key="wt"): st.markdown(f'<div style="background:{C_NAVY_L};border:1px solid {C_GOLD};border-radius:12px;padding:16px;color:#D0D8E8;white-space:pre-wrap;line-height:1.7">{WASSCE_TIPS}</div>',unsafe_allow_html=True)

    st.markdown(f'<div class="ft">{ico(16)} <strong>Teacher Pehpeh by IBT</strong><br>Built by <strong>Rodney L. Bollie, PhD</strong> · <a href="https://www.institutebasictechnology.org">Institute of Basic Technology</a><br><a href="https://www.institutebasictechnology.org/index.php" style="color:{C_BLUE}">Visit our website →</a></div>',unsafe_allow_html=True)


@st.dialog("📋 WASSCE Answer Sheet Guide", width="large")
def wassce_shading_modal():
    _wtab = st.radio("", ["🎯 Shading Guide","❌ Common Errors","📋 Practice Sheet","💡 Exam Tips"],
                     horizontal=True, key="wassce_modal_tab", label_visibility="collapsed")

    if _wtab == "🎯 Shading Guide":
        st.markdown('''
<div style="background:rgba(212,168,67,.07);border:1px solid #D4A84344;border-radius:14px;padding:20px;margin:8px 0">
  <div style="text-align:center;color:#D4A843;font-weight:700;font-size:1rem;margin-bottom:16px">✅ The 3 Rules of Correct Shading</div>
  <div style="display:flex;flex-direction:column;gap:10px">
    <div style="display:flex;align-items:center;gap:12px;background:rgba(212,168,67,.06);border-left:3px solid #D4A843;border-radius:8px;padding:10px 14px">
      <div style="width:30px;height:30px;border-radius:50%;background:#8B1A1A;color:#D4A843;display:flex;align-items:center;justify-content:center;font-weight:900;font-size:15px;flex-shrink:0">1</div>
      <div><strong style="color:#D4A843">Fill completely</strong><br><span style="color:#D0D8E8;font-size:.85rem">Shade the entire bubble — edge to edge. No white space showing inside the circle.</span></div>
    </div>
    <div style="display:flex;align-items:center;gap:12px;background:rgba(212,168,67,.06);border-left:3px solid #D4A843;border-radius:8px;padding:10px 14px">
      <div style="width:30px;height:30px;border-radius:50%;background:#8B1A1A;color:#D4A843;display:flex;align-items:center;justify-content:center;font-weight:900;font-size:15px;flex-shrink:0">2</div>
      <div><strong style="color:#D4A843">One bubble per question</strong><br><span style="color:#D0D8E8;font-size:.85rem">Shade exactly ONE option. Two shaded bubbles = zero marks automatically.</span></div>
    </div>
    <div style="display:flex;align-items:center;gap:12px;background:rgba(212,168,67,.06);border-left:3px solid #D4A843;border-radius:8px;padding:10px 14px">
      <div style="width:30px;height:30px;border-radius:50%;background:#8B1A1A;color:#D4A843;display:flex;align-items:center;justify-content:center;font-weight:900;font-size:15px;flex-shrink:0">3</div>
      <div><strong style="color:#D4A843">Erase cleanly if changing</strong><br><span style="color:#D0D8E8;font-size:.85rem">Use a proper eraser. Ghost marks from incomplete erasure trigger a scanning error.</span></div>
    </div>
  </div>
</div>
<div style="background:rgba(15,34,71,.8);border:1px solid #D4A84344;border-radius:14px;padding:20px;margin:8px 0">
  <div style="text-align:center;color:#D4A843;font-weight:700;font-size:.95rem;margin-bottom:14px">✏️ The 3 Stages of Correct Shading</div>
  <div style="display:flex;justify-content:space-around;align-items:center;flex-wrap:wrap;gap:16px">
    <div style="text-align:center">
      <div style="width:44px;height:44px;border-radius:50%;border:2.5px solid #555;margin:0 auto 8px"></div>
      <div style="color:#EF9A9A;font-size:.75rem;font-weight:700">Empty — NO marks</div>
    </div>
    <div style="text-align:center">
      <div style="width:44px;height:44px;border-radius:50%;border:2.5px solid #888;background:linear-gradient(180deg,transparent 50%,#888 50%);margin:0 auto 8px"></div>
      <div style="color:#FFB74D;font-size:.75rem;font-weight:700">Half — Won't scan</div>
    </div>
    <div style="text-align:center">
      <div style="width:44px;height:44px;border-radius:50%;border:2.5px solid #D4A843;background:#1a1a2e;margin:0 auto 8px"></div>
      <div style="color:#81C784;font-size:.75rem;font-weight:700">Full — CORRECT ✅</div>
    </div>
  </div>
</div>
''', unsafe_allow_html=True)

    elif _wtab == "❌ Common Errors":
        mistakes = [
            ("✓", "22px", "Tick / Checkmark", "❌", "#EF9A9A", "rgba(139,26,26,.2)", "#8B1A1A44",
             "Most common mistake in Liberian schools — the scanner does not recognise a tick"),
            ("✗", "22px", "X Mark", "❌", "#EF9A9A", "rgba(139,26,26,.2)", "#8B1A1A44",
             "Same problem as a tick — not detected by optical scanner"),
            ("·", "28px", "Dot only", "❌", "#EF9A9A", "rgba(139,26,26,.2)", "#8B1A1A44",
             "Too small and too light — the scanner will read it as empty"),
            ("◑", "22px", "Partial shading", "❌", "#EF9A9A", "rgba(139,26,26,.2)", "#8B1A1A44",
             "Half-filled bubbles may or may not scan — too risky. Always complete the fill"),
        ]
        for icon, sz, label, badge, col, bg, bdr, desc in mistakes:
            st.markdown(
                f'<div style="display:flex;align-items:center;gap:14px;background:{bg};border:1px solid {bdr};border-radius:10px;padding:12px 16px;margin-bottom:6px">' +
                f'<div style="width:42px;height:42px;border-radius:50%;border:2px solid #c0392b;display:flex;align-items:center;justify-content:center;font-size:{sz};color:#aaa;flex-shrink:0">{icon}</div>' +
                f'<div><div style="color:{col};font-weight:700;font-size:.88rem">{badge} {label}</div><div style="color:#D0D8E8;font-size:.8rem;margin-top:2px">{desc}</div></div></div>',
                unsafe_allow_html=True
            )
        st.markdown(
            '<div style="display:flex;align-items:center;gap:14px;background:rgba(129,199,132,.1);border:1px solid #81C78444;border-radius:10px;padding:12px 16px;margin-bottom:6px">' +
            '<div style="width:42px;height:42px;border-radius:50%;border:2px solid #D4A843;background:#1a1a2e;flex-shrink:0"></div>' +
            '<div><div style="color:#81C784;font-weight:700;font-size:.88rem">✅ Fully shaded bubble</div><div style="color:#D0D8E8;font-size:.8rem;margin-top:2px">Complete fill from edge to edge using HB pencil — the only method the scanner accepts</div></div></div>',
            unsafe_allow_html=True
        )
        st.error("🚨 **#1 mistake in Liberian schools:** Students shade a tick (✓) instead of filling the bubble. Drill this until it's automatic — *fill, don't tick.*")

    elif _wtab == "📋 Practice Sheet":
        st.info("💡 **For teachers:** Click any bubble to shade it — just like the real exam. Click again to unshade (erase). Two shaded bubbles = zero marks warning shown automatically.")
        import streamlit.components.v1 as _components
        _sheet_html = """
<!DOCTYPE html>
<html>
<head>
<style>
  * { box-sizing: border-box; margin: 0; padding: 0; }
  body { background: #0a0e1a; font-family: 'Courier New', monospace; padding: 12px; }
  .header { background: #0F2247; border-radius: 8px; padding: 10px; text-align: center; margin-bottom: 12px; border: 1px solid #D4A84344; }
  .header-title { color: #D4A843; font-weight: 700; font-size: 13px; letter-spacing: 2px; }
  .header-sub { color: #D0D8E8; font-size: 10px; margin-top: 3px; opacity: .8; }
  .grid { display: grid; grid-template-columns: repeat(2, 1fr); gap: 2px 16px; }
  .row { display: flex; align-items: center; gap: 5px; padding: 2px 4px; border-radius: 4px; transition: background .15s; }
  .row:hover { background: rgba(212,168,67,.05); }
  .qnum { width: 22px; text-align: right; font-size: 10px; color: #666; font-weight: 700; flex-shrink: 0; }
  .bubble {
    width: 22px; height: 22px; border-radius: 50%;
    border: 1.5px solid #999;
    display: flex; align-items: center; justify-content: center;
    font-size: 8px; color: #888; font-weight: 700;
    cursor: pointer; transition: all .18s; user-select: none;
    flex-shrink: 0;
  }
  .bubble:hover { border-color: #D4A843; color: #D4A843; transform: scale(1.15); }
  .bubble.shaded {
    background: #1a1a2e; border-color: #D4A843;
    color: transparent; box-shadow: inset 0 0 0 3px #1a1a2e;
  }
  .bubble.double-error {
    background: #8B1A1A; border-color: #EF5350;
    animation: pulse-err .6s ease-in-out infinite alternate;
  }
  @keyframes pulse-err { from { box-shadow: 0 0 4px #EF535088; } to { box-shadow: 0 0 10px #EF5350cc; } }
  .warning { display:none; background:rgba(139,26,26,.25); border:1px solid #EF535088;
    border-radius:6px; padding:6px 10px; font-size:10px; color:#EF9A9A;
    margin-bottom:8px; text-align:center; }
  .warning.show { display:block; }
  .score-bar { background:#0F2247; border-radius:6px; padding:7px 12px; margin-top:10px;
    display:flex; justify-content:space-between; align-items:center; flex-wrap:wrap; gap:6px; }
  .score-label { color:#D4A843; font-size:11px; font-weight:700; }
  .score-val { color:#81C784; font-size:12px; font-weight:700; }
  .clear-btn { background:#8B1A1A; color:#F5D98E; border:1px solid #D4A843;
    border-radius:6px; padding:4px 12px; font-size:11px; font-weight:700;
    cursor:pointer; font-family:inherit; }
  .clear-btn:hover { background:#B22234; }
  .footer { margin-top:8px; font-size:9px; color:#555; text-align:center; }
</style>
</head>
<body>
<div class="header">
  <div class="header-title">WASSCE INTERACTIVE ANSWER SHEET</div>
  <div class="header-sub">Click a bubble to shade it &nbsp;·&nbsp; Click again to erase &nbsp;·&nbsp; Shade ONE per question</div>
</div>
<div id="warning" class="warning">⚠️ Two shaded bubbles detected! That question scores ZERO marks — erase one.</div>
<div class="grid" id="sheet"></div>
<div class="score-bar">
  <span class="score-label">Questions answered: <span id="count" class="score-val">0 / 60</span></span>
  <button class="clear-btn" onclick="clearAll()">🗑️ Clear Sheet</button>
</div>
<div class="footer">⚠️ Real exam uses HB pencil only · Erase completely · Do not fold sheet</div>

<script>
  const options = ['A','B','C','D','E'];
  const state = {}; // q -> Set of selected options
  const sheet = document.getElementById('sheet');
  const warning = document.getElementById('warning');
  const countEl = document.getElementById('count');

  for (let q = 1; q <= 60; q++) {
    state[q] = new Set();
    const row = document.createElement('div');
    row.className = 'row';
    const qnum = document.createElement('span');
    qnum.className = 'qnum';
    qnum.textContent = q + '.';
    row.appendChild(qnum);
    options.forEach(opt => {
      const b = document.createElement('div');
      b.className = 'bubble';
      b.textContent = opt;
      b.id = 'b_' + q + '_' + opt;
      b.onclick = () => toggle(q, opt);
      row.appendChild(b);
    });
    sheet.appendChild(row);
  }

  function toggle(q, opt) {
    const b = document.getElementById('b_' + q + '_' + opt);
    if (state[q].has(opt)) {
      state[q].delete(opt);
      b.classList.remove('shaded', 'double-error');
      b.textContent = opt;
    } else {
      state[q].add(opt);
      b.classList.add('shaded');
      b.textContent = '';
    }
    // Check for double-shade error
    let hasError = false;
    for (let qq = 1; qq <= 60; qq++) {
      const isDouble = state[qq].size > 1;
      options.forEach(o => {
        const el = document.getElementById('b_' + qq + '_' + o);
        if (isDouble && state[qq].has(o)) {
          el.classList.add('double-error');
        } else if (!isDouble) {
          el.classList.remove('double-error');
        }
      });
      if (isDouble) hasError = true;
    }
    warning.className = 'warning' + (hasError ? ' show' : '');
    // Update count
    const answered = Object.values(state).filter(s => s.size === 1).length;
    countEl.textContent = answered + ' / 60';
  }

  function clearAll() {
    for (let q = 1; q <= 60; q++) {
      state[q].clear();
      options.forEach(opt => {
        const b = document.getElementById('b_' + q + '_' + opt);
        b.classList.remove('shaded', 'double-error');
        b.textContent = opt;
      });
    }
    warning.className = 'warning';
    countEl.textContent = '0 / 60';
  }
</script>
</body>
</html>
"""
        _components.html(_sheet_html, height=780, scrolling=True)

    elif _wtab == "💡 Exam Tips":
        tips = [
            ("✏️", "Use HB Pencil Only", "Never use pen or biro. The optical scanner reads HB pencil marks only. Ballpoint pen cannot be erased cleanly."),
            ("⭕", "Fill the Entire Bubble", "Shade completely from edge to edge. Ticks, dots, X marks, and partial shading all score zero."),
            ("🧹", "Erase Completely", "If you change an answer, erase the old bubble fully before shading the new one. Two shaded bubbles = zero marks."),
            ("🔢", "Check Row Alignment Every 10 Questions", "A single misaligned row can cost 10+ marks instantly. Every 10 questions, verify your row number matches."),
            ("⏱️", "The 5-Minute Rule", "With 5 minutes left, stop answering new questions. Use that time to check every row on your answer sheet."),
            ("📏", "Keep the Sheet Clean and Flat", "No stray marks. Don't fold or bend the sheet. Creases can interfere with scanner reading."),
        ]
        for icon, title, body in tips:
            st.markdown(
                f'<div style="display:flex;gap:12px;align-items:flex-start;background:rgba(255,255,255,.04);' +
                f'border:1px solid rgba(255,255,255,.08);border-radius:10px;padding:12px 14px;margin-bottom:8px">' +
                f'<div style="font-size:22px;width:36px;height:36px;display:flex;align-items:center;justify-content:center;' +
                f'background:rgba(212,168,67,.1);border-radius:8px;flex-shrink:0">{icon}</div>' +
                f'<div><div style="color:#D4A843;font-weight:700;font-size:.9rem">{title}</div>' +
                f'<div style="color:#D0D8E8;font-size:.83rem;margin-top:3px;line-height:1.5">{body}</div></div></div>',
                unsafe_allow_html=True
            )


if __name__=="__main__": main()
